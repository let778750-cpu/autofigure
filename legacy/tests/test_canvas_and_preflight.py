from __future__ import annotations

import copy
import hashlib
import json
import math
import sys
from pathlib import Path
from typing import Callable

import numpy as np
import pytest
from PIL import Image, ImageDraw
from pptx import Presentation


PROJECT_ROOT = Path(__file__).resolve().parents[1]
TOOLS_DIR = PROJECT_ROOT / "tools"
sys.path.insert(0, str(TOOLS_DIR))

from create_canvas_pptx import create_blank_canvas_pptx  # noqa: E402
from geometry_refinement import run_geometry_refinement  # noqa: E402
from powerpoint_native_math import compile_formula  # noqa: E402
from preflight_scene import main as preflight_main  # noqa: E402
from preflight_scene import preflight_scene, resolve_font_path  # noqa: E402
from tests.test_perception_review import adapter as review_adapter  # noqa: E402
from tests.test_perception_review import build_raw_manifest  # noqa: E402
from validate_host_runtime import validate_runtime  # noqa: E402


def make_png(path: Path, size: tuple[int, int] = (800, 600)) -> Path:
    Image.new("RGB", size, "white").save(path, format="PNG")
    return path


def digest(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def write_json(path: Path, payload: dict) -> None:
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def native_formula_record(
    directory: Path,
    *,
    formula_id: str,
    element_id: str,
    canonical_latex: str,
    mode: str = "display",
) -> dict:
    latex_sha256 = hashlib.sha256(canonical_latex.encode("utf-8")).hexdigest()
    receipt_path = directory / f"{formula_id}-converter-receipt.json"
    receipt = compile_formula(formula_id, canonical_latex, mode)
    write_json(receipt_path, receipt)
    return {
        "id": formula_id,
        "element_id": element_id,
        "canonical_latex": canonical_latex,
        "latex_sha256": latex_sha256,
        "mode": mode,
        "render_kind": "native_office_math",
        "fallback_policy": "strict_no_raster_no_svg",
        "converter_receipt_path": str(receipt_path),
        "converter_receipt_sha256": digest(receipt_path),
        "source_evidence": ["source_text"],
        "disposition": "CONFIRMED",
    }


def mutate_review_receipt(spec: dict, mutation: Callable[[dict], None]) -> dict:
    receipt_path = Path(spec["perception"]["review_receipt_path"])
    receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
    mutation(receipt)
    write_json(receipt_path, receipt)
    spec["perception"]["review_receipt_sha256"] = digest(receipt_path)
    return receipt


def refinalize_review(
    spec: dict,
    *,
    manifest_mutation: Callable[[dict], None] | None = None,
) -> tuple[dict, dict]:
    """Rewrite a schema-valid raw manifest, then produce a new genuine PASS receipt."""
    manifest_path = Path(spec["perception"]["manifest_path"])
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    if manifest_mutation is not None:
        manifest_mutation(manifest)
    write_json(manifest_path, manifest)
    spec["perception"]["manifest_sha256"] = digest(manifest_path)

    decisions_path = manifest_path.with_name("perception-review-decisions.json")
    decisions = review_adapter.initialize_review(manifest_path, decisions_path)
    for decision in decisions["decisions"]:
        decision["evidence"] = {
            "kind": "user_confirmed",
            "detail": "Adversarial test fixture reviewed against its frozen source.",
        }
        if decision["formula_like"]:
            decision["status"] = "FORMULA_CONFIRMED"
            decision["authoritative_latex"] = "x^2+y^2"
        else:
            decision["status"] = "CONFIRMED"
            decision["confirmed_text"] = decision["ocr_text"]
    write_json(decisions_path, decisions)

    receipt_path = Path(spec["perception"]["review_receipt_path"])
    receipt, exit_code = review_adapter.finalize_review(
        manifest_path,
        decisions_path,
        receipt_path,
    )
    assert exit_code == 0
    assert receipt["status"] == "PERCEPTION_REVIEW_PASS"
    spec["perception"]["review_receipt_sha256"] = digest(receipt_path)
    return manifest, receipt


def set_candidate_geometry(
    manifest: dict,
    candidate_id: str,
    bbox: dict[str, float],
) -> None:
    candidate = next(
        item for item in manifest["text_candidates"] if item["candidate_id"] == candidate_id
    )
    candidate["bbox_source"] = dict(bbox)
    candidate["bbox_envelope_source"] = dict(bbox)
    x, y, width, height = (bbox[key] for key in ("x", "y", "w", "h"))
    candidate["polygon_source"] = [
        [x, y],
        [x + width, y],
        [x + width, y + height],
        [x, y + height],
    ]


def edge(
    edge_id: str,
    source_id: str,
    target_id: str,
    *,
    source_anchor: str = "right",
    target_anchor: str = "left",
    route: str = "straight",
    via: list[dict[str, float]] | None = None,
) -> dict:
    result = {
        "id": edge_id,
        "from": source_id,
        "to": target_id,
        "source_anchor": source_anchor,
        "target_anchor": target_anchor,
        "route": route,
        "meaning": "data_flow",
        "clearance_px": 2,
        "allowed_crossings": [],
        "allowed_edge_crossings": [],
    }
    if via is not None:
        result["via"] = via
    return result


def test_create_canvas_matches_measured_png_aspect(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png", (1600, 900))
    output = tmp_path / "canvas.pptx"

    report = create_blank_canvas_pptx(source, output)
    presentation = Presentation(output)

    assert report["status"] == "PASS"
    assert report["source"]["sha256"] == digest(source)
    assert report["output_pptx_sha256"] == digest(output)
    assert len(presentation.slides) == 1
    assert len(presentation.slides[0].shapes) == 0
    assert presentation.slide_width / presentation.slide_height == pytest.approx(
        1600 / 900, rel=1e-6
    )


def usable_font() -> Path:
    resolved = resolve_font_path("Arial")
    if resolved is None:
        pytest.skip("No exact Arial font is installed for Pillow metric testing.")
    return resolved


def base_spec(source: Path, *, elements: list[dict]) -> dict:
    with Image.open(source) as image:
        width, height = image.size
        pixel_mode = image.mode
    source_sha = digest(source)
    raw_manifest = build_raw_manifest()
    raw_manifest["run_id"] = "preflight-review-run-001"
    raw_manifest["source"].update(
        {
            "path": str(source.resolve()),
            "sha256": source_sha.upper(),
            "size_bytes": source.stat().st_size,
            "width_px": width,
            "height_px": height,
            "pixel_mode": pixel_mode,
            "format": "PNG",
        }
    )

    evidence_dir = source.parent / "perception-evidence"
    evidence_dir.mkdir(exist_ok=True)
    config_path = evidence_dir / "ocr-config.json"
    config_path.write_text("{}\n", encoding="utf-8")
    raw_schema_path = PROJECT_ROOT / "schemas" / "perception-manifest.schema.json"
    raw_manifest["configuration"].update(
        {
            "path": str(config_path),
            "sha256": review_adapter.sha256_file(config_path),
            "manifest_schema_path": str(raw_schema_path),
            "manifest_schema_sha256": review_adapter.sha256_file(raw_schema_path),
        }
    )
    script_path = evidence_dir / "perception-script.py"
    script_path.write_text("# frozen test script\n", encoding="utf-8")
    raw_manifest["scripts"] = [
        {
            "path": str(script_path),
            "relative_path": "tools/perception-script.py",
            "size_bytes": script_path.stat().st_size,
            "sha256": review_adapter.sha256_file(script_path),
        }
    ]
    for role, model in raw_manifest["models"].items():
        model["role"] = role
        artifacts = []
        for index in range(3):
            model_path = evidence_dir / f"{role}-{index}.bin"
            model_path.write_bytes(f"{role}:{index}".encode())
            artifacts.append(
                {
                    "filename": model_path.name,
                    "path": str(model_path),
                    "size_bytes": model_path.stat().st_size,
                    "sha256": review_adapter.sha256_file(model_path),
                }
            )
        model["artifacts"] = artifacts

    analysis_dir = evidence_dir / "analysis"
    analysis_dir.mkdir(exist_ok=True)
    inventory_path = analysis_dir / "inventory.json"
    inventory_path.write_text(
        json.dumps(
            {
                "schema_version": "1.0.0",
                "source": {"sha256": source_sha.upper()},
                "canvas": {"w": width, "h": height, "background_hex": "#FFFFFF"},
            }
        ),
        encoding="utf-8",
    )
    raw_manifest["upstream_stages"] = [
        {
            "name": "analysis",
            "path": str(analysis_dir),
            "files": [
                {
                    "relative_path": "inventory.json",
                    "size_bytes": inventory_path.stat().st_size,
                    "sha256": review_adapter.sha256_file(inventory_path),
                }
            ],
        }
    ]
    manifest_path = source.parent / "perception-manifest.json"
    manifest_path.write_text(json.dumps(raw_manifest, indent=2), encoding="utf-8")

    decisions_path = source.parent / "perception-review-decisions.json"
    decisions = review_adapter.initialize_review(manifest_path, decisions_path)
    for decision in decisions["decisions"]:
        decision["evidence"] = {
            "kind": "user_confirmed",
            "detail": "Test fixture reviewed against its frozen source.",
        }
        if decision["formula_like"]:
            decision["status"] = "FORMULA_CONFIRMED"
            decision["authoritative_latex"] = "x^2+y^2"
        else:
            decision["status"] = "CONFIRMED"
            decision["confirmed_text"] = decision["ocr_text"]
    decisions_path.write_text(json.dumps(decisions, indent=2), encoding="utf-8")
    review_receipt_path = source.parent / "perception-review-receipt.json"
    receipt, exit_code = review_adapter.finalize_review(
        manifest_path,
        decisions_path,
        review_receipt_path,
    )
    assert exit_code == 0
    assert receipt["status"] == "PERCEPTION_REVIEW_PASS"

    canvas_path = source.parent / "blank-canvas.pptx"
    canvas_report = create_blank_canvas_pptx(source, canvas_path, overwrite=canvas_path.exists())
    normalized_elements = []
    for element in elements:
        defaults = {
            "parent_id": None,
            "semantic_role": str(element.get("id", "test element")),
            "source_evidence": ["user_confirmed"]
            if element.get("type") == "text"
            else ["manual_measurement"],
            "disposition": "CONFIRMED",
            "confidence": 1.0,
            "uncertainty_px": 0,
            "strategy": "native_editable",
            "allowed_overlap": [],
            "status": "pending",
        }
        if element.get("type") == "text":
            defaults.update({"criticality": "ordinary", "perception_candidate_ids": []})
        normalized_elements.append({**defaults, **element})
    return {
        "schema_version": "3.0",
        "mode": "reconstruct_1to1",
        "source": {
            "path": str(source),
            "sha256": digest(source),
            "width_px": width,
            "height_px": height,
            "pixel_format": pixel_mode,
            "user_confirmed": True,
        },
        "perception": {
            "manifest_path": str(manifest_path),
            "manifest_sha256": digest(manifest_path),
            "review_receipt_path": str(review_receipt_path),
            "review_receipt_sha256": digest(review_receipt_path),
        },
        "coordinate_system": {
            "origin": "top-left",
            "unit": "source_pixel",
            "bbox_order": ["x", "y", "w", "h"],
        },
        "canvas": {
            "width_px": width,
            "height_px": height,
            "background": "#FFFFFF",
            "background_evidence": "measured_reference",
            "pptx_path": str(canvas_path),
            "pptx_sha256": canvas_report["output_pptx_sha256"],
            "slide_width_emu": canvas_report["slide"]["width_emu"],
            "slide_height_emu": canvas_report["slide"]["height_emu"],
        },
        "measurement_dpi": 96,
        "elements": normalized_elements,
        "edges": [],
        "formulas": [],
        "uncertainties": [],
    }


def make_geometry_source(path: Path) -> Path:
    image = Image.new("RGB", (800, 600), "white")
    draw = ImageDraw.Draw(image)
    for x in (46, 61, 76, 91, 106):
        draw.rectangle((x, 46, x + 7, 64), fill="black")
    for x in (246, 261, 276, 291):
        draw.rectangle((x, 46, x + 7, 64), fill="black")
    image.save(path, format="PNG")
    return path


def bind_real_geometry_evidence(
    spec: dict,
    tmp_path: Path,
    *,
    expected_status: str = "GEOMETRY_OBSERVATIONS_READY",
) -> tuple[Path, Path, dict]:
    def bind_primary_boxes(manifest: dict) -> None:
        set_candidate_geometry(
            manifest,
            "T0001",
            {"x": 40, "y": 40, "w": 100, "h": 32},
        )
        set_candidate_geometry(
            manifest,
            "T0002",
            {"x": 240, "y": 40, "w": 100, "h": 32},
        )
        for index, candidate in enumerate(manifest["text_candidates"], start=1):
            candidate["primary_observation_id"] = f"O{index:05d}"

    raw_manifest, _review_receipt = refinalize_review(
        spec,
        manifest_mutation=bind_primary_boxes,
    )
    host_receipt = validate_runtime(
        config_path=PROJECT_ROOT / "host-runtime.json",
        project_root=PROJECT_ROOT,
        run_id=str(raw_manifest["run_id"]),
        source_sha256=str(raw_manifest["source"]["sha256"]),
    )
    assert host_receipt["status"] == "PASS"
    host_receipt_path = tmp_path / "host-runtime-receipt.json"
    write_json(host_receipt_path, host_receipt)

    geometry_dir = tmp_path / "geometry"
    geometry_manifest_path, geometry_manifest = run_geometry_refinement(
        source_path=Path(spec["source"]["path"]),
        ocr_manifest_path=Path(spec["perception"]["manifest_path"]),
        host_runtime_receipt_path=host_receipt_path,
        output_dir=geometry_dir,
        project_root=PROJECT_ROOT,
        require_isolated_runtime=True,
    )
    assert geometry_manifest["status"] == expected_status
    spec["geometry"] = {
        "manifest_path": str(geometry_manifest_path),
        "manifest_sha256": digest(geometry_manifest_path),
    }
    return geometry_manifest_path, host_receipt_path, geometry_manifest


def test_preflight_passes_with_real_pillow_font_metrics(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    font = usable_font()
    spec = base_spec(
        source,
        elements=[
            {
                "id": "P1",
                "type": "panel",
                "bbox": {"x": 0, "y": 0, "w": 800, "h": 600},
                "z_index": 0,
            },
            {
                "id": "T1",
                "type": "text",
                "parent_id": "P1",
                "bbox": {"x": 40, "y": 40, "w": 240, "h": 60},
                "z_index": 1,
                "text": "Cross Attention",
                "text_style": {
                    "font_family": "Arial",
                    "font_path": str(font),
                    "font_size_px": 20,
                    "margin_px": 2,
                    "wrap": True,
                },
            },
        ],
    )

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "PASS"
    assert report["text_measurements"][0]["status"] == "PASS"
    assert report["text_measurements"][0]["font_path"] == str(font)


def test_hash_mismatch_returns_spec_invalid(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "S1",
                "type": "shape",
                "bbox": {"x": 10, "y": 10, "w": 100, "h": 80},
                "z_index": 0,
            }
        ],
    )
    spec["source"]["sha256"] = "0" * 64

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "SPEC_INVALID"
    assert "SOURCE_HASH_MISMATCH" in {finding["code"] for finding in report["findings"]}


def test_duplicate_id_returns_spec_invalid(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    repeated = {
        "id": "S1",
        "type": "shape",
        "bbox": {"x": 10, "y": 10, "w": 100, "h": 80},
        "z_index": 0,
    }
    spec = base_spec(source, elements=[repeated, {**repeated, "z_index": 1}])

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "SPEC_INVALID"
    assert "DUPLICATE_ELEMENT_ID" in {finding["code"] for finding in report["findings"]}


def test_collision_text_overflow_containment_and_z_require_replan(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    font = usable_font()
    spec = base_spec(
        source,
        elements=[
            {
                "id": "P1",
                "type": "panel",
                "bbox": {"x": 0, "y": 0, "w": 300, "h": 200},
                "z_index": 0,
            },
            {
                "id": "S1",
                "type": "shape",
                "parent_id": "P1",
                "bbox": {"x": 20, "y": 20, "w": 120, "h": 100},
                "z_index": 1,
            },
            {
                "id": "S2",
                "type": "shape",
                "parent_id": "P1",
                "bbox": {"x": 100, "y": 70, "w": 120, "h": 100},
                "z_index": 2,
            },
            {
                "id": "T1",
                "type": "text",
                "parent_id": "P1",
                "bbox": {"x": 290, "y": 190, "w": 20, "h": 10},
                "z_index": 1,
                "text": "This scientific label cannot fit",
                "text_style": {
                    "font_family": "Arial",
                    "font_path": str(font),
                    "font_size_px": 20,
                    "wrap": True,
                },
            },
        ],
    )

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)
    codes = {finding["code"] for finding in report["findings"]}

    assert report["status"] == "REGION_REPLAN"
    assert "SHAPE_SHAPE_COLLISION" in codes
    assert "TEXT_OVERFLOW" in codes
    assert "PARENT_CONTAINMENT" in codes
    assert "AMBIGUOUS_Z_INDEX" not in codes


def test_same_z_only_requires_replan_when_siblings_overlap(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {"id": "S1", "type": "shape", "bbox": {"x": 20, "y": 20, "w": 80, "h": 80}, "z_index": 1},
            {"id": "S2", "type": "shape", "bbox": {"x": 70, "y": 70, "w": 80, "h": 80}, "z_index": 1},
        ],
    )

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert "AMBIGUOUS_Z_INDEX" in {finding["code"] for finding in report["findings"]}


def test_explicit_allowed_overlap_suppresses_collision(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "S1",
                "type": "shape",
                "bbox": {"x": 20, "y": 20, "w": 120, "h": 100},
                "z_index": 0,
                "allowed_overlap": ["S2"],
            },
            {
                "id": "S2",
                "type": "shape",
                "bbox": {"x": 100, "y": 70, "w": 120, "h": 100},
                "z_index": 1,
            },
        ],
    )

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "PASS"
    assert not any(finding["code"].endswith("COLLISION") for finding in report["findings"])


def test_missing_font_is_inconclusive_not_pass(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "T1",
                "type": "text",
                "bbox": {"x": 20, "y": 20, "w": 200, "h": 80},
                "z_index": 0,
                "text": "Unmeasured text",
                "text_style": {
                    "font_family": "Definitely Missing AutoFigure Font",
                    "font_path": str(tmp_path / "missing-font.ttf"),
                    "font_size_px": 18,
                },
            }
        ],
    )

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "INCONCLUSIVE"
    assert report["passed"] is False
    assert report["text_measurements"][0]["status"] == "INCONCLUSIVE"


def test_unknown_edge_endpoint_is_spec_invalid(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "S1",
                "type": "shape",
                "bbox": {"x": 20, "y": 20, "w": 100, "h": 60},
                "z_index": 0,
            },
        ],
    )
    spec["edges"] = [
        {
            "id": "E1",
            "from": "S1",
            "to": "missing-node",
            "source_anchor": "right",
            "target_anchor": "left",
            "route": "straight",
            "meaning": "data_flow",
            "clearance_px": 2,
            "allowed_crossings": [],
            "allowed_edge_crossings": [],
        }
    ]

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "SPEC_INVALID"
    assert "UNKNOWN_EDGE_ENDPOINT" in {finding["code"] for finding in report["findings"]}


def test_typo_text_cannot_bypass_text_measurement_contract(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "T1",
                "type": "typo_text",
                "bbox": {"x": 20, "y": 20, "w": 100, "h": 40},
                "z_index": 0,
            }
        ],
    )

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "SPEC_INVALID"
    assert "SCHEMA_INVALID" in {finding["code"] for finding in report["findings"]}


def test_perception_manifest_hash_and_source_are_enforced(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "BG1",
                "type": "background",
                "bbox": {"x": 0, "y": 0, "w": 800, "h": 600},
                "z_index": 0,
            },
        ],
    )
    manifest_path = Path(spec["perception"]["manifest_path"])
    manifest_path.write_text(
        json.dumps({"schema_version": "1.0.0", "source": {"sha256": "f" * 64}}),
        encoding="utf-8",
    )

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)
    codes = {finding["code"] for finding in report["findings"]}

    assert report["status"] == "SPEC_INVALID"
    assert "PERCEPTION_MANIFEST_HASH_MISMATCH" in codes
    assert "PERCEPTION_SOURCE_MISMATCH" in codes


def test_schema_invalid_raw_manifest_cannot_pass_with_old_review_receipt(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "BG1",
                "type": "background",
                "bbox": {"x": 0, "y": 0, "w": 800, "h": 600},
                "z_index": 0,
            },
        ],
    )
    manifest_path = Path(spec["perception"]["manifest_path"])
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    manifest["policy"]["ocr_is_ground_truth"] = True
    write_json(manifest_path, manifest)
    spec["perception"]["manifest_sha256"] = digest(manifest_path)

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)
    codes = {finding["code"] for finding in report["findings"]}

    assert report["status"] == "SPEC_INVALID"
    assert "PERCEPTION_MANIFEST_SCHEMA_INVALID" in codes


def test_schema_valid_nonpass_review_receipt_is_inconclusive(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "BG1",
                "type": "background",
                "bbox": {"x": 0, "y": 0, "w": 800, "h": 600},
                "z_index": 0,
            },
        ],
    )
    mutate_review_receipt(spec, lambda receipt: receipt.__setitem__("status", "INCONCLUSIVE"))

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "INCONCLUSIVE"
    assert "PERCEPTION_REVIEW_NOT_PASS" in {finding["code"] for finding in report["findings"]}


def test_forged_review_counts_cannot_hide_missing_candidate_coverage(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "BG1",
                "type": "background",
                "bbox": {"x": 0, "y": 0, "w": 800, "h": 600},
                "z_index": 0,
            },
        ],
    )

    def drop_decision_but_keep_pass_counts(receipt: dict) -> None:
        receipt["decisions"] = receipt["decisions"][:1]

    receipt = mutate_review_receipt(spec, drop_decision_but_keep_pass_counts)
    assert receipt["status"] == "PERCEPTION_REVIEW_PASS"
    assert receipt["counts"]["terminal_count"] == 2

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)
    codes = {finding["code"] for finding in report["findings"]}

    assert report["status"] == "SPEC_INVALID"
    assert "PERCEPTION_REVIEW_COVERAGE_MISMATCH" in codes


def test_local_ocr_text_requires_reviewed_candidate_id(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    font = usable_font()
    spec = base_spec(
        source,
        elements=[
            {
                "id": "T1",
                "type": "text",
                "bbox": {"x": 40, "y": 40, "w": 180, "h": 50},
                "z_index": 0,
                "text": "Mamba",
                "source_evidence": ["local_ocr"],
                "criticality": "ordinary",
                "perception_candidate_ids": [],
                "text_style": {"font_family": "Arial", "font_path": str(font), "font_size_px": 20},
            },
        ],
    )

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)
    codes = {finding["code"] for finding in report["findings"]}

    assert report["status"] == "SPEC_INVALID"
    assert "TEXT_OCR_CANDIDATE_MISSING" in codes


@pytest.mark.parametrize(
    ("element_text", "element_bbox", "expected_code"),
    [
        ("Marnba", {"x": 35, "y": 35, "w": 140, "h": 50}, "TEXT_CANDIDATE_VALUE_MISMATCH"),
        ("Mamba", {"x": 500, "y": 35, "w": 140, "h": 50}, "TEXT_CANDIDATE_BBOX_MISMATCH"),
    ],
)
def test_local_ocr_text_must_match_reviewed_value_and_geometry(
    tmp_path: Path,
    element_text: str,
    element_bbox: dict[str, float],
    expected_code: str,
) -> None:
    source = make_png(tmp_path / "reference.png")
    font = usable_font()
    spec = base_spec(
        source,
        elements=[
            {
                "id": "T1",
                "type": "text",
                "bbox": element_bbox,
                "z_index": 0,
                "text": element_text,
                "source_evidence": ["local_ocr"],
                "criticality": "ordinary",
                "perception_candidate_ids": ["T0001"],
                "text_style": {"font_family": "Arial", "font_path": str(font), "font_size_px": 20},
            },
        ],
    )
    refinalize_review(
        spec,
        manifest_mutation=lambda manifest: set_candidate_geometry(
            manifest,
            "T0001",
            {"x": 40, "y": 40, "w": 120, "h": 30},
        ),
    )

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "SPEC_INVALID"
    assert expected_code in {finding["code"] for finding in report["findings"]}


def test_conflict_envelope_cannot_replace_primary_ocr_geometry(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    font = usable_font()
    spec = base_spec(
        source,
        elements=[
            {
                "id": "T1",
                "type": "text",
                "bbox": {"x": 40, "y": 40, "w": 180, "h": 50},
                "z_index": 0,
                "text": "Mamba",
                "source_evidence": ["local_ocr"],
                "criticality": "ordinary",
                "perception_candidate_ids": ["T0001"],
                "text_style": {
                    "font_family": "Arial",
                    "font_path": str(font),
                    "font_size_px": 20,
                },
            },
        ],
    )

    def move_only_primary_geometry(manifest: dict) -> None:
        set_candidate_geometry(
            manifest,
            "T0001",
            {"x": 500, "y": 40, "w": 120, "h": 30},
        )
        candidate = next(
            item for item in manifest["text_candidates"] if item["candidate_id"] == "T0001"
        )
        candidate["bbox_envelope_source"] = {"x": 30, "y": 30, "w": 600, "h": 60}

    refinalize_review(spec, manifest_mutation=move_only_primary_geometry)
    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "SPEC_INVALID"
    mismatch = next(
        finding
        for finding in report["findings"]
        if finding["code"] == "TEXT_CANDIDATE_BBOX_MISMATCH"
    )
    assert mismatch["evidence"]["candidate_id"] == "T0001"


def test_formula_conflict_envelope_cannot_replace_primary_ocr_geometry(
    tmp_path: Path,
) -> None:
    pytest.importorskip("matplotlib")
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "F1",
                "type": "formula",
                "bbox": {"x": 40, "y": 40, "w": 240, "h": 80},
                "z_index": 0,
                "formula_id": "EQ1",
                "formula_style": {"font_size_px": 24, "color": "#000000", "margin_px": 2},
            },
        ],
    )
    formula = native_formula_record(
        tmp_path,
        formula_id="EQ1",
        element_id="F1",
        canonical_latex=r"x^2+y^2",
    )
    formula["source_evidence"] = ["source_text", "local_ocr"]
    formula["perception_candidate_id"] = "T0002"
    spec["formulas"] = [formula]

    def move_only_primary_geometry(manifest: dict) -> None:
        set_candidate_geometry(
            manifest,
            "T0002",
            {"x": 500, "y": 40, "w": 120, "h": 30},
        )
        candidate = next(
            item for item in manifest["text_candidates"] if item["candidate_id"] == "T0002"
        )
        candidate["bbox_envelope_source"] = {"x": 30, "y": 30, "w": 600, "h": 100}

    refinalize_review(spec, manifest_mutation=move_only_primary_geometry)
    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "SPEC_INVALID"
    assert "FORMULA_CANDIDATE_BBOX_MISMATCH" in {finding["code"] for finding in report["findings"]}


def test_geometry_binding_is_diagnostic_only_and_receipt_hash_bound(
    tmp_path: Path,
    capsys: pytest.CaptureFixture[str],
) -> None:
    source = make_geometry_source(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "BG1",
                "type": "background",
                "bbox": {"x": 0, "y": 0, "w": 800, "h": 600},
                "z_index": 0,
            },
        ],
    )
    geometry_manifest_path, _host_receipt_path, _manifest = bind_real_geometry_evidence(
        spec, tmp_path
    )

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "PASS"
    assert report["geometry_manifest"]["sha256"] == digest(geometry_manifest_path)
    assert report["geometry_manifest"]["mode"] == "observation_only"
    assert report["geometry_manifest"]["promotion_allowed"] is False
    assert report["geometry_manifest"]["contributes_to_drawer_authorization"] is False
    boundary = report["geometry_diagnostics"][0]
    assert boundary["diagnostic_only"] is True
    assert boundary["used_as_scene_geometry"] is False
    assert boundary["authorizes_drawer"] is False

    spec_path = tmp_path / "figure-spec.json"
    output_path = tmp_path / "preflight-report.json"
    write_json(spec_path, spec)
    exit_code = preflight_main(
        [
            str(spec_path),
            "--source",
            str(source),
            "--output",
            str(output_path),
        ]
    )
    capsys.readouterr()
    cli_report = json.loads(output_path.read_text(encoding="utf-8"))

    assert exit_code == 0
    assert cli_report["receipt"]["authorized_for_drawer"] is True
    assert cli_report["receipt"]["geometry_authorized_for_drawer"] is False
    assert cli_report["receipt"]["geometry_manifest_sha256"] == digest(geometry_manifest_path)
    assert cli_report["receipt"]["geometry_manifest_schema_sha256"] == digest(
        PROJECT_ROOT / "schemas" / "geometry-manifest.schema.json"
    )


def test_geometry_inconclusive_blocks_drawer_without_claiming_truth(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "BG1",
                "type": "background",
                "bbox": {"x": 0, "y": 0, "w": 800, "h": 600},
                "z_index": 0,
            },
        ],
    )
    bind_real_geometry_evidence(
        spec,
        tmp_path,
        expected_status="GEOMETRY_INCONCLUSIVE",
    )

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "INCONCLUSIVE"
    assert report["passed"] is False
    assert report["geometry_manifest"]["contributes_to_drawer_authorization"] is False
    assert "GEOMETRY_OBSERVATIONS_INCONCLUSIVE" in {
        finding["code"] for finding in report["findings"]
    }


def test_geometry_binding_rejects_rebound_and_tampered_evidence(tmp_path: Path) -> None:
    source = make_geometry_source(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "BG1",
                "type": "background",
                "bbox": {"x": 0, "y": 0, "w": 800, "h": 600},
                "z_index": 0,
            },
        ],
    )
    geometry_path, host_receipt_path, genuine = bind_real_geometry_evidence(spec, tmp_path)
    host_receipt_bytes = host_receipt_path.read_bytes()
    artifact_paths = {
        name: geometry_path.parent / str(record["relative_path"])
        for name, record in genuine["artifacts"].items()
    }
    artifact_bytes = {name: path.read_bytes() for name, path in artifact_paths.items()}
    cases = {
        "manifest_hash": "GEOMETRY_MANIFEST_HASH_MISMATCH",
        "source": "GEOMETRY_SOURCE_MISMATCH",
        "run": "GEOMETRY_RUN_MISMATCH",
        "ocr": "GEOMETRY_OCR_BINDING_MISMATCH",
        "host_receipt": "GEOMETRY_HOST_RECEIPT_HASH_MISMATCH",
        "runtime": "GEOMETRY_RUNTIME_BINDING_MISMATCH",
        "script": "GEOMETRY_IMPLEMENTATION_HASH_MISMATCH",
        "schema": "GEOMETRY_IMPLEMENTATION_HASH_MISMATCH",
        "artifact": "GEOMETRY_ARTIFACT_HASH_MISMATCH",
        "promotion": "GEOMETRY_AUTHORITY_BOUNDARY_VIOLATION",
        "status": "GEOMETRY_STATUS_INVALID",
    }

    for case, expected_code in cases.items():
        host_receipt_path.write_bytes(host_receipt_bytes)
        for name, path in artifact_paths.items():
            path.write_bytes(artifact_bytes[name])
        candidate = copy.deepcopy(genuine)

        if case == "source":
            candidate["source"]["sha256"] = "A" * 64
        elif case == "run":
            candidate["run_id"] = "different-run-001"
        elif case == "ocr":
            candidate["inputs"]["ocr_manifest"]["sha256"] = "A" * 64
        elif case == "runtime":
            candidate["runtime"]["python_version"] = "0.0.0"
        elif case == "script":
            candidate["implementation"]["script"]["sha256"] = "A" * 64
        elif case == "schema":
            candidate["implementation"]["schema"]["sha256"] = "A" * 64
        elif case == "promotion":
            candidate["policy"]["promotion_allowed"] = True
        elif case == "status":
            candidate["status"] = "PASS"

        write_json(geometry_path, candidate)
        spec["geometry"]["manifest_sha256"] = digest(geometry_path)
        if case == "manifest_hash":
            spec["geometry"]["manifest_sha256"] = "0" * 64
        elif case == "host_receipt":
            host_receipt_path.write_bytes(host_receipt_bytes + b"\n")
        elif case == "artifact":
            artifact_paths["overlay"].write_bytes(artifact_bytes["overlay"] + b"tamper")

        report = preflight_scene(spec, source_path=source, base_dir=tmp_path)
        codes = {finding["code"] for finding in report["findings"]}

        assert report["status"] == "SPEC_INVALID", case
        assert expected_code in codes, (case, sorted(codes))


@pytest.mark.parametrize(
    ("case", "expected_code"),
    [
        ("label_pixels", "GEOMETRY_LABEL_ATLAS_LABEL_SET_MISMATCH"),
        ("ambiguity_values", "GEOMETRY_AMBIGUITY_MASK_NOT_BINARY"),
        ("overlay_mode", "GEOMETRY_OVERLAY_ENCODING_MISMATCH"),
        ("overlay_format", "GEOMETRY_ARTIFACT_FORMAT_MISMATCH"),
    ],
)
def test_geometry_artifact_pixels_must_match_manifest_semantics(
    tmp_path: Path,
    case: str,
    expected_code: str,
) -> None:
    source = make_geometry_source(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "BG1",
                "type": "background",
                "bbox": {"x": 0, "y": 0, "w": 800, "h": 600},
                "z_index": 0,
            },
        ],
    )
    geometry_path, _host_receipt_path, genuine = bind_real_geometry_evidence(spec, tmp_path)
    candidate = copy.deepcopy(genuine)

    if case == "label_pixels":
        artifact_name = "label_atlas"
        artifact_path = (
            geometry_path.parent / candidate["artifacts"][artifact_name]["relative_path"]
        )
        with Image.open(artifact_path) as image:
            atlas = np.asarray(image).copy()
        measured = next(item for item in candidate["text_geometry"] if item["status"] == "MEASURED")
        atlas[atlas == measured["mask_label"]] = 0
        Image.fromarray(atlas).save(artifact_path, format="PNG")
    elif case == "ambiguity_values":
        artifact_name = "ambiguity_mask"
        artifact_path = (
            geometry_path.parent / candidate["artifacts"][artifact_name]["relative_path"]
        )
        with Image.open(artifact_path) as image:
            ambiguity = np.asarray(image).copy()
        ambiguity[0, 0] = 127
        Image.fromarray(ambiguity).save(artifact_path, format="PNG")
    elif case == "overlay_mode":
        artifact_name = "overlay"
        artifact_path = (
            geometry_path.parent / candidate["artifacts"][artifact_name]["relative_path"]
        )
        with Image.open(artifact_path) as image:
            rgba = image.convert("RGBA")
        rgba.save(artifact_path, format="PNG")
    else:
        artifact_name = "overlay"
        artifact_path = (
            geometry_path.parent / candidate["artifacts"][artifact_name]["relative_path"]
        )
        with Image.open(artifact_path) as image:
            rgb = image.convert("RGB")
        rgb.save(artifact_path, format="BMP")

    candidate["artifacts"][artifact_name]["size_bytes"] = artifact_path.stat().st_size
    candidate["artifacts"][artifact_name]["sha256"] = digest(artifact_path)
    write_json(geometry_path, candidate)
    spec["geometry"]["manifest_sha256"] = digest(geometry_path)

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)
    codes = {finding["code"] for finding in report["findings"]}

    assert report["status"] == "SPEC_INVALID"
    assert expected_code in codes, sorted(codes)


@pytest.mark.parametrize("case", ["duplicate_key", "nan"])
def test_geometry_manifest_requires_strict_json(tmp_path: Path, case: str) -> None:
    source = make_geometry_source(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "BG1",
                "type": "background",
                "bbox": {"x": 0, "y": 0, "w": 800, "h": 600},
                "z_index": 0,
            },
        ],
    )
    geometry_path, _host_receipt_path, genuine = bind_real_geometry_evidence(spec, tmp_path)
    if case == "duplicate_key":
        payload = geometry_path.read_text(encoding="utf-8")
        needle = '\n  "schema_version": "1.0.0",\n'
        assert payload.count(needle) == 1
        geometry_path.write_text(payload.replace(needle, needle + needle, 1), encoding="utf-8")
    else:
        genuine["summary"]["ambiguous_pixel_count"] = math.nan
        write_json(geometry_path, genuine)
    spec["geometry"]["manifest_sha256"] = digest(geometry_path)

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)
    codes = {finding["code"] for finding in report["findings"]}

    assert report["status"] == "SPEC_INVALID"
    assert "GEOMETRY_EVIDENCE_UNREADABLE" in codes


def test_geometry_evidence_change_during_preflight_is_rejected(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = make_geometry_source(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "BG1",
                "type": "background",
                "bbox": {"x": 0, "y": 0, "w": 800, "h": 600},
                "z_index": 0,
            },
        ],
    )
    geometry_path, _host_receipt_path, _genuine = bind_real_geometry_evidence(spec, tmp_path)
    target = geometry_path.resolve()
    original_read_bytes = Path.read_bytes
    changed = False

    def racing_read_bytes(path: Path) -> bytes:
        nonlocal changed
        payload = original_read_bytes(path)
        if path.resolve() == target and not changed:
            changed = True
            path.write_bytes(payload + b"\n")
        return payload

    monkeypatch.setattr(Path, "read_bytes", racing_read_bytes)
    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)
    codes = {finding["code"] for finding in report["findings"]}

    assert changed
    assert report["status"] == "SPEC_INVALID"
    assert "EVIDENCE_CHANGED_DURING_PREFLIGHT" in codes


def test_critical_text_without_user_or_source_authority_is_blocked(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    font = usable_font()
    spec = base_spec(
        source,
        elements=[
            {
                "id": "T1",
                "type": "text",
                "bbox": {"x": 40, "y": 40, "w": 180, "h": 50},
                "z_index": 0,
                "text": "0.5 mg mL⁻¹",
                "source_evidence": ["target_visual"],
                "criticality": "critical",
                "perception_candidate_ids": [],
                "text_style": {"font_family": "Arial", "font_path": str(font), "font_size_px": 20},
            },
        ],
    )

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)
    codes = {finding["code"] for finding in report["findings"]}

    assert report["status"] == "SPEC_INVALID"
    assert "CRITICAL_TEXT_AUTHORITY_MISSING" in codes
    assert "MATH_SYNTAX_IN_PLAIN_TEXT" in codes


def test_rebound_wrong_ratio_canvas_is_rejected_by_pagesetup(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png", (800, 600))
    spec = base_spec(
        source,
        elements=[
            {
                "id": "BG1",
                "type": "background",
                "bbox": {"x": 0, "y": 0, "w": 800, "h": 600},
                "z_index": 0,
            },
        ],
    )
    wrong_reference = make_png(tmp_path / "wrong-reference.png", (1600, 900))
    wrong_canvas = tmp_path / "wrong-ratio-canvas.pptx"
    create_blank_canvas_pptx(wrong_reference, wrong_canvas)
    spec["canvas"]["pptx_path"] = str(wrong_canvas)
    spec["canvas"]["pptx_sha256"] = digest(wrong_canvas)

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)
    codes = {finding["code"] for finding in report["findings"]}

    assert report["status"] == "SPEC_INVALID"
    assert "CANVAS_PAGESETUP_MISMATCH" in codes
    assert "CANVAS_ASPECT_MISMATCH" in codes


def test_hash_rebound_nonblank_canvas_is_still_rejected(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "BG1",
                "type": "background",
                "bbox": {"x": 0, "y": 0, "w": 800, "h": 600},
                "z_index": 0,
            },
        ],
    )
    canvas_path = Path(spec["canvas"]["pptx_path"])
    presentation = Presentation(canvas_path)
    presentation.slides[0].shapes.add_textbox(0, 0, 914400, 914400).text = "premature draw"
    presentation.save(canvas_path)
    spec["canvas"]["pptx_sha256"] = digest(canvas_path)

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "SPEC_INVALID"
    assert "CANVAS_NOT_BLANK" in {finding["code"] for finding in report["findings"]}


def test_self_reported_source_pixel_mode_is_measured_not_trusted(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "BG1",
                "type": "background",
                "bbox": {"x": 0, "y": 0, "w": 800, "h": 600},
                "z_index": 0,
            },
        ],
    )
    spec["source"]["pixel_format"] = "RGBA"

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "SPEC_INVALID"
    assert "SOURCE_PIXEL_MODE_MISMATCH" in {finding["code"] for finding in report["findings"]}


def test_declared_background_must_match_hash_bound_measurement(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "BG1",
                "type": "background",
                "bbox": {"x": 0, "y": 0, "w": 800, "h": 600},
                "z_index": 0,
            },
        ],
    )
    spec["canvas"]["background"] = "#000000"

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "SPEC_INVALID"
    assert "BACKGROUND_COLOR_MISMATCH" in {finding["code"] for finding in report["findings"]}


def test_noncontainer_parent_cannot_hide_a_collision(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "S1",
                "type": "shape",
                "bbox": {"x": 40, "y": 40, "w": 200, "h": 180},
                "z_index": 0,
            },
            {
                "id": "S2",
                "type": "shape",
                "parent_id": "S1",
                "bbox": {"x": 100, "y": 100, "w": 100, "h": 80},
                "z_index": 1,
            },
        ],
    )

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)
    codes = {finding["code"] for finding in report["findings"]}

    assert report["status"] == "SPEC_INVALID"
    assert "INVALID_PARENT_TYPE" in codes
    assert "SHAPE_SHAPE_COLLISION" in codes


def test_straight_edge_crossing_text_requires_replan(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    font = usable_font()
    spec = base_spec(
        source,
        elements=[
            {
                "id": "S1",
                "type": "shape",
                "bbox": {"x": 50, "y": 250, "w": 80, "h": 60},
                "z_index": 0,
            },
            {
                "id": "S2",
                "type": "shape",
                "bbox": {"x": 670, "y": 250, "w": 80, "h": 60},
                "z_index": 1,
            },
            {
                "id": "T1",
                "type": "text",
                "bbox": {"x": 320, "y": 240, "w": 160, "h": 80},
                "z_index": 2,
                "text": "Do not cross",
                "text_style": {"font_family": "Arial", "font_path": str(font), "font_size_px": 18},
            },
        ],
    )
    spec["edges"] = [edge("E1", "S1", "S2")]

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "REGION_REPLAN"
    assert "EDGE_TEXT_COLLISION" in {finding["code"] for finding in report["findings"]}


def test_edge_route_point_outside_source_pixel_canvas_requires_replan(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "S1",
                "type": "shape",
                "bbox": {"x": 50, "y": 250, "w": 80, "h": 60},
                "z_index": 0,
            },
            {
                "id": "S2",
                "type": "shape",
                "bbox": {"x": 670, "y": 250, "w": 80, "h": 60},
                "z_index": 1,
            },
        ],
    )
    spec["edges"] = [
        edge(
            "E1",
            "S1",
            "S2",
            route="polyline",
            via=[{"x": -20, "y": 100}, {"x": 400, "y": 100}],
        )
    ]

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "REGION_REPLAN"
    assert "EDGE_ROUTE_OUTSIDE_CANVAS" in {finding["code"] for finding in report["findings"]}


def test_unapproved_edge_crossing_requires_replan(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "A",
                "type": "shape",
                "bbox": {"x": 50, "y": 50, "w": 40, "h": 40},
                "z_index": 0,
            },
            {
                "id": "B",
                "type": "shape",
                "bbox": {"x": 710, "y": 510, "w": 40, "h": 40},
                "z_index": 1,
            },
            {
                "id": "C",
                "type": "shape",
                "bbox": {"x": 50, "y": 510, "w": 40, "h": 40},
                "z_index": 2,
            },
            {
                "id": "D",
                "type": "shape",
                "bbox": {"x": 710, "y": 50, "w": 40, "h": 40},
                "z_index": 3,
            },
        ],
    )
    spec["edges"] = [
        edge("E1", "A", "B", source_anchor="center", target_anchor="center"),
        edge("E2", "C", "D", source_anchor="center", target_anchor="center"),
    ]

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "REGION_REPLAN"
    assert "EDGE_EDGE_CROSSING" in {finding["code"] for finding in report["findings"]}


@pytest.mark.parametrize(
    ("latex", "bbox", "expected_status", "expected_code"),
    [
        (r"x^2+y^2", {"x": 40, "y": 40, "w": 240, "h": 80}, "PASS", None),
        (r"x^2+y^2", {"x": 40, "y": 40, "w": 8, "h": 8}, "REGION_REPLAN", "FORMULA_OVERFLOW"),
    ],
)
def test_mathtext_formula_fit_is_measured_before_drawing(
    tmp_path: Path,
    latex: str,
    bbox: dict[str, float],
    expected_status: str,
    expected_code: str | None,
) -> None:
    pytest.importorskip("matplotlib")
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "F1",
                "type": "formula",
                "bbox": bbox,
                "z_index": 0,
                "formula_id": "EQ1",
                "formula_style": {"font_size_px": 24, "color": "#000000", "margin_px": 2},
            },
        ],
    )
    spec["formulas"] = [
        native_formula_record(
            tmp_path,
            formula_id="EQ1",
            element_id="F1",
            canonical_latex=latex,
        )
    ]

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == expected_status
    assert report["formula_measurements"][0]["status"] == expected_status
    if expected_code is not None:
        assert expected_code in {finding["code"] for finding in report["findings"]}


def test_unparseable_mathtext_formula_is_inconclusive(tmp_path: Path) -> None:
    pytest.importorskip("matplotlib")
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "F1",
                "type": "formula",
                "bbox": {"x": 40, "y": 40, "w": 240, "h": 80},
                "z_index": 0,
                "formula_id": "EQ1",
                "formula_style": {"font_size_px": 24, "color": "#000000"},
            },
        ],
    )
    spec["formulas"] = [
        native_formula_record(
            tmp_path,
            formula_id="EQ1",
            element_id="F1",
            canonical_latex=r"\begin{matrix}a&b\\c&d\end{matrix}",
        )
    ]

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "INCONCLUSIVE"
    assert report["formula_measurements"][0]["status"] == "INCONCLUSIVE"
    assert "FORMULA_METRICS_INCONCLUSIVE" in {finding["code"] for finding in report["findings"]}


def test_plain_text_math_syntax_is_rejected_but_scientific_entities_are_not(tmp_path: Path) -> None:
    source = make_png(tmp_path / "reference.png")
    font = usable_font()
    entity_spec = base_spec(
        source,
        elements=[
            {
                "id": "T1",
                "type": "text",
                "bbox": {"x": 40, "y": 40, "w": 300, "h": 60},
                "z_index": 0,
                "text": "IL-6, p53 and α-SMA",
                "text_style": {
                    "font_family": "Arial",
                    "font_path": str(font),
                    "font_size_px": 20,
                },
            }
        ],
    )

    entity_report = preflight_scene(entity_spec, source_path=source, base_dir=tmp_path)

    assert entity_report["status"] == "PASS"
    assert "MATH_SYNTAX_IN_PLAIN_TEXT" not in {
        finding["code"] for finding in entity_report["findings"]
    }

    math_spec = json.loads(json.dumps(entity_spec))
    math_spec["elements"][0]["text"] = "Loss = x^2"
    math_report = preflight_scene(math_spec, source_path=source, base_dir=tmp_path)

    assert math_report["status"] == "SPEC_INVALID"
    assert "MATH_SYNTAX_IN_PLAIN_TEXT" in {finding["code"] for finding in math_report["findings"]}


def test_inline_formula_requires_one_structured_math_run_and_native_receipt(tmp_path: Path) -> None:
    pytest.importorskip("matplotlib")
    source = make_png(tmp_path / "reference.png")
    font = usable_font()
    spec = base_spec(
        source,
        elements=[
            {
                "id": "T1",
                "type": "text",
                "bbox": {"x": 40, "y": 40, "w": 320, "h": 80},
                "z_index": 0,
                "content_runs": [
                    {"kind": "text", "text": "Loss "},
                    {"kind": "math", "formula_id": "EQ1"},
                ],
                "text_style": {
                    "font_family": "Arial",
                    "font_path": str(font),
                    "font_size_px": 24,
                },
            }
        ],
    )
    spec["formulas"] = [
        native_formula_record(
            tmp_path,
            formula_id="EQ1",
            element_id="T1",
            canonical_latex=r"x_t^2",
            mode="inline",
        )
    ]

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "PASS"
    assert report["formula_converter_receipts"][0]["status"] == "PASS"
    assert report["formula_measurements"][0]["diagnostic_only"] is True
    assert report["formula_measurements"][0]["proves_native_office_math"] is False


def test_duplicate_inline_formula_reference_is_spec_invalid(tmp_path: Path) -> None:
    pytest.importorskip("matplotlib")
    source = make_png(tmp_path / "reference.png")
    font = usable_font()
    spec = base_spec(
        source,
        elements=[
            {
                "id": "T1",
                "type": "text",
                "bbox": {"x": 40, "y": 40, "w": 320, "h": 80},
                "z_index": 0,
                "content_runs": [
                    {"kind": "math", "formula_id": "EQ1"},
                    {"kind": "text", "text": " and "},
                    {"kind": "math", "formula_id": "EQ1"},
                ],
                "text_style": {
                    "font_family": "Arial",
                    "font_path": str(font),
                    "font_size_px": 24,
                },
            }
        ],
    )
    spec["formulas"] = [
        native_formula_record(
            tmp_path,
            formula_id="EQ1",
            element_id="T1",
            canonical_latex="x",
            mode="inline",
        )
    ]

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "SPEC_INVALID"
    assert "FORMULA_REFERENCE_NOT_UNIQUE" in {finding["code"] for finding in report["findings"]}


def test_parseable_mathtext_cannot_replace_native_converter_receipt(tmp_path: Path) -> None:
    pytest.importorskip("matplotlib")
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "F1",
                "type": "formula",
                "bbox": {"x": 40, "y": 40, "w": 240, "h": 80},
                "z_index": 0,
                "formula_id": "EQ1",
                "formula_style": {"font_size_px": 24, "color": "#000000"},
            }
        ],
    )
    formula = native_formula_record(
        tmp_path,
        formula_id="EQ1",
        element_id="F1",
        canonical_latex="x^2+y^2",
    )
    formula["converter_receipt_path"] = str(tmp_path / "missing-receipt.json")
    formula["converter_receipt_sha256"] = "d" * 64
    spec["formulas"] = [formula]

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "INCONCLUSIVE"
    assert report["formula_measurements"][0]["status"] == "PASS"
    assert "FORMULA_CONVERTER_RECEIPT_UNREADABLE" in {
        finding["code"] for finding in report["findings"]
    }


def test_self_signed_formula_receipt_cannot_authorize_different_omml(tmp_path: Path) -> None:
    pytest.importorskip("matplotlib")
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "F1",
                "type": "formula",
                "bbox": {"x": 40, "y": 40, "w": 240, "h": 80},
                "z_index": 0,
                "formula_id": "EQ1",
                "formula_style": {"font_size_px": 24, "color": "#000000"},
            }
        ],
    )
    formula = native_formula_record(
        tmp_path,
        formula_id="EQ1",
        element_id="F1",
        canonical_latex="x",
    )
    receipt_path = Path(formula["converter_receipt_path"])

    # The forged receipt is internally self-consistent: XML, all output
    # hashes, target, and real runtime provenance describe y.  The attacker
    # then rebinds only the claimed canonical source to x and rehashes the
    # receipt in the spec.  Deterministic recompilation must still catch it.
    forged = compile_formula("EQ1", "y", "display")
    forged["canonical_latex"] = "x"
    forged["latex_sha256"] = hashlib.sha256(b"x").hexdigest()
    write_json(receipt_path, forged)
    formula["converter_receipt_sha256"] = digest(receipt_path)
    spec["formulas"] = [formula]

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "SPEC_INVALID"
    assert report["formula_converter_receipts"][0]["status"] == "SPEC_INVALID"
    assert "FORMULA_CONVERTER_RECEIPT_INVALID" in {
        finding["code"] for finding in report["findings"]
    }


def test_formula_receipt_missing_embedded_xml_cannot_pass(tmp_path: Path) -> None:
    pytest.importorskip("matplotlib")
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "F1",
                "type": "formula",
                "bbox": {"x": 40, "y": 40, "w": 240, "h": 80},
                "z_index": 0,
                "formula_id": "EQ1",
                "formula_style": {"font_size_px": 24, "color": "#000000"},
            }
        ],
    )
    formula = native_formula_record(
        tmp_path,
        formula_id="EQ1",
        element_id="F1",
        canonical_latex="x",
    )
    receipt_path = Path(formula["converter_receipt_path"])
    receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
    del receipt["artifacts"]
    write_json(receipt_path, receipt)
    formula["converter_receipt_sha256"] = digest(receipt_path)
    spec["formulas"] = [formula]

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "SPEC_INVALID"
    assert report["formula_converter_receipts"][0]["status"] == "SPEC_INVALID"
    assert "FORMULA_CONVERTER_RECEIPT_INVALID" in {
        finding["code"] for finding in report["findings"]
    }


def test_formula_receipt_bound_fields_and_all_hashes_are_hard_gates(tmp_path: Path) -> None:
    pytest.importorskip("matplotlib")
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "F1",
                "type": "formula",
                "bbox": {"x": 40, "y": 40, "w": 240, "h": 80},
                "z_index": 0,
                "formula_id": "EQ1",
                "formula_style": {"font_size_px": 24, "color": "#000000"},
            }
        ],
    )
    formula = native_formula_record(
        tmp_path,
        formula_id="EQ1",
        element_id="F1",
        canonical_latex="x",
    )
    receipt_path = Path(formula["converter_receipt_path"])
    genuine = json.loads(receipt_path.read_text(encoding="utf-8"))
    spec["formulas"] = [formula]
    mutations = [
        (("formula_id",), "OTHER"),
        (("canonical_latex",), "y"),
        (("mode",), "inline"),
        (("native_target", "omml_root"), "m:oMath"),
        (("semantic_omml_profile",), "office-math-semantic-v1"),
        (("latex_sha256",), "0" * 64),
        (("mathml_sha256",), "1" * 64),
        (("omml_sha256",), "2" * 64),
        (("semantic_omml_sha256",), "3" * 64),
        (("converter", "latex2mathml_version"), "forged-version"),
        (("converter", "lxml_version"), "forged-version"),
        (("converter", "xsl_sha256"), "4" * 64),
    ]

    for field_path, forged_value in mutations:
        forged = json.loads(json.dumps(genuine))
        target = forged
        for key in field_path[:-1]:
            target = target[key]
        target[field_path[-1]] = forged_value
        write_json(receipt_path, forged)
        formula["converter_receipt_sha256"] = digest(receipt_path)

        report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

        assert report["status"] == "SPEC_INVALID", field_path
        assert report["formula_converter_receipts"][0]["status"] == "SPEC_INVALID", field_path
        assert {finding["code"] for finding in report["findings"]}.intersection(
            {
                "FORMULA_CONVERTER_RECEIPT_INVALID",
                "FORMULA_CONVERTER_RECEIPT_BINDING_MISMATCH",
            }
        ), field_path


def test_formula_receipt_sha_is_verified_even_when_receipt_is_genuine(tmp_path: Path) -> None:
    pytest.importorskip("matplotlib")
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "F1",
                "type": "formula",
                "bbox": {"x": 40, "y": 40, "w": 240, "h": 80},
                "z_index": 0,
                "formula_id": "EQ1",
                "formula_style": {"font_size_px": 24, "color": "#000000"},
            }
        ],
    )
    formula = native_formula_record(
        tmp_path,
        formula_id="EQ1",
        element_id="F1",
        canonical_latex="x",
    )
    formula["converter_receipt_sha256"] = "0" * 64
    spec["formulas"] = [formula]

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)

    assert report["status"] == "SPEC_INVALID"
    assert report["formula_converter_receipts"][0]["status"] == "SPEC_INVALID"
    assert "FORMULA_CONVERTER_RECEIPT_HASH_MISMATCH" in {
        finding["code"] for finding in report["findings"]
    }


def test_formula_hash_and_no_fallback_policy_are_hard_gates(tmp_path: Path) -> None:
    pytest.importorskip("matplotlib")
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "F1",
                "type": "formula",
                "bbox": {"x": 40, "y": 40, "w": 240, "h": 80},
                "z_index": 0,
                "formula_id": "EQ1",
                "formula_style": {"font_size_px": 24, "color": "#000000"},
            }
        ],
    )
    formula = native_formula_record(
        tmp_path,
        formula_id="EQ1",
        element_id="F1",
        canonical_latex="x",
    )
    formula["latex_sha256"] = "0" * 64
    formula["render_kind"] = "svg"
    formula["fallback_policy"] = "allow_png"
    spec["formulas"] = [formula]

    report = preflight_scene(spec, source_path=source, base_dir=tmp_path)
    codes = {finding["code"] for finding in report["findings"]}

    assert report["status"] == "SPEC_INVALID"
    assert {
        "FORMULA_LATEX_HASH_MISMATCH",
        "FORMULA_RENDER_KIND_INVALID",
        "FORMULA_FALLBACK_POLICY_INVALID",
    } <= codes


def test_cli_emits_hash_bound_pass_receipt(
    tmp_path: Path, capsys: pytest.CaptureFixture[str]
) -> None:
    source = make_png(tmp_path / "reference.png")
    spec = base_spec(
        source,
        elements=[
            {
                "id": "BG1",
                "type": "background",
                "bbox": {"x": 0, "y": 0, "w": 800, "h": 600},
                "z_index": 0,
            },
        ],
    )
    spec_path = tmp_path / "figure-spec.json"
    spec_path.write_text(json.dumps(spec), encoding="utf-8")

    exit_code = preflight_main(
        [
            str(spec_path),
            "--canvas-pptx",
            spec["canvas"]["pptx_path"],
            "--pretty",
        ]
    )
    report = json.loads(capsys.readouterr().out)

    assert exit_code == 0
    assert report["status"] == "PASS"
    assert report["receipt"]["spec_sha256"] == digest(spec_path)
    assert report["receipt"]["source_sha256"] == digest(source)
    assert report["receipt"]["canvas_pptx_path"] == str(Path(spec["canvas"]["pptx_path"]).resolve())
    assert report["receipt"]["canvas_pptx_sha256"] == spec["canvas"]["pptx_sha256"]
    assert report["receipt"]["canvas_slide_width_emu"] == spec["canvas"]["slide_width_emu"]
    assert report["receipt"]["canvas_slide_height_emu"] == spec["canvas"]["slide_height_emu"]
