from __future__ import annotations

import hashlib
import json
import os
import subprocess
import zipfile
from pathlib import Path

import pytest
from lxml import etree
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt

from tools.powerpoint_native_math import (
    A_NS,
    A14_NS,
    MC_NS,
    M_NS,
    NS,
    NativeMathError,
    _atomic_write_json_fresh,
    _counterfactual_pixel_evidence,
    _decode_metadata,
    _load_mml2omml_transform,
    _metadata_value,
    _powershell_executable,
    _semantic_omml_sha256,
    audit_pptx,
    compile_formula,
    finalize_native_math,
    inject_plan,
    main,
)


def test_mml2omml_stylesheet_is_cached_per_resolved_path(tmp_path: Path) -> None:
    stylesheet = tmp_path / "identity.xsl"
    stylesheet.write_text(
        """<?xml version="1.0"?>
<xsl:stylesheet version="1.0" xmlns:xsl="http://www.w3.org/1999/XSL/Transform">
  <xsl:template match="/"><result/></xsl:template>
</xsl:stylesheet>
""",
        encoding="utf-8",
    )
    _load_mml2omml_transform.cache_clear()
    first = _load_mml2omml_transform(str(stylesheet.resolve()))
    second = _load_mml2omml_transform(str(stylesheet.resolve()))
    assert first is second
    assert _load_mml2omml_transform.cache_info().hits == 1


def write_json(path: Path, payload: object) -> Path:
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return path


def sha256_file(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def make_placeholder_deck(path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    before = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(1.0), Inches(0.4))
    before.name = "before"
    before.text = "before"

    display = slide.shapes.add_textbox(Inches(1.5), Inches(0.8), Inches(4.0), Inches(1.0))
    display.name = "formula-EQ1"
    display.text = "FORMULA_PLACEHOLDER"
    display.text_frame.paragraphs[0].font.size = Pt(20)

    mixed = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(6.0), Inches(0.8))
    mixed.name = "mixed-T1"
    mixed.text = "MIXED PLACEHOLDER"
    mixed.text_frame.paragraphs[0].font.size = Pt(16)

    after = slide.shapes.add_textbox(Inches(0.3), Inches(3.5), Inches(1.0), Inches(0.4))
    after.name = "after"
    after.text = "after"
    presentation.save(path)
    return path


def make_reordered_placeholder_deck(path: Path) -> Path:
    presentation = Presentation()
    distractor = presentation.slides.add_slide(presentation.slide_layouts[6])
    distractor_box = distractor.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.5))
    distractor_box.name = "not-the-formula"
    distractor_box.text = "second in UI order"
    target = presentation.slides.add_slide(presentation.slide_layouts[6])
    placeholder = target.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    placeholder.name = "formula-EQ1"
    placeholder.text = "FORMULA_PLACEHOLDER"
    slide_ids = presentation.slides._sldIdLst  # noqa: SLF001 - fixture must create reordered parts.
    slide_ids.insert(0, slide_ids[-1])
    presentation.save(path)
    return path


def slide_sequence(path: Path) -> list[str]:
    with zipfile.ZipFile(path) as package:
        root = etree.fromstring(package.read("ppt/slides/slide1.xml"))
    sequence: list[str] = []
    tree = root.find(".//p:spTree", namespaces=NS)
    assert tree is not None
    for child in tree:
        if child.tag.endswith("nvGrpSpPr") or child.tag.endswith("grpSpPr"):
            continue
        shape = child
        if child.tag == f"{{{MC_NS}}}AlternateContent":
            choices = child.xpath("./mc:Choice/p:sp", namespaces=NS)
            assert len(choices) == 1
            shape = choices[0]
        c_nv_pr = shape.find("./p:nvSpPr/p:cNvPr", namespaces=NS)
        sequence.append("" if c_nv_pr is None else str(c_nv_pr.get("name")))
    return sequence


def rewrite_first_slide(source: Path, output: Path, mutate: object) -> Path:
    with zipfile.ZipFile(source, "r") as input_zip, zipfile.ZipFile(
        output, "w", compression=zipfile.ZIP_DEFLATED
    ) as output_zip:
        for info in input_zip.infolist():
            value = input_zip.read(info.filename)
            if info.filename == "ppt/slides/slide1.xml":
                root = etree.fromstring(value)
                assert callable(mutate)
                mutate(root)
                value = etree.tostring(
                    root, xml_declaration=True, encoding="UTF-8", standalone=True
                )
            output_zip.writestr(info, value)
    return output


def inject_single_formula(tmp_path: Path, latex: str = "x") -> tuple[Path, Path]:
    source = make_placeholder_deck(tmp_path / "source.pptx")
    receipt_path = write_json(tmp_path / "eq.json", compile_formula("EQ1", latex, "display"))
    plan_path = write_json(
        tmp_path / "plan.json",
        {
            "schema_version": "1.0",
            "operations": [
                {
                    "slide_index": 1,
                    "placeholder_name": "formula-EQ1",
                    "formula_id": "EQ1",
                    "receipt_path": receipt_path.name,
                    "receipt_sha256": sha256_file(receipt_path),
                }
            ]
        },
    )
    output = tmp_path / "native.pptx"
    inject_plan(source, plan_path, output)
    return output, plan_path


def test_compile_formula_produces_hash_bound_native_receipt() -> None:
    receipt = compile_formula("EQ1", r"\frac{E}{mc^2}+\alpha_i", "display")
    assert receipt["document_type"] == "NATIVE_OFFICE_MATH_CONVERTER_RECEIPT"
    assert receipt["status"] == "PASS"
    assert receipt["native_target"] == {
        "kind": "office_math",
        "wrapper": "a14:m",
        "omml_root": "m:oMathPara",
    }
    assert len(receipt["latex_sha256"]) == 64
    assert len(receipt["mathml_sha256"]) == 64
    assert len(receipt["omml_sha256"]) == 64
    assert receipt["semantic_omml_profile"] == "office-math-semantic-v2"
    assert len(receipt["semantic_omml_sha256"]) == 64
    omml = etree.fromstring(receipt["artifacts"]["omml_xml"].encode())
    assert omml.tag == f"{{{M_NS}}}oMathPara"
    assert omml.find(".//m:f", namespaces=NS) is not None


def test_semantic_hash_tolerates_only_powerpoint_equivalent_normalization() -> None:
    receipt = compile_formula("EQ", "p<0.05", "inline")
    original = etree.fromstring(receipt["artifacts"]["omml_xml"].encode())
    normalized = etree.fromstring(receipt["artifacts"]["omml_xml"].encode())
    run = normalized.find("m:r", namespaces=NS)
    assert run is not None
    text_node = run.find("m:t", namespaces=NS)
    assert text_node is not None
    text_node.text = "𝑝"
    run_properties = etree.Element(f"{{{A_NS}}}rPr")
    etree.SubElement(run_properties, f"{{{A_NS}}}latin", typeface="Cambria Math")
    run.insert(0, run_properties)
    second_run = etree.Element(f"{{{M_NS}}}r")
    second_run_properties = etree.SubElement(second_run, f"{{{A_NS}}}rPr")
    etree.SubElement(second_run_properties, f"{{{A_NS}}}latin", typeface="Cambria Math")
    etree.SubElement(second_run, f"{{{M_NS}}}t").text = "<0.05"
    normalized.append(second_run)

    assert _semantic_omml_sha256(normalized) == _semantic_omml_sha256(original)
    second_run.find("m:t", namespaces=NS).text = "<0.01"
    assert _semantic_omml_sha256(normalized) != _semantic_omml_sha256(original)


def test_semantic_hash_preserves_math_variants_and_compatibility_characters() -> None:
    italic = compile_formula("I1", "x", "inline")
    explicit_italic = compile_formula("I2", r"\mathit{x}", "inline")
    bold = compile_formula("B", r"\mathbf{x}", "inline")
    script = compile_formula("S", r"\mathcal{x}", "inline")
    double_struck = compile_formula("D", r"\mathbb{x}", "inline")

    def semantic(receipt: dict[str, object]) -> str:
        artifacts = receipt["artifacts"]
        assert isinstance(artifacts, dict)
        root = etree.fromstring(str(artifacts["omml_xml"]).encode())
        return _semantic_omml_sha256(root)

    assert semantic(italic) == semantic(explicit_italic)
    assert len({semantic(italic), semantic(bold), semantic(script), semantic(double_struck)}) == 4

    def literal(value: str) -> etree._Element:
        root = etree.Element(f"{{{M_NS}}}oMath", nsmap={"m": M_NS})
        run = etree.SubElement(root, f"{{{M_NS}}}r")
        etree.SubElement(run, f"{{{M_NS}}}t").text = value
        return root

    assert _semantic_omml_sha256(literal("²")) != _semantic_omml_sha256(literal("2"))
    assert _semantic_omml_sha256(literal("①")) != _semantic_omml_sha256(literal("1"))
    assert _semantic_omml_sha256(literal("½")) != _semantic_omml_sha256(literal("1⁄2"))


def test_semantic_hash_accepts_powerpoint_plain_style_on_encoded_script_glyph() -> None:
    receipt = compile_formula("ELL", r"\ell", "inline")
    original = etree.fromstring(receipt["artifacts"]["omml_xml"].encode())
    normalized = etree.fromstring(receipt["artifacts"]["omml_xml"].encode())
    run = normalized.find("m:r", namespaces=NS)
    assert run is not None
    properties = etree.Element(f"{{{M_NS}}}rPr")
    etree.SubElement(properties, f"{{{M_NS}}}sty").set(f"{{{M_NS}}}val", "p")
    run.insert(0, properties)
    assert _semantic_omml_sha256(normalized) == _semantic_omml_sha256(original)

    etree.SubElement(properties, f"{{{M_NS}}}scr").set(f"{{{M_NS}}}val", "roman")
    assert _semantic_omml_sha256(normalized) != _semantic_omml_sha256(original)


@pytest.mark.parametrize(
    "latex",
    [
        r"\input{secret}",
        r"\write18{cmd}",
        r"\newcommand{\x}{1}",
        r"x\tag{1}",
        r"x\label{eq:x}",
        r"\operatorname{custom}",
        r"\definitelyNotMath{x}",
    ],
)
def test_compile_formula_rejects_dangerous_lossy_or_unknown_commands(latex: str) -> None:
    with pytest.raises(NativeMathError):
        compile_formula("BAD", latex, "inline")


def test_single_character_mathrm_preserves_upright_style() -> None:
    receipt = compile_formula("UPRIGHT", r"\mathrm{x}", "inline")
    omml = etree.fromstring(receipt["artifacts"]["omml_xml"].encode())
    style = omml.find(".//m:rPr/m:sty", namespaces=NS)
    assert style is not None
    assert style.get(f"{{{M_NS}}}val") == "p"


def test_inject_display_and_mixed_math_preserves_order_and_audits(tmp_path: Path) -> None:
    source = make_placeholder_deck(tmp_path / "source.pptx")
    source_before = source.read_bytes()
    display_receipt = compile_formula("EQ1", r"\frac{E}{mc^2}+\alpha_i", "display")
    alpha_receipt = compile_formula("EQ2", r"\alpha_i", "inline")
    p_receipt = compile_formula("EQ3", r"p<0.05", "inline")
    display_path = write_json(tmp_path / "eq1.json", display_receipt)
    alpha_path = write_json(tmp_path / "eq2.json", alpha_receipt)
    p_path = write_json(tmp_path / "eq3.json", p_receipt)
    plan = {
        "schema_version": "1.0",
        "operations": [
            {
                "slide_index": 1,
                "placeholder_name": "formula-EQ1",
                "formula_id": "EQ1",
                "receipt_path": display_path.name,
                "receipt_sha256": sha256_file(display_path),
            },
            {
                "slide_index": 1,
                "placeholder_name": "mixed-T1",
                "runs": [
                    {"kind": "text", "text": "其中 "},
                    {
                        "kind": "math",
                        "formula_id": "EQ2",
                        "receipt_path": alpha_path.name,
                        "receipt_sha256": sha256_file(alpha_path),
                    },
                    {"kind": "text", "text": " 是学习率，且 "},
                    {
                        "kind": "math",
                        "formula_id": "EQ3",
                        "receipt_path": p_path.name,
                        "receipt_sha256": sha256_file(p_path),
                    },
                ],
            },
        ],
    }
    plan_path = write_json(tmp_path / "plan.json", plan)
    output = tmp_path / "native.pptx"

    report = inject_plan(source, plan_path, output)

    assert report["status"] == "INJECTED_REQUIRES_POWERPOINT_ROUNDTRIP"
    assert source.read_bytes() == source_before
    assert output.is_file()
    assert slide_sequence(source) == ["before", "formula-EQ1", "mixed-T1", "after"]
    assert slide_sequence(output) == ["before", "formula-EQ1", "mixed-T1", "after"]

    with zipfile.ZipFile(output) as package:
        root = etree.fromstring(package.read("ppt/slides/slide1.xml"))
    assert len(root.xpath(".//a14:m", namespaces=NS)) == 3
    assert len(root.xpath(".//m:oMathPara", namespaces=NS)) == 1
    assert len(root.xpath(".//m:oMath", namespaces=NS)) >= 3
    assert not root.xpath(".//mc:AlternateContent/mc:Fallback//p:pic", namespaces=NS)
    assert not root.xpath(".//mc:AlternateContent/mc:Fallback//a:blip", namespaces=NS)
    assert root.xpath(".//a14:m//m:r/a:rPr/a:latin[@typeface='Cambria Math']", namespaces=NS)
    assert A14_NS.encode() in etree.tostring(root)

    audit = audit_pptx(output, plan_path)
    assert audit["status"] == "STRUCTURE_PASS_REQUIRES_POWERPOINT_FINALIZE"
    assert not audit["findings"]
    assert audit["native_formula_count"] == 3
    assert audit["native_shape_count"] == 2


def test_audit_accepts_powerpoint_standalone_inline_wrapper_normalization(
    tmp_path: Path,
) -> None:
    source = make_placeholder_deck(tmp_path / "source.pptx")
    receipt_path = write_json(
        tmp_path / "inline.json", compile_formula("INLINE", r"f_{\mathrm{map}}", "inline")
    )
    plan_path = write_json(
        tmp_path / "plan.json",
        {
            "schema_version": "1.0",
            "operations": [
                {
                    "slide_index": 1,
                    "placeholder_name": "formula-EQ1",
                    "formula_id": "INLINE",
                    "receipt_path": receipt_path.name,
                    "receipt_sha256": sha256_file(receipt_path),
                }
            ],
        },
    )
    native = tmp_path / "native.pptx"
    inject_plan(source, plan_path, native)

    def wrap_as_powerpoint_paragraph(root: etree._Element) -> None:
        wrapper = root.xpath(
            ".//mc:Choice[@Requires='a14']/p:sp"
            "[p:nvSpPr/p:cNvPr[@name='formula-EQ1']]"
            "/p:txBody/a:p/a14:m",
            namespaces=NS,
        )[0]
        inline_root = wrapper.find("m:oMath", namespaces=NS)
        assert inline_root is not None
        wrapper.remove(inline_root)
        paragraph_root = etree.SubElement(wrapper, f"{{{M_NS}}}oMathPara")
        paragraph_root.append(inline_root)

    normalized = rewrite_first_slide(
        native, tmp_path / "powerpoint-normalized.pptx", wrap_as_powerpoint_paragraph
    )
    audit = audit_pptx(normalized, plan_path)
    assert audit["status"] == "STRUCTURE_PASS_REQUIRES_POWERPOINT_FINALIZE"
    assert not audit["findings"]
    formula = audit["inventory"][0]["formulas"][0]
    assert formula["root"] == "m:oMathPara"
    assert formula["powerpoint_mode_normalized"] is True
    receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
    assert formula["semantic_omml_sha256"] == receipt["semantic_omml_sha256"]


def test_plain_textbox_cannot_masquerade_as_native_formula(tmp_path: Path) -> None:
    source = make_placeholder_deck(tmp_path / "plain.pptx")
    audit = audit_pptx(source)
    assert audit["status"] == "FAIL"
    assert audit["native_formula_count"] == 0
    assert {finding["code"] for finding in audit["findings"]} == {"NO_NATIVE_OFFICE_MATH"}


def test_slide_index_follows_presentation_order_not_part_number(tmp_path: Path) -> None:
    source = make_reordered_placeholder_deck(tmp_path / "reordered.pptx")
    receipt_path = write_json(
        tmp_path / "eq.json", compile_formula("EQ1", r"x^2+y^2", "display")
    )
    plan_path = write_json(
        tmp_path / "plan.json",
        {
            "schema_version": "1.0",
            "operations": [
                {
                    "slide_index": 1,
                    "placeholder_name": "formula-EQ1",
                    "formula_id": "EQ1",
                    "receipt_path": receipt_path.name,
                    "receipt_sha256": sha256_file(receipt_path),
                }
            ]
        },
    )
    output = tmp_path / "native.pptx"
    inject_plan(source, plan_path, output)

    audit = audit_pptx(output, plan_path)
    assert audit["status"] == "STRUCTURE_PASS_REQUIRES_POWERPOINT_FINALIZE"
    assert not audit["findings"]
    assert audit["inventory"][0]["slide_index"] == 1


def test_injection_refuses_duplicate_or_missing_placeholder(tmp_path: Path) -> None:
    source = make_placeholder_deck(tmp_path / "source.pptx")
    receipt = compile_formula("EQ1", "x^2", "display")
    receipt_path = write_json(tmp_path / "eq.json", receipt)
    plan = {
        "schema_version": "1.0",
        "operations": [
            {
                "slide_index": 1,
                "placeholder_name": "missing",
                "formula_id": "EQ1",
                "receipt_path": receipt_path.name,
                "receipt_sha256": sha256_file(receipt_path),
            }
        ]
    }
    plan_path = write_json(tmp_path / "plan.json", plan)
    with pytest.raises(NativeMathError, match="matched 0"):
        inject_plan(source, plan_path, tmp_path / "out.pptx")


def test_converter_receipt_tampering_is_rejected(tmp_path: Path) -> None:
    source = make_placeholder_deck(tmp_path / "source.pptx")
    receipt = compile_formula("EQ1", "x^2", "display")
    receipt["canonical_latex"] = "x^3"
    receipt_path = write_json(tmp_path / "tampered.json", receipt)
    plan_path = write_json(
        tmp_path / "plan.json",
        {
            "schema_version": "1.0",
            "operations": [
                {
                    "slide_index": 1,
                    "placeholder_name": "formula-EQ1",
                    "formula_id": "EQ1",
                    "receipt_path": receipt_path.name,
                    "receipt_sha256": sha256_file(receipt_path),
                }
            ]
        },
    )
    with pytest.raises(NativeMathError, match="LaTeX hash mismatch"):
        inject_plan(source, plan_path, tmp_path / "out.pptx")


def test_audit_uses_external_receipt_not_self_signed_shape_metadata(tmp_path: Path) -> None:
    native, plan_path = inject_single_formula(tmp_path, "x")

    def self_sign_changed_math(root: etree._Element) -> None:
        choice = root.xpath(
            ".//mc:Choice[@Requires='a14']/p:sp[p:nvSpPr/p:cNvPr[@name='formula-EQ1']]",
            namespaces=NS,
        )[0]
        text_node = choice.find(".//m:t", namespaces=NS)
        assert text_node is not None
        text_node.text = "y"
        omml = choice.find(".//m:oMathPara", namespaces=NS)
        assert omml is not None
        forged_semantic_hash = _semantic_omml_sha256(omml)
        c_nv_pr = choice.find("./p:nvSpPr/p:cNvPr", namespaces=NS)
        assert c_nv_pr is not None
        metadata = _decode_metadata(str(c_nv_pr.get("descr", "")))
        assert metadata is not None
        metadata["formulas"][0]["semantic_omml_sha256"] = forged_semantic_hash
        metadata["content_runs"][0]["semantic_omml_sha256"] = forged_semantic_hash
        c_nv_pr.set("descr", _metadata_value(metadata))

    forged = rewrite_first_slide(native, tmp_path / "forged.pptx", self_sign_changed_math)
    audit = audit_pptx(forged, plan_path)
    assert audit["status"] == "FAIL"
    assert "NATIVE_SEMANTIC_OMML_HASH_MISMATCH" in {
        finding["code"] for finding in audit["findings"]
    }


def test_audit_rejects_math_outside_direct_paragraph_position(tmp_path: Path) -> None:
    native, plan_path = inject_single_formula(tmp_path)

    def move_math(root: etree._Element) -> None:
        choice = root.xpath(".//mc:Choice[@Requires='a14']/p:sp", namespaces=NS)[0]
        paragraph = choice.find("./p:txBody/a:p", namespaces=NS)
        shape_properties = choice.find("./p:spPr", namespaces=NS)
        wrapper = choice.find("./p:txBody/a:p/a14:m", namespaces=NS)
        assert paragraph is not None and shape_properties is not None and wrapper is not None
        paragraph.remove(wrapper)
        shape_properties.append(wrapper)

    invalid = rewrite_first_slide(native, tmp_path / "invalid-location.pptx", move_math)
    audit = audit_pptx(invalid, plan_path)
    assert audit["status"] == "FAIL"
    assert "NATIVE_MATH_INVALID_OOXML_LOCATION" in {
        finding["code"] for finding in audit["findings"]
    }


def test_audit_requires_external_expected_plan(tmp_path: Path) -> None:
    native, _ = inject_single_formula(tmp_path)
    audit = audit_pptx(native)
    assert audit["status"] == "FAIL"
    assert "EXPECTED_PLAN_REQUIRED" in {finding["code"] for finding in audit["findings"]}


def test_detached_roundtrip_receipt_can_never_authorize_final_pass(tmp_path: Path) -> None:
    native, plan_path = inject_single_formula(tmp_path)
    receipt = {
        "document_type": "POWERPOINT_NATIVE_MATH_ROUNDTRIP_RECEIPT",
        "schema_version": "2.0",
        "status": "OBSERVED_PASS",
        "challenge": "0" * 64,
        "output_pptx": str(native.resolve()),
        "output_sha256": sha256_file(native),
        "expected_plan_sha256": sha256_file(plan_path),
        "powerpoint_version": "fabricated-no-COM",
        "math_shapes": [],
    }
    receipt_path = write_json(tmp_path / "roundtrip.json", receipt)
    audit = audit_pptx(native, plan_path, receipt_path)
    assert audit["status"] == "STRUCTURE_PASS_REQUIRES_POWERPOINT_FINALIZE"
    assert not audit["findings"]
    assert audit["powerpoint_roundtrip"]["trusted_for_final_pass"] is False


def test_audit_rejects_deleted_mixed_text_run(tmp_path: Path) -> None:
    source = make_placeholder_deck(tmp_path / "source.pptx")
    receipt_path = write_json(tmp_path / "eq.json", compile_formula("EQ2", r"\alpha_i", "inline"))
    plan_path = write_json(
        tmp_path / "plan.json",
        {
            "schema_version": "1.0",
            "operations": [
                {
                    "slide_index": 1,
                    "placeholder_name": "mixed-T1",
                    "runs": [
                        {"kind": "text", "text": "其中 "},
                        {
                            "kind": "math",
                            "formula_id": "EQ2",
                            "receipt_path": receipt_path.name,
                            "receipt_sha256": sha256_file(receipt_path),
                        },
                        {"kind": "text", "text": " 是学习率"},
                    ],
                }
            ]
        },
    )
    native = tmp_path / "native.pptx"
    inject_plan(source, plan_path, native)

    def delete_first_text(root: etree._Element) -> None:
        paragraph = root.xpath(
            ".//mc:Choice[@Requires='a14']/p:sp[p:nvSpPr/p:cNvPr[@name='mixed-T1']]"
            "/p:txBody/a:p",
            namespaces=NS,
        )[0]
        first_text = paragraph.find("a:r", namespaces=NS)
        assert first_text is not None
        paragraph.remove(first_text)

    broken = rewrite_first_slide(native, tmp_path / "missing-text.pptx", delete_first_text)
    audit = audit_pptx(broken, plan_path)
    assert audit["status"] == "FAIL"
    assert "NATIVE_CONTENT_RUN_SEQUENCE_MISMATCH" in {
        finding["code"] for finding in audit["findings"]
    }


def test_audit_accepts_powerpoint_fallback_metadata_normalization_but_not_text_drift(
    tmp_path: Path,
) -> None:
    native, plan_path = inject_single_formula(tmp_path)

    def copy_choice_metadata_to_fallback(root: etree._Element) -> None:
        alternate = root.xpath(".//mc:AlternateContent", namespaces=NS)[0]
        choice_properties = alternate.find(
            "./mc:Choice/p:sp/p:nvSpPr/p:cNvPr", namespaces=NS
        )
        fallback_properties = alternate.find(
            "./mc:Fallback/p:sp/p:nvSpPr/p:cNvPr", namespaces=NS
        )
        fallback_text = alternate.find(
            "./mc:Fallback/p:sp/p:txBody/a:p/a:r/a:t", namespaces=NS
        )
        assert (
            choice_properties is not None
            and fallback_properties is not None
            and fallback_text is not None
        )
        fallback_properties.set("descr", str(choice_properties.get("descr")))
        fallback_text.text = "\u00a0"

    normalized = rewrite_first_slide(
        native, tmp_path / "normalized-fallback.pptx", copy_choice_metadata_to_fallback
    )
    normalized_audit = audit_pptx(normalized, plan_path)
    assert normalized_audit["status"] == "STRUCTURE_PASS_REQUIRES_POWERPOINT_FINALIZE"
    assert normalized_audit["inventory"][0][
        "powerpoint_fallback_metadata_normalized"
    ] is True

    def corrupt_fallback_text(root: etree._Element) -> None:
        alternate = root.xpath(".//mc:AlternateContent", namespaces=NS)[0]
        choice_properties = alternate.find(
            "./mc:Choice/p:sp/p:nvSpPr/p:cNvPr", namespaces=NS
        )
        fallback_properties = alternate.find(
            "./mc:Fallback/p:sp/p:nvSpPr/p:cNvPr", namespaces=NS
        )
        fallback_text = alternate.find("./mc:Fallback/p:sp/p:txBody/a:p/a:r/a:t", namespaces=NS)
        assert (
            choice_properties is not None
            and fallback_properties is not None
            and fallback_text is not None
        )
        fallback_properties.set("descr", str(choice_properties.get("descr")))
        fallback_text.text = "forged fallback"

    corrupted = rewrite_first_slide(
        native, tmp_path / "corrupt-fallback.pptx", corrupt_fallback_text
    )
    corrupted_audit = audit_pptx(corrupted, plan_path)
    assert corrupted_audit["status"] == "FAIL"
    assert "NATIVE_MATH_FALLBACK_INVALID" in {
        finding["code"] for finding in corrupted_audit["findings"]
    }


def test_plan_rejects_duplicate_json_keys_and_duplicate_targets(tmp_path: Path) -> None:
    source = make_placeholder_deck(tmp_path / "source.pptx")
    duplicate_key_plan = tmp_path / "duplicate-key.json"
    duplicate_key_plan.write_text(
        '{"schema_version":"1.0","schema_version":"1.0","operations":[]}',
        encoding="utf-8",
    )
    with pytest.raises(NativeMathError, match="duplicate JSON object key"):
        inject_plan(source, duplicate_key_plan, tmp_path / "out-key.pptx")

    receipt_path = write_json(tmp_path / "eq.json", compile_formula("EQ1", "x", "display"))
    operation = {
        "slide_index": 1,
        "placeholder_name": "formula-EQ1",
        "formula_id": "EQ1",
        "receipt_path": receipt_path.name,
        "receipt_sha256": sha256_file(receipt_path),
    }
    duplicate_target_plan = write_json(
        tmp_path / "duplicate-target.json",
        {"schema_version": "1.0", "operations": [operation, operation]},
    )
    with pytest.raises(NativeMathError, match="duplicate plan target"):
        inject_plan(source, duplicate_target_plan, tmp_path / "out-target.pptx")


def test_plan_binds_runtime_target_font_style_and_rejects_partial_style(
    tmp_path: Path,
) -> None:
    source = make_placeholder_deck(tmp_path / "source.pptx")
    receipt_path = write_json(
        tmp_path / "eq.json", compile_formula("EQ1", r"z_t^\tau", "display")
    )
    operation = {
        "slide_index": 1,
        "placeholder_name": "formula-EQ1",
        "formula_id": "EQ1",
        "receipt_path": receipt_path.name,
        "receipt_sha256": sha256_file(receipt_path),
        "target_font_size_pt": 9.75,
        "target_font_color": "#17324D",
    }
    plan = write_json(
        tmp_path / "styled-plan.json",
        {"schema_version": "1.0", "operations": [operation]},
    )
    styled_native = tmp_path / "styled-native.pptx"
    report = inject_plan(source, plan, styled_native)
    assert report["operations"][0]["target_font_size_pt"] == 9.75
    assert report["operations"][0]["target_font_color"] == "#17324D"
    with zipfile.ZipFile(styled_native, "r") as package:
        root = etree.fromstring(package.read("ppt/slides/slide1.xml"))
    body_properties = root.xpath(
        ".//mc:AlternateContent[mc:Choice/p:sp/p:nvSpPr/p:cNvPr[@name='formula-EQ1']]"
        "//a:bodyPr|"
        ".//mc:AlternateContent[mc:Fallback/p:sp/p:nvSpPr/p:cNvPr[@name='formula-EQ1']]"
        "//a:bodyPr",
        namespaces=NS,
    )
    assert len(body_properties) == 2
    assert all(
        properties.get(inset_name) == "0"
        for properties in body_properties
        for inset_name in ("lIns", "rIns", "tIns", "bIns")
    )
    styled_properties = root.xpath(
        ".//mc:AlternateContent[mc:Choice/p:sp/p:nvSpPr/p:cNvPr[@name='formula-EQ1']]"
        "//a:rPr|"
        ".//mc:AlternateContent[mc:Choice/p:sp/p:nvSpPr/p:cNvPr[@name='formula-EQ1']]"
        "//a:endParaRPr",
        namespaces=NS,
    )
    assert styled_properties
    assert {properties.get("sz") for properties in styled_properties} == {"975"}
    assert {
        value
        for properties in styled_properties
        for value in properties.xpath("./a:solidFill/a:srgbClr/@val", namespaces=NS)
    } == {"17324D"}

    partial_plan = write_json(
        tmp_path / "partial-style-plan.json",
        {
            "schema_version": "1.0",
            "operations": [
                {
                    key: value
                    for key, value in operation.items()
                    if key != "target_font_color"
                }
            ],
        },
    )
    with pytest.raises(NativeMathError, match="both target font style fields"):
        inject_plan(source, partial_plan, tmp_path / "partial-style-native.pptx")


def test_injection_and_finalize_require_fresh_output_paths(tmp_path: Path) -> None:
    source = make_placeholder_deck(tmp_path / "source.pptx")
    plan = tmp_path / "plan.json"
    report = write_json(tmp_path / "report.json", {})
    output = tmp_path / "existing.pptx"
    output.write_bytes(b"do-not-overwrite")
    plan.write_text('{"schema_version":"1.0","operations":[]}', encoding="utf-8")

    with pytest.raises(NativeMathError, match="output already exists"):
        inject_plan(source, plan, output)
    with pytest.raises(NativeMathError, match="fresh finalization refuses overwrite"):
        finalize_native_math(
            source,
            plan,
            report,
            output,
            tmp_path / "receipt.json",
            tmp_path / "renders",
            overwrite=True,
        )


def test_audit_rejects_hidden_or_off_canvas_native_math(tmp_path: Path) -> None:
    native, plan_path = inject_single_formula(tmp_path)

    def hide_and_move(root: etree._Element) -> None:
        choice = root.xpath(".//mc:Choice[@Requires='a14']/p:sp", namespaces=NS)[0]
        nonvisual = choice.find("./p:nvSpPr/p:cNvPr", namespaces=NS)
        offset = choice.find("./p:spPr/a:xfrm/a:off", namespaces=NS)
        assert nonvisual is not None and offset is not None
        nonvisual.set("hidden", "1")
        offset.set("x", "-1")

    broken = rewrite_first_slide(native, tmp_path / "hidden.pptx", hide_and_move)
    audit = audit_pptx(broken, plan_path)
    codes = {finding["code"] for finding in audit["findings"]}
    assert "NATIVE_MATH_HIDDEN" in codes
    assert "NATIVE_MATH_OUTSIDE_CANVAS" in codes


def test_audit_rejects_plain_text_formula_candidate_elsewhere_on_slide(tmp_path: Path) -> None:
    native, plan_path = inject_single_formula(tmp_path)

    def add_plain_formula(root: etree._Element) -> None:
        shape = root.xpath(".//p:sp[p:nvSpPr/p:cNvPr[@name='mixed-T1']]", namespaces=NS)[0]
        text = shape.find("./p:txBody/a:p/a:r/a:t", namespaces=NS)
        assert text is not None
        text.text = "p < 0.05"

    broken = rewrite_first_slide(native, tmp_path / "plain-formula.pptx", add_plain_formula)
    audit = audit_pptx(broken, plan_path)
    assert "PLAIN_TEXT_FORMULA_MASQUERADE_RISK" in {
        finding["code"] for finding in audit["findings"]
    }


def test_counterfactual_pixel_evidence_requires_local_formula_ink_delta(
    tmp_path: Path,
) -> None:
    baseline_path = tmp_path / "baseline.png"
    control_path = tmp_path / "control.png"
    outside_path = tmp_path / "outside.png"
    Image.new("RGB", (100, 100), "white").save(control_path)

    baseline = Image.new("RGB", (100, 100), "white")
    ImageDraw.Draw(baseline).rectangle((22, 22, 29, 29), fill="black")
    baseline.save(baseline_path)
    evidence = _counterfactual_pixel_evidence(
        baseline_path,
        control_path,
        left=20,
        top=20,
        shape_width=20,
        shape_height=20,
        slide_width=100,
        slide_height=100,
    )
    assert evidence["pass"] is True
    assert evidence["inside_changed_pixels"] >= evidence["required_changed_pixels"]
    assert evidence["outside_changed_pixels"] == 0

    identical = _counterfactual_pixel_evidence(
        control_path,
        control_path,
        left=20,
        top=20,
        shape_width=20,
        shape_height=20,
        slide_width=100,
        slide_height=100,
    )
    assert identical["pass"] is False

    outside = Image.new("RGB", (100, 100), "white")
    ImageDraw.Draw(outside).rectangle((70, 70, 79, 79), fill="black")
    outside.save(outside_path)
    outside_evidence = _counterfactual_pixel_evidence(
        outside_path,
        control_path,
        left=20,
        top=20,
        shape_width=20,
        shape_height=20,
        slide_width=100,
        slide_height=100,
    )
    assert outside_evidence["pass"] is False
    assert outside_evidence["inside_changed_pixels"] == 0
    assert outside_evidence["outside_changed_pixels"] > 0


def test_counterfactual_pixel_evidence_allows_resolution_scaled_glyph_overhang(
    tmp_path: Path,
) -> None:
    baseline_path = tmp_path / "baseline-large.png"
    control_path = tmp_path / "control-large.png"
    Image.new("RGB", (1000, 1000), "white").save(control_path)
    baseline = Image.new("RGB", (1000, 1000), "white")
    # The ink starts six pixels left of the nominal text box.  That is outside
    # the former fixed two-pixel allowance but within the 1% render halo.
    ImageDraw.Draw(baseline).rectangle((194, 240, 201, 249), fill="black")
    baseline.save(baseline_path)
    evidence = _counterfactual_pixel_evidence(
        baseline_path,
        control_path,
        left=200,
        top=200,
        shape_width=100,
        shape_height=100,
        slide_width=1000,
        slide_height=1000,
    )
    assert evidence["padding_pixels"] == 10
    assert evidence["pass"] is True
    assert evidence["outside_changed_pixels"] == 0


@pytest.mark.skipif(os.name != "nt", reason="PowerPoint finalizer is Windows-only")
def test_roundtrip_render_commit_refuses_raced_destination_and_preserves_sentinel(
    tmp_path: Path,
) -> None:
    script_text = (
        Path(__file__).parents[1] / "tools" / "powerpoint_native_math_roundtrip.ps1"
    ).read_text(encoding="utf-8")
    assert "[System.IO.Directory]::Move($stagingRender, $renderFull)" in script_text
    assert "Remove-Item -LiteralPath $renderFull -Recurse" not in script_text

    staging = tmp_path / "staging"
    destination = tmp_path / "raced-destination"
    staging.mkdir()
    destination.mkdir()
    (staging / "render.png").write_bytes(b"transaction-owned")
    sentinel = destination / "user-sentinel.txt"
    sentinel.write_text("must-survive", encoding="utf-8")
    probe = tmp_path / "directory-move.ps1"
    probe.write_text(
        "param([string]$Source, [string]$Destination)\n"
        "$ErrorActionPreference = 'Stop'\n"
        "[System.IO.Directory]::Move($Source, $Destination)\n",
        encoding="utf-8",
    )
    result = subprocess.run(
        [
            _powershell_executable(),
            "-NoProfile",
            "-NonInteractive",
            "-ExecutionPolicy",
            "Bypass",
            "-File",
            str(probe),
            str(staging),
            str(destination),
        ],
        capture_output=True,
        text=True,
        check=False,
    )
    assert result.returncode != 0
    assert sentinel.read_text(encoding="utf-8") == "must-survive"
    assert (staging / "render.png").read_bytes() == b"transaction-owned"


def test_roundtrip_underlay_resolution_skips_only_tiny_opaque_boundary_contacts() -> None:
    script_text = (
        Path(__file__).parents[1] / "tools" / "powerpoint_native_math_roundtrip.ps1"
    ).read_text(encoding="utf-8")

    picture_guard = "if ($lowerShape.is_picture -or $lowerShape.is_ole)"
    boundary_guard = (
        "if ($isReliableSolidFill -and $underlayOverlap -lt 0.02) { continue }"
    )
    material_partial_guard = (
        "if ($underlayOverlap -lt 0.995 -or -not $isReliableSolidFill)"
    )

    assert picture_guard in script_text
    assert boundary_guard in script_text
    assert material_partial_guard in script_text
    assert script_text.index(picture_guard) < script_text.index(boundary_guard)
    assert script_text.index(boundary_guard) < script_text.index(material_partial_guard)


def test_roundtrip_verifies_hash_bound_formula_styles_without_mutating_mathzones() -> None:
    script_text = (
        Path(__file__).parents[1] / "tools" / "powerpoint_native_math_roundtrip.ps1"
    ).read_text(encoding="utf-8")

    assert "$observedSize = [double]$styleZones.Font.Size" in script_text
    assert "$styleZones.Font.Size =" not in script_text
    assert "$styleZones.Font.Fill.ForeColor.RGB =" not in script_text
    assert "formula_styles_verified_at_open = $false" in script_text
    assert "NATIVE_MATH_STYLE_MISMATCH" in script_text


@pytest.mark.skipif(os.name != "nt", reason="PowerShell output policy is Windows-only")
@pytest.mark.parametrize("forbidden_parameter", ["OutputPath", "ReceiptPath", "RenderDirectory"])
def test_roundtrip_helper_rejects_project_outputs_before_writing(
    tmp_path: Path,
    forbidden_parameter: str,
) -> None:
    project_root = Path(__file__).parents[1]
    script = project_root / "tools" / "powerpoint_native_math_roundtrip.ps1"
    values = {
        "InputPath": str(tmp_path / "missing-input.pptx"),
        "OutputPath": str(tmp_path / "output.pptx"),
        "ExpectedPlanPath": str(tmp_path / "missing-plan.json"),
        "InjectionReportPath": str(tmp_path / "missing-report.json"),
        "ReceiptPath": str(tmp_path / "receipt.json"),
        "RenderDirectory": str(tmp_path / "renders"),
    }
    forbidden = project_root / f"forbidden-{forbidden_parameter.lower()}"
    values[forbidden_parameter] = str(forbidden)
    command = [
        _powershell_executable(),
        "-NoProfile",
        "-NonInteractive",
        "-ExecutionPolicy",
        "Bypass",
        "-File",
        str(script),
    ]
    for key, value in values.items():
        command.extend([f"-{key}", value])
    command.extend(["-Challenge", "0" * 64, "-ParentProcessId", str(os.getpid())])

    result = subprocess.run(command, capture_output=True, text=True, check=False)
    assert result.returncode != 0
    assert "must be under examples\\generated" in (result.stdout + result.stderr)
    assert not forbidden.exists()


@pytest.mark.skipif(os.name != "nt", reason="Win32 device aliases are Windows-only")
def test_roundtrip_helper_rejects_device_alias_without_touching_fixture(tmp_path: Path) -> None:
    project_root = Path(__file__).parents[1]
    fixture = project_root / "examples" / "target_figure.png"
    before = sha256_file(fixture)
    device_alias = "\\\\?\\" + str(fixture)
    script = project_root / "tools" / "powerpoint_native_math_roundtrip.ps1"
    result = subprocess.run(
        [
            _powershell_executable(),
            "-NoProfile",
            "-NonInteractive",
            "-ExecutionPolicy",
            "Bypass",
            "-File",
            str(script),
            "-InputPath",
            str(tmp_path / "missing-input.pptx"),
            "-OutputPath",
            device_alias,
            "-ExpectedPlanPath",
            str(tmp_path / "missing-plan.json"),
            "-InjectionReportPath",
            str(tmp_path / "missing-report.json"),
            "-ReceiptPath",
            str(tmp_path / "receipt.json"),
            "-RenderDirectory",
            str(tmp_path / "renders"),
            "-Challenge",
            "0" * 64,
            "-ParentProcessId",
            str(os.getpid()),
        ],
        capture_output=True,
        text=True,
        check=False,
    )
    assert result.returncode != 0
    assert "Win32 or NT device namespace" in (result.stdout + result.stderr)
    assert sha256_file(fixture) == before


def test_fresh_audit_commit_preserves_destination_created_during_race(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    destination = tmp_path / "audit.json"
    real_link = os.link

    def create_sentinel_then_link(source: object, target: object) -> None:
        Path(target).write_bytes(b"concurrent-user-sentinel")
        real_link(source, target)

    monkeypatch.setattr(os, "link", create_sentinel_then_link)
    with pytest.raises(NativeMathError, match="fresh final audit output already exists"):
        _atomic_write_json_fresh(destination, {"status": "mechanical"}, pretty=True)
    assert destination.read_bytes() == b"concurrent-user-sentinel"
    assert not list(tmp_path.glob(".audit.json.*.tmp"))


def test_injection_commit_preserves_pptx_created_during_race(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = make_placeholder_deck(tmp_path / "source.pptx")
    receipt_path = write_json(tmp_path / "eq.json", compile_formula("EQ1", "x", "display"))
    plan_path = write_json(
        tmp_path / "plan.json",
        {
            "schema_version": "1.0",
            "operations": [
                {
                    "slide_index": 1,
                    "placeholder_name": "formula-EQ1",
                    "formula_id": "EQ1",
                    "receipt_path": receipt_path.name,
                    "receipt_sha256": sha256_file(receipt_path),
                }
            ],
        },
    )
    output = tmp_path / "raced-output.pptx"
    sentinel_bytes = b"concurrent-pptx-sentinel"
    real_link = os.link

    def create_sentinel_then_link(source_path: object, target_path: object) -> None:
        if Path(target_path) == output:
            output.write_bytes(sentinel_bytes)
        real_link(source_path, target_path)

    monkeypatch.setattr(os, "link", create_sentinel_then_link)
    with pytest.raises(NativeMathError, match="output appeared during rewrite"):
        inject_plan(source, plan_path, output)
    assert output.read_bytes() == sentinel_bytes
    assert not list(tmp_path.glob(".raced-output.pptx.*.tmp"))


def test_cli_evidence_outputs_are_fresh_by_default_and_require_explicit_overwrite(
    tmp_path: Path,
) -> None:
    compiled = tmp_path / "compiled.json"
    compiled.write_bytes(b"compile-sentinel")
    assert (
        main(
            [
                "compile",
                "--formula-id",
                "EQ1",
                "--latex",
                "x",
                "--mode",
                "display",
                "--output",
                str(compiled),
            ]
        )
        == 3
    )
    assert compiled.read_bytes() == b"compile-sentinel"
    assert (
        main(
            [
                "compile",
                "--formula-id",
                "EQ1",
                "--latex",
                "x",
                "--mode",
                "display",
                "--output",
                str(compiled),
                "--overwrite",
            ]
        )
        == 0
    )

    case_directory = tmp_path / "case"
    case_directory.mkdir()
    native, plan_path = inject_single_formula(case_directory)
    audit_output = tmp_path / "audit.json"
    audit_output.write_bytes(b"audit-sentinel")
    audit_arguments = [
        "audit",
        "--input",
        str(native),
        "--plan",
        str(plan_path),
        "--output",
        str(audit_output),
    ]
    assert main(audit_arguments) == 3
    assert audit_output.read_bytes() == b"audit-sentinel"
    assert main([*audit_arguments, "--overwrite"]) == 0

    injection_output = tmp_path / "cli-native.pptx"
    report_output = tmp_path / "injection-report.json"
    report_output.write_bytes(b"report-sentinel")
    injection_arguments = [
        "inject",
        "--input",
        str(tmp_path / "case" / "source.pptx"),
        "--plan",
        str(plan_path),
        "--output",
        str(injection_output),
        "--report",
        str(report_output),
    ]
    assert main(injection_arguments) == 3
    assert not injection_output.exists()
    assert report_output.read_bytes() == b"report-sentinel"
    assert main([*injection_arguments, "--overwrite-report"]) == 0
