"""autofigure math — 把 convert 产出的公式文本框批量升级为原生 Office Math（OMML）。

检测（两路命中都收）：
- 强信号：任一 run 的 a:rPr@baseline ∈ {30000(super), -25000(sub)}（convert 对 baseline-shift 的落盘约定）
- 弱信号：全部 run 斜体、去空格后 1..4 字符、至少含一个非 ASCII 数学字母（τ π 𝔼 ℒ ∇ ƒ Î θ 等）；
  纯 ASCII 短标签（B / I / GT Answers / Questions）必须排除

流程：检测 → run 序列重建 LaTeX（baseline → ^{}/_{} 分组，Unicode → LaTeX 命令映射）→
逐公式 compile_formula（薄封装 legacy 引擎 tools/powerpoint_native_math.py，单个失败只
warn 跳过、保留原文本框）→ 在临时副本中把命中形状改名 math:NNN → inject_plan 到临时
pptx → os.replace 原子覆盖 redraw.pptx → 保存重开并刷新 v3 scene/bindings 哈希 → 刷新
fresh render（COM 失败只 warn）。
--dry-run 只检测与重建并写 qa/math-summary.json，不改 PPTX。
"""

from __future__ import annotations

import argparse
import json
import os
import shutil
import sys
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

from tools.v2 import common

BASELINE_SUPER = "30000"  # convert 对 baseline-shift="super" 的落盘值（tools/v2/convert.py）
BASELINE_SUB = "-25000"  # baseline-shift="sub"
NAME_PREFIX = "math:"  # 命中形状的改名前缀（inject_plan 按 cNvPr@name 精确定位）
WEAK_MAX_CHARS = 4  # 弱信号：去空格后的最大字符数

RunInfo = tuple[str, bool, str | None]  # (text, italic, baseline)

# Unicode → LaTeX 映射。查不到的非 ASCII 字符保留原字符（latex2mathml 能吃 Unicode）。
_UNICODE_MAP = {
    # 希腊小写
    "α": r"\alpha", "β": r"\beta", "γ": r"\gamma", "δ": r"\delta",
    "ε": r"\varepsilon", "ζ": r"\zeta", "η": r"\eta", "θ": r"\theta",
    "ι": r"\iota", "κ": r"\kappa", "λ": r"\lambda", "μ": r"\mu",
    "ν": r"\nu", "ξ": r"\xi", "π": r"\pi", "ρ": r"\rho",
    "σ": r"\sigma", "τ": r"\tau", "υ": r"\upsilon", "φ": r"\phi",
    "χ": r"\chi", "ψ": r"\psi", "ω": r"\omega",
    # 希腊大写
    "Γ": r"\Gamma", "Δ": r"\Delta", "Θ": r"\Theta", "Λ": r"\Lambda",
    "Ξ": r"\Xi", "Π": r"\Pi", "Σ": r"\Sigma", "Υ": r"\Upsilon",
    "Φ": r"\Phi", "Ψ": r"\Psi", "Ω": r"\Omega",
    # 运算符与符号
    "∇": r"\nabla", "∂": r"\partial", "∞": r"\infty",
    "±": r"\pm", "×": r"\times", "÷": r"\div",
    "≤": r"\leq", "≥": r"\geq", "≠": r"\neq", "≈": r"\approx",
    "∈": r"\in", "∉": r"\notin", "⊂": r"\subset", "⊆": r"\subseteq",
    "∪": r"\cup", "∩": r"\cap", "·": r"\cdot",
    "∑": r"\sum", "∏": r"\prod", "∫": r"\int",
    "⟨": r"\langle", "⟩": r"\rangle", "…": r"\dots",
    "−": "-",  # U+2212 减号 → ASCII 连字符
    # 花体/双线/重音与异形字母
    "𝔼": r"\mathbb{E}", "ℒ": r"\mathcal{L}", "Î": r"\hat{I}", "ƒ": "f",
}

# LaTeX 特殊字符转义（公式框里 _ ^ 只应来自 baseline 结构；run 文本里出现就转义）
_LATEX_ESCAPES = {
    "\\": r"\backslash",
    "_": r"\_", "^": r"\^{}",
    "%": r"\%", "&": r"\&", "#": r"\#", "$": r"\$",
    "{": r"\{", "}": r"\}",
}


# ---------------------------------------------------------------- 纯函数：检测与 LaTeX 重建


def classify_runs(runs: list[RunInfo]) -> str | None:
    """→ 'strong'（含上下标 run）/ 'weak'（全斜体短公式）/ None（非公式）。"""
    meaningful = [(text, italic, baseline) for text, italic, baseline in runs if text.strip()]
    if not meaningful:
        return None
    if any(baseline in (BASELINE_SUPER, BASELINE_SUB) for _, _, baseline in meaningful):
        return "strong"
    text = "".join(text for text, _, _ in meaningful).replace(" ", "")
    if (
        len(text) <= WEAK_MAX_CHARS
        and all(italic for _, italic, _ in meaningful)
        and any(not ch.isascii() for ch in text)
    ):
        return "weak"
    return None


def _map_char(ch: str, nxt: str | None) -> str:
    """单字符 → LaTeX 片段；字母结尾的反斜杠命令后若紧跟 ASCII 字母需补空格（防命令名粘连）。"""
    token = _UNICODE_MAP.get(ch)
    if token is None:
        token = _LATEX_ESCAPES.get(ch, ch)
    if (
        token.startswith("\\")
        and token[-1:].isalpha()
        and token[-1:].isascii()
        and nxt
        and nxt.isalpha()
        and nxt.isascii()
    ):
        return token + " "
    return token


def _latex_text(text: str) -> str:
    return "".join(
        _map_char(ch, text[index + 1] if index + 1 < len(text) else None)
        for index, ch in enumerate(text)
    )


def rebuild_latex(runs: list[RunInfo]) -> str:
    """run 序列 → LaTeX：baseline run 进 ^{}/_{} 分组（连续同向合并进同一组），其余进主体。"""
    parts: list[str] = []
    open_group: str | None = None
    for text, _italic, baseline in runs:
        group = None
        if baseline == BASELINE_SUPER:
            group = "super"
        elif baseline == BASELINE_SUB:
            group = "sub"
        if group != open_group:
            if open_group is not None:
                parts.append("}")
            if group is not None:
                parts.append("^{" if group == "super" else "_{")
            open_group = group
        parts.append(_latex_text(text))
    if open_group is not None:
        parts.append("}")
    return "".join(parts)


# ---------------------------------------------------------------- pptx 侧提取


def _shape_runs(shape) -> list[RunInfo]:
    runs: list[RunInfo] = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            rpr = run._r.find(qn("a:rPr"))
            baseline = rpr.get("baseline") if rpr is not None else None
            runs.append((run.text, run.font.italic is True, baseline))
    return runs


def _target_font(shape) -> tuple[float, str] | None:
    """首个非空 run 的字号/颜色 → (pt, '#RRGGBB')（保持注入后视觉一致）。

    plan 校验要求 target_font_size_pt 与 target_font_color 成对出现且字号 6..72；
    任一拿不到或越界就都不给。
    """
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text.strip():
                continue
            if run.font.size is None:
                return None
            size = float(run.font.size.pt)
            try:
                rgb = run.font.color.rgb
            except AttributeError:  # 未显式设色（_NoneColor）
                return None
            if rgb is None or not 6.0 <= size <= 72.0:
                return None
            return size, f"#{rgb}"
    return None


def _load_engine():
    try:
        from tools import powerpoint_native_math
    except ModuleNotFoundError as exc:
        raise common.fail(
            f"缺少公式引擎依赖: {exc.name}（在 .venv 中 pip install -r requirements-v2.txt）"
        ) from exc
    return powerpoint_native_math


def _echo(text: str) -> None:
    """stdout 跟随控制台代码页（GBK）：公式里的非 ASCII 字符防 UnicodeEncodeError。"""
    encoding = sys.stdout.encoding or "utf-8"
    sys.stdout.write(text.encode(encoding, errors="backslashreplace").decode(encoding) + "\n")


def _pptx_bound_names_and_omml_count(path: Path) -> tuple[set[str], int]:
    """Read names/OMML directly because python-pptx hides AlternateContent shapes."""
    names: set[str] = set()
    omml_count = 0
    with zipfile.ZipFile(path) as package:
        for member in package.namelist():
            if not member.startswith("ppt/slides/slide") or not member.endswith(".xml"):
                continue
            root = ET.fromstring(package.read(member))
            for node in root.iter():
                local_name = node.tag.rsplit("}", 1)[-1]
                if local_name == "cNvPr" and node.get("name"):
                    names.add(node.get("name"))
                elif local_name == "oMath":
                    omml_count += 1
    return names, omml_count


def _refresh_v3_contracts_after_math(run: common.Run, candidates: list[dict]) -> dict:
    """Bind renamed OMML shapes and the new artifact hash back to Scene v3."""
    from tools.v2.contracts import initialize_contracts, read_json, transition, utc_now, write_json

    meta = initialize_contracts(run)
    if meta["workflow"]["state"] != "repairing":
        transition(run, "repairing", "native-math-upgrade-started")

    # High-level reopen proves package compatibility.  OMML formula text boxes
    # live inside mc:AlternateContent and are intentionally invisible to
    # python-pptx's slide.shapes collection, so binding identity is read from
    # the underlying slide XML instead of being falsely reported as missing.
    Presentation(run.pptx_path)
    readback_names, omml_count = _pptx_bound_names_and_omml_count(run.pptx_path)
    rename_map = {
        candidate["name"]: candidate["placeholder"]
        for candidate in candidates
        if candidate["status"] == "injected"
    }
    formula_by_name = {
        candidate["placeholder"]: candidate["formula_id"]
        for candidate in candidates
        if candidate["status"] == "injected"
    }

    bindings = read_json(run.bindings_path)
    rebound_elements: dict[str, str] = {}
    for binding in bindings.get("bindings", []):
        old_name = binding.get("shape_name")
        if old_name in rename_map:
            binding["shape_name"] = rename_map[old_name]
            binding["object_kind"] = "native-math"
            binding["native_math"] = True
            binding["formula_id"] = formula_by_name[binding["shape_name"]]
            rebound_elements[binding["element_id"]] = binding["formula_id"]
        binding["readback_found"] = binding.get("shape_name") in readback_names

    expected = set(formula_by_name)
    rebound = {
        binding.get("shape_name")
        for binding in bindings.get("bindings", [])
        if binding.get("native_math") is True
    }
    missing_rebound = sorted(expected - rebound)
    if missing_rebound:
        raise common.fail(f"OMML shapes have no v3 binding: {missing_rebound}")
    if omml_count < len(expected):
        raise common.fail(
            f"OMML save/reopen count is incomplete: expected {len(expected)}, got {omml_count}"
        )
    bindings_complete = bool(bindings.get("bindings")) and all(
        binding.get("readback_found") is True for binding in bindings["bindings"]
    )
    if not bindings_complete:
        raise common.fail("OMML save/reopen left incomplete PowerPoint shape bindings")

    pptx_hash = common.sha256_file(run.pptx_path)
    bindings.update(
        {
            "updated_at": utc_now(),
            "backend": "pptx-offline+native-math",
            "artifact_sha256": pptx_hash,
            "saved_reopened": True,
            "bindings_complete": True,
        }
    )
    write_json(run.bindings_path, bindings)

    scene = read_json(run.scene_path)
    for element in scene.get("elements", []):
        formula_id = rebound_elements.get(element.get("id"))
        if formula_id:
            element["native_math"] = True
            element["formula_id"] = formula_id
    scene["updated_at"] = utc_now()
    scene["artifact"] = {
        "backend": "pptx-offline+native-math",
        "path": "redraw.pptx",
        "sha256": pptx_hash,
    }
    write_json(run.scene_path, scene)
    from tools.v2.layout import audit_layout

    layout_report = audit_layout(run)
    transition(
        run,
        "candidate",
        "native-math-upgrade-complete",
        details={
            "pptx_sha256": pptx_hash,
            "formula_count": len(expected),
            "object_count": len(bindings["bindings"]),
            "bindings_complete": True,
            "layout_pass": layout_report["pass"],
            "layout_findings": len(layout_report["findings"]),
        },
    )
    return {
        "pptx_sha256": pptx_hash,
        "object_count": len(bindings["bindings"]),
        "omml_count": omml_count,
        "saved_reopened": True,
        "bindings_complete": True,
        "layout_pass": layout_report["pass"],
        "layout_findings": len(layout_report["findings"]),
    }


# ---------------------------------------------------------------- 主流程


def upgrade(run: common.Run, *, dry_run: bool = False) -> dict:
    """检测公式框并批量注入 OMML，返回 summary dict（同时落 qa/math-summary.json）。"""
    if not run.pptx_path.is_file():
        raise common.fail(f"未找到 PPTX: {run.pptx_path}（请先运行 autofigure convert）")
    run.qa_dir.mkdir(exist_ok=True)
    prs = Presentation(run.pptx_path)

    candidates: list[dict] = []
    for slide_index, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            signal = classify_runs(_shape_runs(shape))
            if signal is None:
                continue
            candidates.append({"slide_index": slide_index, "shape": shape, "signal": signal})

    # 编号避开残留 math:NNN 名（上次注入失败后重跑等场景），保证 cNvPr@name 唯一
    candidate_ids = {id(cand["shape"]) for cand in candidates}
    taken_names = {
        shape.name
        for slide in prs.slides
        for shape in slide.shapes
        if id(shape) not in candidate_ids
    }
    number = 0
    for cand in candidates:
        number += 1
        while f"{NAME_PREFIX}{number:03d}" in taken_names:
            number += 1
        placeholder = f"{NAME_PREFIX}{number:03d}"
        taken_names.add(placeholder)
        shape = cand["shape"]
        cand["name"] = shape.name
        cand["placeholder"] = placeholder
        cand["formula_id"] = f"EQ{number:03d}"
        cand["latex"] = rebuild_latex(_shape_runs(shape))
        cand["bold"] = any(
            run.font.bold
            for paragraph in shape.text_frame.paragraphs
            for run in paragraph.runs
        )
        cand["font"] = _target_font(shape)
        cand["status"] = "detected"

    warnings: list[str] = []
    summary = {
        "pptx": "redraw.pptx",
        "dry_run": dry_run,
        "detected": len(candidates),
        "strong": sum(1 for cand in candidates if cand["signal"] == "strong"),
        "weak": sum(1 for cand in candidates if cand["signal"] == "weak"),
        "injected": 0,
        "failed": 0,
        "warnings": warnings,
        "notes": [],
        "formulas": [],
    }

    def write_summary() -> None:
        rows = []
        for cand in candidates:
            row = {
                key: cand[key]
                for key in ("name", "placeholder", "formula_id", "slide_index", "signal", "latex", "status")
            }
            if cand["bold"]:
                row["bold"] = True
            if cand.get("error"):
                row["error"] = cand["error"]
            rows.append(row)
        summary["formulas"] = rows
        summary["injected"] = sum(1 for cand in candidates if cand["status"] == "injected")
        summary["failed"] = sum(1 for cand in candidates if cand["status"] == "failed")
        if any(cand["bold"] and cand["status"] == "injected" for cand in candidates):
            summary["notes"].append("部分公式原文为粗体：OMML 数学区不保留粗体（可接受差异）")
        (run.qa_dir / "math-summary.json").write_text(
            json.dumps(summary, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
        )

    if not candidates or dry_run:
        write_summary()
        return summary

    engine = _load_engine()
    math_dir = run.qa_dir / "math"
    if math_dir.exists():  # 重跑覆盖当前最佳：清掉上次注入的 receipt/plan
        shutil.rmtree(math_dir)
    math_dir.mkdir(parents=True)

    operations: list[dict] = []
    for cand in candidates:
        try:
            receipt = engine.compile_formula(cand["formula_id"], cand["latex"], "inline")
        except Exception as exc:  # 单公式失败只 warn 跳过该框（保留原样），不整案失败
            cand["status"] = "failed"
            cand["error"] = f"{type(exc).__name__}: {exc}"
            warnings.append(f"{cand['name']} compile 失败（保留原文本框）: {exc}")
            continue
        receipt_path = math_dir / f"{cand['formula_id']}.json"
        receipt_path.write_text(
            json.dumps(receipt, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
        )
        operation = {
            "slide_index": cand["slide_index"],
            "placeholder_name": cand["placeholder"],
            "formula_id": cand["formula_id"],
            "receipt_path": receipt_path.name,
            "receipt_sha256": common.sha256_file(receipt_path),
        }
        if cand["font"] is not None:
            operation["target_font_size_pt"] = cand["font"][0]
            operation["target_font_color"] = cand["font"][1]
        operations.append(operation)
        cand["status"] = "compiled"

    ready = [cand for cand in candidates if cand["status"] == "compiled"]
    if not ready:  # 全部失败/无公式框：正常退出并说明，不改 pptx
        write_summary()
        return summary

    # 先改名落盘（inject_plan 按 cNvPr@name 定位），再注入到临时 pptx，成功后原子替换
    for cand in ready:
        cand["shape"].name = cand["placeholder"]

    plan_path = math_dir / "plan.json"
    plan_path.write_text(
        json.dumps({"schema_version": "1.0", "operations": operations}, ensure_ascii=False, indent=2)
        + "\n",
        encoding="utf-8",
    )
    staged_pptx = run.pptx_path.with_name("redraw.math-source.pptx")
    tmp_pptx = run.pptx_path.with_name("redraw.math-tmp.pptx")
    staged_pptx.unlink(missing_ok=True)
    tmp_pptx.unlink(missing_ok=True)
    prs.save(staged_pptx)
    try:
        engine.inject_plan(staged_pptx, plan_path, tmp_pptx)
    except Exception as exc:
        tmp_pptx.unlink(missing_ok=True)
        raise common.fail(f"OMML 注入失败，redraw.pptx 保持原样: {exc}") from exc
    finally:
        staged_pptx.unlink(missing_ok=True)
    os.replace(tmp_pptx, run.pptx_path)
    for cand in ready:
        cand["status"] = "injected"

    summary.update(_refresh_v3_contracts_after_math(run, candidates))
    write_summary()
    return summary


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(prog="autofigure math", description=__doc__)
    parser.add_argument("run_dir", type=Path, help="run 目录")
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="只检测并重建 LaTeX 写 qa/math-summary.json，不改 PPTX",
    )
    args = parser.parse_args(argv)

    run = common.open_run(args.run_dir)
    summary = upgrade(run, dry_run=args.dry_run)
    _echo(f"检测到 {summary['detected']} 个公式框（强信号 {summary['strong']} / 弱信号 {summary['weak']}）")
    if args.dry_run:
        _echo("dry-run：未改动 PPTX")
    elif summary["injected"]:
        _echo(f"已注入 {summary['injected']} 个原生 Office Math；失败 {summary['failed']} 个（保留原文本框）")
    elif summary["detected"]:
        _echo(f"全部 {summary['failed']} 个公式 compile 失败，PPTX 未改动")
    else:
        _echo("无公式框，PPTX 未改动")
    for warning in summary["warnings"]:
        _echo(f"warning: {warning}")
    _echo(f"明细: {run.qa_dir / 'math-summary.json'}")
    if not args.dry_run and summary["injected"]:
        try:
            from tools.v2 import render_export

            meta = run.load_meta()
            render_export.render(run.pptx_path, run.render_png, int(meta["width"]), int(meta["height"]))
            _echo(f"fresh render: {run.render_png}")
        except Exception as exc:  # COM 渲染失败只 warn，注入结果不受影响
            _echo(f"warning: fresh render 刷新失败: {exc}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
