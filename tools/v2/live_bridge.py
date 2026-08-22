"""Build a fail-closed Scene 2.1 compatibility case for powerpoint-live.

Autofigure's v3 scene remains the source of truth.  The generated case is a
derived adapter packet for the currently installed PowerPoint MCP, which only
accepts the scientific-illustrator project/render 2.0 and scene 2.0/2.1
contracts.  It is never an approval record.
"""

from __future__ import annotations

import re
import shutil
import xml.etree.ElementTree as ET
import zipfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterable

from PIL import Image
from pptx import Presentation

from tools.v2 import common
from tools.v2.contracts import read_json, write_json
from tools.v2.svggeom import parse_path_d

EMU_PER_PX = 914400 / 96
TARGET_ID = "autofigure-pptx"
ADAPTER_SCENE_SCHEMA = "2.1.0"
# powerpoint-live resolves profiles from its bundled, immutable profile
# directory.  A case-local/custom profile id cannot be submitted.  The
# double-column profile is the closest built-in contract for wide scientific
# figures; the exact reference pixel canvas remains authoritative below.
TARGET_PROFILE_ID = "journal-double-column"


def _now() -> str:
    return datetime.now(timezone.utc).isoformat().replace("+00:00", "Z")


def _project_id(case: str) -> str:
    slug = re.sub(r"[^a-z0-9-]+", "-", case.lower()).strip("-")
    slug = re.sub(r"-+", "-", slug)[:52].strip("-") or "case"
    return f"autofigure-{slug}"


def _walk_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        yield shape
        children = getattr(shape, "shapes", None)
        if children is not None:
            yield from _walk_shapes(children)


def _shape_bounds(shape: Any) -> dict[str, float]:
    if isinstance(shape, dict):
        return shape
    return {
        "x": float(shape.left) / EMU_PER_PX,
        "y": float(shape.top) / EMU_PER_PX,
        "width": max(float(shape.width) / EMU_PER_PX, 0.5),
        "height": max(float(shape.height) / EMU_PER_PX, 0.5),
    }


def _union_bounds(shapes: list[Any]) -> dict[str, float]:
    bounds = [_shape_bounds(shape) for shape in shapes]
    left = min(item["x"] for item in bounds)
    top = min(item["y"] for item in bounds)
    right = max(item["x"] + item["width"] for item in bounds)
    bottom = max(item["y"] + item["height"] for item in bounds)
    return {
        "x": round(left, 4),
        "y": round(top, 4),
        "width": round(max(right - left, 0.5), 4),
        "height": round(max(bottom - top, 0.5), 4),
    }


def _text_spec(shape: Any, fallback_text: str = "") -> dict[str, Any]:
    text = getattr(shape, "text", "") or fallback_text
    font_size = 12.0
    font_family = "Arial"
    frame = getattr(shape, "text_frame", None)
    if frame is not None:
        for paragraph in frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size is not None:
                    font_size = float(run.font.size.pt)
                if run.font.name:
                    font_family = run.font.name
                if run.text:
                    break
            if paragraph.runs:
                break
    return {
        "text": text,
        "fontToken": "autofigure-font",
        "fontSize": max(font_size, 1.0),
        "fontWeight": "normal",
        "horizontalAlign": "center",
        "verticalAlign": "middle",
        "margins": {"left": 0, "top": 0, "right": 0, "bottom": 0},
        "autofit": "none",
        "runs": [{"text": text, "fontToken": font_family, "fontSize": max(font_size, 1.0)}],
    }


def _xml_shape_bounds(path: Path) -> tuple[dict[str, dict[str, float]], dict[int, dict[str, float]]]:
    """Read geometry for shapes hidden in mc:AlternateContent (notably OMML)."""
    by_name: dict[str, dict[str, float]] = {}
    by_id: dict[int, dict[str, float]] = {}
    drawing_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    shape_tags = {"sp", "pic", "cxnSp", "graphicFrame", "grpSp"}
    with zipfile.ZipFile(path) as package:
        for member in package.namelist():
            if not member.startswith("ppt/slides/slide") or not member.endswith(".xml"):
                continue
            root = ET.fromstring(package.read(member))
            parents = {child: parent for parent in root.iter() for child in parent}
            for identity in root.iter():
                if identity.tag.rsplit("}", 1)[-1] != "cNvPr":
                    continue
                owner = identity
                while owner in parents and owner.tag.rsplit("}", 1)[-1] not in shape_tags:
                    owner = parents[owner]
                if owner.tag.rsplit("}", 1)[-1] not in shape_tags:
                    continue
                transform = owner.find(f".//{{{drawing_ns}}}xfrm")
                if transform is None:
                    continue
                offset = transform.find(f"{{{drawing_ns}}}off")
                extent = transform.find(f"{{{drawing_ns}}}ext")
                if offset is None or extent is None:
                    continue
                try:
                    bounds = {
                        "x": float(offset.get("x", "0")) / EMU_PER_PX,
                        "y": float(offset.get("y", "0")) / EMU_PER_PX,
                        "width": max(float(extent.get("cx", "0")) / EMU_PER_PX, 0.5),
                        "height": max(float(extent.get("cy", "0")) / EMU_PER_PX, 0.5),
                    }
                    shape_id = int(identity.get("id", "0"))
                except ValueError:
                    continue
                name = identity.get("name")
                if name:
                    by_name[name] = bounds
                if shape_id:
                    by_id[shape_id] = bounds
    return by_name, by_id


def _svg_elements(path: Path) -> dict[str, ET.Element]:
    root = ET.fromstring(path.read_text(encoding="utf-8"))
    return {element_id: element for element in root.iter() if (element_id := element.get("id"))}


def _cubic_path(segments: list[tuple]) -> dict[str, Any] | None:
    start: tuple[float, float] | None = None
    current: tuple[float, float] | None = None
    cubics: list[dict[str, Any]] = []
    for segment in segments:
        if segment[0] == "M":
            start = current = (float(segment[1]), float(segment[2]))
        elif segment[0] == "L" and current is not None:
            end = (float(segment[1]), float(segment[2]))
            cubics.append(
                {
                    "control1": {
                        "x": current[0] + (end[0] - current[0]) / 3,
                        "y": current[1] + (end[1] - current[1]) / 3,
                    },
                    "control2": {
                        "x": current[0] + 2 * (end[0] - current[0]) / 3,
                        "y": current[1] + 2 * (end[1] - current[1]) / 3,
                    },
                    "end": {"x": end[0], "y": end[1]},
                }
            )
            current = end
        elif segment[0] == "C" and current is not None:
            end = (float(segment[5]), float(segment[6]))
            cubics.append(
                {
                    "control1": {"x": float(segment[1]), "y": float(segment[2])},
                    "control2": {"x": float(segment[3]), "y": float(segment[4])},
                    "end": {"x": end[0], "y": end[1]},
                }
            )
            current = end
    if start is None or not cubics:
        return None
    return {
        "kind": "cubic",
        "coordinateSpace": "canvas",
        "start": {"x": start[0], "y": start[1]},
        "segments": cubics,
    }


def _edge_route(element: dict[str, Any]) -> tuple[str, dict[str, Any] | None]:
    geometry = element.get("geometry", {})
    d = geometry.get("d")
    if not d:
        return "straight", None
    segments = parse_path_d(d)
    if any(segment[0] == "C" for segment in segments):
        return "curved", _cubic_path(segments)
    points = [
        {"x": float(segment[1]), "y": float(segment[2])}
        for segment in segments
        if segment[0] in {"M", "L"}
    ]
    if len(points) > 2:
        return "polyline", {"kind": "polyline", "coordinateSpace": "canvas", "points": points}
    return "straight", None


def _crop_assets(
    run: common.Run,
    live_root: Path,
    asset_contract: dict[str, Any],
) -> tuple[dict[str, dict[str, Any]], list[dict[str, Any]]]:
    result: dict[str, dict[str, Any]] = {}
    manifest_assets: list[dict[str, Any]] = []
    destination = live_root / "assets"
    destination.mkdir(parents=True, exist_ok=True)
    with Image.open(run.source_png) as reference:
        for item in asset_contract.get("assets", []):
            if item.get("source") != "reference_crop":
                continue
            if item.get("authorized") is not True:
                raise common.fail(f"atomic asset is not authorized: {item.get('id')}")
            x, y, width, height = (int(value) for value in item["bbox"])
            if width <= 0 or height <= 0:
                raise common.fail(f"invalid atomic asset bbox: {item.get('id')}")
            filename = f"{common.slugify(item['id'], 64)}.png"
            output = destination / filename
            reference.crop((x, y, x + width, y + height)).save(output)
            asset_id = item["id"]
            slot_id = f"slot-{common.slugify(asset_id, 48)}"
            result[asset_id] = {
                "asset_id": asset_id,
                "slot_id": slot_id,
                "path": output,
                "bbox": [x, y, width, height],
            }
            manifest_assets.append(
                {
                    "assetId": asset_id,
                    "responsibilityClass": "verified_source",
                    "selectedFile": f"assets/{output.name}",
                    "sourceFile": "input/reference.png",
                    "sha256": common.sha256_file(output),
                    "mime": "image/png",
                    "width": width,
                    "height": height,
                    "targetSlot": slot_id,
                    "provenance": {
                        "route": "verified_import",
                        "status": "verified_source",
                        "license": None,
                        "rightsBasis": item.get("authorization_basis"),
                    },
                    "reviewStatus": "PENDING",
                }
            )
    return result, manifest_assets


def build_powerpoint_live_bridge(run: common.Run) -> dict[str, Any]:
    """Derive a complete managed-session case without changing v3 truth files."""
    if not run.pptx_path.is_file() or not run.scene_path.is_file() or not run.bindings_path.is_file():
        raise common.fail("PowerPoint-live bridge requires current PPTX, scene.json and bindings.json")
    meta = run.load_meta()
    live_root = run.live_case_dir
    for directory in (live_root / "input", live_root / "design", live_root / "assets"):
        directory.mkdir(parents=True, exist_ok=True)
    reference_copy = live_root / "input" / "reference.png"
    template_copy = live_root / "input" / "candidate.pptx"
    shutil.copy2(run.source_png, reference_copy)
    shutil.copy2(run.pptx_path, template_copy)

    project_id = _project_id(meta["case"])
    scene_v3 = read_json(run.scene_path)
    bindings = read_json(run.bindings_path)
    assets_v3 = read_json(run.assets_path)
    if bindings.get("artifact_sha256") != common.sha256_file(run.pptx_path):
        raise common.fail("PowerPoint-live bridge refuses stale PPTX bindings")

    presentation = Presentation(run.pptx_path)
    shapes = list(_walk_shapes(presentation.slides[0].shapes))
    by_name = {shape.name: shape for shape in shapes}
    by_id = {shape.shape_id: shape for shape in shapes}
    xml_by_name, xml_by_id = _xml_shape_bounds(run.pptx_path)
    binding_rows: dict[str, list[Any]] = {}
    binding_meta: dict[str, dict[str, Any]] = {}
    for binding in bindings.get("bindings", []):
        shape = (
            by_name.get(binding.get("shape_name"))
            or by_id.get(binding.get("shape_id"))
            or xml_by_name.get(binding.get("shape_name"))
            or xml_by_id.get(binding.get("shape_id"))
        )
        if shape is None:
            raise common.fail(f"PowerPoint-live bridge cannot read bound shape: {binding}")
        binding_rows.setdefault(binding["element_id"], []).append(shape)
        binding_meta.setdefault(binding["element_id"], binding)

    atomic_assets, manifest_assets = _crop_assets(run, live_root, assets_v3)
    svg_by_id = _svg_elements(run.redraw_svg)
    nodes: list[dict[str, Any]] = []
    edges: list[dict[str, Any]] = []
    relations: list[dict[str, Any]] = []
    entities: list[dict[str, Any]] = []
    asset_slots: list[dict[str, Any]] = []

    for element in scene_v3.get("elements", []):
        element_id = element["id"]
        bound_shapes = binding_rows.get(element_id)
        if not bound_shapes:
            raise common.fail(f"PowerPoint-live bridge has no shape binding for {element_id}")
        binding = binding_meta[element_id]
        topology = element.get("topology", {})
        is_edge = bool(topology.get("source") and topology.get("target"))
        if is_edge:
            route, path_spec = _edge_route(element)
            relation_id = f"rel:{element_id}"
            svg_element = svg_by_id.get(element_id)
            marker_start = svg_element is not None and (
                svg_element.get("marker-start") is not None
                or "marker-start" in (svg_element.get("style") or "")
            )
            marker_end = svg_element is not None and (
                svg_element.get("marker-end") is not None
                or "marker-end" in (svg_element.get("style") or "")
            )
            direction = "bidirectional" if marker_start and marker_end else "forward"
            edge = {
                "id": element_id,
                "semanticRelationIds": [relation_id],
                "geometryOnly": False,
                "source": topology["source"],
                "target": topology["target"],
                "direction": direction,
                "route": route,
                "styleTokens": {"strokePattern": "solid"},
            }
            if path_spec is not None:
                edge["pathSpec"] = path_spec
            edges.append(edge)
            relations.append(
                {
                    "relationId": relation_id,
                    "source": topology["source"],
                    "target": topology["target"],
                    "kind": "flow",
                    "direction": direction,
                    "sourceIds": ["reference"],
                }
            )
            continue

        object_kind = binding.get("object_kind")
        kind = (
            "text"
            if object_kind in {"text", "native-math"}
            else "line"
            if object_kind in {"line", "connector"}
            else "shape"
        )
        if element_id in atomic_assets:
            kind = "asset"
        bounds = _union_bounds(bound_shapes)
        node: dict[str, Any] = {
            "id": element_id,
            "semanticId": element_id,
            "kind": kind,
            **bounds,
            "editable": bool(element.get("editable", binding.get("editable", True))),
            "zIndex": int(element.get("z_index", len(nodes))),
        }
        if kind == "text":
            svg_element = svg_by_id.get(element_id)
            fallback_text = "" if svg_element is None else "".join(svg_element.itertext()).strip()
            node["textSpec"] = _text_spec(bound_shapes[0], fallback_text)
        elif kind == "asset":
            asset = atomic_assets[element_id]
            node["assetSpec"] = {
                "assetId": element_id,
                "fitMode": "stretch",
                "crop": {"left": 0, "top": 0, "right": 0, "bottom": 0},
                "atomicRasterUnit": True,
                "containsReconstructableContent": False,
            }
            asset_slots.append(
                {
                    "slotId": asset["slot_id"],
                    "semanticAnchor": element_id,
                    **bounds,
                    "overlaySafeArea": "none; formal text and topology remain native",
                    "expectedAspectRatio": bounds["width"] / bounds["height"],
                }
            )
        nodes.append(node)
        label = getattr(bound_shapes[0], "text", "") or node.get("textSpec", {}).get("text") or element_id
        entities.append(
            {
                "entityId": element_id,
                "label": label,
                "kind": "annotation" if kind == "text" else "panel" if "panel" in element_id else "component",
                "sourceIds": ["reference"],
                "required": True,
                "editable": node["editable"],
            }
        )

    node_ids = {item["id"] for item in nodes}
    missing_endpoints = sorted(
        {
            endpoint
            for edge in edges
            for endpoint in (edge["source"], edge["target"])
            if endpoint not in node_ids
        }
    )
    if missing_endpoints:
        raise common.fail(f"PowerPoint-live bridge edge endpoints are not nodes: {missing_endpoints}")

    scene_version = f"autofigure-{common.sha256_file(run.scene_path)[:12]}"
    scene = {
        "schemaVersion": ADAPTER_SCENE_SCHEMA,
        "projectId": project_id,
        "version": scene_version,
        "canvas": {
            "width": meta["width"],
            "height": meta["height"],
            "unit": "px",
            "profileId": TARGET_PROFILE_ID,
            "backgroundToken": "white",
        },
        "nodes": nodes,
        "edges": edges,
        "assetSlots": asset_slots,
    }
    scene_path = live_root / "design" / "scene_graph.json"
    write_json(scene_path, scene)

    if not manifest_assets:
        manifest_assets.append(
            {
                "assetId": "autofigure-native-scene",
                "responsibilityClass": "exact_vector",
                "targetSlot": "native-scene",
                "sceneGraphVersion": scene_version,
                "sceneGraphSha256": common.sha256_file(scene_path),
                "sceneElementIds": sorted([item["id"] for item in nodes + edges]),
                "provenance": {"route": "backend_native", "status": "deterministic", "license": None},
                "reviewStatus": "PENDING",
            }
        )

    element_ids = [item["id"] for item in nodes + edges]
    render_plan = {
        "schemaVersion": "2.0.0",
        "projectId": project_id,
        "sceneGraphVersion": scene_version,
        "targets": [
            {
                "targetId": TARGET_ID,
                "backendId": "powerpoint",
                "adapterId": "powerpoint-live",
                "required": True,
                "deliverables": ["pptx", "png"],
                "requiredCapabilities": [
                    "managed-session",
                    "native-connector",
                    "freeform",
                    "inspect",
                    "audit",
                    "save-reopen",
                ],
            }
        ],
        "elements": [
            {
                "elementId": element_id,
                "mustRemainEditable": next(
                    item.get("editable", True) for item in nodes + edges if item["id"] == element_id
                ),
                "qualityOwner": "asset" if element_id in atomic_assets else "backend",
                "reason": "Derived Autofigure v3 compatibility route for regional live repair.",
                "routes": [
                    {
                        "targetId": TARGET_ID,
                        "renderer": "verified_raster" if element_id in atomic_assets else "backend_native",
                    }
                ],
            }
            for element_id in element_ids
        ],
        "deliverables": ["pptx", "png"],
    }
    scientific_spec = {
        "schemaVersion": "1.0.0",
        "projectId": project_id,
        "figureRole": "system_architecture",
        "language": "en",
        "audience": "Autofigure human review",
        "claims": [
            {
                "claimId": "reference-correspondence",
                "text": "This editable candidate reconstructs the designated reference; scientific interpretation remains source-bound.",
                "evidenceType": "conceptual",
                "sourceIds": ["reference"],
                "mustShow": True,
            }
        ],
        "entities": entities,
        "relations": relations,
        "formulas": [],
        "evidenceDeclarations": [],
        "unknowns": ["Independent scientific-semantic approval has not been supplied."],
        "approval": {
            "status": "PENDING",
            "approvedBy": None,
            "approvedAt": None,
            "notes": "Bridge generation does not approve scientific content.",
        },
    }
    source_manifest = {
        "schemaVersion": "1.0.0",
        "projectId": project_id,
        "taskMode": "RECONSTRUCT_1TO1",
        "designatedReferenceSourceId": "reference",
        "designatedReferenceSha256": meta["source_sha256"],
        "referenceUsage": "measurement_only",
        "referenceEmbedded": False,
        "independentReferenceCase": True,
        "sources": [
            {
                "sourceId": "reference",
                "kind": "reference_image",
                "pathOrUrl": "input/reference.png",
                "sha256": common.sha256_file(reference_copy),
                "authority": "user_supplied",
                "licenseStatus": "unknown",
                "confidentiality": "internal",
                "notes": "Designated measurement-only reconstruction reference.",
            },
            {
                "sourceId": "candidate-pptx",
                "kind": "prior_source",
                "pathOrUrl": "input/candidate.pptx",
                "sha256": common.sha256_file(template_copy),
                "authority": "user_supplied",
                "licenseStatus": "unknown",
                "confidentiality": "internal",
                "notes": "Offline Autofigure candidate copied for managed live repair; never edited in place.",
            },
        ],
    }
    project_state = {
        "schemaVersion": "2.0.0",
        "contractVersion": "2.0.0",
        "projectId": project_id,
        "taskMode": "RECONSTRUCT_1TO1",
        "state": "PLANNED",
        "lastSuccessfulState": "PLANNED",
        "profileId": TARGET_PROFILE_ID,
        "requestedBackendIds": ["powerpoint"],
        "updatedAt": _now(),
        "blockedReason": None,
        "gates": {
            "environment": "PASS",
            "scientificApproval": "PENDING",
            "assetDirection": "PENDING",
            "styleApproval": "PENDING",
            "assetIntegrity": "PENDING",
            "backendAudit": "PENDING",
            "artifactSetIntegrity": "PENDING",
            "backendReadback": "PENDING",
            "crossBackendEquivalence": "PENDING",
            "quality": "PENDING",
            "humanApproval": "PENDING",
            "exportReadback": "PENDING",
        },
        "history": [
            {
                "state": "PLANNED",
                "at": _now(),
                "evidence": ["derived from Autofigure v3 scene/bindings and current PPTX hash"],
                "note": "Compatibility case only; releaseAuthority=NONE.",
            }
        ],
    }
    asset_manifest = {"schemaVersion": "1.1.0", "projectId": project_id, "assets": manifest_assets}
    evidence_provenance = {
        "schemaVersion": "1.0.0",
        "projectId": project_id,
        "taskMode": "RECONSTRUCT_1TO1",
        "designatedReference": {
            "sourceId": "reference",
            "sha256": meta["source_sha256"],
            "usage": "measurement_only",
            "embedded": False,
        },
        "reviewIsolation": {
            "doesNotOverrideAssetManifest": True,
            "doesNotEstablishScientificValidity": True,
            "doesNotEstablishRights": True,
        },
        "reviewStatus": "PENDING",
        "lifecycleStatus": "DRAFT",
        "records": [],
    }

    documents = {
        "project_state.json": project_state,
        "input/source_manifest.json": source_manifest,
        "design/scene_graph.json": scene,
        "design/render_plan.json": render_plan,
        "design/scientific_spec.json": scientific_spec,
        "assets/asset_manifest.json": asset_manifest,
        "assets/evidence_provenance.json": evidence_provenance,
    }
    for relative, document in documents.items():
        write_json(live_root / relative, document)

    manifest = {
        "schema_version": "1.0.0",
        "kind": "powerpoint_live_scene_bridge",
        "ready": True,
        "release_authority": "NONE",
        "path_base": "autofigure-case-root",
        "autofigure_case_root": ".",
        "case_root": "qa/powerpoint-live-case",
        "project_id": project_id,
        "target_id": TARGET_ID,
        "task_mode": "RECONSTRUCT_1TO1",
        "template_path": "qa/powerpoint-live-case/input/candidate.pptx",
        "reference_sha256": meta["source_sha256"],
        "source_scene_schema_version": scene_v3.get("schema_version"),
        "source_scene_sha256": common.sha256_file(run.scene_path),
        "adapter_scene_schema_version": ADAPTER_SCENE_SCHEMA,
        "adapter_scene_sha256": common.sha256_file(live_root / "design" / "scene_graph.json"),
        "element_count": len(element_ids),
        "contract_files": {
            relative: common.sha256_file(live_root / relative) for relative in documents
        },
    }
    write_json(run.live_bridge_path, manifest)
    return manifest
