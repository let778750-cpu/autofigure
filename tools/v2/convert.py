"""autofigure convert — 把 VLM 重绘的 SVG 转换为原生可编辑 PPTX。

映射合同（references/v2-prompt-contract.md）：
- rect/circle/ellipse/line/polyline/polygon/path → 原生形状 / custGeom 自由曲线（保留三次贝塞尔）
- text/tspan → 原生文本框 runs（italic、字号、颜色、baseline-shift → 上下标）
- linearGradient → a:gradFill；stroke-dasharray → prstDash；marker → 自由曲线箭头
- <rect id="atomic:*"> 占位符 → 从参考图裁剪对应 bbox 嵌入为位图
- <image> 容错：按 bbox 从参考图裁剪替代并记 warning；覆盖画布 ≥50% 直接拒绝（防整图截图冒充矢量）
- <g> 无变换时转原生分组，其余情况拍平（当前 VLM 输出实测无 <g>）
"""

from __future__ import annotations

import argparse
import hashlib
import io
import json
import math
import re
import sys
import xml.etree.ElementTree as ET
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Pt

from tools.v2 import common
from tools.v2.layout import Box, collect_svg_boxes
from tools.v2.svggeom import Matrix, parse_path_d, parse_transform

SVG_NS = "{http://www.w3.org/2000/svg}"
PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
EMU_PER_PX = 9525
PT_PER_PX = 0.75  # 96 dpi
BASELINE_ASCENT = 0.95  # 实测标定：文本框顶到首行基线（Arial em 比例，含半行距）
TEXT_BOX_PAD_X = 32.0  # 保持锚点不动的透明选择框余量；覆盖 live 实测最大 28.44 px 横向不足
TEXT_BOX_PAD_Y = 20.0  # 覆盖旋转标签 live 实测最大 19.50 px 纵向不足

NAMED_COLORS = {
    "black": "#000000", "white": "#FFFFFF", "red": "#FF0000", "green": "#008000",
    "blue": "#0000FF", "gray": "#808080", "grey": "#808080", "orange": "#FFA500",
    "yellow": "#FFFF00", "purple": "#800080", "brown": "#A52A2A", "pink": "#FFC0CB",
    "none": None, "transparent": None,
}

INHERITED_STYLE_KEYS = (
    "fill", "fill-opacity", "stroke", "stroke-width", "stroke-opacity",
    "stroke-dasharray", "stroke-linecap", "stroke-linejoin", "opacity",
    "font-size", "font-family", "font-style", "font-weight", "text-anchor",
)


# ---------------------------------------------------------------- 样式与颜色


def _parse_style_attr(style: str | None) -> dict[str, str]:
    result: dict[str, str] = {}
    if style:
        for item in style.split(";"):
            if ":" in item:
                key, value = item.split(":", 1)
                result[key.strip()] = value.strip()
    return result


def _element_style(element: ET.Element, parent_style: dict[str, str]) -> dict[str, str]:
    style = dict(parent_style)
    inline = _parse_style_attr(element.get("style"))
    for key in INHERITED_STYLE_KEYS:
        if element.get(key) is not None:
            style[key] = element.get(key)  # presentation attribute
        if key in inline:
            style[key] = inline[key]
    for key in ("fill-rule", "stroke-linecap", "stroke-linejoin"):
        if key in inline:
            style[key] = inline[key]
    return style


def parse_color(value: str | None) -> tuple[str, float] | None:
    """→ (hex, alpha) 或 None（none/未指定）。"""
    if value is None:
        return None
    value = value.strip()
    if value in NAMED_COLORS:
        named = NAMED_COLORS[value]
        return (named, 1.0) if named else None
    if value.startswith("#"):
        hexpart = value[1:]
        if len(hexpart) == 3:
            hexpart = "".join(ch * 2 for ch in hexpart)
        if len(hexpart) == 6:
            return ("#" + hexpart.upper(), 1.0)
        if len(hexpart) == 8:
            return ("#" + hexpart[:6].upper(), int(hexpart[6:], 16) / 255.0)
    match = re.match(r"rgba?\(([^)]*)\)", value)
    if match:
        parts = [p.strip() for p in match.group(1).split(",")]
        if len(parts) >= 3:
            r, g, b = (int(float(p)) for p in parts[:3])
            alpha = float(parts[3]) if len(parts) > 3 else 1.0
            return (f"#{r:02X}{g:02X}{b:02X}", alpha)
    return None


def _opacity(style: dict[str, str], key: str) -> float:
    try:
        return float(style.get(key, "1"))
    except ValueError:
        return 1.0


# ---------------------------------------------------------------- OOXML 辅助


def _px(value: float) -> Emu:
    return Emu(round(value * EMU_PER_PX))


def _insert_before_any(sp_pr, element, successor_tags: tuple[str, ...]) -> None:
    """按 OOXML spPr 子元素顺序插入：fill < ln < effectLst < scene3d < sp3d < extLst。"""
    for tag in successor_tags:
        sibling = sp_pr.find(qn(tag))
        if sibling is not None:
            sibling.addprevious(element)
            return
    sp_pr.append(element)


_FILL_SUCCESSORS = ("a:ln", "a:effectLst", "a:effectDag", "a:scene3d", "a:sp3d", "a:extLst")
_LINE_SUCCESSORS = ("a:effectLst", "a:effectDag", "a:scene3d", "a:sp3d", "a:extLst")


def _set_solid_fill(sp_pr, hex_color: str, alpha: float) -> None:
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
        for child in sp_pr.findall(qn(tag)):
            sp_pr.remove(child)
    fill = sp_pr.makeelement(qn("a:solidFill"), {})
    color = sp_pr.makeelement(qn("a:srgbClr"), {"val": hex_color[1:]})
    if alpha < 1.0:
        color.append(sp_pr.makeelement(qn("a:alpha"), {"val": str(round(alpha * 100000))}))
    fill.append(color)
    _insert_before_any(sp_pr, fill, _FILL_SUCCESSORS)


def _set_no_fill(sp_pr) -> None:
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
        for child in sp_pr.findall(qn(tag)):
            sp_pr.remove(child)
    _insert_before_any(sp_pr, sp_pr.makeelement(qn("a:noFill"), {}), _FILL_SUCCESSORS)


def _set_gradient_fill(sp_pr, stops: list[tuple[float, str, float]], angle_deg: float) -> None:
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
        for child in sp_pr.findall(qn(tag)):
            sp_pr.remove(child)
    grad = sp_pr.makeelement(qn("a:gradFill"), {})
    gs_lst = grad.makeelement(qn("a:gsLst"), {})
    for pos, hex_color, alpha in stops:
        gs = gs_lst.makeelement(qn("a:gs"), {"pos": str(round(pos * 100000))})
        color = gs_lst.makeelement(qn("a:srgbClr"), {"val": hex_color[1:]})
        if alpha < 1.0:
            color.append(gs_lst.makeelement(qn("a:alpha"), {"val": str(round(alpha * 100000))}))
        gs.append(color)
        gs_lst.append(gs)
    grad.append(gs_lst)
    angle = round(math.degrees(angle_deg) * 60000) % 21600000
    grad.append(grad.makeelement(qn("a:lin"), {"ang": str(angle), "scaled": "1"}))
    _insert_before_any(sp_pr, grad, _FILL_SUCCESSORS)


def _apply_line(sp_pr, style: dict[str, str], default_width: float = 1.0) -> None:
    """按 style 在 spPr 上设置 a:ln（颜色/宽度/虚线/圆角）。无 stroke 则不建 a:ln。"""
    stroke = parse_color(style.get("stroke"))
    if stroke is None:
        old = sp_pr.find(qn("a:ln"))
        if old is not None:
            sp_pr.remove(old)
        line = sp_pr.makeelement(qn("a:ln"), {})
        line.append(line.makeelement(qn("a:noFill"), {}))
        _insert_before_any(sp_pr, line, _LINE_SUCCESSORS)
        return
    hex_color, stroke_alpha = stroke
    alpha = stroke_alpha * _opacity(style, "stroke-opacity") * _opacity(style, "opacity")
    try:
        width = float(style.get("stroke-width", default_width))
    except ValueError:
        width = default_width
    old = sp_pr.find(qn("a:ln"))
    if old is not None:
        sp_pr.remove(old)
    line = sp_pr.makeelement(qn("a:ln"), {"w": str(round(width * EMU_PER_PX))})
    if style.get("stroke-linecap") in ("round", "square"):
        line.set("cap", "rnd" if style["stroke-linecap"] == "round" else "sq")
    fill = line.makeelement(qn("a:solidFill"), {})
    color = line.makeelement(qn("a:srgbClr"), {"val": hex_color[1:]})
    if alpha < 1.0:
        color.append(line.makeelement(qn("a:alpha"), {"val": str(round(alpha * 100000))}))
    fill.append(color)
    line.append(fill)
    dash = _dash_preset(style.get("stroke-dasharray"), width, style.get("stroke-linecap"))
    if dash:
        line.append(line.makeelement(qn("a:prstDash"), {"val": dash}))
    _insert_before_any(sp_pr, line, _LINE_SUCCESSORS)


def _dash_preset(dasharray: str | None, width: float, linecap: str | None) -> str | None:
    if not dasharray or dasharray.strip() in ("none", ""):
        return None
    try:
        parts = [float(p) for p in re.split(r"[\s,]+", dasharray.strip()) if p]
    except ValueError:
        return None
    if not parts or parts[0] <= 0:
        return None
    first = parts[0]
    if first <= width * 1.5:
        # OOXML 无 roundDot/squareDot：圆点用 dot，方点用 sysDot
        return "dot" if linecap == "round" else "sysDot"
    if first <= 4 * width:
        return "sysDash"
    if first <= 8 * width:
        return "dash"
    return "lgDash"


def _disable_shadow(shape) -> None:
    shape.shadow.inherit = False


def _next_shape_id(slide) -> int:
    return slide.shapes._next_shape_id


# ---------------------------------------------------------------- 渐变与 marker 定义


def _collect_defs(tree: ET.Element) -> dict[str, ET.Element]:
    defs: dict[str, ET.Element] = {}
    for element in tree.iter():
        element_id = element.get("id")
        if element_id:
            defs[element_id] = element
    return defs


def _gradient_stops(grad: ET.Element) -> list[tuple[float, str, float]]:
    stops: list[tuple[float, str, float]] = []
    for stop in grad.findall(f"{SVG_NS}stop"):
        style = _parse_style_attr(stop.get("style"))
        offset = stop.get("offset", "0")
        pos = float(offset[:-1]) / 100.0 if offset.endswith("%") else float(offset)
        color_value = stop.get("stop-color") or style.get("stop-color", "#000000")
        parsed = parse_color(color_value) or ("#000000", 1.0)
        opacity = stop.get("stop-opacity") or style.get("stop-opacity")
        alpha = parsed[1] * (float(opacity) if opacity is not None else 1.0)
        stops.append((max(0.0, min(1.0, pos)), parsed[0], alpha))
    return stops


def _gradient_angle(grad: ET.Element) -> float:
    def num(key: str, default: float) -> float:
        raw = grad.get(key)
        if raw is None:
            return default
        return float(raw[:-1]) / 100.0 if raw.endswith("%") else float(raw)

    x1, y1, x2, y2 = num("x1", 0.0), num("y1", 0.0), num("x2", 1.0), num("y2", 0.0)
    return math.atan2(y2 - y1, x2 - x1)


def _resolve_fill_url(style: dict[str, str], defs: dict[str, ET.Element]) -> ET.Element | None:
    fill = style.get("fill", "")
    match = re.match(r"url\(#([^)]+)\)", fill or "")
    if match:
        return defs.get(match.group(1))
    return None


# ---------------------------------------------------------------- 形状发射


class ConvertContext:
    def __init__(
        self,
        slide,
        defs: dict[str, ET.Element],
        source_png: Path | None,
        canvas_width: int,
        canvas_height: int,
        asset_authorizations: dict[str, dict] | None = None,
        layout_boxes: dict[str, Box] | None = None,
    ):
        self.slide = slide
        self.defs = defs
        self.source_png = source_png
        self.canvas_width = canvas_width
        self.canvas_height = canvas_height
        self.canvas_area = float(canvas_width * canvas_height)
        self.asset_authorizations = asset_authorizations or {}
        self.layout_boxes = layout_boxes or {}
        self.warnings: list[str] = []
        self.counts: dict[str, int] = {}
        self.bindings: list[dict] = []
        self.scene_elements: dict[str, dict] = {}
        self.assets: list[dict] = []
        self.current_element_id = ""
        self.current_svg_tag = ""
        self._element_counts: dict[str, int] = {}
        self._binding_counts: dict[str, int] = {}
        self.pending_connections: list[tuple[object, str | None, str | None, int, int]] = []

    def bump(self, kind: str) -> None:
        self.counts[kind] = self.counts.get(kind, 0) + 1

    def warn(self, message: str) -> None:
        self.warnings.append(message)

    def begin_element(self, element: ET.Element, tag: str) -> tuple[str, str]:
        previous = (self.current_element_id, self.current_svg_tag)
        self._element_counts[tag] = self._element_counts.get(tag, 0) + 1
        self.current_element_id = element.get("id") or f"svg-{tag}-{self._element_counts[tag]:04d}"
        self.current_svg_tag = tag
        role = element.get("data-role") or (
            "edge" if tag in ("line", "path", "polyline") and _has_marker(element) else tag
        )
        topology = {
            key: value
            for key, value in (
                ("source", element.get("data-source-id")),
                ("target", element.get("data-target-id")),
                ("source_gap", element.get("data-source-gap")),
                ("target_gap", element.get("data-target-gap")),
                ("attached", element.get("data-attach")),
            )
            if value
        }
        geometry = {
            key: element.get(key)
            for key in (
                "x",
                "y",
                "width",
                "height",
                "cx",
                "cy",
                "r",
                "rx",
                "ry",
                "x1",
                "y1",
                "x2",
                "y2",
                "points",
                "d",
                "transform",
            )
            if element.get(key) is not None
        }
        scene_element = {
                "id": self.current_element_id,
                "kind": role,
                "role": role,
                "svg_tag": tag,
                "geometry": geometry,
                "topology": topology,
                "z_index": len(self.scene_elements),
                "editable": True,
            }
        layout = {
            key: value
            for key, value in (
                ("container_id", element.get("data-layout-container")),
                ("padding", element.get("data-layout-padding")),
                ("tolerance", element.get("data-layout-tolerance")),
                ("repeat_group", element.get("data-repeat-group")),
                ("repeat_axis", element.get("data-repeat-axis")),
                ("repeat_order", element.get("data-repeat-order")),
                ("repeat_size_tolerance", element.get("data-repeat-size-tolerance")),
                ("repeat_axis_tolerance", element.get("data-repeat-axis-tolerance")),
                ("repeat_spacing_tolerance", element.get("data-repeat-spacing-tolerance")),
            )
            if value is not None
        }
        if layout:
            scene_element["layout"] = layout
        self.scene_elements.setdefault(self.current_element_id, scene_element)
        return previous

    def end_element(self, previous: tuple[str, str]) -> None:
        self.current_element_id, self.current_svg_tag = previous

    def _shape_name(self, kind: str) -> str:
        element_id = re.sub(r"[^A-Za-z0-9_.-]+", "-", self.current_element_id).strip("-")
        element_id = element_id[:80] or "element"
        key = f"{element_id}:{kind}"
        self._binding_counts[key] = self._binding_counts.get(key, 0) + 1
        return f"af-{element_id}-{kind}-{self._binding_counts[key]:02d}"

    def register_shape(self, shape, kind: str, *, editable: bool = True) -> str:
        name = self._shape_name(kind)
        c_nv_pr = shape._element.find(f".//{qn('p:cNvPr')}")
        if c_nv_pr is not None:
            c_nv_pr.set("name", name)
        self.register_raw(shape.shape_id, name, kind, editable=editable)
        return name

    def register_xml(self, shape, shape_id: int, kind: str, *, editable: bool = True) -> str:
        name = self._shape_name(kind)
        c_nv_pr = shape.find(f".//{qn('p:cNvPr')}")
        if c_nv_pr is not None:
            c_nv_pr.set("name", name)
        self.register_raw(shape_id, name, kind, editable=editable)
        return name

    def register_raw(self, shape_id: int, name: str, kind: str, *, editable: bool = True) -> None:
        self.bindings.append(
            {
                "element_id": self.current_element_id,
                "shape_id": int(shape_id),
                "shape_name": name,
                "object_kind": kind,
                "editable": editable,
                "semantic_group_id": self.current_element_id,
            }
        )
        self.scene_elements.setdefault(
            self.current_element_id,
            {
                "id": self.current_element_id,
                "kind": "edge" if kind in ("connector", "line", "freeform-arrow") else kind,
                "svg_tag": self.current_svg_tag,
                "editable": editable,
            },
        )
        if not editable:
            self.scene_elements[self.current_element_id]["editable"] = False

    def register_asset(self, asset_id: str, bbox: list[int], source_sha256: str) -> None:
        authorization = self.asset_authorizations.get(asset_id, {})
        authorized = authorization.get("authorized") is True
        raster_reason = authorization.get("raster_reason") or (
            "Reference-faithful irreducible creative microasset authorized by the user."
        )
        decomposition_note = authorization.get("decomposition_note") or (
            "No faithful native decomposition; formal text, formulas, nodes, and topology remain native."
        )
        record = {
            **authorization,
            "id": asset_id,
            "source": "reference_crop",
            "source_sha256": source_sha256,
            "bbox": bbox,
            "editable": False,
            "authorized": authorized,
            "source_tightly_cropped": True,
            "atomic_raster_unit": True,
            "contains_reconstructable_content": False,
            "raster_reason": raster_reason,
            "decomposition_note": decomposition_note,
        }
        self.assets.append(record)
        if not authorized:
            self.warn(
                f"{asset_id}: atomic raster has no explicit assets.json authorization; "
                "strict approval will fail"
            )

    def add_pending_connection(self, shape, element: ET.Element) -> None:
        source = element.get("data-source-id")
        target = element.get("data-target-id")
        if (not source and not target) or element.get("data-attach") == "false":
            return
        self.pending_connections.append(
            (
                shape,
                source,
                target,
                int(element.get("data-source-site", "0")),
                int(element.get("data-target-site", "0")),
            )
        )

    def resolve_connections(self) -> None:
        shape_ids = {
            binding["element_id"]: binding["shape_id"]
            for binding in self.bindings
            if binding["object_kind"] not in ("connector", "line", "freeform-arrow")
        }
        for connector, source, target, source_site, target_site in self.pending_connections:
            c_nv = connector._element.find(f".//{qn('p:cNvCxnSpPr')}")
            if c_nv is None:
                self.warn(f"{self.current_element_id}: connector has no cNvCxnSpPr; attachment skipped")
                continue
            for tag, element_id, site in (
                ("a:stCxn", source, source_site),
                ("a:endCxn", target, target_site),
            ):
                if not element_id:
                    continue
                shape_id = shape_ids.get(element_id)
                if shape_id is None:
                    self.warn(f"connector target id not found: {element_id}")
                    continue
                c_nv.append(c_nv.makeelement(qn(tag), {"id": str(shape_id), "idx": str(site)}))

    def group_arrow_parts(self, parts: list) -> None:
        """Create a real PowerPoint group for a shaft plus custom arrowheads."""
        unique = []
        seen: set[int] = set()
        for part in parts:
            marker = id(part)
            if marker not in seen:
                seen.add(marker)
                unique.append(part)
        if len(unique) < 2:
            return
        shape_id = _next_shape_id(self.slide)
        name = self._shape_name("arrow-group")
        width_emu = round(self.canvas_width * EMU_PER_PX)
        height_emu = round(self.canvas_height * EMU_PER_PX)
        group_xml = (
            '<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
            'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<p:nvGrpSpPr><p:cNvPr id="{shape_id}" name="{name}"/>'
            '<p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
            '<p:grpSpPr><a:xfrm><a:off x="0" y="0"/>'
            f'<a:ext cx="{width_emu}" cy="{height_emu}"/><a:chOff x="0" y="0"/>'
            f'<a:chExt cx="{width_emu}" cy="{height_emu}"/></a:xfrm></p:grpSpPr>'
            '</p:grpSp>'
        )
        group = _parse_sp_xml(group_xml)
        for part in unique:
            group.append(part)
        self.slide.shapes._spTree.append(group)
        child_ids = {
            int(node.get("id"))
            for part in unique
            for node in part.findall(f".//{qn('p:cNvPr')}")
            if node.get("id")
        }
        for binding in self.bindings:
            if binding["shape_id"] in child_ids:
                binding["group_shape_id"] = shape_id
                binding["physically_grouped"] = True
        self.register_raw(shape_id, name, "arrow-group", editable=True)
        self.scene_elements[self.current_element_id]["physically_grouped"] = True
        self.bump("arrow-group")


def _has_marker(element: ET.Element) -> bool:
    inline = _parse_style_attr(element.get("style"))
    return any(element.get(key) or inline.get(key) for key in ("marker-start", "marker-end"))


def _apply_fill_and_line(sp_pr, style: dict[str, str], ctx: ConvertContext) -> None:
    grad = _resolve_fill_url(style, ctx.defs)
    if grad is not None and grad.tag == f"{SVG_NS}linearGradient":
        stops = _gradient_stops(grad)
        if stops:
            _set_gradient_fill(sp_pr, stops, _gradient_angle(grad))
        else:
            _set_no_fill(sp_pr)
    else:
        if grad is not None:
            ctx.warn(f"暂不支持的渐变类型: {grad.tag}，按无填充处理")
        fill = parse_color(style.get("fill", "#000000"))
        if fill is None:
            _set_no_fill(sp_pr)
        else:
            hex_color, fill_alpha = fill
            alpha = fill_alpha * _opacity(style, "fill-opacity") * _opacity(style, "opacity")
            _set_solid_fill(sp_pr, hex_color, alpha)
    _apply_line(sp_pr, style)


def _emit_rect(ctx: ConvertContext, el: ET.Element, style: dict[str, str], matrix: Matrix) -> None:
    x = float(el.get("x", 0))
    y = float(el.get("y", 0))
    w = float(el.get("width", 0))
    h = float(el.get("height", 0))
    if w <= 0 or h <= 0:
        return
    element_id = el.get("id", "")
    if element_id.startswith("atomic:"):
        _emit_atomic(ctx, element_id, x, y, w, h, matrix)
        return
    if not matrix.is_axis_aligned():
        segments = [("M", *matrix.apply(x, y)), ("L", *matrix.apply(x + w, y)),
                    ("L", *matrix.apply(x + w, y + h)), ("L", *matrix.apply(x, y + h)), ("Z",)]
        _emit_freeform(ctx, segments, style)
        return
    (x0, y0), (x1, y1) = matrix.apply(x, y), matrix.apply(x + w, y + h)
    rx = float(el.get("rx", el.get("ry", 0)) or 0)
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rx > 0 else MSO_SHAPE.RECTANGLE
    shape = ctx.slide.shapes.add_shape(kind, _px(min(x0, x1)), _px(min(y0, y1)), _px(abs(x1 - x0)), _px(abs(y1 - y0)))
    if rx > 0:
        try:
            shape.adjustments[0] = min(0.5, rx / min(w, h))
        except (IndexError, AttributeError):
            pass
    _disable_shadow(shape)
    _apply_fill_and_line(shape._element.spPr, style, ctx)
    ctx.register_shape(shape, "rect")
    ctx.bump("rect")


def _emit_ellipse(ctx: ConvertContext, el: ET.Element, style: dict[str, str], matrix: Matrix) -> None:
    if el.tag == f"{SVG_NS}circle":
        cx, cy = float(el.get("cx", 0)), float(el.get("cy", 0))
        rx = ry = float(el.get("r", 0))
    else:
        cx, cy = float(el.get("cx", 0)), float(el.get("cy", 0))
        rx, ry = float(el.get("rx", 0)), float(el.get("ry", 0))
    if rx <= 0 or ry <= 0:
        return
    x0, y0 = matrix.apply(cx - rx, cy - ry)
    x1, y1 = matrix.apply(cx + rx, cy + ry)
    shape = ctx.slide.shapes.add_shape(
        MSO_SHAPE.OVAL, _px(min(x0, x1)), _px(min(y0, y1)), _px(abs(x1 - x0)), _px(abs(y1 - y0))
    )
    _disable_shadow(shape)
    _apply_fill_and_line(shape._element.spPr, style, ctx)
    ctx.register_shape(shape, "ellipse")
    ctx.bump("ellipse")


def _marker_reference(el: ET.Element, side: str) -> str | None:
    attr = f"marker-{side}"
    value = el.get(attr) or _parse_style_attr(el.get("style")).get(attr)
    match = re.match(r"url\(#([^)]+)\)", value or "")
    return match.group(1) if match else None


def _native_marker_spec(
    marker: ET.Element,
    line_style: dict[str, str],
    el: ET.Element,
    side: str,
) -> dict[str, str] | None:
    """Map a simple SVG triangle to a real DrawingML line end.

    Unsupported marker artwork is deliberately rejected here so the caller can
    emit an explicit custom-freeform fallback instead of silently substituting a
    default PowerPoint arrow.
    """
    if marker.get("orient", "auto") not in ("auto", "auto-start-reverse"):
        return None
    paths = [child for child in marker if child.tag == f"{SVG_NS}path"]
    if len(paths) != 1 or paths[0].get("transform"):
        return None
    segments = parse_path_d(paths[0].get("d", ""))
    points = [(part[1], part[2]) for part in segments if part[0] in ("M", "L")]
    if len(points) != 3:
        return None

    marker_style = _element_style(paths[0], _element_style(marker, {}))
    line_color = parse_color(line_style.get("stroke"))
    marker_color = parse_color(marker_style.get("fill")) or parse_color(marker_style.get("stroke"))
    if line_color is None or marker_color is None or line_color[0] != marker_color[0]:
        return None

    explicit_type = el.get(f"data-{side}-arrow-type") or el.get("data-arrow-type")
    closed = any(part[0] == "Z" for part in segments) or marker_style.get("fill") not in (
        None,
        "none",
        "transparent",
    )
    arrow_type = explicit_type or ("triangle" if closed else "arrow")
    if arrow_type not in {"triangle", "arrow", "stealth", "diamond", "oval"}:
        return None

    xs = [point[0] for point in points]
    ys = [point[1] for point in points]
    try:
        stroke_width = max(float(line_style.get("stroke-width", "1")), 0.01)
    except ValueError:
        stroke_width = 1.0
    length_ratio = max(xs) - min(xs)
    width_ratio = max(ys) - min(ys)

    def bucket(value: float) -> str:
        ratio = value / stroke_width
        if ratio <= 2.25:
            return "sm"
        if ratio <= 4.0:
            return "med"
        return "lg"

    width = el.get(f"data-{side}-arrow-width") or el.get("data-arrow-width")
    length = el.get(f"data-{side}-arrow-length") or el.get("data-arrow-length")
    width = width or bucket(width_ratio)
    length = length or bucket(length_ratio)
    if width not in {"sm", "med", "lg"} or length not in {"sm", "med", "lg"}:
        return None
    return {"type": arrow_type, "w": width, "len": length}


def _apply_native_markers(
    ctx: ConvertContext,
    sp_pr,
    el: ET.Element,
    style: dict[str, str],
) -> set[str]:
    line = sp_pr.find(qn("a:ln"))
    if line is None:
        return set()
    handled: set[str] = set()
    for side, tag in (("start", "a:headEnd"), ("end", "a:tailEnd")):
        marker_id = _marker_reference(el, side)
        if not marker_id:
            continue
        marker = ctx.defs.get(marker_id)
        spec = _native_marker_spec(marker, style, el, side) if marker is not None else None
        if spec is None:
            ctx.warn(
                f"{ctx.current_element_id}:{side}: marker {marker_id} cannot be represented "
                "as a native PowerPoint arrow; using grouped custom-freeform fallback"
            )
            continue
        old = line.find(qn(tag))
        if old is not None:
            line.remove(old)
        line.append(line.makeelement(qn(tag), spec))
        handled.add(side)
        ctx.bump("native-arrowhead")
    return handled


def _emit_straight_connector(
    ctx: ConvertContext,
    el: ET.Element,
    style: dict[str, str],
    start: tuple[float, float],
    end: tuple[float, float],
) -> None:
    x1, y1 = start
    x2, y2 = end
    connector = ctx.slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, _px(x1), _px(y1), _px(x2), _px(y2)
    )
    _disable_shadow(connector)
    _apply_line(connector._element.spPr, style)
    is_connector = _has_marker(el) or el.get("data-source-id") or el.get("data-target-id")
    ctx.register_shape(connector, "connector" if is_connector else "line")
    ctx.add_pending_connection(connector, el)
    handled = _apply_native_markers(ctx, connector._element.spPr, el, style)
    ctx.bump("line")
    marker_shapes = _emit_markers(
        ctx,
        el,
        style,
        [("M", x1, y1), ("L", x2, y2)],
        skip=handled,
    )
    if marker_shapes:
        ctx.group_arrow_parts([connector._element, *marker_shapes])


def _emit_line(ctx: ConvertContext, el: ET.Element, style: dict[str, str], matrix: Matrix) -> None:
    start = matrix.apply(float(el.get("x1", 0)), float(el.get("y1", 0)))
    end = matrix.apply(float(el.get("x2", 0)), float(el.get("y2", 0)))
    _emit_straight_connector(ctx, el, style, start, end)


def _emit_poly(ctx: ConvertContext, el: ET.Element, style: dict[str, str], matrix: Matrix, close: bool) -> None:
    points = [float(p) for p in re.split(r"[\s,]+", (el.get("points") or "").strip()) if p]
    if len(points) < 4:
        return
    segments: list[tuple] = []
    first = matrix.apply(points[0], points[1])
    segments.append(("M", *first))
    for i in range(2, len(points) - 1, 2):
        segments.append(("L", *matrix.apply(points[i], points[i + 1])))
    if close:
        segments.append(("Z",))
    shape = _emit_freeform(
        ctx,
        segments,
        style,
        object_kind="freeform-arrow" if _has_marker(el) else "freeform",
    )
    if shape is None:
        return
    handled = _apply_native_markers(ctx, shape.find(qn("p:spPr")), el, style)
    marker_shapes = _emit_markers(ctx, el, style, segments, skip=handled)
    if marker_shapes:
        ctx.group_arrow_parts([shape, *marker_shapes])


def _emit_path(ctx: ConvertContext, el: ET.Element, style: dict[str, str], matrix: Matrix) -> None:
    d = el.get("d")
    if not d:
        return
    segments = parse_path_d(d)
    if matrix != Matrix():
        segments = _transform_segments(segments, matrix)
    if len(segments) == 2 and segments[0][0] == "M" and segments[1][0] == "L":
        _emit_straight_connector(
            ctx,
            el,
            style,
            (segments[0][1], segments[0][2]),
            (segments[1][1], segments[1][2]),
        )
        return
    shape = _emit_freeform(
        ctx,
        segments,
        style,
        object_kind="freeform-arrow" if _has_marker(el) or el.get("data-role") == "arrow" else "freeform",
    )
    if shape is None:
        return
    handled = _apply_native_markers(ctx, shape.find(qn("p:spPr")), el, style)
    marker_shapes = _emit_markers(ctx, el, style, segments, skip=handled)
    if marker_shapes:
        ctx.group_arrow_parts([shape, *marker_shapes])


def _transform_segments(segments: list[tuple], matrix: Matrix) -> list[tuple]:
    result: list[tuple] = []
    for seg in segments:
        if seg[0] == "M" or seg[0] == "L":
            result.append((seg[0], *matrix.apply(seg[1], seg[2])))
        elif seg[0] == "C":
            p1 = matrix.apply(seg[1], seg[2])
            p2 = matrix.apply(seg[3], seg[4])
            p3 = matrix.apply(seg[5], seg[6])
            result.append(("C", p1[0], p1[1], p2[0], p2[1], p3[0], p3[1]))
        else:
            result.append(seg)
    return result


def _segment_vertices(segments: list[tuple]) -> list[tuple[float, float]]:
    vertices: list[tuple[float, float]] = []
    for seg in segments:
        if seg[0] in ("M", "L"):
            vertices.append((seg[1], seg[2]))
        elif seg[0] == "C":
            vertices.append((seg[5], seg[6]))
    return vertices


def _emit_freeform(
    ctx: ConvertContext,
    segments: list[tuple],
    style: dict[str, str],
    *,
    object_kind: str = "freeform",
    editable: bool = True,
):
    # bbox 必须包含贝塞尔控制点：控制多边形永远包住曲线本体，
    # 否则坐标越出声明的 path w/h，PowerPoint 会判定文件损坏。
    xs: list[float] = []
    ys: list[float] = []
    for s in segments:
        if s[0] in ("M", "L"):
            xs.append(s[1])
            ys.append(s[2])
        elif s[0] == "C":
            xs.extend((s[1], s[3], s[5]))
            ys.extend((s[2], s[4], s[6]))
    if not xs:
        return None
    min_x, max_x = min(xs), max(xs)
    min_y, max_y = min(ys), max(ys)
    w = max(max_x - min_x, 0.01)
    h = max(max_y - min_y, 0.01)

    def rel(x: float, y: float) -> tuple[int, int]:
        return (round((x - min_x) * EMU_PER_PX), round((y - min_y) * EMU_PER_PX))

    path_xml = [f'<a:path w="{round(w * EMU_PER_PX)}" h="{round(h * EMU_PER_PX)}">']
    for seg in segments:
        if seg[0] == "M":
            x, y = rel(seg[1], seg[2])
            path_xml.append(f'<a:moveTo><a:pt x="{x}" y="{y}"/></a:moveTo>')
        elif seg[0] == "L":
            x, y = rel(seg[1], seg[2])
            path_xml.append(f'<a:lnTo><a:pt x="{x}" y="{y}"/></a:lnTo>')
        elif seg[0] == "C":
            x1, y1 = rel(seg[1], seg[2])
            x2, y2 = rel(seg[3], seg[4])
            x3, y3 = rel(seg[5], seg[6])
            path_xml.append(
                f'<a:cubicBezTo><a:pt x="{x1}" y="{y1}"/><a:pt x="{x2}" y="{y2}"/>'
                f'<a:pt x="{x3}" y="{y3}"/></a:cubicBezTo>'
            )
        elif seg[0] == "Z":
            path_xml.append("<a:close/>")
    path_xml.append("</a:path>")

    shape_id = _next_shape_id(ctx.slide)
    xml = (
        '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:nvSpPr><p:cNvPr id="{shape_id}" name="freeform-{shape_id}"/>'
        "<p:cNvSpPr/><p:nvPr/></p:nvSpPr>"
        "<p:spPr>"
        f'<a:xfrm><a:off x="{round(min_x * EMU_PER_PX)}" y="{round(min_y * EMU_PER_PX)}"/>'
        f'<a:ext cx="{round(w * EMU_PER_PX)}" cy="{round(h * EMU_PER_PX)}"/></a:xfrm>'
        "<a:custGeom><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>"
        '<a:rect l="0" t="0" r="0" b="0"/>'
        f"<a:pathLst>{''.join(path_xml)}</a:pathLst></a:custGeom>"
        "</p:spPr></p:sp>"
    )
    sp = _parse_sp_xml(xml)
    ctx.slide.shapes._spTree.append(sp)
    sp_pr = sp.find(qn("p:spPr"))
    _apply_fill_and_line(sp_pr, style, ctx)
    ctx.register_xml(sp, shape_id, object_kind, editable=editable)
    ctx.bump("freeform")
    return sp


def _parse_sp_xml(xml: str):
    from lxml import etree

    return etree.fromstring(xml.encode("utf-8"))


def _emit_markers(
    ctx: ConvertContext,
    el: ET.Element,
    style: dict[str, str],
    segments: list[tuple],
    *,
    skip: set[str] | None = None,
) -> list:
    """按 SVG marker 语义放置箭头：marker-start 沿行进方向，marker-end 沿末端切线。

    末端切线取末段真实方向（C 段 = 端点 − 第二控制点），而非首末顶点弦方向——
    弦方向在曲线路径上偏差可达 40°+（01 案例 π/a 圆间曲线箭头横甩脱开的根因）。
    """
    skip = skip or set()
    emitted: list = []
    if len(segments) < 2:
        return emitted
    vertices = _segment_vertices(segments)
    if len(vertices) < 2:
        return emitted
    for side, attr, vertex, direction in (
        (
            "start",
            "marker-start",
            vertices[0],
            _start_tangent(segments) or _chord(vertices, forward=True),
        ),
        (
            "end",
            "marker-end",
            vertices[-1],
            _end_tangent(segments) or _chord(vertices, forward=False),
        ),
    ):
        if side in skip:
            continue
        ref = el.get(attr) or _parse_style_attr(el.get("style")).get(attr)
        match = re.match(r"url\(#([^)]+)\)", ref or "")
        if not match:
            continue
        marker = ctx.defs.get(match.group(1))
        if marker is None:
            continue
        angle = math.atan2(direction[1], direction[0])
        emitted.extend(_draw_marker(ctx, marker, vertex[0], vertex[1], angle, style))
    return emitted


def _start_tangent(segments: list[tuple]) -> tuple[float, float] | None:
    """起点处行进方向：首个绘制段 C → c1 − 起点；L → 终点 − 起点。"""
    start: tuple[float, float] | None = None
    for seg in segments:
        if seg[0] == "M":
            start = (seg[1], seg[2])
        elif start is not None and seg[0] in ("L", "C"):
            other = (seg[1], seg[2])
            direction = (other[0] - start[0], other[1] - start[1])
            return direction if math.hypot(*direction) > 1e-9 else None
    return None


def _end_tangent(segments: list[tuple]) -> tuple[float, float] | None:
    """末端切线方向：末段 C → 端点 − 第二控制点；L → 端点 − 前一 on-curve 点。"""
    prev: tuple[float, float] | None = None
    direction: tuple[float, float] | None = None
    for seg in segments:
        if seg[0] == "M":
            prev = (seg[1], seg[2])
        elif seg[0] == "L":
            if prev is not None:
                direction = (seg[1] - prev[0], seg[2] - prev[1])
            prev = (seg[1], seg[2])
        elif seg[0] == "C":
            direction = (seg[5] - seg[3], seg[6] - seg[4])
            prev = (seg[5], seg[6])
    if direction is not None and math.hypot(*direction) <= 1e-9:
        return None
    return direction


def _chord(vertices: list[tuple[float, float]], forward: bool) -> tuple[float, float]:
    first, last = vertices[0], vertices[-1]
    if forward:
        return (last[0] - first[0], last[1] - first[1])
    return (last[0] - first[0], last[1] - first[1])


def _draw_marker(
    ctx: ConvertContext, marker: ET.Element, x: float, y: float, angle: float, line_style: dict[str, str]
) -> list:
    ref_x = float(marker.get("refX", 0))
    ref_y = float(marker.get("refY", 0))
    placement = (
        Matrix(e=x, f=y)
        .multiply(Matrix(a=math.cos(angle), b=math.sin(angle), c=-math.sin(angle), d=math.cos(angle)))
        .multiply(Matrix(e=-ref_x, f=-ref_y))
    )
    emitted: list = []
    for child in marker:
        if child.tag != f"{SVG_NS}path":
            continue
        marker_style = _element_style(child, {})
        segments = parse_path_d(child.get("d", ""))
        segments = _transform_segments(segments, placement)
        shape = _emit_freeform(
            ctx,
            segments,
            marker_style,
            object_kind="arrowhead-fallback",
        )
        if shape is not None:
            emitted.append(shape)
    ctx.bump("marker")
    return emitted


def _attach_atomic_raster_tags(
    ctx: ConvertContext,
    picture,
    asset_id: str,
    source_sha256: str,
) -> None:
    """Persist the same hash-bound raster metadata read by powerpoint-live."""
    authorization = ctx.asset_authorizations.get(asset_id, {})
    if authorization.get("authorized") is not True:
        return
    tags = {
        "AISCIENTIFICILLUSTRATORASSETID": asset_id,
        "AISCIENTIFICILLUSTRATORSOURCESHA256": source_sha256,
        "AISCIENTIFICILLUSTRATORRASTERREASON": authorization.get("raster_reason")
        or "Reference-faithful irreducible creative microasset authorized by the user.",
        "AISCIENTIFICILLUSTRATORSOURCETIGHTLYCROPPED": "True",
        "AISCIENTIFICILLUSTRATORATOMICRASTERUNIT": "True",
        "AISCIENTIFICILLUSTRATORCONTAINSRECONSTRUCTABLECONTENT": "False",
        "AISCIENTIFICILLUSTRATORDECOMPOSITIONNOTE": authorization.get("decomposition_note")
        or "No faithful native decomposition; formal text, formulas, nodes, and topology remain native.",
    }
    tag_list = ET.Element(f"{{{PML_NS}}}tagLst")
    for name, value in tags.items():
        ET.SubElement(tag_list, f"{{{PML_NS}}}tag", {"name": name, "val": str(value)})
    tag_part = Part(
        ctx.slide.part.package.next_partname("/ppt/tags/tag%d.xml"),
        CT.PML_TAGS,
        ctx.slide.part.package,
        ET.tostring(tag_list, encoding="utf-8", xml_declaration=True),
    )
    relationship_id = ctx.slide.part.relate_to(tag_part, RT.TAGS)
    nv_pr = picture._element.find(f"{qn('p:nvPicPr')}/{qn('p:nvPr')}")
    if nv_pr is None:
        raise common.fail(f"{asset_id}: picture has no p:nvPr for raster metadata")
    old = nv_pr.find(qn("p:custDataLst"))
    if old is not None:
        nv_pr.remove(old)
    custom_data = OxmlElement("p:custDataLst")
    tag_reference = OxmlElement("p:tags")
    tag_reference.set(qn("r:id"), relationship_id)
    custom_data.append(tag_reference)
    nv_pr.append(custom_data)


def _emit_atomic(ctx: ConvertContext, element_id: str, x: float, y: float, w: float, h: float, matrix: Matrix) -> None:
    if ctx.source_png is None:
        ctx.warn(f"{element_id}: 缺少参考图，无法裁剪嵌入")
        return
    from PIL import Image

    x0, y0 = matrix.apply(x, y)
    x1, y1 = matrix.apply(x + w, y + h)
    left, top = round(min(x0, x1)), round(min(y0, y1))
    right, bottom = round(max(x0, x1)), round(max(y0, y1))
    with Image.open(ctx.source_png) as image:
        crop = image.crop((left, top, right, bottom))
        buffer = io.BytesIO()
        crop.save(buffer, format="PNG")
        buffer.seek(0)
    source_sha256 = hashlib.sha256(buffer.getvalue()).hexdigest()
    picture = ctx.slide.shapes.add_picture(
        buffer,
        _px(left),
        _px(top),
        _px(right - left),
        _px(bottom - top),
    )
    asset_id = ctx.current_element_id or element_id
    _attach_atomic_raster_tags(ctx, picture, asset_id, source_sha256)
    ctx.register_shape(picture, "atomic-raster", editable=False)
    ctx.register_asset(
        asset_id,
        [left, top, right - left, bottom - top],
        source_sha256,
    )
    ctx.bump("atomic")


IMAGE_MAX_AREA_RATIO = 0.5  # 护栏：<image> 覆盖画布 ≥ 该比例即拒绝（防整图截图冒充矢量交付）


def _emit_image(ctx: ConvertContext, el: ET.Element, matrix: Matrix) -> None:
    """<image> 合同容错：不读内嵌数据，按 bbox 从参考图裁剪替代（等价 atomic 占位符）。"""
    x = float(el.get("x", 0))
    y = float(el.get("y", 0))
    w = float(el.get("width", 0))
    h = float(el.get("height", 0))
    if w <= 0 or h <= 0:
        ctx.warn("<image> 缺少有效 x/y/width/height，已跳过")
        return
    ratio = w * h / ctx.canvas_area
    if ratio >= IMAGE_MAX_AREA_RATIO:
        raise common.fail(
            f"<image> 覆盖画布 {ratio:.0%}（≥{IMAGE_MAX_AREA_RATIO:.0%}），疑似整图截图冒充矢量，已拒绝"
        )
    ctx.warn('<image> 已按 bbox 从参考图裁剪替代（建议改用 <rect id="atomic:*"> 占位符）')
    _emit_atomic(ctx, "<image>", x, y, w, h, matrix)


# ---------------------------------------------------------------- 文本

_CHAR_WIDTHS = {
    " ": 0.28, "i": 0.28, "l": 0.28, "j": 0.3, "t": 0.33, "f": 0.33, "r": 0.36,
    "m": 0.85, "w": 0.8, "I": 0.3, "M": 0.85, "W": 0.9,
    ".": 0.3, ",": 0.3, ":": 0.33, ";": 0.33, "'": 0.25, "(": 0.36, ")": 0.36,
    "-": 0.36, "…": 0.8, "τ": 0.5, "π": 0.55, "ƒ": 0.4,
}


def _estimate_width(text: str, font_size: float) -> float:
    total = 0.0
    for ch in text:
        if ch in _CHAR_WIDTHS:
            total += _CHAR_WIDTHS[ch]
        elif ch.isupper():
            total += 0.67
        elif ch.isdigit():
            total += 0.56
        elif ord(ch) >= 0x2E80:
            total += 1.0
        else:
            total += 0.52
    return total * font_size


def _collect_text_runs(el: ET.Element, base_style: dict[str, str]) -> list[dict]:
    """把 <text> 的文字与 <tspan> 子节点拍平为 run 序列。"""
    runs: list[dict] = []

    def add_run(text: str, style: dict[str, str], shift: str | None) -> None:
        if text:
            runs.append({"text": text, "style": dict(style), "shift": shift})

    if el.text:
        add_run(el.text, base_style, None)
    for tspan in el:
        if tspan.tag != f"{SVG_NS}tspan":
            continue
        style = dict(base_style)
        inline = _parse_style_attr(tspan.get("style"))
        for key in ("font-size", "font-family", "font-style", "font-weight", "fill"):
            if tspan.get(key) is not None:
                style[key] = tspan.get(key)
            if key in inline:
                style[key] = inline[key]
        shift = tspan.get("baseline-shift") or inline.get("baseline-shift")
        if tspan.text:
            add_run(tspan.text, style, shift)
        if tspan.tail:
            add_run(tspan.tail, base_style, None)
    return runs


def _layout_number(el: ET.Element, name: str, default: float) -> float:
    raw = el.get(name)
    if raw is None:
        return default
    try:
        value = float(raw)
    except ValueError as exc:
        raise common.fail(f"{el.get('id') or 'text'}: {name} 必须是像素数值") from exc
    if not math.isfinite(value) or value < 0:
        raise common.fail(f"{el.get('id') or 'text'}: {name} 必须是有限非负数")
    return value


def _constrain_text_box(
    ctx: ConvertContext,
    el: ET.Element,
    *,
    anchor: str,
    anchor_x: float,
    left: float,
    top: float,
    width: float,
    height: float,
    rotation_deg: float,
) -> tuple[float, float, float, float]:
    """Constrain an explicitly annotated text box to its SVG container.

    The anchor is preserved.  Only transparent selection-box surplus is
    clipped; if the declared anchor itself lies outside the container the
    source contract is invalid and conversion fails instead of silently moving
    visible content.
    """
    container_id = el.get("data-layout-container")
    if not container_id:
        return left, top, width, height
    container = ctx.layout_boxes.get(container_id)
    if container is None:
        raise common.fail(
            f"{el.get('id') or 'text'}: data-layout-container={container_id!r} "
            "未指向可测量的 SVG 形状"
        )
    if rotation_deg:
        raise common.fail(
            f"{el.get('id') or 'text'}: 旋转文字的容器约束必须先展平为明确 bbox"
        )
    padding = _layout_number(el, "data-layout-padding", 0.0)
    inner_left = container.x + padding
    inner_top = container.y + padding
    inner_right = container.right - padding
    inner_bottom = container.bottom - padding
    if inner_right <= inner_left or inner_bottom <= inner_top:
        raise common.fail(f"{el.get('id') or 'text'}: layout padding 抹空了容器 {container_id!r}")
    if not inner_left <= anchor_x <= inner_right:
        raise common.fail(
            f"{el.get('id') or 'text'}: 文字锚点 x={anchor_x:.3f} 位于容器 {container_id!r} 外"
        )
    original = (left, top, width, height)
    if anchor == "middle":
        width = min(width, 2.0 * min(anchor_x - inner_left, inner_right - anchor_x))
        left = anchor_x - width / 2.0
    elif anchor == "end":
        width = min(width, anchor_x - inner_left)
        left = anchor_x - width
    else:
        left = max(left, inner_left)
        width = min(width, inner_right - left)
    if top < inner_top:
        # A top shift would move the visible baseline and hide a source error.
        raise common.fail(
            f"{el.get('id') or 'text'}: 文本框上边界 {top:.3f} 越出容器 {container_id!r}"
        )
    height = min(height, inner_bottom - top)
    if width <= 0 or height <= 0:
        raise common.fail(
            f"{el.get('id') or 'text'}: 容器 {container_id!r} 内没有可用文本区域"
        )
    constrained = (left, top, width, height)
    if any(abs(before - after) > 0.01 for before, after in zip(original, constrained)):
        ctx.warn(
            f"{el.get('id') or 'text'}: transparent text box constrained to "
            f"{container_id} (padding={padding:g}px)"
        )
    return constrained


def _emit_text(ctx: ConvertContext, el: ET.Element, style: dict[str, str], matrix: Matrix) -> None:
    x = float(el.get("x", 0))
    y = float(el.get("y", 0))
    try:
        font_size = float(style.get("font-size", "16"))
    except ValueError:
        font_size = 16.0
    runs = _collect_text_runs(el, style)
    full_text = "".join(r["text"] for r in runs)
    if not full_text.strip():
        return

    anchor = style.get("text-anchor", "start")
    text_w = max(_estimate_width(full_text, font_size), 1.0)
    box_w = text_w + TEXT_BOX_PAD_X
    box_h = font_size * 1.25 + TEXT_BOX_PAD_Y

    rotation_deg = 0.0
    rot_match = re.fullmatch(r"rotate\(\s*(-?[\d.]+)[\s,]+(-?[\d.]+)[\s,]+(-?[\d.]+)\s*\)", el.get("transform", "").strip())
    if rot_match:
        rotation_deg = float(rot_match.group(1))
        cx, cy = float(rot_match.group(2)), float(rot_match.group(3))
        cx, cy = matrix.apply(cx, cy)
        box_w = min(box_w, max(2 * min(cx, ctx.canvas_width - cx), 0.5))
        box_h = min(box_h, max(2 * min(cy, ctx.canvas_height - cy), 0.5))
        left, top = cx - box_w / 2, cy - box_h / 2
        align = PP_ALIGN.CENTER
        vertical = MSO_ANCHOR.MIDDLE
    else:
        x, y = matrix.apply(x, y)
        if anchor == "middle":
            box_w = min(box_w, max(2 * min(x, ctx.canvas_width - x), 0.5))
            left, align = x - box_w / 2, PP_ALIGN.CENTER
        elif anchor == "end":
            box_w = min(box_w, max(x, 0.5))
            left, align = x - box_w, PP_ALIGN.RIGHT
        else:
            box_w = min(box_w, max(ctx.canvas_width - x, 0.5))
            left, align = x, PP_ALIGN.LEFT
        top = y - font_size * BASELINE_ASCENT
        box_h = min(box_h, max(ctx.canvas_height - top, 0.5))
        vertical = MSO_ANCHOR.TOP

    left, top, box_w, box_h = _constrain_text_box(
        ctx,
        el,
        anchor=anchor,
        anchor_x=x if not rotation_deg else left + box_w / 2.0,
        left=left,
        top=top,
        width=box_w,
        height=box_h,
        rotation_deg=rotation_deg,
    )

    textbox = ctx.slide.shapes.add_textbox(_px(left), _px(top), _px(box_w), _px(box_h))
    _disable_shadow(textbox)
    if rotation_deg:
        textbox.rotation = rotation_deg
    frame = textbox.text_frame
    frame.word_wrap = False
    frame.vertical_anchor = vertical
    for margin in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(frame, margin, 0)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align

    for run_info in runs:
        run = paragraph.add_run()
        run.text = run_info["text"]
        rstyle = run_info["style"]
        try:
            size = float(rstyle.get("font-size", font_size))
        except ValueError:
            size = font_size
        run.font.size = Pt(size * PT_PER_PX)
        family = (rstyle.get("font-family") or "Arial").split(",")[0].strip()
        run.font.name = family
        if str(rstyle.get("font-style", "")).lower() in ("italic", "oblique"):
            run.font.italic = True
        if str(rstyle.get("font-weight", "")).lower() in ("bold", "600", "700", "800", "900"):
            run.font.bold = True
        color = parse_color(rstyle.get("fill", "#000000"))
        if color:
            run.font.color.rgb = RGBColor.from_string(color[0][1:])
        if run_info["shift"] in ("super", "sup"):
            run._r.get_or_add_rPr().set("baseline", "30000")
        elif run_info["shift"] == "sub":
            run._r.get_or_add_rPr().set("baseline", "-25000")
    ctx.register_shape(textbox, "text")
    ctx.bump("text")


# ---------------------------------------------------------------- 树遍历与主流程


def _emit_bound(ctx: ConvertContext, element: ET.Element, tag: str, emitter, *args) -> None:
    previous = ctx.begin_element(element, tag)
    try:
        emitter(ctx, element, *args)
    finally:
        ctx.end_element(previous)


def _walk(ctx: ConvertContext, element: ET.Element, style: dict[str, str], matrix: Matrix) -> None:
    tag = element.tag
    if tag in (f"{SVG_NS}defs", f"{SVG_NS}marker", f"{SVG_NS}linearGradient", f"{SVG_NS}radialGradient"):
        return
    own_style = _element_style(element, style)
    own_matrix = matrix.multiply(parse_transform(element.get("transform")))
    if tag == f"{SVG_NS}svg" or tag == f"{SVG_NS}g":
        for child in element:
            _walk(ctx, child, own_style, own_matrix)
    elif tag == f"{SVG_NS}rect":
        _emit_bound(ctx, element, "rect", _emit_rect, own_style, own_matrix)
    elif tag in (f"{SVG_NS}circle", f"{SVG_NS}ellipse"):
        _emit_bound(ctx, element, tag.replace(SVG_NS, ""), _emit_ellipse, own_style, own_matrix)
    elif tag == f"{SVG_NS}line":
        _emit_bound(ctx, element, "line", _emit_line, own_style, own_matrix)
    elif tag == f"{SVG_NS}polyline":
        _emit_bound(ctx, element, "polyline", _emit_poly, own_style, own_matrix, False)
    elif tag == f"{SVG_NS}polygon":
        _emit_bound(ctx, element, "polygon", _emit_poly, own_style, own_matrix, True)
    elif tag == f"{SVG_NS}path":
        _emit_bound(ctx, element, "path", _emit_path, own_style, own_matrix)
    elif tag == f"{SVG_NS}text":
        _emit_bound(ctx, element, "text", _emit_text, own_style, own_matrix)
    elif tag == f"{SVG_NS}image":
        _emit_bound(ctx, element, "image", _emit_image, own_matrix)
    else:
        ctx.warn(f"未知元素已跳过: {tag}")


def _svg_dimension(value: str | None) -> float | None:
    if value is None:
        return None
    match = re.fullmatch(r"\s*([-+]?(?:\d*\.\d+|\d+\.?))(?:px)?\s*", value)
    return float(match.group(1)) if match else None


def _iter_readback_shapes(shapes):
    for shape in shapes:
        yield shape
        children = getattr(shape, "shapes", None)
        if children is not None:
            yield from _iter_readback_shapes(children)


def _write_conversion_contracts(run: common.Run, ctx: ConvertContext, reopened) -> dict:
    from tools.v2.contracts import read_json, utc_now, write_json

    pptx_hash = common.sha256_file(run.pptx_path)
    readback_shapes = list(_iter_readback_shapes(reopened.slides[0].shapes))
    readback_names = {shape.name for shape in readback_shapes}
    for binding in ctx.bindings:
        binding["readback_found"] = binding["shape_name"] in readback_names

    scene = read_json(run.scene_path)
    existing_elements = {item["id"]: item for item in scene.get("elements", [])}
    merged_elements = []
    for element_id, generated in ctx.scene_elements.items():
        # Preserve model/operator annotations that are not regenerated from SVG,
        # but keep rendered geometry, topology and z-order authoritative.  Letting
        # a stale scene record override these fields made a corrected DOM order
        # disagree with the actual PowerPoint shape stack on the next conversion.
        merged_elements.append({**existing_elements.pop(element_id, {}), **generated})
    # The current SVG is the complete offline scene carrier.  Retain manual
    # annotations only for IDs that still exist; carrying unmatched historical
    # elements creates stale scene objects with no PowerPoint binding.
    merged_elements.sort(key=lambda item: (int(item.get("z_index", 0)), item["id"]))
    scene.update(
        {
            "updated_at": utc_now(),
            "elements": merged_elements,
            "edges": [
                {"id": item["id"], **item.get("topology", {})}
                for item in merged_elements
                if item.get("kind") == "edge" or item.get("topology")
            ],
            "artifact": {
                "backend": "pptx-offline",
                "path": str(run.pptx_path),
                "sha256": pptx_hash,
            },
        }
    )
    write_json(run.scene_path, scene)

    assets = read_json(run.assets_path)
    existing_assets = {item["id"]: item for item in assets.get("assets", [])}
    generated_assets = []
    for generated in ctx.assets:
        existing = existing_assets.pop(generated["id"], {})
        generated_assets.append({**generated, **existing, "editable": False})
    generated_assets.extend(existing_assets.values())
    assets.update({"updated_at": utc_now(), "assets": generated_assets})
    write_json(run.assets_path, assets)

    bindings_complete = bool(ctx.bindings) and all(
        binding["readback_found"] for binding in ctx.bindings
    )
    bindings = read_json(run.bindings_path)
    bindings.update(
        {
            "updated_at": utc_now(),
            "backend": "pptx-offline",
            "artifact_sha256": pptx_hash,
            "saved_reopened": True,
            "bindings_complete": bindings_complete,
            "bindings": ctx.bindings,
        }
    )
    write_json(run.bindings_path, bindings)
    return {
        "pptx_sha256": pptx_hash,
        "object_count": len(readback_shapes),
        "bindings_complete": bindings_complete,
    }


def convert(run: common.Run) -> dict:
    if not run.redraw_svg.is_file():
        raise common.fail(f"未找到 SVG: {run.redraw_svg}（先把 GPT 输出保存到这里）")
    from tools.v2.contracts import initialize_contracts, read_json, transition

    meta = initialize_contracts(run)
    current_state = meta["workflow"]["state"]
    if current_state in {"approved", "candidate", "qa_failed"}:
        transition(run, "repairing", "offline-conversion-started")
    width, height = int(meta["width"]), int(meta["height"])

    tree = ET.parse(run.redraw_svg)
    root = tree.getroot()
    required_view_box = root.get("viewBox", "")
    required_parts = [
        float(part) for part in re.split(r"[\s,]+", required_view_box.strip()) if part
    ]
    if len(required_parts) != 4:
        raise common.fail("SVG viewBox is required and must contain four numbers")
    if (
        abs(required_parts[0]) > 1e-6
        or abs(required_parts[1]) > 1e-6
        or round(required_parts[2]) != width
        or round(required_parts[3]) != height
    ):
        raise common.fail(
            f"SVG viewBox {required_parts} must be [0, 0, {width}, {height}] "
            "to preserve reference coordinates"
        )
    svg_width = _svg_dimension(root.get("width"))
    svg_height = _svg_dimension(root.get("height"))
    if svg_width is None or svg_height is None:
        raise common.fail("SVG root width/height must be numeric pixel dimensions")
    if round(svg_width) != width or round(svg_height) != height:
        raise common.fail(
            f"SVG root size {svg_width}x{svg_height} does not match reference {width}x{height}"
        )
    view_box = root.get("viewBox", "")
    if view_box:
        parts = [float(p) for p in re.split(r"[\s,]+", view_box.strip())]
        if len(parts) == 4 and (round(parts[2]) != width or round(parts[3]) != height):
            raise common.fail(
                f"SVG viewBox {parts[2]}x{parts[3]} 与参考图 {width}x{height} 不一致，违反输出合同"
            )

    prs = Presentation()
    prs.slide_width = _px(width)
    prs.slide_height = _px(height)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    defs = _collect_defs(root)
    source_png = run.source_png if run.source_png.is_file() else None
    assets_document = read_json(run.assets_path)
    asset_authorizations = {
        item["id"]: item for item in assets_document.get("assets", []) if item.get("id")
    }
    ctx = ConvertContext(
        slide,
        defs,
        source_png,
        width,
        height,
        asset_authorizations=asset_authorizations,
        layout_boxes=collect_svg_boxes(root),
    )
    _walk(ctx, root, {}, Matrix())
    ctx.resolve_connections()

    run.qa_dir.mkdir(exist_ok=True)
    prs.save(run.pptx_path)

    # 读回统计（机械验收的基础）
    reopened = Presentation(run.pptx_path)
    readback_texts = 0
    for shape in _iter_readback_shapes(reopened.slides[0].shapes):
        if shape.has_text_frame and shape.text_frame.text.strip():
            readback_texts += 1
    contract_summary = _write_conversion_contracts(run, ctx, reopened)
    from tools.v2.layout import audit_layout

    layout_report = audit_layout(run)
    summary = {
        "svg": str(run.redraw_svg),
        "pptx": str(run.pptx_path),
        "slide_count": len(reopened.slides._sldIdLst),
        "shape_count": len(reopened.slides[0].shapes),
        "textbox_with_text": readback_texts,
        "emitted": ctx.counts,
        "warnings": ctx.warnings,
        "layout_pass": layout_report["pass"],
        "layout_findings": len(layout_report["findings"]),
        **contract_summary,
    }
    (run.qa_dir / "convert-summary.json").write_text(
        json.dumps(summary, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )
    transition(run, "candidate", "offline-conversion-complete", details=contract_summary)
    return summary


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(prog="autofigure convert", description=__doc__)
    parser.add_argument("run_dir", type=Path, help="run 目录")
    parser.add_argument("--no-render", action="store_true", help="跳过 PowerPoint fresh render")
    args = parser.parse_args(argv)

    run = common.open_run(args.run_dir)
    summary = convert(run)
    sys.stdout.write(json.dumps(summary, ensure_ascii=False, indent=2) + "\n")
    if summary["warnings"]:
        sys.stdout.write(f"注意：{len(summary['warnings'])} 条 warning，详见 convert-summary.json\n")
    sys.stdout.write(f"PPTX 已生成: {run.pptx_path}\n")

    if not args.no_render:
        from tools.v2 import render_export

        meta = run.load_meta()
        render_export.render(run.pptx_path, run.render_png, int(meta["width"]), int(meta["height"]))
        sys.stdout.write(f"fresh render: {run.render_png}\n")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
