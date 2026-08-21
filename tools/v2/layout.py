"""Explicit layout contracts for SVG sources and saved PowerPoint objects.

The visual model is responsible for measuring a reference, but it must declare
the relationships that deterministic tooling can verify:

``data-layout-container``
    ID of the shape that must contain this element.  ``data-layout-padding``
    and ``data-layout-tolerance`` are optional pixel values.

``data-repeat-group`` / ``data-repeat-axis`` / ``data-repeat-order``
    Declare a repeated horizontal or vertical motif.  Equal size, a shared
    cross-axis centre and regular centre-to-centre spacing are audited in both
    SVG source coordinates and the saved/reopened PowerPoint artifact.

These annotations deliberately use a separate namespace from
``data-owner-id``.  Existing owner IDs describe semantic label/edge ownership
and do not imply geometric containment.

Saved PowerPoint geometry also has one global, annotation-free contract: every
scene-bound object, including OMML hidden in ``mc:AlternateContent``, must stay
inside the slide canvas.  This package-level check avoids false results from
tools that can move only the python-pptx-visible fallback object.
"""

from __future__ import annotations

import json
import math
import re
import xml.etree.ElementTree as ET
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation

from tools.v2 import common
from tools.v2.contracts import read_json
from tools.v2.svggeom import Matrix, parse_path_d, parse_transform

SVG_NS = "{http://www.w3.org/2000/svg}"

DEFAULT_CONTAINMENT_TOLERANCE_PX = 0.25
DEFAULT_CANVAS_TOLERANCE_PX = 0.25
DEFAULT_REPEAT_SIZE_TOLERANCE_PX = 0.25
DEFAULT_REPEAT_AXIS_TOLERANCE_PX = 0.25
# Integer-pixel reference reconstructions commonly distribute an odd remainder
# as N and N+1 pixel steps.  A 1 px range is regular; larger drift is not.
DEFAULT_REPEAT_SPACING_TOLERANCE_PX = 1.0

_DRAWABLE = {
    "rect",
    "circle",
    "ellipse",
    "line",
    "polyline",
    "polygon",
    "path",
    "text",
    "image",
}
_SKIP = {"defs", "marker", "linearGradient", "radialGradient"}


@dataclass(frozen=True)
class Box:
    x: float
    y: float
    width: float
    height: float

    @property
    def right(self) -> float:
        return self.x + self.width

    @property
    def bottom(self) -> float:
        return self.y + self.height

    @property
    def center_x(self) -> float:
        return self.x + self.width / 2.0

    @property
    def center_y(self) -> float:
        return self.y + self.height / 2.0

    def as_dict(self) -> dict[str, float]:
        return {
            "x": round(self.x, 6),
            "y": round(self.y, 6),
            "width": round(self.width, 6),
            "height": round(self.height, 6),
        }


@dataclass(frozen=True)
class SvgRecord:
    element_id: str | None
    tag: str
    element: ET.Element
    matrix: Matrix
    box: Box | None


def _local_name(element: ET.Element) -> str:
    return element.tag.rsplit("}", 1)[-1]


def _number(element: ET.Element, name: str, default: float = 0.0) -> float:
    raw = element.get(name)
    if raw is None:
        return default
    try:
        return float(raw)
    except ValueError as exc:
        raise ValueError(f"{element.get('id') or _local_name(element)}: invalid {name}={raw!r}") from exc


def _points_box(points: Iterable[tuple[float, float]], matrix: Matrix) -> Box | None:
    transformed = [matrix.apply(x, y) for x, y in points]
    if not transformed:
        return None
    xs = [point[0] for point in transformed]
    ys = [point[1] for point in transformed]
    return Box(min(xs), min(ys), max(xs) - min(xs), max(ys) - min(ys))


def _element_box(element: ET.Element, matrix: Matrix) -> Box | None:
    tag = _local_name(element)
    if tag in {"rect", "image"}:
        x, y = _number(element, "x"), _number(element, "y")
        width, height = _number(element, "width"), _number(element, "height")
        return _points_box(
            ((x, y), (x + width, y), (x, y + height), (x + width, y + height)),
            matrix,
        )
    if tag == "circle":
        cx, cy, radius = _number(element, "cx"), _number(element, "cy"), _number(element, "r")
        return _points_box(
            (
                (cx - radius, cy - radius),
                (cx + radius, cy - radius),
                (cx - radius, cy + radius),
                (cx + radius, cy + radius),
            ),
            matrix,
        )
    if tag == "ellipse":
        cx, cy = _number(element, "cx"), _number(element, "cy")
        rx, ry = _number(element, "rx"), _number(element, "ry")
        return _points_box(
            ((cx - rx, cy - ry), (cx + rx, cy - ry), (cx - rx, cy + ry), (cx + rx, cy + ry)),
            matrix,
        )
    if tag == "line":
        return _points_box(
            ((_number(element, "x1"), _number(element, "y1")), (_number(element, "x2"), _number(element, "y2"))),
            matrix,
        )
    if tag in {"polyline", "polygon"}:
        values = [float(item) for item in re.findall(r"[-+]?(?:\d*\.\d+|\d+\.?)", element.get("points", ""))]
        return _points_box(zip(values[0::2], values[1::2]), matrix)
    if tag == "path":
        points: list[tuple[float, float]] = []
        for segment in parse_path_d(element.get("d", "")):
            if segment[0] in {"M", "L"}:
                points.append((segment[1], segment[2]))
            elif segment[0] == "C":
                points.extend(
                    ((segment[1], segment[2]), (segment[3], segment[4]), (segment[5], segment[6]))
                )
        return _points_box(points, matrix)
    # Text bounds depend on the actual PowerPoint font/OMML renderer.  They are
    # intentionally judged from the saved artifact, not guessed from SVG text.
    return None


def collect_svg_records(root: ET.Element) -> list[SvgRecord]:
    """Return drawable records with the same inherited transform semantics as convert."""
    records: list[SvgRecord] = []

    def walk(element: ET.Element, parent_matrix: Matrix) -> None:
        tag = _local_name(element)
        if tag in _SKIP:
            return
        matrix = parent_matrix.multiply(parse_transform(element.get("transform")))
        if tag in {"svg", "g"}:
            for child in element:
                walk(child, matrix)
            return
        if tag not in _DRAWABLE:
            return
        records.append(
            SvgRecord(
                element_id=element.get("id"),
                tag=tag,
                element=element,
                matrix=matrix,
                box=_element_box(element, matrix),
            )
        )

    walk(root, Matrix())
    return records


def collect_svg_boxes(root: ET.Element) -> dict[str, Box]:
    """Boxes for explicitly identified SVG elements, used by the converter."""
    return {
        record.element_id: record.box
        for record in collect_svg_records(root)
        if record.element_id and record.box is not None
    }


def _float_attr(element: ET.Element, name: str, default: float) -> float:
    value = element.get(name)
    if value is None:
        return default
    try:
        parsed = float(value)
    except ValueError as exc:
        raise ValueError(f"{element.get('id') or _local_name(element)}: invalid {name}={value!r}") from exc
    if not math.isfinite(parsed) or parsed < 0:
        raise ValueError(f"{element.get('id') or _local_name(element)}: {name} must be finite and >= 0")
    return parsed


def _overflow(child: Box, container: Box, padding: float) -> dict[str, float]:
    inner_left = container.x + padding
    inner_top = container.y + padding
    inner_right = container.right - padding
    inner_bottom = container.bottom - padding
    return {
        "left": max(0.0, inner_left - child.x),
        "top": max(0.0, inner_top - child.y),
        "right": max(0.0, child.right - inner_right),
        "bottom": max(0.0, child.bottom - inner_bottom),
    }


def _union(boxes: list[Box]) -> Box:
    left = min(box.x for box in boxes)
    top = min(box.y for box in boxes)
    right = max(box.right for box in boxes)
    bottom = max(box.bottom for box in boxes)
    return Box(left, top, right - left, bottom - top)


def _backend_boxes(run: common.Run) -> tuple[dict[str, Box], list[str]]:
    """Resolve scene element IDs to saved PowerPoint shape bounds."""
    from tools.v2.live_bridge import _shape_bounds, _walk_shapes, _xml_shape_bounds

    presentation = Presentation(run.pptx_path)
    normal_by_name: dict[str, Box] = {}
    normal_by_id: dict[int, Box] = {}
    for shape in _walk_shapes(presentation.slides[0].shapes):
        raw = _shape_bounds(shape)
        box = Box(raw["x"], raw["y"], raw["width"], raw["height"])
        normal_by_name[shape.name] = box
        normal_by_id[int(shape.shape_id)] = box
    xml_by_name_raw, xml_by_id_raw = _xml_shape_bounds(run.pptx_path)
    xml_by_name = {
        name: Box(value["x"], value["y"], value["width"], value["height"])
        for name, value in xml_by_name_raw.items()
    }
    xml_by_id = {
        shape_id: Box(value["x"], value["y"], value["width"], value["height"])
        for shape_id, value in xml_by_id_raw.items()
    }

    resolved: dict[str, list[Box]] = {}
    missing: list[str] = []
    bindings = read_json(run.bindings_path)
    for binding in bindings.get("bindings", []):
        name = binding.get("shape_name")
        shape_id = binding.get("shape_id")
        box = (
            normal_by_name.get(name)
            or normal_by_id.get(shape_id)
            or xml_by_name.get(name)
            or xml_by_id.get(shape_id)
        )
        element_id = binding.get("element_id")
        if not element_id:
            continue
        if box is None:
            missing.append(element_id)
        else:
            resolved.setdefault(element_id, []).append(box)
    return {element_id: _union(boxes) for element_id, boxes in resolved.items()}, sorted(set(missing))


def _backend_canvas(run: common.Run) -> Box:
    """Return the saved PowerPoint slide canvas in the project's 96-DPI pixels."""
    from tools.v2.live_bridge import EMU_PER_PX

    presentation = Presentation(run.pptx_path)
    return Box(
        0.0,
        0.0,
        float(presentation.slide_width) / EMU_PER_PX,
        float(presentation.slide_height) / EMU_PER_PX,
    )


def _finding(
    code: str,
    stage: str,
    target: str,
    message: str,
    *,
    metrics: dict[str, Any] | None = None,
) -> dict[str, Any]:
    result: dict[str, Any] = {
        "code": code,
        "stage": stage,
        "severity": "error",
        "target": target,
        "message": message,
    }
    if metrics:
        result["metrics"] = metrics
    return result


def _repeat_metrics(boxes: list[tuple[int, str, Box]], axis: str) -> dict[str, Any]:
    ordered = sorted(boxes, key=lambda item: (item[0], item[1]))
    widths = [item[2].width for item in ordered]
    heights = [item[2].height for item in ordered]
    centers = [
        item[2].center_y if axis == "vertical" else item[2].center_x
        for item in ordered
    ]
    cross_centers = [
        item[2].center_x if axis == "vertical" else item[2].center_y
        for item in ordered
    ]
    steps = [centers[index + 1] - centers[index] for index in range(len(centers) - 1)]
    return {
        "members": [item[1] for item in ordered],
        "widths_px": [round(value, 6) for value in widths],
        "heights_px": [round(value, 6) for value in heights],
        "centers_px": [round(value, 6) for value in centers],
        "cross_axis_centers_px": [round(value, 6) for value in cross_centers],
        "steps_px": [round(value, 6) for value in steps],
        "width_range_px": round(max(widths) - min(widths), 6),
        "height_range_px": round(max(heights) - min(heights), 6),
        "cross_axis_range_px": round(max(cross_centers) - min(cross_centers), 6),
        "spacing_range_px": round(max(steps) - min(steps), 6) if steps else 0.0,
    }


def audit_layout(run: common.Run) -> dict[str, Any]:
    """Audit explicit layout annotations in SVG and saved PowerPoint geometry."""
    root = ET.parse(run.redraw_svg).getroot()
    records = collect_svg_records(root)
    source_boxes = {
        record.element_id: record.box
        for record in records
        if record.element_id and record.box is not None
    }
    backend_boxes, missing_backend = _backend_boxes(run)
    backend_canvas = _backend_canvas(run)
    findings: list[dict[str, Any]] = []

    canvas_rows: list[dict[str, Any]] = []
    for element_id, box in sorted(backend_boxes.items()):
        overflow = _overflow(box, backend_canvas, 0.0)
        row = {
            "element": element_id,
            "box": box.as_dict(),
            "overflow_px": {key: round(value, 6) for key, value in overflow.items()},
        }
        canvas_rows.append(row)
        if max(overflow.values()) > DEFAULT_CANVAS_TOLERANCE_PX:
            findings.append(
                _finding(
                    "L10",
                    "backend",
                    element_id,
                    "saved PowerPoint object exceeds the slide canvas",
                    metrics=row,
                )
            )

    annotated = [
        record
        for record in records
        if record.element.get("data-layout-container") or record.element.get("data-repeat-group")
    ]
    for record in annotated:
        if not record.element_id:
            findings.append(
                _finding(
                    "L1",
                    "source",
                    f"anonymous-{record.tag}",
                    "layout-annotated elements require an explicit stable id",
                )
            )

    containment_rows: list[dict[str, Any]] = []
    for record in annotated:
        child_id = record.element_id
        container_id = record.element.get("data-layout-container")
        if not child_id or not container_id:
            continue
        padding = _float_attr(record.element, "data-layout-padding", 0.0)
        tolerance = _float_attr(
            record.element,
            "data-layout-tolerance",
            DEFAULT_CONTAINMENT_TOLERANCE_PX,
        )
        row: dict[str, Any] = {
            "element": child_id,
            "container": container_id,
            "padding_px": padding,
            "tolerance_px": tolerance,
        }
        source_child = source_boxes.get(child_id)
        source_container = source_boxes.get(container_id)
        if source_container is None:
            findings.append(
                _finding("L2", "source", child_id, f"layout container {container_id!r} has no measurable SVG geometry")
            )
        elif source_child is not None:
            overflow = _overflow(source_child, source_container, padding)
            row["source"] = {
                "child": source_child.as_dict(),
                "container": source_container.as_dict(),
                "overflow_px": {key: round(value, 6) for key, value in overflow.items()},
            }
            if max(overflow.values()) > tolerance:
                findings.append(
                    _finding(
                        "L3",
                        "source",
                        child_id,
                        f"element exceeds SVG container {container_id!r}",
                        metrics=row["source"],
                    )
                )

        backend_child = backend_boxes.get(child_id)
        backend_container = backend_boxes.get(container_id)
        if backend_child is None or backend_container is None:
            missing = [
                element_id
                for element_id, box in ((child_id, backend_child), (container_id, backend_container))
                if box is None
            ]
            findings.append(
                _finding(
                    "L4",
                    "backend",
                    child_id,
                    f"PowerPoint binding/bounds missing for {', '.join(missing)}",
                )
            )
        else:
            overflow = _overflow(backend_child, backend_container, padding)
            row["backend"] = {
                "child": backend_child.as_dict(),
                "container": backend_container.as_dict(),
                "overflow_px": {key: round(value, 6) for key, value in overflow.items()},
            }
            if max(overflow.values()) > tolerance:
                findings.append(
                    _finding(
                        "L5",
                        "backend",
                        child_id,
                        f"saved PowerPoint object exceeds container {container_id!r}",
                        metrics=row["backend"],
                    )
                )
        containment_rows.append(row)

    repeat_records: dict[str, list[SvgRecord]] = {}
    for record in records:
        group = record.element.get("data-repeat-group")
        if group:
            repeat_records.setdefault(group, []).append(record)

    repeat_rows: list[dict[str, Any]] = []
    for group, members in sorted(repeat_records.items()):
        axes = {member.element.get("data-repeat-axis") for member in members}
        if len(axes) != 1 or next(iter(axes)) not in {"vertical", "horizontal"}:
            findings.append(
                _finding("L6", "source", group, "repeat group must declare one consistent vertical/horizontal axis")
            )
            continue
        axis = next(iter(axes))
        if len(members) < 2:
            findings.append(_finding("L6", "source", group, "repeat group requires at least two members"))
            continue
        try:
            orders = [int(member.element.get("data-repeat-order", "0")) for member in members]
        except ValueError:
            findings.append(_finding("L6", "source", group, "repeat order must be an integer"))
            continue
        if len(set(orders)) != len(orders):
            findings.append(_finding("L6", "source", group, "repeat order values must be unique"))
            continue
        size_tolerance = max(
            _float_attr(member.element, "data-repeat-size-tolerance", DEFAULT_REPEAT_SIZE_TOLERANCE_PX)
            for member in members
        )
        axis_tolerance = max(
            _float_attr(member.element, "data-repeat-axis-tolerance", DEFAULT_REPEAT_AXIS_TOLERANCE_PX)
            for member in members
        )
        spacing_tolerance = max(
            _float_attr(member.element, "data-repeat-spacing-tolerance", DEFAULT_REPEAT_SPACING_TOLERANCE_PX)
            for member in members
        )
        row: dict[str, Any] = {
            "id": group,
            "axis": axis,
            "tolerances_px": {
                "size": size_tolerance,
                "cross_axis": axis_tolerance,
                "spacing": spacing_tolerance,
            },
        }
        for stage, boxes_by_id in (("source", source_boxes), ("backend", backend_boxes)):
            measured: list[tuple[int, str, Box]] = []
            missing: list[str] = []
            for order, member in zip(orders, members):
                if not member.element_id or member.element_id not in boxes_by_id:
                    missing.append(member.element_id or f"anonymous-{member.tag}")
                else:
                    measured.append((order, member.element_id, boxes_by_id[member.element_id]))
            if missing:
                findings.append(
                    _finding(
                        "L4" if stage == "backend" else "L6",
                        stage,
                        group,
                        f"repeat member geometry missing for {', '.join(missing)}",
                    )
                )
                continue
            metrics = _repeat_metrics(measured, axis)
            row[stage] = metrics
            size_range = max(metrics["width_range_px"], metrics["height_range_px"])
            if size_range > size_tolerance:
                findings.append(
                    _finding("L7", stage, group, "repeated elements do not have equal size", metrics=metrics)
                )
            if metrics["cross_axis_range_px"] > axis_tolerance:
                findings.append(
                    _finding("L8", stage, group, "repeated elements are not on one cross-axis", metrics=metrics)
                )
            if metrics["spacing_range_px"] > spacing_tolerance:
                findings.append(
                    _finding("L9", stage, group, "repeated element spacing is irregular", metrics=metrics)
                )
        repeat_rows.append(row)

    # Missing unrelated shapes are already handled by the global bindings gate;
    # keep this report scoped to annotated layout objects.
    annotated_ids = {record.element_id for record in annotated if record.element_id}
    missing_annotated = sorted(annotated_ids.intersection(missing_backend))
    report = {
        "schema_version": "1.0.0",
        "svg": str(run.redraw_svg),
        "pptx": str(run.pptx_path),
        "defaults_px": {
            "containment": DEFAULT_CONTAINMENT_TOLERANCE_PX,
            "canvas": DEFAULT_CANVAS_TOLERANCE_PX,
            "repeat_size": DEFAULT_REPEAT_SIZE_TOLERANCE_PX,
            "repeat_axis": DEFAULT_REPEAT_AXIS_TOLERANCE_PX,
            "repeat_spacing": DEFAULT_REPEAT_SPACING_TOLERANCE_PX,
        },
        "canvas": {
            "backend": backend_canvas.as_dict(),
            "checked_elements": len(canvas_rows),
            "objects": canvas_rows,
        },
        "containment": containment_rows,
        "repeat_groups": repeat_rows,
        "missing_annotated_backend_objects": missing_annotated,
        "findings": findings,
        "pass": not findings,
    }
    run.qa_dir.mkdir(exist_ok=True)
    run.layout_audit_path.write_text(
        json.dumps(report, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    return report


def strict_blockers(report: dict[str, Any]) -> list[str]:
    return [
        f"layout:{finding['code']}:{finding['stage']}:{finding['target']}"
        for finding in report.get("findings", [])
        if finding.get("severity") == "error"
    ]


def main(argv: list[str] | None = None) -> int:
    import argparse

    parser = argparse.ArgumentParser(prog="autofigure layout", description=__doc__)
    parser.add_argument("run_dir", type=Path)
    args = parser.parse_args(argv)
    report = audit_layout(common.open_run(args.run_dir))
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return 0 if report["pass"] else 2


if __name__ == "__main__":
    raise SystemExit(main())
