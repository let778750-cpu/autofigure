"""Compile a frozen scene declaration into a hash-bound Figure Spec.

The compiler never interprets pixels or invents topology. It only validates and
joins a FROZEN source authority, a PASS perception-review receipt, a blank
canvas, native-math converter receipts, and Designer-authored scene objects.
"""

from __future__ import annotations

import argparse
import copy
import hashlib
import json
import sys
from collections import Counter
from pathlib import Path
from typing import Any, Mapping, Sequence

from jsonschema import Draft202012Validator, FormatChecker
from pptx import Presentation

TOOLS_DIRECTORY = Path(__file__).resolve().parent
if str(TOOLS_DIRECTORY) not in sys.path:
    sys.path.insert(0, str(TOOLS_DIRECTORY))

from finalize_perception_review import atomic_write_json, sha256_file  # noqa: E402
from validate_source_authority import (  # noqa: E402
    SourceAuthorityError,
    validate_authority,
)
from migrate_figure_spec_v3_to_v4 import upgrade_edges, upgrade_elements  # noqa: E402
from render_strategy import RenderStrategyError, validate_render_strategy_contract  # noqa: E402


PROJECT_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_SCENE_SCHEMA = PROJECT_ROOT / "schemas" / "scene-declaration.schema.json"
DEFAULT_FIGURE_SCHEMA = PROJECT_ROOT / "schemas" / "figure-spec.schema.json"
DEFAULT_AUTHORITY_SCHEMA = PROJECT_ROOT / "schemas" / "source-authority.schema.json"
DEFAULT_REVIEW_SCHEMA = PROJECT_ROOT / "schemas" / "perception-review.schema.json"


class FigureSpecCompileError(RuntimeError):
    """Raised when frozen inputs do not close into one deterministic spec."""


def _load_object(path: Path, label: str) -> dict[str, Any]:
    try:
        value = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, UnicodeError, json.JSONDecodeError) as exc:
        raise FigureSpecCompileError(f"cannot read {label}: {path}: {exc}") from exc
    if not isinstance(value, dict):
        raise FigureSpecCompileError(f"{label} must be a JSON object")
    return value


def _validate_json(instance: Any, schema_path: Path, label: str) -> None:
    schema = _load_object(schema_path, f"{label} schema")
    Draft202012Validator.check_schema(schema)
    errors = sorted(
        Draft202012Validator(schema, format_checker=FormatChecker()).iter_errors(instance),
        key=lambda error: list(error.absolute_path),
    )
    if errors:
        first = errors[0]
        location = "$" + "".join(
            f"[{part}]" if isinstance(part, int) else f".{part}"
            for part in first.absolute_path
        )
        raise FigureSpecCompileError(f"{label} rejected at {location}: {first.message}")


def _resolve_path(value: str, *, base: Path = PROJECT_ROOT) -> Path:
    requested = Path(value)
    return requested.resolve(strict=True) if requested.is_absolute() else (base / requested).resolve(strict=True)


def _require_hash(path: Path, expected: str, label: str) -> str:
    actual = sha256_file(path)
    if actual.upper() != expected.upper():
        raise FigureSpecCompileError(
            f"{label} SHA-256 mismatch: expected {expected}, got {actual}"
        )
    return actual


def _authority_evidence(item: Mapping[str, Any]) -> list[str]:
    kinds = [str(evidence["kind"]) for evidence in item["source_evidence"]]
    if not kinds:
        raise FigureSpecCompileError(
            f"{item['authority_item_id']} lacks authoritative source evidence"
        )
    return list(dict.fromkeys(kinds))


def _latex_sequence_terms(expression: str) -> list[str]:
    """Return top-level terms from a parenthesized LaTeX sequence.

    Source authority may deliberately freeze a visual row as one semantic
    formula, while PowerPoint needs one native Office Math object per spatially
    separated label.  This parser supports only that narrow, deterministic
    decomposition; it does not accept arbitrary formula rewrites.
    """

    value = expression.strip()
    if len(value) < 2 or value[0] != "(" or value[-1] != ")":
        raise FigureSpecCompileError(
            "sequence_term_index requires a parenthesized authoritative LaTeX sequence"
        )
    body = value[1:-1]
    terms: list[str] = []
    current: list[str] = []
    brace_depth = 0
    paren_depth = 0
    escaped = False
    for character in body:
        if escaped:
            current.append(character)
            escaped = False
            continue
        if character == "\\":
            current.append(character)
            escaped = True
            continue
        if character == "{":
            brace_depth += 1
        elif character == "}":
            brace_depth -= 1
        elif character == "(":
            paren_depth += 1
        elif character == ")":
            paren_depth -= 1
        if brace_depth < 0 or paren_depth < 0:
            raise FigureSpecCompileError("unbalanced authoritative LaTeX sequence")
        if character == "," and brace_depth == 0 and paren_depth == 0:
            term = "".join(current).strip()
            if not term:
                raise FigureSpecCompileError("authoritative LaTeX sequence contains an empty term")
            terms.append(term)
            current = []
        else:
            current.append(character)
    if brace_depth or paren_depth or escaped:
        raise FigureSpecCompileError("unbalanced authoritative LaTeX sequence")
    final = "".join(current).strip()
    if not final:
        raise FigureSpecCompileError("authoritative LaTeX sequence contains an empty term")
    terms.append(final)
    if len(terms) < 2:
        raise FigureSpecCompileError(
            "sequence_term_index requires an authoritative sequence with at least two terms"
        )
    return terms


def _validate_scene_structure(
    scene: Mapping[str, Any], authority_items: Mapping[str, Mapping[str, Any]]
) -> dict[str, Mapping[str, Any]]:
    elements = scene["elements"]
    identifiers = [str(element.get("id", "")) for element in elements]
    duplicates = sorted(
        identifier for identifier, count in Counter(identifiers).items() if count > 1
    )
    if duplicates or any(not identifier for identifier in identifiers):
        raise FigureSpecCompileError(
            f"scene element IDs must be nonempty and unique: duplicates={duplicates}"
        )
    by_element = {str(element["id"]): element for element in elements}
    for element in elements:
        parent_id = element.get("parent_id")
        if parent_id is not None and parent_id not in by_element:
            raise FigureSpecCompileError(
                f"{element['id']} references unknown parent {parent_id}"
            )
        authority_item_id = element.get("authority_item_id")
        if authority_item_id is None:
            continue
        item = authority_items.get(str(authority_item_id))
        if item is None or item["disposition"] != "CONFIRMED":
            raise FigureSpecCompileError(
                f"{element['id']} references unavailable CONFIRMED authority {authority_item_id}"
            )
        expected_kinds = {
            "text": {"TEXT", "SEMANTIC_REGION", "FORMULA"},
            "formula": {"FORMULA"},
            "manual_asset_slot": {"MANUAL_ASSET"},
            "panel": {"SEMANTIC_REGION"},
        }.get(str(element.get("type")))
        if expected_kinds is not None and item["kind"] not in expected_kinds:
            raise FigureSpecCompileError(
                f"{element['id']} type {element.get('type')} cannot bind {item['kind']} authority"
            )

    edge_ids = [str(edge.get("id", "")) for edge in scene["edges"]]
    if len(edge_ids) != len(set(edge_ids)) or any(not edge_id for edge_id in edge_ids):
        raise FigureSpecCompileError("scene edge IDs must be nonempty and unique")
    for edge in scene["edges"]:
        if edge.get("from") not in by_element or edge.get("to") not in by_element:
            raise FigureSpecCompileError(f"{edge['id']} has an unknown endpoint")
        authority_item_id = edge.get("authority_item_id")
        if authority_item_id is not None:
            item = authority_items.get(str(authority_item_id))
            if item is None or item["kind"] != "RELATION" or item["disposition"] != "CONFIRMED":
                raise FigureSpecCompileError(
                    f"{edge['id']} lacks a CONFIRMED RELATION authority"
                )
    return by_element


def _formula_records(
    scene: Mapping[str, Any],
    authority_items: Mapping[str, Mapping[str, Any]],
    by_element: Mapping[str, Mapping[str, Any]],
) -> list[dict[str, Any]]:
    formula_authority_ids = {
        item_id
        for item_id, item in authority_items.items()
        if item["kind"] == "FORMULA" and item["disposition"] == "CONFIRMED"
    }
    bindings = scene["formula_bindings"]
    bound_ids = [str(binding["authority_item_id"]) for binding in bindings]
    if set(bound_ids) != formula_authority_ids:
        missing = sorted(formula_authority_ids - set(bound_ids))
        extra = sorted(set(bound_ids) - formula_authority_ids)
        raise FigureSpecCompileError(
            f"formula bindings must exhaust CONFIRMED authority formulas: missing={missing}, extra={extra}"
        )
    bindings_by_authority: dict[str, list[Mapping[str, Any]]] = {}
    for binding in bindings:
        bindings_by_authority.setdefault(str(binding["authority_item_id"]), []).append(binding)
    for authority_id, authority_bindings in bindings_by_authority.items():
        indexed = [binding for binding in authority_bindings if "sequence_term_index" in binding]
        if indexed and len(indexed) != len(authority_bindings):
            raise FigureSpecCompileError(
                f"{authority_id} cannot mix whole-formula and sequence-term bindings"
            )
        if indexed:
            terms = _latex_sequence_terms(str(authority_items[authority_id]["canonical_latex"]))
            indexes = [int(binding["sequence_term_index"]) for binding in indexed]
            if sorted(indexes) != list(range(len(terms))):
                raise FigureSpecCompileError(
                    f"{authority_id} sequence bindings must exhaust indexes 0..{len(terms) - 1}"
                )
    formula_ids = [str(binding["formula_id"]) for binding in bindings]
    element_ids = [str(binding["element_id"]) for binding in bindings]
    if len(formula_ids) != len(set(formula_ids)):
        raise FigureSpecCompileError("formula instance IDs must be unique")
    if len(element_ids) != len(set(element_ids)):
        raise FigureSpecCompileError("formula element owners must be unique")

    records: list[dict[str, Any]] = []
    for binding in bindings:
        item = authority_items[str(binding["authority_item_id"])]
        canonical_latex = str(item["canonical_latex"])
        if "sequence_term_index" in binding:
            terms = _latex_sequence_terms(canonical_latex)
            canonical_latex = terms[int(binding["sequence_term_index"])]
        latex_sha256 = hashlib.sha256(canonical_latex.encode("utf-8")).hexdigest()
        element_id = str(binding["element_id"])
        element = by_element.get(element_id)
        if element is None:
            raise FigureSpecCompileError(
                f"{binding['formula_id']} owner {element_id} does not exist"
            )
        owner_type = element.get("type")
        if owner_type == "formula":
            references = [element.get("formula_id")]
        elif owner_type == "text":
            references = [
                run.get("formula_id")
                for run in element.get("content_runs", [])
                if isinstance(run, Mapping) and run.get("kind") == "math"
            ]
        else:
            references = []
        if references != [binding["formula_id"]]:
            raise FigureSpecCompileError(
                f"{element_id} must own exactly one reference to {binding['formula_id']}"
            )
        receipt_path = _resolve_path(str(binding["converter_receipt_path"]))
        receipt_hash = _require_hash(
            receipt_path,
            str(binding["converter_receipt_sha256"]),
            f"{binding['formula_id']} converter receipt",
        )
        receipt = _load_object(receipt_path, "native-math converter receipt")
        expected = {
            "document_type": "NATIVE_OFFICE_MATH_CONVERTER_RECEIPT",
            "status": "PASS",
            "formula_id": binding["formula_id"],
            "canonical_latex": canonical_latex,
            "latex_sha256": latex_sha256,
            "mode": item["formula_mode"],
        }
        for field, value in expected.items():
            actual = receipt.get(field)
            if field == "latex_sha256" and isinstance(actual, str):
                actual = actual.lower()
            if actual != value:
                raise FigureSpecCompileError(
                    f"{binding['formula_id']} converter receipt {field} mismatch"
                )
        native_target = receipt.get("native_target")
        if native_target != {
            "kind": "office_math",
            "wrapper": "a14:m",
            "omml_root": "m:oMath" if item["formula_mode"] == "inline" else "m:oMathPara",
        }:
            raise FigureSpecCompileError(
                f"{binding['formula_id']} converter receipt has the wrong native target"
            )
        source_evidence = _authority_evidence(item)
        if binding.get("perception_candidate_id") is not None:
            source_evidence.append("local_ocr")
        record = {
            "id": binding["formula_id"],
            "element_id": element_id,
            "canonical_latex": canonical_latex,
            "latex_sha256": latex_sha256,
            "mode": item["formula_mode"],
            "render_kind": "native_office_math",
            "fallback_policy": "strict_no_raster_no_svg",
            "converter_receipt_path": str(receipt_path),
            "converter_receipt_sha256": receipt_hash,
            "source_evidence": list(dict.fromkeys(source_evidence)),
            "disposition": "CONFIRMED",
            "authority_item_id": binding["authority_item_id"],
        }
        if "sequence_term_index" in binding:
            record["authority_derivation"] = {
                "kind": "sequence_term",
                "term_index": int(binding["sequence_term_index"]),
                "source_latex_sha256": item["latex_sha256"],
            }
        if binding.get("perception_candidate_id") is not None:
            record["perception_candidate_id"] = binding["perception_candidate_id"]
        records.append(record)
    return records


def _materialize_elements(scene: Mapping[str, Any]) -> list[dict[str, Any]]:
    """Compile parent-local z order into a strict global order.

    Designers specify z_index within each parent scope.  The frozen Figure Spec
    stores both that local value and a depth-prefixed global value so the
    deterministic preflight and backend agree that every child is above its
    container without losing the authored sibling order.
    """

    elements = copy.deepcopy(scene["elements"])
    by_id = {str(element["id"]): element for element in elements}
    depths: dict[str, int] = {}

    def depth(element_id: str, trail: tuple[str, ...] = ()) -> int:
        if element_id in depths:
            return depths[element_id]
        if element_id in trail:
            raise FigureSpecCompileError(
                f"scene parent cycle detected: {' -> '.join((*trail, element_id))}"
            )
        parent_id = by_id[element_id].get("parent_id")
        value = 0 if parent_id is None else depth(str(parent_id), (*trail, element_id)) + 1
        depths[element_id] = value
        return value

    for element in elements:
        local_z = element.get("z_index")
        if not isinstance(local_z, int) or isinstance(local_z, bool) or not 0 <= local_z < 1000:
            raise FigureSpecCompileError(
                f"{element['id']} z_index must be an integer in [0, 999]"
            )
        element["scene_z_index"] = local_z
        element["z_index"] = depth(str(element["id"])) * 1000 + local_z
    return elements


def compile_figure_spec(
    scene_path: Path,
    output_path: Path,
    *,
    scene_schema_path: Path = DEFAULT_SCENE_SCHEMA,
    figure_schema_path: Path = DEFAULT_FIGURE_SCHEMA,
    authority_schema_path: Path = DEFAULT_AUTHORITY_SCHEMA,
    review_schema_path: Path = DEFAULT_REVIEW_SCHEMA,
) -> dict[str, Any]:
    if output_path.exists():
        raise FigureSpecCompileError(f"output already exists: {output_path}")
    resolved_scene = scene_path.resolve(strict=True)
    scene = _load_object(resolved_scene, "scene declaration")
    _validate_json(scene, scene_schema_path.resolve(strict=True), "scene declaration")

    authority_path = _resolve_path(str(scene["source_authority"]["path"]))
    _require_hash(authority_path, str(scene["source_authority"]["sha256"]), "source authority")
    try:
        authority_validation = validate_authority(
            authority_path,
            schema_path=authority_schema_path.resolve(strict=True),
            project_root=PROJECT_ROOT,
        )
    except (SourceAuthorityError, OSError) as exc:
        raise FigureSpecCompileError(str(exc)) from exc
    if authority_validation["authority_status"] != "FROZEN":
        raise FigureSpecCompileError("source authority must be FROZEN")
    authority = _load_object(authority_path, "source authority")
    authority_items = {
        str(item["authority_item_id"]): item for item in authority["items"]
    }

    review_path = _resolve_path(str(scene["perception_review"]["path"]))
    review_hash = _require_hash(
        review_path, str(scene["perception_review"]["sha256"]), "perception review"
    )
    review = _load_object(review_path, "perception review receipt")
    _validate_json(review, review_schema_path.resolve(strict=True), "perception review receipt")
    if review.get("document_type") != "PERCEPTION_REVIEW_RECEIPT" or review.get("status") != "PERCEPTION_REVIEW_PASS":
        raise FigureSpecCompileError("perception review must be PERCEPTION_REVIEW_PASS")
    review_authority = review.get("source_authority")
    if review_authority is None or str(review_authority["sha256"]).upper() != authority_validation["authority_sha256"].upper():
        raise FigureSpecCompileError("perception review is not bound to this source authority")
    if str(review["raw_manifest"]["source_sha256"]).upper() != authority_validation["source_sha256"].upper():
        raise FigureSpecCompileError("perception review source differs from source authority")

    canvas_path = _resolve_path(str(scene["canvas"]["path"]))
    canvas_hash = _require_hash(canvas_path, str(scene["canvas"]["sha256"]), "canvas")
    presentation = Presentation(canvas_path)
    if len(presentation.slides) != 1 or len(presentation.slides[0].shapes) != 0:
        raise FigureSpecCompileError("canvas must contain exactly one blank slide")
    source = authority["source"]
    if (
        int(scene["canvas"]["width_px"]) != int(source["width_px"])
        or int(scene["canvas"]["height_px"]) != int(source["height_px"])
    ):
        raise FigureSpecCompileError("scene canvas pixel dimensions differ from authority")
    if abs(
        presentation.slide_width / presentation.slide_height
        - int(source["width_px"]) / int(source["height_px"])
    ) > 1e-6:
        raise FigureSpecCompileError("PowerPoint canvas aspect ratio differs from source")

    by_element = _validate_scene_structure(scene, authority_items)
    formulas = _formula_records(scene, authority_items, by_element)
    source_path = _resolve_path(str(source["relative_path"]))
    manifest_path = Path(str(review["raw_manifest"]["manifest_path"])).resolve(strict=True)

    spec: dict[str, Any] = {
        "schema_version": "4.0",
        "mode": scene["mode"],
        "source": {
            "path": str(source_path),
            "sha256": source["sha256"],
            "width_px": source["width_px"],
            "height_px": source["height_px"],
            "pixel_format": source["pixel_format"],
            "user_confirmed": True,
        },
        "perception": {
            "manifest_path": str(manifest_path),
            "manifest_sha256": review["raw_manifest"]["manifest_sha256"],
            "review_receipt_path": str(review_path),
            "review_receipt_sha256": review_hash,
        },
        "coordinate_system": {
            "origin": "top-left",
            "unit": "source_pixel",
            "bbox_order": ["x", "y", "w", "h"],
        },
        "canvas": {
            "width_px": scene["canvas"]["width_px"],
            "height_px": scene["canvas"]["height_px"],
            "background": scene["canvas"]["background"],
            "background_evidence": scene["canvas"]["background_evidence"],
            "pptx_path": str(canvas_path),
            "pptx_sha256": canvas_hash,
            "slide_width_emu": presentation.slide_width,
            "slide_height_emu": presentation.slide_height,
        },
        "measurement_dpi": scene["measurement_dpi"],
        "policy_profile": scene.get("policy_profile", "standard"),
        "elements": upgrade_elements(_materialize_elements(scene)),
        "edges": upgrade_edges(scene["edges"]),
        "formulas": formulas,
        "uncertainties": scene["uncertainties"],
        "authority": {
            "path": str(authority_path),
            "sha256": authority_validation["authority_sha256"],
            "authority_id": authority_validation["authority_id"],
        },
        "scene_declaration": {
            "path": str(resolved_scene),
            "sha256": sha256_file(resolved_scene),
            "schema_path": str(scene_schema_path.resolve(strict=True)),
            "schema_sha256": sha256_file(scene_schema_path.resolve(strict=True)),
        },
    }
    if scene.get("geometry") is not None:
        geometry_path = _resolve_path(str(scene["geometry"]["path"]))
        spec["geometry"] = {
            "manifest_path": str(geometry_path),
            "manifest_sha256": _require_hash(
                geometry_path, str(scene["geometry"]["sha256"]), "geometry manifest"
            ),
        }
    if scene["canvas"].get("background_reason") is not None:
        spec["canvas"]["background_reason"] = scene["canvas"]["background_reason"]
    try:
        validate_render_strategy_contract(spec["elements"], spec["edges"])
    except RenderStrategyError as exc:
        raise FigureSpecCompileError(str(exc)) from exc
    _validate_json(spec, figure_schema_path.resolve(strict=True), "compiled Figure Spec")
    atomic_write_json(output_path, spec)
    return spec


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--scene", required=True, type=Path)
    parser.add_argument("--output", required=True, type=Path)
    parser.add_argument("--scene-schema", type=Path, default=DEFAULT_SCENE_SCHEMA)
    parser.add_argument("--figure-schema", type=Path, default=DEFAULT_FIGURE_SCHEMA)
    parser.add_argument("--authority-schema", type=Path, default=DEFAULT_AUTHORITY_SCHEMA)
    parser.add_argument("--review-schema", type=Path, default=DEFAULT_REVIEW_SCHEMA)
    return parser


def main(argv: Sequence[str] | None = None) -> int:
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8")
    args = build_parser().parse_args(argv)
    try:
        spec = compile_figure_spec(
            args.scene,
            args.output,
            scene_schema_path=args.scene_schema,
            figure_schema_path=args.figure_schema,
            authority_schema_path=args.authority_schema,
            review_schema_path=args.review_schema,
        )
    except (FigureSpecCompileError, OSError) as exc:
        print(f"FIGURE_SPEC_COMPILE_REJECTED: {exc}", file=sys.stderr)
        return 2
    print(
        json.dumps(
            {
                "status": "SPEC_FROZEN",
                "elements": len(spec["elements"]),
                "edges": len(spec["edges"]),
                "formulas": len(spec["formulas"]),
                "output": str(args.output.resolve()),
            },
            ensure_ascii=False,
        )
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
