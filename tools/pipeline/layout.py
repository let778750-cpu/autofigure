"""Explicit layout contracts for SVG sources and saved PowerPoint objects.

The visual model is responsible for measuring a reference, but it must declare
the relationships that deterministic tooling can verify:

``data-layout-container``
    ID of the shape that must contain this element.  ``data-layout-padding``
    and ``data-layout-tolerance`` are optional pixel values.

``data-repeat-group`` / ``data-repeat-axis`` / ``data-repeat-order``
    Declare a repeated horizontal, vertical, or two-dimensional translated
    motif.  Equal size, the declared axis and regular centre-to-centre spacing
    are audited in both SVG source coordinates and the saved/reopened
    PowerPoint artifact.  An annotated SVG ``g`` is measured as the union of
    its bound drawable descendants; group annotations are never discarded.

``data-peer-group`` / ``data-peer-axis`` / ``data-peer-order``
    Declare ordered layout peers.  Horizontal/vertical peers must stay on the
    declared cross-axis with regular centre spacing.  ``xy`` peers are treated
    as translated semantic equivalents and must retain equal dimensions and a
    regular translation vector.  Incomplete or contradictory declarations are
    hard findings, rather than silently ignored metadata.

``data-peer-size-group``
    Declare semantic peers whose width and height must remain equal within an
    explicit pixel tolerance, even when they are not regularly spaced.

``data-gap-source-id`` / ``data-gap-target-id``
    Bind a straight arrow to the open gap between two shapes.  Insets derive
    its along-axis endpoints from the peer boundaries, so moving either peer
    cannot leave a stale fixed-length arrow behind.

``data-z-above`` / ``data-z-below``
    Declare paint-order constraints in both SVG and the saved PowerPoint.

``data-text-flow="stacked-characters"``
    Require upright one-character-per-paragraph text, rather than a rotated
    word that only happens to occupy a vertical box.

These annotations deliberately use a separate namespace from
``data-owner-id``.  Existing owner IDs describe semantic label/edge ownership
and do not imply geometric containment.

Saved PowerPoint geometry also has one global, annotation-free contract: every
scene-bound object, including OMML hidden in ``mc:AlternateContent``, must stay
inside the slide canvas.  This package-level check avoids false results from
tools that can move only the python-pptx-visible fallback object.
"""

from __future__ import annotations

import json
import math
import re
import unicodedata
import xml.etree.ElementTree as ET
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation

from tools.core import common
from tools.core.contracts import read_json
from tools.core.svggeom import Matrix, parse_path_d, parse_transform

SVG_NS = "{http://www.w3.org/2000/svg}"

DEFAULT_CONTAINMENT_TOLERANCE_PX = 0.25
DEFAULT_CANVAS_TOLERANCE_PX = 0.25
DEFAULT_REPEAT_SIZE_TOLERANCE_PX = 0.25
DEFAULT_REPEAT_AXIS_TOLERANCE_PX = 0.25
# Integer-pixel reference reconstructions commonly distribute an odd remainder
# as N and N+1 pixel steps.  A 1 px range is regular; larger drift is not.
DEFAULT_REPEAT_SPACING_TOLERANCE_PX = 1.0
DEFAULT_PEER_SIZE_TOLERANCE_PX = 0.25
MAX_PEER_SIZE_TOLERANCE_PX = 1.0
DEFAULT_PEER_AXIS_TOLERANCE_PX = 0.25
MAX_PEER_AXIS_TOLERANCE_PX = 1.0
DEFAULT_PEER_SPACING_TOLERANCE_PX = 1.0
MAX_PEER_SPACING_TOLERANCE_PX = 1.0
DEFAULT_GAP_TOLERANCE_PX = 0.75
MAX_GAP_TOLERANCE_PX = 1.0
PT_PER_PX = 0.75
TEXT_STACK_SPACING_TOLERANCE_PT = 0.01
MAX_TEXT_FRAME_OVERFLOW_TOLERANCE_PX = 3.0
# OOXML stores geometry as integer EMUs.  A round-trip can therefore introduce
# sub-thousandth-pixel noise even when the declared frame is exact.
OOXML_GEOMETRY_EPSILON_PX = 0.001

_DRAWABLE = {
    "rect",
    "circle",
    "ellipse",
    "line",
    "polyline",
    "polygon",
    "path",
    "text",
    "image",
}
_SKIP = {"defs", "marker", "linearGradient", "radialGradient"}


@dataclass(frozen=True)
class Box:
    x: float
    y: float
    width: float
    height: float

    @property
    def right(self) -> float:
        return self.x + self.width

    @property
    def bottom(self) -> float:
        return self.y + self.height

    @property
    def center_x(self) -> float:
        return self.x + self.width / 2.0

    @property
    def center_y(self) -> float:
        return self.y + self.height / 2.0

    def as_dict(self) -> dict[str, float]:
        return {
            "x": round(self.x, 6),
            "y": round(self.y, 6),
            "width": round(self.width, 6),
            "height": round(self.height, 6),
        }


@dataclass(frozen=True)
class SvgRecord:
    element_id: str | None
    tag: str
    element: ET.Element
    matrix: Matrix
    box: Box | None
    descendant_ids: tuple[str, ...] = ()
    logical_group: bool = False


def _local_name(element: ET.Element) -> str:
    return element.tag.rsplit("}", 1)[-1]


def _number(element: ET.Element, name: str, default: float = 0.0) -> float:
    raw = element.get(name)
    if raw is None:
        return default
    try:
        return float(raw)
    except ValueError as exc:
        raise ValueError(f"{element.get('id') or _local_name(element)}: invalid {name}={raw!r}") from exc


def _points_box(points: Iterable[tuple[float, float]], matrix: Matrix) -> Box | None:
    transformed = [matrix.apply(x, y) for x, y in points]
    if not transformed:
        return None
    xs = [point[0] for point in transformed]
    ys = [point[1] for point in transformed]
    return Box(min(xs), min(ys), max(xs) - min(xs), max(ys) - min(ys))


def _element_box(element: ET.Element, matrix: Matrix) -> Box | None:
    tag = _local_name(element)
    if tag in {"rect", "image"}:
        x, y = _number(element, "x"), _number(element, "y")
        width, height = _number(element, "width"), _number(element, "height")
        return _points_box(
            ((x, y), (x + width, y), (x, y + height), (x + width, y + height)),
            matrix,
        )
    if tag == "circle":
        cx, cy, radius = _number(element, "cx"), _number(element, "cy"), _number(element, "r")
        return _points_box(
            (
                (cx - radius, cy - radius),
                (cx + radius, cy - radius),
                (cx - radius, cy + radius),
                (cx + radius, cy + radius),
            ),
            matrix,
        )
    if tag == "ellipse":
        cx, cy = _number(element, "cx"), _number(element, "cy")
        rx, ry = _number(element, "rx"), _number(element, "ry")
        return _points_box(
            ((cx - rx, cy - ry), (cx + rx, cy - ry), (cx - rx, cy + ry), (cx + rx, cy + ry)),
            matrix,
        )
    if tag == "line":
        return _points_box(
            ((_number(element, "x1"), _number(element, "y1")), (_number(element, "x2"), _number(element, "y2"))),
            matrix,
        )
    if tag in {"polyline", "polygon"}:
        values = [float(item) for item in re.findall(r"[-+]?(?:\d*\.\d+|\d+\.?)", element.get("points", ""))]
        return _points_box(zip(values[0::2], values[1::2]), matrix)
    if tag == "path":
        points: list[tuple[float, float]] = []
        for segment in parse_path_d(element.get("d", "")):
            if segment[0] in {"M", "L"}:
                points.append((segment[1], segment[2]))
            elif segment[0] == "C":
                points.extend(
                    ((segment[1], segment[2]), (segment[3], segment[4]), (segment[5], segment[6]))
                )
        return _points_box(points, matrix)
    # Text bounds depend on the actual PowerPoint font/OMML renderer.  They are
    # intentionally judged from the saved artifact, not guessed from SVG text.
    return None


def collect_svg_records(root: ET.Element) -> list[SvgRecord]:
    """Return drawable records with the same inherited transform semantics as convert."""
    records: list[SvgRecord] = []

    def walk(element: ET.Element, parent_matrix: Matrix) -> None:
        tag = _local_name(element)
        if tag in _SKIP:
            return
        matrix = parent_matrix.multiply(parse_transform(element.get("transform")))
        if tag in {"svg", "g"}:
            for child in element:
                walk(child, matrix)
            return
        if tag not in _DRAWABLE:
            return
        records.append(
            SvgRecord(
                element_id=element.get("id"),
                tag=tag,
                element=element,
                matrix=matrix,
                box=_element_box(element, matrix),
            )
        )

    walk(root, Matrix())
    return records


def collect_svg_boxes(root: ET.Element) -> dict[str, Box]:
    """Boxes for explicitly identified SVG elements, used by the converter."""
    return {
        record.element_id: record.box
        for record in collect_svg_records(root)
        if record.element_id and record.box is not None
    }


_LAYOUT_ANNOTATION_NAMES = (
    "data-layout-container",
    "data-layout-padding",
    "data-layout-tolerance",
    "data-repeat-group",
    "data-repeat-axis",
    "data-repeat-order",
    "data-peer-size-group",
    "data-peer-group",
    "data-peer-axis",
    "data-peer-order",
    "data-gap-source-id",
    "data-gap-target-id",
    "data-z-above",
    "data-z-below",
    "data-text-flow",
    "data-text-container",
)


def _collect_layout_records(root: ET.Element) -> list[SvgRecord]:
    """Include annotation-bearing SVG groups as measurable logical objects.

    The converter intentionally flattens SVG groups into native PowerPoint
    objects.  Layout contracts nevertheless belong to the logical group, so
    source and backend geometry are both reconstructed from the same stable
    descendant IDs.
    """

    drawable_records = collect_svg_records(root)
    records = list(drawable_records)
    for element in root.iter():
        if _local_name(element) != "g" or not any(
            element.get(name) is not None for name in _LAYOUT_ANNOTATION_NAMES
        ):
            continue
        descendants = set(element.iter())
        members = [
            record
            for record in drawable_records
            if record.element is not element and record.element in descendants
        ]
        boxes = [record.box for record in members if record.box is not None]
        records.append(
            SvgRecord(
                element_id=element.get("id"),
                tag="g",
                element=element,
                matrix=Matrix(),
                box=_union(boxes) if boxes else None,
                descendant_ids=tuple(
                    record.element_id for record in members if record.element_id
                ),
                logical_group=True,
            )
        )
    return records


def _add_logical_backend_boxes(
    records: list[SvgRecord],
    backend_boxes: dict[str, Box],
) -> list[str]:
    """Resolve flattened SVG group bounds from all bound descendants."""

    missing: list[str] = []
    for record in records:
        if not record.logical_group or not record.element_id:
            continue
        member_ids = list(record.descendant_ids)
        member_boxes = [backend_boxes[item] for item in member_ids if item in backend_boxes]
        if not member_ids or len(member_boxes) != len(member_ids):
            missing.append(record.element_id)
            continue
        backend_boxes[record.element_id] = _union(member_boxes)
    return missing


def _float_attr(element: ET.Element, name: str, default: float) -> float:
    value = element.get(name)
    if value is None:
        return default
    try:
        parsed = float(value)
    except ValueError as exc:
        raise ValueError(f"{element.get('id') or _local_name(element)}: invalid {name}={value!r}") from exc
    if not math.isfinite(parsed) or parsed < 0:
        raise ValueError(f"{element.get('id') or _local_name(element)}: {name} must be finite and >= 0")
    return parsed


def _overflow(child: Box, container: Box, padding: float) -> dict[str, float]:
    inner_left = container.x + padding
    inner_top = container.y + padding
    inner_right = container.right - padding
    inner_bottom = container.bottom - padding
    return {
        "left": max(0.0, inner_left - child.x),
        "top": max(0.0, inner_top - child.y),
        "right": max(0.0, child.right - inner_right),
        "bottom": max(0.0, child.bottom - inner_bottom),
    }


def _union(boxes: list[Box]) -> Box:
    left = min(box.x for box in boxes)
    top = min(box.y for box in boxes)
    right = max(box.right for box in boxes)
    bottom = max(box.bottom for box in boxes)
    return Box(left, top, right - left, bottom - top)


def _backend_boxes(
    run: common.Run,
) -> tuple[dict[str, Box], list[str], dict[str, list[dict[str, Any]]]]:
    """Resolve scene IDs to saved PowerPoint bounds and object-level details."""
    from tools.repair.live_bridge import _shape_bounds, _walk_shapes, _xml_shape_bound_indexes

    presentation = Presentation(run.pptx_path)
    normal_by_identity: dict[tuple[int, str], Box] = {}
    detail_by_identity: dict[tuple[int, str], dict[str, Any]] = {}
    for z_order, shape in enumerate(_walk_shapes(presentation.slides[0].shapes), start=1):
        raw = _shape_bounds(shape)
        box = Box(raw["x"], raw["y"], raw["width"], raw["height"])
        identity = (int(shape.shape_id), shape.name)
        normal_by_identity[identity] = box
        text = None
        paragraphs = None
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text
            paragraphs = [paragraph.text for paragraph in shape.text_frame.paragraphs]
        line_spacings_pt = None
        if getattr(shape, "has_text_frame", False):
            line_spacings_pt = []
            for paragraph in shape.text_frame.paragraphs:
                spacing = paragraph.line_spacing
                points = getattr(spacing, "pt", None)
                line_spacings_pt.append(
                    None if points is None else round(float(points), 6)
                )
        detail = {
            "shape_name": shape.name,
            "shape_id": int(shape.shape_id),
            "z_order": z_order,
            "rotation_deg": round(float(getattr(shape, "rotation", 0.0) or 0.0), 6),
            "text": text,
            "paragraphs": paragraphs,
            "line_spacings_pt": line_spacings_pt,
        }
        detail_by_identity[identity] = detail
    _, _, xml_by_identity_raw = _xml_shape_bound_indexes(run.pptx_path)
    xml_by_identity = {
        identity: Box(value["x"], value["y"], value["width"], value["height"])
        for identity, value in xml_by_identity_raw.items()
    }

    resolved: dict[str, list[Box]] = {}
    resolved_details: dict[str, list[dict[str, Any]]] = {}
    missing: list[str] = []
    bindings = read_json(run.bindings_path)
    for binding in bindings.get("bindings", []):
        element_id = binding.get("element_id")
        if not element_id:
            continue
        shape_id = binding.get("shape_id")
        shape_name = binding.get("shape_name")
        if (
            isinstance(shape_id, bool)
            or not isinstance(shape_id, int)
            or shape_id <= 0
            or not isinstance(shape_name, str)
            or not shape_name
        ):
            missing.append(element_id)
            continue
        identity = (shape_id, shape_name)
        box = normal_by_identity.get(identity) or xml_by_identity.get(identity)
        if box is None:
            missing.append(element_id)
        else:
            resolved.setdefault(element_id, []).append(box)
            detail = detail_by_identity.get(identity)
            if detail is not None:
                resolved_details.setdefault(element_id, []).append(detail)
    for binding in bindings.get("logical_group_bindings", []):
        element_id = binding.get("element_id")
        backend_ids = binding.get("backend_object_ids")
        backend_names = binding.get("backend_object_names")
        if (
            not isinstance(element_id, str)
            or not element_id
            or not isinstance(backend_ids, list)
            or not isinstance(backend_names, list)
            or not backend_ids
            or len(backend_ids) != len(backend_names)
        ):
            if isinstance(element_id, str) and element_id:
                missing.append(element_id)
            continue
        group_boxes: list[Box] = []
        group_details: list[dict[str, Any]] = []
        for shape_id, shape_name in zip(backend_ids, backend_names, strict=True):
            if (
                not isinstance(shape_id, int)
                or isinstance(shape_id, bool)
                or shape_id <= 0
                or not isinstance(shape_name, str)
                or not shape_name
            ):
                continue
            identity = (shape_id, shape_name)
            box = normal_by_identity.get(identity) or xml_by_identity.get(identity)
            if box is not None:
                group_boxes.append(box)
            detail = detail_by_identity.get(identity)
            if detail is not None:
                group_details.append(detail)
        if len(group_boxes) != len(backend_ids):
            missing.append(element_id)
            continue
        resolved[element_id] = group_boxes
        if group_details:
            resolved_details[element_id] = group_details
    return (
        {element_id: _union(boxes) for element_id, boxes in resolved.items()},
        sorted(set(missing)),
        resolved_details,
    )


def _backend_canvas(run: common.Run) -> Box:
    """Return the saved PowerPoint slide canvas in the project's 96-DPI pixels."""
    from tools.repair.live_bridge import EMU_PER_PX

    presentation = Presentation(run.pptx_path)
    return Box(
        0.0,
        0.0,
        float(presentation.slide_width) / EMU_PER_PX,
        float(presentation.slide_height) / EMU_PER_PX,
    )


def _finding(
    code: str,
    stage: str,
    target: str,
    message: str,
    *,
    metrics: dict[str, Any] | None = None,
) -> dict[str, Any]:
    result: dict[str, Any] = {
        "code": code,
        "stage": stage,
        "severity": "error",
        "target": target,
        "message": message,
    }
    if metrics:
        result["metrics"] = metrics
    return result


def _repeat_metrics(boxes: list[tuple[int, str, Box]], axis: str) -> dict[str, Any]:
    ordered = sorted(boxes, key=lambda item: (item[0], item[1]))
    widths = [item[2].width for item in ordered]
    heights = [item[2].height for item in ordered]
    centers_xy = [(item[2].center_x, item[2].center_y) for item in ordered]
    if axis == "xy":
        steps_xy = [
            (
                centers_xy[index + 1][0] - centers_xy[index][0],
                centers_xy[index + 1][1] - centers_xy[index][1],
            )
            for index in range(len(centers_xy) - 1)
        ]
        x_steps = [item[0] for item in steps_xy]
        y_steps = [item[1] for item in steps_xy]
        return {
            "members": [item[1] for item in ordered],
            "widths_px": [round(value, 6) for value in widths],
            "heights_px": [round(value, 6) for value in heights],
            "centers_xy_px": [
                [round(x, 6), round(y, 6)] for x, y in centers_xy
            ],
            "steps_xy_px": [
                [round(x, 6), round(y, 6)] for x, y in steps_xy
            ],
            "width_range_px": round(max(widths) - min(widths), 6),
            "height_range_px": round(max(heights) - min(heights), 6),
            "cross_axis_range_px": 0.0,
            "spacing_range_px": round(
                max(
                    (max(x_steps) - min(x_steps)) if x_steps else 0.0,
                    (max(y_steps) - min(y_steps)) if y_steps else 0.0,
                ),
                6,
            ),
        }
    centers = [
        item[2].center_y if axis == "vertical" else item[2].center_x
        for item in ordered
    ]
    cross_centers = [
        item[2].center_x if axis == "vertical" else item[2].center_y
        for item in ordered
    ]
    steps = [centers[index + 1] - centers[index] for index in range(len(centers) - 1)]
    return {
        "members": [item[1] for item in ordered],
        "widths_px": [round(value, 6) for value in widths],
        "heights_px": [round(value, 6) for value in heights],
        "centers_px": [round(value, 6) for value in centers],
        "cross_axis_centers_px": [round(value, 6) for value in cross_centers],
        "steps_px": [round(value, 6) for value in steps],
        "width_range_px": round(max(widths) - min(widths), 6),
        "height_range_px": round(max(heights) - min(heights), 6),
        "cross_axis_range_px": round(max(cross_centers) - min(cross_centers), 6),
        "spacing_range_px": round(max(steps) - min(steps), 6) if steps else 0.0,
    }


def _normalized_axis(value: str | None) -> str | None:
    return {
        "x": "horizontal",
        "horizontal": "horizontal",
        "y": "vertical",
        "vertical": "vertical",
        "xy": "xy",
    }.get((value or "").strip().lower())


def _ordered_contract(
    members: list[SvgRecord],
    *,
    group: str,
    code: str,
) -> tuple[list[int] | None, list[dict[str, Any]]]:
    findings: list[dict[str, Any]] = []
    raw_orders = [member.element.get("data-repeat-order" if code == "L6" else "data-peer-order") for member in members]
    if any(value is None for value in raw_orders):
        findings.append(_finding(code, "source", group, "every group member requires an explicit order"))
        return None, findings
    try:
        orders = [int(value or "") for value in raw_orders]
    except ValueError:
        findings.append(_finding(code, "source", group, "group order must be an integer"))
        return None, findings
    if len(set(orders)) != len(orders) or sorted(orders) != list(range(1, len(orders) + 1)):
        findings.append(
            _finding(
                code,
                "source",
                group,
                "group order values must be unique and contiguous from 1",
                metrics={"orders": orders},
            )
        )
        return None, findings
    return orders, findings


def _text_content(element: ET.Element) -> str:
    return "".join(element.itertext())


def _positioned_text_lines(element: ET.Element) -> list[str]:
    """Return semantic lines for explicitly positioned SVG tspans."""

    lines = [
        (child.text or "").strip()
        for child in element
        if _local_name(child) == "tspan"
        and any(child.get(name) is not None for name in ("x", "y", "dy"))
        and (child.text or "").strip()
    ]
    if lines:
        return lines
    content = _text_content(element).strip()
    return [content] if content else []


def _graphemes(text: str) -> list[str]:
    result: list[str] = []
    for character in text:
        if character.isspace():
            continue
        if unicodedata.combining(character) and result:
            result[-1] += character
        else:
            result.append(character)
    return result


def _record_endpoints(record: SvgRecord) -> tuple[tuple[float, float], tuple[float, float]] | None:
    element = record.element
    centerline = element.get("data-arrow-centerline")
    if centerline:
        segments = parse_path_d(centerline)
        points: list[tuple[float, float]] = []
        for segment in segments:
            if segment[0] in {"M", "L"}:
                points.append((segment[1], segment[2]))
            elif segment[0] == "C":
                points.append((segment[5], segment[6]))
            else:
                return None
        if len(points) != 2:
            return None
        return record.matrix.apply(*points[0]), record.matrix.apply(*points[1])
    if record.tag == "line":
        return (
            record.matrix.apply(_number(element, "x1"), _number(element, "y1")),
            record.matrix.apply(_number(element, "x2"), _number(element, "y2")),
        )
    if record.tag != "path":
        return None
    segments = parse_path_d(element.get("d", ""))
    points: list[tuple[float, float]] = []
    for segment in segments:
        if segment[0] in {"M", "L"}:
            points.append((segment[1], segment[2]))
        elif segment[0] == "C":
            points.append((segment[5], segment[6]))
        else:
            return None
    if len(points) != 2:
        return None
    return record.matrix.apply(*points[0]), record.matrix.apply(*points[1])


def _id_list(value: str | None) -> list[str]:
    return [item for item in re.split(r"[\s,]+", value or "") if item]


def _point_delta(
    actual: tuple[float, float],
    expected: tuple[float, float],
) -> dict[str, float]:
    return {
        "x": round(actual[0] - expected[0], 6),
        "y": round(actual[1] - expected[1], 6),
        "distance": round(math.dist(actual, expected), 6),
    }


def _gap_expected_endpoints(
    axis: str,
    source_box: Box,
    target_box: Box,
    start_inset: float,
    end_inset: float,
    cross_value: str | None,
) -> tuple[tuple[float, float], tuple[float, float]]:
    def cross_position(overlap_start: float, overlap_end: float) -> float:
        if cross_value is None:
            return (overlap_start + overlap_end) / 2.0
        cross = float(cross_value)
        if not math.isfinite(cross):
            raise ValueError("gap cross-position must be finite")
        if cross < overlap_start or cross > overlap_end:
            raise ValueError("gap cross-position must lie within the peer cross-axis overlap")
        return cross

    if axis == "horizontal":
        if source_box.right >= target_box.x:
            raise ValueError("horizontal gap source must be left of target")
        overlap_start = max(source_box.y, target_box.y)
        overlap_end = min(source_box.bottom, target_box.bottom)
        if overlap_start >= overlap_end:
            raise ValueError("horizontal gap peers must overlap on the vertical axis")
        cross = cross_position(overlap_start, overlap_end)
        return (
            (source_box.right + start_inset, cross),
            (target_box.x - end_inset, cross),
        )
    if axis == "vertical":
        if source_box.bottom >= target_box.y:
            raise ValueError("vertical gap source must be above target")
        overlap_start = max(source_box.x, target_box.x)
        overlap_end = min(source_box.right, target_box.right)
        if overlap_start >= overlap_end:
            raise ValueError("vertical gap peers must overlap on the horizontal axis")
        cross = cross_position(overlap_start, overlap_end)
        return (
            (cross, source_box.bottom + start_inset),
            (cross, target_box.y - end_inset),
        )
    raise ValueError("gap axis must be horizontal or vertical")


def audit_layout(run: common.Run) -> dict[str, Any]:
    """Audit explicit layout annotations in SVG and saved PowerPoint geometry."""
    root = ET.parse(run.redraw_svg).getroot()
    records = _collect_layout_records(root)
    source_boxes = {
        record.element_id: record.box
        for record in records
        if record.element_id and record.box is not None
    }
    backend_boxes, missing_backend, backend_details = _backend_boxes(run)
    missing_backend.extend(_add_logical_backend_boxes(records, backend_boxes))
    backend_canvas = _backend_canvas(run)
    findings: list[dict[str, Any]] = []

    canvas_rows: list[dict[str, Any]] = []
    for element_id, box in sorted(backend_boxes.items()):
        overflow = _overflow(box, backend_canvas, 0.0)
        row = {
            "element": element_id,
            "box": box.as_dict(),
            "overflow_px": {key: round(value, 6) for key, value in overflow.items()},
        }
        canvas_rows.append(row)
        if max(overflow.values()) > DEFAULT_CANVAS_TOLERANCE_PX:
            findings.append(
                _finding(
                    "L10",
                    "backend",
                    element_id,
                    "saved PowerPoint object exceeds the slide canvas",
                    metrics=row,
                )
            )

    annotated = [
        record
        for record in records
        if any(record.element.get(name) is not None for name in _LAYOUT_ANNOTATION_NAMES)
    ]
    for record in annotated:
        if not record.element_id:
            findings.append(
                _finding(
                    "L1",
                    "source",
                    f"anonymous-{record.tag}",
                    "layout-annotated elements require an explicit stable id",
                )
            )
        target = record.element_id or f"anonymous-{record.tag}"
        if record.element.get("data-layout-container") is None and (
            record.element.get("data-layout-padding") is not None
            or record.element.get("data-layout-tolerance") is not None
        ):
            findings.append(
                _finding(
                    "L2",
                    "source",
                    target,
                    "layout padding/tolerance requires data-layout-container",
                )
            )
        if record.element.get("data-repeat-group") is None and (
            record.element.get("data-repeat-axis") is not None
            or record.element.get("data-repeat-order") is not None
        ):
            findings.append(
                _finding(
                    "L6",
                    "source",
                    target,
                    "repeat axis/order requires data-repeat-group",
                )
            )

    containment_rows: list[dict[str, Any]] = []
    for record in annotated:
        child_id = record.element_id
        container_id = record.element.get("data-layout-container")
        if not child_id or not container_id:
            continue
        padding = _float_attr(record.element, "data-layout-padding", 0.0)
        tolerance = _float_attr(
            record.element,
            "data-layout-tolerance",
            DEFAULT_CONTAINMENT_TOLERANCE_PX,
        )
        row: dict[str, Any] = {
            "element": child_id,
            "container": container_id,
            "padding_px": padding,
            "tolerance_px": tolerance,
        }
        if child_id == container_id:
            findings.append(
                _finding(
                    "L2",
                    "source",
                    child_id,
                    "layout container must not reference the element itself",
                )
            )
            containment_rows.append(row)
            continue
        source_child = source_boxes.get(child_id)
        source_container = source_boxes.get(container_id)
        if source_container is None:
            findings.append(
                _finding("L2", "source", child_id, f"layout container {container_id!r} has no measurable SVG geometry")
            )
        elif padding * 2 >= min(source_container.width, source_container.height):
            findings.append(
                _finding(
                    "L2",
                    "source",
                    child_id,
                    "layout padding collapses the declared container interior",
                    metrics={
                        "padding_px": padding,
                        "container": source_container.as_dict(),
                    },
                )
            )
        elif source_child is not None:
            overflow = _overflow(source_child, source_container, padding)
            row["source"] = {
                "child": source_child.as_dict(),
                "container": source_container.as_dict(),
                "overflow_px": {key: round(value, 6) for key, value in overflow.items()},
            }
            if max(overflow.values()) > tolerance:
                findings.append(
                    _finding(
                        "L3",
                        "source",
                        child_id,
                        f"element exceeds SVG container {container_id!r}",
                        metrics=row["source"],
                    )
                )

        backend_child = backend_boxes.get(child_id)
        backend_container = backend_boxes.get(container_id)
        if backend_child is None or backend_container is None:
            missing = [
                element_id
                for element_id, box in ((child_id, backend_child), (container_id, backend_container))
                if box is None
            ]
            findings.append(
                _finding(
                    "L4",
                    "backend",
                    child_id,
                    f"PowerPoint binding/bounds missing for {', '.join(missing)}",
                )
            )
        else:
            overflow = _overflow(backend_child, backend_container, padding)
            row["backend"] = {
                "child": backend_child.as_dict(),
                "container": backend_container.as_dict(),
                "overflow_px": {key: round(value, 6) for key, value in overflow.items()},
            }
            if max(overflow.values()) > tolerance:
                findings.append(
                    _finding(
                        "L5",
                        "backend",
                        child_id,
                        f"saved PowerPoint object exceeds container {container_id!r}",
                        metrics=row["backend"],
                    )
                )
        containment_rows.append(row)

    repeat_records: dict[str, list[SvgRecord]] = {}
    for record in records:
        group = record.element.get("data-repeat-group")
        if group:
            repeat_records.setdefault(group, []).append(record)

    repeat_rows: list[dict[str, Any]] = []
    for group, members in sorted(repeat_records.items()):
        raw_axes = {member.element.get("data-repeat-axis") for member in members}
        axes = {_normalized_axis(value) for value in raw_axes}
        if len(axes) != 1 or None in axes:
            findings.append(
                _finding(
                    "L6",
                    "source",
                    group,
                    "repeat group must declare one consistent x/y/xy axis",
                    metrics={"declared_axes": sorted(str(value) for value in raw_axes)},
                )
            )
            continue
        axis = next(iter(axes))
        if len(members) < 2:
            findings.append(_finding("L6", "source", group, "repeat group requires at least two members"))
            continue
        orders, order_findings = _ordered_contract(members, group=group, code="L6")
        findings.extend(order_findings)
        if orders is None:
            continue
        size_tolerance = max(
            _float_attr(member.element, "data-repeat-size-tolerance", DEFAULT_REPEAT_SIZE_TOLERANCE_PX)
            for member in members
        )
        axis_tolerance = max(
            _float_attr(member.element, "data-repeat-axis-tolerance", DEFAULT_REPEAT_AXIS_TOLERANCE_PX)
            for member in members
        )
        spacing_tolerance = max(
            _float_attr(member.element, "data-repeat-spacing-tolerance", DEFAULT_REPEAT_SPACING_TOLERANCE_PX)
            for member in members
        )
        row: dict[str, Any] = {
            "id": group,
            "axis": axis,
            "tolerances_px": {
                "size": size_tolerance,
                "cross_axis": axis_tolerance,
                "spacing": spacing_tolerance,
            },
        }
        for stage, boxes_by_id in (("source", source_boxes), ("backend", backend_boxes)):
            measured: list[tuple[int, str, Box]] = []
            missing: list[str] = []
            for order, member in zip(orders, members):
                if not member.element_id or member.element_id not in boxes_by_id:
                    missing.append(member.element_id or f"anonymous-{member.tag}")
                else:
                    measured.append((order, member.element_id, boxes_by_id[member.element_id]))
            if missing:
                findings.append(
                    _finding(
                        "L4" if stage == "backend" else "L6",
                        stage,
                        group,
                        f"repeat member geometry missing for {', '.join(missing)}",
                    )
                )
                continue
            metrics = _repeat_metrics(measured, axis)
            row[stage] = metrics
            size_range = max(metrics["width_range_px"], metrics["height_range_px"])
            if size_range > size_tolerance:
                findings.append(
                    _finding("L7", stage, group, "repeated elements do not have equal size", metrics=metrics)
                )
            if metrics["cross_axis_range_px"] > axis_tolerance:
                findings.append(
                    _finding("L8", stage, group, "repeated elements are not on one cross-axis", metrics=metrics)
                )
            if metrics["spacing_range_px"] > spacing_tolerance:
                findings.append(
                    _finding("L9", stage, group, "repeated element spacing is irregular", metrics=metrics)
                )
        repeat_rows.append(row)

    peer_records: dict[str, list[SvgRecord]] = {}
    for record in records:
        group = record.element.get("data-peer-size-group")
        if group:
            peer_records.setdefault(group, []).append(record)
    peer_rows: list[dict[str, Any]] = []
    for group, members in sorted(peer_records.items()):
        if len(members) < 2:
            findings.append(
                _finding("L11", "source", group, "peer-size group requires at least two members")
            )
            continue
        tolerances = [
            _float_attr(
                member.element,
                "data-peer-size-tolerance",
                DEFAULT_PEER_SIZE_TOLERANCE_PX,
            )
            for member in members
        ]
        tolerance = tolerances[0]
        row: dict[str, Any] = {
            "id": group,
            "members": [member.element_id for member in members],
            "tolerance_px": tolerance,
        }
        if any(
            not math.isclose(value, tolerance, abs_tol=1e-9)
            for value in tolerances[1:]
        ) or tolerance > MAX_PEER_SIZE_TOLERANCE_PX:
            findings.append(
                _finding(
                    "L11",
                    "source",
                    group,
                    "peer-size tolerance must be identical for every member and <= 1 px",
                    metrics={"member_tolerances_px": tolerances},
                )
            )
        for stage, boxes_by_id in (("source", source_boxes), ("backend", backend_boxes)):
            measured = [
                (member.element_id, boxes_by_id.get(member.element_id or ""))
                for member in members
            ]
            missing = [element_id or "anonymous" for element_id, box in measured if box is None]
            if missing:
                findings.append(
                    _finding(
                        "L11",
                        stage,
                        group,
                        f"peer-size geometry missing for {', '.join(missing)}",
                    )
                )
                continue
            widths = [box.width for _, box in measured if box is not None]
            heights = [box.height for _, box in measured if box is not None]
            metrics = {
                "widths_px": [round(value, 6) for value in widths],
                "heights_px": [round(value, 6) for value in heights],
                "width_range_px": round(max(widths) - min(widths), 6),
                "height_range_px": round(max(heights) - min(heights), 6),
            }
            row[stage] = metrics
            if max(metrics["width_range_px"], metrics["height_range_px"]) > tolerance:
                findings.append(
                    _finding(
                        "L12",
                        stage,
                        group,
                        "semantic peer shapes do not have equal size",
                        metrics=metrics,
                    )
                )
        peer_rows.append(row)

    layout_peer_records: dict[str, list[SvgRecord]] = {}
    for record in records:
        group = record.element.get("data-peer-group")
        if group:
            layout_peer_records.setdefault(group, []).append(record)
        elif record.element.get("data-peer-axis") is not None or record.element.get(
            "data-peer-order"
        ) is not None:
            target = record.element_id or f"anonymous-{record.tag}"
            findings.append(
                _finding(
                    "L17",
                    "source",
                    target,
                    "peer axis/order requires data-peer-group",
                )
            )

    layout_peer_rows: list[dict[str, Any]] = []
    for group, members in sorted(layout_peer_records.items()):
        row: dict[str, Any] = {
            "id": group,
            "members": [member.element_id for member in members],
        }
        if len(members) < 2:
            findings.append(
                _finding("L17", "source", group, "peer group requires at least two members")
            )
            layout_peer_rows.append(row)
            continue
        raw_axes = {member.element.get("data-peer-axis") for member in members}
        axes = {_normalized_axis(value) for value in raw_axes}
        if len(axes) != 1 or None in axes:
            findings.append(
                _finding(
                    "L17",
                    "source",
                    group,
                    "peer group must declare one consistent x/y/xy axis",
                    metrics={"declared_axes": sorted(str(value) for value in raw_axes)},
                )
            )
            layout_peer_rows.append(row)
            continue
        axis = next(iter(axes))
        orders, order_findings = _ordered_contract(members, group=group, code="L17")
        findings.extend(order_findings)
        if orders is None:
            layout_peer_rows.append(row)
            continue
        size_tolerances = [
            _float_attr(
                member.element,
                "data-peer-size-tolerance",
                DEFAULT_PEER_SIZE_TOLERANCE_PX,
            )
            for member in members
        ]
        axis_tolerances = [
            _float_attr(
                member.element,
                "data-peer-axis-tolerance",
                DEFAULT_PEER_AXIS_TOLERANCE_PX,
            )
            for member in members
        ]
        spacing_tolerances = [
            _float_attr(
                member.element,
                "data-peer-spacing-tolerance",
                DEFAULT_PEER_SPACING_TOLERANCE_PX,
            )
            for member in members
        ]
        tolerances_valid = all(
            math.isclose(value, values[0], abs_tol=1e-9)
            for values in (size_tolerances, axis_tolerances, spacing_tolerances)
            for value in values[1:]
        ) and (
            size_tolerances[0] <= MAX_PEER_SIZE_TOLERANCE_PX
            and axis_tolerances[0] <= MAX_PEER_AXIS_TOLERANCE_PX
            and spacing_tolerances[0] <= MAX_PEER_SPACING_TOLERANCE_PX
        )
        row.update(
            {
                "axis": axis,
                "orders": orders,
                "tolerances_px": {
                    "size": size_tolerances[0],
                    "cross_axis": axis_tolerances[0],
                    "spacing": spacing_tolerances[0],
                },
            }
        )
        if not tolerances_valid:
            findings.append(
                _finding(
                    "L17",
                    "source",
                    group,
                    "peer tolerances must be identical per member and within hard maxima",
                    metrics={
                        "size": size_tolerances,
                        "cross_axis": axis_tolerances,
                        "spacing": spacing_tolerances,
                    },
                )
            )
        for stage, boxes_by_id in (("source", source_boxes), ("backend", backend_boxes)):
            measured: list[tuple[int, str, Box]] = []
            missing: list[str] = []
            for order, member in zip(orders, members):
                if not member.element_id or member.element_id not in boxes_by_id:
                    missing.append(member.element_id or f"anonymous-{member.tag}")
                else:
                    measured.append((order, member.element_id, boxes_by_id[member.element_id]))
            if missing:
                findings.append(
                    _finding(
                        "L17",
                        stage,
                        group,
                        f"peer member geometry missing for {', '.join(missing)}",
                    )
                )
                continue
            metrics = _repeat_metrics(measured, axis)
            row[stage] = metrics
            if axis == "xy":
                if max(metrics["width_range_px"], metrics["height_range_px"]) > size_tolerances[0]:
                    findings.append(
                        _finding(
                            "L18",
                            stage,
                            group,
                            "xy peers do not retain equal dimensions",
                            metrics=metrics,
                        )
                    )
            else:
                centers = metrics["centers_px"]
                if any(
                    centers[index + 1] <= centers[index]
                    for index in range(len(centers) - 1)
                ) or metrics["cross_axis_range_px"] > axis_tolerances[0]:
                    findings.append(
                        _finding(
                            "L19",
                            stage,
                            group,
                            "peers violate declared order or cross-axis alignment",
                            metrics=metrics,
                        )
                    )
            if metrics["spacing_range_px"] > spacing_tolerances[0]:
                findings.append(
                    _finding(
                        "L20",
                        stage,
                        group,
                        "peer centre spacing/translation vector is irregular",
                        metrics=metrics,
                    )
                )
        layout_peer_rows.append(row)

    arrow_readback_path = run.qa_dir / "powerpoint-arrow-readback.json"
    arrow_readback: dict[str, Any] = {}
    if arrow_readback_path.is_file():
        report = read_json(arrow_readback_path)
        if report.get("artifact_sha256") == common.sha256_file(run.pptx_path):
            arrow_readback = {
                item.get("element_id"): item
                for item in report.get("records", [])
                if item.get("element_id")
            }

    gap_rows: list[dict[str, Any]] = []
    for record in records:
        source_id = record.element.get("data-gap-source-id")
        target_id = record.element.get("data-gap-target-id")
        if not source_id and not target_id:
            continue
        target = record.element_id or f"anonymous-{record.tag}"
        axis = record.element.get("data-gap-axis")
        row: dict[str, Any] = {
            "element": target,
            "source_id": source_id,
            "target_id": target_id,
            "axis": axis,
        }
        topology_match = (
            record.element.get("data-source-id") == source_id
            and record.element.get("data-target-id") == target_id
        )
        if (
            not record.element_id
            or not source_id
            or not target_id
            or axis not in {"horizontal", "vertical"}
            or not topology_match
        ):
            findings.append(
                _finding(
                    "L13",
                    "source",
                    target,
                    "gap arrow requires one stable source/target truth and a horizontal/vertical axis",
                )
            )
            gap_rows.append(row)
            continue
        source_box = source_boxes.get(source_id)
        target_box = source_boxes.get(target_id)
        endpoints = _record_endpoints(record)
        if source_box is None or target_box is None or endpoints is None:
            findings.append(
                _finding(
                    "L13",
                    "source",
                    target,
                    "gap arrow source/target geometry or straight endpoints are missing",
                )
            )
            gap_rows.append(row)
            continue
        start_inset = _float_attr(record.element, "data-gap-start-inset", 0.0)
        end_inset = _float_attr(record.element, "data-gap-end-inset", 0.0)
        tolerance = _float_attr(
            record.element,
            "data-gap-tolerance",
            DEFAULT_GAP_TOLERANCE_PX,
        )
        if tolerance > MAX_GAP_TOLERANCE_PX:
            findings.append(
                _finding(
                    "L13",
                    "source",
                    target,
                    f"gap tolerance must be <= {MAX_GAP_TOLERANCE_PX:g} px",
                    metrics={"declared_tolerance_px": tolerance},
                )
            )
            gap_rows.append(row)
            continue
        cross_value = record.element.get("data-gap-cross-position")
        try:
            expected = _gap_expected_endpoints(
                axis,
                source_box,
                target_box,
                start_inset,
                end_inset,
                cross_value,
            )
        except ValueError as exc:
            findings.append(_finding("L13", "source", target, str(exc)))
            gap_rows.append(row)
            continue
        source_metric = {
            "actual_start": [round(value, 6) for value in endpoints[0]],
            "actual_end": [round(value, 6) for value in endpoints[1]],
            "expected_start": [round(value, 6) for value in expected[0]],
            "expected_end": [round(value, 6) for value in expected[1]],
            "start_delta": _point_delta(endpoints[0], expected[0]),
            "end_delta": _point_delta(endpoints[1], expected[1]),
            "tolerance_px": tolerance,
        }
        row["source"] = source_metric
        if max(
            source_metric["start_delta"]["distance"],
            source_metric["end_delta"]["distance"],
        ) > tolerance:
            findings.append(
                _finding(
                    "L13",
                    "source",
                    target,
                    "gap arrow endpoints do not adapt to the declared peer boundaries",
                    metrics=source_metric,
                )
            )
        readback = arrow_readback.get(record.element_id)
        topology_expected = (
            {} if readback is None else readback.get("topology", {}).get("expected", {})
        )
        backend_topology_match = bool(
            topology_expected.get("source_id") == source_id
            and topology_expected.get("target_id") == target_id
        )
        path_geometry = {} if readback is None else readback.get("path_geometry", {})
        semantic_centerline = (
            {} if readback is None else readback.get("semantic_centerline") or {}
        )
        endpoint_geometry = path_geometry
        endpoint_evidence = "saved-pptx-native-path-geometry"
        if not (
            isinstance(endpoint_geometry.get("actual_start"), list)
            and isinstance(endpoint_geometry.get("actual_end"), list)
        ) and semantic_centerline.get("status") == "PASS":
            endpoint_geometry = semantic_centerline
            endpoint_evidence = semantic_centerline.get("evidence")
        actual_start = endpoint_geometry.get("actual_start")
        actual_end = endpoint_geometry.get("actual_end")
        backend_source_box = backend_boxes.get(source_id)
        backend_target_box = backend_boxes.get(target_id)
        backend_metric: dict[str, Any] | None = None
        if (
            backend_source_box is not None
            and backend_target_box is not None
            and isinstance(actual_start, list)
            and len(actual_start) == 2
            and isinstance(actual_end, list)
            and len(actual_end) == 2
        ):
            try:
                backend_expected = _gap_expected_endpoints(
                    axis,
                    backend_source_box,
                    backend_target_box,
                    start_inset,
                    end_inset,
                    cross_value,
                )
                backend_metric = {
                    "actual_start": actual_start,
                    "actual_end": actual_end,
                    "expected_start": [round(value, 6) for value in backend_expected[0]],
                    "expected_end": [round(value, 6) for value in backend_expected[1]],
                    "start_delta": _point_delta(tuple(actual_start), backend_expected[0]),
                    "end_delta": _point_delta(tuple(actual_end), backend_expected[1]),
                    "tolerance_px": tolerance,
                }
            except ValueError:
                backend_metric = None
        backend_endpoint_pass = bool(
            backend_metric
            and max(
                backend_metric["start_delta"]["distance"],
                backend_metric["end_delta"]["distance"],
            )
            <= tolerance
        )
        backend_pass = bool(
            readback
            and readback.get("status") == "PASS"
            and path_geometry.get("status") == "PASS"
            and endpoint_geometry.get("status") == "PASS"
            and readback.get("end_head", {}).get("actual", {}).get("type") not in {None, "none"}
            and backend_topology_match
            and backend_endpoint_pass
        )
        row["backend"] = {
            "shape_id": None if readback is None else readback.get("shape_id"),
            "shape_name": None if readback is None else readback.get("shape_name"),
            "path_and_head_readback_pass": backend_pass,
            "topology_matches_gap_contract": backend_topology_match,
            "endpoint_evidence": endpoint_evidence,
            "endpoint_metrics": backend_metric,
        }
        if not backend_pass:
            findings.append(
                _finding(
                    "L14",
                    "backend",
                    target,
                    "saved PowerPoint gap arrow lacks exact path/head readback evidence",
                )
            )
        gap_rows.append(row)

    source_order = {
        record.element_id: index
        for index, record in enumerate(records, start=1)
        if record.element_id
    }
    z_order_rows: list[dict[str, Any]] = []
    for record in records:
        above = _id_list(record.element.get("data-z-above"))
        below = _id_list(record.element.get("data-z-below"))
        if not above and not below:
            continue
        target = record.element_id or f"anonymous-{record.tag}"
        row: dict[str, Any] = {"element": target, "above": above, "below": below}
        source_pass = bool(record.element_id)
        for other_id in above:
            source_pass = source_pass and other_id in source_order and source_order[target] > source_order[other_id]
        for other_id in below:
            source_pass = source_pass and other_id in source_order and source_order[target] < source_order[other_id]
        row["source_pass"] = source_pass
        if not source_pass:
            findings.append(
                _finding("L15", "source", target, "SVG paint order violates relative z-order contract")
            )
        current_details = backend_details.get(record.element_id or "", [])
        backend_pass = len(current_details) == 1
        if backend_pass:
            current_z = current_details[0]["z_order"]
            for other_id in above:
                peers = backend_details.get(other_id, [])
                backend_pass = backend_pass and bool(peers) and current_z > max(item["z_order"] for item in peers)
            for other_id in below:
                peers = backend_details.get(other_id, [])
                backend_pass = backend_pass and bool(peers) and current_z < min(item["z_order"] for item in peers)
        row["backend_pass"] = backend_pass
        row["backend_objects"] = current_details
        if not backend_pass:
            findings.append(
                _finding(
                    "L15",
                    "backend",
                    target,
                    "saved PowerPoint paint order violates relative z-order contract",
                )
            )
        z_order_rows.append(row)

    text_flow_rows: list[dict[str, Any]] = []
    for record in records:
        flow = record.element.get("data-text-flow")
        if not flow:
            continue
        target = record.element_id or f"anonymous-{record.tag}"
        content = _text_content(record.element).strip()
        graphemes = _graphemes(content)
        container_id = record.element.get("data-text-container")
        if flow == "multiline":
            container_id = record.element.get("data-layout-container")
        frame_overflow_tolerance = _float_attr(
            record.element,
            "data-text-frame-overflow-tolerance",
            1.0,
        )
        row: dict[str, Any] = {
            "element": target,
            "mode": flow,
            "text": content,
            "graphemes": graphemes,
            "container": container_id,
            "frame_overflow_tolerance_px": frame_overflow_tolerance,
        }
        source_pass = bool(
            record.element_id
            and len(graphemes) >= 2
            and container_id
            and frame_overflow_tolerance <= MAX_TEXT_FRAME_OVERFLOW_TOLERANCE_PX
        )
        expected_rotation = round(math.degrees(math.atan2(record.matrix.b, record.matrix.a)), 6)
        if flow == "stacked-characters":
            stack_step_raw = record.element.get("data-text-stack-step")
            source_pass = (
                source_pass
                and record.matrix.is_axis_aligned()
                and stack_step_raw is not None
            )
            expected_stack_step_px = (
                None
                if stack_step_raw is None
                else _float_attr(record.element, "data-text-stack-step", 0.0)
            )
        elif flow == "rotated-word":
            expected_stack_step_px = None
            source_pass = source_pass and math.isclose(
                abs(expected_rotation) % 180.0,
                90.0,
                abs_tol=1e-6,
            )
        elif flow == "multiline":
            expected_stack_step_px = None
            expected_lines = _positioned_text_lines(record.element)
            source_pass = (
                source_pass
                and len(expected_lines) >= 2
                and record.matrix.is_axis_aligned()
            )
        else:
            expected_stack_step_px = None
            source_pass = False
        container_box = source_boxes.get(container_id or "")
        if container_id and container_box is None:
            source_pass = False
        row["source"] = {
            "rotation_deg": expected_rotation,
            "stack_step_px": expected_stack_step_px,
            "contract_pass": source_pass,
        }
        if not source_pass:
            findings.append(
                _finding("L16", "source", target, "invalid explicit text-flow contract")
            )

        details = backend_details.get(record.element_id or "", [])
        box = backend_boxes.get(record.element_id or "")
        backend_pass = len(details) == 1 and box is not None
        if backend_pass:
            detail = details[0]
            actual_text = detail.get("text") or ""
            actual_paragraphs = [item for item in (detail.get("paragraphs") or []) if item]
            actual_rotation = float(detail.get("rotation_deg") or 0.0)
            actual_line_spacings_pt = detail.get("line_spacings_pt") or []
            if flow == "stacked-characters":
                expected_spacing_pt = (
                    None
                    if expected_stack_step_px is None
                    else expected_stack_step_px * PT_PER_PX
                )
                backend_pass = (
                    actual_paragraphs == graphemes
                    and math.isclose(actual_rotation % 360.0, 0.0, abs_tol=1e-6)
                    and expected_spacing_pt is not None
                    and len(actual_line_spacings_pt) == len(graphemes)
                    and all(
                        value is not None
                        and math.isclose(
                            float(value),
                            expected_spacing_pt,
                            abs_tol=TEXT_STACK_SPACING_TOLERANCE_PT,
                        )
                        for value in actual_line_spacings_pt
                    )
                )
            elif flow == "multiline":
                backend_pass = (
                    actual_paragraphs == expected_lines
                    and math.isclose(actual_rotation % 360.0, 0.0, abs_tol=1e-6)
                )
            else:
                backend_pass = (
                    actual_text == content
                    and len(actual_paragraphs) == 1
                    and math.isclose(
                        (actual_rotation - expected_rotation) % 360.0,
                        0.0,
                        abs_tol=1e-6,
                    )
                )
            if backend_pass:
                backend_container = backend_boxes.get(container_id)
                if backend_container is None:
                    backend_pass = False
                else:
                    if flow == "multiline":
                        backend_pass = (
                            max(_overflow(box, backend_container, 0.0).values())
                            <= frame_overflow_tolerance + OOXML_GEOMETRY_EPSILON_PX
                        )
                    else:
                        center_delta = max(
                            abs(box.center_x - backend_container.center_x),
                            abs(box.center_y - backend_container.center_y),
                        )
                        visual_width = box.height if flow == "rotated-word" else box.width
                        visual_height = box.width if flow == "rotated-word" else box.height
                        visual_box = Box(
                            box.center_x - visual_width / 2.0,
                            box.center_y - visual_height / 2.0,
                            visual_width,
                            visual_height,
                        )
                        backend_pass = (
                            center_delta <= 2.0
                            and max(_overflow(visual_box, backend_container, 0.0).values())
                            <= frame_overflow_tolerance + OOXML_GEOMETRY_EPSILON_PX
                        )
            row["backend"] = {
                "text": actual_text,
                "paragraphs": actual_paragraphs,
                "rotation_deg": actual_rotation,
                "line_spacings_pt": actual_line_spacings_pt,
                "contract_pass": backend_pass,
            }
        else:
            row["backend"] = {"contract_pass": False}
        if not backend_pass:
            findings.append(
                _finding(
                    "L16",
                    "backend",
                    target,
                    "saved PowerPoint text flow, rotation, or container alignment drifted",
                )
            )
        text_flow_rows.append(row)

    # Missing unrelated shapes are already handled by the global bindings gate;
    # keep this report scoped to annotated layout objects.
    annotated_ids = {record.element_id for record in annotated if record.element_id}
    missing_annotated = sorted(annotated_ids.intersection(missing_backend))
    report = {
        "schema_version": "1.2.0",
        "svg": "redraw.svg",
        "pptx": "redraw.pptx",
        "defaults_px": {
            "containment": DEFAULT_CONTAINMENT_TOLERANCE_PX,
            "canvas": DEFAULT_CANVAS_TOLERANCE_PX,
            "repeat_size": DEFAULT_REPEAT_SIZE_TOLERANCE_PX,
            "repeat_axis": DEFAULT_REPEAT_AXIS_TOLERANCE_PX,
            "repeat_spacing": DEFAULT_REPEAT_SPACING_TOLERANCE_PX,
            "peer_size": DEFAULT_PEER_SIZE_TOLERANCE_PX,
            "peer_axis": DEFAULT_PEER_AXIS_TOLERANCE_PX,
            "peer_spacing": DEFAULT_PEER_SPACING_TOLERANCE_PX,
            "gap": DEFAULT_GAP_TOLERANCE_PX,
        },
        "canvas": {
            "backend": backend_canvas.as_dict(),
            "checked_elements": len(canvas_rows),
            "objects": canvas_rows,
        },
        "containment": containment_rows,
        "repeat_groups": repeat_rows,
        "peer_size_groups": peer_rows,
        "peer_groups": layout_peer_rows,
        "annotation_coverage": {
            "annotated_records": len(annotated),
            "logical_svg_groups": sum(record.logical_group for record in annotated),
            "containment_contracts": len(containment_rows),
            "repeat_groups": len(repeat_rows),
            "peer_size_groups": len(peer_rows),
            "peer_groups": len(layout_peer_rows),
        },
        "gap_arrows": gap_rows,
        "z_order_contracts": z_order_rows,
        "text_flow_contracts": text_flow_rows,
        "missing_annotated_backend_objects": missing_annotated,
        "findings": findings,
        "pass": not findings,
    }
    return report


def persist_layout_audit(run: common.Run, report: dict[str, Any]) -> None:
    """Explicitly write the layout report into the case's qa/ evidence tree.

    ``audit_layout`` 本身保持只读,便于回归测试在不改写正式案例证据的前提下
    复核报告;落盘由本函数显式承担。
    """

    run.layout_audit_path.write_text(
        json.dumps(report, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )


def strict_blockers(report: dict[str, Any]) -> list[str]:
    return [
        f"layout:{finding['code']}:{finding['stage']}:{finding['target']}"
        for finding in report.get("findings", [])
        if finding.get("severity") == "error"
    ]


def main(argv: list[str] | None = None) -> int:
    import argparse

    parser = argparse.ArgumentParser(prog="autofigure layout", description=__doc__)
    parser.add_argument("run_dir", type=Path)
    args = parser.parse_args(argv)
    run = common.open_run(args.run_dir)
    report = audit_layout(run)
    run.qa_dir.mkdir(exist_ok=True)
    persist_layout_audit(run, report)
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return 0 if report["pass"] else 2


if __name__ == "__main__":
    raise SystemExit(main())
