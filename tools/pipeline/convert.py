"""autofigure convert — 把 VLM 重绘的 SVG 转换为原生可编辑 PPTX。

映射合同（references/prompt-contract.md）：
- rect/circle/ellipse/line/polyline/polygon/path → 原生形状 / custGeom 自由曲线（保留三次贝塞尔）
- text/tspan → 原生文本框 runs（italic、字号、颜色、baseline-shift → 上下标）
- linearGradient → a:gradFill；stroke-dasharray → prstDash；marker → 原生线端，无法单对象表达时拒绝
- <rect id="atomic:*"> 占位符 → 从参考图裁剪对应 bbox 嵌入为位图；assets.json 登记了
  同原子 id 的 source=vtracer-trace 条目时改走 atomic-vector 分支：哈希核验登记的
  vector_source_svg 片段，每个 path 编译为原生 custGeom 成员并包成单个 group
- <image> 容错：按 bbox 从参考图裁剪替代并记 warning；覆盖画布 ≥50% 直接拒绝（防整图截图冒充矢量）
- <g id> 保留为 scene 中的 logical_group，并以成员对象身份建立复合绑定；
  PowerPoint 中仍只生成子对象，不伪造第二个可见 group 对象
"""

from __future__ import annotations

import argparse
import hashlib
import io
import json
import math
import re
import sys
import unicodedata
import xml.etree.ElementTree as ET
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Pt

from tools.core import common
from tools.arrows.arrow_spec import (
    ARROW_SPEC_VERSION,
    DASH_TO_OOXML,
    compiler_strategy,
    head as arrow_head,
    path_from_segments,
    semantic_dash_from_ooxml,
    silhouette_from_segments,
    spec_sha256,
    validate_arrow_spec,
)
from tools.pipeline.layout import Box, collect_svg_boxes
from tools.core.svggeom import Matrix, parse_path_d, parse_transform

SVG_NS = "{http://www.w3.org/2000/svg}"
PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
EMU_PER_PX = 9525
PT_PER_PX = 0.75  # 96 dpi
BASELINE_ASCENT = 0.95  # 实测标定：文本框顶到首行基线（Arial em 比例，含半行距）
TEXT_BOX_PAD_X = 32.0  # 保持锚点不动的透明选择框余量；覆盖 live 实测最大 28.44 px 横向不足
TEXT_BOX_PAD_Y = 20.0  # 覆盖旋转标签 live 实测最大 19.50 px 纵向不足

NAMED_COLORS = {
    "black": "#000000", "white": "#FFFFFF", "red": "#FF0000", "green": "#008000",
    "blue": "#0000FF", "gray": "#808080", "grey": "#808080", "orange": "#FFA500",
    "yellow": "#FFFF00", "purple": "#800080", "brown": "#A52A2A", "pink": "#FFC0CB",
    "none": None, "transparent": None,
}

INHERITED_STYLE_KEYS = (
    "fill", "fill-opacity", "stroke", "stroke-width", "stroke-opacity",
    "stroke-dasharray", "stroke-linecap", "stroke-linejoin", "opacity",
    "data-ppt-dash",
    "font-size", "font-family", "font-style", "font-weight", "text-anchor",
    "text-decoration",
)

BLOCK_AUTOSHAPE_TYPES = {
    "leftRightArrow": MSO_SHAPE.LEFT_RIGHT_ARROW,
}

# Scene projection is compiler-owned.  Repair conversions may retain only
# explicitly namespaced human/model review notes; geometry, topology, backend
# state and every other implementation field must be regenerated from the
# current canonical SVG so a stable ID cannot resurrect stale compiler truth.
SCENE_ELEMENT_ANNOTATION_FIELDS = frozenset(
    {
        "annotations",
        "model_annotations",
        "operator_annotations",
        "review_annotations",
    }
)


def _brace_spec_hash_fields(scene_element: dict) -> dict[str, str]:
    """Return canonical and migration-alias BraceSpec hashes.

    Schema 4 treats ``brace_spec`` as compiler truth.  ``primitive_spec`` is a
    temporary compatibility alias only; when both are present their canonical
    JSON hashes must be identical or conversion fails before publication.
    """

    brace_spec = scene_element.get("brace_spec")
    primitive_spec = scene_element.get("primitive_spec")
    if brace_spec is not None and not isinstance(brace_spec, dict):
        raise common.fail(f"{scene_element.get('id', 'brace')}: brace_spec must be an object")
    if primitive_spec is not None and not isinstance(primitive_spec, dict):
        raise common.fail(
            f"{scene_element.get('id', 'brace')}: primitive_spec migration alias must be an object"
        )
    fields: dict[str, str] = {}
    if isinstance(brace_spec, dict):
        fields["brace_spec_sha256"] = spec_sha256(brace_spec)
    if isinstance(primitive_spec, dict):
        fields["primitive_spec_sha256"] = spec_sha256(primitive_spec)
    if (
        "brace_spec_sha256" in fields
        and "primitive_spec_sha256" in fields
        and fields["brace_spec_sha256"] != fields["primitive_spec_sha256"]
    ):
        raise common.fail(
            f"{scene_element.get('id', 'brace')}: brace_spec and primitive_spec hashes differ"
        )
    return fields


# ---------------------------------------------------------------- 样式与颜色


def _parse_style_attr(style: str | None) -> dict[str, str]:
    result: dict[str, str] = {}
    if style:
        for item in style.split(";"):
            if ":" in item:
                key, value = item.split(":", 1)
                result[key.strip()] = value.strip()
    return result


def _element_style(element: ET.Element, parent_style: dict[str, str]) -> dict[str, str]:
    style = dict(parent_style)
    inline = _parse_style_attr(element.get("style"))
    for key in INHERITED_STYLE_KEYS:
        if element.get(key) is not None:
            style[key] = element.get(key)  # presentation attribute
        if key in inline:
            style[key] = inline[key]
    for key in ("fill-rule", "stroke-linecap", "stroke-linejoin"):
        if key in inline:
            style[key] = inline[key]
    return style


def parse_color(value: str | None) -> tuple[str, float] | None:
    """→ (hex, alpha) 或 None（none/未指定）。"""
    if value is None:
        return None
    value = value.strip()
    if value in NAMED_COLORS:
        named = NAMED_COLORS[value]
        return (named, 1.0) if named else None
    if value.startswith("#"):
        hexpart = value[1:]
        if len(hexpart) == 3:
            hexpart = "".join(ch * 2 for ch in hexpart)
        if len(hexpart) == 6:
            return ("#" + hexpart.upper(), 1.0)
        if len(hexpart) == 8:
            return ("#" + hexpart[:6].upper(), int(hexpart[6:], 16) / 255.0)
    match = re.match(r"rgba?\(([^)]*)\)", value)
    if match:
        parts = [p.strip() for p in match.group(1).split(",")]
        if len(parts) >= 3:
            r, g, b = (int(float(p)) for p in parts[:3])
            alpha = float(parts[3]) if len(parts) > 3 else 1.0
            return (f"#{r:02X}{g:02X}{b:02X}", alpha)
    return None


def _opacity(style: dict[str, str], key: str) -> float:
    raw = style.get(key)
    if raw is None:
        return 1.0
    try:
        value = float(raw)
    except (TypeError, ValueError) as exc:
        raise common.fail(f"invalid SVG {key}: {raw!r}") from exc
    if not math.isfinite(value) or not 0.0 <= value <= 1.0:
        raise common.fail(f"SVG {key} must be finite and between 0 and 1: {raw!r}")
    return value


# ---------------------------------------------------------------- OOXML 辅助


def _px(value: float) -> Emu:
    return Emu(round(value * EMU_PER_PX))


def _insert_before_any(sp_pr, element, successor_tags: tuple[str, ...]) -> None:
    """按 OOXML spPr 子元素顺序插入：fill < ln < effectLst < scene3d < sp3d < extLst。"""
    for tag in successor_tags:
        sibling = sp_pr.find(qn(tag))
        if sibling is not None:
            sibling.addprevious(element)
            return
    sp_pr.append(element)


_FILL_SUCCESSORS = ("a:ln", "a:effectLst", "a:effectDag", "a:scene3d", "a:sp3d", "a:extLst")
_LINE_SUCCESSORS = ("a:effectLst", "a:effectDag", "a:scene3d", "a:sp3d", "a:extLst")


def _set_solid_fill(sp_pr, hex_color: str, alpha: float) -> None:
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
        for child in sp_pr.findall(qn(tag)):
            sp_pr.remove(child)
    fill = sp_pr.makeelement(qn("a:solidFill"), {})
    color = sp_pr.makeelement(qn("a:srgbClr"), {"val": hex_color[1:]})
    if alpha < 1.0:
        color.append(sp_pr.makeelement(qn("a:alpha"), {"val": str(round(alpha * 100000))}))
    fill.append(color)
    _insert_before_any(sp_pr, fill, _FILL_SUCCESSORS)


def _set_no_fill(sp_pr) -> None:
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
        for child in sp_pr.findall(qn(tag)):
            sp_pr.remove(child)
    _insert_before_any(sp_pr, sp_pr.makeelement(qn("a:noFill"), {}), _FILL_SUCCESSORS)


def _set_gradient_fill(sp_pr, stops: list[tuple[float, str, float]], angle_deg: float) -> None:
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
        for child in sp_pr.findall(qn(tag)):
            sp_pr.remove(child)
    grad = sp_pr.makeelement(qn("a:gradFill"), {})
    gs_lst = grad.makeelement(qn("a:gsLst"), {})
    for pos, hex_color, alpha in stops:
        gs = gs_lst.makeelement(qn("a:gs"), {"pos": str(round(pos * 100000))})
        color = gs_lst.makeelement(qn("a:srgbClr"), {"val": hex_color[1:]})
        if alpha < 1.0:
            color.append(gs_lst.makeelement(qn("a:alpha"), {"val": str(round(alpha * 100000))}))
        gs.append(color)
        gs_lst.append(gs)
    grad.append(gs_lst)
    angle = round(math.degrees(angle_deg) * 60000) % 21600000
    grad.append(grad.makeelement(qn("a:lin"), {"ang": str(angle), "scaled": "1"}))
    _insert_before_any(sp_pr, grad, _FILL_SUCCESSORS)


def _apply_line(sp_pr, style: dict[str, str], default_width: float = 1.0) -> None:
    """按 style 在 spPr 上设置 a:ln（颜色/宽度/虚线/圆角）。无 stroke 则不建 a:ln。"""
    stroke = parse_color(style.get("stroke"))
    if stroke is None:
        old = sp_pr.find(qn("a:ln"))
        if old is not None:
            sp_pr.remove(old)
        line = sp_pr.makeelement(qn("a:ln"), {})
        line.append(line.makeelement(qn("a:noFill"), {}))
        _insert_before_any(sp_pr, line, _LINE_SUCCESSORS)
        return
    hex_color, stroke_alpha = stroke
    alpha = stroke_alpha * _opacity(style, "stroke-opacity") * _opacity(style, "opacity")
    try:
        width = float(style.get("stroke-width", default_width))
    except ValueError:
        width = default_width
    old = sp_pr.find(qn("a:ln"))
    if old is not None:
        sp_pr.remove(old)
    line = sp_pr.makeelement(qn("a:ln"), {"w": str(round(width * EMU_PER_PX))})
    line_cap = style.get("stroke-linecap")
    if style.get("data-ppt-dash") == "round_dot":
        line_cap = "round"
    elif style.get("data-ppt-dash") == "square_dot":
        line_cap = "butt"
    if line_cap in ("round", "square"):
        line.set("cap", "rnd" if line_cap == "round" else "sq")
    fill = line.makeelement(qn("a:solidFill"), {})
    color = line.makeelement(qn("a:srgbClr"), {"val": hex_color[1:]})
    if alpha < 1.0:
        color.append(line.makeelement(qn("a:alpha"), {"val": str(round(alpha * 100000))}))
    fill.append(color)
    line.append(fill)
    dash = _dash_preset(
        style.get("stroke-dasharray"),
        width,
        style.get("stroke-linecap"),
        style.get("data-ppt-dash"),
    )
    if dash:
        line.append(line.makeelement(qn("a:prstDash"), {"val": dash}))
    join = style.get("stroke-linejoin", "miter")
    if join == "round":
        line.append(line.makeelement(qn("a:round"), {}))
    elif join == "bevel":
        line.append(line.makeelement(qn("a:bevel"), {}))
    else:
        line.append(line.makeelement(qn("a:miter"), {}))
    _insert_before_any(sp_pr, line, _LINE_SUCCESSORS)


def _dash_preset(
    dasharray: str | None,
    width: float,
    linecap: str | None,
    explicit: str | None = None,
) -> str | None:
    if explicit is not None:
        if explicit not in DASH_TO_OOXML:
            raise ValueError(f"unsupported data-ppt-dash: {explicit}")
        return DASH_TO_OOXML[explicit]
    if not dasharray or dasharray.strip() in ("none", ""):
        return None
    try:
        parts = [float(p) for p in re.split(r"[\s,]+", dasharray.strip()) if p]
    except ValueError as exc:
        raise ValueError(f"invalid stroke-dasharray: {dasharray!r}") from exc
    if not parts or any(not math.isfinite(value) or value <= 0 for value in parts):
        raise ValueError(f"stroke-dasharray values must be finite and positive: {dasharray!r}")
    first = parts[0]
    if first <= width * 1.5:
        # OOXML 无 roundDot/squareDot：圆点用 dot，方点用 sysDot
        return "dot" if linecap == "round" else "sysDot"
    if first <= 4 * width:
        return "sysDash"
    if first <= 8 * width:
        return "dash"
    return "lgDash"


def _disable_shadow(shape) -> None:
    shape.shadow.inherit = False


def _next_shape_id(slide) -> int:
    return slide.shapes._next_shape_id


# ---------------------------------------------------------------- 渐变与 marker 定义


def _collect_defs(tree: ET.Element) -> dict[str, ET.Element]:
    defs: dict[str, ET.Element] = {}
    for element in tree.iter():
        element_id = element.get("id")
        if element_id:
            defs[element_id] = element
    return defs


def _gradient_stops(grad: ET.Element) -> list[tuple[float, str, float]]:
    stops: list[tuple[float, str, float]] = []
    for stop in grad.findall(f"{SVG_NS}stop"):
        style = _parse_style_attr(stop.get("style"))
        offset = stop.get("offset", "0")
        pos = float(offset[:-1]) / 100.0 if offset.endswith("%") else float(offset)
        color_value = stop.get("stop-color") or style.get("stop-color", "#000000")
        parsed = parse_color(color_value) or ("#000000", 1.0)
        opacity = stop.get("stop-opacity") or style.get("stop-opacity")
        alpha = parsed[1] * (float(opacity) if opacity is not None else 1.0)
        stops.append((max(0.0, min(1.0, pos)), parsed[0], alpha))
    return stops


def _gradient_angle(grad: ET.Element) -> float:
    def num(key: str, default: float) -> float:
        raw = grad.get(key)
        if raw is None:
            return default
        return float(raw[:-1]) / 100.0 if raw.endswith("%") else float(raw)

    x1, y1, x2, y2 = num("x1", 0.0), num("y1", 0.0), num("x2", 1.0), num("y2", 0.0)
    return math.atan2(y2 - y1, x2 - x1)


def _resolve_fill_url(style: dict[str, str], defs: dict[str, ET.Element]) -> ET.Element | None:
    fill = style.get("fill", "")
    match = re.match(r"url\(#([^)]+)\)", fill or "")
    if match:
        return defs.get(match.group(1))
    return None


# ---------------------------------------------------------------- 形状发射


class ConvertContext:
    def __init__(
        self,
        slide,
        defs: dict[str, ET.Element],
        source_png: Path | None,
        canvas_width: int,
        canvas_height: int,
        asset_authorizations: dict[str, dict] | None = None,
        layout_boxes: dict[str, Box] | None = None,
        input_route: str = "svg-seeded",
        reference_sha256: str = "",
        case_root: Path | None = None,
    ):
        self.slide = slide
        self.defs = defs
        self.source_png = source_png
        self.canvas_width = canvas_width
        self.canvas_height = canvas_height
        self.canvas_area = float(canvas_width * canvas_height)
        self.asset_authorizations = asset_authorizations or {}
        self.layout_boxes = layout_boxes or {}
        self.input_route = input_route
        self.reference_sha256 = reference_sha256
        self.case_root = case_root
        self.warnings: list[str] = []
        self.counts: dict[str, int] = {}
        self.bindings: list[dict] = []
        self.logical_group_bindings: list[dict] = []
        self.scene_elements: dict[str, dict] = {}
        self.assets: list[dict] = []
        self.asset_spec_count = 0
        self.current_element_id = ""
        self.current_svg_tag = ""
        self._element_counts: dict[str, int] = {}
        self._binding_counts: dict[str, int] = {}
        self._logical_group_stack: list[str] = []
        self.pending_connections: list[tuple[object, str | None, str | None, int, int]] = []

    def bump(self, kind: str) -> None:
        self.counts[kind] = self.counts.get(kind, 0) + 1

    def warn(self, message: str) -> None:
        self.warnings.append(message)

    def begin_element(
        self,
        element: ET.Element,
        tag: str,
        matrix: Matrix | None = None,
        computed_style: dict[str, str] | None = None,
    ) -> tuple[str, str]:
        previous = (self.current_element_id, self.current_svg_tag)
        brace_spec = None
        if tag == "path" and element.get("data-primitive-kind") == "brace":
            from tools.assets.primitives import materialize_brace, transform_brace_spec

            brace_spec = materialize_brace(element)
            if brace_spec is not None:
                brace_spec = transform_brace_spec(brace_spec, matrix or Matrix())
        self._element_counts[tag] = self._element_counts.get(tag, 0) + 1
        self.current_element_id = element.get("id") or f"svg-{tag}-{self._element_counts[tag]:04d}"
        self.current_svg_tag = tag
        edge_semantics = tag in ("line", "path", "polyline", "polygon") and (
            _has_explicit_edge_semantics(element)
        )
        role = "edge" if edge_semantics else element.get("data-role") or tag
        topology = {
            key: value
            for key, value in (
                ("source", element.get("data-source-id")),
                ("target", element.get("data-target-id")),
                ("relation", element.get("data-topology-relation")),
                ("pair_with", element.get("data-pair-with")),
                ("source_gap", element.get("data-source-gap")),
                ("target_gap", element.get("data-target-gap")),
                ("attached", element.get("data-attach")),
            )
            if value
        }
        geometry = {
            key: element.get(key)
            for key in (
                "x",
                "y",
                "width",
                "height",
                "cx",
                "cy",
                "r",
                "rx",
                "ry",
                "x1",
                "y1",
                "x2",
                "y2",
                "points",
                "d",
                "transform",
            )
            if element.get(key) is not None
        }
        scene_element = {
                "id": self.current_element_id,
                "kind": role,
                "role": role,
                "svg_tag": tag,
                "geometry": geometry,
                "topology": topology,
                "z_index": len(self.scene_elements),
                "editable": True,
            }
        if computed_style:
            scene_element["style"] = dict(sorted(computed_style.items()))
        source_attributes = {
            key: value
            for key, value in sorted(element.attrib.items())
            if key != "id"
        }
        if source_attributes:
            scene_element["source_attributes"] = source_attributes
        if tag == "text":
            scene_element["text"] = "".join(element.itertext())
        layout = {
            key: value
            for key, value in (
                ("container_id", element.get("data-layout-container")),
                ("padding", element.get("data-layout-padding")),
                ("tolerance", element.get("data-layout-tolerance")),
                ("repeat_group", element.get("data-repeat-group")),
                ("repeat_axis", element.get("data-repeat-axis")),
                ("repeat_order", element.get("data-repeat-order")),
                 ("repeat_size_tolerance", element.get("data-repeat-size-tolerance")),
                 ("repeat_axis_tolerance", element.get("data-repeat-axis-tolerance")),
                 ("repeat_spacing_tolerance", element.get("data-repeat-spacing-tolerance")),
                 ("peer_size_group", element.get("data-peer-size-group")),
                 ("peer_size_tolerance", element.get("data-peer-size-tolerance")),
                 ("gap_source_id", element.get("data-gap-source-id")),
                 ("gap_target_id", element.get("data-gap-target-id")),
                 ("gap_axis", element.get("data-gap-axis")),
                 ("gap_start_inset", element.get("data-gap-start-inset")),
                 ("gap_end_inset", element.get("data-gap-end-inset")),
                 ("gap_cross_position", element.get("data-gap-cross-position")),
                 ("gap_tolerance", element.get("data-gap-tolerance")),
                 ("z_above", element.get("data-z-above")),
                 ("z_below", element.get("data-z-below")),
                 ("text_flow", element.get("data-text-flow")),
                 ("text_container", element.get("data-text-container")),
                 ("text_stack_step", element.get("data-text-stack-step")),
                 (
                     "text_frame_overflow_tolerance",
                     element.get("data-text-frame-overflow-tolerance"),
                 ),
             )
            if value is not None
        }
        if layout:
            scene_element["layout"] = layout
        if brace_spec is not None:
            # Schema 4 promotes BraceSpec to a first-class scene contract.  The
            # generic PrimitiveSpec field remains a byte-equivalent migration
            # alias for existing consumers and must never carry separate truth.
            scene_element["brace_spec"] = brace_spec
            scene_element["primitive_spec"] = brace_spec
        topology_relation = _topology_relation_metadata(element)
        if topology_relation is not None:
            scene_element["topology_relation"] = topology_relation
        self.scene_elements.setdefault(self.current_element_id, scene_element)
        return previous

    def end_element(self, previous: tuple[str, str]) -> None:
        self.current_element_id, self.current_svg_tag = previous

    def begin_logical_group(
        self,
        element: ET.Element,
        style: dict[str, str],
        matrix: Matrix,
    ) -> dict[str, object] | None:
        """Record an SVG semantic group without manufacturing a PPT shape.

        ID-less ``g`` elements are transform/style carriers only.  An explicit
        ID, however, is a stable logical object that may be named by a frozen
        inventory or by ArrowSpec topology.  Its PowerPoint representation is
        the ordered set of native descendant shapes, not a second visible
        object and not a python-pptx group wrapper.
        """

        group_id = element.get("id")
        if not group_id:
            return None
        if (
            _has_explicit_edge_semantics(element)
        ):
            raise common.fail(
                f"{group_id}: a logical arrow cannot be authored as an SVG group; "
                "one logical arrow must compile to one visible object"
            )
        if group_id in self.scene_elements:
            raise common.fail(f"duplicate SVG semantic id: {group_id}")

        source_attributes = {
            key: value
            for key, value in sorted(element.attrib.items())
            if key != "id"
        }
        geometry = {
            "transform": element.get("transform")
        } if element.get("transform") is not None else {}
        scene_element: dict[str, object] = {
            "id": group_id,
            "kind": "logical_group",
            "role": element.get("data-role") or "logical_group",
            "svg_tag": "g",
            "geometry": geometry,
            "member_ids": [],
            "logical_descendant_group_ids": [],
            "z_index": len(self.scene_elements),
            "editable": True,
            "composite": True,
        }
        if style:
            scene_element["style"] = dict(sorted(style.items()))
        if source_attributes:
            scene_element["source_attributes"] = source_attributes
        self.scene_elements[group_id] = scene_element
        self._logical_group_stack.append(group_id)
        return {
            "group_id": group_id,
            "binding_start": len(self.bindings),
            "logical_group_start": len(self.logical_group_bindings),
            "matrix": matrix,
        }

    def end_logical_group(self, token: dict[str, object] | None) -> None:
        if token is None:
            return
        group_id = str(token["group_id"])
        if not self._logical_group_stack or self._logical_group_stack[-1] != group_id:
            raise common.fail(f"logical SVG group stack drift: {group_id}")
        self._logical_group_stack.pop()

        binding_start = int(token["binding_start"])
        group_start = int(token["logical_group_start"])
        member_rows = self.bindings[binding_start:]
        member_ids = list(
            dict.fromkeys(
                row["element_id"]
                for row in member_rows
                if isinstance(row.get("element_id"), str) and row["element_id"]
            )
        )
        descendant_group_ids = list(
            dict.fromkeys(
                row["element_id"]
                for row in self.logical_group_bindings[group_start:]
                if isinstance(row.get("element_id"), str) and row["element_id"]
            )
        )
        backend_ids = [int(row["shape_id"]) for row in member_rows]
        backend_names = [str(row["shape_name"]) for row in member_rows]
        identities = [
            {"shape_id": shape_id, "shape_name": shape_name}
            for shape_id, shape_name in zip(backend_ids, backend_names, strict=True)
        ]

        attachment_shape_id: int | None = None
        attachment_shape_name: str | None = None
        by_id = {
            int(shape.shape_id): shape
            for shape in _iter_readback_shapes(self.slide.shapes)
        }
        ranked: list[tuple[int, int, str]] = []
        for order, (shape_id, shape_name) in enumerate(
            zip(backend_ids, backend_names, strict=True)
        ):
            shape = by_id.get(shape_id)
            area = 0 if shape is None else int(shape.width) * int(shape.height)
            ranked.append((area, -order, shape_name))
        if ranked:
            _, neg_order, _ = max(ranked)
            representative_index = -neg_order
            attachment_shape_id = backend_ids[representative_index]
            attachment_shape_name = backend_names[representative_index]

        scene_group = self.scene_elements[group_id]
        scene_group["member_ids"] = member_ids
        scene_group["logical_descendant_group_ids"] = descendant_group_ids
        self.logical_group_bindings.append(
            {
                "element_id": group_id,
                "binding_kind": "logical-group-composite",
                "object_kind": "logical-group",
                "editable": True,
                "member_element_ids": member_ids,
                "logical_descendant_group_ids": descendant_group_ids,
                "backend_object_ids": backend_ids,
                "backend_object_names": backend_names,
                "backend_object_identities": identities,
                "visible_object_count": len(identities),
                "attachment_shape_id": attachment_shape_id,
                "attachment_shape_name": attachment_shape_name,
                "attachment_policy": "largest-member-area-then-source-order",
                "readback_found": False,
            }
        )

    def _shape_name(self, kind: str) -> str:
        element_id = re.sub(r"[^A-Za-z0-9_.-]+", "-", self.current_element_id).strip("-")
        element_id = element_id[:80] or "element"
        key = f"{element_id}:{kind}"
        self._binding_counts[key] = self._binding_counts.get(key, 0) + 1
        return f"af-{element_id}-{kind}-{self._binding_counts[key]:02d}"

    def register_shape(self, shape, kind: str, *, editable: bool = True) -> str:
        name = self._shape_name(kind)
        c_nv_pr = shape._element.find(f".//{qn('p:cNvPr')}")
        if c_nv_pr is not None:
            c_nv_pr.set("name", name)
            self._write_semantic_description(c_nv_pr)
        self.register_raw(shape.shape_id, name, kind, editable=editable)
        return name

    def register_xml(self, shape, shape_id: int, kind: str, *, editable: bool = True) -> str:
        name = self._shape_name(kind)
        c_nv_pr = shape.find(f".//{qn('p:cNvPr')}")
        if c_nv_pr is not None:
            c_nv_pr.set("name", name)
            self._write_semantic_description(c_nv_pr)
        self.register_raw(shape_id, name, kind, editable=editable)
        return name

    def register_raw(self, shape_id: int, name: str, kind: str, *, editable: bool = True) -> None:
        scene_element = self.scene_elements.get(self.current_element_id, {})
        binding = {
            "element_id": self.current_element_id,
            "shape_id": int(shape_id),
            "shape_name": name,
            "object_kind": kind,
            "editable": editable,
            "semantic_group_id": self.current_element_id,
        }
        if self._logical_group_stack:
            binding["logical_group_ids"] = list(self._logical_group_stack)
        arrow_spec = scene_element.get("arrow_spec")
        if isinstance(arrow_spec, dict):
            binding["arrow_spec_sha256"] = spec_sha256(arrow_spec)
            binding["single_visible_object"] = arrow_spec.get("single_visible_object") is True
        binding.update(_brace_spec_hash_fields(scene_element))
        topology_relation = scene_element.get("topology_relation")
        if isinstance(topology_relation, dict):
            binding["topology_relation"] = dict(topology_relation)
        self.bindings.append(binding)
        self.scene_elements.setdefault(
            self.current_element_id,
            {
                "id": self.current_element_id,
                "kind": "edge" if kind in ("connector", "line", "freeform-arrow") else kind,
                "svg_tag": self.current_svg_tag,
                "editable": editable,
            },
        )
        if not editable:
            self.scene_elements[self.current_element_id]["editable"] = False

    def _write_semantic_description(self, c_nv_pr) -> None:
        scene_element = self.scene_elements.get(self.current_element_id, {})
        payload: dict[str, object] = {"autofigure_element_id": self.current_element_id}
        arrow_spec = scene_element.get("arrow_spec")
        if isinstance(arrow_spec, dict):
            payload["arrow_spec_sha256"] = spec_sha256(arrow_spec)
            # A closed block-arrow freeform has no native Office centerline:
            # the visible object is its single filled silhouette.  Persist the
            # declared semantic centerline beside the spec hash so saved-PPTX
            # readback can verify the logical route without misusing the
            # silhouette's first/last contour points as arrow endpoints.
            if arrow_spec.get("representation") == "block_arrow":
                payload["arrow_semantic_centerline"] = arrow_spec.get("path")
        payload.update(_brace_spec_hash_fields(scene_element))
        topology_relation = scene_element.get("topology_relation")
        if isinstance(topology_relation, dict):
            payload["topology_relation"] = dict(topology_relation)
        c_nv_pr.set("descr", json.dumps(payload, sort_keys=True, separators=(",", ":")))

    def attach_asset_contracts(
        self,
        scene: dict,
        assets: dict,
        frozen_inventory: dict,
    ) -> int:
        """Attach AssetSpecs after every logical group has closed.

        ``asset_spec.py`` owns all frozen-contract interpretation.  This method
        only projects the resulting single-owner identity into compiler
        bindings and the already-created native PowerPoint members.
        """

        from tools.assets.asset_spec import AssetSpecError, asset_spec_sha256, attach_asset_specs

        carrier = {
            "schema_version": scene.get("schema_version"),
            "kind": "scene",
            "reference_sha256": scene.get("reference_sha256"),
            "canvas": scene.get("canvas"),
            # Values are intentionally shared with ``self.scene_elements`` so
            # the canonical attach operation mutates compiler-owned records.
            "elements": list(self.scene_elements.values()),
            "edges": [],
        }
        attach_asset_specs(carrier, assets, frozen_inventory)
        contracted_groups = {
            element_id: element
            for element_id, element in self.scene_elements.items()
            if isinstance(element.get("asset_spec"), dict)
        }
        if not contracted_groups:
            self.asset_spec_count = 0
            return 0

        group_bindings = {
            binding.get("element_id"): binding
            for binding in self.logical_group_bindings
            if isinstance(binding.get("element_id"), str)
        }
        shapes_by_identity = {
            (int(shape.shape_id), shape.name): shape
            for shape in _iter_readback_shapes(self.slide.shapes)
        }
        member_owner: dict[str, str] = {}
        errors: list[str] = []
        for asset_group_id, group in sorted(contracted_groups.items()):
            spec = group["asset_spec"]
            digest = asset_spec_sha256(spec)
            if group.get("asset_spec_sha256") != digest:
                errors.append(f"asset-spec-hash:{asset_group_id}")
                continue
            logical_binding = group_bindings.get(asset_group_id)
            if logical_binding is None:
                errors.append(f"asset-spec-binding-logical-missing:{asset_group_id}")
                continue
            logical_binding["asset_spec"] = spec
            logical_binding["asset_spec_sha256"] = digest
            for member_id in spec["member_ids"]:
                existing_owner = member_owner.get(member_id)
                if existing_owner is not None and existing_owner != asset_group_id:
                    errors.append(
                        "asset-spec-member-multiple-assets:"
                        f"{member_id}:{existing_owner}:{asset_group_id}"
                    )
                    continue
                member_owner[member_id] = asset_group_id
                member_bindings = [
                    binding
                    for binding in self.bindings
                    if binding.get("element_id") == member_id
                    and asset_group_id in binding.get("logical_group_ids", [])
                ]
                if not member_bindings:
                    errors.append(
                        f"asset-spec-binding-member-missing:{asset_group_id}:{member_id}"
                    )
                    continue
                for binding in member_bindings:
                    binding["asset_group_id"] = asset_group_id
                    binding["asset_spec_sha256"] = digest
                    identity = (int(binding["shape_id"]), binding["shape_name"])
                    shape = shapes_by_identity.get(identity)
                    if shape is None:
                        errors.append(
                            f"asset-spec-binding-shape-missing:{asset_group_id}:{member_id}"
                        )
                        continue
                    c_nv_pr = shape._element.find(f".//{qn('p:cNvPr')}")
                    if c_nv_pr is None:
                        errors.append(
                            f"asset-spec-binding-description-missing:{asset_group_id}:{member_id}"
                        )
                        continue
                    try:
                        description = json.loads(c_nv_pr.get("descr") or "{}")
                    except json.JSONDecodeError:
                        description = {}
                    if not isinstance(description, dict):
                        description = {}
                    description.update(
                        {
                            "asset_group_id": asset_group_id,
                            "asset_spec_sha256": digest,
                        }
                    )
                    c_nv_pr.set(
                        "descr",
                        json.dumps(
                            description,
                            sort_keys=True,
                            separators=(",", ":"),
                        ),
                    )
        if errors:
            raise AssetSpecError(errors)
        self.asset_spec_count = len(contracted_groups)
        return self.asset_spec_count

    def register_asset(self, asset_id: str, bbox: list[int], source_sha256: str) -> None:
        authorization = self.asset_authorizations.get(asset_id, {})
        authorized = authorization.get("authorized") is True
        raster_reason = authorization.get("raster_reason") or (
            "Reference-faithful irreducible creative microasset authorized by the user."
        )
        decomposition_note = authorization.get("decomposition_note") or (
            "No faithful native decomposition; formal text, formulas, nodes, and topology remain native."
        )
        record = {
            **authorization,
            "id": asset_id,
            "source": "reference_crop",
            "source_sha256": source_sha256,
            "bbox": bbox,
            "editable": False,
            "authorized": authorized,
            "source_tightly_cropped": True,
            "atomic_raster_unit": True,
            "contains_reconstructable_content": False,
            "raster_reason": raster_reason,
            "decomposition_note": decomposition_note,
        }
        self.assets.append(record)
        if not authorized:
            self.warn(
                f"{asset_id}: atomic raster has no explicit assets.json authorization; "
                "strict approval will fail"
            )

    def register_vector_asset(self, asset_id: str) -> None:
        """Record one compiled atomic-vector asset for the assets.json merge.

        The derived assets.json entry already carries the full atomic-vector
        contract fields (``tools/asset_spec.py``).  The generated record only
        asserts identity, source, and editability so the merge keeps the entry
        inside its closed field set.
        """
        self.assets.append(
            {
                "id": asset_id,
                "editable": True,
                "source": "vtracer-trace",
            }
        )

    def add_pending_connection(self, shape, element: ET.Element) -> None:
        if not _has_explicit_edge_semantics(element):
            return
        source = element.get("data-source-id")
        target = element.get("data-target-id")
        if (not source and not target) or element.get("data-attach") == "false":
            return
        self.pending_connections.append(
            (
                shape,
                source,
                target,
                int(element.get("data-source-site", "0")),
                int(element.get("data-target-site", "0")),
            )
        )

    def resolve_connections(self) -> None:
        shape_ids = {
            binding["element_id"]: binding["shape_id"]
            for binding in self.bindings
            if binding["object_kind"] not in ("connector", "line", "freeform-arrow")
        }
        shape_ids.update(
            {
                binding["element_id"]: int(binding["attachment_shape_id"])
                for binding in self.logical_group_bindings
                if isinstance(binding.get("attachment_shape_id"), int)
            }
        )
        for connector, source, target, source_site, target_site in self.pending_connections:
            c_nv = connector._element.find(f".//{qn('p:cNvCxnSpPr')}")
            if c_nv is None:
                self.warn(f"{self.current_element_id}: connector has no cNvCxnSpPr; attachment skipped")
                continue
            for tag, element_id, site in (
                ("a:stCxn", source, source_site),
                ("a:endCxn", target, target_site),
            ):
                if not element_id:
                    continue
                shape_id = shape_ids.get(element_id)
                if shape_id is None:
                    self.warn(f"connector target id not found: {element_id}")
                    continue
                c_nv.append(c_nv.makeelement(qn(tag), {"id": str(shape_id), "idx": str(site)}))

def _has_marker(element: ET.Element) -> bool:
    inline = _parse_style_attr(element.get("style"))
    return any(element.get(key) or inline.get(key) for key in ("marker-start", "marker-end"))


def _has_explicit_edge_semantics(element: ET.Element) -> bool:
    """Return whether an SVG primitive explicitly requests ArrowSpec/attachment.

    ``data-source-id`` and ``data-target-id`` are shared by scientific topology
    relations.  They are therefore evidence about relation endpoints, not by
    themselves permission to manufacture an ArrowSpec or PowerPoint connector
    attachment.
    """

    role = (element.get("data-role") or "").strip().lower()
    if role in {"arrow", "edge", "connector"} or _has_marker(element):
        return True
    if any(name.startswith("data-arrow-") for name in element.attrib):
        return True
    if any(
        element.get(name) is not None
        for name in (
            "data-start-head-type",
            "data-end-head-type",
            "data-start-arrow-width",
            "data-start-arrow-length",
            "data-end-arrow-width",
            "data-end-arrow-length",
            "data-source-site",
            "data-target-site",
            "data-attach",
            "data-routing",
        )
    ):
        return True
    return False


def _topology_relation_metadata(element: ET.Element) -> dict[str, str] | None:
    relation = element.get("data-topology-relation")
    if relation is None:
        return None
    values = {
        "element_id": element.get("id"),
        "source_id": element.get("data-source-id"),
        "target_id": element.get("data-target-id"),
        "relation": relation,
    }
    return {
        key: value
        for key, value in values.items()
        if isinstance(value, str) and value
    }


def _apply_fill_and_line(sp_pr, style: dict[str, str], ctx: ConvertContext) -> None:
    grad = _resolve_fill_url(style, ctx.defs)
    if grad is not None and grad.tag == f"{SVG_NS}linearGradient":
        stops = _gradient_stops(grad)
        if stops:
            _set_gradient_fill(sp_pr, stops, _gradient_angle(grad))
        else:
            _set_no_fill(sp_pr)
    else:
        if grad is not None:
            ctx.warn(f"暂不支持的渐变类型: {grad.tag}，按无填充处理")
        fill = parse_color(style.get("fill", "#000000"))
        if fill is None:
            _set_no_fill(sp_pr)
        else:
            hex_color, fill_alpha = fill
            alpha = fill_alpha * _opacity(style, "fill-opacity") * _opacity(style, "opacity")
            _set_solid_fill(sp_pr, hex_color, alpha)
    _apply_line(sp_pr, style)


def _emit_rect(ctx: ConvertContext, el: ET.Element, style: dict[str, str], matrix: Matrix) -> None:
    x = float(el.get("x", 0))
    y = float(el.get("y", 0))
    w = float(el.get("width", 0))
    h = float(el.get("height", 0))
    if w <= 0 or h <= 0:
        return
    element_id = el.get("id", "")
    if element_id.startswith("atomic:"):
        _emit_atomic(ctx, element_id, x, y, w, h, matrix)
        return
    if not matrix.is_axis_aligned():
        segments = [("M", *matrix.apply(x, y)), ("L", *matrix.apply(x + w, y)),
                    ("L", *matrix.apply(x + w, y + h)), ("L", *matrix.apply(x, y + h)), ("Z",)]
        _emit_freeform(ctx, segments, style)
        return
    (x0, y0), (x1, y1) = matrix.apply(x, y), matrix.apply(x + w, y + h)
    rx = float(el.get("rx", el.get("ry", 0)) or 0)
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rx > 0 else MSO_SHAPE.RECTANGLE
    shape = ctx.slide.shapes.add_shape(kind, _px(min(x0, x1)), _px(min(y0, y1)), _px(abs(x1 - x0)), _px(abs(y1 - y0)))
    if rx > 0:
        try:
            shape.adjustments[0] = min(0.5, rx / min(w, h))
        except (IndexError, AttributeError):
            pass
    _disable_shadow(shape)
    _apply_fill_and_line(shape._element.spPr, style, ctx)
    ctx.register_shape(shape, "rect")
    ctx.bump("rect")


def _emit_ellipse(ctx: ConvertContext, el: ET.Element, style: dict[str, str], matrix: Matrix) -> None:
    if el.tag == f"{SVG_NS}circle":
        cx, cy = float(el.get("cx", 0)), float(el.get("cy", 0))
        rx = ry = float(el.get("r", 0))
    else:
        cx, cy = float(el.get("cx", 0)), float(el.get("cy", 0))
        rx, ry = float(el.get("rx", 0)), float(el.get("ry", 0))
    if rx <= 0 or ry <= 0:
        return
    x0, y0 = matrix.apply(cx - rx, cy - ry)
    x1, y1 = matrix.apply(cx + rx, cy + ry)
    shape = ctx.slide.shapes.add_shape(
        MSO_SHAPE.OVAL, _px(min(x0, x1)), _px(min(y0, y1)), _px(abs(x1 - x0)), _px(abs(y1 - y0))
    )
    _disable_shadow(shape)
    _apply_fill_and_line(shape._element.spPr, style, ctx)
    ctx.register_shape(shape, "ellipse")
    ctx.bump("ellipse")


def _marker_reference(el: ET.Element, side: str) -> str | None:
    attr = f"marker-{side}"
    value = el.get(attr) or _parse_style_attr(el.get("style")).get(attr)
    match = re.match(r"url\(#([^)]+)\)", value or "")
    return match.group(1) if match else None


def _native_marker_spec(
    marker: ET.Element,
    line_style: dict[str, str],
    el: ET.Element,
    side: str,
) -> dict[str, str] | None:
    """Map a simple SVG triangle to a real DrawingML line end.

    Unsupported marker artwork is deliberately rejected here so the caller can
    emit an explicit custom-freeform fallback instead of silently substituting a
    default PowerPoint arrow.
    """
    if marker.get("orient", "auto") not in ("auto", "auto-start-reverse"):
        return None
    paths = [child for child in marker if child.tag == f"{SVG_NS}path"]
    if len(paths) != 1:
        return None
    segments = parse_path_d(paths[0].get("d", ""))
    child_matrix = parse_transform(paths[0].get("transform"))
    uniform_scale = math.hypot(child_matrix.a, child_matrix.b)
    perpendicular_scale = math.hypot(child_matrix.c, child_matrix.d)
    orthogonality = child_matrix.a * child_matrix.c + child_matrix.b * child_matrix.d
    if (
        uniform_scale <= 1e-9
        or perpendicular_scale <= 1e-9
        or abs(uniform_scale - perpendicular_scale) > 1e-6
        or abs(orthogonality) > 1e-6
    ):
        # A non-uniform/skewed marker silhouette cannot be represented by a
        # discrete native PowerPoint line end without fidelity loss.
        return None
    points = [
        child_matrix.apply(part[1], part[2])
        for part in segments
        if part[0] in ("M", "L")
    ]
    if len(points) != 3:
        return None
    reverse_start = side == "start" and marker.get("orient", "auto") == "auto"
    tip_x = min(point[0] for point in points) if reverse_start else max(point[0] for point in points)
    tip_indices = [index for index, point in enumerate(points) if abs(point[0] - tip_x) <= 1e-6]
    if len(tip_indices) != 1:
        return None
    base = [point for index, point in enumerate(points) if index != tip_indices[0]]
    tip = points[tip_indices[0]]
    if (
        abs(base[0][0] - base[1][0]) > 1e-6
        or abs(tip[1] - (base[0][1] + base[1][1]) / 2.0) > 1e-6
    ):
        return None

    marker_style = _element_style(paths[0], _element_style(marker, {}))
    line_color = parse_color(line_style.get("stroke"))
    marker_color = parse_color(marker_style.get("fill")) or parse_color(marker_style.get("stroke"))
    if line_color is None or marker_color is None or line_color[0] != marker_color[0]:
        return None

    explicit_type = el.get(f"data-{side}-arrow-type") or el.get("data-arrow-type")
    if explicit_type == "open":
        explicit_type = "arrow"
    closed = any(part[0] == "Z" for part in segments) or marker_style.get("fill") not in (
        None,
        "none",
        "transparent",
    )
    arrow_type = explicit_type or ("triangle" if closed else "arrow")
    if arrow_type not in {"triangle", "arrow", "stealth", "diamond", "oval"}:
        return None

    xs = [point[0] for point in points]
    ys = [point[1] for point in points]
    try:
        stroke_width = max(float(line_style.get("stroke-width", "1")), 0.01)
    except ValueError:
        stroke_width = 1.0
    length_ratio = max(xs) - min(xs)
    width_ratio = max(ys) - min(ys)
    view_box_values = [
        float(value)
        for value in re.split(r"[\s,]+", marker.get("viewBox", "").strip())
        if value
    ]
    if len(view_box_values) == 4 and view_box_values[2] > 0 and view_box_values[3] > 0:
        try:
            viewport_width = float(marker.get("markerWidth", "3"))
            viewport_height = float(marker.get("markerHeight", "3"))
        except ValueError:
            return None
        view_scale = min(
            viewport_width / view_box_values[2],
            viewport_height / view_box_values[3],
        )
        length_ratio *= view_scale
        width_ratio *= view_scale
    marker_units = marker.get("markerUnits", "strokeWidth")
    if marker_units == "strokeWidth":
        # Physical marker size scales with stroke width, so the ratio below is
        # already expressed in line-width units.
        length_ratio *= stroke_width
        width_ratio *= stroke_width
    elif marker_units != "userSpaceOnUse":
        return None

    def bucket(value: float) -> str:
        ratio = value / stroke_width
        if ratio <= 2.25:
            return "sm"
        if ratio <= 4.0:
            return "med"
        return "lg"

    width = el.get(f"data-{side}-arrow-width") or el.get("data-arrow-width")
    length = el.get(f"data-{side}-arrow-length") or el.get("data-arrow-length")
    width = width or bucket(width_ratio)
    length = length or bucket(length_ratio)
    if width not in {"sm", "med", "lg"} or length not in {"sm", "med", "lg"}:
        return None
    return {"type": arrow_type, "w": width, "len": length}


def _marker_custom_payload(marker_id: str, marker: ET.Element | None) -> dict[str, object]:
    if marker is None:
        return {"marker_id": marker_id, "missing": True}
    paths = []
    for child in marker:
        if child.tag != f"{SVG_NS}path":
            continue
        paths.append(
            {
                "d": child.get("d", ""),
                "transform": child.get("transform"),
                "fill": child.get("fill"),
                "stroke": child.get("stroke"),
            }
        )
    return {
        "marker_id": marker_id,
        "marker_units": marker.get("markerUnits", "strokeWidth"),
        "marker_width": marker.get("markerWidth"),
        "marker_height": marker.get("markerHeight"),
        "ref_x": marker.get("refX", "0"),
        "ref_y": marker.get("refY", "0"),
        "orient": marker.get("orient", "0"),
        "view_box": marker.get("viewBox"),
        "paths": paths,
    }


def _marker_color(marker: ET.Element | None) -> str | None:
    if marker is None:
        return None
    for child in marker:
        if child.tag != f"{SVG_NS}path":
            continue
        style = _element_style(child, _element_style(marker, {}))
        color = parse_color(style.get("fill")) or parse_color(style.get("stroke"))
        if color:
            return color[0]
    return None


def _arrow_head_from_svg(
    ctx: ConvertContext,
    el: ET.Element,
    style: dict[str, str],
    side: str,
) -> dict[str, object]:
    marker_id = _marker_reference(el, side)
    if not marker_id:
        return arrow_head()
    marker = ctx.defs.get(marker_id)
    native = _native_marker_spec(marker, style, el, side) if marker is not None else None
    color = _marker_color(marker)
    if native is not None:
        return arrow_head(
            "open" if native["type"] == "arrow" else native["type"],
            width=native["w"],
            length=native["len"],
            color=color,
        )
    width = el.get(f"data-{side}-arrow-width") or el.get("data-arrow-width") or "med"
    length = el.get(f"data-{side}-arrow-length") or el.get("data-arrow-length") or "med"
    if width not in {"sm", "med", "lg"}:
        raise common.fail(
            f"{el.get('id') or 'arrow'}:{side}: unsupported arrow width {width!r}"
        )
    if length not in {"sm", "med", "lg"}:
        raise common.fail(
            f"{el.get('id') or 'arrow'}:{side}: unsupported arrow length {length!r}"
        )
    return arrow_head(
        "custom",
        width=width,
        length=length,
        color=color,
        custom_path=_marker_custom_payload(marker_id, marker),
    )


def _block_head_from_svg(
    el: ET.Element,
    style: dict[str, str],
    side: str,
) -> dict[str, object]:
    """Read semantic heads already embodied by a closed block silhouette.

    A block arrow must not attach SVG markers to its outline: doing so creates
    a second visible head.  These data attributes preserve the directional
    semantics while the compiler emits exactly one filled PowerPoint shape.
    """

    head_type = el.get(f"data-{side}-head-type", "none")
    if head_type == "none":
        return arrow_head()
    if head_type not in {
        "open",
        "triangle",
        "stealth",
        "diamond",
        "oval",
        "custom",
    }:
        raise common.fail(
            f"{el.get('id') or 'arrow'}:{side}: unsupported semantic block head "
            f"{head_type!r}"
        )
    width = el.get(f"data-{side}-arrow-width") or "lg"
    length = el.get(f"data-{side}-arrow-length") or "sm"
    color = parse_color(style.get("fill")) or parse_color(style.get("stroke"))
    custom_path = None
    if head_type == "custom":
        # For a block arrow the custom head is already part of the one closed
        # silhouette.  Preserve that semantic fact without manufacturing a
        # second PowerPoint object or pretending it is a line-end enum.
        custom_path = {
            "kind": "embedded-silhouette",
            "side": side,
            "element_id": el.get("id"),
        }
    return arrow_head(
        head_type,
        width=width,
        length=length,
        color=None if color is None else color[0],
        custom_path=custom_path,
    )


def _segments_bbox(segments: list[tuple]) -> list[float]:
    coordinates: list[tuple[float, float]] = []
    for segment in segments:
        values = segment[1:]
        coordinates.extend(
            (float(values[index]), float(values[index + 1]))
            for index in range(0, len(values) - 1, 2)
        )
    if not coordinates:
        raise common.fail("block arrow has no measurable silhouette")
    min_x = min(point[0] for point in coordinates)
    min_y = min(point[1] for point in coordinates)
    max_x = max(point[0] for point in coordinates)
    max_y = max(point[1] for point in coordinates)
    return [min_x, min_y, max_x - min_x, max_y - min_y]


def _block_autoshape_spec(el: ET.Element, segments: list[tuple]) -> dict[str, object] | None:
    subtype = el.get("data-ppt-autoshape")
    if subtype is None:
        return None
    if subtype not in BLOCK_AUTOSHAPE_TYPES:
        raise common.fail(
            f"{el.get('id') or 'arrow'}: unsupported block AutoShape {subtype!r}"
        )
    raw_adjustments = el.get("data-ppt-adjustments", "0.5 0.5")
    try:
        adjustments = [
            float(value)
            for value in re.split(r"[\s,]+", raw_adjustments.strip())
            if value
        ]
    except ValueError as exc:
        raise common.fail(
            f"{el.get('id') or 'arrow'}: invalid block AutoShape adjustments"
        ) from exc
    if len(adjustments) != 2 or any(
        not math.isfinite(value) or value < 0 or value > 1
        for value in adjustments
    ):
        raise common.fail(
            f"{el.get('id') or 'arrow'}: block AutoShape requires two adjustments in [0,1]"
        )
    return {
        "subtype": subtype,
        "adjustments": adjustments,
        "bbox": _segments_bbox(segments),
    }


def _left_right_block_centerline(
    el: ET.Element,
    segments: list[tuple],
    autoshape: dict[str, object] | None,
) -> dict[str, object] | None:
    """Derive a semantic centerline from a canonical left-right silhouette.

    The closed silhouette remains the visible-geometry truth.  This centerline
    exists only to bind start/end direction, endpoint, and tangent evidence; it
    is never substituted for the block-arrow shape itself.
    """

    if autoshape is None:
        return None
    if autoshape.get("subtype") != "leftRightArrow":
        raise common.fail(
            f"{el.get('id') or 'arrow'}: block AutoShape has no supported "
            "left-right centerline semantics"
        )
    if any(segment[0] not in {"M", "L", "Z"} for segment in segments):
        raise common.fail(
            f"{el.get('id') or 'arrow'}: leftRightArrow silhouette must be an "
            "axis-aligned closed polygon"
        )
    # Validate the one-subpath closed silhouette before deriving anything.
    silhouette_from_segments(segments)
    points = [
        (float(segment[1]), float(segment[2]))
        for segment in segments
        if segment[0] in {"M", "L"}
    ]
    left, top, width, height = (float(value) for value in autoshape["bbox"])
    if width <= height or width <= 0 or height <= 0:
        raise common.fail(
            f"{el.get('id') or 'arrow'}: leftRightArrow silhouette must have a "
            "positive horizontal major axis"
        )
    center_x = left + width / 2.0
    center_y = top + height / 2.0
    tolerance = max(width, height, 1.0) * 1e-6

    def has_point(expected_x: float, expected_y: float) -> bool:
        return any(
            math.isclose(x, expected_x, abs_tol=tolerance)
            and math.isclose(y, expected_y, abs_tol=tolerance)
            for x, y in points
        )

    if not has_point(left, center_y) or not has_point(left + width, center_y):
        raise common.fail(
            f"{el.get('id') or 'arrow'}: leftRightArrow silhouette lacks "
            "axis-aligned left/right tips"
        )
    if any(
        not has_point(2.0 * center_x - x, y)
        or not has_point(x, 2.0 * center_y - y)
        for x, y in points
    ):
        raise common.fail(
            f"{el.get('id') or 'arrow'}: leftRightArrow silhouette is not "
            "axis-aligned and bilaterally symmetric"
        )
    if any(
        el.get(f"data-{side}-head-type", "none") == "none"
        for side in ("start", "end")
    ):
        raise common.fail(
            f"{el.get('id') or 'arrow'}: leftRightArrow requires explicit "
            "start and end head semantics"
        )
    return {
        "kind": "straight",
        "coordinate_space": "canvas",
        "points": [
            {"x": left, "y": center_y},
            {"x": left + width, "y": center_y},
        ],
    }


def _explicit_block_centerline(el: ET.Element) -> dict[str, object] | None:
    """Parse an explicit canvas-space SVG path used only as semantic evidence.

    Arbitrary closed freeforms do not expose a trustworthy centerline from their
    silhouette alone.  Their author must therefore provide
    ``data-arrow-centerline="M ... L/C ..."``.  This path is never rendered as
    a second object; it is the deterministic ArrowSpec endpoint/tangent truth.
    """

    raw = el.get("data-arrow-centerline")
    if raw is None:
        return None
    try:
        return path_from_segments(parse_path_d(raw))
    except (IndexError, TypeError, ValueError) as exc:
        raise common.fail(
            f"{el.get('id') or 'arrow'}: invalid data-arrow-centerline"
        ) from exc


def _block_centerline(
    el: ET.Element,
    segments: list[tuple],
    autoshape: dict[str, object] | None,
) -> dict[str, object]:
    explicit = _explicit_block_centerline(el)
    derived = _left_right_block_centerline(el, segments, autoshape)
    if derived is not None:
        if explicit is not None and explicit != derived:
            raise common.fail(
                f"{el.get('id') or 'arrow'}: explicit block centerline drifts "
                "from the deterministic AutoShape centerline"
            )
        return derived
    if explicit is None:
        raise common.fail(
            f"{el.get('id') or 'arrow'}: single closed block freeform requires "
            "data-arrow-centerline"
        )
    return explicit


def _record_arrow_spec(
    ctx: ConvertContext,
    el: ET.Element,
    style: dict[str, str],
    segments: list[tuple],
) -> dict[str, object] | None:
    representation = el.get("data-arrow-representation")
    is_arrow = _has_explicit_edge_semantics(el)
    if not is_arrow:
        return None
    representation = representation or "line_arrow"
    if representation not in {"line_arrow", "block_arrow"}:
        raise common.fail(
            f"{el.get('id') or 'arrow'}: unsupported data-arrow-representation="
            f"{representation!r}"
        )
    autoshape = (
        _block_autoshape_spec(el, segments)
        if representation == "block_arrow"
        else None
    )
    normalized_path = (
        path_from_segments(segments)
        if representation == "line_arrow"
        else _block_centerline(el, segments, autoshape)
    )
    source_id = el.get("data-source-id")
    target_id = el.get("data-target-id")
    gap_source_id = el.get("data-gap-source-id")
    gap_target_id = el.get("data-gap-target-id")
    if (gap_source_id or gap_target_id) and (
        gap_source_id != source_id or gap_target_id != target_id
    ):
        raise common.fail(
            f"{el.get('id') or 'arrow'}: gap source/target must match ArrowSpec topology"
        )
    requested_routing = el.get("data-arrow-routing") or el.get("data-routing")
    requested_topology = el.get("data-arrow-topology")
    if requested_routing not in {None, "fixed", "host"}:
        raise common.fail(
            f"{el.get('id') or 'arrow'}: unsupported arrow routing={requested_routing!r}"
        )
    if requested_topology not in {None, "attached", "declared", "none"}:
        raise common.fail(
            f"{el.get('id') or 'arrow'}: unsupported arrow topology={requested_topology!r}"
        )
    host_routing = requested_routing == "host" or (
        requested_routing is None
        and representation == "line_arrow"
        and normalized_path is not None
        and normalized_path.get("kind") == "straight"
        and el.get("data-attach", "true") != "false"
    )
    if requested_topology == "declared":
        host_routing = False
    elif requested_topology == "attached":
        host_routing = True
    attached = bool(source_id and target_id and host_routing)
    topology_mode = (
        "attached"
        if attached
        else "declared"
        if source_id and target_id
        else "none"
    )
    body_color = (
        parse_color(style.get("fill"))
        if representation == "block_arrow"
        else parse_color(style.get("stroke"))
    )
    raw_width = (
        el.get("data-arrow-body-width")
        if representation == "block_arrow"
        else style.get("stroke-width", "1")
    ) or "1"
    try:
        width_px = float(raw_width)
    except ValueError as exc:
        raise common.fail(
            f"{el.get('id') or 'arrow'}: invalid stroke-width={raw_width!r}"
        ) from exc
    if not math.isfinite(width_px) or width_px <= 0:
        raise common.fail(
            f"{el.get('id') or 'arrow'}: stroke-width must be finite and positive"
        )
    try:
        dash = _dash_preset(
            style.get("stroke-dasharray"),
            width_px,
            style.get("stroke-linecap"),
            style.get("data-ppt-dash"),
        )
    except ValueError as exc:
        raise common.fail(f"{el.get('id') or 'arrow'}: {exc}") from exc
    source_evidence: dict[str, object] = {
        "input_route": ctx.input_route,
        "reference_sha256": ctx.reference_sha256,
        "reference_bbox": None,
        "confidence": None,
    }
    bbox = el.get("data-reference-bbox")
    if bbox:
        try:
            values = [float(item) for item in re.split(r"[\s,]+", bbox.strip()) if item]
            if len(values) == 4:
                source_evidence["reference_bbox"] = values
        except ValueError:
            pass
    confidence = el.get("data-inference-confidence")
    if confidence:
        try:
            parsed_confidence = float(confidence)
            if 0 <= parsed_confidence <= 1:
                source_evidence["confidence"] = parsed_confidence
        except ValueError:
            pass
    if ctx.input_route == "reference-only" and source_evidence["reference_bbox"] is None:
        coordinates: list[tuple[float, float]] = []
        for segment in segments:
            values = segment[1:]
            coordinates.extend(
                (float(values[index]), float(values[index + 1]))
                for index in range(0, len(values) - 1, 2)
            )
        if coordinates:
            margin = max(6.0, width_px * 4.0)
            min_x = max(0.0, min(point[0] for point in coordinates) - margin)
            min_y = max(0.0, min(point[1] for point in coordinates) - margin)
            max_x = min(float(ctx.canvas_width), max(point[0] for point in coordinates) + margin)
            max_y = min(float(ctx.canvas_height), max(point[1] for point in coordinates) + margin)
            source_evidence["reference_bbox"] = [
                round(min_x, 4),
                round(min_y, 4),
                round(max_x - min_x, 4),
                round(max_y - min_y, 4),
            ]
            source_evidence["bbox_basis"] = "candidate-centerline-expanded"
    if ctx.input_route == "reference-only" and source_evidence["confidence"] is None:
        source_evidence["confidence"] = 0.5
        source_evidence["confidence_basis"] = "unspecified-in-source"

    spec: dict[str, object] = {
        "schema_version": ARROW_SPEC_VERSION,
        "representation": representation,
        "path": normalized_path,
        "routing": "host" if attached else "fixed",
        "topology": {
            "mode": topology_mode,
            "source_id": source_id,
            "target_id": target_id,
            "source_site": int(el.get("data-source-site", "0")) if source_id else None,
            "target_site": int(el.get("data-target-site", "0")) if target_id else None,
        },
        "body": {
            "color": None if body_color is None else body_color[0],
            "width_px": width_px,
            "dash": style.get("data-ppt-dash")
            or semantic_dash_from_ooxml(dash, style.get("stroke-linecap")),
            "line_cap": (
                "round"
                if style.get("data-ppt-dash") == "round_dot"
                else "butt"
                if style.get("data-ppt-dash") == "square_dot"
                else style.get("stroke-linecap", "butt")
            ),
            "line_join": style.get("stroke-linejoin", "miter"),
        },
        "start_head": (
            _block_head_from_svg(el, style, "start")
            if representation == "block_arrow"
            else _arrow_head_from_svg(ctx, el, style, "start")
        ),
        "end_head": (
            _block_head_from_svg(el, style, "end")
            if representation == "block_arrow"
            else _arrow_head_from_svg(ctx, el, style, "end")
        ),
        "silhouette_path": (
            silhouette_from_segments(segments) if representation == "block_arrow" else None
        ),
        "autoshape": autoshape,
        "fallback_policy": "strict_fail",
        "single_visible_object": True,
        "source_evidence": source_evidence,
    }
    errors = validate_arrow_spec(
        spec,
        expected_input_route=ctx.input_route,
        expected_reference_sha256=ctx.reference_sha256,
    )
    if errors:
        raise common.fail(
            f"{el.get('id') or 'arrow'}: invalid ArrowSpec: {', '.join(errors)}"
        )
    strategy = compiler_strategy(spec)
    if strategy not in {
        "native-block-autoshape",
        "single-closed-freeform",
        "native-connector-line-end",
        "native-line-line-end",
        "native-freeform-line-end",
    }:
        raise common.fail(
            f"{el.get('id') or 'arrow'}: ArrowSpec cannot be compiled as exactly "
            "one visible PowerPoint object; use a supported native line end or "
            "provide one closed silhouette"
        )
    ctx.scene_elements[ctx.current_element_id]["arrow_spec"] = spec
    return spec


def _apply_native_markers(
    ctx: ConvertContext,
    sp_pr,
    el: ET.Element,
    style: dict[str, str],
) -> set[str]:
    line = sp_pr.find(qn("a:ln"))
    if line is None:
        return set()
    handled: set[str] = set()
    for side, tag in (("start", "a:headEnd"), ("end", "a:tailEnd")):
        marker_id = _marker_reference(el, side)
        if not marker_id:
            continue
        marker = ctx.defs.get(marker_id)
        spec = _native_marker_spec(marker, style, el, side) if marker is not None else None
        if spec is None:
            raise common.fail(
                f"{ctx.current_element_id}:{side}: marker {marker_id} cannot be represented "
                "as one native PowerPoint line end; grouped shaft-plus-head fallback "
                "is forbidden"
            )
        old = line.find(qn(tag))
        if old is not None:
            line.remove(old)
        line.append(line.makeelement(qn(tag), spec))
        handled.add(side)
        ctx.bump("native-arrowhead")
    return handled


def _emit_straight_connector(
    ctx: ConvertContext,
    el: ET.Element,
    style: dict[str, str],
    start: tuple[float, float],
    end: tuple[float, float],
) -> None:
    x1, y1 = start
    x2, y2 = end
    _record_arrow_spec(ctx, el, style, [("M", x1, y1), ("L", x2, y2)])
    connector = ctx.slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, _px(x1), _px(y1), _px(x2), _px(y2)
    )
    _disable_shadow(connector)
    _apply_line(connector._element.spPr, style)
    is_connector = _has_explicit_edge_semantics(el)
    ctx.register_shape(connector, "connector" if is_connector else "line")
    ctx.add_pending_connection(connector, el)
    _apply_native_markers(ctx, connector._element.spPr, el, style)
    ctx.bump("line")


def _emit_line(ctx: ConvertContext, el: ET.Element, style: dict[str, str], matrix: Matrix) -> None:
    start = matrix.apply(float(el.get("x1", 0)), float(el.get("y1", 0)))
    end = matrix.apply(float(el.get("x2", 0)), float(el.get("y2", 0)))
    _emit_straight_connector(ctx, el, style, start, end)


def _emit_block_autoshape(
    ctx: ConvertContext,
    spec: dict[str, object],
    style: dict[str, str],
) -> None:
    autoshape = spec.get("autoshape")
    if not isinstance(autoshape, dict):
        raise common.fail(f"{ctx.current_element_id}: missing block AutoShape contract")
    subtype = autoshape.get("subtype")
    shape_type = BLOCK_AUTOSHAPE_TYPES.get(str(subtype))
    bbox = autoshape.get("bbox")
    adjustments = autoshape.get("adjustments")
    if shape_type is None or not isinstance(bbox, list) or not isinstance(adjustments, list):
        raise common.fail(f"{ctx.current_element_id}: invalid block AutoShape contract")
    x, y, width, height = (float(value) for value in bbox)
    shape = ctx.slide.shapes.add_shape(
        shape_type,
        _px(x),
        _px(y),
        _px(width),
        _px(height),
    )
    if len(shape.adjustments) != len(adjustments):
        raise common.fail(
            f"{ctx.current_element_id}: PowerPoint AutoShape adjustment count drift"
        )
    for index, value in enumerate(adjustments):
        shape.adjustments[index] = float(value)
    _disable_shadow(shape)
    _apply_fill_and_line(shape._element.spPr, style, ctx)
    ctx.register_shape(shape, "block-arrow-autoshape")
    ctx.bump("block-arrow-autoshape")


def _emit_poly(ctx: ConvertContext, el: ET.Element, style: dict[str, str], matrix: Matrix, close: bool) -> None:
    points = [float(p) for p in re.split(r"[\s,]+", (el.get("points") or "").strip()) if p]
    if len(points) < 4:
        return
    segments: list[tuple] = []
    first = matrix.apply(points[0], points[1])
    segments.append(("M", *first))
    for i in range(2, len(points) - 1, 2):
        segments.append(("L", *matrix.apply(points[i], points[i + 1])))
    if close:
        segments.append(("Z",))
    spec = _record_arrow_spec(ctx, el, style, segments)
    if (
        close
        and isinstance(spec, dict)
        and spec.get("representation") == "block_arrow"
        and spec.get("autoshape") is not None
    ):
        _emit_block_autoshape(ctx, spec, style)
        return
    shape = _emit_freeform(
        ctx,
        segments,
        style,
        object_kind="freeform-arrow" if _has_explicit_edge_semantics(el) else "freeform",
    )
    if shape is None:
        return
    _apply_native_markers(ctx, shape.find(qn("p:spPr")), el, style)


def _emit_path(ctx: ConvertContext, el: ET.Element, style: dict[str, str], matrix: Matrix) -> None:
    d = el.get("d")
    if not d:
        return
    segments = parse_path_d(d)
    if matrix != Matrix():
        segments = _transform_segments(segments, matrix)
    _record_arrow_spec(ctx, el, style, segments)
    if len(segments) == 2 and segments[0][0] == "M" and segments[1][0] == "L":
        _emit_straight_connector(
            ctx,
            el,
            style,
            (segments[0][1], segments[0][2]),
            (segments[1][1], segments[1][2]),
        )
        return
    shape = _emit_freeform(
        ctx,
        segments,
        style,
        object_kind="freeform-arrow" if _has_explicit_edge_semantics(el) else "freeform",
    )
    if shape is None:
        return
    _apply_native_markers(ctx, shape.find(qn("p:spPr")), el, style)


def _transform_segments(segments: list[tuple], matrix: Matrix) -> list[tuple]:
    result: list[tuple] = []
    for seg in segments:
        if seg[0] == "M" or seg[0] == "L":
            result.append((seg[0], *matrix.apply(seg[1], seg[2])))
        elif seg[0] == "C":
            p1 = matrix.apply(seg[1], seg[2])
            p2 = matrix.apply(seg[3], seg[4])
            p3 = matrix.apply(seg[5], seg[6])
            result.append(("C", p1[0], p1[1], p2[0], p2[1], p3[0], p3[1]))
        else:
            result.append(seg)
    return result


def _build_freeform_sp(ctx: ConvertContext, segments: list[tuple]):
    """Build one custGeom freeform ``p:sp`` element without placing it.

    Returns ``(element, shape_id)``, or ``None`` when the segments carry no
    points.  The caller owns tree placement, fill/line, and registration, so
    the same geometry path serves top-level shapes and atomic-vector group
    members alike.
    """
    # bbox 必须包含贝塞尔控制点：控制多边形永远包住曲线本体，
    # 否则坐标越出声明的 path w/h，PowerPoint 会判定文件损坏。
    xs: list[float] = []
    ys: list[float] = []
    for s in segments:
        if s[0] in ("M", "L"):
            xs.append(s[1])
            ys.append(s[2])
        elif s[0] == "C":
            xs.extend((s[1], s[3], s[5]))
            ys.extend((s[2], s[4], s[6]))
    if not xs:
        return None
    min_x, max_x = min(xs), max(xs)
    min_y, max_y = min(ys), max(ys)
    w = max(max_x - min_x, 0.01)
    h = max(max_y - min_y, 0.01)

    def rel(x: float, y: float) -> tuple[int, int]:
        return (round((x - min_x) * EMU_PER_PX), round((y - min_y) * EMU_PER_PX))

    path_xml = [f'<a:path w="{round(w * EMU_PER_PX)}" h="{round(h * EMU_PER_PX)}">']
    for seg in segments:
        if seg[0] == "M":
            x, y = rel(seg[1], seg[2])
            path_xml.append(f'<a:moveTo><a:pt x="{x}" y="{y}"/></a:moveTo>')
        elif seg[0] == "L":
            x, y = rel(seg[1], seg[2])
            path_xml.append(f'<a:lnTo><a:pt x="{x}" y="{y}"/></a:lnTo>')
        elif seg[0] == "C":
            x1, y1 = rel(seg[1], seg[2])
            x2, y2 = rel(seg[3], seg[4])
            x3, y3 = rel(seg[5], seg[6])
            path_xml.append(
                f'<a:cubicBezTo><a:pt x="{x1}" y="{y1}"/><a:pt x="{x2}" y="{y2}"/>'
                f'<a:pt x="{x3}" y="{y3}"/></a:cubicBezTo>'
            )
        elif seg[0] == "Z":
            path_xml.append("<a:close/>")
    path_xml.append("</a:path>")

    shape_id = _next_shape_id(ctx.slide)
    xml = (
        '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:nvSpPr><p:cNvPr id="{shape_id}" name="freeform-{shape_id}"/>'
        "<p:cNvSpPr/><p:nvPr/></p:nvSpPr>"
        "<p:spPr>"
        f'<a:xfrm><a:off x="{round(min_x * EMU_PER_PX)}" y="{round(min_y * EMU_PER_PX)}"/>'
        f'<a:ext cx="{round(w * EMU_PER_PX)}" cy="{round(h * EMU_PER_PX)}"/></a:xfrm>'
        "<a:custGeom><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>"
        '<a:rect l="0" t="0" r="0" b="0"/>'
        f"<a:pathLst>{''.join(path_xml)}</a:pathLst></a:custGeom>"
        "</p:spPr></p:sp>"
    )
    return _parse_sp_xml(xml), shape_id


def _emit_freeform(
    ctx: ConvertContext,
    segments: list[tuple],
    style: dict[str, str],
    *,
    object_kind: str = "freeform",
    editable: bool = True,
):
    built = _build_freeform_sp(ctx, segments)
    if built is None:
        return None
    sp, shape_id = built
    ctx.slide.shapes._spTree.append(sp)
    sp_pr = sp.find(qn("p:spPr"))
    _apply_fill_and_line(sp_pr, style, ctx)
    ctx.register_xml(sp, shape_id, object_kind, editable=editable)
    ctx.bump("freeform")
    return sp


def _parse_sp_xml(xml: str):
    from lxml import etree

    return etree.fromstring(xml.encode("utf-8"))


def _attach_atomic_raster_tags(
    ctx: ConvertContext,
    picture,
    asset_id: str,
    source_sha256: str,
) -> None:
    """Persist the same hash-bound raster metadata read by powerpoint-live."""
    authorization = ctx.asset_authorizations.get(asset_id, {})
    if authorization.get("authorized") is not True:
        return
    tags = {
        "AISCIENTIFICILLUSTRATORASSETID": asset_id,
        "AISCIENTIFICILLUSTRATORSOURCESHA256": source_sha256,
        "AISCIENTIFICILLUSTRATORRASTERREASON": authorization.get("raster_reason")
        or "Reference-faithful irreducible creative microasset authorized by the user.",
        "AISCIENTIFICILLUSTRATORSOURCETIGHTLYCROPPED": "True",
        "AISCIENTIFICILLUSTRATORATOMICRASTERUNIT": "True",
        "AISCIENTIFICILLUSTRATORCONTAINSRECONSTRUCTABLECONTENT": "False",
        "AISCIENTIFICILLUSTRATORDECOMPOSITIONNOTE": authorization.get("decomposition_note")
        or "No faithful native decomposition; formal text, formulas, nodes, and topology remain native.",
    }
    tag_list = ET.Element(f"{{{PML_NS}}}tagLst")
    for name, value in tags.items():
        ET.SubElement(tag_list, f"{{{PML_NS}}}tag", {"name": name, "val": str(value)})
    tag_part = Part(
        ctx.slide.part.package.next_partname("/ppt/tags/tag%d.xml"),
        CT.PML_TAGS,
        ctx.slide.part.package,
        ET.tostring(tag_list, encoding="utf-8", xml_declaration=True),
    )
    relationship_id = ctx.slide.part.relate_to(tag_part, RT.TAGS)
    nv_pr = picture._element.find(f"{qn('p:nvPicPr')}/{qn('p:nvPr')}")
    if nv_pr is None:
        raise common.fail(f"{asset_id}: picture has no p:nvPr for raster metadata")
    old = nv_pr.find(qn("p:custDataLst"))
    if old is not None:
        nv_pr.remove(old)
    custom_data = OxmlElement("p:custDataLst")
    tag_reference = OxmlElement("p:tags")
    tag_reference.set(qn("r:id"), relationship_id)
    custom_data.append(tag_reference)
    nv_pr.append(custom_data)


def _attach_atomic_vector_tags(
    ctx: ConvertContext,
    group,
    asset_id: str,
    entry: dict,
    source_sha256: str,
) -> None:
    """Persist hash-bound atomic-vector metadata in the raster tag style."""
    tags = {
        "AISCIENTIFICILLUSTRATORASSETID": asset_id,
        "AISCIENTIFICILLUSTRATORSOURCESHA256": source_sha256,
        "AISCIENTIFICILLUSTRATORREPRESENTATION": "atomic-vector",
        "AISCIENTIFICILLUSTRATOREDITABLE": "True",
        "AISCIENTIFICILLUSTRATORORIGIN": "vtracer-provider",
        "AISCIENTIFICILLUSTRATORTRACEMETHOD": entry["trace_method"],
        "AISCIENTIFICILLUSTRATORTRACEENGINEVERSION": entry["trace_engine_version"],
    }
    tag_list = ET.Element(f"{{{PML_NS}}}tagLst")
    for name, value in tags.items():
        ET.SubElement(tag_list, f"{{{PML_NS}}}tag", {"name": name, "val": str(value)})
    tag_part = Part(
        ctx.slide.part.package.next_partname("/ppt/tags/tag%d.xml"),
        CT.PML_TAGS,
        ctx.slide.part.package,
        ET.tostring(tag_list, encoding="utf-8", xml_declaration=True),
    )
    relationship_id = ctx.slide.part.relate_to(tag_part, RT.TAGS)
    nv_pr = group.find(f"{qn('p:nvGrpSpPr')}/{qn('p:nvPr')}")
    if nv_pr is None:
        raise common.fail(f"{asset_id}: group has no p:nvPr for vector metadata")
    old = nv_pr.find(qn("p:custDataLst"))
    if old is not None:
        nv_pr.remove(old)
    custom_data = OxmlElement("p:custDataLst")
    tag_reference = OxmlElement("p:tags")
    tag_reference.set(qn("r:id"), relationship_id)
    custom_data.append(tag_reference)
    nv_pr.append(custom_data)


def _vector_fragment_view_box(fragment: ET.Element, asset_id: str) -> tuple[float, float, float, float]:
    """Return the fragment root geometry as ``(x, y, width, height)``."""
    parts = [
        float(part)
        for part in re.split(r"[\s,]+", (fragment.get("viewBox") or "").strip())
        if part
    ]
    if len(parts) == 4 and parts[2] > 0 and parts[3] > 0:
        return parts[0], parts[1], parts[2], parts[3]
    width = _svg_dimension(fragment.get("width"))
    height = _svg_dimension(fragment.get("height"))
    if width is None or height is None or width <= 0 or height <= 0:
        raise common.fail(f"{asset_id}: vector_source_svg 缺少有效 viewBox/尺寸")
    return 0.0, 0.0, width, height


def _emit_atomic_vector(
    ctx: ConvertContext,
    asset_id: str,
    entry: dict,
    x: float,
    y: float,
    w: float,
    h: float,
    matrix: Matrix,
) -> None:
    """Compile one authorized traced SVG fragment into a native freeform group.

    The fragment replaces the reference-crop picture for this atomic id:
    every path keeps its curve geometry as a custGeom member of a single
    PowerPoint group, mapped deterministically from the SVG root geometry
    onto the element bbox.  The declared ``vector_source_svg`` bytes are
    hash-verified and contract-subset-checked before any shape is emitted
    (fail closed).
    """
    from tools.assets.asset_spec import validate_atomic_vector_asset
    from tools.assets.asset_trace import check_svg_contract_subset

    errors = validate_atomic_vector_asset(entry)
    if errors:
        raise common.fail(f"{asset_id}: atomic-vector 资产条目无效: {', '.join(errors)}")
    if ctx.case_root is None:
        raise common.fail(f"{asset_id}: 缺少案例根目录，无法核验 vector_source_svg")
    declared = entry["vector_source_svg"]
    case_root = ctx.case_root.resolve()
    svg_path = (case_root / declared["path"]).resolve()
    if svg_path != case_root and case_root not in svg_path.parents:
        raise common.fail(f"{asset_id}: vector_source_svg 越出案例目录: {declared['path']}")
    if not svg_path.is_file():
        raise common.fail(f"{asset_id}: vector_source_svg 不存在: {declared['path']}")
    source_sha256 = common.sha256_file(svg_path)
    if source_sha256 != declared["sha256"]:
        raise common.fail(
            f"{asset_id}: vector_source_svg SHA-256 与 assets.json 登记不一致，已拒绝编译"
        )
    violations = check_svg_contract_subset(svg_path)
    if violations:
        raise common.fail(
            f"{asset_id}: vector_source_svg 违反 SVG 合同子集: {', '.join(violations)}"
        )
    try:
        fragment = ET.parse(svg_path).getroot()
    except ET.ParseError as exc:
        raise common.fail(f"{asset_id}: vector_source_svg 解析失败: {exc}") from exc

    vb_x, vb_y, vb_w, vb_h = _vector_fragment_view_box(fragment, asset_id)
    scale_x, scale_y = w / vb_w, h / vb_h
    placement = Matrix(a=scale_x, d=scale_y, e=x - vb_x * scale_x, f=y - vb_y * scale_y)
    members: list[tuple[list[tuple], dict[str, str]]] = []

    def collect(element: ET.Element, style: dict[str, str], fragment_matrix: Matrix) -> None:
        own_style = _element_style(element, style)
        own_matrix = fragment_matrix.multiply(parse_transform(element.get("transform")))
        if element.tag.rsplit("}", 1)[-1] == "path":
            d = (element.get("d") or "").strip()
            if d:
                members.append((_transform_segments(parse_path_d(d), own_matrix), own_style))
            return
        for child in element:
            collect(child, own_style, own_matrix)

    collect(fragment, {}, placement)
    if matrix != Matrix():
        members = [(_transform_segments(segments, matrix), style) for segments, style in members]
    if not members:
        raise common.fail(f"{asset_id}: vector_source_svg 不含任何 path 几何")

    group_id = _next_shape_id(ctx.slide)
    group = _parse_sp_xml(
        '<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:nvGrpSpPr><p:cNvPr id="{group_id}" name="group-{group_id}"/>'
        "<p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>"
        "<p:grpSpPr><a:xfrm>"
        '<a:off x="0" y="0"/><a:ext cx="0" cy="0"/>'
        '<a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/>'
        "</a:xfrm></p:grpSpPr></p:grpSp>"
    )
    ctx.slide.shapes._spTree.append(group)
    member_boxes: list[tuple[int, int, int, int]] = []
    for segments, style in members:
        built = _build_freeform_sp(ctx, segments)
        if built is None:
            continue
        member, _member_shape_id = built
        group.append(member)
        sp_pr = member.find(qn("p:spPr"))
        _apply_fill_and_line(sp_pr, style, ctx)
        xfrm = sp_pr.find(qn("a:xfrm"))
        off = xfrm.find(qn("a:off"))
        ext = xfrm.find(qn("a:ext"))
        member_boxes.append(
            (int(off.get("x")), int(off.get("y")), int(ext.get("cx")), int(ext.get("cy")))
        )
    if not member_boxes:
        ctx.slide.shapes._spTree.remove(group)
        raise common.fail(f"{asset_id}: vector_source_svg 未产生任何可编译几何")
    # 子坐标系与父坐标系 1:1（chOff/chExt == off/ext），成员保持画布绝对 EMU 坐标。
    min_x = min(box[0] for box in member_boxes)
    min_y = min(box[1] for box in member_boxes)
    max_x = max(box[0] + box[2] for box in member_boxes)
    max_y = max(box[1] + box[3] for box in member_boxes)
    group_xfrm = group.find(f"{qn('p:grpSpPr')}/{qn('a:xfrm')}")
    for tag, attributes in (
        ("a:off", {"x": str(min_x), "y": str(min_y)}),
        ("a:ext", {"cx": str(max_x - min_x), "cy": str(max_y - min_y)}),
        ("a:chOff", {"x": str(min_x), "y": str(min_y)}),
        ("a:chExt", {"cx": str(max_x - min_x), "cy": str(max_y - min_y)}),
    ):
        node = group_xfrm.find(qn(tag))
        for name, value in attributes.items():
            node.set(name, value)
    _attach_atomic_vector_tags(ctx, group, asset_id, entry, source_sha256)
    ctx.register_xml(group, group_id, "atomic-vector", editable=True)
    ctx.register_vector_asset(asset_id)
    ctx.bump("atomic-vector")


def _match_atomic_vector_entry(ctx: ConvertContext, element_id: str) -> dict | None:
    """Resolve the authorized vtracer-trace entry linked to one atomic element id.

    An entry matches when its id equals the element id, or when it is a
    ``vtracer-trace`` entry whose ``fallback_atomic_raster`` equals the
    element id (vector entry id ``atomic:<slug>-vector`` linked back to the
    raster entry ``atomic:<slug>``).  Both rules hitting different entries
    is ambiguous and fails closed instead of silently picking one.
    """
    from tools.assets.asset_spec import ATOMIC_VECTOR_SOURCE

    exact = ctx.asset_authorizations.get(element_id)
    if not (isinstance(exact, dict) and exact.get("source") == ATOMIC_VECTOR_SOURCE):
        exact = None
    linked = [
        item
        for item in ctx.asset_authorizations.values()
        if isinstance(item, dict)
        and item.get("source") == ATOMIC_VECTOR_SOURCE
        and item.get("fallback_atomic_raster") == element_id
    ]
    candidates = ([exact] if exact is not None else []) + [
        item for item in linked if item is not exact
    ]
    if not candidates:
        return None
    matched_ids = sorted({str(item.get("id")) for item in candidates})
    if len(matched_ids) > 1:
        raise common.fail(
            f"{element_id}: atomic-vector 条目匹配歧义（{', '.join(matched_ids)}），已拒绝编译"
        )
    return candidates[0]


def _emit_atomic(ctx: ConvertContext, element_id: str, x: float, y: float, w: float, h: float, matrix: Matrix) -> None:
    asset_id = ctx.current_element_id or element_id
    vector_entry = _match_atomic_vector_entry(ctx, asset_id)
    if vector_entry is not None:
        _emit_atomic_vector(ctx, vector_entry["id"], vector_entry, x, y, w, h, matrix)
        return
    if ctx.source_png is None:
        ctx.warn(f"{element_id}: 缺少参考图，无法裁剪嵌入")
        return
    from PIL import Image

    x0, y0 = matrix.apply(x, y)
    x1, y1 = matrix.apply(x + w, y + h)
    left, top = round(min(x0, x1)), round(min(y0, y1))
    right, bottom = round(max(x0, x1)), round(max(y0, y1))
    with Image.open(ctx.source_png) as image:
        crop = image.crop((left, top, right, bottom))
        buffer = io.BytesIO()
        crop.save(buffer, format="PNG")
        buffer.seek(0)
    source_sha256 = hashlib.sha256(buffer.getvalue()).hexdigest()
    picture = ctx.slide.shapes.add_picture(
        buffer,
        _px(left),
        _px(top),
        _px(right - left),
        _px(bottom - top),
    )
    _attach_atomic_raster_tags(ctx, picture, asset_id, source_sha256)
    ctx.register_shape(picture, "atomic-raster", editable=False)
    ctx.register_asset(
        asset_id,
        [left, top, right - left, bottom - top],
        source_sha256,
    )
    ctx.bump("atomic")


IMAGE_MAX_AREA_RATIO = 0.5  # 护栏：<image> 覆盖画布 ≥ 该比例即拒绝（防整图截图冒充矢量交付）


def _emit_image(ctx: ConvertContext, el: ET.Element, matrix: Matrix) -> None:
    """<image> 合同容错：不读内嵌数据，按 bbox 从参考图裁剪替代（等价 atomic 占位符）。"""
    x = float(el.get("x", 0))
    y = float(el.get("y", 0))
    w = float(el.get("width", 0))
    h = float(el.get("height", 0))
    if w <= 0 or h <= 0:
        ctx.warn("<image> 缺少有效 x/y/width/height，已跳过")
        return
    ratio = w * h / ctx.canvas_area
    if ratio >= IMAGE_MAX_AREA_RATIO:
        raise common.fail(
            f"<image> 覆盖画布 {ratio:.0%}（≥{IMAGE_MAX_AREA_RATIO:.0%}），疑似整图截图冒充矢量，已拒绝"
        )
    ctx.warn('<image> 已按 bbox 从参考图裁剪替代（建议改用 <rect id="atomic:*"> 占位符）')
    _emit_atomic(ctx, "<image>", x, y, w, h, matrix)


# ---------------------------------------------------------------- 文本

_CHAR_WIDTHS = {
    " ": 0.28, "i": 0.28, "l": 0.28, "j": 0.3, "t": 0.33, "f": 0.33, "r": 0.36,
    "m": 0.85, "w": 0.8, "I": 0.3, "M": 0.85, "W": 0.9,
    ".": 0.3, ",": 0.3, ":": 0.33, ";": 0.33, "'": 0.25, "(": 0.36, ")": 0.36,
    "-": 0.36, "…": 0.8, "τ": 0.5, "π": 0.55, "ƒ": 0.4,
}


def _estimate_width(text: str, font_size: float) -> float:
    total = 0.0
    for ch in text:
        if ch in _CHAR_WIDTHS:
            total += _CHAR_WIDTHS[ch]
        elif ch.isupper():
            total += 0.67
        elif ch.isdigit():
            total += 0.56
        elif ord(ch) >= 0x2E80:
            total += 1.0
        else:
            total += 0.52
    return total * font_size


def _tspan_style(tspan: ET.Element, base_style: dict[str, str]) -> dict[str, str]:
    style = dict(base_style)
    inline = _parse_style_attr(tspan.get("style"))
    for key in (
        "font-size",
        "font-family",
        "font-style",
        "font-weight",
        "fill",
        "text-decoration",
    ):
        if tspan.get(key) is not None:
            style[key] = tspan.get(key)
        if key in inline:
            style[key] = inline[key]
    return style


def _collect_text_lines(
    el: ET.Element,
    base_style: dict[str, str],
) -> list[dict[str, object]]:
    """Collect SVG text as baseline-bound editable lines.

    Positioned ``tspan`` elements are lines, while unpositioned spans remain
    inline runs (for formula superscripts/subscripts).  Older code flattened
    every tspan into one PowerPoint paragraph, which silently collapsed
    multi-line labels and made their font boxes collide with nearby objects.
    """

    lines: list[dict[str, object]] = []
    runs: list[dict] = []
    baseline = float(el.get("y", 0))

    def add_run(text: str, style: dict[str, str], shift: str | None) -> None:
        if text and text.strip():
            runs.append({"text": text, "style": dict(style), "shift": shift})

    def flush() -> None:
        nonlocal runs
        if runs:
            lines.append({"baseline": baseline, "runs": runs})
            runs = []

    if el.text:
        add_run(el.text, base_style, None)
    for tspan in el:
        if tspan.tag != f"{SVG_NS}tspan":
            continue
        inline = _parse_style_attr(tspan.get("style"))
        positioned = any(tspan.get(name) is not None for name in ("x", "y", "dy"))
        if positioned and runs:
            flush()
        if tspan.get("y") is not None:
            baseline = float(tspan.get("y", baseline))
        elif tspan.get("dy") is not None:
            baseline += float(tspan.get("dy", 0))
        style = _tspan_style(tspan, base_style)
        shift = tspan.get("baseline-shift") or inline.get("baseline-shift")
        if tspan.text:
            add_run(tspan.text, style, shift)
        if tspan.tail:
            add_run(tspan.tail, base_style, None)
    flush()
    return lines


def _collect_text_runs(el: ET.Element, base_style: dict[str, str]) -> list[dict]:
    """Return a compatibility flat run list for inline and formula text."""

    return [
        run
        for line in _collect_text_lines(el, base_style)
        for run in line["runs"]
    ]


def _stacked_character_runs(runs: list[dict]) -> list[dict]:
    """Expand inline SVG runs into one editable PowerPoint run per character.

    ``data-text-flow=stacked-characters`` is an explicit semantic contract,
    not a visual guess.  Whitespace only separates source runs; every visible
    Unicode code point becomes one paragraph in the saved presentation.
    """
    characters: list[dict] = []
    for run in runs:
        for character in run["text"]:
            if character.isspace():
                continue
            if unicodedata.combining(character) and characters:
                characters[-1]["text"] += character
                continue
            characters.append(
                {
                    "text": character,
                    "style": dict(run["style"]),
                    "shift": run["shift"],
                }
            )
    return characters


def _format_text_run(run, run_info: dict, default_font_size: float) -> None:
    run.text = run_info["text"]
    rstyle = run_info["style"]
    try:
        size = float(rstyle.get("font-size", default_font_size))
    except ValueError:
        size = default_font_size
    run.font.size = Pt(size * PT_PER_PX)
    family = (rstyle.get("font-family") or "Arial").split(",")[0].strip()
    run.font.name = family
    if str(rstyle.get("font-style", "")).lower() in ("italic", "oblique"):
        run.font.italic = True
    if str(rstyle.get("font-weight", "")).lower() in ("bold", "600", "700", "800", "900"):
        run.font.bold = True
    if "underline" in str(rstyle.get("text-decoration", "")).lower().split():
        run.font.underline = True
    color = parse_color(rstyle.get("fill", "#000000"))
    if color:
        run.font.color.rgb = RGBColor.from_string(color[0][1:])
    if run_info["shift"] in ("super", "sup"):
        run._r.get_or_add_rPr().set("baseline", "30000")
    elif run_info["shift"] == "sub":
        run._r.get_or_add_rPr().set("baseline", "-25000")


def _text_box_dimension(el: ET.Element, name: str, default: float) -> float:
    raw = el.get(name)
    if raw is None:
        return default
    try:
        value = float(raw)
    except ValueError as exc:
        raise common.fail(
            f"{el.get('id') or 'text'}: {name} must be a pixel value"
        ) from exc
    if not math.isfinite(value) or value <= 0:
        raise common.fail(f"{el.get('id') or 'text'}: {name} must be finite and positive")
    return value


def _layout_number(el: ET.Element, name: str, default: float) -> float:
    raw = el.get(name)
    if raw is None:
        return default
    try:
        value = float(raw)
    except ValueError as exc:
        raise common.fail(f"{el.get('id') or 'text'}: {name} 必须是像素数值") from exc
    if not math.isfinite(value) or value < 0:
        raise common.fail(f"{el.get('id') or 'text'}: {name} 必须是有限非负数")
    return value


def _constrain_text_box(
    ctx: ConvertContext,
    el: ET.Element,
    *,
    anchor: str,
    anchor_x: float,
    left: float,
    top: float,
    width: float,
    height: float,
    rotation_deg: float,
) -> tuple[float, float, float, float]:
    """Constrain an explicitly annotated text box to its SVG container.

    The anchor is preserved.  Only transparent selection-box surplus is
    clipped; if the declared anchor itself lies outside the container the
    source contract is invalid and conversion fails instead of silently moving
    visible content.
    """
    container_id = el.get("data-layout-container")
    if not container_id:
        return left, top, width, height
    container = ctx.layout_boxes.get(container_id)
    if container is None:
        raise common.fail(
            f"{el.get('id') or 'text'}: data-layout-container={container_id!r} "
            "未指向可测量的 SVG 形状"
        )
    if rotation_deg:
        raise common.fail(
            f"{el.get('id') or 'text'}: 旋转文字的容器约束必须先展平为明确 bbox"
        )
    padding = _layout_number(el, "data-layout-padding", 0.0)
    inner_left = container.x + padding
    inner_top = container.y + padding
    inner_right = container.right - padding
    inner_bottom = container.bottom - padding
    if inner_right <= inner_left or inner_bottom <= inner_top:
        raise common.fail(f"{el.get('id') or 'text'}: layout padding 抹空了容器 {container_id!r}")
    if not inner_left <= anchor_x <= inner_right:
        raise common.fail(
            f"{el.get('id') or 'text'}: 文字锚点 x={anchor_x:.3f} 位于容器 {container_id!r} 外"
        )
    original = (left, top, width, height)
    if anchor == "middle":
        width = min(width, 2.0 * min(anchor_x - inner_left, inner_right - anchor_x))
        left = anchor_x - width / 2.0
    elif anchor == "end":
        width = min(width, anchor_x - inner_left)
        left = anchor_x - width
    else:
        left = max(left, inner_left)
        width = min(width, inner_right - left)
    if top < inner_top:
        # A top shift would move the visible baseline and hide a source error.
        raise common.fail(
            f"{el.get('id') or 'text'}: 文本框上边界 {top:.3f} 越出容器 {container_id!r}"
        )
    height = min(height, inner_bottom - top)
    if width <= 0 or height <= 0:
        raise common.fail(
            f"{el.get('id') or 'text'}: 容器 {container_id!r} 内没有可用文本区域"
        )
    constrained = (left, top, width, height)
    if any(abs(before - after) > 0.01 for before, after in zip(original, constrained)):
        ctx.warn(
            f"{el.get('id') or 'text'}: transparent text box constrained to "
            f"{container_id} (padding={padding:g}px)"
        )
    return constrained


def _emit_text(ctx: ConvertContext, el: ET.Element, style: dict[str, str], matrix: Matrix) -> None:
    positioned_tspans = [
        child
        for child in el
        if child.tag == f"{SVG_NS}tspan"
        and any(child.get(name) is not None for name in ("x", "y", "dy"))
    ]
    x = float(
        el.get("x")
        or (positioned_tspans[0].get("x") if positioned_tspans else None)
        or 0
    )
    y = float(
        el.get("y")
        or (positioned_tspans[0].get("y") if positioned_tspans else None)
        or 0
    )
    try:
        font_size = float(style.get("font-size", "16"))
    except ValueError:
        font_size = 16.0
    text_lines = _collect_text_lines(el, style)
    runs = [run for line in text_lines for run in line["runs"]]
    declared_text_flow = el.get("data-text-flow")
    text_flow = declared_text_flow or ("multiline" if len(text_lines) > 1 else "inline")
    if text_flow not in {"inline", "multiline", "rotated-word", "stacked-characters"}:
        raise common.fail(
            f"{el.get('id') or 'text'}: unsupported data-text-flow={text_flow!r}"
        )
    if text_flow in {"rotated-word", "stacked-characters"}:
        container_id = el.get("data-text-container")
        if not container_id or container_id not in ctx.layout_boxes:
            raise common.fail(
                f"{el.get('id') or 'text'}: {text_flow} requires a measurable "
                "data-text-container"
            )
    if text_flow == "stacked-characters" and el.get("data-text-stack-step") is None:
        raise common.fail(
            f"{el.get('id') or 'text'}: stacked-characters requires "
            "data-text-stack-step"
        )
    stacked_runs = _stacked_character_runs(runs) if text_flow == "stacked-characters" else []
    if text_flow == "stacked-characters" and len(stacked_runs) < 2:
        raise common.fail(
            f"{el.get('id') or 'text'}: stacked-characters requires at least two visible characters"
        )
    if text_flow == "multiline" and len(text_lines) < 2:
        raise common.fail(
            f"{el.get('id') or 'text'}: multiline requires at least two positioned text lines"
        )
    full_text = "".join(r["text"] for r in runs)
    if not full_text.strip():
        return

    anchor = style.get("text-anchor", "start")
    if text_flow == "stacked-characters":
        text_w = max(
            max(_estimate_width(item["text"], font_size) for item in stacked_runs),
            1.0,
        )
        stack_step = _layout_number(el, "data-text-stack-step", font_size * 1.15)
        if stack_step <= 0:
            raise common.fail(
                f"{el.get('id') or 'text'}: data-text-stack-step must be positive"
            )
        default_box_height = (
            font_size * 1.25
            + stack_step * (len(stacked_runs) - 1)
            + TEXT_BOX_PAD_Y
        )
    elif text_flow == "multiline":
        text_w = max(
            sum(
                _estimate_width(str(run["text"]), float(run["style"].get("font-size", font_size)))
                for run in line["runs"]
            )
            for line in text_lines
        )
        baselines = [float(line["baseline"]) for line in text_lines]
        if any(right <= left for left, right in zip(baselines, baselines[1:], strict=False)):
            raise common.fail(
                f"{el.get('id') or 'text'}: multiline baselines must increase from top to bottom"
            )
        stack_step = baselines[1] - baselines[0]
        default_box_height = (
            baselines[-1] - baselines[0] + font_size * 1.25 + TEXT_BOX_PAD_Y
        )
    else:
        text_w = max(_estimate_width(full_text, font_size), 1.0)
        stack_step = 0.0
        default_box_height = font_size * 1.25 + TEXT_BOX_PAD_Y
    box_w = _text_box_dimension(el, "data-text-box-width", text_w + TEXT_BOX_PAD_X)
    box_h = _text_box_dimension(
        el,
        "data-text-box-height",
        default_box_height,
    )

    rotation_deg = 0.0
    rot_match = re.fullmatch(r"rotate\(\s*(-?[\d.]+)[\s,]+(-?[\d.]+)[\s,]+(-?[\d.]+)\s*\)", el.get("transform", "").strip())
    if rot_match:
        if text_flow == "stacked-characters":
            raise common.fail(
                f"{el.get('id') or 'text'}: stacked-characters must use upright glyphs, not rotation"
            )
        rotation_deg = float(rot_match.group(1))
        if text_flow == "rotated-word" and not math.isclose(
            abs(rotation_deg) % 180.0,
            90.0,
            abs_tol=1e-6,
        ):
            raise common.fail(
                f"{el.get('id') or 'text'}: rotated-word requires a 90 degree rotation"
            )
        cx, cy = float(rot_match.group(2)), float(rot_match.group(3))
        cx, cy = matrix.apply(cx, cy)
        box_w = min(box_w, max(2 * min(cx, ctx.canvas_width - cx), 0.5))
        box_h = min(box_h, max(2 * min(cy, ctx.canvas_height - cy), 0.5))
        left, top = cx - box_w / 2, cy - box_h / 2
        align = PP_ALIGN.CENTER
        vertical = MSO_ANCHOR.MIDDLE
    else:
        if text_flow == "rotated-word":
            raise common.fail(
                f"{el.get('id') or 'text'}: rotated-word requires an explicit rotate(90 cx cy) transform"
            )
        x, y = matrix.apply(x, y)
        if anchor == "middle":
            box_w = min(box_w, max(2 * min(x, ctx.canvas_width - x), 0.5))
            left, align = x - box_w / 2, PP_ALIGN.CENTER
        elif anchor == "end":
            box_w = min(box_w, max(x, 0.5))
            left, align = x - box_w, PP_ALIGN.RIGHT
        else:
            box_w = min(box_w, max(ctx.canvas_width - x, 0.5))
            left, align = x, PP_ALIGN.LEFT
        top = y - font_size * BASELINE_ASCENT
        box_h = min(box_h, max(ctx.canvas_height - top, 0.5))
        vertical = MSO_ANCHOR.TOP

    left, top, box_w, box_h = _constrain_text_box(
        ctx,
        el,
        anchor=anchor,
        anchor_x=x if not rotation_deg else left + box_w / 2.0,
        left=left,
        top=top,
        width=box_w,
        height=box_h,
        rotation_deg=rotation_deg,
    )

    textbox = ctx.slide.shapes.add_textbox(_px(left), _px(top), _px(box_w), _px(box_h))
    _disable_shadow(textbox)
    if rotation_deg:
        textbox.rotation = rotation_deg
    frame = textbox.text_frame
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.word_wrap = False
    frame.vertical_anchor = vertical
    for margin in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(frame, margin, 0)
    if text_flow == "stacked-characters":
        for index, run_info in enumerate(stacked_runs):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.alignment = align
            paragraph.space_before = Pt(0)
            paragraph.space_after = Pt(0)
            paragraph.line_spacing = Pt(stack_step * PT_PER_PX)
            _format_text_run(paragraph.add_run(), run_info, font_size)
    elif text_flow == "multiline":
        baselines = [float(line["baseline"]) for line in text_lines]
        steps = [
            baselines[index + 1] - baselines[index]
            for index in range(len(baselines) - 1)
        ]
        for index, line in enumerate(text_lines):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.alignment = align
            paragraph.space_before = Pt(0)
            paragraph.space_after = Pt(0)
            step = steps[index] if index < len(steps) else steps[-1]
            paragraph.line_spacing = Pt(step * PT_PER_PX)
            for run_info in line["runs"]:
                _format_text_run(paragraph.add_run(), run_info, font_size)
    else:
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        for run_info in runs:
            _format_text_run(paragraph.add_run(), run_info, font_size)
    ctx.register_shape(textbox, "text")
    ctx.bump("text")


# ---------------------------------------------------------------- 树遍历与主流程


def _emit_bound(ctx: ConvertContext, element: ET.Element, tag: str, emitter, *args) -> None:
    matrix = next((arg for arg in reversed(args) if isinstance(arg, Matrix)), Matrix())
    computed_style = next((arg for arg in args if isinstance(arg, dict)), None)
    previous = ctx.begin_element(element, tag, matrix, computed_style)
    try:
        emitter(ctx, element, *args)
    finally:
        ctx.end_element(previous)


def _walk(ctx: ConvertContext, element: ET.Element, style: dict[str, str], matrix: Matrix) -> None:
    tag = element.tag
    if tag in (f"{SVG_NS}defs", f"{SVG_NS}marker", f"{SVG_NS}linearGradient", f"{SVG_NS}radialGradient"):
        return
    own_style = _element_style(element, style)
    own_matrix = matrix.multiply(parse_transform(element.get("transform")))
    if tag == f"{SVG_NS}svg":
        for child in element:
            _walk(ctx, child, own_style, own_matrix)
    elif tag == f"{SVG_NS}g":
        token = ctx.begin_logical_group(element, own_style, own_matrix)
        try:
            for child in element:
                _walk(ctx, child, own_style, own_matrix)
        finally:
            ctx.end_logical_group(token)
    elif tag == f"{SVG_NS}rect":
        _emit_bound(ctx, element, "rect", _emit_rect, own_style, own_matrix)
    elif tag in (f"{SVG_NS}circle", f"{SVG_NS}ellipse"):
        _emit_bound(ctx, element, tag.replace(SVG_NS, ""), _emit_ellipse, own_style, own_matrix)
    elif tag == f"{SVG_NS}line":
        _emit_bound(ctx, element, "line", _emit_line, own_style, own_matrix)
    elif tag == f"{SVG_NS}polyline":
        _emit_bound(ctx, element, "polyline", _emit_poly, own_style, own_matrix, False)
    elif tag == f"{SVG_NS}polygon":
        _emit_bound(ctx, element, "polygon", _emit_poly, own_style, own_matrix, True)
    elif tag == f"{SVG_NS}path":
        _emit_bound(ctx, element, "path", _emit_path, own_style, own_matrix)
    elif tag == f"{SVG_NS}text":
        _emit_bound(ctx, element, "text", _emit_text, own_style, own_matrix)
    elif tag == f"{SVG_NS}image":
        _emit_bound(ctx, element, "image", _emit_image, own_matrix)
    else:
        ctx.warn(f"未知元素已跳过: {tag}")


def _svg_dimension(value: str | None) -> float | None:
    if value is None:
        return None
    match = re.fullmatch(r"\s*([-+]?(?:\d*\.\d+|\d+\.?))(?:px)?\s*", value)
    return float(match.group(1)) if match else None


def _iter_readback_shapes(shapes):
    for shape in shapes:
        yield shape
        children = getattr(shape, "shapes", None)
        if children is not None:
            yield from _iter_readback_shapes(children)


def _shape_description(shape) -> dict[str, object]:
    c_nv_pr = shape._element.find(f".//{qn('p:cNvPr')}")
    if c_nv_pr is None:
        return {}
    try:
        parsed = json.loads(c_nv_pr.get("descr") or "{}")
    except json.JSONDecodeError:
        return {}
    return parsed if isinstance(parsed, dict) else {}


def _write_conversion_contracts(run: common.Run, ctx: ConvertContext, reopened) -> dict:
    from tools.core.contracts import read_json, utc_now, write_json

    pptx_hash = common.sha256_file(run.pptx_path)
    readback_shapes = list(_iter_readback_shapes(reopened.slides[0].shapes))
    readback_identities = {
        (int(shape.shape_id), shape.name) for shape in readback_shapes
    }
    readback_by_identity = {
        (int(shape.shape_id), shape.name): shape for shape in readback_shapes
    }
    for binding in ctx.bindings:
        identity = (
            int(binding["shape_id"]), binding["shape_name"]
        )
        binding["readback_found"] = identity in readback_identities
        topology_relation = binding.get("topology_relation")
        brace_spec_hash = binding.get("brace_spec_sha256")
        primitive_spec_hash = binding.get("primitive_spec_sha256")
        asset_group_id = binding.get("asset_group_id")
        asset_spec_hash = binding.get("asset_spec_sha256")
        embedded: dict[str, object] = {}
        if (
            isinstance(topology_relation, dict)
            or isinstance(brace_spec_hash, str)
            or isinstance(primitive_spec_hash, str)
            or isinstance(asset_group_id, str)
            or isinstance(asset_spec_hash, str)
        ):
            shape = readback_by_identity.get(identity)
            if shape is not None:
                embedded = _shape_description(shape)
        if isinstance(topology_relation, dict):
            binding["topology_relation_readback_found"] = (
                embedded.get("topology_relation") == topology_relation
            )
        if isinstance(brace_spec_hash, str):
            binding["brace_spec_readback_found"] = (
                embedded.get("brace_spec_sha256") == brace_spec_hash
            )
        if isinstance(primitive_spec_hash, str):
            binding["primitive_spec_readback_found"] = (
                embedded.get("primitive_spec_sha256") == primitive_spec_hash
            )
        if isinstance(asset_group_id, str) and isinstance(asset_spec_hash, str):
            binding["asset_spec_readback_found"] = (
                embedded.get("asset_group_id") == asset_group_id
                and embedded.get("asset_spec_sha256") == asset_spec_hash
            )
    for binding in ctx.logical_group_bindings:
        backend_ids = binding.get("backend_object_ids")
        backend_names = binding.get("backend_object_names")
        well_formed = (
            isinstance(backend_ids, list)
            and isinstance(backend_names, list)
            and bool(backend_ids)
            and len(backend_ids) == len(backend_names)
        )
        identities = (
            {
                (int(shape_id), shape_name)
                for shape_id, shape_name in zip(
                    backend_ids, backend_names, strict=True
                )
                if isinstance(shape_id, int)
                and not isinstance(shape_id, bool)
                and isinstance(shape_name, str)
                and shape_name
            }
            if well_formed
            else set()
        )
        binding["readback_found"] = bool(identities) and len(identities) == len(
            backend_ids
        ) and identities.issubset(readback_identities)
        asset_spec_hash = binding.get("asset_spec_sha256")
        if isinstance(asset_spec_hash, str):
            asset_group_id = binding.get("element_id")
            binding["asset_spec_readback_found"] = bool(identities) and all(
                _shape_description(readback_by_identity[identity]).get(
                    "asset_group_id"
                )
                == asset_group_id
                and _shape_description(readback_by_identity[identity]).get(
                    "asset_spec_sha256"
                )
                == asset_spec_hash
                for identity in identities
                if identity in readback_by_identity
            ) and identities.issubset(readback_identities)

    scene = read_json(run.scene_path)
    carrier = scene.get("canonical_svg")
    source_role = carrier.get("source_role") if isinstance(carrier, dict) else None
    # A fresh reference-only reconstruction is an isolation boundary: generated
    # elements must not inherit annotations or implementation hints from a
    # previous candidate that happened to reuse the same stable IDs. Repair
    # candidates retain the merge behavior because their explicit purpose is to
    # refine the current canonical scene.
    fresh_reference_reconstruction = (
        ctx.input_route == "reference-only"
        and source_role == "reconstruction-candidate"
    )
    existing_elements = (
        {}
        if fresh_reference_reconstruction
        else {item["id"]: item for item in scene.get("elements", [])}
    )
    merged_elements = []
    for element_id, generated in ctx.scene_elements.items():
        existing = existing_elements.pop(element_id, {})
        annotations = {
            field: dict(value)
            for field in SCENE_ELEMENT_ANNOTATION_FIELDS
            if isinstance((value := existing.get(field)), dict)
        }
        # The generated record replaces every compiler-owned field.  In
        # particular, omission is meaningful: when the current SVG no longer
        # declares an edge, a prior ArrowSpec/topology/backend record must not
        # survive merely because the stable element ID was reused.
        merged_elements.append({**annotations, **generated})
    # The current SVG is the complete offline scene carrier.  Retain manual
    # annotations only for IDs that still exist; carrying unmatched historical
    # elements creates stale scene objects with no PowerPoint binding.
    merged_elements.sort(key=lambda item: (int(item.get("z_index", 0)), item["id"]))
    scene.update(
        {
            "updated_at": utc_now(),
            "elements": merged_elements,
            "edges": [
                {
                    "id": item["id"],
                    **item.get("topology", {}),
                    **(
                        {"arrow_spec": item["arrow_spec"]}
                        if isinstance(item.get("arrow_spec"), dict)
                        else {}
                    ),
                }
                for item in merged_elements
                if item.get("kind") == "edge"
                or isinstance(item.get("arrow_spec"), dict)
            ],
            "artifact": {
                "backend": "pptx-offline",
                "path": "redraw.pptx",
                "sha256": pptx_hash,
            },
        }
    )
    write_json(run.scene_path, scene)

    assets = read_json(run.assets_path)
    existing_assets = {item["id"]: item for item in assets.get("assets", [])}
    generated_assets = []
    for generated in ctx.assets:
        existing = existing_assets.pop(generated["id"], {})
        # editable 以编译器生成记录为准：atomic-raster 恒为 False（位图层），
        # atomic-vector 恒为 True（原生 group），既有条目的其余字段原样保留。
        generated_assets.append({**generated, **existing, "editable": generated.get("editable", False)})
    generated_assets.extend(existing_assets.values())
    assets.update({"updated_at": utc_now(), "assets": generated_assets})
    write_json(run.assets_path, assets)

    bindings_complete = (
        bool(ctx.bindings)
        and all(binding["readback_found"] for binding in ctx.bindings)
        and all(
            binding.get("topology_relation_readback_found", True)
            for binding in ctx.bindings
        )
        and all(
            binding.get("brace_spec_readback_found", True)
            for binding in ctx.bindings
        )
        and all(
            binding.get("primitive_spec_readback_found", True)
            for binding in ctx.bindings
        )
        and all(
            binding.get("asset_spec_readback_found", True)
            for binding in ctx.bindings
        )
        and all(
            binding["readback_found"] for binding in ctx.logical_group_bindings
        )
        and all(
            binding.get("asset_spec_readback_found", True)
            for binding in ctx.logical_group_bindings
        )
    )
    bindings = read_json(run.bindings_path)
    bindings.update(
        {
            "updated_at": utc_now(),
            "backend": "pptx-offline",
            "artifact_sha256": pptx_hash,
            "package_reopened": True,
            "saved_reopened": False,
            "bindings_complete": bindings_complete,
            "bindings": ctx.bindings,
            "logical_group_bindings": ctx.logical_group_bindings,
        }
    )
    write_json(run.bindings_path, bindings)
    return {
        "pptx_sha256": pptx_hash,
        "object_count": len(readback_shapes),
        "bindings_complete": bindings_complete,
    }


def write_asset_spec_audit(run: common.Run) -> dict:
    """Regenerate the frozen-contract and saved-PPTX AssetSpec audit.

    The report never infers an asset.  With no opportunity map it is a strict
    no-op unless stale AssetSpec projections are already present.  When a map
    exists, semantic truth comes exclusively from ``audit_asset_specs`` and
    this wrapper closes its declared hashes over bindings and OOXML readback.
    """

    from tools.assets.asset_spec import (
        AssetSpecError,
        asset_contract_sha256,
        asset_spec_sha256,
        audit_asset_specs,
        canonical_asset_contract_payload,
    )
    from tools.core.contracts import read_json, write_json
    from tools.assets.reference_inventory import canonical_sha256
    from tools.core.revisions import compiler_fingerprint, revision_id, scene_sha256

    scene = read_json(run.scene_path)
    assets = read_json(run.assets_path)
    regions = read_json(run.regions_path)
    frozen_inventory = regions.get("reference_inventory")
    inventory = frozen_inventory if isinstance(frozen_inventory, dict) else {}
    blockers = list(audit_asset_specs(scene, assets, inventory))
    has_opportunity_map = "microasset_opportunity_map" in assets
    frozen_asset_contract = None
    frozen_asset_contract_sha256 = None
    if has_opportunity_map:
        try:
            frozen_asset_contract = canonical_asset_contract_payload(assets)
            frozen_asset_contract_sha256 = asset_contract_sha256(assets)
        except AssetSpecError as exc:
            blockers.extend(exc.errors)

    collections = [
        value
        for value in (scene.get("elements"), scene.get("nodes"))
        if isinstance(value, list)
    ]
    elements = [item for values in collections for item in values if isinstance(item, dict)]
    spec_by_group = {
        element["id"]: element["asset_spec"]
        for element in elements
        if isinstance(element.get("id"), str)
        and isinstance(element.get("asset_spec"), dict)
    }
    digest_by_group: dict[str, str] = {}
    member_owner: dict[str, str] = {}
    for asset_group_id, spec in sorted(spec_by_group.items()):
        try:
            digest_by_group[asset_group_id] = asset_spec_sha256(spec)
        except (TypeError, ValueError):
            continue
        member_ids = spec.get("member_ids", [])
        if not isinstance(member_ids, list):
            continue
        for member_id in member_ids:
            if not isinstance(member_id, str):
                continue
            existing_owner = member_owner.get(member_id)
            if existing_owner is not None and existing_owner != asset_group_id:
                blockers.append(
                    "asset-spec-member-multiple-assets:"
                    f"{member_id}:{existing_owner}:{asset_group_id}"
                )
            else:
                member_owner[member_id] = asset_group_id

    bindings_document = read_json(run.bindings_path)
    physical_bindings = (
        bindings_document.get("bindings", [])
        if isinstance(bindings_document.get("bindings"), list)
        else []
    )
    logical_bindings = (
        bindings_document.get("logical_group_bindings", [])
        if isinstance(bindings_document.get("logical_group_bindings"), list)
        else []
    )
    for asset_group_id, spec in sorted(spec_by_group.items()):
        digest = digest_by_group.get(asset_group_id)
        rows = [
            row
            for row in logical_bindings
            if isinstance(row, dict) and row.get("element_id") == asset_group_id
        ]
        if len(rows) != 1:
            blockers.append(
                f"asset-spec-binding-logical-count:{asset_group_id}:{len(rows)}"
            )
            continue
        row = rows[0]
        if row.get("asset_spec") != spec:
            blockers.append(f"asset-spec-binding-logical-spec:{asset_group_id}")
        if row.get("asset_spec_sha256") != digest:
            blockers.append(f"asset-spec-binding-logical-hash:{asset_group_id}")
        if row.get("asset_spec_readback_found") is not True:
            blockers.append(f"asset-spec-binding-logical-readback:{asset_group_id}")
    for row in logical_bindings:
        if not isinstance(row, dict):
            continue
        element_id = row.get("element_id")
        carries_spec = "asset_spec" in row or "asset_spec_sha256" in row
        if carries_spec and element_id not in spec_by_group:
            blockers.append(f"asset-spec-binding-logical-orphan:{element_id}")

    expected_by_identity: dict[tuple[int, str], tuple[str, str, str]] = {}
    member_binding_count = 0
    for member_id, asset_group_id in sorted(member_owner.items()):
        digest = digest_by_group.get(asset_group_id)
        rows = [
            row
            for row in physical_bindings
            if isinstance(row, dict) and row.get("element_id") == member_id
        ]
        if not rows:
            blockers.append(
                f"asset-spec-binding-member-missing:{asset_group_id}:{member_id}"
            )
            continue
        for row in rows:
            member_binding_count += 1
            if row.get("asset_group_id") != asset_group_id:
                blockers.append(f"asset-spec-binding-member-group:{member_id}")
            if asset_group_id not in row.get("logical_group_ids", []):
                blockers.append(f"asset-spec-binding-member-scope:{member_id}")
            if row.get("asset_spec_sha256") != digest:
                blockers.append(f"asset-spec-binding-member-hash:{member_id}")
            if "asset_spec" in row:
                blockers.append(f"asset-spec-binding-member-expanded:{member_id}")
            if row.get("asset_spec_readback_found") is not True:
                blockers.append(f"asset-spec-binding-member-readback:{member_id}")
            shape_id = row.get("shape_id")
            shape_name = row.get("shape_name")
            if (
                not isinstance(shape_id, int)
                or isinstance(shape_id, bool)
                or not isinstance(shape_name, str)
                or not shape_name
            ):
                blockers.append(f"asset-spec-binding-member-identity:{member_id}")
                continue
            identity = (shape_id, shape_name)
            if identity in expected_by_identity:
                blockers.append(f"asset-spec-binding-identity-duplicate:{member_id}")
            expected_by_identity[identity] = (member_id, asset_group_id, digest or "")
    for row in physical_bindings:
        if not isinstance(row, dict):
            continue
        member_id = row.get("element_id")
        carries_spec = "asset_group_id" in row or "asset_spec_sha256" in row
        if carries_spec and member_id not in member_owner:
            blockers.append(f"asset-spec-binding-member-orphan:{member_id}")

    readback_count = 0
    if run.pptx_path.is_file():
        try:
            reopened = Presentation(run.pptx_path)
        except Exception:
            reopened = None
            blockers.append("asset-spec-readback-pptx-invalid")
        readback_shapes = (
            list(_iter_readback_shapes(reopened.slides[0].shapes))
            if reopened is not None and reopened.slides
            else []
        )
        shapes_by_identity = {
            (int(shape.shape_id), shape.name): shape for shape in readback_shapes
        }
        for identity, (member_id, asset_group_id, digest) in sorted(
            expected_by_identity.items()
        ):
            shape = shapes_by_identity.get(identity)
            if shape is None:
                blockers.append(f"asset-spec-readback-missing:{member_id}")
                continue
            description = _shape_description(shape)
            if (
                description.get("asset_group_id") != asset_group_id
                or description.get("asset_spec_sha256") != digest
            ):
                blockers.append(f"asset-spec-readback-hash:{member_id}")
                continue
            if "asset_spec" in description:
                blockers.append(f"asset-spec-readback-expanded:{member_id}")
                continue
            readback_count += 1
        for identity, shape in shapes_by_identity.items():
            description = _shape_description(shape)
            carries_spec = (
                "asset_group_id" in description
                or "asset_spec_sha256" in description
            )
            if carries_spec and identity not in expected_by_identity:
                blockers.append(f"asset-spec-readback-orphan:{shape.name}")
    elif spec_by_group:
        blockers.append("asset-spec-readback-pptx-missing")

    frozen_opportunities = (
        frozen_asset_contract["microasset_opportunity_map"]
        if isinstance(frozen_asset_contract, dict)
        else None
    )
    opportunity_count = (
        len(frozen_opportunities) if isinstance(frozen_opportunities, list) else 0
    )
    opportunity_map_sha256 = (
        canonical_sha256(frozen_opportunities)
        if isinstance(frozen_opportunities, list)
        else None
    )
    policy_sha256 = (
        canonical_sha256(frozen_asset_contract["policy"])
        if isinstance(frozen_asset_contract, dict)
        else None
    )
    blockers = list(dict.fromkeys(blockers))
    report = {
        "schema_version": "4.0.0",
        "kind": "asset_spec_audit",
        "case": run.load_meta()["case"],
        "reference_sha256": run.load_meta()["source_sha256"],
        "revision_id": revision_id(scene),
        "scene_sha256": scene_sha256(scene),
        "compiler_fingerprint": compiler_fingerprint(),
        "assets_sha256": common.sha256_file(run.assets_path),
        "regions_sha256": common.sha256_file(run.regions_path),
        "inventory_sha256": canonical_sha256(inventory),
        "bindings_sha256": common.sha256_file(run.bindings_path),
        "pptx_sha256": (
            common.sha256_file(run.pptx_path) if run.pptx_path.is_file() else None
        ),
        "asset_contract_sha256": frozen_asset_contract_sha256,
        "policy_sha256": policy_sha256,
        "microasset_opportunity_map_sha256": opportunity_map_sha256,
        "opportunity_count": opportunity_count,
        "asset_spec_count": len(spec_by_group),
        "logical_group_binding_count": sum(
            isinstance(row, dict)
            and ("asset_spec" in row or "asset_spec_sha256" in row)
            for row in logical_bindings
        ),
        "member_binding_count": member_binding_count,
        "pptx_readback_count": readback_count,
        "no_op": (
            not has_opportunity_map and not spec_by_group
        ),
        "findings": blockers,
        "blockers": blockers,
        "pass": not blockers,
    }
    write_json(run.qa_dir / "asset-spec-audit.json", report)
    return report


def _convert_in_place(run: common.Run) -> dict:
    from tools.core.contracts import initialize_contracts, read_json, transition
    from tools.core.revisions import bind_canonical_svg, materialize_svg, read_svg_text_exact

    meta = initialize_contracts(run)
    current_state = meta["workflow"]["state"]
    width, height = int(meta["width"]), int(meta["height"])

    # v4 construction truth lives in scene.json.  Legacy/tests may still place
    # a first carrier at redraw.svg; import it once, then all later rebuilds are
    # materialized from the scene-bound bytes.
    scene = read_json(run.scene_path)
    if isinstance(scene.get("canonical_svg"), dict):
        try:
            materialize_svg(run, scene)
        except ValueError as exc:
            raise common.fail(str(exc)) from exc
    elif run.redraw_svg.is_file():
        source_role = (
            "external-seed-normalized"
            if meta["processing_mode"] in {"svg_import", "svg_repair"}
            else "reference-reconstruction"
        )
        bind_canonical_svg(
            scene,
            read_svg_text_exact(run.redraw_svg),
            source_role=source_role,
        )
        from tools.core.contracts import write_json

        write_json(run.scene_path, scene)
        materialize_svg(run, scene)
    else:
        raise common.fail(
            f"scene.json 尚无 canonical_svg，且未找到兼容输入 {run.redraw_svg}；"
            "请先运行 autofigure ingest"
        )

    gate_path = run.source_gate_report_path
    if gate_path.is_file():
        gate = read_json(gate_path)
        decision = gate.get("decision")
        gate_role = gate.get("route_gate", {}).get("candidate_role")
        if decision == "reject" and meta["processing_mode"] != "png_reconstruct":
            raise common.fail("source gate rejected the seed; processing_mode must be png_reconstruct")
        if (
            decision == "accept"
            and gate_role == "external-seed"
            and meta["processing_mode"] != "svg_import"
        ):
            raise common.fail(
                "an accepted external seed must use processing_mode=svg_import"
            )
        if (
            decision == "repair"
            and meta["input_route"] == "svg-seeded"
            and meta["processing_mode"] != "svg_repair"
        ):
            raise common.fail("a repairable SVG source must use processing_mode=svg_repair")

    tree = ET.parse(run.redraw_svg)
    root = tree.getroot()
    required_view_box = root.get("viewBox", "")
    required_parts = [
        float(part) for part in re.split(r"[\s,]+", required_view_box.strip()) if part
    ]
    if len(required_parts) != 4:
        raise common.fail("SVG viewBox is required and must contain four numbers")
    if (
        abs(required_parts[0]) > 1e-6
        or abs(required_parts[1]) > 1e-6
        or round(required_parts[2]) != width
        or round(required_parts[3]) != height
    ):
        raise common.fail(
            f"SVG viewBox {required_parts} must be [0, 0, {width}, {height}] "
            "to preserve reference coordinates"
        )
    svg_width = _svg_dimension(root.get("width"))
    svg_height = _svg_dimension(root.get("height"))
    if svg_width is None or svg_height is None:
        raise common.fail("SVG root width/height must be numeric pixel dimensions")
    if round(svg_width) != width or round(svg_height) != height:
        raise common.fail(
            f"SVG root size {svg_width}x{svg_height} does not match reference {width}x{height}"
        )
    view_box = root.get("viewBox", "")
    if view_box:
        parts = [float(p) for p in re.split(r"[\s,]+", view_box.strip())]
        if len(parts) == 4 and (round(parts[2]) != width or round(parts[3]) != height):
            raise common.fail(
                f"SVG viewBox {parts[2]}x{parts[3]} 与参考图 {width}x{height} 不一致，违反输出合同"
            )

    prs = Presentation()
    prs.slide_width = _px(width)
    prs.slide_height = _px(height)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    defs = _collect_defs(root)
    source_png = run.source_png if run.source_png.is_file() else None
    assets_document = read_json(run.assets_path)
    asset_authorizations = {
        item["id"]: item for item in assets_document.get("assets", []) if item.get("id")
    }
    ctx = ConvertContext(
        slide,
        defs,
        source_png,
        width,
        height,
        asset_authorizations=asset_authorizations,
        layout_boxes=collect_svg_boxes(root),
        input_route=meta["input_route"],
        reference_sha256=meta["source_sha256"],
        case_root=run.root,
    )
    _walk(ctx, root, {}, Matrix())
    regions_document = read_json(run.regions_path)
    frozen_inventory = regions_document.get("reference_inventory")
    ctx.attach_asset_contracts(
        scene,
        assets_document,
        frozen_inventory if isinstance(frozen_inventory, dict) else {},
    )
    ctx.resolve_connections()

    # Conversion preflight is intentionally side-effect free.  In particular,
    # an unsupported marker must fail before changing workflow state or touching
    # the existing PPTX/bindings.  Only a fully compiled in-memory slide enters
    # the repairing state and proceeds to package serialization.
    if current_state in {"approved", "candidate", "qa_failed"}:
        transition(run, "repairing", "offline-conversion-started")

    run.qa_dir.mkdir(exist_ok=True)
    prs.save(run.pptx_path)

    # 读回统计（机械验收的基础）
    reopened = Presentation(run.pptx_path)
    readback_texts = 0
    for shape in _iter_readback_shapes(reopened.slides[0].shapes):
        if shape.has_text_frame and shape.text_frame.text.strip():
            readback_texts += 1
    contract_summary = _write_conversion_contracts(run, ctx, reopened)
    from tools.arrows.pptx_arrows import write_arrow_reports
    from tools.assets.primitives import audit_primitives
    from tools.providers.providers import write_case_capabilities

    arrow_compile_report, arrow_readback_report = write_arrow_reports(run)
    primitive_report = audit_primitives(run)
    provider_report = write_case_capabilities(run)
    from tools.pipeline.layout import audit_layout

    layout_report = audit_layout(run)
    from tools.core.revisions import stamp_active_revision

    revision = stamp_active_revision(run)
    asset_spec_report = write_asset_spec_audit(run)
    if not asset_spec_report["pass"]:
        raise common.fail(
            "AssetSpec conversion audit failed: "
            + ", ".join(asset_spec_report["blockers"])
        )
    summary = {
        "svg": "redraw.svg",
        "pptx": "redraw.pptx",
        "slide_count": len(reopened.slides._sldIdLst),
        "shape_count": len(reopened.slides[0].shapes),
        "textbox_with_text": readback_texts,
        "emitted": ctx.counts,
        "warnings": ctx.warnings,
        "layout_pass": layout_report["pass"],
        "layout_findings": len(layout_report["findings"]),
        "arrow_compile_pass": arrow_compile_report["pass"],
        "arrow_readback_pass": arrow_readback_report["pass"],
        "primitive_pass": primitive_report["pass"],
        "asset_spec_pass": asset_spec_report["pass"],
        "asset_spec_count": asset_spec_report["asset_spec_count"],
        "asset_spec_readback_count": asset_spec_report["pptx_readback_count"],
        "asset_spec_findings": len(asset_spec_report["findings"]),
        "arrow_authoring_allowed": provider_report["powerpoint_live"]["arrow_authoring_allowed"],
        "revision_id": revision["revision_id"],
        "scene_sha256": revision["scene_sha256"],
        "compiler_fingerprint": revision["compiler_fingerprint"],
        **contract_summary,
    }
    (run.qa_dir / "convert-summary.json").write_text(
        json.dumps(summary, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )
    transition(run, "candidate", "offline-conversion-complete", details=contract_summary)
    return summary


_CONVERT_ROOT_PUBLICATIONS = (
    "run.json",
    "scene.json",
    "assets.json",
    "regions.json",
    "bindings.json",
    "provenance.json",
    "redraw.svg",
    "redraw.pptx",
)

_CONVERT_QA_PUBLICATIONS = (
    "convert-summary.json",
    "arrow-compile-report.json",
    "powerpoint-arrow-readback.json",
    "arrow-composition-audit.json",
    "primitive-audit.json",
    "asset-spec-audit.json",
    "provider-capabilities.json",
    "layout-audit.json",
    "revision-receipt.json",
    "asset-contract-receipt.json",
)


def _convert_qa_publications(run: common.Run) -> tuple[str, ...]:
    """Return generated reports plus any pre-frozen asset receipt.

    Convert never manufactures the receipt.  Missing opportunity maps remain
    a no-op, while a receipt already frozen before conversion is preserved by
    the shadow publication transaction.
    """

    return tuple(
        name
        for name in _CONVERT_QA_PUBLICATIONS
        if name != "asset-contract-receipt.json"
        or (run.qa_dir / name).is_file()
    )


def _validate_shadow_conversion(run: common.Run, summary: dict) -> None:
    """Fail closed unless the shadow projection is complete and self-consistent."""

    from tools.core.contracts import read_json

    qa_publications = _convert_qa_publications(run)
    required = [
        *(run.root / name for name in _CONVERT_ROOT_PUBLICATIONS),
        *(run.qa_dir / name for name in qa_publications),
    ]
    missing = [path.name for path in required if not path.is_file()]
    if missing:
        raise common.fail(
            "shadow conversion did not produce required files: "
            + ", ".join(sorted(missing))
        )

    reopened = Presentation(run.pptx_path)
    if not reopened.slides:
        raise common.fail("shadow conversion produced a PowerPoint with no slides")

    disk_summary = read_json(run.qa_dir / "convert-summary.json")
    if disk_summary != summary:
        raise common.fail("shadow convert-summary does not match the compiler result")
    pptx_hash = common.sha256_file(run.pptx_path)
    if summary.get("pptx_sha256") != pptx_hash:
        raise common.fail("shadow PowerPoint hash does not match convert-summary")

    meta = run.load_meta()
    scene = read_json(run.scene_path)
    bindings = read_json(run.bindings_path)
    receipt = read_json(run.revision_receipt_path)
    active = meta.get("active_revision")
    scene_revision = scene.get("revision")
    binding_revision = bindings.get("scene_revision")
    if not all(isinstance(item, dict) for item in (active, scene_revision, binding_revision)):
        raise common.fail("shadow conversion is missing scene revision closure")
    for key in ("revision_id", "scene_sha256", "compiler_fingerprint"):
        values = {
            active.get(key),
            scene_revision.get(key),
            binding_revision.get(key),
            receipt.get(key),
            summary.get(key),
        }
        if len(values) != 1 or None in values:
            raise common.fail(f"shadow conversion revision mismatch: {key}")
    if active.get("artifacts", {}).get("redraw_pptx") != pptx_hash:
        raise common.fail("shadow active revision does not bind the PowerPoint bytes")

    asset_spec_audit = read_json(run.qa_dir / "asset-spec-audit.json")
    if asset_spec_audit.get("pass") is not True:
        raise common.fail("shadow AssetSpec audit did not pass")
    for key in ("revision_id", "scene_sha256", "compiler_fingerprint"):
        if asset_spec_audit.get(key) != summary.get(key):
            raise common.fail(f"shadow AssetSpec audit revision mismatch: {key}")
    expected_asset_hashes = {
        "assets_sha256": common.sha256_file(run.assets_path),
        "regions_sha256": common.sha256_file(run.regions_path),
        "bindings_sha256": common.sha256_file(run.bindings_path),
        "pptx_sha256": pptx_hash,
    }
    for key, digest in expected_asset_hashes.items():
        if asset_spec_audit.get(key) != digest:
            raise common.fail(f"shadow AssetSpec audit artifact mismatch: {key}")
    from tools.assets.asset_spec import asset_contract_sha256, canonical_asset_contract_payload

    assets_document = read_json(run.assets_path)
    expected_asset_contract_sha256 = (
        asset_contract_sha256(assets_document)
        if "microasset_opportunity_map" in assets_document
        else None
    )
    if (
        asset_spec_audit.get("asset_contract_sha256")
        != expected_asset_contract_sha256
    ):
        raise common.fail("shadow AssetSpec audit asset-contract hash mismatch")
    if "microasset_opportunity_map" in assets_document:
        from tools.assets.reference_inventory import canonical_sha256

        frozen_asset_contract = canonical_asset_contract_payload(assets_document)
        expected_contract_fields = {
            "policy_sha256": canonical_sha256(frozen_asset_contract["policy"]),
            "microasset_opportunity_map_sha256": canonical_sha256(
                frozen_asset_contract["microasset_opportunity_map"]
            ),
            "opportunity_count": len(
                frozen_asset_contract["microasset_opportunity_map"]
            ),
        }
        for key, value in expected_contract_fields.items():
            if asset_spec_audit.get(key) != value:
                raise common.fail(f"shadow AssetSpec audit contract mismatch: {key}")
    if asset_spec_audit.get("asset_spec_count") != summary.get("asset_spec_count"):
        raise common.fail("shadow AssetSpec count does not match convert-summary")
    if asset_spec_audit.get("pptx_readback_count") != summary.get(
        "asset_spec_readback_count"
    ):
        raise common.fail("shadow AssetSpec readback count does not match convert-summary")

    for name in qa_publications:
        read_json(run.qa_dir / name)


def convert(run: common.Run) -> dict:
    """Compile and validate in a shadow case, then atomically publish outputs.

    No compiler, audit, or revision writer receives the formal case path.  A
    failed build therefore cannot partially update ``run.json``, ``scene.json``,
    PowerPoint, bindings, or conversion QA.  The short publication phase keeps
    a recovery journal and restores all previous bytes if any replacement fails.
    """

    from tools.core.transactions import publish_staged_files, staged_case_copy

    staging_root = common.PROJECT_ROOT / ".autofigure-staging"
    label = f"convert-{run.root.name}"
    with staged_case_copy(run.root, staging_root=staging_root, label=label) as shadow_root:
        shadow = common.Run(shadow_root)
        summary = _convert_in_place(shadow)
        _validate_shadow_conversion(shadow, summary)
        qa_publications = _convert_qa_publications(shadow)
        publications = [
            *((shadow.root / name, run.root / name) for name in _CONVERT_ROOT_PUBLICATIONS),
            *((shadow.qa_dir / name, run.qa_dir / name) for name in qa_publications),
        ]
        publish_staged_files(
            publications,
            staging_root=staging_root,
            label=label,
        )
    return summary


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(prog="autofigure convert", description=__doc__)
    parser.add_argument("run_dir", type=Path, help="run 目录")
    parser.add_argument("--no-render", action="store_true", help="跳过 PowerPoint fresh render")
    args = parser.parse_args(argv)

    run = common.open_run(args.run_dir)
    summary = convert(run)
    sys.stdout.write(json.dumps(summary, ensure_ascii=False, indent=2) + "\n")
    if summary["warnings"]:
        sys.stdout.write(f"注意：{len(summary['warnings'])} 条 warning，详见 convert-summary.json\n")
    sys.stdout.write(f"PPTX 已生成: {run.pptx_path}\n")

    if not args.no_render:
        from tools.pipeline import render_export

        meta = run.load_meta()
        render_export.render(run.pptx_path, run.render_png, int(meta["width"]), int(meta["height"]))
        from tools.core.revisions import stamp_active_revision

        stamp_active_revision(run)
        sys.stdout.write(f"fresh render: {run.render_png}\n")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
