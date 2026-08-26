"""autofigure math — 把 convert 产出的公式文本框批量升级为原生 Office Math（OMML）。

检测（两路命中都收）：
- 强信号：任一 run 的 a:rPr@baseline ∈ {30000(super), -25000(sub)}（convert 对 baseline-shift 的落盘约定）
- 弱信号：全部 run 斜体、去空格后 1..4 字符、至少含一个非 ASCII 数学字母（τ π 𝔼 ℒ ∇ ƒ Î θ 等）；
  纯 ASCII 短标签（B / I / GT Answers / Questions）必须排除

流程：检测 → run 序列重建 LaTeX（baseline → ^{}/_{} 分组，Unicode → LaTeX 命令映射）→
逐公式 compile_formula（薄封装 legacy 引擎 tools/powerpoint_native_math.py，单个失败只
warn 跳过、保留原文本框）→ 在临时副本中把命中形状改名 math:NNN → inject_plan 到临时
pptx → os.replace 原子覆盖 redraw.pptx → 保存重开并刷新 v3 scene/bindings 哈希 → 刷新
fresh render（COM 失败只 warn）。
--dry-run 只检测与重建并写 qa/math-summary.json，不改 PPTX。
"""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import shutil
import sys
import unicodedata
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

from tools.core import common

BASELINE_SUPER = "30000"  # convert 对 baseline-shift="super" 的落盘值（tools/convert.py）
BASELINE_SUB = "-25000"  # baseline-shift="sub"
NAME_PREFIX = "math:"  # 命中形状的改名前缀（inject_plan 按 cNvPr@name 精确定位）
WEAK_MAX_CHARS = 4  # 弱信号：去空格后的最大字符数

RunInfo = tuple[str, bool, str | None]  # (text, italic, baseline)

# Unicode → LaTeX 映射。查不到的非 ASCII 字符保留原字符（latex2mathml 能吃 Unicode）。
_UNICODE_MAP = {
    # 希腊小写
    "α": r"\alpha", "β": r"\beta", "γ": r"\gamma", "δ": r"\delta",
    "ε": r"\varepsilon", "ζ": r"\zeta", "η": r"\eta", "θ": r"\theta",
    "ι": r"\iota", "κ": r"\kappa", "λ": r"\lambda", "μ": r"\mu",
    "ν": r"\nu", "ξ": r"\xi", "π": r"\pi", "ρ": r"\rho",
    "σ": r"\sigma", "τ": r"\tau", "υ": r"\upsilon", "φ": r"\phi",
    "χ": r"\chi", "ψ": r"\psi", "ω": r"\omega",
    # 希腊大写
    "Γ": r"\Gamma", "Δ": r"\Delta", "Θ": r"\Theta", "Λ": r"\Lambda",
    "Ξ": r"\Xi", "Π": r"\Pi", "Σ": r"\Sigma", "Υ": r"\Upsilon",
    "Φ": r"\Phi", "Ψ": r"\Psi", "Ω": r"\Omega",
    # 运算符与符号
    "∇": r"\nabla", "∂": r"\partial", "∞": r"\infty",
    "±": r"\pm", "×": r"\times", "÷": r"\div",
    "≤": r"\leq", "≥": r"\geq", "≠": r"\neq", "≈": r"\approx",
    "∈": r"\in", "∉": r"\notin", "⊂": r"\subset", "⊆": r"\subseteq",
    "∪": r"\cup", "∩": r"\cap", "·": r"\cdot",
    "∑": r"\sum", "∏": r"\prod", "∫": r"\int",
    "⟨": r"\langle", "⟩": r"\rangle", "…": r"\dots",
    "−": "-",  # U+2212 减号 → ASCII 连字符
    # 花体/双线/重音与异形字母
    "𝔼": r"\mathbb{E}", "ℒ": r"\mathcal{L}", "Î": r"\hat{I}", "ƒ": "f",
}

# LaTeX 特殊字符转义（公式框里 _ ^ 只应来自 baseline 结构；run 文本里出现就转义）
_LATEX_ESCAPES = {
    "\\": r"\backslash",
    "_": r"\_", "^": r"\^{}",
    "%": r"\%", "&": r"\&", "#": r"\#", "$": r"\$",
    "{": r"\{", "}": r"\}",
}


# ---------------------------------------------------------------- 纯函数：检测与 LaTeX 重建


def classify_runs(runs: list[RunInfo]) -> str | None:
    """→ 'strong'（含上下标 run）/ 'weak'（全斜体短公式）/ None（非公式）。"""
    meaningful = [(text, italic, baseline) for text, italic, baseline in runs if text.strip()]
    if not meaningful:
        return None
    if any(baseline in (BASELINE_SUPER, BASELINE_SUB) for _, _, baseline in meaningful):
        return "strong"
    text = "".join(text for text, _, _ in meaningful).replace(" ", "")
    if (
        len(text) <= WEAK_MAX_CHARS
        and all(italic for _, italic, _ in meaningful)
        and any(not ch.isascii() for ch in text)
    ):
        return "weak"
    return None


def _map_char(ch: str, nxt: str | None) -> str:
    """单字符 → LaTeX 片段；字母结尾的反斜杠命令后若紧跟 ASCII 字母需补空格（防命令名粘连）。"""
    token = _UNICODE_MAP.get(ch)
    if token is None:
        token = _LATEX_ESCAPES.get(ch, ch)
    if (
        token.startswith("\\")
        and token[-1:].isalpha()
        and token[-1:].isascii()
        and nxt
        and nxt.isalpha()
        and nxt.isascii()
    ):
        return token + " "
    return token


def _latex_text(text: str) -> str:
    return "".join(
        _map_char(ch, text[index + 1] if index + 1 < len(text) else None)
        for index, ch in enumerate(text)
    )


def rebuild_latex(runs: list[RunInfo]) -> str:
    """run 序列 → LaTeX：baseline run 进 ^{}/_{} 分组（连续同向合并进同一组），其余进主体。"""
    parts: list[str] = []
    open_group: str | None = None
    for text, _italic, baseline in runs:
        group = None
        if baseline == BASELINE_SUPER:
            group = "super"
        elif baseline == BASELINE_SUB:
            group = "sub"
        if group != open_group:
            if open_group is not None:
                parts.append("}")
            if group is not None:
                parts.append("^{" if group == "super" else "_{")
            open_group = group
        parts.append(_latex_text(text))
    if open_group is not None:
        parts.append("}")
    return "".join(parts)


# ---------------------------------------------------------------- pptx 侧提取


def _shape_runs(shape) -> list[RunInfo]:
    runs: list[RunInfo] = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            rpr = run._r.find(qn("a:rPr"))
            baseline = rpr.get("baseline") if rpr is not None else None
            runs.append((run.text, run.font.italic is True, baseline))
    return runs


def _target_font(shape) -> tuple[float, str] | None:
    """首个非空 run 的字号/颜色 → (pt, '#RRGGBB')（保持注入后视觉一致）。

    plan 校验要求 target_font_size_pt 与 target_font_color 成对出现且字号 6..72；
    任一拿不到或越界就都不给。
    """
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text.strip():
                continue
            if run.font.size is None:
                return None
            size = float(run.font.size.pt)
            try:
                rgb = run.font.color.rgb
            except AttributeError:  # 未显式设色（_NoneColor）
                return None
            if rgb is None or not 6.0 <= size <= 72.0:
                return None
            return size, f"#{rgb}"
    return None


def _load_engine():
    try:
        from tools import powerpoint_native_math
    except ModuleNotFoundError as exc:
        raise common.fail(
            f"缺少公式引擎依赖: {exc.name}（在 .venv 中 pip install -r requirements.txt）"
        ) from exc
    return powerpoint_native_math


def _echo(text: str) -> None:
    """stdout 跟随控制台代码页（GBK）：公式里的非 ASCII 字符防 UnicodeEncodeError。"""
    encoding = sys.stdout.encoding or "utf-8"
    sys.stdout.write(text.encode(encoding, errors="backslashreplace").decode(encoding) + "\n")


def _pptx_bound_names_and_omml_count(path: Path) -> tuple[set[str], int]:
    """Read names/OMML directly because python-pptx hides AlternateContent shapes."""
    names: set[str] = set()
    omml_count = 0
    with zipfile.ZipFile(path) as package:
        for member in package.namelist():
            if not member.startswith("ppt/slides/slide") or not member.endswith(".xml"):
                continue
            root = ET.fromstring(package.read(member))
            for node in root.iter():
                local_name = node.tag.rsplit("}", 1)[-1]
                if local_name == "cNvPr" and node.get("name"):
                    names.add(node.get("name"))
                elif local_name == "oMath":
                    omml_count += 1
    return names, omml_count


def _canonical_sha256(value: object) -> str:
    payload = json.dumps(
        value, ensure_ascii=False, sort_keys=True, separators=(",", ":")
    ).encode("utf-8")
    return hashlib.sha256(payload).hexdigest()


def _native_math_inventory(path: Path) -> tuple[list[dict], int]:
    """Read logical Office Math objects with branch-normalized semantics."""

    records: list[dict] = []
    raw_omml_count = 0
    with zipfile.ZipFile(path) as package:
        slide_entries = sorted(
            member
            for member in package.namelist()
            if member.startswith("ppt/slides/slide") and member.endswith(".xml")
        )
        for member in slide_entries:
            root = ET.fromstring(package.read(member))
            raw_omml_count += sum(
                node.tag.rsplit("}", 1)[-1] == "oMath" for node in root.iter()
            )
            for alternate in root.iter():
                if alternate.tag.rsplit("}", 1)[-1] != "AlternateContent":
                    continue
                formulas = [
                    node
                    for node in alternate.iter()
                    if node.tag.rsplit("}", 1)[-1] == "oMath"
                ]
                if not formulas:
                    continue
                carriers = list(alternate)
                carrier_names = [
                    carrier.tag.rsplit("}", 1)[-1] for carrier in carriers
                ]
                if (
                    carrier_names.count("Choice") != 1
                    or carrier_names.count("Fallback") > 1
                    or any(
                        name not in {"Choice", "Fallback"}
                        for name in carrier_names
                    )
                ):
                    raise common.fail(
                        "native-math AlternateContent carrier profile is invalid"
                    )
                branch_formulas: list[ET.Element] = []
                formula_branch_identities: list[tuple[int, str]] = []
                for carrier in carriers:
                    carrier_formulas = [
                        node
                        for node in carrier.iter()
                        if node.tag.rsplit("}", 1)[-1] == "oMath"
                    ]
                    if len(carrier_formulas) > 1:
                        raise common.fail(
                            "native-math branch has duplicate oMath objects"
                        )
                    if carrier_formulas:
                        carrier_identities = {
                            (int(node.get("id")), node.get("name", ""))
                            for node in carrier.iter()
                            if node.tag.rsplit("}", 1)[-1] == "cNvPr"
                            and node.get("id")
                            and node.get("name")
                        }
                        if len(carrier_identities) != 1:
                            raise common.fail(
                                "native-math branch has ambiguous shape identity"
                            )
                        formula_branch_identities.append(
                            next(iter(carrier_identities))
                        )
                    branch_formulas.extend(carrier_formulas)
                choice = next(
                    carrier
                    for carrier in carriers
                    if carrier.tag.rsplit("}", 1)[-1] == "Choice"
                )
                if not any(
                    node.tag.rsplit("}", 1)[-1] == "oMath"
                    for node in choice.iter()
                ):
                    raise common.fail("native-math Choice branch has no oMath object")
                if len(set(formula_branch_identities)) != 1:
                    raise common.fail(
                        "native-math branches disagree on shape identity"
                    )
                from lxml import etree
                from tools.powerpoint_native_math import _semantic_omml_sha256

                identities = {
                    (int(node.get("id")), node.get("name", ""))
                    for node in alternate.iter()
                    if node.tag.rsplit("}", 1)[-1] == "cNvPr"
                    and node.get("id")
                    and node.get("name")
                }
                if len(identities) != 1:
                    raise common.fail(
                        "native-math AlternateContent has ambiguous shape identity"
                    )
                shape_id, shape_name = next(iter(identities))
                semantic_hashes = {
                    _semantic_omml_sha256(
                        etree.fromstring(ET.tostring(formula, encoding="utf-8"))
                    )
                    for formula in branch_formulas
                }
                if len(semantic_hashes) != 1:
                    raise common.fail(
                        f"native-math branches disagree for {shape_name}"
                    )
                normalized_texts = {
                    "".join(
                        unicodedata.normalize("NFKC", node.text or "")
                        for node in formula.iter()
                        if node.tag.rsplit("}", 1)[-1] == "t"
                    ).strip()
                    for formula in branch_formulas
                }
                normalized_texts.discard("")
                if not normalized_texts:
                    raise common.fail(f"native-math formula text is empty for {shape_name}")
                core = {
                    "slide": member,
                    "shape_id": shape_id,
                    "shape_name": shape_name,
                    "normalized_texts": sorted(normalized_texts),
                    "semantic_omml_sha256": next(iter(semantic_hashes)),
                }
                records.append(
                    {**core, "formula_signature_sha256": _canonical_sha256(core)}
                )
    records.sort(key=lambda item: (item["slide"], item["shape_id"], item["shape_name"]))
    return records, raw_omml_count


def _attach_math_artifact_identity(run: common.Run, summary: dict) -> None:
    records, raw_omml_count = _native_math_inventory(run.pptx_path)
    summary.update(
        {
            "pptx_sha256": common.sha256_file(run.pptx_path),
            "omml_count": len(records),
            "logical_formula_count": len(records),
            "raw_omml_count": raw_omml_count,
            "formula_inventory_sha256": _canonical_sha256(records),
            "native_math_inventory": records,
        }
    )


def verify_existing_native_math(run: common.Run) -> dict:
    """Read back existing OMML against bindings, plan, and converter receipts."""

    from tools.core.contracts import read_json, write_json

    records, raw_omml_count = _native_math_inventory(run.pptx_path)
    if not records:
        raise common.fail("verify-existing requires native Office Math objects")
    bindings = read_json(run.bindings_path)
    native_bindings = [
        item
        for item in bindings.get("bindings", [])
        if item.get("object_kind") == "native-math"
        or item.get("native_math") is True
    ]
    binding_by_identity = {
        (int(item.get("shape_id", -1)), item.get("shape_name")): item
        for item in native_bindings
    }
    if len(binding_by_identity) != len(records):
        raise common.fail("native-math binding inventory does not match PPTX readback")

    plan_path = run.qa_dir / "math" / "plan.json"
    if not plan_path.is_file():
        raise common.fail("native-math declaration plan is missing")
    plan = read_json(plan_path)
    operations = plan.get("operations", [])
    operation_by_placeholder = {
        item.get("placeholder_name"): item
        for item in operations
        if item.get("placeholder_name")
    }
    if len(operation_by_placeholder) != len(records):
        raise common.fail("native-math declaration list does not match PPTX readback")

    formulas: list[dict] = []
    for record in records:
        identity = (record["shape_id"], record["shape_name"])
        binding = binding_by_identity.get(identity)
        operation = operation_by_placeholder.get(record["shape_name"])
        if binding is None or operation is None:
            raise common.fail(
                f"native-math declaration is missing for {record['shape_name']}"
            )
        formula_id = binding.get("formula_id")
        if formula_id != operation.get("formula_id"):
            raise common.fail(
                f"native-math formula identity mismatch for {record['shape_name']}"
            )
        receipt_path = run.qa_dir / "math" / str(operation.get("receipt_path", ""))
        if (
            not receipt_path.is_file()
            or common.sha256_file(receipt_path) != operation.get("receipt_sha256")
        ):
            raise common.fail(f"native-math receipt mismatch for {formula_id}")
        receipt = read_json(receipt_path)
        if (
            receipt.get("status") != "PASS"
            or receipt.get("formula_id") != formula_id
            or receipt.get("semantic_omml_sha256")
            != record["semantic_omml_sha256"]
        ):
            raise common.fail(
                f"native-math semantic readback mismatch for {formula_id}"
            )
        formulas.append(
            {
                "name": record["shape_name"],
                "placeholder": record["shape_name"],
                "formula_id": formula_id,
                "element_id": binding.get("element_id"),
                "slide_index": int(Path(record["slide"]).stem.removeprefix("slide")),
                "shape_id": record["shape_id"],
                "latex": receipt.get("canonical_latex"),
                "status": "verified",
                "normalized_texts": record["normalized_texts"],
                "semantic_omml_sha256": record["semantic_omml_sha256"],
                "formula_signature_sha256": record["formula_signature_sha256"],
                "receipt_sha256": operation.get("receipt_sha256"),
            }
        )

    summary = {
        "pptx": "redraw.pptx",
        "dry_run": True,
        "verify_existing": True,
        "audit_mode": "existing-native-math-readback",
        "detected": len(records),
        "strong": 0,
        "weak": 0,
        "injected": len(records),
        "verified": len(records),
        "failed": 0,
        "warnings": [],
        "notes": ["Existing native Office Math verified read-only against converter receipts."],
        "formulas": formulas,
        "pptx_sha256": common.sha256_file(run.pptx_path),
        "omml_count": len(records),
        "logical_formula_count": len(records),
        "raw_omml_count": raw_omml_count,
        "formula_inventory_sha256": _canonical_sha256(records),
        "native_math_inventory": records,
        "bindings_sha256": common.sha256_file(run.bindings_path),
        "bindings_complete": bindings.get("bindings_complete") is True,
        "saved_reopened": bindings.get("saved_reopened") is True,
        "plan_sha256": common.sha256_file(plan_path),
    }
    write_json(run.qa_dir / "math-summary.json", summary)
    return summary


def math_summary_blockers(run: common.Run) -> list[str]:
    """Return fail-closed blockers for any declared math-summary evidence."""

    summary_path = run.qa_dir / "math-summary.json"
    try:
        records, raw_omml_count = _native_math_inventory(run.pptx_path)
    except (Exception, SystemExit):
        return ["math-summary:invalid"]
    if not summary_path.is_file():
        return ["math-summary:missing"] if records else []
    try:
        summary = json.loads(summary_path.read_text(encoding="utf-8"))
    except (Exception, SystemExit):
        return ["math-summary:invalid"]
    blockers: list[str] = []
    if summary.get("pptx_sha256") != common.sha256_file(run.pptx_path):
        blockers.append("math-summary:pptx-hash-mismatch")
    if summary.get("omml_count") != len(records):
        blockers.append("math-summary:omml-count-mismatch")
    if summary.get("logical_formula_count") != len(records):
        blockers.append("math-summary:logical-count-mismatch")
    if summary.get("raw_omml_count") != raw_omml_count:
        blockers.append("math-summary:raw-omml-count-mismatch")
    if summary.get("formula_inventory_sha256") != _canonical_sha256(records):
        blockers.append("math-summary:formula-signature-mismatch")
    declared_inventory = summary.get("native_math_inventory")
    if declared_inventory != records:
        blockers.append("math-summary:formula-inventory-mismatch")
    bindings_payload = json.loads(run.bindings_path.read_text(encoding="utf-8"))
    native_bindings = [
        item
        for item in bindings_payload.get("bindings", [])
        if item.get("object_kind") == "native-math"
        or item.get("native_math") is True
    ]
    declared_formulas = summary.get("formulas")
    declaration_incomplete = not isinstance(declared_formulas, list)
    if records:
        declaration_incomplete = declaration_incomplete or (
            summary.get("detected") != len(records)
            or summary.get("injected") != len(records)
            or summary.get("verified") != len(records)
            or len(declared_formulas) != len(records)
            or len(native_bindings) != len(records)
        )
    else:
        # A detector-only dry run is useful diagnostic output, but it is not
        # proof that editable native math exists in the artifact.  Fail closed
        # when a stale or detector-only summary declares formulas while the
        # current PPTX has no logical OMML records.
        declaration_incomplete = declaration_incomplete or (
            summary.get("detected") not in (0, None)
            or summary.get("injected") not in (0, None)
            or summary.get("verified") not in (0, None)
            or bool(declared_formulas)
        )
    if declaration_incomplete:
        blockers.append("math-summary:declaration-empty-or-incomplete")
    if records and summary.get("verify_existing") is not True:
        blockers.append("math-summary:existing-readback-unverified")
    if records and summary.get("saved_reopened") is not True:
        blockers.append("math-summary:save-reopen-unverified")
    if records:
        if summary.get("bindings_sha256") != common.sha256_file(run.bindings_path):
            blockers.append("math-summary:bindings-hash-mismatch")
        plan_path = run.qa_dir / "math" / "plan.json"
        if (
            not plan_path.is_file()
            or summary.get("plan_sha256") != common.sha256_file(plan_path)
        ):
            blockers.append("math-summary:plan-hash-mismatch")
        else:
            try:
                plan = json.loads(plan_path.read_text(encoding="utf-8"))
                operation_by_placeholder = {
                    item.get("placeholder_name"): item
                    for item in plan.get("operations", [])
                    if item.get("placeholder_name")
                }
                binding_by_identity = {
                    (int(item.get("shape_id", -1)), item.get("shape_name")): item
                    for item in native_bindings
                }
                expected_formulas: list[dict] = []
                for record in records:
                    binding = binding_by_identity[
                        (record["shape_id"], record["shape_name"])
                    ]
                    operation = operation_by_placeholder[record["shape_name"]]
                    formula_id = binding.get("formula_id")
                    if formula_id != operation.get("formula_id"):
                        raise KeyError(record["shape_name"])
                    receipt_path = run.qa_dir / "math" / str(
                        operation.get("receipt_path", "")
                    )
                    if (
                        not receipt_path.is_file()
                        or common.sha256_file(receipt_path)
                        != operation.get("receipt_sha256")
                    ):
                        raise KeyError(formula_id)
                    receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
                    if (
                        receipt.get("formula_id") != formula_id
                        or receipt.get("status") != "PASS"
                        or receipt.get("semantic_omml_sha256")
                        != record["semantic_omml_sha256"]
                    ):
                        raise KeyError(formula_id)
                    expected_formulas.append(
                        {
                            "name": record["shape_name"],
                            "placeholder": record["shape_name"],
                            "formula_id": formula_id,
                            "element_id": binding.get("element_id"),
                            "slide_index": int(
                                Path(record["slide"]).stem.removeprefix("slide")
                            ),
                            "shape_id": record["shape_id"],
                            "latex": receipt.get("canonical_latex"),
                            "status": "verified",
                            "normalized_texts": record["normalized_texts"],
                            "semantic_omml_sha256": record[
                                "semantic_omml_sha256"
                            ],
                            "formula_signature_sha256": record[
                                "formula_signature_sha256"
                            ],
                            "receipt_sha256": operation.get("receipt_sha256"),
                        }
                    )
                if summary.get("formulas") != expected_formulas:
                    blockers.append("math-summary:formula-declaration-mismatch")
            except (KeyError, TypeError, ValueError, json.JSONDecodeError):
                blockers.append("math-summary:formula-declaration-mismatch")
    return list(dict.fromkeys(blockers))


def _refresh_v3_contracts_after_math(run: common.Run, candidates: list[dict]) -> dict:
    """Bind renamed OMML shapes and the new artifact hash back to Scene v3."""
    from tools.core.contracts import initialize_contracts, read_json, transition, utc_now, write_json

    meta = initialize_contracts(run)
    if meta["workflow"]["state"] != "repairing":
        transition(run, "repairing", "native-math-upgrade-started")

    # High-level reopen proves package compatibility.  OMML formula text boxes
    # live inside mc:AlternateContent and are intentionally invisible to
    # python-pptx's slide.shapes collection, so binding identity is read from
    # the underlying slide XML instead of being falsely reported as missing.
    Presentation(run.pptx_path)
    readback_names, omml_count = _pptx_bound_names_and_omml_count(run.pptx_path)
    rename_map = {
        candidate["name"]: candidate["placeholder"]
        for candidate in candidates
        if candidate["status"] == "injected"
    }
    formula_by_name = {
        candidate["placeholder"]: candidate["formula_id"]
        for candidate in candidates
        if candidate["status"] == "injected"
    }

    bindings = read_json(run.bindings_path)
    rebound_elements: dict[str, str] = {}
    for binding in bindings.get("bindings", []):
        old_name = binding.get("shape_name")
        if old_name in rename_map:
            binding["shape_name"] = rename_map[old_name]
            binding["object_kind"] = "native-math"
            binding["native_math"] = True
            binding["formula_id"] = formula_by_name[binding["shape_name"]]
            rebound_elements[binding["element_id"]] = binding["formula_id"]
        binding["readback_found"] = binding.get("shape_name") in readback_names

    for binding in bindings.get("logical_group_bindings", []):
        backend_names = binding.get("backend_object_names")
        backend_ids = binding.get("backend_object_ids")
        if not isinstance(backend_names, list) or not isinstance(backend_ids, list):
            binding["readback_found"] = False
            continue
        rebound_names = [rename_map.get(name, name) for name in backend_names]
        binding["backend_object_names"] = rebound_names
        binding["backend_object_identities"] = [
            {"shape_id": shape_id, "shape_name": shape_name}
            for shape_id, shape_name in zip(backend_ids, rebound_names, strict=True)
        ] if len(backend_ids) == len(rebound_names) else []
        attachment_name = binding.get("attachment_shape_name")
        if isinstance(attachment_name, str):
            binding["attachment_shape_name"] = rename_map.get(
                attachment_name, attachment_name
            )
        binding["readback_found"] = bool(rebound_names) and all(
            name in readback_names for name in rebound_names
        )

    expected = set(formula_by_name)
    rebound = {
        binding.get("shape_name")
        for binding in bindings.get("bindings", [])
        if binding.get("native_math") is True
    }
    missing_rebound = sorted(expected - rebound)
    if missing_rebound:
        raise common.fail(f"OMML shapes have no v3 binding: {missing_rebound}")
    if omml_count < len(expected):
        raise common.fail(
            f"OMML save/reopen count is incomplete: expected {len(expected)}, got {omml_count}"
        )
    bindings_complete = (
        bool(bindings.get("bindings"))
        and all(
            binding.get("readback_found") is True
            for binding in bindings["bindings"]
        )
        and all(
            binding.get("readback_found") is True
            for binding in bindings.get("logical_group_bindings", [])
        )
    )
    if not bindings_complete:
        raise common.fail("OMML save/reopen left incomplete PowerPoint shape bindings")

    pptx_hash = common.sha256_file(run.pptx_path)
    bindings.update(
        {
            "updated_at": utc_now(),
            "backend": "pptx-offline+native-math",
            "artifact_sha256": pptx_hash,
            "package_reopened": True,
            "saved_reopened": False,
            "bindings_complete": True,
        }
    )
    write_json(run.bindings_path, bindings)

    scene = read_json(run.scene_path)
    for element in scene.get("elements", []):
        formula_id = rebound_elements.get(element.get("id"))
        if formula_id:
            element["native_math"] = True
            element["formula_id"] = formula_id
    scene["updated_at"] = utc_now()
    scene["artifact"] = {
        "backend": "pptx-offline+native-math",
        "path": "redraw.pptx",
        "sha256": pptx_hash,
    }
    write_json(run.scene_path, scene)
    from tools.arrows.pptx_arrows import write_arrow_reports
    from tools.assets.primitives import audit_primitives
    from tools.providers.providers import write_case_capabilities

    write_arrow_reports(run)
    audit_primitives(run)
    write_case_capabilities(run)
    from tools.pipeline.layout import audit_layout

    layout_report = audit_layout(run)
    transition(
        run,
        "candidate",
        "native-math-upgrade-complete",
        details={
            "pptx_sha256": pptx_hash,
            "formula_count": len(expected),
            "object_count": len(bindings["bindings"]),
            "bindings_complete": True,
            "layout_pass": layout_report["pass"],
            "layout_findings": len(layout_report["findings"]),
        },
    )
    return {
        "pptx_sha256": pptx_hash,
        "object_count": len(bindings["bindings"]),
        "omml_count": omml_count,
        "package_reopened": True,
        "saved_reopened": False,
        "bindings_complete": True,
        "layout_pass": layout_report["pass"],
        "layout_findings": len(layout_report["findings"]),
    }


# ---------------------------------------------------------------- 主流程


def upgrade(
    run: common.Run, *, dry_run: bool = False, verify_existing: bool = False
) -> dict:
    """检测公式框并批量注入 OMML，返回 summary dict（同时落 qa/math-summary.json）。"""
    if not run.pptx_path.is_file():
        raise common.fail(f"未找到 PPTX: {run.pptx_path}（请先运行 autofigure convert）")
    run.qa_dir.mkdir(exist_ok=True)
    existing_native_math, _ = _native_math_inventory(run.pptx_path)
    if verify_existing or (dry_run and existing_native_math):
        return verify_existing_native_math(run)
    prs = Presentation(run.pptx_path)

    candidates: list[dict] = []
    for slide_index, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            signal = classify_runs(_shape_runs(shape))
            if signal is None:
                continue
            candidates.append({"slide_index": slide_index, "shape": shape, "signal": signal})
    if existing_native_math and not candidates:
        return verify_existing_native_math(run)

    # 编号避开残留 math:NNN 名（上次注入失败后重跑等场景），保证 cNvPr@name 唯一
    candidate_ids = {id(cand["shape"]) for cand in candidates}
    taken_names = {
        shape.name
        for slide in prs.slides
        for shape in slide.shapes
        if id(shape) not in candidate_ids
    }
    number = 0
    for cand in candidates:
        number += 1
        while f"{NAME_PREFIX}{number:03d}" in taken_names:
            number += 1
        placeholder = f"{NAME_PREFIX}{number:03d}"
        taken_names.add(placeholder)
        shape = cand["shape"]
        cand["name"] = shape.name
        cand["placeholder"] = placeholder
        cand["formula_id"] = f"EQ{number:03d}"
        cand["latex"] = rebuild_latex(_shape_runs(shape))
        cand["bold"] = any(
            run.font.bold
            for paragraph in shape.text_frame.paragraphs
            for run in paragraph.runs
        )
        cand["font"] = _target_font(shape)
        cand["status"] = "detected"

    warnings: list[str] = []
    summary = {
        "pptx": "redraw.pptx",
        "dry_run": dry_run,
        "detected": len(candidates),
        "strong": sum(1 for cand in candidates if cand["signal"] == "strong"),
        "weak": sum(1 for cand in candidates if cand["signal"] == "weak"),
        "injected": 0,
        "failed": 0,
        "warnings": warnings,
        "notes": [],
        "formulas": [],
    }

    def write_summary() -> None:
        rows = []
        for cand in candidates:
            row = {
                key: cand[key]
                for key in ("name", "placeholder", "formula_id", "slide_index", "signal", "latex", "status")
            }
            if cand["bold"]:
                row["bold"] = True
            if cand.get("error"):
                row["error"] = cand["error"]
            rows.append(row)
        summary["formulas"] = rows
        summary["injected"] = sum(1 for cand in candidates if cand["status"] == "injected")
        summary["failed"] = sum(1 for cand in candidates if cand["status"] == "failed")
        if any(cand["bold"] and cand["status"] == "injected" for cand in candidates):
            summary["notes"].append("部分公式原文为粗体：OMML 数学区不保留粗体（可接受差异）")
        _attach_math_artifact_identity(run, summary)
        (run.qa_dir / "math-summary.json").write_text(
            json.dumps(summary, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
        )

    if not candidates or dry_run:
        write_summary()
        return summary

    engine = _load_engine()
    math_dir = run.qa_dir / "math"
    if math_dir.exists():  # 重跑覆盖当前最佳：清掉上次注入的 receipt/plan
        shutil.rmtree(math_dir)
    math_dir.mkdir(parents=True)

    operations: list[dict] = []
    for cand in candidates:
        try:
            receipt = engine.compile_formula(cand["formula_id"], cand["latex"], "inline")
        except Exception as exc:  # 单公式失败只 warn 跳过该框（保留原样），不整案失败
            cand["status"] = "failed"
            cand["error"] = f"{type(exc).__name__}: {exc}"
            warnings.append(f"{cand['name']} compile 失败（保留原文本框）: {exc}")
            continue
        receipt_path = math_dir / f"{cand['formula_id']}.json"
        receipt_path.write_text(
            json.dumps(receipt, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
        )
        operation = {
            "slide_index": cand["slide_index"],
            "placeholder_name": cand["placeholder"],
            "formula_id": cand["formula_id"],
            "receipt_path": receipt_path.name,
            "receipt_sha256": common.sha256_file(receipt_path),
        }
        if cand["font"] is not None:
            operation["target_font_size_pt"] = cand["font"][0]
            operation["target_font_color"] = cand["font"][1]
        operations.append(operation)
        cand["status"] = "compiled"

    ready = [cand for cand in candidates if cand["status"] == "compiled"]
    if not ready:  # 全部失败/无公式框：正常退出并说明，不改 pptx
        write_summary()
        return summary

    # 先改名落盘（inject_plan 按 cNvPr@name 定位），再注入到临时 pptx，成功后原子替换
    for cand in ready:
        cand["shape"].name = cand["placeholder"]

    plan_path = math_dir / "plan.json"
    plan_path.write_text(
        json.dumps({"schema_version": "1.0", "operations": operations}, ensure_ascii=False, indent=2)
        + "\n",
        encoding="utf-8",
    )
    staged_pptx = run.pptx_path.with_name("redraw.math-source.pptx")
    tmp_pptx = run.pptx_path.with_name("redraw.math-tmp.pptx")
    staged_pptx.unlink(missing_ok=True)
    tmp_pptx.unlink(missing_ok=True)
    prs.save(staged_pptx)
    try:
        engine.inject_plan(staged_pptx, plan_path, tmp_pptx)
    except Exception as exc:
        tmp_pptx.unlink(missing_ok=True)
        raise common.fail(f"OMML 注入失败，redraw.pptx 保持原样: {exc}") from exc
    finally:
        staged_pptx.unlink(missing_ok=True)
    os.replace(tmp_pptx, run.pptx_path)
    for cand in ready:
        cand["status"] = "injected"

    summary.update(_refresh_v3_contracts_after_math(run, candidates))
    write_summary()
    from tools.core.revisions import stamp_active_revision

    stamp_active_revision(run)
    return summary


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(prog="autofigure math", description=__doc__)
    parser.add_argument("run_dir", type=Path, help="run 目录")
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="只检测并重建 LaTeX 写 qa/math-summary.json，不改 PPTX",
    )
    parser.add_argument(
        "--verify-existing",
        action="store_true",
        help="只读验证当前原生 Office Math、bindings、plan 与 receipt，不改 PPTX",
    )
    args = parser.parse_args(argv)

    if args.dry_run and args.verify_existing:
        raise common.fail("--dry-run 与 --verify-existing 不能同时使用")

    run = common.open_run(args.run_dir)
    summary = upgrade(
        run, dry_run=args.dry_run, verify_existing=args.verify_existing
    )
    _echo(f"检测到 {summary['detected']} 个公式框（强信号 {summary['strong']} / 弱信号 {summary['weak']}）")
    if args.verify_existing:
        _echo(
            f"verify-existing：只读验证 {summary['verified']} 个原生 Office Math，PPTX 未改动"
        )
    elif args.dry_run:
        _echo("dry-run：未改动 PPTX")
    elif summary["injected"]:
        _echo(f"已注入 {summary['injected']} 个原生 Office Math；失败 {summary['failed']} 个（保留原文本框）")
    elif summary["detected"]:
        _echo(f"全部 {summary['failed']} 个公式 compile 失败，PPTX 未改动")
    else:
        _echo("无公式框，PPTX 未改动")
    for warning in summary["warnings"]:
        _echo(f"warning: {warning}")
    _echo(f"明细: {run.qa_dir / 'math-summary.json'}")
    if not args.dry_run and not args.verify_existing and summary["injected"]:
        try:
            from tools.pipeline import render_export

            meta = run.load_meta()
            render_export.render(run.pptx_path, run.render_png, int(meta["width"]), int(meta["height"]))
            from tools.core.revisions import stamp_active_revision

            stamp_active_revision(run)
            _echo(f"fresh render: {run.render_png}")
        except Exception as exc:  # COM 渲染失败只 warn，注入结果不受影响
            _echo(f"warning: fresh render 刷新失败: {exc}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
