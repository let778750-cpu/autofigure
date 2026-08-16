#!/usr/bin/env python3
"""Create a truly blank PPTX canvas whose aspect ratio is measured from a PNG."""

from __future__ import annotations

import argparse
import hashlib
import json
import math
import os
import sys
import tempfile
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

TOOLS_DIRECTORY = Path(__file__).resolve().parent
if str(TOOLS_DIRECTORY) not in sys.path:
    sys.path.insert(0, str(TOOLS_DIRECTORY))

try:
    from output_policy import resolve_output_path
except ModuleNotFoundError:  # Support: python -m tools.create_canvas_pptx
    try:
        from .output_policy import resolve_output_path
    except ImportError:  # Support importlib loading a standalone file from the project root.
        from tools.output_policy import resolve_output_path


DEFAULT_LONG_EDGE_IN = 13.333333
MIN_SLIDE_EDGE_IN = 1.0
MAX_SLIDE_EDGE_IN = 56.0


def sha256_file(path: Path) -> str:
    """Return the SHA-256 digest of *path* without loading it all into memory."""
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def inspect_png(path: str | Path) -> dict[str, Any]:
    """Measure a PNG and bind the measurement to its content hash."""
    source = Path(path).expanduser().resolve()
    if not source.is_file():
        raise ValueError(f"PNG source does not exist: {source}")
    try:
        with Image.open(source) as image:
            if image.format != "PNG":
                raise ValueError(f"Expected a PNG source, got {image.format or 'unknown'}: {source}")
            width_px, height_px = image.size
            mode = image.mode
    except OSError as exc:
        raise ValueError(f"Cannot read PNG source {source}: {exc}") from exc
    if width_px <= 0 or height_px <= 0:
        raise ValueError(f"PNG has invalid dimensions: {width_px}x{height_px}")
    return {
        "path": str(source),
        "sha256": sha256_file(source),
        "width_px": width_px,
        "height_px": height_px,
        "pixel_mode": mode,
    }


def slide_size_for_aspect(
    width_px: int,
    height_px: int,
    *,
    long_edge_in: float = DEFAULT_LONG_EDGE_IN,
    min_edge_in: float = MIN_SLIDE_EDGE_IN,
    max_edge_in: float = MAX_SLIDE_EDGE_IN,
) -> dict[str, float | int]:
    """Calculate an Office-safe slide size while preserving the PNG aspect exactly."""
    values = (width_px, height_px, long_edge_in, min_edge_in, max_edge_in)
    if not all(isinstance(value, (int, float)) and math.isfinite(float(value)) and value > 0 for value in values):
        raise ValueError("Image dimensions and slide-edge limits must be finite positive numbers.")
    if min_edge_in > max_edge_in:
        raise ValueError("min_edge_in cannot exceed max_edge_in.")

    aspect = float(width_px) / float(height_px)
    if aspect >= 1.0:
        width_in = float(long_edge_in)
        height_in = width_in / aspect
        if height_in < min_edge_in:
            height_in = min_edge_in
            width_in = height_in * aspect
    else:
        height_in = float(long_edge_in)
        width_in = height_in * aspect
        if width_in < min_edge_in:
            width_in = min_edge_in
            height_in = width_in / aspect

    if max(width_in, height_in) > max_edge_in + 1e-9:
        raise ValueError(
            f"Aspect ratio {aspect:.6f} cannot fit within the configured {max_edge_in:g}-inch Office limit."
        )

    width_emu = int(Inches(width_in))
    height_emu = int(Inches(height_in))
    return {
        "width_in": width_in,
        "height_in": height_in,
        "width_emu": width_emu,
        "height_emu": height_emu,
        "aspect_ratio": aspect,
    }


def create_blank_canvas_pptx(
    source_png: str | Path,
    output_pptx: str | Path,
    *,
    long_edge_in: float = DEFAULT_LONG_EDGE_IN,
    overwrite: bool = False,
) -> dict[str, Any]:
    """Create one empty blank slide, returning hash-bound source and canvas metadata."""
    source = inspect_png(source_png)
    slide_size = slide_size_for_aspect(
        source["width_px"],
        source["height_px"],
        long_edge_in=long_edge_in,
    )
    output = resolve_output_path(output_pptx)
    if output.suffix.lower() != ".pptx":
        raise ValueError(f"Output must use the .pptx extension: {output}")
    if output.exists() and not overwrite:
        raise ValueError(f"Refusing to overwrite existing PPTX without overwrite=True: {output}")
    output.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()
    presentation.slide_width = slide_size["width_emu"]
    presentation.slide_height = slide_size["height_emu"]
    blank_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank_layout)
    if len(slide.shapes) != 0:
        raise RuntimeError("The selected PowerPoint layout is not blank.")

    temporary_path: Path | None = None
    try:
        with tempfile.NamedTemporaryFile(
            prefix=f".{output.stem}-",
            suffix=".pptx",
            dir=output.parent,
            delete=False,
        ) as temporary:
            temporary_path = Path(temporary.name)
        presentation.save(temporary_path)
        os.replace(temporary_path, output)
        temporary_path = None
    finally:
        if temporary_path is not None and temporary_path.exists():
            temporary_path.unlink()

    return {
        "status": "PASS",
        "source": source,
        "output_pptx": str(output),
        "output_pptx_sha256": sha256_file(output),
        "slide": {
            **slide_size,
            "slide_count": 1,
            "shape_count": 0,
        },
    }


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Measure a PNG and create a one-slide blank PPTX with the identical aspect ratio."
    )
    parser.add_argument("source_png", help="Reference PNG whose measured aspect ratio controls the slide.")
    parser.add_argument("output_pptx", help="New blank .pptx path.")
    parser.add_argument("--long-edge-in", type=float, default=DEFAULT_LONG_EDGE_IN)
    parser.add_argument("--overwrite", action="store_true")
    parser.add_argument("--pretty", action="store_true")
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    try:
        result = create_blank_canvas_pptx(
            args.source_png,
            args.output_pptx,
            long_edge_in=args.long_edge_in,
            overwrite=args.overwrite,
        )
        exit_code = 0
    except (OSError, RuntimeError, ValueError) as exc:
        result = {"status": "SPEC_INVALID", "passed": False, "errors": [str(exc)]}
        exit_code = 2
    print(json.dumps(result, ensure_ascii=False, indent=2 if args.pretty else None))
    return exit_code


if __name__ == "__main__":
    sys.exit(main())
