#!/usr/bin/env python3
"""Hash-bound, deterministic scene preflight before any PowerPoint mutation."""

from __future__ import annotations

import argparse
import hashlib
import io
import json
import math
import os
import re
import sys
from collections import Counter, defaultdict
from pathlib import Path
from typing import Any, Iterable, Mapping, Sequence

import numpy as np
from jsonschema import Draft202012Validator, FormatChecker
from PIL import Image, ImageFont
from pptx import Presentation
from pptx.exc import PackageNotFoundError

TOOLS_DIRECTORY = Path(__file__).resolve().parent
if str(TOOLS_DIRECTORY) not in sys.path:
    sys.path.insert(0, str(TOOLS_DIRECTORY))

try:
    from create_canvas_pptx import inspect_png, sha256_file, slide_size_for_aspect
    from powerpoint_native_math import NativeMathError, _validated_receipt
except ModuleNotFoundError:  # Support: python -m tools.preflight_scene
    from .create_canvas_pptx import inspect_png, sha256_file, slide_size_for_aspect
    from .powerpoint_native_math import NativeMathError, _validated_receipt

try:
    from output_policy import resolve_output_path
except ModuleNotFoundError:  # Support: python -m tools.preflight_scene
    try:
        from .output_policy import resolve_output_path
    except ImportError:  # Support importlib loading a standalone file from the project root.
        from tools.output_policy import resolve_output_path


DEFAULT_SCHEMA_PATH = Path(__file__).resolve().parents[1] / "schemas" / "figure-spec.schema.json"
DEFAULT_PERCEPTION_SCHEMA_PATH = (
    Path(__file__).resolve().parents[1] / "schemas" / "perception-manifest.schema.json"
)
DEFAULT_REVIEW_SCHEMA_PATH = (
    Path(__file__).resolve().parents[1] / "schemas" / "perception-review.schema.json"
)
DEFAULT_GEOMETRY_SCHEMA_PATH = (
    Path(__file__).resolve().parents[1] / "schemas" / "geometry-manifest.schema.json"
)
DEFAULT_HOST_RUNTIME_RECEIPT_SCHEMA_PATH = (
    Path(__file__).resolve().parents[1] / "schemas" / "host-runtime-receipt.schema.json"
)
DEFAULT_GEOMETRY_SCRIPT_PATH = Path(__file__).resolve().parent / "geometry_refinement.py"
TEXT_TYPES = {"text", "formula"}
NON_COLLIDING_TYPES = {"background"}
CONTAINER_TYPES = {"background", "panel", "plot", "legend", "micro_asset"}
STATUS_PRECEDENCE = {"PASS": 0, "INCONCLUSIVE": 1, "REGION_REPLAN": 2, "SPEC_INVALID": 3}
FONT_ALIASES = {
    "arial": ["arial.ttf"],
    "calibri": ["calibri.ttf"],
    "microsoftyahei": ["msyh.ttc", "msyh.ttf"],
    "notosanscjksc": ["NotoSansCJK-Regular.ttc", "NotoSansCJKsc-Regular.otf"],
    "sourcesanshansc": ["SourceHanSansSC-Regular.otf"],
    "timesnewroman": ["times.ttf"],
}
MATH_SEMANTIC_ROLES = {
    "equation",
    "formula",
    "math",
    "math_label",
    "parameter_expression",
    "variable_expression",
}
LATEX_COMMAND_RE = re.compile(r"\\(?:[A-Za-z]+|[()[\]])")
LATEX_DELIMITER_RE = re.compile(r"(?:\$\$[^$]+\$\$|(?<!\$)\$[^$]+\$(?!\$)|\\\(|\\\)|\\\[|\\\])")
SCRIPT_NOTATION_RE = re.compile(
    r"(?<![\w])(?:[A-Za-z\u0370-\u03ff])(?:_(?:\{[^{}]+\}|[A-Za-z0-9]+)|\^(?:\{[^{}]+\}|[-+A-Za-z0-9]+))(?![\w])"
)
UNICODE_MATH_STRUCTURE_RE = re.compile(
    r"[=\u2248\u2260\u2264\u2265\u00b1\u2213\u00d7\u00f7\u2211\u220f\u222b\u221a\u2202\u2207]"
)
UNICODE_SCRIPT_RE = re.compile(r"[\u2070\u00b9\u00b2\u00b3\u2074-\u2079\u2080-\u2089]")


def sha256_text(value: str) -> str:
    """Hash the exact UTF-8 canonical source retained in the frozen spec."""
    return hashlib.sha256(value.encode("utf-8")).hexdigest()


def _strict_json_loads(payload: str, *, label: str) -> Any:
    """Parse standards-compliant JSON and reject ambiguous duplicate keys."""

    def reject_constant(value: str) -> Any:
        raise ValueError(f"{label} contains non-finite JSON number {value!r}")

    def finite_float(value: str) -> float:
        result = float(value)
        if not math.isfinite(result):
            raise ValueError(f"{label} contains out-of-range JSON number {value!r}")
        return result

    def unique_object(pairs: list[tuple[str, Any]]) -> dict[str, Any]:
        result: dict[str, Any] = {}
        for key, value in pairs:
            if key in result:
                raise ValueError(f"{label} contains duplicate object key {key!r}")
            result[key] = value
        return result

    try:
        return json.loads(
            payload,
            object_pairs_hook=unique_object,
            parse_constant=reject_constant,
            parse_float=finite_float,
        )
    except json.JSONDecodeError as exc:
        raise ValueError(f"{label} is not valid JSON: {exc}") from exc


def _strict_json_file(path: Path, *, label: str) -> Any:
    return _strict_json_loads(path.read_text(encoding="utf-8"), label=label)


def _strict_json_bytes(payload: bytes, *, label: str) -> Any:
    try:
        text_payload = payload.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise ValueError(f"{label} is not UTF-8 JSON: {exc}") from exc
    return _strict_json_loads(text_payload, label=label)


def _sha256_bytes(payload: bytes) -> str:
    return hashlib.sha256(payload).hexdigest()


def _non_finite_path(value: Any, path: str = "$") -> str | None:
    """Find non-finite floats in mappings supplied directly to the Python API."""
    if isinstance(value, float) and not math.isfinite(value):
        return path
    if isinstance(value, Mapping):
        for key, item in value.items():
            found = _non_finite_path(item, f"{path}.{key}")
            if found is not None:
                return found
    elif isinstance(value, (list, tuple)):
        for index, item in enumerate(value):
            found = _non_finite_path(item, f"{path}[{index}]")
            if found is not None:
                return found
    return None


def _content_runs(element: Mapping[str, Any]) -> list[Mapping[str, Any]]:
    value = element.get("content_runs")
    if not isinstance(value, list):
        return []
    return [item for item in value if isinstance(item, Mapping)]


def _text_projection(element: Mapping[str, Any]) -> str:
    """Return only prose from structured content; math never becomes ordinary text."""
    if isinstance(element.get("text"), str):
        return str(element["text"])
    return "".join(
        str(run.get("text", "")) for run in _content_runs(element) if run.get("kind") == "text"
    )


def _text_measurement_projection(element: Mapping[str, Any]) -> str:
    """Use a visible placeholder only for approximate capacity diagnostics."""
    if isinstance(element.get("text"), str):
        return str(element["text"])
    return "".join(
        str(run.get("text", "")) if run.get("kind") == "text" else "\u25a1"
        for run in _content_runs(element)
    )


def math_like_text_reasons(element: Mapping[str, Any]) -> list[str]:
    """Conservatively find math syntax that must be represented by a math run.

    Scientific entity labels such as IL-6 and alpha-SMA are intentionally not
    classified from hyphens, digits, or Greek letters alone.
    """
    text = str(element.get("text", ""))
    reasons: list[str] = []
    if LATEX_DELIMITER_RE.search(text):
        reasons.append("latex_delimiter")
    if LATEX_COMMAND_RE.search(text):
        reasons.append("latex_command")
    if SCRIPT_NOTATION_RE.search(text):
        reasons.append("subscript_or_superscript_notation")
    if UNICODE_MATH_STRUCTURE_RE.search(text):
        reasons.append("math_operator")
    if UNICODE_SCRIPT_RE.search(text):
        reasons.append("unicode_script")
    semantic_role = str(element.get("semantic_role", "")).strip().casefold()
    if text.strip() and semantic_role in MATH_SEMANTIC_ROLES:
        reasons.append("math_semantic_role")
    return sorted(set(reasons))


def _normalized_font_name(value: str) -> str:
    return re.sub(r"[^a-z0-9]+", "", value.casefold())


def default_font_search_paths() -> list[Path]:
    roots: list[Path] = []
    windows_dir = os.environ.get("WINDIR")
    if windows_dir:
        roots.append(Path(windows_dir) / "Fonts")
    roots.extend(
        [
            Path("C:/Windows/Fonts"),
            Path("/Library/Fonts"),
            Path("/System/Library/Fonts"),
            Path("/usr/share/fonts"),
            Path.home() / ".fonts",
        ]
    )
    unique: list[Path] = []
    seen: set[str] = set()
    for root in roots:
        key = str(root).casefold()
        if key not in seen and root.is_dir():
            unique.append(root)
            seen.add(key)
    return unique


def resolve_font_path(
    font_family: str,
    *,
    font_path: str | None = None,
    base_dir: Path | None = None,
    search_paths: Sequence[str | Path] | None = None,
) -> Path | None:
    """Resolve the requested font without silently substituting a different family."""
    if font_path:
        candidate = Path(font_path).expanduser()
        if not candidate.is_absolute() and base_dir is not None:
            candidate = base_dir / candidate
        candidate = candidate.resolve()
        return candidate if candidate.is_file() else None

    try:
        loaded = ImageFont.truetype(font_family, 12)
        loaded_path = Path(str(getattr(loaded, "path", "")))
        if loaded_path.is_file():
            return loaded_path.resolve()
    except OSError:
        pass

    family_key = _normalized_font_name(font_family)
    candidate_names = FONT_ALIASES.get(family_key, [])
    roots = [Path(value).expanduser() for value in (search_paths or default_font_search_paths())]
    for root in roots:
        if not root.is_dir():
            continue
        for filename in candidate_names:
            candidate = root / filename
            if candidate.is_file():
                return candidate.resolve()

    extensions = {".ttf", ".ttc", ".otf"}
    for root in roots:
        if not root.is_dir():
            continue
        try:
            candidates = root.rglob("*")
            for candidate in candidates:
                if candidate.suffix.casefold() not in extensions:
                    continue
                if _normalized_font_name(candidate.stem) == family_key:
                    return candidate.resolve()
        except OSError:
            continue
    return None


def _text_width(font: ImageFont.FreeTypeFont, text: str) -> float:
    return float(font.getlength(text)) if text else 0.0


def _break_long_token(token: str, font: ImageFont.FreeTypeFont, max_width: float) -> list[str]:
    chunks: list[str] = []
    current = ""
    for character in token:
        candidate = current + character
        if current and _text_width(font, candidate) > max_width:
            chunks.append(current)
            current = character
        else:
            current = candidate
    if current or not chunks:
        chunks.append(current)
    return chunks


def wrap_text_pillow(
    text: str, font: ImageFont.FreeTypeFont, max_width: float, *, wrap: bool = True
) -> list[str]:
    """Greedily wrap text using real Pillow glyph advances, including CJK/long-token fallback."""
    if not wrap:
        return text.split("\n") or [""]
    if max_width <= 0:
        return text.split("\n") or [""]

    output: list[str] = []
    for paragraph in text.split("\n"):
        if paragraph == "":
            output.append("")
            continue
        tokens = re.findall(r"\S+\s*", paragraph)
        current = ""
        for token in tokens:
            candidate = current + token
            if current and _text_width(font, candidate.rstrip()) > max_width:
                output.append(current.rstrip())
                current = ""
            if _text_width(font, token.rstrip()) <= max_width:
                current += token
                continue
            chunks = _break_long_token(token.rstrip(), font, max_width)
            for chunk in chunks[:-1]:
                if current:
                    combined = current + chunk
                    if _text_width(font, combined) <= max_width:
                        output.append(combined.rstrip())
                    else:
                        output.append(current.rstrip())
                        output.append(chunk.rstrip())
                    current = ""
                else:
                    output.append(chunk.rstrip())
            current += chunks[-1]
            if token.endswith(" "):
                current += " "
        output.append(current.rstrip())
    return output or [""]


def _margin_values(value: Any) -> tuple[float, float, float, float]:
    if value is None:
        return 0.0, 0.0, 0.0, 0.0
    if isinstance(value, (int, float)):
        margin = float(value)
        return margin, margin, margin, margin
    if isinstance(value, Mapping):
        return tuple(float(value.get(side, 0.0)) for side in ("left", "right", "top", "bottom"))
    raise ValueError("margin_px must be a number or an object")


def measure_text_fit(
    element: Mapping[str, Any],
    *,
    measurement_dpi: float = 96.0,
    base_dir: Path | None = None,
    font_search_paths: Sequence[str | Path] | None = None,
) -> dict[str, Any]:
    """Measure one text box using its requested font; never substitute on lookup failure."""
    style = element.get("text_style") or {}
    family = str(style.get("font_family", "")).strip()
    resolved_font = resolve_font_path(
        family,
        font_path=style.get("font_path"),
        base_dir=base_dir,
        search_paths=font_search_paths,
    )
    if resolved_font is None:
        return {
            "status": "INCONCLUSIVE",
            "element_id": element.get("id"),
            "font_family": family,
            "message": "Requested font could not be resolved; text fit was not marked PASS.",
        }

    if "font_size_px" in style:
        font_size_px = float(style["font_size_px"])
    else:
        font_size_px = float(style["font_size_pt"]) * float(measurement_dpi) / 72.0
    font = ImageFont.truetype(str(resolved_font), max(1, int(round(font_size_px))))
    bbox = element["bbox"]
    left, right, top, bottom = _margin_values(style.get("margin_px"))
    rotation = float(style.get("rotation_deg", 0.0)) % 360.0
    if not any(math.isclose(rotation, value, abs_tol=1e-6) for value in (0.0, 90.0, 180.0, 270.0)):
        return {
            "status": "INCONCLUSIVE",
            "element_id": element.get("id"),
            "font_family": family,
            "font_path": str(resolved_font),
            "message": f"Arbitrary text rotation {rotation:g}° requires renderer readback.",
        }

    box_width = float(bbox["w"])
    box_height = float(bbox["h"])
    if math.isclose(rotation, 90.0, abs_tol=1e-6) or math.isclose(rotation, 270.0, abs_tol=1e-6):
        box_width, box_height = box_height, box_width
    available_width = box_width - left - right
    available_height = box_height - top - bottom
    if available_width <= 0 or available_height <= 0:
        return {
            "status": "REGION_REPLAN",
            "element_id": element.get("id"),
            "font_family": family,
            "font_path": str(resolved_font),
            "available_width_px": available_width,
            "available_height_px": available_height,
            "message": "Text margins consume the entire text box.",
        }

    measurement_text = _text_measurement_projection(element)
    lines = wrap_text_pillow(
        measurement_text,
        font,
        available_width,
        wrap=bool(style.get("wrap", True)),
    )
    ascent, descent = font.getmetrics()
    line_height = float(ascent + descent) * float(style.get("line_spacing", 1.0))
    measured_width = max((_text_width(font, line) for line in lines), default=0.0)
    measured_height = line_height * max(1, len(lines))
    max_lines = int(math.floor(available_height / line_height)) if line_height > 0 else 0
    fits = measured_width <= available_width + 0.5 and measured_height <= available_height + 0.5
    return {
        "status": "PASS" if fits else "REGION_REPLAN",
        "element_id": element.get("id"),
        "font_family": family,
        "font_path": str(resolved_font),
        "font_size_px": font_size_px,
        "available_width_px": round(available_width, 4),
        "available_height_px": round(available_height, 4),
        "measured_width_px": round(measured_width, 4),
        "measured_height_px": round(measured_height, 4),
        "line_height_px": round(line_height, 4),
        "line_count": len(lines),
        "max_lines": max_lines,
        "wrapped_lines": lines,
        "structured_math_placeholder_count": sum(
            1 for run in _content_runs(element) if run.get("kind") == "math"
        ),
        "diagnostic_only": bool(_content_runs(element)),
        "message": "Text fits using Pillow glyph metrics."
        if fits
        else "Text exceeds the measured box capacity.",
    }


def measure_formula_fit(
    element: Mapping[str, Any],
    formula: Mapping[str, Any],
    *,
    measurement_dpi: float = 96.0,
) -> dict[str, Any]:
    """Approximate capacity with MathText; never prove Office Math insertability."""
    try:
        from matplotlib.font_manager import FontProperties
        from matplotlib.mathtext import MathTextParser
    except ImportError as exc:
        return {
            "status": "INCONCLUSIVE",
            "element_id": element.get("id"),
            "diagnostic_only": True,
            "proves_native_office_math": False,
            "message": f"Matplotlib MathText is unavailable: {exc}",
        }

    style = (
        element.get("formula_style") if isinstance(element.get("formula_style"), Mapping) else {}
    )
    if "font_size_px" in style:
        font_size_px = float(style["font_size_px"])
        font_size_pt = font_size_px * 72.0 / float(measurement_dpi)
    else:
        font_size_pt = float(style.get("font_size_pt", 12.0))
        font_size_px = font_size_pt * float(measurement_dpi) / 72.0
    latex = str(formula.get("canonical_latex", "")).strip()
    if latex.startswith(r"\(") and latex.endswith(r"\)"):
        latex = latex[2:-2]
    expression = latex if latex.startswith("$") and latex.endswith("$") else f"${latex}$"
    try:
        parsed = MathTextParser("path").parse(
            expression,
            dpi=float(measurement_dpi),
            prop=FontProperties(size=font_size_pt),
        )
    except (TypeError, ValueError) as exc:
        return {
            "status": "INCONCLUSIVE",
            "element_id": element.get("id"),
            "canonical_latex": str(formula.get("canonical_latex", "")),
            "diagnostic_only": True,
            "proves_native_office_math": False,
            "message": f"Formula cannot be deterministically parsed by MathText: {exc}",
        }

    left, right, top, bottom = _margin_values(style.get("margin_px"))
    bbox = element["bbox"]
    available_width = float(bbox["w"]) - left - right
    available_height = float(bbox["h"]) - top - bottom
    rotation = float(style.get("rotation_deg", 0.0)) % 360.0
    if not any(math.isclose(rotation, value, abs_tol=1e-6) for value in (0.0, 90.0, 180.0, 270.0)):
        return {
            "status": "INCONCLUSIVE",
            "element_id": element.get("id"),
            "diagnostic_only": True,
            "proves_native_office_math": False,
            "message": f"Arbitrary formula rotation {rotation:g} degrees requires renderer readback.",
        }
    measured_width = float(parsed.width)
    measured_height = float(parsed.height) + max(0.0, float(parsed.depth))
    if math.isclose(rotation, 90.0, abs_tol=1e-6) or math.isclose(rotation, 270.0, abs_tol=1e-6):
        measured_width, measured_height = measured_height, measured_width
    fits = (
        available_width > 0
        and available_height > 0
        and measured_width <= available_width + 0.5
        and measured_height <= available_height + 0.5
    )
    return {
        "status": "PASS" if fits else "REGION_REPLAN",
        "element_id": element.get("id"),
        "formula_id": formula.get("id"),
        "canonical_latex": str(formula.get("canonical_latex", "")),
        "measurement_engine": "matplotlib_mathtext",
        "diagnostic_only": True,
        "proves_native_office_math": False,
        "font_size_px": round(font_size_px, 4),
        "available_width_px": round(available_width, 4),
        "available_height_px": round(available_height, 4),
        "measured_width_px": round(measured_width, 4),
        "measured_height_px": round(measured_height, 4),
        "message": (
            "MathText approximation fits; native Office Math still requires converter receipt and readback."
            if fits
            else "MathText approximation exceeds the declared box."
        ),
    }


def inspect_formula_converter_receipt(
    formula: Mapping[str, Any],
    *,
    base_dir: Path,
) -> tuple[dict[str, Any] | None, list[dict[str, Any]]]:
    """Trust a compile receipt only after current-runtime deterministic recompilation."""
    issues: list[dict[str, Any]] = []
    path_value = str(formula.get("converter_receipt_path", "")).strip()
    if not path_value:
        issues.append(
            {
                "disposition": "INCONCLUSIVE",
                "code": "FORMULA_CONVERTER_RECEIPT_MISSING",
                "message": "Native Office Math converter receipt is missing.",
            }
        )
        return None, issues
    path = Path(path_value).expanduser()
    if not path.is_absolute():
        path = base_dir / path
    path = path.resolve()
    if not path.is_file():
        issues.append(
            {
                "disposition": "INCONCLUSIVE",
                "code": "FORMULA_CONVERTER_RECEIPT_UNREADABLE",
                "message": f"Native Office Math converter receipt does not exist: {path}",
            }
        )
        return {"path": str(path), "status": "INCONCLUSIVE"}, issues

    actual_hash = sha256_file(path)
    expected_hash = str(formula.get("converter_receipt_sha256", ""))
    if actual_hash.casefold() != expected_hash.casefold():
        issues.append(
            {
                "disposition": "SPEC_INVALID",
                "code": "FORMULA_CONVERTER_RECEIPT_HASH_MISMATCH",
                "message": "Formula converter receipt hash differs from the frozen spec.",
                "evidence": {"expected": expected_hash, "actual": actual_hash},
            }
        )

    try:
        # This shared validator does not trust fields merely because they are
        # self-consistent.  It parses and hashes the embedded MathML/OMML,
        # binds the trusted Office XSL, then recompiles canonical LaTeX with
        # the currently installed pinned converter and compares every output.
        payload = _validated_receipt(path)
    except NativeMathError as exc:
        message = str(exc)
        unavailable = "conversion is INCONCLUSIVE" in message or (
            "was not found" in message and "MML2OMML.XSL" in message
        )
        issues.append(
            {
                "disposition": "INCONCLUSIVE" if unavailable else "SPEC_INVALID",
                "code": (
                    "FORMULA_CONVERTER_RUNTIME_UNAVAILABLE"
                    if unavailable
                    else "FORMULA_CONVERTER_RECEIPT_INVALID"
                ),
                "message": (
                    "Native Office Math receipt could not be verified by deterministic "
                    f"recompilation: {message}"
                ),
            }
        )
        record_status = (
            "SPEC_INVALID"
            if any(issue["disposition"] == "SPEC_INVALID" for issue in issues)
            else "INCONCLUSIVE"
        )
        return {
            "path": str(path),
            "sha256": actual_hash,
            "status": record_status,
        }, issues

    canonical_latex = str(formula.get("canonical_latex", ""))
    expected_values = {
        "document_type": "NATIVE_OFFICE_MATH_CONVERTER_RECEIPT",
        "formula_id": formula.get("id"),
        "mode": formula.get("mode"),
        "canonical_latex": canonical_latex,
        "latex_sha256": formula.get("latex_sha256"),
        "semantic_omml_profile": "office-math-semantic-v2",
        "native_target": {
            "kind": "office_math",
            "wrapper": "a14:m",
            "omml_root": "m:oMath" if formula.get("mode") == "inline" else "m:oMathPara",
        },
    }
    mismatches = {
        key: {"expected": expected, "actual": payload.get(key)}
        for key, expected in expected_values.items()
        if payload.get(key) != expected
    }
    if mismatches:
        issues.append(
            {
                "disposition": "SPEC_INVALID",
                "code": "FORMULA_CONVERTER_RECEIPT_BINDING_MISMATCH",
                "message": "Converter receipt is bound to different canonical formula content.",
                "evidence": {"mismatches": mismatches},
            }
        )
    target = (
        payload.get("native_target") if isinstance(payload.get("native_target"), Mapping) else {}
    )
    return (
        {
            "path": str(path),
            "sha256": actual_hash,
            "formula_id": payload.get("formula_id"),
            "status": "PASS" if not issues else "SPEC_INVALID",
            "latex_sha256": payload.get("latex_sha256"),
            "mathml_sha256": payload.get("mathml_sha256"),
            "omml_sha256": payload.get("omml_sha256"),
            "semantic_omml_profile": payload.get("semantic_omml_profile"),
            "semantic_omml_sha256": payload.get("semantic_omml_sha256"),
            "native_target": dict(target),
        },
        issues,
    )


def _bbox(element: Mapping[str, Any]) -> tuple[float, float, float, float] | None:
    value = element.get("bbox")
    if not isinstance(value, Mapping):
        return None
    try:
        x, y, width, height = (float(value[key]) for key in ("x", "y", "w", "h"))
    except (KeyError, TypeError, ValueError):
        return None
    if (
        not all(math.isfinite(number) for number in (x, y, width, height))
        or width <= 0
        or height <= 0
    ):
        return None
    return x, y, width, height


def _contains(
    parent: tuple[float, float, float, float], child: tuple[float, float, float, float]
) -> bool:
    px, py, pw, ph = parent
    cx, cy, cw, ch = child
    return (
        cx >= px - 1e-6
        and cy >= py - 1e-6
        and cx + cw <= px + pw + 1e-6
        and cy + ch <= py + ph + 1e-6
    )


def _intersection(
    first: tuple[float, float, float, float], second: tuple[float, float, float, float]
) -> tuple[float, float, float, float] | None:
    ax, ay, aw, ah = first
    bx, by, bw, bh = second
    left, top = max(ax, bx), max(ay, by)
    right, bottom = min(ax + aw, bx + bw), min(ay + ah, by + bh)
    if right - left <= 1e-6 or bottom - top <= 1e-6:
        return None
    return left, top, right - left, bottom - top


def _resolve_path(value: str | Path, base_dir: Path) -> Path:
    path = Path(value).expanduser()
    if not path.is_absolute():
        path = base_dir / path
    return path.resolve()


def _validate_instance(
    instance: Any,
    schema_path: Path,
    *,
    schema: Mapping[str, Any] | None = None,
) -> list[str]:
    loaded_schema = (
        schema
        if schema is not None
        else _strict_json_file(schema_path, label=f"JSON schema {schema_path}")
    )
    validator = Draft202012Validator(loaded_schema, format_checker=FormatChecker())
    messages: list[str] = []
    for error in sorted(validator.iter_errors(instance), key=lambda item: list(item.absolute_path)):
        location = "/".join(str(part) for part in error.absolute_path) or "$"
        messages.append(f"{location}: {error.message}")
    return messages


def _inspect_blank_canvas_pptx(path: Path) -> dict[str, Any]:
    presentation = Presentation(path)
    slide_count = len(presentation.slides)
    shape_count = sum(len(slide.shapes) for slide in presentation.slides)
    return {
        "path": str(path),
        "sha256": sha256_file(path),
        "slide_width_emu": int(presentation.slide_width),
        "slide_height_emu": int(presentation.slide_height),
        "slide_count": slide_count,
        "shape_count": shape_count,
        "aspect_ratio": float(presentation.slide_width) / float(presentation.slide_height),
    }


def _anchor_point(
    box: tuple[float, float, float, float], anchor: str
) -> tuple[float, float] | None:
    x, y, width, height = box
    anchors = {
        "top": (x + width / 2.0, y),
        "right": (x + width, y + height / 2.0),
        "bottom": (x + width / 2.0, y + height),
        "left": (x, y + height / 2.0),
        "center": (x + width / 2.0, y + height / 2.0),
    }
    return anchors.get(anchor)


def _inflate_rect(
    box: tuple[float, float, float, float], amount: float
) -> tuple[float, float, float, float]:
    x, y, width, height = box
    return x - amount, y - amount, width + 2.0 * amount, height + 2.0 * amount


def _segment_intersects_rect(
    first: tuple[float, float],
    second: tuple[float, float],
    box: tuple[float, float, float, float],
) -> bool:
    """Liang-Barsky segment/closed-rectangle intersection."""
    x0, y0 = first
    x1, y1 = second
    x, y, width, height = box
    dx, dy = x1 - x0, y1 - y0
    p = (-dx, dx, -dy, dy)
    q = (x0 - x, x + width - x0, y0 - y, y + height - y0)
    lower, upper = 0.0, 1.0
    for denominator, numerator in zip(p, q, strict=True):
        if math.isclose(denominator, 0.0, abs_tol=1e-12):
            if numerator < 0:
                return False
            continue
        ratio = numerator / denominator
        if denominator < 0:
            lower = max(lower, ratio)
        else:
            upper = min(upper, ratio)
        if lower > upper:
            return False
    return True


def _segments_cross(
    first_start: tuple[float, float],
    first_end: tuple[float, float],
    second_start: tuple[float, float],
    second_end: tuple[float, float],
) -> bool:
    if {first_start, first_end}.intersection({second_start, second_end}):
        return False

    def orientation(
        a: tuple[float, float], b: tuple[float, float], c: tuple[float, float]
    ) -> float:
        return (b[0] - a[0]) * (c[1] - a[1]) - (b[1] - a[1]) * (c[0] - a[0])

    o1 = orientation(first_start, first_end, second_start)
    o2 = orientation(first_start, first_end, second_end)
    o3 = orientation(second_start, second_end, first_start)
    o4 = orientation(second_start, second_end, first_end)
    return (o1 * o2 < -1e-9) and (o3 * o4 < -1e-9)


def _candidate_bbox_matches(
    element_box: tuple[float, float, float, float],
    candidate_box: Mapping[str, Any],
) -> tuple[bool, dict[str, float]]:
    candidate = tuple(float(candidate_box[key]) for key in ("x", "y", "w", "h"))
    overlap = _intersection(element_box, candidate)
    candidate_area = candidate[2] * candidate[3]
    covered = (
        0.0 if overlap is None or candidate_area <= 0 else overlap[2] * overlap[3] / candidate_area
    )
    center = (candidate[0] + candidate[2] / 2.0, candidate[1] + candidate[3] / 2.0)
    ex, ey, ew, eh = element_box
    center_inside = ex <= center[0] <= ex + ew and ey <= center[1] <= ey + eh
    return center_inside and covered >= 0.5, {"candidate_coverage": round(covered, 6)}


def _manifest_hashed_files(
    manifest: Mapping[str, Any], manifest_dir: Path
) -> list[tuple[str, Path, str]]:
    records: list[tuple[str, Path, str]] = []
    configuration = manifest.get("configuration")
    if isinstance(configuration, Mapping):
        for prefix, path_key, hash_key in (
            ("ocr_config", "path", "sha256"),
            ("perception_schema", "manifest_schema_path", "manifest_schema_sha256"),
        ):
            if configuration.get(path_key):
                records.append(
                    (
                        prefix,
                        _resolve_path(str(configuration[path_key]), manifest_dir),
                        str(configuration.get(hash_key, "")),
                    )
                )
    for script in manifest.get("scripts", []) or []:
        if isinstance(script, Mapping) and script.get("path"):
            records.append(
                (
                    f"script:{script.get('relative_path', script.get('path'))}",
                    _resolve_path(str(script["path"]), manifest_dir),
                    str(script.get("sha256", "")),
                )
            )
    models = manifest.get("models")
    if isinstance(models, Mapping):
        for role, model in models.items():
            if not isinstance(model, Mapping):
                continue
            for artifact in model.get("artifacts", []) or []:
                if isinstance(artifact, Mapping) and artifact.get("path"):
                    records.append(
                        (
                            f"model:{role}:{artifact.get('filename', '')}",
                            _resolve_path(str(artifact["path"]), manifest_dir),
                            str(artifact.get("sha256", "")),
                        )
                    )
    for stage in manifest.get("upstream_stages", []) or []:
        if not isinstance(stage, Mapping) or not stage.get("path"):
            continue
        stage_root = _resolve_path(str(stage["path"]), manifest_dir)
        for artifact in stage.get("files", []) or []:
            if isinstance(artifact, Mapping) and artifact.get("relative_path"):
                records.append(
                    (
                        f"upstream:{stage.get('name')}:{artifact['relative_path']}",
                        (stage_root / str(artifact["relative_path"])).resolve(),
                        str(artifact.get("sha256", "")),
                    )
                )
    return records


def _load_schema(
    schema: Mapping[str, Any] | None, schema_path: str | Path | None
) -> Mapping[str, Any]:
    if schema is not None:
        return schema
    path = Path(schema_path or DEFAULT_SCHEMA_PATH)
    loaded = _strict_json_file(path, label=f"figure schema {path}")
    if not isinstance(loaded, Mapping):
        raise ValueError(f"figure schema {path} root must be an object")
    return loaded


def preflight_scene(
    spec: Mapping[str, Any],
    *,
    source_path: str | Path | None = None,
    canvas_pptx_path: str | Path | None = None,
    schema: Mapping[str, Any] | None = None,
    schema_path: str | Path | None = None,
    base_dir: str | Path | None = None,
    font_search_paths: Sequence[str | Path] | None = None,
) -> dict[str, Any]:
    """Validate one render-ready scene and return a deterministic disposition."""
    findings: list[dict[str, Any]] = []
    text_measurements: list[dict[str, Any]] = []
    formula_measurements: list[dict[str, Any]] = []
    formula_converter_receipts: list[dict[str, Any]] = []
    finding_counter = 0
    evidence_snapshots: list[tuple[str, Path, bytes]] = []

    def snapshot_bytes(path: Path, *, label: str) -> bytes:
        payload = path.read_bytes()
        evidence_snapshots.append((label, path, payload))
        return payload

    schema_documents: dict[Path, Mapping[str, Any]] = {}
    schema_payloads: dict[Path, bytes] = {}

    def snapshot_schema(path: Path, *, label: str) -> Mapping[str, Any]:
        resolved = path.expanduser().resolve()
        cached = schema_documents.get(resolved)
        if cached is not None:
            return cached
        payload = snapshot_bytes(resolved, label=label)
        loaded = _strict_json_bytes(payload, label=label)
        if not isinstance(loaded, Mapping):
            raise ValueError(f"{label} root must be an object")
        schema_payloads[resolved] = payload
        schema_documents[resolved] = loaded
        return loaded

    def add_finding(
        disposition: str,
        code: str,
        message: str,
        *,
        element_ids: Iterable[str] = (),
        evidence: Mapping[str, Any] | None = None,
        repair: str | None = None,
    ) -> None:
        nonlocal finding_counter
        finding_counter += 1
        findings.append(
            {
                "id": f"PF{finding_counter:04d}",
                "severity": "INCONCLUSIVE" if disposition == "INCONCLUSIVE" else "MAJOR",
                "disposition": disposition,
                "code": code,
                "element_ids": list(element_ids),
                "message": message,
                "evidence": dict(evidence or {}),
                "repair": repair,
            }
        )

    try:
        loaded_schema = (
            schema
            if schema is not None
            else snapshot_schema(Path(schema_path or DEFAULT_SCHEMA_PATH), label="figure schema")
        )
        validator = Draft202012Validator(loaded_schema)
        for error in sorted(validator.iter_errors(spec), key=lambda item: list(item.absolute_path)):
            location = "/".join(str(part) for part in error.absolute_path) or "$"
            add_finding(
                "SPEC_INVALID",
                "SCHEMA_INVALID",
                f"{location}: {error.message}",
                evidence={"schema_path": location},
                repair="Correct the figure specification before planning or drawing.",
            )
    except (OSError, ValueError, json.JSONDecodeError) as exc:
        add_finding("SPEC_INVALID", "SCHEMA_UNAVAILABLE", str(exc))

    for canonical_schema_path, canonical_label in (
        (DEFAULT_PERCEPTION_SCHEMA_PATH, "perception manifest schema"),
        (DEFAULT_REVIEW_SCHEMA_PATH, "perception review schema"),
        (DEFAULT_GEOMETRY_SCHEMA_PATH, "geometry manifest schema"),
        (
            DEFAULT_HOST_RUNTIME_RECEIPT_SCHEMA_PATH,
            "host runtime receipt schema",
        ),
    ):
        try:
            snapshot_schema(canonical_schema_path, label=canonical_label)
        except (OSError, ValueError) as exc:
            add_finding(
                "SPEC_INVALID",
                "SCHEMA_UNAVAILABLE",
                str(exc),
                evidence={"schema_path": str(canonical_schema_path)},
            )

    non_finite_location = _non_finite_path(spec)
    if non_finite_location is not None:
        add_finding(
            "SPEC_INVALID",
            "NON_FINITE_JSON_NUMBER",
            "The figure specification contains a non-finite numeric value.",
            evidence={"path": non_finite_location},
            repair="Use finite standards-compliant JSON numbers only.",
        )

    resolved_base = Path(base_dir).resolve() if base_dir is not None else Path.cwd().resolve()
    source_spec = spec.get("source") if isinstance(spec.get("source"), Mapping) else {}
    requested_source = source_path if source_path is not None else source_spec.get("path")
    measured_source: dict[str, Any] | None = None
    if requested_source:
        requested_path = Path(str(requested_source)).expanduser()
        if not requested_path.is_absolute():
            requested_path = resolved_base / requested_path
        try:
            measured_source = inspect_png(requested_path)
            expected_hash = str(source_spec.get("sha256", "")).casefold()
            if measured_source["sha256"].casefold() != expected_hash:
                add_finding(
                    "SPEC_INVALID",
                    "SOURCE_HASH_MISMATCH",
                    "The reference PNG hash does not match the frozen figure specification.",
                    evidence={"expected": expected_hash, "actual": measured_source["sha256"]},
                    repair="Freeze a new spec for this exact PNG; do not reuse old perception or geometry.",
                )
            expected_size = (source_spec.get("width_px"), source_spec.get("height_px"))
            actual_size = (measured_source["width_px"], measured_source["height_px"])
            if expected_size != actual_size:
                add_finding(
                    "SPEC_INVALID",
                    "SOURCE_SIZE_MISMATCH",
                    "The measured PNG dimensions do not match the frozen figure specification.",
                    evidence={"expected": expected_size, "actual": actual_size},
                    repair="Regenerate the spec from script-measured source dimensions.",
                )
            expected_mode = str(source_spec.get("pixel_format", ""))
            if measured_source["pixel_mode"] != expected_mode:
                add_finding(
                    "SPEC_INVALID",
                    "SOURCE_PIXEL_MODE_MISMATCH",
                    "The measured PNG pixel mode does not match the frozen specification.",
                    evidence={"expected": expected_mode, "actual": measured_source["pixel_mode"]},
                    repair="Regenerate perception and the figure spec for the exact PNG pixel mode.",
                )
            canvas = spec.get("canvas")
            if isinstance(canvas, Mapping) and ("width_px" in canvas or "height_px" in canvas):
                canvas_size = (canvas.get("width_px"), canvas.get("height_px"))
                if canvas_size != actual_size:
                    add_finding(
                        "SPEC_INVALID",
                        "CANVAS_SOURCE_SIZE_MISMATCH",
                        "Canvas pixel dimensions must match the bound source dimensions.",
                        evidence={"canvas": canvas_size, "source": actual_size},
                    )
        except (OSError, PackageNotFoundError, ValueError) as exc:
            add_finding("SPEC_INVALID", "SOURCE_UNREADABLE", str(exc))
    else:
        add_finding(
            "SPEC_INVALID", "SOURCE_PATH_MISSING", "No source PNG path is available for preflight."
        )

    canvas_record: dict[str, Any] | None = None
    canvas_spec = spec.get("canvas") if isinstance(spec.get("canvas"), Mapping) else {}
    requested_canvas = (
        canvas_pptx_path if canvas_pptx_path is not None else canvas_spec.get("pptx_path")
    )
    if requested_canvas:
        canvas_path = _resolve_path(str(requested_canvas), resolved_base)
        try:
            canvas_record = _inspect_blank_canvas_pptx(canvas_path)
            expected_canvas_hash = str(canvas_spec.get("pptx_sha256", ""))
            if canvas_record["sha256"].casefold() != expected_canvas_hash.casefold():
                add_finding(
                    "SPEC_INVALID",
                    "CANVAS_PPTX_HASH_MISMATCH",
                    "The actual blank canvas deck is not the deck frozen in the figure spec.",
                    evidence={"expected": expected_canvas_hash, "actual": canvas_record["sha256"]},
                    repair="Use the exact create_canvas_pptx.py output bound to this spec.",
                )
            if canvas_record["slide_count"] != 1 or canvas_record["shape_count"] != 0:
                add_finding(
                    "SPEC_INVALID",
                    "CANVAS_NOT_BLANK",
                    "The pre-drawing PowerPoint canvas must contain exactly one empty slide.",
                    evidence={
                        "slide_count": canvas_record["slide_count"],
                        "shape_count": canvas_record["shape_count"],
                    },
                    repair="Regenerate the blank canvas before any Drawer mutation.",
                )
            expected_emu = (
                canvas_spec.get("slide_width_emu"),
                canvas_spec.get("slide_height_emu"),
            )
            actual_emu = (
                canvas_record["slide_width_emu"],
                canvas_record["slide_height_emu"],
            )
            if expected_emu != actual_emu:
                add_finding(
                    "SPEC_INVALID",
                    "CANVAS_PAGESETUP_MISMATCH",
                    "PowerPoint PageSetup dimensions do not match the frozen canvas dimensions.",
                    evidence={"expected_emu": expected_emu, "actual_emu": actual_emu},
                    repair="Recreate and rebind the canvas deck; do not resize it manually.",
                )
            if measured_source is not None:
                source_aspect = float(measured_source["width_px"]) / float(
                    measured_source["height_px"]
                )
                if not math.isclose(
                    canvas_record["aspect_ratio"], source_aspect, rel_tol=1e-7, abs_tol=1e-9
                ):
                    add_finding(
                        "SPEC_INVALID",
                        "CANVAS_ASPECT_MISMATCH",
                        "PowerPoint PageSetup aspect ratio does not match the measured source PNG.",
                        evidence={
                            "source_aspect": source_aspect,
                            "canvas_aspect": canvas_record["aspect_ratio"],
                        },
                        repair="Create the deck with create_canvas_pptx.py for this exact PNG.",
                    )
                canonical = slide_size_for_aspect(
                    measured_source["width_px"], measured_source["height_px"]
                )
                canonical_emu = (canonical["width_emu"], canonical["height_emu"])
                if actual_emu != canonical_emu:
                    add_finding(
                        "SPEC_INVALID",
                        "CANVAS_NONCANONICAL_SIZE",
                        "PowerPoint PageSetup is proportional but not the canonical AutoFigure size.",
                        evidence={"expected_emu": canonical_emu, "actual_emu": actual_emu},
                        repair="Regenerate the deck with the default create_canvas_pptx.py settings.",
                    )
        except (OSError, ValueError) as exc:
            add_finding(
                "SPEC_INVALID",
                "CANVAS_PPTX_UNREADABLE",
                f"Cannot read the bound blank canvas deck: {exc}",
            )
    else:
        add_finding(
            "SPEC_INVALID",
            "CANVAS_PPTX_MISSING",
            "Render-ready spec has no hash-bound blank PowerPoint canvas.",
        )

    perception_record: dict[str, Any] | None = None
    perception_manifest: Mapping[str, Any] | None = None
    perception_manifest_path: Path | None = None
    perception_manifest_sha256: str | None = None
    candidate_by_id: dict[str, Mapping[str, Any]] = {}
    decision_by_id: dict[str, Mapping[str, Any]] = {}
    perception_spec = spec.get("perception") if isinstance(spec.get("perception"), Mapping) else {}
    manifest_value = perception_spec.get("manifest_path")
    if manifest_value:
        manifest_path = _resolve_path(str(manifest_value), resolved_base)
        perception_manifest_path = manifest_path
        try:
            manifest_bytes = snapshot_bytes(manifest_path, label="perception manifest")
            actual_manifest_hash = _sha256_bytes(manifest_bytes)
            perception_manifest_sha256 = actual_manifest_hash
            expected_manifest_hash = str(perception_spec.get("manifest_sha256", ""))
            if actual_manifest_hash.casefold() != expected_manifest_hash.casefold():
                add_finding(
                    "SPEC_INVALID",
                    "PERCEPTION_MANIFEST_HASH_MISMATCH",
                    "Perception manifest hash does not match the frozen figure specification.",
                    evidence={"expected": expected_manifest_hash, "actual": actual_manifest_hash},
                    repair="Bind the spec to the exact reviewed perception manifest.",
                )
            loaded_manifest = _strict_json_bytes(manifest_bytes, label="perception manifest")
            if not isinstance(loaded_manifest, Mapping):
                raise ValueError("perception manifest root is not an object")
            perception_manifest = loaded_manifest
            schema_errors = _validate_instance(
                loaded_manifest,
                DEFAULT_PERCEPTION_SCHEMA_PATH,
                schema=snapshot_schema(
                    DEFAULT_PERCEPTION_SCHEMA_PATH,
                    label="perception manifest schema",
                ),
            )
            if schema_errors:
                add_finding(
                    "SPEC_INVALID",
                    "PERCEPTION_MANIFEST_SCHEMA_INVALID",
                    "The raw perception manifest does not satisfy its checked-in schema.",
                    evidence={"errors": schema_errors[:8], "error_count": len(schema_errors)},
                    repair="Rerun the pinned perception gate; do not hand-author the manifest.",
                )
            manifest_source = (
                loaded_manifest.get("source")
                if isinstance(loaded_manifest.get("source"), Mapping)
                else {}
            )
            manifest_source_hash = str(manifest_source.get("sha256", ""))
            spec_source_hash = str(source_spec.get("sha256", ""))
            if manifest_source_hash.casefold() != spec_source_hash.casefold():
                add_finding(
                    "SPEC_INVALID",
                    "PERCEPTION_SOURCE_MISMATCH",
                    "Perception manifest is bound to a different reference image.",
                    evidence={
                        "spec_source": spec_source_hash,
                        "manifest_source": manifest_source_hash,
                    },
                    repair="Rerun perception for this exact source and invalidate the old spec.",
                )
            for field, expected in (
                ("width_px", source_spec.get("width_px")),
                ("height_px", source_spec.get("height_px")),
                ("pixel_mode", source_spec.get("pixel_format")),
            ):
                if manifest_source.get(field) != expected:
                    add_finding(
                        "SPEC_INVALID",
                        "PERCEPTION_SOURCE_METADATA_MISMATCH",
                        f"Perception source field {field!r} differs from the frozen source.",
                        evidence={
                            "field": field,
                            "expected": expected,
                            "actual": manifest_source.get(field),
                        },
                    )
            if loaded_manifest.get("status") != "OCR_HYPOTHESES_REVIEW_REQUIRED":
                add_finding(
                    "SPEC_INVALID",
                    "PERCEPTION_RAW_GATE_NOT_READY",
                    "Only a non-degraded OCR_HYPOTHESES_REVIEW_REQUIRED manifest may enter review.",
                    evidence={"status": loaded_manifest.get("status")},
                )
            acceptance = loaded_manifest.get("acceptance_checks")
            if not isinstance(acceptance, Mapping) or acceptance.get("passed") is not True:
                add_finding(
                    "SPEC_INVALID",
                    "PERCEPTION_ACCEPTANCE_FAILED",
                    "The raw perception run did not pass its deterministic acceptance checks.",
                )
            raw_candidates = loaded_manifest.get("text_candidates", [])
            if isinstance(raw_candidates, list):
                raw_ids = [
                    str(item.get("candidate_id"))
                    for item in raw_candidates
                    if isinstance(item, Mapping)
                ]
                if len(raw_ids) != len(set(raw_ids)):
                    add_finding(
                        "SPEC_INVALID",
                        "PERCEPTION_DUPLICATE_CANDIDATE_ID",
                        "Raw perception candidate IDs are not unique.",
                    )
                candidate_by_id = {
                    str(item.get("candidate_id")): item
                    for item in raw_candidates
                    if isinstance(item, Mapping)
                }
                summary = loaded_manifest.get("summary")
                if not isinstance(summary, Mapping) or summary.get("candidate_count") != len(
                    raw_candidates
                ):
                    add_finding(
                        "SPEC_INVALID",
                        "PERCEPTION_CANDIDATE_COUNT_MISMATCH",
                        "Raw perception candidate count is internally inconsistent.",
                    )
            for label, artifact_path, expected_hash in _manifest_hashed_files(
                loaded_manifest, manifest_path.parent
            ):
                try:
                    actual_hash = sha256_file(artifact_path)
                except OSError as exc:
                    add_finding(
                        "SPEC_INVALID",
                        "PERCEPTION_ARTIFACT_UNREADABLE",
                        f"Cannot read a hash-bound perception artifact: {exc}",
                        evidence={"label": label, "path": str(artifact_path)},
                    )
                    continue
                if actual_hash.casefold() != expected_hash.casefold():
                    add_finding(
                        "SPEC_INVALID",
                        "PERCEPTION_ARTIFACT_HASH_MISMATCH",
                        "A script, model, config, schema, or upstream artifact changed after perception.",
                        evidence={
                            "label": label,
                            "path": str(artifact_path),
                            "expected": expected_hash,
                            "actual": actual_hash,
                        },
                        repair="Invalidate this run and rerun perception with the pinned files.",
                    )

            receipt_value = perception_spec.get("review_receipt_path")
            if not receipt_value:
                add_finding(
                    "SPEC_INVALID",
                    "PERCEPTION_REVIEW_RECEIPT_MISSING",
                    "Raw OCR hypotheses cannot authorize drawing without a separate review receipt.",
                )
            else:
                receipt_path = _resolve_path(str(receipt_value), resolved_base)
                receipt_bytes = snapshot_bytes(receipt_path, label="perception review receipt")
                actual_receipt_hash = _sha256_bytes(receipt_bytes)
                expected_receipt_hash = str(perception_spec.get("review_receipt_sha256", ""))
                if actual_receipt_hash.casefold() != expected_receipt_hash.casefold():
                    add_finding(
                        "SPEC_INVALID",
                        "PERCEPTION_REVIEW_HASH_MISMATCH",
                        "Perception review receipt hash does not match the figure spec.",
                        evidence={"expected": expected_receipt_hash, "actual": actual_receipt_hash},
                    )
                loaded_receipt = _strict_json_bytes(
                    receipt_bytes, label="perception review receipt"
                )
                if not isinstance(loaded_receipt, Mapping):
                    raise ValueError("perception review receipt root is not an object")
                receipt_errors = _validate_instance(
                    loaded_receipt,
                    DEFAULT_REVIEW_SCHEMA_PATH,
                    schema=snapshot_schema(
                        DEFAULT_REVIEW_SCHEMA_PATH,
                        label="perception review schema",
                    ),
                )
                if receipt_errors:
                    add_finding(
                        "SPEC_INVALID",
                        "PERCEPTION_REVIEW_SCHEMA_INVALID",
                        "The perception review receipt does not satisfy its checked-in schema.",
                        evidence={"errors": receipt_errors[:8], "error_count": len(receipt_errors)},
                    )
                if loaded_receipt.get("document_type") != "PERCEPTION_REVIEW_RECEIPT":
                    add_finding(
                        "SPEC_INVALID",
                        "PERCEPTION_REVIEW_DOCUMENT_TYPE_INVALID",
                        "The bound review document is not a finalized receipt.",
                    )
                binding = (
                    loaded_receipt.get("raw_manifest")
                    if isinstance(loaded_receipt.get("raw_manifest"), Mapping)
                    else {}
                )
                expected_raw_schema_hash = _sha256_bytes(
                    schema_payloads[DEFAULT_PERCEPTION_SCHEMA_PATH.resolve()]
                )
                binding_mismatches = {
                    "manifest_sha256": (binding.get("manifest_sha256"), actual_manifest_hash),
                    "manifest_schema_sha256": (
                        binding.get("manifest_schema_sha256"),
                        expected_raw_schema_hash,
                    ),
                    "run_id": (binding.get("run_id"), loaded_manifest.get("run_id")),
                    "source_sha256": (binding.get("source_sha256"), manifest_source_hash),
                }
                bad_binding = {
                    key: {"receipt": values[0], "expected": values[1]}
                    for key, values in binding_mismatches.items()
                    if str(values[0]).casefold() != str(values[1]).casefold()
                }
                if bad_binding:
                    add_finding(
                        "SPEC_INVALID",
                        "PERCEPTION_REVIEW_BINDING_MISMATCH",
                        "Review receipt is bound to a different raw perception run.",
                        evidence=bad_binding,
                    )
                receipt_decisions = loaded_receipt.get("decisions", [])
                if isinstance(receipt_decisions, list):
                    decision_ids = [
                        str(item.get("candidate_id"))
                        for item in receipt_decisions
                        if isinstance(item, Mapping)
                    ]
                    decision_by_id = {
                        str(item.get("candidate_id")): item
                        for item in receipt_decisions
                        if isinstance(item, Mapping)
                    }
                    if len(decision_ids) != len(set(decision_ids)) or set(decision_ids) != set(
                        candidate_by_id
                    ):
                        add_finding(
                            "SPEC_INVALID",
                            "PERCEPTION_REVIEW_COVERAGE_MISMATCH",
                            "Review receipt must contain exactly one decision for every raw candidate.",
                            evidence={
                                "raw_candidate_count": len(candidate_by_id),
                                "decision_count": len(decision_ids),
                            },
                        )
                    for candidate_id, candidate in candidate_by_id.items():
                        decision = decision_by_id.get(candidate_id)
                        if decision is None:
                            continue
                        expected_formula_like = any(
                            "FORMULA_LIKE" in str(flag).upper()
                            for flag in candidate.get("review_flags", []) or []
                        )
                        snapshot_pairs = {
                            "ocr_text": (decision.get("ocr_text"), candidate.get("text")),
                            "ocr_confidence": (
                                decision.get("ocr_confidence"),
                                candidate.get("ocr_confidence"),
                            ),
                            "review_flags": (
                                decision.get("review_flags"),
                                candidate.get("review_flags"),
                            ),
                            "formula_like": (decision.get("formula_like"), expected_formula_like),
                        }
                        snapshot_mismatch = {
                            key: {"receipt": values[0], "raw": values[1]}
                            for key, values in snapshot_pairs.items()
                            if values[0] != values[1]
                        }
                        if snapshot_mismatch:
                            add_finding(
                                "SPEC_INVALID",
                                "PERCEPTION_REVIEW_SNAPSHOT_MISMATCH",
                                "A review decision snapshot differs from its raw OCR candidate.",
                                evidence={
                                    "candidate_id": candidate_id,
                                    "fields": snapshot_mismatch,
                                },
                            )
                        status_value = str(decision.get("status", ""))
                        evidence_value = decision.get("evidence")
                        evidence_kind = (
                            evidence_value.get("kind")
                            if isinstance(evidence_value, Mapping)
                            else None
                        )
                        if status_value in {
                            "CONFIRMED",
                            "CORRECTED",
                            "NOT_TEXT",
                            "FORMULA_CONFIRMED",
                        } and evidence_kind not in {
                            "user_confirmed",
                            "source_text",
                        }:
                            add_finding(
                                "SPEC_INVALID",
                                "PERCEPTION_REVIEW_AUTHORITY_MISSING",
                                "Terminal review decisions require user_confirmed or source_text evidence.",
                                evidence={"candidate_id": candidate_id, "status": status_value},
                            )
                        if status_value == "CONFIRMED" and decision.get(
                            "confirmed_text"
                        ) != candidate.get("text"):
                            add_finding(
                                "SPEC_INVALID",
                                "PERCEPTION_REVIEW_CONFIRMED_TEXT_CHANGED",
                                "Changed OCR text must use CORRECTED, not CONFIRMED.",
                                evidence={"candidate_id": candidate_id},
                            )
                        if expected_formula_like and status_value not in {
                            "FORMULA_CONFIRMED",
                            "NOT_TEXT",
                            "PENDING",
                            "INCONCLUSIVE",
                        }:
                            add_finding(
                                "SPEC_INVALID",
                                "PERCEPTION_REVIEW_FORMULA_POLICY_VIOLATION",
                                "Formula-like candidates cannot be promoted as ordinary text.",
                                evidence={"candidate_id": candidate_id, "status": status_value},
                            )
                counts = (
                    loaded_receipt.get("counts")
                    if isinstance(loaded_receipt.get("counts"), Mapping)
                    else {}
                )
                unresolved = [
                    item_id
                    for item_id, item in decision_by_id.items()
                    if item.get("status") in {"PENDING", "INCONCLUSIVE"}
                ]
                terminal = [
                    item_id
                    for item_id, item in decision_by_id.items()
                    if item.get("status")
                    in {"CONFIRMED", "CORRECTED", "NOT_TEXT", "FORMULA_CONFIRMED"}
                ]
                status_counts = Counter(
                    str(item.get("status", "")) for item in decision_by_id.values()
                )
                recomputed_counts = {
                    "total_candidates": len(candidate_by_id),
                    "terminal_count": len(terminal),
                    "confirmed_count": status_counts["CONFIRMED"],
                    "corrected_count": status_counts["CORRECTED"],
                    "not_text_count": status_counts["NOT_TEXT"],
                    "formula_confirmed_count": status_counts["FORMULA_CONFIRMED"],
                    "pending_count": status_counts["PENDING"],
                    "inconclusive_count": status_counts["INCONCLUSIVE"],
                    "unresolved_count": len(unresolved),
                }
                count_mismatches = {
                    key: {"receipt": counts.get(key), "recomputed": value}
                    for key, value in recomputed_counts.items()
                    if counts.get(key) != value
                }
                if count_mismatches:
                    add_finding(
                        "SPEC_INVALID",
                        "PERCEPTION_REVIEW_COUNT_MISMATCH",
                        "Review receipt counts do not match its exact decision set.",
                        evidence=count_mismatches,
                    )
                receipt_pass = (
                    loaded_receipt.get("status") == "PERCEPTION_REVIEW_PASS"
                    and not unresolved
                    and not loaded_receipt.get("missing_candidate_ids")
                    and not loaded_receipt.get("unresolved_candidate_ids")
                    and not loaded_receipt.get("gate_blockers")
                    and counts.get("unresolved_count") == 0
                    and counts.get("missing_count") == 0
                    and counts.get("gate_blocker_count") == 0
                    and counts.get("terminal_count") == len(terminal) == len(candidate_by_id)
                    and counts.get("total_candidates") == len(candidate_by_id)
                )
                if not receipt_pass:
                    add_finding(
                        "INCONCLUSIVE",
                        "PERCEPTION_REVIEW_NOT_PASS",
                        "Perception review has unresolved candidates or raw-gate blockers.",
                        evidence={
                            "receipt_status": loaded_receipt.get("status"),
                            "unresolved_ids": unresolved,
                            "counts": dict(counts),
                        },
                        repair="Resolve every candidate using authoritative evidence and regenerate the receipt.",
                    )
                perception_record = {
                    "path": str(manifest_path),
                    "sha256": actual_manifest_hash,
                    "manifest_schema_sha256": _sha256_bytes(
                        schema_payloads[DEFAULT_PERCEPTION_SCHEMA_PATH.resolve()]
                    ),
                    "source_sha256": manifest_source_hash,
                    "run_id": loaded_manifest.get("run_id"),
                    "review_receipt_path": str(receipt_path),
                    "review_receipt_sha256": actual_receipt_hash,
                    "review_schema_sha256": _sha256_bytes(
                        schema_payloads[DEFAULT_REVIEW_SCHEMA_PATH.resolve()]
                    ),
                    "review_status": loaded_receipt.get("status"),
                }
        except (OSError, ValueError, json.JSONDecodeError) as exc:
            add_finding(
                "SPEC_INVALID",
                "PERCEPTION_EVIDENCE_UNREADABLE",
                f"Cannot read or validate the bound perception evidence: {exc}",
                repair="Provide the exact raw manifest and final review receipt before preflight.",
            )
    else:
        add_finding(
            "SPEC_INVALID",
            "PERCEPTION_MANIFEST_MISSING",
            "Render-ready spec has no bound perception manifest.",
        )

    geometry_record: dict[str, Any] | None = None
    geometry_diagnostics: list[dict[str, Any]] = []
    geometry_by_candidate: dict[str, Mapping[str, Any]] = {}
    geometry_spec_value = spec.get("geometry")
    if geometry_spec_value is not None and isinstance(geometry_spec_value, Mapping):
        geometry_manifest_value = geometry_spec_value.get("manifest_path")
        if not geometry_manifest_value:
            add_finding(
                "SPEC_INVALID",
                "GEOMETRY_MANIFEST_MISSING",
                "The geometry binding has no manifest path.",
            )
        else:
            geometry_manifest_path = _resolve_path(str(geometry_manifest_value), resolved_base)
            try:
                geometry_manifest_bytes = snapshot_bytes(
                    geometry_manifest_path, label="geometry manifest"
                )
                geometry_manifest_hash = _sha256_bytes(geometry_manifest_bytes)
                expected_geometry_hash = str(geometry_spec_value.get("manifest_sha256", ""))
                if geometry_manifest_hash.casefold() != expected_geometry_hash.casefold():
                    add_finding(
                        "SPEC_INVALID",
                        "GEOMETRY_MANIFEST_HASH_MISMATCH",
                        "Geometry manifest hash does not match the frozen figure specification.",
                        evidence={
                            "expected": expected_geometry_hash,
                            "actual": geometry_manifest_hash,
                        },
                        repair="Bind the exact Phase-1 geometry manifest or rerun the stage.",
                    )
                loaded_geometry = _strict_json_bytes(
                    geometry_manifest_bytes, label="geometry manifest"
                )
                if not isinstance(loaded_geometry, Mapping):
                    raise ValueError("geometry manifest root is not an object")

                geometry_schema_errors = _validate_instance(
                    loaded_geometry,
                    DEFAULT_GEOMETRY_SCHEMA_PATH,
                    schema=snapshot_schema(
                        DEFAULT_GEOMETRY_SCHEMA_PATH,
                        label="geometry manifest schema",
                    ),
                )
                if geometry_schema_errors:
                    add_finding(
                        "SPEC_INVALID",
                        "GEOMETRY_MANIFEST_SCHEMA_INVALID",
                        "Geometry observations do not satisfy the checked-in schema.",
                        evidence={
                            "errors": geometry_schema_errors[:8],
                            "error_count": len(geometry_schema_errors),
                        },
                        repair="Rerun geometry_refinement.py; do not hand-author observations.",
                    )

                text_geometry = loaded_geometry.get("text_geometry")
                if isinstance(text_geometry, list):
                    geometry_ids = [
                        str(item.get("candidate_id"))
                        for item in text_geometry
                        if isinstance(item, Mapping)
                    ]
                    geometry_by_candidate = {
                        str(item.get("candidate_id")): item
                        for item in text_geometry
                        if isinstance(item, Mapping)
                    }
                    if len(geometry_ids) != len(set(geometry_ids)) or set(geometry_ids) != set(
                        candidate_by_id
                    ):
                        add_finding(
                            "SPEC_INVALID",
                            "GEOMETRY_CANDIDATE_COVERAGE_MISMATCH",
                            "Geometry must contain exactly one record for every OCR candidate.",
                            evidence={
                                "geometry_ids": sorted(geometry_ids),
                                "ocr_ids": sorted(candidate_by_id),
                            },
                        )
                    snapshot_mismatches = {}
                    for candidate_id, geometry_item in geometry_by_candidate.items():
                        ocr_candidate = candidate_by_id.get(candidate_id)
                        if ocr_candidate is None:
                            continue
                        fields = {
                            "primary_observation_id": (
                                geometry_item.get("primary_observation_id"),
                                ocr_candidate.get("primary_observation_id"),
                            ),
                            "text": (
                                geometry_item.get("text"),
                                ocr_candidate.get("text"),
                            ),
                            "ocr_confidence": (
                                geometry_item.get("ocr_confidence"),
                                ocr_candidate.get("ocr_confidence"),
                            ),
                            "input_bbox_source": (
                                geometry_item.get("input_bbox_source"),
                                ocr_candidate.get("bbox_source"),
                            ),
                            "input_polygon_source": (
                                geometry_item.get("input_polygon_source"),
                                ocr_candidate.get("polygon_source"),
                            ),
                        }
                        bad_fields = {
                            field: {"geometry": values[0], "ocr": values[1]}
                            for field, values in fields.items()
                            if values[0] != values[1]
                        }
                        if bad_fields:
                            snapshot_mismatches[candidate_id] = bad_fields
                    if snapshot_mismatches:
                        add_finding(
                            "SPEC_INVALID",
                            "GEOMETRY_CANDIDATE_SNAPSHOT_MISMATCH",
                            "Geometry candidate snapshots differ from the bound OCR manifest.",
                            evidence=snapshot_mismatches,
                        )

                geometry_summary = (
                    loaded_geometry.get("summary")
                    if isinstance(loaded_geometry.get("summary"), Mapping)
                    else {}
                )
                geometry_pairs = (
                    loaded_geometry.get("neighbor_pairs")
                    if isinstance(loaded_geometry.get("neighbor_pairs"), list)
                    else []
                )
                geometry_frames = (
                    loaded_geometry.get("frame_candidates")
                    if isinstance(loaded_geometry.get("frame_candidates"), list)
                    else []
                )
                measured_ink_count = sum(
                    item.get("status") == "MEASURED" for item in geometry_by_candidate.values()
                )
                measured_alignment_count = sum(
                    isinstance(item.get("baseline"), Mapping)
                    and item["baseline"].get("status") == "MEASURED"
                    for item in geometry_by_candidate.values()
                )
                measured_frame_count = sum(
                    isinstance(item, Mapping) and item.get("status") == "MEASURED"
                    for item in geometry_frames
                )
                recomputed_geometry_summary = {
                    "candidate_count": len(geometry_by_candidate),
                    "measured_ink_count": measured_ink_count,
                    "inconclusive_ink_count": len(geometry_by_candidate) - measured_ink_count,
                    "reliable_ink_bottom_alignment_count": measured_alignment_count,
                    "neighbor_pair_count": len(geometry_pairs),
                    "frame_candidate_count": len(geometry_frames),
                    "measured_frame_count": measured_frame_count,
                    "degradations": loaded_geometry.get("degradations"),
                }
                geometry_summary_mismatches = {
                    field: {
                        "manifest": geometry_summary.get(field),
                        "recomputed": expected,
                    }
                    for field, expected in recomputed_geometry_summary.items()
                    if geometry_summary.get(field) != expected
                }
                if geometry_summary_mismatches:
                    add_finding(
                        "SPEC_INVALID",
                        "GEOMETRY_SUMMARY_MISMATCH",
                        "Geometry summary does not match its observation arrays.",
                        evidence=geometry_summary_mismatches,
                    )
                pair_binding_issues = []
                pair_ids = []
                for pair in geometry_pairs:
                    if not isinstance(pair, Mapping):
                        continue
                    pair_ids.append(str(pair.get("pair_id")))
                    endpoints = [
                        str(pair.get("candidate_a_id")),
                        str(pair.get("candidate_b_id")),
                    ]
                    invalid_endpoints = []
                    for candidate_id in endpoints:
                        item = geometry_by_candidate.get(candidate_id)
                        baseline = item.get("baseline") if isinstance(item, Mapping) else None
                        if (
                            not isinstance(item, Mapping)
                            or item.get("status") != "MEASURED"
                            or not isinstance(baseline, Mapping)
                            or baseline.get("status") != "MEASURED"
                        ):
                            invalid_endpoints.append(candidate_id)
                    if invalid_endpoints:
                        pair_binding_issues.append(
                            {
                                "pair_id": pair.get("pair_id"),
                                "invalid_endpoints": invalid_endpoints,
                            }
                        )
                if len(pair_ids) != len(set(pair_ids)) or pair_binding_issues:
                    add_finding(
                        "SPEC_INVALID",
                        "GEOMETRY_PAIR_BINDING_MISMATCH",
                        "Neighbor gaps must bind unique, measured, reliable-alignment candidates.",
                        evidence={
                            "duplicate_pair_ids": len(pair_ids) != len(set(pair_ids)),
                            "issues": pair_binding_issues,
                        },
                    )

                geometry_status = str(loaded_geometry.get("status", ""))
                if geometry_status not in {
                    "GEOMETRY_OBSERVATIONS_READY",
                    "GEOMETRY_INCONCLUSIVE",
                }:
                    add_finding(
                        "SPEC_INVALID",
                        "GEOMETRY_STATUS_INVALID",
                        "Geometry manifest has an unknown Phase-1 status.",
                        evidence={"status": geometry_status},
                    )
                recomputed_geometry_status = (
                    "GEOMETRY_OBSERVATIONS_READY"
                    if measured_ink_count or measured_frame_count
                    else "GEOMETRY_INCONCLUSIVE"
                )
                if geometry_status != recomputed_geometry_status:
                    add_finding(
                        "SPEC_INVALID",
                        "GEOMETRY_STATUS_COUNT_MISMATCH",
                        "Geometry status is inconsistent with its measured observation counts.",
                        evidence={
                            "status": geometry_status,
                            "recomputed": recomputed_geometry_status,
                            "measured_ink_count": measured_ink_count,
                            "measured_frame_count": measured_frame_count,
                        },
                    )
                geometry_mode = str(loaded_geometry.get("mode", ""))
                geometry_policy = (
                    loaded_geometry.get("policy")
                    if isinstance(loaded_geometry.get("policy"), Mapping)
                    else {}
                )
                promotion_allowed = geometry_policy.get("promotion_allowed")
                if geometry_mode != "observation_only" or promotion_allowed is not False:
                    add_finding(
                        "SPEC_INVALID",
                        "GEOMETRY_AUTHORITY_BOUNDARY_VIOLATION",
                        "Phase-1 geometry must remain observation-only and cannot authorize promotion.",
                        evidence={
                            "mode": geometry_mode,
                            "promotion_allowed": promotion_allowed,
                        },
                        repair="Regenerate the manifest with promotion_allowed=false.",
                    )

                geometry_source = (
                    loaded_geometry.get("source")
                    if isinstance(loaded_geometry.get("source"), Mapping)
                    else {}
                )
                source_mismatches = {}
                for field, expected in (
                    ("sha256", source_spec.get("sha256")),
                    ("width_px", source_spec.get("width_px")),
                    ("height_px", source_spec.get("height_px")),
                    ("pixel_mode", source_spec.get("pixel_format")),
                ):
                    actual = geometry_source.get(field)
                    matches = (
                        str(actual).casefold() == str(expected).casefold()
                        if field == "sha256"
                        else actual == expected
                    )
                    if not matches:
                        source_mismatches[field] = {
                            "geometry": actual,
                            "expected": expected,
                        }
                if source_mismatches:
                    add_finding(
                        "SPEC_INVALID",
                        "GEOMETRY_SOURCE_MISMATCH",
                        "Geometry observations belong to a different source image.",
                        evidence=source_mismatches,
                        repair="Rerun Phase-1 geometry for the exact frozen PNG.",
                    )

                geometry_inputs = (
                    loaded_geometry.get("inputs")
                    if isinstance(loaded_geometry.get("inputs"), Mapping)
                    else {}
                )
                ocr_binding = (
                    geometry_inputs.get("ocr_manifest")
                    if isinstance(geometry_inputs.get("ocr_manifest"), Mapping)
                    else {}
                )
                expected_ocr_binding = {
                    "sha256": perception_manifest_sha256,
                    "schema_version": (
                        perception_manifest.get("schema_version")
                        if isinstance(perception_manifest, Mapping)
                        else None
                    ),
                    "run_id": (
                        perception_manifest.get("run_id")
                        if isinstance(perception_manifest, Mapping)
                        else None
                    ),
                    "source_sha256": (
                        (perception_manifest.get("source") or {}).get("sha256")
                        if isinstance(perception_manifest, Mapping)
                        and isinstance(perception_manifest.get("source"), Mapping)
                        else None
                    ),
                }
                ocr_binding_mismatches = {}
                for field, expected in expected_ocr_binding.items():
                    actual = ocr_binding.get(field)
                    if str(actual).casefold() != str(expected).casefold():
                        ocr_binding_mismatches[field] = {
                            "geometry": actual,
                            "expected": expected,
                        }
                if ocr_binding_mismatches:
                    add_finding(
                        "SPEC_INVALID",
                        "GEOMETRY_OCR_BINDING_MISMATCH",
                        "Geometry observations are not bound to the current OCR manifest.",
                        evidence=ocr_binding_mismatches,
                        repair="Rerun geometry after the final OCR manifest is frozen.",
                    )
                if str(loaded_geometry.get("run_id", "")) != str(expected_ocr_binding["run_id"]):
                    add_finding(
                        "SPEC_INVALID",
                        "GEOMETRY_RUN_MISMATCH",
                        "Geometry and OCR manifests have different run IDs.",
                        evidence={
                            "geometry": loaded_geometry.get("run_id"),
                            "ocr": expected_ocr_binding["run_id"],
                        },
                    )

                if ocr_binding.get("path"):
                    bound_ocr_path = _resolve_path(
                        str(ocr_binding["path"]), geometry_manifest_path.parent
                    )
                    try:
                        bound_ocr_bytes = snapshot_bytes(
                            bound_ocr_path, label="geometry-bound OCR manifest"
                        )
                        bound_ocr_hash = _sha256_bytes(bound_ocr_bytes)
                        bound_ocr_size = len(bound_ocr_bytes)
                        if bound_ocr_hash.casefold() != str(
                            ocr_binding.get("sha256", "")
                        ).casefold() or bound_ocr_size != ocr_binding.get("size_bytes"):
                            add_finding(
                                "SPEC_INVALID",
                                "GEOMETRY_OCR_FILE_BINDING_MISMATCH",
                                "The OCR file read by geometry changed after refinement.",
                                evidence={
                                    "path": str(bound_ocr_path),
                                    "sha256": bound_ocr_hash,
                                    "size_bytes": bound_ocr_size,
                                },
                            )
                        if (
                            perception_manifest_path is not None
                            and bound_ocr_path != perception_manifest_path
                        ):
                            add_finding(
                                "SPEC_INVALID",
                                "GEOMETRY_OCR_PATH_MISMATCH",
                                "Geometry is bound to a different OCR manifest path.",
                                evidence={
                                    "geometry": str(bound_ocr_path),
                                    "figure_spec": str(perception_manifest_path),
                                },
                            )
                    except OSError as exc:
                        add_finding(
                            "SPEC_INVALID",
                            "GEOMETRY_OCR_FILE_UNREADABLE",
                            f"Cannot read geometry's bound OCR manifest: {exc}",
                        )

                host_binding = (
                    geometry_inputs.get("host_runtime_receipt")
                    if isinstance(geometry_inputs.get("host_runtime_receipt"), Mapping)
                    else {}
                )
                host_receipt_path: Path | None = None
                host_receipt_hash: str | None = None
                host_receipt: Mapping[str, Any] | None = None
                if host_binding.get("path"):
                    host_receipt_path = _resolve_path(
                        str(host_binding["path"]), geometry_manifest_path.parent
                    )
                    try:
                        host_receipt_bytes = snapshot_bytes(
                            host_receipt_path, label="host runtime receipt"
                        )
                        host_receipt_hash = _sha256_bytes(host_receipt_bytes)
                        host_receipt_size = len(host_receipt_bytes)
                        if host_receipt_hash.casefold() != str(
                            host_binding.get("sha256", "")
                        ).casefold() or host_receipt_size != host_binding.get("size_bytes"):
                            add_finding(
                                "SPEC_INVALID",
                                "GEOMETRY_HOST_RECEIPT_HASH_MISMATCH",
                                "The host runtime receipt changed after geometry refinement.",
                                evidence={
                                    "path": str(host_receipt_path),
                                    "sha256": host_receipt_hash,
                                    "size_bytes": host_receipt_size,
                                },
                            )
                        loaded_host_receipt = _strict_json_bytes(
                            host_receipt_bytes, label="host runtime receipt"
                        )
                        if not isinstance(loaded_host_receipt, Mapping):
                            raise ValueError("host runtime receipt root is not an object")
                        host_receipt = loaded_host_receipt
                        host_schema_errors = _validate_instance(
                            host_receipt,
                            DEFAULT_HOST_RUNTIME_RECEIPT_SCHEMA_PATH,
                            schema=snapshot_schema(
                                DEFAULT_HOST_RUNTIME_RECEIPT_SCHEMA_PATH,
                                label="host runtime receipt schema",
                            ),
                        )
                        if host_schema_errors:
                            add_finding(
                                "SPEC_INVALID",
                                "GEOMETRY_HOST_RECEIPT_SCHEMA_INVALID",
                                "Geometry's host runtime receipt is not schema-valid.",
                                evidence={
                                    "errors": host_schema_errors[:8],
                                    "error_count": len(host_schema_errors),
                                },
                            )
                        receipt_bindings = (
                            host_receipt.get("bindings")
                            if isinstance(host_receipt.get("bindings"), Mapping)
                            else {}
                        )
                        for binding_name, binding in receipt_bindings.items():
                            if not isinstance(binding, Mapping) or not binding.get("path"):
                                continue
                            bound_path = _resolve_path(
                                str(binding["path"]), host_receipt_path.parent
                            )
                            try:
                                bound_bytes = snapshot_bytes(
                                    bound_path,
                                    label=f"host runtime binding {binding_name}",
                                )
                                bound_hash = _sha256_bytes(bound_bytes)
                                bound_size = len(bound_bytes)
                                if bound_hash.casefold() != str(
                                    binding.get("sha256", "")
                                ).casefold() or bound_size != binding.get("size_bytes"):
                                    add_finding(
                                        "SPEC_INVALID",
                                        "GEOMETRY_HOST_RECEIPT_INTERNAL_BINDING_MISMATCH",
                                        "A file bound by the host runtime receipt changed.",
                                        evidence={
                                            "binding": str(binding_name),
                                            "path": str(bound_path),
                                            "sha256": bound_hash,
                                            "size_bytes": bound_size,
                                        },
                                    )
                                if (
                                    binding_name == "receipt_schema"
                                    and bound_hash.casefold()
                                    != _sha256_bytes(
                                        schema_payloads[
                                            DEFAULT_HOST_RUNTIME_RECEIPT_SCHEMA_PATH.resolve()
                                        ]
                                    ).casefold()
                                ):
                                    add_finding(
                                        "SPEC_INVALID",
                                        "GEOMETRY_HOST_SCHEMA_NOT_CURRENT",
                                        "Host receipt uses a different runtime-receipt schema.",
                                    )
                            except OSError as exc:
                                add_finding(
                                    "SPEC_INVALID",
                                    "GEOMETRY_HOST_RECEIPT_INTERNAL_BINDING_UNREADABLE",
                                    f"Cannot read a host receipt binding: {exc}",
                                    evidence={"binding": str(binding_name)},
                                )
                    except (OSError, ValueError, json.JSONDecodeError) as exc:
                        add_finding(
                            "SPEC_INVALID",
                            "GEOMETRY_HOST_RECEIPT_UNREADABLE",
                            f"Cannot read geometry's bound host runtime receipt: {exc}",
                        )
                else:
                    add_finding(
                        "SPEC_INVALID",
                        "GEOMETRY_HOST_RECEIPT_MISSING",
                        "Geometry observations do not bind a host runtime receipt.",
                    )

                geometry_runtime = (
                    loaded_geometry.get("runtime")
                    if isinstance(loaded_geometry.get("runtime"), Mapping)
                    else {}
                )
                if host_receipt is not None:
                    actual_host_context = (
                        host_receipt.get("context")
                        if isinstance(host_receipt.get("context"), Mapping)
                        else {}
                    )
                    actual_host_runtime = (
                        host_receipt.get("runtime")
                        if isinstance(host_receipt.get("runtime"), Mapping)
                        else {}
                    )
                    snap_context = (
                        host_binding.get("context")
                        if isinstance(host_binding.get("context"), Mapping)
                        else {}
                    )
                    snap_runtime = (
                        host_binding.get("runtime")
                        if isinstance(host_binding.get("runtime"), Mapping)
                        else {}
                    )
                    expected_host_values = {
                        "schema_version": host_receipt.get("schema_version"),
                        "status": host_receipt.get("status"),
                        "context.run_id": actual_host_context.get("run_id"),
                        "context.source_sha256": actual_host_context.get("source_sha256"),
                        "runtime.runtime_id": actual_host_runtime.get("runtime_id"),
                        "runtime.python_executable": actual_host_runtime.get("python_executable"),
                        "runtime.python_version": actual_host_runtime.get("python_version"),
                    }
                    snapshot_values = {
                        "schema_version": host_binding.get("schema_version"),
                        "status": host_binding.get("status"),
                        "context.run_id": snap_context.get("run_id"),
                        "context.source_sha256": snap_context.get("source_sha256"),
                        "runtime.runtime_id": snap_runtime.get("runtime_id"),
                        "runtime.python_executable": snap_runtime.get("python_executable"),
                        "runtime.python_version": snap_runtime.get("python_version"),
                    }
                    snapshot_mismatches = {
                        field: {"geometry": snapshot_values[field], "receipt": expected}
                        for field, expected in expected_host_values.items()
                        if str(snapshot_values[field]).casefold() != str(expected).casefold()
                    }
                    for field, expected in (
                        ("context.run_id", loaded_geometry.get("run_id")),
                        ("context.source_sha256", geometry_source.get("sha256")),
                        ("status", "PASS"),
                    ):
                        if str(expected_host_values[field]).casefold() != str(expected).casefold():
                            snapshot_mismatches[f"receipt.{field}"] = {
                                "receipt": expected_host_values[field],
                                "expected": expected,
                            }
                    for field in ("runtime_id", "python_executable", "python_version"):
                        manifest_value = geometry_runtime.get(field)
                        receipt_value = actual_host_runtime.get(field)
                        if str(manifest_value).casefold() != str(receipt_value).casefold():
                            snapshot_mismatches[f"geometry_runtime.{field}"] = {
                                "geometry": manifest_value,
                                "receipt": receipt_value,
                            }
                    host_isolation = (
                        host_receipt.get("isolation")
                        if isinstance(host_receipt.get("isolation"), Mapping)
                        else {}
                    )
                    for field in (
                        "required",
                        "isolated",
                        "ignore_environment",
                        "no_user_site",
                        "safe_path",
                    ):
                        if host_isolation.get(field) is not True:
                            snapshot_mismatches[f"receipt.isolation.{field}"] = {
                                "receipt": host_isolation.get(field),
                                "expected": True,
                            }
                    for collection_name in ("smoke_tests", "checks"):
                        collection = host_receipt.get(collection_name)
                        if (
                            not isinstance(collection, list)
                            or not collection
                            or any(
                                not isinstance(item, Mapping) or item.get("passed") is not True
                                for item in collection
                            )
                        ):
                            snapshot_mismatches[f"receipt.{collection_name}"] = {
                                "receipt": collection,
                                "expected": "non-empty and all passed",
                            }
                    if snapshot_mismatches:
                        add_finding(
                            "SPEC_INVALID",
                            "GEOMETRY_RUNTIME_BINDING_MISMATCH",
                            "Geometry runtime provenance differs from its bound host receipt.",
                            evidence=snapshot_mismatches,
                            repair="Rerun geometry using the canonical Host CV runtime.",
                        )
                if geometry_runtime.get("isolated") is not True:
                    add_finding(
                        "SPEC_INVALID",
                        "GEOMETRY_RUNTIME_NOT_ISOLATED",
                        "Geometry did not record isolated Host CV execution.",
                    )

                geometry_implementation = (
                    loaded_geometry.get("implementation")
                    if isinstance(loaded_geometry.get("implementation"), Mapping)
                    else {}
                )
                implementation_records = (
                    ("script", DEFAULT_GEOMETRY_SCRIPT_PATH),
                    ("schema", DEFAULT_GEOMETRY_SCHEMA_PATH),
                )
                implementation_summary: dict[str, Any] = {}
                for label, current_path in implementation_records:
                    binding = (
                        geometry_implementation.get(label)
                        if isinstance(geometry_implementation.get(label), Mapping)
                        else {}
                    )
                    if not binding.get("path"):
                        add_finding(
                            "SPEC_INVALID",
                            "GEOMETRY_IMPLEMENTATION_BINDING_MISSING",
                            f"Geometry manifest does not bind its {label} file.",
                            evidence={"binding": label},
                        )
                        continue
                    implementation_path = _resolve_path(
                        str(binding["path"]), geometry_manifest_path.parent
                    )
                    try:
                        implementation_bytes = snapshot_bytes(
                            implementation_path,
                            label=f"geometry bound {label}",
                        )
                        current_bytes = snapshot_bytes(
                            current_path.resolve(),
                            label=f"checked-in geometry {label}",
                        )
                        actual_hash = _sha256_bytes(implementation_bytes)
                        actual_size = len(implementation_bytes)
                        current_hash = _sha256_bytes(current_bytes)
                        expected_size = binding.get("size_bytes")
                        recorded_hash = str(binding.get("sha256", ""))
                        if (
                            actual_hash.casefold() != recorded_hash.casefold()
                            or actual_size != expected_size
                            or recorded_hash.casefold() != current_hash.casefold()
                        ):
                            add_finding(
                                "SPEC_INVALID",
                                "GEOMETRY_IMPLEMENTATION_HASH_MISMATCH",
                                f"Geometry {label} binding is stale or not the checked-in implementation.",
                                evidence={
                                    "binding": label,
                                    "path": str(implementation_path),
                                    "recorded_sha256": recorded_hash,
                                    "actual_sha256": actual_hash,
                                    "checked_in_sha256": current_hash,
                                    "recorded_size_bytes": expected_size,
                                    "actual_size_bytes": actual_size,
                                },
                                repair="Rerun geometry after the implementation contract changes.",
                            )
                        implementation_summary[label] = {
                            "path": str(implementation_path),
                            "sha256": actual_hash,
                        }
                    except OSError as exc:
                        add_finding(
                            "SPEC_INVALID",
                            "GEOMETRY_IMPLEMENTATION_UNREADABLE",
                            f"Cannot read geometry's bound {label}: {exc}",
                            evidence={"binding": label},
                        )

                geometry_artifacts = (
                    loaded_geometry.get("artifacts")
                    if isinstance(loaded_geometry.get("artifacts"), Mapping)
                    else {}
                )
                artifact_summary: dict[str, Any] = {}
                expected_source_size = (
                    source_spec.get("width_px"),
                    source_spec.get("height_px"),
                )
                decoded_ambiguous_pixel_count: int | None = None
                decoded_label_atlas: np.ndarray | None = None
                decoded_ambiguity_mask: np.ndarray | None = None
                for label, artifact in geometry_artifacts.items():
                    if not isinstance(artifact, Mapping) or not artifact.get("relative_path"):
                        continue
                    artifact_path = _resolve_path(
                        str(artifact["relative_path"]), geometry_manifest_path.parent
                    )
                    try:
                        artifact_bytes = snapshot_bytes(
                            artifact_path, label=f"geometry artifact {label}"
                        )
                        artifact_hash = _sha256_bytes(artifact_bytes)
                        artifact_size = len(artifact_bytes)
                        if artifact_hash.casefold() != str(
                            artifact.get("sha256", "")
                        ).casefold() or artifact_size != artifact.get("size_bytes"):
                            add_finding(
                                "SPEC_INVALID",
                                "GEOMETRY_ARTIFACT_HASH_MISMATCH",
                                "A geometry evidence artifact changed after refinement.",
                                evidence={
                                    "artifact": str(label),
                                    "path": str(artifact_path),
                                    "sha256": artifact_hash,
                                    "size_bytes": artifact_size,
                                },
                            )
                        if artifact_path.parent != geometry_manifest_path.parent:
                            add_finding(
                                "SPEC_INVALID",
                                "GEOMETRY_ARTIFACT_OUTSIDE_STAGE",
                                "Geometry artifacts must remain in the manifest's stage directory.",
                                evidence={
                                    "artifact": str(label),
                                    "path": str(artifact_path),
                                },
                            )
                        with Image.open(io.BytesIO(artifact_bytes)) as image:
                            artifact_format = image.format
                            artifact_mode = image.mode
                            artifact_array = np.asarray(image)
                            artifact_dimensions = image.size
                            if artifact_format != "PNG":
                                add_finding(
                                    "SPEC_INVALID",
                                    "GEOMETRY_ARTIFACT_FORMAT_MISMATCH",
                                    "A geometry image artifact is not encoded as PNG.",
                                    evidence={
                                        "artifact": str(label),
                                        "decoded_format": artifact_format,
                                    },
                                )
                            if label == "overlay" and (
                                artifact_mode != "RGB"
                                or artifact_array.dtype != np.uint8
                                or artifact_array.ndim != 3
                                or artifact_array.shape[2] != 3
                            ):
                                add_finding(
                                    "SPEC_INVALID",
                                    "GEOMETRY_OVERLAY_ENCODING_MISMATCH",
                                    "Geometry overlay is not a native RGB8 PNG.",
                                    evidence={
                                        "mode": artifact_mode,
                                        "dtype": str(artifact_array.dtype),
                                        "shape": list(artifact_array.shape),
                                    },
                                )
                            elif label == "label_atlas":
                                if (
                                    artifact_mode != "I;16"
                                    or artifact_array.dtype != np.uint16
                                    or artifact_array.ndim != 2
                                ):
                                    add_finding(
                                        "SPEC_INVALID",
                                        "GEOMETRY_LABEL_ATLAS_ENCODING_MISMATCH",
                                        "Geometry label atlas is not the canonical uint16 lossless PNG.",
                                        evidence={
                                            "mode": artifact_mode,
                                            "dtype": str(artifact_array.dtype),
                                            "shape": list(artifact_array.shape),
                                        },
                                    )
                                else:
                                    decoded_label_atlas = artifact_array.copy()
                            if label == "ambiguity_mask":
                                if (
                                    artifact_mode != "L"
                                    or artifact_array.dtype != np.uint8
                                    or artifact_array.ndim != 2
                                ):
                                    add_finding(
                                        "SPEC_INVALID",
                                        "GEOMETRY_AMBIGUITY_MASK_ENCODING_MISMATCH",
                                        "Geometry ambiguity mask is not the canonical uint8 grayscale PNG.",
                                        evidence={
                                            "mode": artifact_mode,
                                            "dtype": str(artifact_array.dtype),
                                            "shape": list(artifact_array.shape),
                                        },
                                    )
                                else:
                                    decoded_ambiguity_mask = artifact_array.copy()
                                    ambiguity_values = {
                                        int(value) for value in np.unique(artifact_array)
                                    }
                                    if not ambiguity_values.issubset({0, 255}):
                                        add_finding(
                                            "SPEC_INVALID",
                                            "GEOMETRY_AMBIGUITY_MASK_NOT_BINARY",
                                            "Geometry ambiguity mask contains values other than 0 and 255.",
                                            evidence={"values": sorted(ambiguity_values)},
                                        )
                                    decoded_ambiguous_pixel_count = int(
                                        np.count_nonzero(artifact_array == 255)
                                    )
                        if artifact_dimensions != expected_source_size:
                            add_finding(
                                "SPEC_INVALID",
                                "GEOMETRY_ARTIFACT_DIMENSION_MISMATCH",
                                "Geometry evidence image dimensions differ from the frozen source.",
                                evidence={
                                    "artifact": str(label),
                                    "expected": expected_source_size,
                                    "actual": artifact_dimensions,
                                },
                            )
                        if (
                            artifact.get("width_px") != artifact_dimensions[0]
                            or artifact.get("height_px") != artifact_dimensions[1]
                        ):
                            add_finding(
                                "SPEC_INVALID",
                                "GEOMETRY_ARTIFACT_METADATA_MISMATCH",
                                "Geometry artifact metadata differs from the decoded PNG.",
                                evidence={
                                    "artifact": str(label),
                                    "declared": (
                                        artifact.get("width_px"),
                                        artifact.get("height_px"),
                                    ),
                                    "decoded": artifact_dimensions,
                                },
                            )
                        artifact_summary[str(label)] = {
                            "path": str(artifact_path),
                            "sha256": artifact_hash,
                            "size_bytes": artifact_size,
                            "width_px": artifact_dimensions[0],
                            "height_px": artifact_dimensions[1],
                            "decoded_format": artifact_format,
                        }
                    except (OSError, ValueError) as exc:
                        add_finding(
                            "SPEC_INVALID",
                            "GEOMETRY_ARTIFACT_UNREADABLE",
                            f"Cannot read a hash-bound geometry artifact: {exc}",
                            evidence={"artifact": str(label), "path": str(artifact_path)},
                        )
                if (
                    decoded_ambiguous_pixel_count is not None
                    and geometry_summary.get("ambiguous_pixel_count")
                    != decoded_ambiguous_pixel_count
                ):
                    add_finding(
                        "SPEC_INVALID",
                        "GEOMETRY_AMBIGUITY_COUNT_MISMATCH",
                        "Geometry summary ambiguity count differs from its lossless mask.",
                        evidence={
                            "manifest": geometry_summary.get("ambiguous_pixel_count"),
                            "decoded": decoded_ambiguous_pixel_count,
                        },
                    )

                if decoded_label_atlas is not None and isinstance(text_geometry, list):
                    measured_by_label: dict[int, Mapping[str, Any]] = {}
                    duplicate_labels: list[int] = []
                    invalid_item_ids: list[str] = []
                    for item in text_geometry:
                        if not isinstance(item, Mapping):
                            continue
                        item_id = str(item.get("candidate_id", ""))
                        item_status = str(item.get("status", ""))
                        label_value = item.get("mask_label")
                        if item_status != "MEASURED":
                            ink_area_value = item.get("ink_area_px")
                            if (
                                label_value is not None
                                or item.get("ink_bbox") is not None
                                or not (
                                    ink_area_value is None
                                    or (
                                        isinstance(ink_area_value, int)
                                        and not isinstance(ink_area_value, bool)
                                        and ink_area_value == 0
                                    )
                                )
                            ):
                                invalid_item_ids.append(item_id)
                            continue
                        if not isinstance(label_value, int) or isinstance(label_value, bool):
                            invalid_item_ids.append(item_id)
                            continue
                        if label_value in measured_by_label:
                            duplicate_labels.append(label_value)
                            continue
                        measured_by_label[label_value] = item

                    if invalid_item_ids:
                        add_finding(
                            "SPEC_INVALID",
                            "GEOMETRY_TEXT_LABEL_CONTRACT_MISMATCH",
                            "Measured and inconclusive text records violate the label-atlas contract.",
                            evidence={"candidate_ids": sorted(set(invalid_item_ids))},
                        )
                    if duplicate_labels:
                        add_finding(
                            "SPEC_INVALID",
                            "GEOMETRY_LABEL_ATLAS_DUPLICATE_LABEL",
                            "Multiple measured text records claim the same atlas label.",
                            evidence={"labels": sorted(set(duplicate_labels))},
                        )

                    decoded_labels = {
                        int(value) for value in np.unique(decoded_label_atlas) if int(value) != 0
                    }
                    expected_labels = set(measured_by_label)
                    if decoded_labels != expected_labels:
                        add_finding(
                            "SPEC_INVALID",
                            "GEOMETRY_LABEL_ATLAS_LABEL_SET_MISMATCH",
                            "Label-atlas values do not exactly match measured text records.",
                            evidence={
                                "missing_labels": sorted(expected_labels - decoded_labels),
                                "unexpected_labels": sorted(decoded_labels - expected_labels),
                            },
                        )

                    for label_value in sorted(decoded_labels & expected_labels):
                        item = measured_by_label[label_value]
                        label_mask = decoded_label_atlas == label_value
                        ys, xs = np.nonzero(label_mask)
                        decoded_count = int(xs.size)
                        decoded_box = (
                            {
                                "x0": int(xs.min()),
                                "y0": int(ys.min()),
                                "x1": int(xs.max()) + 1,
                                "y1": int(ys.max()) + 1,
                            }
                            if decoded_count
                            else None
                        )
                        if decoded_count != item.get("ink_area_px") or decoded_box != item.get(
                            "ink_bbox"
                        ):
                            add_finding(
                                "SPEC_INVALID",
                                "GEOMETRY_LABEL_ATLAS_MEASUREMENT_MISMATCH",
                                "A label's decoded pixels differ from its text-geometry measurement.",
                                evidence={
                                    "candidate_id": item.get("candidate_id"),
                                    "label": label_value,
                                    "declared_pixel_count": item.get("ink_area_px"),
                                    "decoded_pixel_count": decoded_count,
                                    "declared_bbox": item.get("ink_bbox"),
                                    "decoded_bbox": decoded_box,
                                },
                            )

                if (
                    decoded_ambiguity_mask is not None
                    and decoded_label_atlas is not None
                    and decoded_ambiguity_mask.shape != decoded_label_atlas.shape
                ):
                    add_finding(
                        "SPEC_INVALID",
                        "GEOMETRY_MASK_SHAPE_MISMATCH",
                        "Geometry label and ambiguity masks do not share one source-pixel grid.",
                    )

                geometry_record = {
                    "path": str(geometry_manifest_path),
                    "sha256": geometry_manifest_hash,
                    "schema_path": str(DEFAULT_GEOMETRY_SCHEMA_PATH),
                    "schema_sha256": _sha256_bytes(
                        schema_payloads[DEFAULT_GEOMETRY_SCHEMA_PATH.resolve()]
                    ),
                    "run_id": loaded_geometry.get("run_id"),
                    "source_sha256": geometry_source.get("sha256"),
                    "ocr_manifest_sha256": ocr_binding.get("sha256"),
                    "host_runtime_receipt_path": (
                        str(host_receipt_path) if host_receipt_path is not None else None
                    ),
                    "host_runtime_receipt_sha256": host_receipt_hash,
                    "status": geometry_status,
                    "mode": geometry_mode,
                    "promotion_allowed": promotion_allowed,
                    "observation_only": geometry_mode == "observation_only",
                    "contributes_to_drawer_authorization": False,
                    "implementation": implementation_summary,
                    "artifacts": artifact_summary,
                }
                geometry_diagnostics.append(
                    {
                        "kind": "PHASE1_GEOMETRY_BOUNDARY",
                        "status": geometry_status,
                        "diagnostic_only": True,
                        "used_as_scene_geometry": False,
                        "authorizes_drawer": False,
                        "text_observation_count": len(
                            loaded_geometry.get("text_geometry", []) or []
                        ),
                        "neighbor_pair_count": len(loaded_geometry.get("neighbor_pairs", []) or []),
                        "frame_candidate_count": len(
                            loaded_geometry.get("frame_candidates", []) or []
                        ),
                    }
                )
                if geometry_status == "GEOMETRY_INCONCLUSIVE":
                    add_finding(
                        "INCONCLUSIVE",
                        "GEOMETRY_OBSERVATIONS_INCONCLUSIVE",
                        "Phase-1 geometry did not produce a complete observation set.",
                        evidence={
                            "degradations": list(loaded_geometry.get("degradations", []) or [])
                        },
                        repair="Resolve the reported CV degradation or keep the scene blocked.",
                    )
            except (OSError, ValueError, json.JSONDecodeError) as exc:
                add_finding(
                    "SPEC_INVALID",
                    "GEOMETRY_EVIDENCE_UNREADABLE",
                    f"Cannot read or validate the bound geometry evidence: {exc}",
                    repair="Provide the exact hash-bound Phase-1 geometry manifest.",
                )

    background_evidence = str(canvas_spec.get("background_evidence", ""))
    if background_evidence == "measured_reference" and perception_manifest is not None:
        analysis_stage = next(
            (
                stage
                for stage in perception_manifest.get("upstream_stages", []) or []
                if isinstance(stage, Mapping) and stage.get("name") == "analysis"
            ),
            None,
        )
        inventory_record = None
        if isinstance(analysis_stage, Mapping):
            inventory_record = next(
                (
                    item
                    for item in analysis_stage.get("files", []) or []
                    if isinstance(item, Mapping) and item.get("relative_path") == "inventory.json"
                ),
                None,
            )
        if not isinstance(analysis_stage, Mapping) or not isinstance(inventory_record, Mapping):
            add_finding(
                "SPEC_INVALID",
                "BACKGROUND_MEASUREMENT_MISSING",
                "measured_reference background requires a hash-bound analysis/inventory.json artifact.",
            )
        else:
            try:
                analysis_root = _resolve_path(str(analysis_stage["path"]), manifest_path.parent)
                inventory_path = (analysis_root / "inventory.json").resolve()
                inventory_bytes = snapshot_bytes(inventory_path, label="analysis inventory")
                actual_inventory_hash = _sha256_bytes(inventory_bytes)
                expected_inventory_hash = str(inventory_record.get("sha256", ""))
                if actual_inventory_hash.casefold() != expected_inventory_hash.casefold():
                    add_finding(
                        "SPEC_INVALID",
                        "BACKGROUND_INVENTORY_HASH_MISMATCH",
                        "The analysis inventory changed after the perception manifest was frozen.",
                    )
                inventory = _strict_json_bytes(inventory_bytes, label="analysis inventory")
                inventory_source_hash = str((inventory.get("source") or {}).get("sha256", ""))
                if (
                    inventory_source_hash.casefold()
                    != str(source_spec.get("sha256", "")).casefold()
                ):
                    add_finding(
                        "SPEC_INVALID",
                        "BACKGROUND_INVENTORY_SOURCE_MISMATCH",
                        "Background inventory belongs to a different source image.",
                    )
                measured_background = str((inventory.get("canvas") or {}).get("background_hex", ""))
                declared_background = str(canvas_spec.get("background", ""))
                if measured_background.casefold() != declared_background.casefold():
                    add_finding(
                        "SPEC_INVALID",
                        "BACKGROUND_COLOR_MISMATCH",
                        "Declared canvas background differs from the hash-bound measured reference.",
                        evidence={"declared": declared_background, "measured": measured_background},
                        repair="Use the measured background or record an approved intentional deviation.",
                    )
            except (OSError, ValueError, json.JSONDecodeError) as exc:
                add_finding(
                    "SPEC_INVALID",
                    "BACKGROUND_INVENTORY_UNREADABLE",
                    f"Cannot read the hash-bound background inventory: {exc}",
                )

    raw_elements = spec.get("elements")
    elements = (
        [item for item in raw_elements if isinstance(item, Mapping)]
        if isinstance(raw_elements, list)
        else []
    )
    ids = [str(element.get("id")) for element in elements if element.get("id") is not None]
    for element_id, count in Counter(ids).items():
        if count > 1:
            add_finding(
                "SPEC_INVALID",
                "DUPLICATE_ELEMENT_ID",
                f"Element id {element_id!r} occurs {count} times.",
                element_ids=[element_id],
                repair="Assign one stable unique id to every element.",
            )
    by_id = {
        str(element.get("id")): element for element in elements if element.get("id") is not None
    }
    valid_boxes: dict[str, tuple[float, float, float, float]] = {}

    raw_edges = spec.get("edges", [])
    edges = (
        [item for item in raw_edges if isinstance(item, Mapping)]
        if isinstance(raw_edges, list)
        else []
    )
    edge_ids = [str(edge.get("id")) for edge in edges if edge.get("id") is not None]
    for edge_id, count in Counter(edge_ids).items():
        if count > 1:
            add_finding(
                "SPEC_INVALID",
                "DUPLICATE_EDGE_ID",
                f"Edge id {edge_id!r} occurs {count} times.",
                repair="Assign one stable unique id to every edge.",
            )
    for edge in edges:
        edge_id = str(edge.get("id", "<missing>"))
        source_id = str(edge.get("from", ""))
        target_id = str(edge.get("to", ""))
        unknown = [endpoint for endpoint in (source_id, target_id) if endpoint not in by_id]
        if unknown:
            add_finding(
                "SPEC_INVALID",
                "UNKNOWN_EDGE_ENDPOINT",
                f"Edge {edge_id!r} references unknown endpoint(s).",
                element_ids=unknown,
                evidence={"edge_id": edge_id, "from": source_id, "to": target_id},
                repair="Bind both edge endpoints to existing element ids.",
            )
        elif source_id == target_id and edge.get("allow_self_loop") is not True:
            add_finding(
                "SPEC_INVALID",
                "SELF_EDGE_UNDECLARED",
                f"Edge {edge_id!r} binds both endpoints to the same element.",
                element_ids=[source_id],
                evidence={"edge_id": edge_id},
                repair="Correct the endpoint or declare a supported feedback-loop contract.",
            )

    raw_formulas = spec.get("formulas", [])
    formulas = (
        [item for item in raw_formulas if isinstance(item, Mapping)]
        if isinstance(raw_formulas, list)
        else []
    )
    formula_ids = [str(formula.get("id")) for formula in formulas if formula.get("id") is not None]
    for formula_id, count in Counter(formula_ids).items():
        if count > 1:
            add_finding(
                "SPEC_INVALID",
                "DUPLICATE_FORMULA_ID",
                f"Formula id {formula_id!r} occurs {count} times.",
            )
    formula_by_id = {
        str(formula.get("id")): formula for formula in formulas if formula.get("id") is not None
    }
    formula_references: dict[str, list[dict[str, Any]]] = defaultdict(list)
    for element_id, element in by_id.items():
        element_type = str(element.get("type", "")).casefold()
        if element_type == "formula":
            formula_references[str(element.get("formula_id", ""))].append(
                {"kind": "display_element", "element_id": element_id}
            )
        if element_type != "text":
            continue
        for run_index, run in enumerate(_content_runs(element)):
            if run.get("kind") != "math":
                continue
            formula_references[str(run.get("formula_id", ""))].append(
                {"kind": "inline_run", "element_id": element_id, "run_index": run_index}
            )

    for referenced_id, references in formula_references.items():
        if referenced_id not in formula_by_id:
            add_finding(
                "SPEC_INVALID",
                "FORMULA_RECORD_MISSING",
                "Formula element or math run has no matching canonical formula record.",
                element_ids=[str(item["element_id"]) for item in references],
                evidence={"formula_id": referenced_id, "references": references},
            )

    trusted_formula_evidence = {"user_confirmed", "source_text"}
    for formula in formulas:
        formula_id = str(formula.get("id", "<missing>"))
        element_id = str(formula.get("element_id", ""))
        canonical_latex = str(formula.get("canonical_latex", ""))
        if not canonical_latex:
            add_finding(
                "SPEC_INVALID",
                "FORMULA_CANONICAL_LATEX_MISSING",
                f"Formula {formula_id!r} has no canonical LaTeX source.",
                element_ids=[element_id],
            )
        actual_latex_hash = sha256_text(canonical_latex)
        expected_latex_hash = str(formula.get("latex_sha256", ""))
        if actual_latex_hash.casefold() != expected_latex_hash.casefold():
            add_finding(
                "SPEC_INVALID",
                "FORMULA_LATEX_HASH_MISMATCH",
                "Canonical LaTeX does not match its frozen UTF-8 SHA-256.",
                element_ids=[element_id],
                evidence={"expected": expected_latex_hash, "actual": actual_latex_hash},
            )
        if formula.get("render_kind") != "native_office_math":
            add_finding(
                "SPEC_INVALID",
                "FORMULA_RENDER_KIND_INVALID",
                "Every formula must declare render_kind=native_office_math.",
                element_ids=[element_id],
                evidence={"actual": formula.get("render_kind")},
            )
        if formula.get("fallback_policy") != "strict_no_raster_no_svg":
            add_finding(
                "SPEC_INVALID",
                "FORMULA_FALLBACK_POLICY_INVALID",
                "Raster, SVG, image, and ordinary-text formula fallbacks are forbidden.",
                element_ids=[element_id],
                evidence={"actual": formula.get("fallback_policy")},
            )

        references = formula_references.get(formula_id, [])
        if len(references) != 1:
            add_finding(
                "SPEC_INVALID",
                "FORMULA_REFERENCE_NOT_UNIQUE",
                "Each formula record must be consumed exactly once by one display element or math run.",
                element_ids=[str(item["element_id"]) for item in references] or [element_id],
                evidence={"formula_id": formula_id, "references": references},
            )
        else:
            reference = references[0]
            expected_mode = "inline" if reference["kind"] == "inline_run" else "display"
            if formula.get("mode") != expected_mode:
                add_finding(
                    "SPEC_INVALID",
                    "FORMULA_MODE_REFERENCE_MISMATCH",
                    "Formula mode must match its structured inline/display reference.",
                    element_ids=[str(reference["element_id"])],
                    evidence={"expected": expected_mode, "actual": formula.get("mode")},
                )
            if element_id != str(reference["element_id"]):
                add_finding(
                    "SPEC_INVALID",
                    "FORMULA_ELEMENT_MISMATCH",
                    "Formula element_id must equal the owner of its sole formula reference.",
                    element_ids=[element_id, str(reference["element_id"])],
                    evidence={"formula_id": formula_id, "reference": reference},
                )
        if element_id not in by_id:
            add_finding(
                "SPEC_INVALID",
                "FORMULA_ELEMENT_MISMATCH",
                f"Formula {formula_id!r} binds an unknown owner element.",
                element_ids=[element_id],
            )

        converter_record, converter_issues = inspect_formula_converter_receipt(
            formula,
            base_dir=resolved_base,
        )
        if converter_record is not None:
            formula_converter_receipts.append(converter_record)
        for issue in converter_issues:
            add_finding(
                str(issue["disposition"]),
                str(issue["code"]),
                str(issue["message"]),
                element_ids=[element_id],
                evidence=issue.get("evidence"),
                repair=(
                    "Compile canonical LaTeX with the pinned native Office Math converter and freeze its PASS receipt."
                ),
            )

        evidence = {str(item) for item in formula.get("source_evidence", []) or []}
        candidate_id = formula.get("perception_candidate_id")
        authoritative = bool(evidence.intersection(trusted_formula_evidence))
        if candidate_id:
            if "local_ocr" not in evidence:
                add_finding(
                    "SPEC_INVALID",
                    "FORMULA_OCR_EVIDENCE_TAG_MISSING",
                    "A formula candidate binding requires local_ocr in source_evidence.",
                    element_ids=[element_id],
                )
            candidate = candidate_by_id.get(str(candidate_id))
            decision = decision_by_id.get(str(candidate_id))
            if candidate is None or decision is None:
                add_finding(
                    "SPEC_INVALID",
                    "FORMULA_CANDIDATE_BINDING_UNKNOWN",
                    "Formula references a candidate absent from the bound perception evidence.",
                    element_ids=[element_id],
                    evidence={"candidate_id": candidate_id},
                )
            elif decision.get("status") != "FORMULA_CONFIRMED":
                add_finding(
                    "SPEC_INVALID",
                    "FORMULA_CANDIDATE_DECISION_INVALID",
                    "Formula candidates require a FORMULA_CONFIRMED review decision.",
                    element_ids=[element_id],
                    evidence={"candidate_id": candidate_id, "status": decision.get("status")},
                )
            elif str(decision.get("authoritative_latex", "")) != canonical_latex:
                add_finding(
                    "SPEC_INVALID",
                    "FORMULA_LATEX_MISMATCH",
                    "Formula LaTeX differs from the authoritative reviewed candidate decision.",
                    element_ids=[element_id],
                    evidence={"candidate_id": candidate_id},
                )
        elif "local_ocr" in evidence:
            add_finding(
                "SPEC_INVALID",
                "FORMULA_OCR_CANDIDATE_MISSING",
                "Formula claiming local_ocr evidence must bind one reviewed formula candidate.",
                element_ids=[element_id],
            )
        elif not authoritative:
            add_finding(
                "INCONCLUSIVE",
                "FORMULA_SOURCE_UNCONFIRMED",
                f"Formula {formula_id!r} lacks user-confirmed LaTeX or authoritative source text.",
                element_ids=[element_id],
            )

    for element_id, element in by_id.items():
        box = _bbox(element)
        if box is None:
            add_finding(
                "SPEC_INVALID",
                "INVALID_BBOX",
                "bbox must contain finite x/y and strictly positive w/h values.",
                element_ids=[element_id],
            )
            continue
        valid_boxes[element_id] = box
        disposition = str(element.get("disposition", ""))
        if (
            disposition in {"INCONCLUSIVE", "UNREADABLE"}
            or element.get("strategy") == "source_ambiguity"
        ):
            add_finding(
                "INCONCLUSIVE",
                "UNRESOLVED_ELEMENT_EVIDENCE",
                "Element evidence is unresolved and cannot enter a render-ready scene.",
                element_ids=[element_id],
                evidence={"disposition": disposition, "strategy": element.get("strategy")},
                repair="Resolve the source evidence or remove the unsupported element from the frozen scene.",
            )
        if str(element.get("type", "")).casefold() == "text" and disposition != "CONFIRMED":
            add_finding(
                "INCONCLUSIVE",
                "TEXT_NOT_CONFIRMED",
                "Text objects must be confirmed before Drawer execution.",
                element_ids=[element_id],
            )
        if str(element.get("type", "")).casefold() == "text" and disposition == "CONFIRMED":
            evidence = {str(item) for item in element.get("source_evidence", []) or []}
            authoritative = bool(evidence.intersection({"user_confirmed", "source_text"}))
            candidate_ids = [
                str(item) for item in element.get("perception_candidate_ids", []) or []
            ]
            criticality = str(element.get("criticality", ""))
            math_reasons = math_like_text_reasons(element)
            if math_reasons:
                add_finding(
                    "SPEC_INVALID",
                    "MATH_SYNTAX_IN_PLAIN_TEXT",
                    "Mathematical syntax cannot masquerade as ordinary PowerPoint text.",
                    element_ids=[element_id],
                    evidence={"reasons": math_reasons, "text": element.get("text")},
                    repair=(
                        "Split the content into content_runs and bind every math run to one canonical formula_id."
                    ),
                )
            if criticality == "critical" and not authoritative:
                add_finding(
                    "INCONCLUSIVE",
                    "CRITICAL_TEXT_AUTHORITY_MISSING",
                    "Critical titles, numbers, units, and arrow labels require explicit user/source authority.",
                    element_ids=[element_id],
                    evidence={"source_evidence": sorted(evidence)},
                )
            if "local_ocr" in evidence and not candidate_ids:
                add_finding(
                    "SPEC_INVALID",
                    "TEXT_OCR_CANDIDATE_MISSING",
                    "Text claiming local_ocr evidence must bind at least one reviewed candidate id.",
                    element_ids=[element_id],
                )
            if candidate_ids and "local_ocr" not in evidence:
                add_finding(
                    "SPEC_INVALID",
                    "TEXT_OCR_EVIDENCE_TAG_MISSING",
                    "perception_candidate_ids require local_ocr in source_evidence.",
                    element_ids=[element_id],
                )
            if not authoritative and not candidate_ids:
                add_finding(
                    "INCONCLUSIVE",
                    "TEXT_EVIDENCE_INSUFFICIENT",
                    "Confirmed text needs authoritative source evidence or a reviewed OCR candidate binding.",
                    element_ids=[element_id],
                    evidence={"source_evidence": sorted(evidence)},
                )
            for candidate_id in candidate_ids:
                candidate = candidate_by_id.get(candidate_id)
                decision = decision_by_id.get(candidate_id)
                if candidate is None or decision is None:
                    add_finding(
                        "SPEC_INVALID",
                        "TEXT_CANDIDATE_BINDING_UNKNOWN",
                        "Text element references a candidate absent from the bound raw/review evidence.",
                        element_ids=[element_id],
                        evidence={"candidate_id": candidate_id},
                    )
                    continue
                if decision.get("status") not in {"CONFIRMED", "CORRECTED"}:
                    add_finding(
                        "SPEC_INVALID",
                        "TEXT_CANDIDATE_DECISION_INVALID",
                        "A text element may bind only a CONFIRMED or CORRECTED text decision.",
                        element_ids=[element_id],
                        evidence={"candidate_id": candidate_id, "status": decision.get("status")},
                    )
                    continue
                confirmed_text = str(decision.get("confirmed_text", ""))
                prose_projection = _text_projection(element)
                run_text_values = [
                    str(run.get("text", ""))
                    for run in _content_runs(element)
                    if run.get("kind") == "text"
                ]
                value_matches = (
                    confirmed_text == prose_projection or confirmed_text in run_text_values
                )
                if not value_matches:
                    add_finding(
                        "SPEC_INVALID",
                        "TEXT_CANDIDATE_VALUE_MISMATCH",
                        "Text prose differs from its authoritative reviewed candidate decision.",
                        element_ids=[element_id],
                        evidence={
                            "candidate_id": candidate_id,
                            "element_prose": prose_projection,
                            "text_runs": run_text_values,
                            "reviewed_text": confirmed_text,
                        },
                    )
                # The envelope spans every conflicting/duplicate OCR observation and can be
                # arbitrarily larger than the detector's selected primary evidence.  It is an
                # uncertainty diagnostic, never the geometry contract for a scene element.
                candidate_box = candidate.get("bbox_source")
                if isinstance(candidate_box, Mapping):
                    bbox_matches, bbox_evidence = _candidate_bbox_matches(box, candidate_box)
                    if not bbox_matches:
                        add_finding(
                            "SPEC_INVALID",
                            "TEXT_CANDIDATE_BBOX_MISMATCH",
                            "Reviewed OCR candidate is not geometrically bound to this text box.",
                            element_ids=[element_id],
                            evidence={"candidate_id": candidate_id, **bbox_evidence},
                        )
                geometry_item = geometry_by_candidate.get(candidate_id)
                if isinstance(geometry_item, Mapping):
                    ink_box_value = geometry_item.get("ink_bbox")
                    ink_box = None
                    if isinstance(ink_box_value, Mapping):
                        try:
                            x0 = float(ink_box_value["x0"])
                            y0 = float(ink_box_value["y0"])
                            x1 = float(ink_box_value["x1"])
                            y1 = float(ink_box_value["y1"])
                            if x1 > x0 and y1 > y0:
                                ink_box = (x0, y0, x1 - x0, y1 - y0)
                        except (KeyError, TypeError, ValueError):
                            ink_box = None
                    geometry_diagnostics.append(
                        {
                            "kind": "TEXT_INK_OBSERVATION",
                            "element_id": element_id,
                            "candidate_id": candidate_id,
                            "observation_status": geometry_item.get("status"),
                            "ink_bbox_source": dict(ink_box_value)
                            if isinstance(ink_box_value, Mapping)
                            else None,
                            "edge_uncertainty_px": geometry_item.get("edge_uncertainty_px"),
                            "ink_bottom_alignment": geometry_item.get("baseline"),
                            "element_contains_observed_ink": (
                                _contains(box, ink_box) if ink_box is not None else None
                            ),
                            "diagnostic_only": True,
                            "used_as_scene_geometry": False,
                            "authorizes_drawer": False,
                        }
                    )
        if measured_source is not None:
            canvas_box = (
                0.0,
                0.0,
                float(measured_source["width_px"]),
                float(measured_source["height_px"]),
            )
            if not _contains(canvas_box, box):
                add_finding(
                    "REGION_REPLAN",
                    "BBOX_OUTSIDE_CANVAS",
                    "Element bbox extends outside the measured PNG canvas.",
                    element_ids=[element_id],
                    evidence={"bbox": box, "canvas": canvas_box},
                    repair="Replan this region inside the source-pixel canvas.",
                )

        parent_id = element.get("parent_id")
        if parent_id is not None and str(parent_id) not in by_id:
            add_finding(
                "SPEC_INVALID",
                "UNKNOWN_PARENT",
                f"Parent {parent_id!r} does not exist.",
                element_ids=[element_id],
            )
        elif parent_id is not None:
            parent_type = str(by_id[str(parent_id)].get("type", "")).casefold()
            if parent_type not in CONTAINER_TYPES:
                add_finding(
                    "SPEC_INVALID",
                    "INVALID_PARENT_TYPE",
                    "Only declared container element types may own child elements.",
                    element_ids=[str(parent_id), element_id],
                    evidence={"parent_type": parent_type, "allowed_types": sorted(CONTAINER_TYPES)},
                )
        for allowed_id in element.get("allowed_overlap", []) or []:
            if str(allowed_id) == element_id or str(allowed_id) not in by_id:
                add_finding(
                    "SPEC_INVALID",
                    "INVALID_ALLOWED_OVERLAP",
                    f"allowed_overlap target {allowed_id!r} is self-referential or unknown.",
                    element_ids=[element_id],
                    repair="Reference an existing different element id, or remove the declaration.",
                )

    parent_of = {
        element_id: str(element["parent_id"])
        for element_id, element in by_id.items()
        if element.get("parent_id") is not None
        and str(element.get("parent_id")) in by_id
        and str(by_id[str(element.get("parent_id"))].get("type", "")).casefold() in CONTAINER_TYPES
    }
    cycle_keys: set[tuple[str, ...]] = set()
    for element_id in by_id:
        path: list[str] = []
        seen: set[str] = set()
        current = element_id
        while current in parent_of:
            if current in seen:
                cycle = tuple(sorted(set(path[path.index(current) :])))
                if cycle not in cycle_keys:
                    cycle_keys.add(cycle)
                    add_finding(
                        "SPEC_INVALID",
                        "PARENT_CYCLE",
                        "Parent relationships contain a cycle.",
                        element_ids=cycle,
                    )
                break
            seen.add(current)
            path.append(current)
            current = parent_of[current]

    for child_id, parent_id in parent_of.items():
        if (
            child_id in valid_boxes
            and parent_id in valid_boxes
            and not _contains(valid_boxes[parent_id], valid_boxes[child_id])
        ):
            add_finding(
                "REGION_REPLAN",
                "PARENT_CONTAINMENT",
                "Child bbox is not fully contained by its parent bbox.",
                element_ids=[parent_id, child_id],
                evidence={
                    "parent_bbox": valid_boxes[parent_id],
                    "child_bbox": valid_boxes[child_id],
                },
                repair="Resize or reposition the child/parent before building the region.",
            )
        parent_z = by_id[parent_id].get("z_index")
        child_z = by_id[child_id].get("z_index")
        if isinstance(parent_z, int) and isinstance(child_z, int) and parent_z >= child_z:
            add_finding(
                "REGION_REPLAN",
                "PARENT_Z_ORDER",
                "A parent container must be behind its child in z-order.",
                element_ids=[parent_id, child_id],
                evidence={"parent_z_index": parent_z, "child_z_index": child_z},
                repair="Assign a greater z_index to the child.",
            )

    z_by_scope: dict[str | None, dict[int, list[str]]] = defaultdict(lambda: defaultdict(list))
    for element_id, element in by_id.items():
        z_index = element.get("z_index")
        if isinstance(z_index, int):
            scope = str(element["parent_id"]) if element.get("parent_id") is not None else None
            z_by_scope[scope][z_index].append(element_id)
    for scope, z_values in z_by_scope.items():
        for z_index, scoped_ids in z_values.items():
            if len(scoped_ids) > 1:
                add_finding(
                    "REGION_REPLAN",
                    "AMBIGUOUS_Z_INDEX",
                    "Sibling elements share the same z_index.",
                    element_ids=scoped_ids,
                    evidence={"parent_id": scope, "z_index": z_index},
                    repair="Give siblings deterministic, distinct z_index values.",
                )

    def ancestors(element_id: str) -> set[str]:
        result: set[str] = set()
        current = element_id
        while current in parent_of and parent_of[current] not in result:
            current = parent_of[current]
            result.add(current)
        return result

    allowed_pairs: set[frozenset[str]] = set()
    for element_id, element in by_id.items():
        for target in element.get("allowed_overlap", []) or []:
            if str(target) in by_id and str(target) != element_id:
                allowed_pairs.add(frozenset((element_id, str(target))))

    collidable = [
        element_id
        for element_id, element in by_id.items()
        if element_id in valid_boxes
        and str(element.get("type", "")).casefold() not in NON_COLLIDING_TYPES
    ]
    for index, first_id in enumerate(collidable):
        for second_id in collidable[index + 1 :]:
            if first_id in ancestors(second_id) or second_id in ancestors(first_id):
                continue
            if frozenset((first_id, second_id)) in allowed_pairs:
                continue
            overlap = _intersection(valid_boxes[first_id], valid_boxes[second_id])
            if overlap is None:
                continue
            first_type = str(by_id[first_id].get("type", "")).casefold()
            second_type = str(by_id[second_id].get("type", "")).casefold()
            if first_type in TEXT_TYPES and second_type in TEXT_TYPES:
                code = "TEXT_TEXT_COLLISION"
            elif first_type in TEXT_TYPES or second_type in TEXT_TYPES:
                code = "TEXT_SHAPE_COLLISION"
            else:
                code = "SHAPE_SHAPE_COLLISION"
            add_finding(
                "REGION_REPLAN",
                code,
                "Objects overlap without an explicit allowed_overlap declaration.",
                element_ids=[first_id, second_id],
                evidence={"intersection_bbox": overlap},
                repair="Reposition the objects or explicitly declare the intentional overlap.",
            )

    edge_routes: dict[str, list[tuple[float, float]]] = {}
    for edge in edges:
        edge_id = str(edge.get("id", "<missing>"))
        source_id = str(edge.get("from", ""))
        target_id = str(edge.get("to", ""))
        if source_id not in valid_boxes or target_id not in valid_boxes:
            continue
        route_kind = str(edge.get("route", ""))
        if route_kind == "curve":
            add_finding(
                "INCONCLUSIVE",
                "CURVED_EDGE_REQUIRES_RENDERER",
                "Curved connector clearance cannot be proven by the deterministic polyline gate.",
                element_ids=[source_id, target_id],
                evidence={"edge_id": edge_id},
                repair="Use an explicit polyline/orthogonal route or verify the curve from renderer evidence.",
            )
            continue
        source_anchor = str(edge.get("source_anchor", ""))
        target_anchor = str(edge.get("target_anchor", ""))
        source_point_value = edge.get("source_point")
        target_point_value = edge.get("target_point")
        if source_anchor == "free":
            start = (
                (float(source_point_value["x"]), float(source_point_value["y"]))
                if isinstance(source_point_value, Mapping)
                else None
            )
        else:
            start = _anchor_point(valid_boxes[source_id], source_anchor)
        if target_anchor == "free":
            end = (
                (float(target_point_value["x"]), float(target_point_value["y"]))
                if isinstance(target_point_value, Mapping)
                else None
            )
        else:
            end = _anchor_point(valid_boxes[target_id], target_anchor)
        if start is None or end is None:
            add_finding(
                "SPEC_INVALID",
                "EDGE_ANCHOR_POINT_MISSING",
                "Free edge anchors require explicit source_point/target_point coordinates.",
                element_ids=[source_id, target_id],
                evidence={"edge_id": edge_id},
            )
            continue
        via = [
            (float(point["x"]), float(point["y"]))
            for point in edge.get("via", []) or []
            if isinstance(point, Mapping) and "x" in point and "y" in point
        ]
        points = [start, *via, end]
        if route_kind == "straight" and via:
            add_finding(
                "SPEC_INVALID",
                "STRAIGHT_EDGE_HAS_VIA",
                "Straight edges cannot contain via points.",
                evidence={"edge_id": edge_id},
            )
            continue
        if route_kind in {"orthogonal", "polyline"} and not via:
            add_finding(
                "SPEC_INVALID",
                "EDGE_VIA_REQUIRED",
                "Orthogonal and polyline routes require explicit via points.",
                evidence={"edge_id": edge_id, "route": route_kind},
            )
            continue
        if route_kind == "orthogonal":
            diagonal_segments = [
                index
                for index, (first, second) in enumerate(zip(points, points[1:], strict=False))
                if not (
                    math.isclose(first[0], second[0], abs_tol=1e-6)
                    or math.isclose(first[1], second[1], abs_tol=1e-6)
                )
            ]
            if diagonal_segments:
                add_finding(
                    "REGION_REPLAN",
                    "ORTHOGONAL_EDGE_DIAGONAL",
                    "An orthogonal edge contains diagonal route segments.",
                    evidence={"edge_id": edge_id, "segment_indexes": diagonal_segments},
                )
        if measured_source is not None:
            outside_points = [
                {"x": point[0], "y": point[1]}
                for point in points
                if not (
                    0 <= point[0] <= float(measured_source["width_px"])
                    and 0 <= point[1] <= float(measured_source["height_px"])
                )
            ]
            if outside_points:
                add_finding(
                    "REGION_REPLAN",
                    "EDGE_ROUTE_OUTSIDE_CANVAS",
                    "Edge route contains points outside the source-pixel canvas.",
                    evidence={"edge_id": edge_id, "outside_points": outside_points},
                )
        clearance = float(edge.get("clearance_px", 2.0))
        allowed_crossings = {str(item) for item in edge.get("allowed_crossings", []) or []}
        obstacle_types_ignored = {"background", "panel", "connector", "line"}
        for obstacle_id, obstacle_box in valid_boxes.items():
            if obstacle_id in {source_id, target_id} or obstacle_id in allowed_crossings:
                continue
            obstacle_type = str(by_id[obstacle_id].get("type", "")).casefold()
            if obstacle_type in obstacle_types_ignored:
                continue
            expanded = _inflate_rect(obstacle_box, clearance)
            if any(
                _segment_intersects_rect(first, second, expanded)
                for first, second in zip(points, points[1:], strict=False)
            ):
                add_finding(
                    "REGION_REPLAN",
                    "EDGE_TEXT_COLLISION"
                    if obstacle_type in TEXT_TYPES
                    else "EDGE_OBSTACLE_COLLISION",
                    "Connector route violates declared obstacle/text clearance.",
                    element_ids=[obstacle_id],
                    evidence={"edge_id": edge_id, "clearance_px": clearance},
                    repair="Reroute the connector or explicitly declare a scientifically valid crossing.",
                )
        edge_routes[edge_id] = points

    for index, first_edge in enumerate(edges):
        first_id = str(first_edge.get("id", ""))
        if first_id not in edge_routes:
            continue
        first_allowed = {str(item) for item in first_edge.get("allowed_edge_crossings", []) or []}
        first_segments = list(zip(edge_routes[first_id], edge_routes[first_id][1:], strict=False))
        for second_edge in edges[index + 1 :]:
            second_id = str(second_edge.get("id", ""))
            if second_id not in edge_routes:
                continue
            second_allowed = {
                str(item) for item in second_edge.get("allowed_edge_crossings", []) or []
            }
            if second_id in first_allowed or first_id in second_allowed:
                continue
            second_segments = list(
                zip(edge_routes[second_id], edge_routes[second_id][1:], strict=False)
            )
            if any(
                _segments_cross(first_start, first_end, second_start, second_end)
                for first_start, first_end in first_segments
                for second_start, second_end in second_segments
            ):
                add_finding(
                    "REGION_REPLAN",
                    "EDGE_EDGE_CROSSING",
                    "Connector routes cross without an explicit crossing allowance.",
                    evidence={"edge_ids": [first_id, second_id]},
                    repair="Reroute one connector or explicitly justify the crossing.",
                )

    measurement_dpi = (
        float(spec.get("measurement_dpi", 96.0))
        if isinstance(spec.get("measurement_dpi", 96.0), (int, float))
        else 96.0
    )
    for element_id, element in by_id.items():
        if str(element.get("type", "")).casefold() != "text" or element_id not in valid_boxes:
            continue
        try:
            measurement = measure_text_fit(
                element,
                measurement_dpi=measurement_dpi,
                base_dir=resolved_base,
                font_search_paths=font_search_paths,
            )
        except (OSError, TypeError, ValueError) as exc:
            measurement = {
                "status": "INCONCLUSIVE",
                "element_id": element_id,
                "message": f"Text measurement failed: {exc}",
            }
        text_measurements.append(measurement)
        if measurement["status"] == "REGION_REPLAN":
            add_finding(
                "REGION_REPLAN",
                "TEXT_OVERFLOW",
                measurement["message"],
                element_ids=[element_id],
                evidence=measurement,
                repair="Rewrap, enlarge/reposition the box, or reduce font size without crossing the profile minimum.",
            )
        elif measurement["status"] == "INCONCLUSIVE":
            add_finding(
                "INCONCLUSIVE",
                "FONT_OR_TEXT_METRICS_INCONCLUSIVE",
                measurement["message"],
                element_ids=[element_id],
                evidence=measurement,
                repair="Resolve the exact font or verify this box using renderer readback before drawing is accepted.",
            )

    for formula in formulas:
        element_id = str(formula.get("element_id", ""))
        element = by_id.get(element_id)
        if element is None or element_id not in valid_boxes:
            continue
        candidate_id = formula.get("perception_candidate_id")
        if candidate_id and str(candidate_id) in candidate_by_id:
            # Formula placement is bound to the selected primary OCR observation.  A conflict
            # envelope may be retained for review, but cannot make an unrelated box look valid.
            candidate_box = candidate_by_id[str(candidate_id)].get("bbox_source")
            if isinstance(candidate_box, Mapping):
                bbox_matches, bbox_evidence = _candidate_bbox_matches(
                    valid_boxes[element_id], candidate_box
                )
                if not bbox_matches:
                    add_finding(
                        "SPEC_INVALID",
                        "FORMULA_CANDIDATE_BBOX_MISMATCH",
                        "Reviewed formula candidate is not geometrically bound to this formula box.",
                        element_ids=[element_id],
                        evidence={"candidate_id": candidate_id, **bbox_evidence},
                    )
            geometry_item = geometry_by_candidate.get(str(candidate_id))
            if isinstance(geometry_item, Mapping):
                geometry_diagnostics.append(
                    {
                        "kind": "FORMULA_INK_OBSERVATION",
                        "element_id": element_id,
                        "candidate_id": str(candidate_id),
                        "observation_status": geometry_item.get("status"),
                        "ink_bbox_source": geometry_item.get("ink_bbox"),
                        "edge_uncertainty_px": geometry_item.get("edge_uncertainty_px"),
                        "ink_bottom_alignment": geometry_item.get("baseline"),
                        "diagnostic_only": True,
                        "used_as_scene_geometry": False,
                        "authorizes_drawer": False,
                    }
                )
        measurement_element = element
        if str(element.get("type", "")).casefold() == "text":
            text_style = (
                element.get("text_style") if isinstance(element.get("text_style"), Mapping) else {}
            )
            inline_formula_style = {
                key: text_style[key]
                for key in ("font_size_px", "font_size_pt", "rotation_deg")
                if key in text_style
            }
            inline_formula_style["margin_px"] = 0
            measurement_element = {**element, "formula_style": inline_formula_style}
        measurement = measure_formula_fit(
            measurement_element,
            formula,
            measurement_dpi=measurement_dpi,
        )
        formula_measurements.append(measurement)
        if measurement["status"] == "REGION_REPLAN":
            add_finding(
                "REGION_REPLAN",
                "FORMULA_OVERFLOW",
                measurement["message"],
                element_ids=[element_id],
                evidence=measurement,
                repair="Resize/reposition the formula box or reduce the confirmed formula size.",
            )
        elif measurement["status"] == "INCONCLUSIVE":
            add_finding(
                "INCONCLUSIVE",
                "FORMULA_METRICS_INCONCLUSIVE",
                measurement["message"],
                element_ids=[element_id],
                evidence=measurement,
                repair="Use renderer readback for this confirmed formula before Drawer authorization.",
            )

    checked_snapshots: set[tuple[str, str]] = set()
    for label, evidence_path, payload in evidence_snapshots:
        snapshot_key = (str(evidence_path), _sha256_bytes(payload))
        if snapshot_key in checked_snapshots:
            continue
        checked_snapshots.add(snapshot_key)
        try:
            current_payload = evidence_path.read_bytes()
        except OSError as exc:
            add_finding(
                "SPEC_INVALID",
                "EVIDENCE_CHANGED_DURING_PREFLIGHT",
                "A hash-bound evidence file became unreadable during preflight.",
                evidence={"label": label, "path": str(evidence_path), "error": str(exc)},
            )
            continue
        if current_payload != payload:
            add_finding(
                "SPEC_INVALID",
                "EVIDENCE_CHANGED_DURING_PREFLIGHT",
                "A hash-bound evidence file changed while preflight was validating it.",
                evidence={
                    "label": label,
                    "path": str(evidence_path),
                    "snapshot_sha256": _sha256_bytes(payload),
                    "current_sha256": _sha256_bytes(current_payload),
                },
                repair="Rerun preflight against immutable run evidence.",
            )

    status = "PASS"
    for finding in findings:
        if STATUS_PRECEDENCE[finding["disposition"]] > STATUS_PRECEDENCE[status]:
            status = finding["disposition"]
    counts = Counter(finding["disposition"] for finding in findings)
    return {
        "status": status,
        "passed": status == "PASS",
        "source": measured_source,
        "canvas_pptx": canvas_record,
        "perception_manifest": perception_record,
        "geometry_manifest": geometry_record,
        "geometry_diagnostics": geometry_diagnostics,
        "schema_bindings": {
            str(path): {
                "path": str(path),
                "sha256": _sha256_bytes(payload),
                "size_bytes": len(payload),
            }
            for path, payload in sorted(schema_payloads.items(), key=lambda item: str(item[0]))
        },
        "formula_converter_receipts": formula_converter_receipts,
        "summary": {
            "element_count": len(elements),
            "edge_count": len(edges),
            "formula_count": len(formulas),
            "finding_count": len(findings),
            "spec_invalid": counts["SPEC_INVALID"],
            "region_replan": counts["REGION_REPLAN"],
            "inconclusive": counts["INCONCLUSIVE"],
        },
        "findings": findings,
        "text_measurements": text_measurements,
        "formula_measurements": formula_measurements,
    }


def _write_json(path: Path, payload: Mapping[str, Any], *, pretty: bool) -> None:
    path = resolve_output_path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    temporary = path.with_name(f".{path.name}.tmp")
    temporary.write_text(
        json.dumps(payload, ensure_ascii=False, indent=2 if pretty else None),
        encoding="utf-8",
    )
    os.replace(temporary, path)


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description=(
            "Preflight a hash-bound figure scene before PowerPoint mutation. "
            "Exit codes: PASS=0, REGION_REPLAN=3, INCONCLUSIVE=4, SPEC_INVALID=2."
        )
    )
    parser.add_argument("spec", help="Render-ready figure-spec JSON path.")
    parser.add_argument(
        "--source", help="Override source PNG path; metadata must still match the spec."
    )
    parser.add_argument(
        "--canvas-pptx",
        help="Override the blank canvas PPTX path; bytes and PageSetup must match the spec.",
    )
    parser.add_argument("--schema", default=str(DEFAULT_SCHEMA_PATH), help="JSON Schema path.")
    parser.add_argument(
        "--font-dir", action="append", default=[], help="Additional exact-font search directory."
    )
    parser.add_argument("--output", help="Optional JSON report path.")
    parser.add_argument("--pretty", action="store_true")
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    spec_path = Path(args.spec).expanduser().resolve()
    try:
        preflight_script_path = Path(__file__).resolve()
        preflight_script_bytes = preflight_script_path.read_bytes()
        spec_bytes = spec_path.read_bytes()
        spec = _strict_json_bytes(spec_bytes, label="figure specification")
        if not isinstance(spec, Mapping):
            raise ValueError("Figure spec root must be a JSON object.")
        source_override = str(Path(args.source).expanduser().resolve()) if args.source else None
        canvas_override = (
            str(Path(args.canvas_pptx).expanduser().resolve()) if args.canvas_pptx else None
        )
        report = preflight_scene(
            spec,
            source_path=source_override,
            canvas_pptx_path=canvas_override,
            schema_path=args.schema,
            base_dir=spec_path.parent,
            font_search_paths=args.font_dir or None,
        )
        schema_path = Path(args.schema).expanduser().resolve()
        report_schema_bindings = report.get("schema_bindings")
        if not isinstance(report_schema_bindings, Mapping):
            raise ValueError("preflight report is missing schema snapshot bindings")

        def bound_schema_hash(path: Path) -> str:
            binding = report_schema_bindings.get(str(path.resolve()))
            if not isinstance(binding, Mapping):
                raise ValueError(f"preflight did not snapshot schema {path.resolve()}")
            return str(binding.get("sha256", ""))

        report["receipt"] = {
            "authorized_for_drawer": report.get("status") == "PASS",
            "preflight_status": report.get("status"),
            "spec_path": str(spec_path),
            "spec_sha256": _sha256_bytes(spec_bytes),
            "schema_path": str(schema_path),
            "schema_sha256": bound_schema_hash(schema_path),
            "preflight_script_sha256": _sha256_bytes(preflight_script_bytes),
            "source_sha256": (report.get("source") or {}).get("sha256"),
            "perception_manifest_sha256": (report.get("perception_manifest") or {}).get("sha256"),
            "perception_manifest_schema_sha256": bound_schema_hash(DEFAULT_PERCEPTION_SCHEMA_PATH),
            "perception_review_receipt_sha256": (report.get("perception_manifest") or {}).get(
                "review_receipt_sha256"
            ),
            "perception_review_schema_sha256": bound_schema_hash(DEFAULT_REVIEW_SCHEMA_PATH),
            "geometry_manifest_sha256": (report.get("geometry_manifest") or {}).get("sha256"),
            "geometry_manifest_schema_sha256": bound_schema_hash(DEFAULT_GEOMETRY_SCHEMA_PATH),
            "geometry_status": (report.get("geometry_manifest") or {}).get("status"),
            "geometry_mode": (report.get("geometry_manifest") or {}).get("mode"),
            "geometry_promotion_allowed": (report.get("geometry_manifest") or {}).get(
                "promotion_allowed"
            ),
            "geometry_authorized_for_drawer": False,
            "geometry_artifacts": (report.get("geometry_manifest") or {}).get("artifacts", {}),
            "formula_converter_receipts": [
                {
                    "formula_id": item.get("formula_id"),
                    "path": item.get("path"),
                    "sha256": item.get("sha256"),
                    "latex_sha256": item.get("latex_sha256"),
                    "omml_sha256": item.get("omml_sha256"),
                    "semantic_omml_profile": item.get("semantic_omml_profile"),
                    "semantic_omml_sha256": item.get("semantic_omml_sha256"),
                }
                for item in report.get("formula_converter_receipts", [])
            ],
            "canvas_pptx_path": (report.get("canvas_pptx") or {}).get("path"),
            "canvas_pptx_sha256": (report.get("canvas_pptx") or {}).get("sha256"),
            "canvas_slide_width_emu": (report.get("canvas_pptx") or {}).get("slide_width_emu"),
            "canvas_slide_height_emu": (report.get("canvas_pptx") or {}).get("slide_height_emu"),
        }
        if spec_path.read_bytes() != spec_bytes:
            raise ValueError("figure specification changed during preflight")
        if preflight_script_path.read_bytes() != preflight_script_bytes:
            raise ValueError("preflight implementation changed during execution")
        for schema_key, binding in report_schema_bindings.items():
            if not isinstance(binding, Mapping):
                raise ValueError(f"invalid schema binding for {schema_key}")
            schema_file = Path(str(binding.get("path", schema_key)))
            if _sha256_bytes(schema_file.read_bytes()) != str(binding.get("sha256", "")):
                raise ValueError(f"schema changed during preflight: {schema_file}")
    except (OSError, ValueError, json.JSONDecodeError) as exc:
        report = {
            "status": "SPEC_INVALID",
            "passed": False,
            "summary": {
                "finding_count": 1,
                "spec_invalid": 1,
                "region_replan": 0,
                "inconclusive": 0,
            },
            "findings": [
                {
                    "id": "PF0001",
                    "severity": "MAJOR",
                    "disposition": "SPEC_INVALID",
                    "code": "SPEC_UNREADABLE",
                    "element_ids": [],
                    "message": str(exc),
                    "evidence": {},
                    "repair": "Provide a readable schema-valid figure spec.",
                }
            ],
            "text_measurements": [],
        }
    if args.output:
        _write_json(args.output, report, pretty=args.pretty)
    print(json.dumps(report, ensure_ascii=False, indent=2 if args.pretty else None))
    return {"PASS": 0, "SPEC_INVALID": 2, "REGION_REPLAN": 3, "INCONCLUSIVE": 4}[report["status"]]


if __name__ == "__main__":
    sys.exit(main())
