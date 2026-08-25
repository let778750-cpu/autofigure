"""Object-level visual contracts for typography, scale, and collisions.

Pixel similarity is deliberately insufficient for small scientific figure
objects: a large white background can hide a wrong icon scale, a displaced
label, or an oversized glyph.  This module binds a frozen reference inventory
to both the source SVG and the saved PowerPoint artifact.

The inventory is authored from the designated reference before construction.
It never infers expected geometry from the candidate being audited.
"""

from __future__ import annotations

import math
import re
import xml.etree.ElementTree as ET
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation

from tools import common
from tools.contracts import read_json, write_json
from tools.layout import Box, collect_svg_records
from tools.reference_inventory import (
    DEFAULT_SOURCE_ANISOTROPY_TOLERANCE,
    canonical_sha256,
    normalize_topology_contract,
    topology_contracts_sha256,
)
from tools.svggeom import Matrix

SVG_NS = "{http://www.w3.org/2000/svg}"
PT_PER_PX = 0.75
BASELINE_ASCENT = 0.95
DEFAULT_BBOX_TOLERANCE_PX = 2.0
DEFAULT_ASPECT_TOLERANCE = 0.03
DEFAULT_COLLISION_TOLERANCE_PX = 0.5

_STYLE_KEYS = {
    "fill",
    "font-family",
    "font-size",
    "font-style",
    "font-weight",
    "text-anchor",
    "text-decoration",
}
_COLLISION_KINDS = {"text", "formula", "arrow", "icon", "plot", "brace"}


@dataclass(frozen=True)
class TextSource:
    element_id: str
    lines: list[str]
    line_runs: list[list[dict[str, Any]]]
    anchor: tuple[float, float]
    alignment: str
    visual_box: Box | None


@dataclass(frozen=True)
class SourceIndex:
    elements: dict[str, ET.Element]
    descendant_ids: dict[str, list[str]]
    drawable_ids: set[str]
    boxes: dict[str, Box]
    texts: dict[str, TextSource]
    matrices: dict[str, Matrix]


@dataclass(frozen=True)
class BackendIndex:
    boxes: dict[str, Box]
    shapes: dict[str, list[Any]]
    binding_ids: set[str]
    readback_ids: set[str]


def _local_name(element: ET.Element) -> str:
    return element.tag.rsplit("}", 1)[-1]


def _parse_style(raw: str | None) -> dict[str, str]:
    result: dict[str, str] = {}
    if not raw:
        return result
    for item in raw.split(";"):
        if ":" not in item:
            continue
        key, value = item.split(":", 1)
        result[key.strip()] = value.strip()
    return result


def _style(element: ET.Element, inherited: dict[str, str]) -> dict[str, str]:
    result = dict(inherited)
    inline = _parse_style(element.get("style"))
    for key in _STYLE_KEYS:
        if element.get(key) is not None:
            result[key] = str(element.get(key))
        if key in inline:
            result[key] = inline[key]
    return result


def _box(raw: Any) -> Box | None:
    if not isinstance(raw, list) or len(raw) != 4:
        return None
    try:
        values = [float(item) for item in raw]
    except (TypeError, ValueError):
        return None
    if (
        not all(math.isfinite(item) for item in values)
        or values[2] < 0
        or values[3] < 0
    ):
        return None
    return Box(*values)


def _attribute_box(element: ET.Element) -> Box | None:
    raw = element.get("data-visual-bbox")
    if not raw:
        return None
    try:
        values = [float(item) for item in re.split(r"[\s,]+", raw.strip()) if item]
    except ValueError:
        return None
    return _box(values)


def _union(boxes: Iterable[Box]) -> Box | None:
    materialized = list(boxes)
    if not materialized:
        return None
    left = min(box.x for box in materialized)
    top = min(box.y for box in materialized)
    right = max(box.right for box in materialized)
    bottom = max(box.bottom for box in materialized)
    return Box(left, top, right - left, bottom - top)


def _number(raw: str | None, default: float) -> float:
    if raw is None:
        return default
    match = re.match(r"\s*([-+]?(?:\d*\.\d+|\d+\.?))", raw)
    return default if match is None else float(match.group(1))


def _source_text(
    element: ET.Element,
    inherited_style: dict[str, str],
    child_styles: dict[int, dict[str, str]],
) -> TextSource:
    lines: list[list[dict[str, Any]]] = []
    current: list[dict[str, Any]] = []
    baseline = _number(element.get("y"), 0.0)

    def append(text: str | None, style: dict[str, str]) -> None:
        if text and text.strip():
            current.append({"text": text, "style": style})

    def flush() -> None:
        nonlocal current
        if current:
            lines.append(current)
            current = []

    append(element.text, inherited_style)
    first_positioned_x: float | None = None
    first_positioned_y: float | None = None
    for child in element:
        if _local_name(child) != "tspan":
            continue
        positioned = any(child.get(name) is not None for name in ("x", "y", "dy"))
        if positioned and current:
            flush()
        if child.get("x") is not None and first_positioned_x is None:
            first_positioned_x = _number(child.get("x"), 0.0)
        if child.get("y") is not None:
            baseline = _number(child.get("y"), baseline)
            if first_positioned_y is None:
                first_positioned_y = baseline
        elif child.get("dy") is not None:
            baseline += _number(child.get("dy"), 0.0)
            if first_positioned_y is None:
                first_positioned_y = baseline
        append(child.text, child_styles.get(id(child), inherited_style))
        append(child.tail, inherited_style)
    flush()
    rendered_lines = [
        "".join(str(run["text"]) for run in line).strip()
        for line in lines
    ]
    anchor_x = _number(element.get("x"), first_positioned_x or 0.0)
    anchor_y = _number(element.get("y"), first_positioned_y or baseline)
    alignment = {
        "middle": "center",
        "end": "right",
        "start": "left",
    }.get(inherited_style.get("text-anchor", "start"), "left")
    return TextSource(
        element_id=element.get("id") or "",
        lines=rendered_lines,
        line_runs=lines,
        anchor=(anchor_x, anchor_y),
        alignment=alignment,
        visual_box=_attribute_box(element),
    )


def _build_source_index(svg_path: Path) -> SourceIndex:
    root = ET.parse(svg_path).getroot()
    elements = {
        element.get("id"): element
        for element in root.iter()
        if element.get("id")
    }
    records = collect_svg_records(root)
    primitive_boxes = {
        record.element_id: record.box
        for record in records
        if record.element_id and record.box is not None
    }
    matrices = {
        record.element_id: record.matrix
        for record in records
        if record.element_id
    }
    drawable_ids = {
        record.element_id
        for record in records
        if record.element_id
    }
    styles: dict[int, dict[str, str]] = {}

    def walk_style(element: ET.Element, inherited: dict[str, str]) -> None:
        current = _style(element, inherited)
        styles[id(element)] = current
        if _local_name(element) in {"defs", "marker"}:
            return
        for child in element:
            walk_style(child, current)

    walk_style(root, {})
    texts: dict[str, TextSource] = {}
    for element_id, element in elements.items():
        if _local_name(element) == "text":
            texts[element_id] = _source_text(
                element,
                styles.get(id(element), {}),
                styles,
            )

    boxes = dict(primitive_boxes)
    boxes.update(
        {
            element_id: text.visual_box
            for element_id, text in texts.items()
            if text.visual_box is not None
        }
    )
    descendant_ids: dict[str, list[str]] = {}
    for element_id, element in elements.items():
        descendants = [
            child.get("id")
            for child in element.iter()
            if child is not element
            and child.get("id")
            and _local_name(child) not in {"title", "desc", "metadata", "defs", "marker"}
        ]
        descendant_ids[element_id] = [item for item in descendants if item]
    # Groups do not become PowerPoint objects in the offline compiler.  Their
    # semantic bounds are therefore the union of their bound descendants.
    changed = True
    while changed:
        changed = False
        for element_id, descendants in descendant_ids.items():
            union = _union(boxes[item] for item in descendants if item in boxes)
            if union is not None and boxes.get(element_id) != union:
                boxes[element_id] = union
                changed = True
    return SourceIndex(elements, descendant_ids, drawable_ids, boxes, texts, matrices)


def _build_backend_index(run: common.Run) -> BackendIndex:
    from tools.live_bridge import _shape_bounds, _walk_shapes

    presentation = Presentation(run.pptx_path)
    by_identity = {
        (int(shape.shape_id), shape.name): shape
        for shape in _walk_shapes(presentation.slides[0].shapes)
    }
    shapes: dict[str, list[Any]] = {}
    bindings = read_json(run.bindings_path)
    binding_ids: set[str] = set()
    readback_ids: set[str] = set()
    for binding in bindings.get("bindings", []):
        element_id = binding.get("element_id")
        shape_id = binding.get("shape_id")
        shape_name = binding.get("shape_name")
        if not isinstance(element_id, str):
            continue
        binding_ids.add(element_id)
        if binding.get("readback_found") is True:
            readback_ids.add(element_id)
        shape = by_identity.get((shape_id, shape_name))
        if shape is not None:
            shapes.setdefault(element_id, []).append(shape)
    for binding in bindings.get("logical_group_bindings", []):
        element_id = binding.get("element_id")
        backend_ids = binding.get("backend_object_ids")
        backend_names = binding.get("backend_object_names")
        if not isinstance(element_id, str):
            continue
        binding_ids.add(element_id)
        if binding.get("readback_found") is True:
            readback_ids.add(element_id)
        if (
            not isinstance(backend_ids, list)
            or not isinstance(backend_names, list)
            or not backend_ids
            or len(backend_ids) != len(backend_names)
        ):
            continue
        bound_shapes = [
            by_identity.get((shape_id, shape_name))
            for shape_id, shape_name in zip(
                backend_ids, backend_names, strict=True
            )
        ]
        if all(shape is not None for shape in bound_shapes):
            shapes[element_id] = [shape for shape in bound_shapes if shape is not None]
    boxes: dict[str, Box] = {}
    for element_id, bound_shapes in shapes.items():
        values = []
        for shape in bound_shapes:
            raw = _shape_bounds(shape)
            values.append(Box(raw["x"], raw["y"], raw["width"], raw["height"]))
        union = _union(values)
        if union is not None:
            boxes[element_id] = union
    return BackendIndex(boxes, shapes, binding_ids, readback_ids)


def _expanded_ids(source: SourceIndex, element_ids: list[str]) -> list[str]:
    expanded: list[str] = []
    for element_id in element_ids:
        if element_id in source.boxes and element_id not in source.descendant_ids:
            expanded.append(element_id)
            continue
        descendants = source.descendant_ids.get(element_id, [])
        bound_descendants = [item for item in descendants if item in source.boxes]
        if bound_descendants:
            expanded.extend(bound_descendants)
        elif element_id in source.boxes:
            expanded.append(element_id)
        else:
            expanded.append(element_id)
    return list(dict.fromkeys(expanded))


def _object_box(
    index: SourceIndex | BackendIndex,
    source: SourceIndex,
    element_ids: list[str],
) -> Box | None:
    expanded = _expanded_ids(source, element_ids)
    return _union(index.boxes[item] for item in expanded if item in index.boxes)


def _box_delta(actual: Box, expected: Box) -> dict[str, float]:
    return {
        "x": round(actual.x - expected.x, 6),
        "y": round(actual.y - expected.y, 6),
        "width": round(actual.width - expected.width, 6),
        "height": round(actual.height - expected.height, 6),
    }


def _bbox_pass(actual: Box, expected: Box, tolerance: float) -> bool:
    return max(abs(item) for item in _box_delta(actual, expected).values()) <= tolerance


def _aspect_error(actual: Box, expected: Box) -> float:
    if min(actual.height, expected.height) <= 0:
        return math.inf
    expected_ratio = expected.width / expected.height
    if expected_ratio == 0:
        return math.inf
    return abs((actual.width / actual.height) / expected_ratio - 1.0)


def _matrix_anisotropy(matrix: Matrix) -> tuple[float, float, float]:
    """Return (relative anisotropy, major scale, minor scale) for an affine map.

    Singular values make the result invariant to rotation and reflection.  A
    uniform scale therefore reports zero while non-uniform scale and shear
    remain visible even when the final object bbox happens to match.
    """

    trace = matrix.a**2 + matrix.b**2 + matrix.c**2 + matrix.d**2
    determinant = matrix.a * matrix.d - matrix.b * matrix.c
    discriminant = max(0.0, trace**2 - 4.0 * determinant**2)
    major_squared = max(0.0, (trace + math.sqrt(discriminant)) / 2.0)
    minor_squared = max(0.0, (trace - math.sqrt(discriminant)) / 2.0)
    major = math.sqrt(major_squared)
    minor = math.sqrt(minor_squared)
    anisotropy = math.inf if minor <= 1e-12 else major / minor - 1.0
    return anisotropy, major, minor


def _audit_source_transform_anisotropy(
    object_spec: dict[str, Any],
    source: SourceIndex,
    element_ids: list[str],
) -> dict[str, Any] | None:
    """Audit source-only affine distortion for semantic microassets."""

    if object_spec.get("kind") not in {"icon", "plot"}:
        return None
    visual = object_spec.get("visual", {})
    allow = bool(visual.get("allow_source_anisotropic_scale", False))
    basis = visual.get("source_anisotropy_basis") if allow else None
    expanded_ids = _expanded_ids(source, element_ids)
    element_rows: list[dict[str, Any]] = []
    for element_id in expanded_ids:
        matrix = source.matrices.get(element_id)
        if matrix is None:
            continue
        anisotropy, major, minor = _matrix_anisotropy(matrix)
        element_rows.append(
            {
                "element_id": element_id,
                "matrix": [
                    round(matrix.a, 9),
                    round(matrix.b, 9),
                    round(matrix.c, 9),
                    round(matrix.d, 9),
                    round(matrix.e, 9),
                    round(matrix.f, 9),
                ],
                "major_scale": round(major, 9),
                "minor_scale": round(minor, 9),
                "anisotropy": (
                    None if not math.isfinite(anisotropy) else round(anisotropy, 9)
                ),
                "pass": allow
                or anisotropy <= DEFAULT_SOURCE_ANISOTROPY_TOLERANCE,
            }
        )
    offenders = [row for row in element_rows if not row["pass"]]
    return {
        "policy": "explicitly_allowed" if allow else "forbid",
        "allow_source_anisotropic_scale": allow,
        "source_anisotropy_basis": basis,
        "tolerance": DEFAULT_SOURCE_ANISOTROPY_TOLERANCE,
        "elements": element_rows,
        "offenders": offenders,
        "pass": not offenders,
    }


def _font_family(raw: Any) -> str:
    return str(raw or "").split(",", 1)[0].strip().strip("'\"").casefold()


def _weight(raw: Any) -> int:
    value = str(raw or "400").strip().casefold()
    if value == "bold":
        return 700
    if value == "normal":
        return 400
    try:
        return int(float(value))
    except ValueError:
        return 400


def _backend_alignment(paragraph: Any) -> str:
    raw = getattr(paragraph, "alignment", None)
    value = getattr(raw, "value", raw)
    return {1: "left", 2: "center", 3: "right"}.get(value, "left")


def _backend_text(shape: Any) -> dict[str, Any]:
    frame = shape.text_frame
    paragraphs = [paragraph for paragraph in frame.paragraphs if paragraph.text]
    lines = [paragraph.text for paragraph in paragraphs]
    line_runs: list[list[dict[str, Any]]] = []
    for paragraph in paragraphs:
        runs = []
        for run in paragraph.runs:
            size = getattr(run.font.size, "pt", None)
            runs.append(
                {
                    "text": run.text,
                    "font_family": run.font.name,
                    "font_size_pt": None if size is None else float(size),
                    "bold": run.font.bold is True,
                    "italic": run.font.italic is True,
                    "underline": bool(run.font.underline),
                }
            )
        line_runs.append(runs)
    auto_size = frame.auto_size
    auto_size_value = getattr(auto_size, "value", auto_size)
    return {
        "lines": lines,
        "line_runs": line_runs,
        "alignment": _backend_alignment(paragraphs[0]) if paragraphs else "left",
        "auto_size": auto_size_value,
        "word_wrap": frame.word_wrap,
    }


def _expected_text(typography: dict[str, Any]) -> str:
    exact = typography.get("exact_text", "")
    if isinstance(exact, list):
        return "\n".join(str(item) for item in exact)
    return str(exact)


def _expected_font_px(typography: dict[str, Any]) -> float | None:
    raw = typography.get("font_size_px")
    if raw is not None:
        return float(raw)
    raw = typography.get("font_size_pt")
    return None if raw is None else float(raw) / PT_PER_PX


def _line_weights(typography: dict[str, Any], line_count: int) -> list[int]:
    raw = typography.get("font_weights")
    if isinstance(raw, list) and len(raw) == line_count:
        return [_weight(item) for item in raw]
    return [_weight(typography.get("font_weight", 400))] * line_count


def _finding(code: str, target: str, message: str, metrics: dict[str, Any] | None = None) -> dict[str, Any]:
    value: dict[str, Any] = {
        "code": code,
        "target": target,
        "severity": "error",
        "message": message,
    }
    if metrics:
        value["metrics"] = metrics
    return value


def _audit_typography(
    object_spec: dict[str, Any],
    source: SourceIndex,
    backend: BackendIndex,
    findings: list[dict[str, Any]],
) -> dict[str, Any]:
    object_id = str(object_spec["id"])
    typography = object_spec.get("typography")
    if not isinstance(typography, dict):
        findings.append(_finding("V10", object_id, "text/formula object lacks typography contract"))
        return {"id": object_id, "pass": False}
    element_ids = [str(item) for item in object_spec.get("element_ids", [])]
    text_ids = [item for item in element_ids if item in source.texts]
    if len(text_ids) != 1:
        findings.append(
            _finding("V11", object_id, "typography object must bind exactly one SVG text element")
        )
        return {"id": object_id, "pass": False}
    text_id = text_ids[0]
    source_text = source.texts[text_id]
    expected_text = _expected_text(typography)
    expected_lines = expected_text.split("\n") if expected_text else []
    expected_line_count = int(typography.get("line_count", len(expected_lines)))
    expected_font_px = _expected_font_px(typography)
    tolerance_px = float(typography.get("font_size_tolerance_px", 0.5))
    expected_family = _font_family(typography.get("font_family"))
    expected_weights = _line_weights(typography, expected_line_count)
    expected_italic = str(typography.get("font_style", "normal")).casefold() in {
        "italic",
        "oblique",
    }
    expected_underline = typography.get("underline") is True
    expected_alignment = str(typography.get("alignment", source_text.alignment)).casefold()

    source_text_value = "\n".join(source_text.lines)
    source_pass = source_text_value == expected_text and len(source_text.lines) == expected_line_count
    source_line_metrics: list[dict[str, Any]] = []
    for index, runs in enumerate(source_text.line_runs):
        line_weight = expected_weights[min(index, len(expected_weights) - 1)] if expected_weights else 400
        run_rows = []
        for run in runs:
            style = run["style"]
            size_px = _number(style.get("font-size"), 16.0)
            family = _font_family(style.get("font-family") or "Arial")
            weight = _weight(style.get("font-weight"))
            italic = str(style.get("font-style", "normal")).casefold() in {"italic", "oblique"}
            underline = "underline" in str(style.get("text-decoration", "")).casefold().split()
            run_pass = (
                (expected_font_px is None or abs(size_px - expected_font_px) <= tolerance_px)
                and (not expected_family or family == expected_family)
                and ((weight >= 600) == (line_weight >= 600))
                and italic == expected_italic
                and underline == expected_underline
            )
            source_pass = source_pass and run_pass
            run_rows.append(
                {
                    "text": run["text"],
                    "font_family": family,
                    "font_size_px": size_px,
                    "font_weight": weight,
                    "italic": italic,
                    "underline": underline,
                    "pass": run_pass,
                }
            )
        source_line_metrics.append({"runs": run_rows})
    source_pass = source_pass and source_text.alignment == expected_alignment
    if not source_pass:
        findings.append(_finding("V12", object_id, "SVG typography differs from frozen reference inventory"))

    bound_shapes = []
    for expanded in _expanded_ids(source, element_ids):
        bound_shapes.extend(backend.shapes.get(expanded, []))
    text_shapes = [shape for shape in bound_shapes if getattr(shape, "has_text_frame", False)]
    backend_pass = len(text_shapes) == 1
    backend_metrics: dict[str, Any] | None = None
    if backend_pass:
        backend_metrics = _backend_text(text_shapes[0])
        backend_value = "\n".join(backend_metrics["lines"])
        backend_pass = (
            backend_value == expected_text
            and len(backend_metrics["lines"]) == expected_line_count
            and backend_metrics["alignment"] == expected_alignment
            and backend_metrics["auto_size"] in (None, 0)
            and backend_metrics["word_wrap"] is not True
        )
        for index, runs in enumerate(backend_metrics["line_runs"]):
            line_weight = expected_weights[min(index, len(expected_weights) - 1)] if expected_weights else 400
            if not runs:
                backend_pass = False
                continue
            for run in runs:
                size_pt = run["font_size_pt"]
                backend_pass = backend_pass and (
                    size_pt is not None
                    and (
                        expected_font_px is None
                        or abs(size_pt - expected_font_px * PT_PER_PX)
                        <= tolerance_px * PT_PER_PX
                    )
                    and (not expected_family or _font_family(run["font_family"]) == expected_family)
                    and (run["bold"] == (line_weight >= 600))
                    and (run["italic"] == expected_italic)
                    and (run["underline"] == expected_underline)
                )
    if not backend_pass:
        findings.append(
            _finding(
                "V13",
                object_id,
                "saved PowerPoint text, font, line count, alignment, or autofit differs from contract",
            )
        )
    return {
        "id": object_id,
        "source": {
            "text": source_text_value,
            "lines": source_text.lines,
            "alignment": source_text.alignment,
            "line_metrics": source_line_metrics,
            "pass": source_pass,
        },
        "backend": {**(backend_metrics or {}), "object_count": len(text_shapes), "pass": backend_pass},
        "pass": source_pass and backend_pass,
    }


def _translated_text_backend_box(
    object_spec: dict[str, Any],
    source: SourceIndex,
    backend: BackendIndex,
) -> Box | None:
    expected = _box(object_spec.get("bbox"))
    if expected is None:
        return None
    element_ids = [str(item) for item in object_spec.get("element_ids", [])]
    text_ids = [item for item in element_ids if item in source.texts]
    if len(text_ids) != 1:
        return None
    source_text = source.texts[text_ids[0]]
    shapes = backend.shapes.get(text_ids[0], [])
    shapes = [shape for shape in shapes if getattr(shape, "has_text_frame", False)]
    if len(shapes) != 1:
        return None
    from tools.live_bridge import _shape_bounds

    raw = _shape_bounds(shapes[0])
    selection = Box(raw["x"], raw["y"], raw["width"], raw["height"])
    if source_text.alignment == "center":
        backend_anchor_x = selection.center_x
    elif source_text.alignment == "right":
        backend_anchor_x = selection.right
    else:
        backend_anchor_x = selection.x
    typography = object_spec.get("typography", {})
    font_px = _expected_font_px(typography) or 16.0
    backend_anchor_y = selection.y + font_px * BASELINE_ASCENT
    dx = backend_anchor_x - source_text.anchor[0]
    dy = backend_anchor_y - source_text.anchor[1]
    return Box(expected.x + dx, expected.y + dy, expected.width, expected.height)


def _gap(a: Box, b: Box, axis: str) -> float:
    if axis == "x":
        if a.right <= b.x:
            return b.x - a.right
        if b.right <= a.x:
            return a.x - b.right
        return -min(a.right, b.right) + max(a.x, b.x)
    if axis == "y":
        if a.bottom <= b.y:
            return b.y - a.bottom
        if b.bottom <= a.y:
            return a.y - b.bottom
        return -min(a.bottom, b.bottom) + max(a.y, b.y)
    raise ValueError("clearance axis must be x or y")


def _intersection(a: Box, b: Box) -> tuple[float, float]:
    return max(0.0, min(a.right, b.right) - max(a.x, b.x)), max(
        0.0,
        min(a.bottom, b.bottom) - max(a.y, b.y),
    )


def _topology_scope_ids(
    object_spec: dict[str, Any],
    contract: dict[str, Any],
    source: SourceIndex,
) -> set[str]:
    root_id = contract.get("scope_element_id")
    if root_id is None and str(object_spec.get("id", "")) in source.descendant_ids:
        root_id = str(object_spec["id"])
    if isinstance(root_id, str) and root_id in source.elements:
        scope = {
            element_id
            for element_id in source.descendant_ids.get(root_id, [])
            if element_id in source.drawable_ids
        }
        if root_id in source.drawable_ids:
            scope.add(root_id)
    else:
        scope = {
            str(element_id)
            for element_id in object_spec.get("element_ids", [])
            if str(element_id) in source.drawable_ids
        }
    for element_id in source.drawable_ids:
        if element_id in contract["role_mapping"] or any(
            re.fullmatch(pattern, element_id)
            for pattern in contract["role_patterns"].values()
        ):
            scope.add(element_id)
    return scope


def _classify_topology_ids(
    element_ids: set[str],
    contract: dict[str, Any],
) -> tuple[dict[str, set[str]], list[str], list[str]]:
    roles = {role: set() for role in contract["role_counts"]}
    ambiguous: list[str] = []
    unclassified: list[str] = []
    for element_id in sorted(element_ids):
        resolved: set[str] = set()
        mapped = contract["role_mapping"].get(element_id)
        if mapped is not None:
            resolved.add(mapped)
        for role, pattern in contract["role_patterns"].items():
            if re.fullmatch(pattern, element_id):
                resolved.add(role)
        if not resolved:
            unclassified.append(element_id)
        elif len(resolved) > 1:
            ambiguous.append(element_id)
        else:
            role = next(iter(resolved))
            roles.setdefault(role, set()).add(element_id)
    return roles, ambiguous, unclassified


def _topology_components(
    element_ids: set[str],
    pairs: list[dict[str, str]],
    relations: list[dict[str, Any]],
) -> int:
    if not element_ids:
        return 0
    adjacency = {element_id: set() for element_id in element_ids}

    def connect(left: Any, right: Any) -> None:
        if left in adjacency and right in adjacency and left != right:
            adjacency[left].add(right)
            adjacency[right].add(left)

    for pair in pairs:
        connect(pair["a"], pair["b"])
    for relation in relations:
        element_id = relation.get("element_id")
        if element_id in adjacency:
            connect(element_id, relation["source_id"])
            connect(element_id, relation["target_id"])
        else:
            connect(relation["source_id"], relation["target_id"])

    remaining = set(element_ids)
    count = 0
    while remaining:
        count += 1
        pending = [remaining.pop()]
        while pending:
            current = pending.pop()
            neighbors = adjacency[current] & remaining
            remaining.difference_update(neighbors)
            pending.extend(neighbors)
    return count


def _missing_topology_links(
    available_ids: set[str],
    contract: dict[str, Any],
) -> tuple[list[str], list[str]]:
    missing_pairs = [
        pair["id"]
        for pair in contract["required_pairs"]
        if pair["a"] not in available_ids or pair["b"] not in available_ids
    ]
    missing_relations = []
    for relation in contract["relations"]:
        participants = {relation["source_id"], relation["target_id"]}
        if relation.get("element_id") is not None:
            participants.add(relation["element_id"])
        if not participants.issubset(available_ids):
            missing_relations.append(str(relation["id"]))
    return missing_pairs, missing_relations


def _source_relation_drift(
    source: SourceIndex,
    scope_ids: set[str],
    contract: dict[str, Any],
) -> tuple[list[str], list[str], list[str]]:
    expected_by_element = {
        relation["element_id"]: relation
        for relation in contract["relations"]
        if relation.get("element_id") is not None
    }
    missing: list[str] = []
    malformed: list[str] = []
    unexpected: list[str] = []
    for element_id, expected in sorted(expected_by_element.items()):
        element = source.elements.get(element_id)
        if element is None:
            continue
        source_id = element.get("data-source-id")
        target_id = element.get("data-target-id")
        relation = element.get("data-topology-relation")
        if source_id is None or target_id is None or relation is None:
            missing.append(element_id)
            continue
        if (
            source_id != expected["source_id"]
            or target_id != expected["target_id"]
            or relation != expected["relation"]
        ):
            malformed.append(element_id)
    for element_id in sorted(scope_ids - set(expected_by_element)):
        element = source.elements.get(element_id)
        if element is None:
            continue
        if any(
            element.get(name) is not None
            for name in (
                "data-source-id",
                "data-target-id",
                "data-topology-relation",
                "data-relation",
            )
        ):
            unexpected.append(element_id)
    return missing, malformed, unexpected


def _source_pair_drift(
    source: SourceIndex,
    scope_ids: set[str],
    contract: dict[str, Any],
) -> tuple[list[str], list[list[str]]]:
    expected = {
        tuple(sorted((pair["a"], pair["b"]))): pair["id"]
        for pair in contract["required_pairs"]
    }
    actual: set[tuple[str, str]] = set()
    declarations: dict[str, set[str]] = {}
    for element_id in sorted(scope_ids):
        element = source.elements.get(element_id)
        if element is None:
            continue
        raw = element.get("data-pair-with")
        if raw is None:
            continue
        declared = declarations.setdefault(element_id, set())
        for paired_id in re.split(r"[\s,]+", raw.strip()):
            if paired_id and paired_id != element_id:
                declared.add(paired_id)
                actual.add(tuple(sorted((element_id, paired_id))))
    missing = [
        pair_id
        for pair, pair_id in expected.items()
        if pair[1] not in declarations.get(pair[0], set())
        or pair[0] not in declarations.get(pair[1], set())
    ]
    unexpected = [list(pair) for pair in sorted(actual - set(expected))]
    return missing, unexpected


def _audit_topology(
    object_spec: dict[str, Any],
    source: SourceIndex,
    backend: BackendIndex,
    findings: list[dict[str, Any]],
) -> dict[str, Any] | None:
    contract, errors = normalize_topology_contract(object_spec)
    if contract is None and not errors:
        return None
    object_id = str(object_spec.get("id", "unknown"))
    if contract is None:
        findings.append(
            _finding(
                "V30",
                object_id,
                "object topology contract is invalid",
                {"errors": errors},
            )
        )
        return {
            "id": object_id,
            "contract_sha256": canonical_sha256(object_spec.get("topology_contract")),
            "errors": errors,
            "pass": False,
        }

    expected_roles = {
        role: set(element_ids)
        for role, element_ids in contract["expected_roles"].items()
    }
    expected_ids = set().union(*expected_roles.values())
    source_scope = _topology_scope_ids(object_spec, contract, source)
    source_roles, source_ambiguous, source_unclassified = _classify_topology_ids(
        source_scope,
        contract,
    )
    source_missing = {
        role: sorted(expected_roles[role] - source_roles.get(role, set()))
        for role in expected_roles
    }
    source_extra = {
        role: sorted(source_roles.get(role, set()) - expected_roles[role])
        for role in expected_roles
    }
    source_counts = {
        role: len(source_roles.get(role, set()))
        for role in contract["role_counts"]
    }
    source_role_pass = (
        source_counts == contract["role_counts"]
        and not any(source_missing.values())
        and not any(source_extra.values())
        and not source_ambiguous
        and not source_unclassified
    )
    if not source_role_pass:
        findings.append(
            _finding(
                "V31",
                object_id,
                "SVG topology role counts or closed-world membership differ from contract",
                {
                    "role_counts": source_counts,
                    "missing_ids": source_missing,
                    "extra_ids": source_extra,
                    "ambiguous_ids": source_ambiguous,
                    "unclassified_ids": source_unclassified,
                },
            )
        )

    def backend_scope_for(element_ids: set[str]) -> set[str]:
        return {
            element_id
            for element_id in element_ids
            if element_id in source_scope
            or element_id in contract["role_mapping"]
            or any(
                re.fullmatch(pattern, element_id)
                for pattern in contract["role_patterns"].values()
            )
        }

    backend_available = set(backend.shapes) & backend.readback_ids
    backend_scope = backend_scope_for(backend_available)
    backend_roles, backend_ambiguous, backend_unclassified = _classify_topology_ids(
        backend_scope,
        contract,
    )
    backend_missing = {
        role: sorted(expected_roles[role] - backend_roles.get(role, set()))
        for role in expected_roles
    }
    backend_extra = {
        role: sorted(backend_roles.get(role, set()) - expected_roles[role])
        for role in expected_roles
    }
    backend_counts = {
        role: len(backend_roles.get(role, set()))
        for role in contract["role_counts"]
    }
    binding_missing = sorted(expected_ids - backend.binding_ids)
    readback_missing = sorted(expected_ids - backend.readback_ids)
    shape_missing = sorted(expected_ids - set(backend.shapes))
    binding_roles, _, _ = _classify_topology_ids(
        backend_scope_for(backend.binding_ids),
        contract,
    )
    readback_roles, _, _ = _classify_topology_ids(
        backend_scope_for(backend.readback_ids),
        contract,
    )
    binding_extra = {
        role: sorted(binding_roles.get(role, set()) - expected_roles[role])
        for role in expected_roles
    }
    readback_extra = {
        role: sorted(readback_roles.get(role, set()) - expected_roles[role])
        for role in expected_roles
    }
    backend_role_pass = (
        backend_counts == contract["role_counts"]
        and not any(backend_missing.values())
        and not any(backend_extra.values())
        and not backend_ambiguous
        and not backend_unclassified
    )
    if not backend_role_pass:
        findings.append(
            _finding(
                "V32",
                object_id,
                "PowerPoint topology role counts or closed-world membership differ from contract",
                {
                    "role_counts": backend_counts,
                    "missing_ids": backend_missing,
                    "extra_ids": backend_extra,
                    "ambiguous_ids": backend_ambiguous,
                    "unclassified_ids": backend_unclassified,
                },
            )
        )
    backend_evidence_pass = not (
        binding_missing
        or readback_missing
        or shape_missing
        or any(binding_extra.values())
        or any(readback_extra.values())
    )
    if not backend_evidence_pass:
        findings.append(
            _finding(
                "V33",
                object_id,
                "PowerPoint topology bindings/readback are incomplete",
                {
                    "binding_missing": binding_missing,
                    "readback_missing": readback_missing,
                    "shape_missing": shape_missing,
                    "binding_extra": binding_extra,
                    "readback_extra": readback_extra,
                },
            )
        )

    source_ids = set().union(*source_roles.values())
    backend_ids = set().union(*backend_roles.values())
    source_missing_pairs, source_missing_relations = _missing_topology_links(
        source_ids,
        contract,
    )
    missing_relation_metadata, malformed_relations, unexpected_relations = (
        _source_relation_drift(
            source,
            source_scope,
            contract,
        )
    )
    missing_pair_metadata, unexpected_pairs = _source_pair_drift(
        source,
        source_scope,
        contract,
    )
    source_link_pass = not (
        source_missing_pairs
        or source_missing_relations
        or missing_pair_metadata
        or missing_relation_metadata
        or malformed_relations
        or unexpected_relations
        or unexpected_pairs
    )
    if not source_link_pass:
        findings.append(
            _finding(
                "V34",
                object_id,
                "SVG required pair/relation closure differs from topology contract",
                {
                    "missing_pairs": source_missing_pairs,
                    "missing_relations": source_missing_relations,
                    "missing_pair_metadata": missing_pair_metadata,
                    "missing_relation_metadata": missing_relation_metadata,
                    "malformed_relations": malformed_relations,
                    "unexpected_relations": unexpected_relations,
                    "unexpected_pairs": unexpected_pairs,
                },
            )
        )
    backend_missing_pairs, backend_missing_relations = _missing_topology_links(
        backend_ids,
        contract,
    )
    backend_link_pass = not backend_missing_pairs and not backend_missing_relations
    if not backend_link_pass:
        findings.append(
            _finding(
                "V35",
                object_id,
                "PowerPoint required pair/relation bindings are incomplete",
                {
                    "missing_pairs": backend_missing_pairs,
                    "missing_relations": backend_missing_relations,
                },
            )
        )

    source_components = _topology_components(
        source_ids,
        contract["required_pairs"],
        contract["relations"],
    )
    backend_components = _topology_components(
        backend_ids,
        contract["required_pairs"],
        contract["relations"],
    )
    source_component_pass = source_components == contract["component_count"]
    backend_component_pass = backend_components == contract["component_count"]
    if not source_component_pass:
        findings.append(
            _finding(
                "V36",
                object_id,
                "SVG topology component count differs from contract",
                {
                    "expected": contract["component_count"],
                    "actual": source_components,
                },
            )
        )
    if not backend_component_pass:
        findings.append(
            _finding(
                "V37",
                object_id,
                "PowerPoint topology component count differs from contract",
                {
                    "expected": contract["component_count"],
                    "actual": backend_components,
                },
            )
        )

    source_pass = source_role_pass and source_link_pass and source_component_pass
    backend_pass = (
        backend_role_pass
        and backend_evidence_pass
        and backend_link_pass
        and backend_component_pass
    )
    return {
        "id": object_id,
        "contract_sha256": canonical_sha256(object_spec["topology_contract"]),
        "expected": {
            "role_counts": contract["role_counts"],
            "role_ids": {
                role: sorted(element_ids)
                for role, element_ids in expected_roles.items()
            },
            "required_pair_count": len(contract["required_pairs"]),
            "relation_count": len(contract["relations"]),
            "component_count": contract["component_count"],
        },
        "source": {
            "role_counts": source_counts,
            "role_ids": {
                role: sorted(element_ids)
                for role, element_ids in source_roles.items()
            },
            "missing_ids": source_missing,
            "extra_ids": source_extra,
            "unclassified_ids": source_unclassified,
            "ambiguous_ids": source_ambiguous,
            "missing_pairs": source_missing_pairs,
            "missing_relations": source_missing_relations,
            "missing_pair_metadata": missing_pair_metadata,
            "missing_relation_metadata": missing_relation_metadata,
            "malformed_relations": malformed_relations,
            "unexpected_relations": unexpected_relations,
            "unexpected_pairs": unexpected_pairs,
            "component_count": source_components,
            "pass": source_pass,
        },
        "backend": {
            "role_counts": backend_counts,
            "role_ids": {
                role: sorted(element_ids)
                for role, element_ids in backend_roles.items()
            },
            "missing_ids": backend_missing,
            "extra_ids": backend_extra,
            "unclassified_ids": backend_unclassified,
            "ambiguous_ids": backend_ambiguous,
            "binding_missing": binding_missing,
            "readback_missing": readback_missing,
            "shape_missing": shape_missing,
            "binding_extra": binding_extra,
            "readback_extra": readback_extra,
            "missing_pairs": backend_missing_pairs,
            "missing_relations": backend_missing_relations,
            "component_count": backend_components,
            "pass": backend_pass,
        },
        "pass": source_pass and backend_pass,
    }


def evaluate_visual_contracts(run: common.Run) -> dict[str, Any]:
    """Evaluate the current candidate against its frozen object inventory."""

    regions = read_json(run.regions_path)
    inventory = regions.get("reference_inventory")
    required = isinstance(inventory, dict) and inventory.get("required") is True
    if not isinstance(inventory, dict):
        report = {
            "schema_version": "1.0.0",
            "status": "not-declared",
            "required": False,
            "pass": True,
            "blockers": [],
            "object_count": 0,
        }
        write_json(run.qa_dir / "visual-contracts-report.json", report)
        return report

    findings: list[dict[str, Any]] = []
    if required and inventory.get("status") != "frozen":
        findings.append(_finding("V1", "reference_inventory", "required inventory is not frozen"))
    source = _build_source_index(run.redraw_svg)
    backend = _build_backend_index(run)
    objects = inventory.get("objects", [])
    if not isinstance(objects, list):
        objects = []
        findings.append(_finding("V2", "reference_inventory", "objects must be a list"))

    rows: list[dict[str, Any]] = []
    object_source_boxes: dict[str, Box] = {}
    object_backend_boxes: dict[str, Box] = {}
    object_specs: dict[str, dict[str, Any]] = {}
    typography_rows: list[dict[str, Any]] = []
    topology_rows: list[dict[str, Any]] = []
    source_transform_rows: list[dict[str, Any]] = []
    for object_spec in objects:
        if not isinstance(object_spec, dict) or not object_spec.get("id"):
            findings.append(_finding("V2", "reference_inventory", "invalid object entry"))
            continue
        object_id = str(object_spec["id"])
        object_specs[object_id] = object_spec
        kind = str(object_spec.get("kind", "shape"))
        expected = _box(object_spec.get("bbox"))
        element_ids = [str(item) for item in object_spec.get("element_ids", [])]
        row: dict[str, Any] = {
            "id": object_id,
            "kind": kind,
            "element_ids": element_ids,
            "expected_bbox": None if expected is None else expected.as_dict(),
        }
        if expected is None or not element_ids:
            findings.append(_finding("V3", object_id, "object requires bbox and element_ids"))
            row["pass"] = False
            rows.append(row)
            continue
        source_box = _object_box(source, source, element_ids)
        backend_box = _object_box(backend, source, element_ids)
        if kind in {"text", "formula"}:
            backend_visual_box = _translated_text_backend_box(object_spec, source, backend)
            if backend_visual_box is not None:
                backend_box = backend_visual_box
        tolerance = float(
            object_spec.get("visual", {}).get(
                "bbox_tolerance_px",
                object_spec.get("typography", {}).get(
                    "bbox_tolerance_px",
                    DEFAULT_BBOX_TOLERANCE_PX,
                ),
            )
        )
        aspect_tolerance = float(
            object_spec.get("visual", {}).get(
                "aspect_ratio_tolerance",
                DEFAULT_ASPECT_TOLERANCE,
            )
        )
        delegated_arrow_geometry = (
            kind == "arrow"
            and isinstance(object_spec.get("contract_refs"), dict)
            and isinstance(
                object_spec["contract_refs"].get("arrow_visual"), dict
            )
        )
        source_pass = source_box is not None and (
            delegated_arrow_geometry or _bbox_pass(source_box, expected, tolerance)
        )
        backend_pass = backend_box is not None and (
            delegated_arrow_geometry or _bbox_pass(backend_box, expected, tolerance)
        )
        aspect_source = math.inf if source_box is None else _aspect_error(source_box, expected)
        aspect_backend = math.inf if backend_box is None else _aspect_error(backend_box, expected)
        if (
            kind in {"icon", "plot", "shape", "arrow", "brace"}
            and not delegated_arrow_geometry
        ):
            source_pass = source_pass and aspect_source <= aspect_tolerance
            backend_pass = backend_pass and aspect_backend <= aspect_tolerance
        source_geometry_pass = source_pass
        source_transform = _audit_source_transform_anisotropy(
            object_spec,
            source,
            element_ids,
        )
        if source_transform is not None:
            source_transform = {"id": object_id, **source_transform}
            source_transform_rows.append(source_transform)
            if not source_transform["pass"]:
                findings.append(
                    _finding(
                        "V38",
                        object_id,
                        "SVG semantic microasset uses undeclared anisotropic transform",
                        {
                            "tolerance": source_transform["tolerance"],
                            "offenders": source_transform["offenders"],
                        },
                    )
                )
            source_pass = source_pass and source_transform["pass"]
        if source_box is None:
            findings.append(_finding("V4", object_id, "SVG object has no measurable bound elements"))
        elif not source_geometry_pass:
            findings.append(
                _finding(
                    "V5",
                    object_id,
                    "SVG object scale or position differs from reference inventory",
                    {"delta": _box_delta(source_box, expected)},
                )
            )
            object_source_boxes[object_id] = source_box
        else:
            object_source_boxes[object_id] = source_box
        if backend_box is None:
            findings.append(_finding("V6", object_id, "PowerPoint object binding/bounds are missing"))
        elif not backend_pass:
            findings.append(
                _finding(
                    "V7",
                    object_id,
                    "saved PowerPoint object scale or position differs from reference inventory",
                    {"delta": _box_delta(backend_box, expected)},
                )
            )
            object_backend_boxes[object_id] = backend_box
        else:
            object_backend_boxes[object_id] = backend_box
        row.update(
            {
                "source_bbox": None if source_box is None else source_box.as_dict(),
                "backend_bbox": None if backend_box is None else backend_box.as_dict(),
                "bbox_tolerance_px": tolerance,
                "aspect_tolerance": aspect_tolerance,
                "geometry_authority": (
                    "arrow_visual_physical_gate"
                    if delegated_arrow_geometry
                    else "reference_inventory_bbox"
                ),
                "source_aspect_error": None if not math.isfinite(aspect_source) else round(aspect_source, 6),
                "backend_aspect_error": None if not math.isfinite(aspect_backend) else round(aspect_backend, 6),
                "source_geometry_pass": source_geometry_pass,
                "source_transform_pass": (
                    None if source_transform is None else source_transform["pass"]
                ),
                "source_pass": source_pass,
                "backend_pass": backend_pass,
                "pass": source_pass and backend_pass,
            }
        )
        topology_row = _audit_topology(object_spec, source, backend, findings)
        if topology_row is not None:
            topology_rows.append(topology_row)
            row["topology_pass"] = topology_row["pass"]
            row["pass"] = row["pass"] and topology_row["pass"]
        rows.append(row)
        if kind in {"text", "formula"}:
            typography_rows.append(_audit_typography(object_spec, source, backend, findings))

    visual_contracts = regions.get("visual_contracts", {})
    clearances = inventory.get("clearance_contracts", [])
    if isinstance(visual_contracts, dict) and visual_contracts.get("clearances") is not None:
        clearances = visual_contracts.get("clearances")
    clearance_rows: list[dict[str, Any]] = []
    for contract in clearances if isinstance(clearances, list) else []:
        if not isinstance(contract, dict):
            continue
        contract_id = str(contract.get("id", f"clearance-{len(clearance_rows) + 1}"))
        left_id, right_id = str(contract.get("a", "")), str(contract.get("b", ""))
        axis = str(contract.get("axis", ""))
        minimum = float(contract.get("min_px", 0.0))
        source_a, source_b = object_source_boxes.get(left_id), object_source_boxes.get(right_id)
        backend_a, backend_b = object_backend_boxes.get(left_id), object_backend_boxes.get(right_id)
        source_gap = None if source_a is None or source_b is None else _gap(source_a, source_b, axis)
        backend_gap = None if backend_a is None or backend_b is None else _gap(backend_a, backend_b, axis)
        passed = (
            source_gap is not None
            and backend_gap is not None
            and source_gap >= minimum
            and backend_gap >= minimum
        )
        if not passed:
            findings.append(
                _finding(
                    "V20",
                    contract_id,
                    "declared visual clearance collapsed or became an overlap",
                    {"source_gap_px": source_gap, "backend_gap_px": backend_gap, "min_px": minimum},
                )
            )
        clearance_rows.append(
            {
                "id": contract_id,
                "a": left_id,
                "b": right_id,
                "axis": axis,
                "min_px": minimum,
                "source_gap_px": source_gap,
                "backend_gap_px": backend_gap,
                "pass": passed,
            }
        )

    collision_rows: list[dict[str, Any]] = []
    collision_tolerance = float(
        visual_contracts.get("collision_tolerance_px", DEFAULT_COLLISION_TOLERANCE_PX)
        if isinstance(visual_contracts, dict)
        else DEFAULT_COLLISION_TOLERANCE_PX
    )
    object_ids = sorted(object_specs)
    for index, left_id in enumerate(object_ids):
        left_spec = object_specs[left_id]
        left_kind = str(left_spec.get("kind", "shape"))
        if left_kind not in _COLLISION_KINDS:
            continue
        for right_id in object_ids[index + 1 :]:
            right_spec = object_specs[right_id]
            right_kind = str(right_spec.get("kind", "shape"))
            if right_kind not in _COLLISION_KINDS:
                continue
            allowed = set(str(item) for item in left_spec.get("allow_overlap_with", []))
            allowed.update(str(item) for item in right_spec.get("allow_overlap_with", []))
            if right_id in allowed or left_id in allowed:
                continue
            left_expected, right_expected = _box(left_spec.get("bbox")), _box(right_spec.get("bbox"))
            if left_expected is None or right_expected is None:
                continue
            expected_overlap = _intersection(left_expected, right_expected)
            if expected_overlap[0] > collision_tolerance and expected_overlap[1] > collision_tolerance:
                continue
            source_a, source_b = object_source_boxes.get(left_id), object_source_boxes.get(right_id)
            backend_a, backend_b = object_backend_boxes.get(left_id), object_backend_boxes.get(right_id)
            source_overlap = (0.0, 0.0) if source_a is None or source_b is None else _intersection(source_a, source_b)
            backend_overlap = (0.0, 0.0) if backend_a is None or backend_b is None else _intersection(backend_a, backend_b)
            passed = not (
                source_overlap[0] > collision_tolerance
                and source_overlap[1] > collision_tolerance
            ) and not (
                backend_overlap[0] > collision_tolerance
                and backend_overlap[1] > collision_tolerance
            )
            if not passed:
                target = f"{left_id}|{right_id}"
                findings.append(
                    _finding(
                        "V21",
                        target,
                        "objects disjoint in the reference now intersect",
                        {
                            "source_overlap_px": list(source_overlap),
                            "backend_overlap_px": list(backend_overlap),
                        },
                    )
                )
            collision_rows.append(
                {
                    "a": left_id,
                    "b": right_id,
                    "source_overlap_px": list(source_overlap),
                    "backend_overlap_px": list(backend_overlap),
                    "pass": passed,
                }
            )

    blockers = [
        f"visual-contract:{finding['code']}:{finding['target']}"
        for finding in findings
    ]
    report = {
        "schema_version": "1.0.0",
        "status": "PASS" if not blockers else "FAIL",
        "required": required,
        "reference_sha256": inventory.get("reference_sha256"),
        "object_count": len(rows),
        "objects": rows,
        "typography": typography_rows,
        "topology_contract_count": len(topology_rows),
        "topology_contracts_sha256": topology_contracts_sha256(inventory),
        "topology": topology_rows,
        "source_transform_anisotropy": source_transform_rows,
        "clearances": clearance_rows,
        "collision_checks": collision_rows,
        "findings": findings,
        "blockers": list(dict.fromkeys(blockers)),
        "pass": not blockers,
    }
    write_json(run.qa_dir / "visual-contracts-report.json", report)
    return report


def strict_blockers(report: dict[str, Any]) -> list[str]:
    return list(report.get("blockers", []))
