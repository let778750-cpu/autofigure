"""Hash-bound PowerPoint arrow compilation and OOXML readback evidence."""

from __future__ import annotations

import json
import math
import re
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation

from tools.core import common
from tools.arrows.arrow_composition import find_reciprocal_arrow_overlaps
from tools.arrows.arrow_spec import (
    compiler_strategy,
    path_endpoints,
    path_from_segments,
    path_tangents,
    semantic_dash_from_ooxml,
    spec_sha256,
    validate_scene_arrow_specs,
    validate_arrow_spec,
)
from tools.core.contracts import read_json, utc_now, write_json
from tools.core.svggeom import Matrix, parse_path_d, parse_transform

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
SVG_NS = "http://www.w3.org/2000/svg"
NS = {"p": P_NS, "a": A_NS}
EMU_PER_PX = 9525.0
SOURCE_PATH_TOLERANCE_PX = 1e-4
OOXML_TO_SPEC_HEAD = {
    "none": "none",
    "arrow": "open",
    "triangle": "triangle",
    "stealth": "stealth",
    "diamond": "diamond",
    "oval": "oval",
}


def _walk_shapes(shapes) -> Iterable[Any]:
    for shape in shapes:
        yield shape
        children = getattr(shape, "shapes", None)
        if children is not None:
            yield from _walk_shapes(children)


def _ooxml_shape_identities(pptx_path: Path) -> set[tuple[int, str]]:
    """Return exact shape identities, including native math AlternateContent.

    ``python-pptx`` intentionally ignores Office Math shapes stored in
    ``mc:AlternateContent``.  PowerPoint preserves those objects during a real
    save/reopen, so treating the library's partial enumeration as the full
    slide inventory creates false missing-binding failures.  ``p:cNvPr`` is
    the package-level identity source PowerPoint itself preserves; duplicate
    Choice/Fallback records collapse to the same ``(id, name)`` pair.
    """

    identities: set[tuple[int, str]] = set()
    with zipfile.ZipFile(pptx_path) as package:
        slide_entries = sorted(
            entry
            for entry in package.namelist()
            if re.fullmatch(r"ppt/slides/slide\d+\.xml", entry)
        )
        for entry in slide_entries:
            root = ET.fromstring(package.read(entry))
            for node in root.iter():
                if node.tag.rsplit("}", 1)[-1] != "cNvPr":
                    continue
                raw_id = node.get("id")
                name = node.get("name")
                if raw_id is None or not name:
                    continue
                identities.add((int(raw_id), name))
    return identities


def refresh_bindings(run: common.Run, *, host_saved_reopened: bool) -> dict[str, Any]:
    """Rebind the current root artifact without manufacturing host evidence."""

    presentation = Presentation(run.pptx_path)
    shapes = list(_walk_shapes(presentation.slides[0].shapes))
    library_identities = {(int(shape.shape_id), shape.name) for shape in shapes}
    identities = _ooxml_shape_identities(run.pptx_path) or library_identities
    bindings = read_json(run.bindings_path)
    physical_bindings = bindings.get("bindings", [])
    composite_bindings = bindings.get("logical_group_bindings", [])
    bound_identities = {
        (int(binding.get("shape_id", -1)), binding.get("shape_name"))
        for binding in physical_bindings
    }
    identities_by_element: dict[str, list[tuple[int, str]]] = {}
    for binding in physical_bindings:
        # A name-only or id-only match can accidentally bind a different shape
        # after PowerPoint deletes or reorders objects.  Exact identity is the
        # minimum safe save/reopen contract.
        identity = (
            int(binding.get("shape_id", -1)), binding.get("shape_name")
        )
        binding["readback_found"] = identity in identities
        element_id = binding.get("element_id")
        if isinstance(element_id, str) and element_id:
            identities_by_element.setdefault(element_id, []).append(identity)

    for binding in composite_bindings:
        backend_ids = binding.get("backend_object_ids")
        backend_names = binding.get("backend_object_names")
        well_formed = (
            isinstance(backend_ids, list)
            and isinstance(backend_names, list)
            and bool(backend_ids)
            and len(backend_ids) == len(backend_names)
            and all(
                isinstance(shape_id, int)
                and not isinstance(shape_id, bool)
                and shape_id > 0
                for shape_id in backend_ids
            )
            and all(isinstance(name, str) and name for name in backend_names)
        )
        composite_identities = (
            list(zip(backend_ids, backend_names, strict=True))
            if well_formed
            else []
        )
        member_ids = binding.get("member_element_ids")
        member_identities = (
            [
                identity
                for member_id in member_ids
                for identity in identities_by_element.get(member_id, [])
            ]
            if isinstance(member_ids, list)
            else []
        )
        binding["readback_found"] = (
            bool(composite_identities)
            and composite_identities == member_identities
            and len(set(composite_identities)) == len(composite_identities)
            and all(identity in identities for identity in composite_identities)
        )
    complete = (
        bool(physical_bindings)
        and all(item.get("readback_found") is True for item in physical_bindings)
        and all(item.get("readback_found") is True for item in composite_bindings)
        and len(bound_identities) == len(physical_bindings)
        and bound_identities == identities
    )
    artifact_hash = common.sha256_file(run.pptx_path)
    bindings.update(
        {
            "updated_at": utc_now(),
            "artifact_sha256": artifact_hash,
            "package_reopened": True,
            "saved_reopened": bool(host_saved_reopened),
            "bindings_complete": complete,
        }
    )
    write_json(run.bindings_path, bindings)
    return {
        "artifact_sha256": artifact_hash,
        "bindings_complete": complete,
        "saved_reopened": bool(host_saved_reopened),
        "object_count": len(identities),
        "binding_count": len(physical_bindings),
        "logical_group_binding_count": len(composite_bindings),
        "unbound_object_count": len(identities - bound_identities),
        "missing_bound_object_count": len(bound_identities - identities),
    }


def _identity_and_properties(node: ET.Element) -> tuple[ET.Element | None, ET.Element | None, str]:
    local = node.tag.rsplit("}", 1)[-1]
    if local == "sp":
        return node.find("./p:nvSpPr/p:cNvPr", NS), node.find("./p:spPr", NS), "shape"
    if local == "cxnSp":
        return node.find("./p:nvCxnSpPr/p:cNvPr", NS), node.find("./p:spPr", NS), "connector"
    if local == "grpSp":
        return node.find("./p:nvGrpSpPr/p:cNvPr", NS), node.find("./p:grpSpPr", NS), "group"
    return None, None, local


def _head(line: ET.Element | None, tag: str) -> dict[str, Any]:
    node = None if line is None else line.find(f"a:{tag}", NS)
    if node is None:
        return {"type": "none", "width": None, "length": None}
    raw_type = node.get("type", "none")
    return {
        "type": OOXML_TO_SPEC_HEAD.get(raw_type, raw_type),
        "ooxml_type": raw_type,
        "width": node.get("w", "med"),
        "length": node.get("len", "med"),
    }


def _line_readback(properties: ET.Element | None) -> dict[str, Any] | None:
    line = None if properties is None else properties.find("a:ln", NS)
    if line is None:
        return None
    color = line.find("./a:solidFill/a:srgbClr", NS)
    dash = line.find("a:prstDash", NS)
    width = line.get("w")
    raw_cap = line.get("cap", "flat")
    line_cap = {"rnd": "round", "sq": "square", "flat": "butt"}.get(raw_cap, raw_cap)
    if line.find("a:round", NS) is not None:
        line_join = "round"
    elif line.find("a:bevel", NS) is not None:
        line_join = "bevel"
    else:
        line_join = "miter"
    raw_dash = None if dash is None else dash.get("val")
    return {
        "color": None if color is None else f"#{color.get('val', '').upper()}",
        "width_px": None if width is None else float(width) / EMU_PER_PX,
        "dash": semantic_dash_from_ooxml(raw_dash, line_cap),
        "ooxml_dash": raw_dash or "solid",
        "line_cap": line_cap,
        "line_join": line_join,
        "start_head": _head(line, "headEnd"),
        "end_head": _head(line, "tailEnd"),
    }


def _fill_readback(properties: ET.Element | None) -> dict[str, Any]:
    color = (
        None
        if properties is None
        else properties.find("./a:solidFill/a:srgbClr", NS)
    )
    return {
        "type": "solid" if color is not None else "none",
        "color": None if color is None else f"#{color.get('val', '').upper()}",
    }


def _autoshape_readback(properties: ET.Element | None) -> dict[str, Any] | None:
    preset = None if properties is None else properties.find("a:prstGeom", NS)
    if preset is None:
        return None
    adjustments: list[float] = []
    names: list[str] = []
    guide_list = preset.find("a:avLst", NS)
    if guide_list is not None:
        for guide in guide_list.findall("a:gd", NS):
            formula = guide.get("fmla", "")
            match = re.fullmatch(r"val\s+([-+]?\d+(?:\.\d+)?)", formula)
            if match is None:
                continue
            names.append(guide.get("name", ""))
            adjustments.append(float(match.group(1)) / 100000.0)
    return {
        "subtype": preset.get("prst"),
        "adjustment_names": names,
        "adjustments": adjustments,
        "bbox": _transform(properties),
    }


def _transform(properties: ET.Element | None) -> dict[str, float | bool] | None:
    node = None if properties is None else properties.find("a:xfrm", NS)
    if node is None:
        return None
    offset = node.find("a:off", NS)
    extent = node.find("a:ext", NS)
    if offset is None or extent is None:
        return None
    return {
        "x": float(offset.get("x", "0")) / EMU_PER_PX,
        "y": float(offset.get("y", "0")) / EMU_PER_PX,
        "width": float(extent.get("cx", "0")) / EMU_PER_PX,
        "height": float(extent.get("cy", "0")) / EMU_PER_PX,
        "flip_h": node.get("flipH") in {"1", "true"},
        "flip_v": node.get("flipV") in {"1", "true"},
    }


def _map_local(
    point: ET.Element,
    *,
    transform: dict[str, float | bool],
    path_width: float,
    path_height: float,
) -> tuple[float, float]:
    raw_x = float(point.get("x", "0"))
    raw_y = float(point.get("y", "0"))
    x = 0.0 if path_width == 0 else raw_x / path_width * float(transform["width"])
    y = 0.0 if path_height == 0 else raw_y / path_height * float(transform["height"])
    if transform["flip_h"]:
        x = float(transform["width"]) - x
    if transform["flip_v"]:
        y = float(transform["height"]) - y
    return float(transform["x"]) + x, float(transform["y"]) + y


def _shape_segments(properties: ET.Element | None) -> list[tuple] | None:
    transform = _transform(properties)
    if properties is None or transform is None:
        return None
    preset = properties.find("a:prstGeom", NS)
    if preset is not None and preset.get("prst") == "line":
        start_x = float(transform["x"])
        start_y = float(transform["y"])
        end_x = start_x + float(transform["width"])
        end_y = start_y + float(transform["height"])
        if transform["flip_h"]:
            start_x, end_x = end_x, start_x
        if transform["flip_v"]:
            start_y, end_y = end_y, start_y
        return [("M", start_x, start_y), ("L", end_x, end_y)]

    path_node = properties.find("./a:custGeom/a:pathLst/a:path", NS)
    if path_node is None:
        return None
    path_width = float(path_node.get("w", "0"))
    path_height = float(path_node.get("h", "0"))
    segments: list[tuple] = []
    for command in path_node:
        local = command.tag.rsplit("}", 1)[-1]
        points = command.findall("a:pt", NS)
        if local == "moveTo" and len(points) == 1:
            segments.append(("M", *_map_local(points[0], transform=transform, path_width=path_width, path_height=path_height)))
        elif local == "lnTo" and len(points) == 1:
            segments.append(("L", *_map_local(points[0], transform=transform, path_width=path_width, path_height=path_height)))
        elif local == "cubicBezTo" and len(points) == 3:
            mapped = [
                _map_local(point, transform=transform, path_width=path_width, path_height=path_height)
                for point in points
            ]
            segments.append(("C", *mapped[0], *mapped[1], *mapped[2]))
        elif local == "close":
            segments.append(("Z",))
    return segments or None


def _shape_path(properties: ET.Element | None) -> dict[str, Any] | None:
    segments = _shape_segments(properties)
    if segments is None:
        return None
    try:
        return path_from_segments(segments)
    except ValueError:
        return None


def _parse_inventory(path: Path) -> list[dict[str, Any]]:
    with zipfile.ZipFile(path) as package:
        root = ET.fromstring(package.read("ppt/slides/slide1.xml"))
    tree = root.find("./p:cSld/p:spTree", NS)
    if tree is None:
        return []
    records: list[dict[str, Any]] = []

    def visit(node: ET.Element, parent_group_id: int | None = None) -> None:
        identity, properties, kind = _identity_and_properties(node)
        if identity is None:
            return
        shape_id = int(identity.get("id", "0"))
        record = {
            "shape_id": shape_id,
            "shape_name": identity.get("name", ""),
            "description": identity.get("descr", ""),
            "ooxml_kind": kind,
            "parent_group_id": parent_group_id,
            "line": _line_readback(properties),
            "fill": _fill_readback(properties),
            "autoshape": _autoshape_readback(properties),
            "path": _shape_path(properties),
            "segments": _shape_segments(properties),
            "connections": None,
        }
        if kind == "connector":
            connection_properties = node.find("./p:nvCxnSpPr/p:cNvCxnSpPr", NS)
            start_connection = (
                None
                if connection_properties is None
                else connection_properties.find("a:stCxn", NS)
            )
            end_connection = (
                None
                if connection_properties is None
                else connection_properties.find("a:endCxn", NS)
            )
            record["connections"] = {
                "source_shape_id": (
                    None if start_connection is None else int(start_connection.get("id", "0"))
                ),
                "source_site": (
                    None if start_connection is None else int(start_connection.get("idx", "0"))
                ),
                "target_shape_id": (
                    None if end_connection is None else int(end_connection.get("id", "0"))
                ),
                "target_site": (
                    None if end_connection is None else int(end_connection.get("idx", "0"))
                ),
            }
        records.append(record)
        if kind == "group":
            for child in node:
                if child.tag in {f"{{{P_NS}}}sp", f"{{{P_NS}}}cxnSp", f"{{{P_NS}}}grpSp"}:
                    visit(child, shape_id)

    for child in tree:
        if child.tag in {f"{{{P_NS}}}sp", f"{{{P_NS}}}cxnSp", f"{{{P_NS}}}grpSp"}:
            visit(child)
    return records


def read_pptx_inventory(path: Path) -> list[dict[str, Any]]:
    """Return deterministic slide-one OOXML geometry for QA consumers."""

    return _parse_inventory(path)


def _angle(vector: tuple[float, float]) -> float | None:
    if math.hypot(*vector) <= 1e-9:
        return None
    return math.degrees(math.atan2(vector[1], vector[0]))


def _angle_difference(left: float | None, right: float | None) -> float | None:
    if left is None or right is None:
        return None
    return abs((left - right + 180.0) % 360.0 - 180.0)


def _path_control_points(
    path: dict[str, Any] | None,
) -> tuple[tuple[str, int], list[tuple[float, float]]]:
    """Return a fail-closed path signature and every geometry-defining point."""

    if not isinstance(path, dict):
        raise ValueError("path unavailable")
    kind = path.get("kind")
    if kind in {"straight", "polyline"}:
        points = path.get("points")
        if not isinstance(points, list):
            raise ValueError("path points unavailable")
        return (kind, len(points)), [
            (float(point["x"]), float(point["y"])) for point in points
        ]
    if kind == "cubic":
        start = path.get("start")
        segments = path.get("segments")
        if not isinstance(start, dict) or not isinstance(segments, list):
            raise ValueError("cubic path controls unavailable")
        points = [(float(start["x"]), float(start["y"]))]
        for segment in segments:
            for key in ("control1", "control2", "end"):
                point = segment[key]
                points.append((float(point["x"]), float(point["y"])))
        return (kind, len(segments)), points
    raise ValueError(f"unsupported path kind: {kind}")


def _source_path_comparison(
    expected: dict[str, Any] | None,
    actual: dict[str, Any] | None,
) -> dict[str, Any]:
    """Compare source geometry without the looser PowerPoint readback tolerance."""

    try:
        expected_signature, expected_points = _path_control_points(expected)
        actual_signature, actual_points = _path_control_points(actual)
    except (KeyError, TypeError, ValueError) as exc:
        return {
            "status": "INCONCLUSIVE",
            "reason": "source-path-unavailable",
            "detail": str(exc),
        }
    if expected_signature != actual_signature or len(expected_points) != len(actual_points):
        return {
            "status": "FAIL",
            "reason": "source-path-command-signature",
            "expected_signature": list(expected_signature),
            "actual_signature": list(actual_signature),
        }
    point_errors = [
        math.dist(expected_point, actual_point)
        for expected_point, actual_point in zip(expected_points, actual_points, strict=True)
    ]
    maximum = max(point_errors, default=math.inf)
    return {
        "status": "PASS" if maximum <= SOURCE_PATH_TOLERANCE_PX else "FAIL",
        "reason": None if maximum <= SOURCE_PATH_TOLERANCE_PX else "source-path-coordinate-drift",
        "path_signature": list(expected_signature),
        "max_point_error_px": round(maximum, 6),
        "tolerance_px": SOURCE_PATH_TOLERANCE_PX,
    }


def _transform_source_segments(segments: list[tuple], matrix: Matrix) -> list[tuple]:
    transformed: list[tuple] = []
    for segment in segments:
        if segment[0] in {"M", "L"}:
            transformed.append((segment[0], *matrix.apply(segment[1], segment[2])))
        elif segment[0] == "C":
            control1 = matrix.apply(segment[1], segment[2])
            control2 = matrix.apply(segment[3], segment[4])
            end = matrix.apply(segment[5], segment[6])
            transformed.append(("C", *control1, *control2, *end))
        elif segment[0] == "Z":
            transformed.append(("Z",))
        else:
            raise ValueError(f"unsupported source segment: {segment[0]}")
    return transformed


def _source_segments(
    tag: str,
    geometry: Any,
    matrix: Matrix,
) -> list[tuple]:
    get = geometry.get
    if tag == "line":
        segments = [
            ("M", float(get("x1") or 0), float(get("y1") or 0)),
            ("L", float(get("x2") or 0), float(get("y2") or 0)),
        ]
    elif tag in {"polyline", "polygon"}:
        values = [
            float(value)
            for value in re.split(r"[\s,]+", (get("points") or "").strip())
            if value
        ]
        if len(values) < 4 or len(values) % 2:
            raise ValueError(f"invalid {tag} points")
        segments = [("M", values[0], values[1])]
        segments.extend(
            ("L", values[index], values[index + 1])
            for index in range(2, len(values), 2)
        )
        if tag == "polygon":
            segments.append(("Z",))
    elif tag == "path":
        path_data = get("d")
        if not path_data:
            raise ValueError("path d unavailable")
        segments = parse_path_d(path_data)
    else:
        raise ValueError(f"unsupported arrow source tag: {tag}")
    return _transform_source_segments(segments, matrix)


def _svg_geometry_inventory(svg_path: Path) -> dict[str, list[dict[str, Any]]]:
    root = ET.parse(svg_path).getroot()
    inventory: dict[str, list[dict[str, Any]]] = {}
    element_counts: dict[str, int] = {}
    skipped_tags = {"defs", "marker", "linearGradient", "radialGradient"}
    renderable_tags = {
        "rect",
        "circle",
        "ellipse",
        "line",
        "polyline",
        "polygon",
        "path",
        "text",
        "image",
    }

    def walk(element: ET.Element, parent_matrix: Matrix) -> None:
        tag = element.tag.rsplit("}", 1)[-1]
        if tag in skipped_tags:
            return
        matrix = parent_matrix.multiply(parse_transform(element.get("transform")))
        if tag in renderable_tags:
            # Mirror ConvertContext.begin_element exactly.  SVG-seeded inputs may
            # legitimately omit explicit ids; conversion still assigns stable
            # per-tag ids (svg-line-0001, ...), so source QA must resolve the
            # same identity instead of treating every id-less arrow as missing.
            element_counts[tag] = element_counts.get(tag, 0) + 1
            element_id = element.get("id") or f"svg-{tag}-{element_counts[tag]:04d}"
            inventory.setdefault(element_id, []).append(
                {
                    "element": element,
                    "matrix": matrix,
                    "tag": tag,
                }
            )
        for child in element:
            walk(child, matrix)

    walk(root, Matrix())
    return inventory


def _source_silhouette_comparison(
    expected: dict[str, Any] | None,
    actual: list[tuple] | None,
) -> dict[str, Any]:
    expected_segments = None if expected is None else expected.get("segments")
    if not isinstance(expected_segments, list) or not isinstance(actual, list):
        return {"status": "INCONCLUSIVE", "reason": "source-silhouette-unavailable"}
    if len(expected_segments) != len(actual):
        return {
            "status": "FAIL",
            "reason": "source-silhouette-command-signature",
        }
    maximum = 0.0
    for expected_segment, actual_segment in zip(expected_segments, actual, strict=True):
        if (
            not expected_segment
            or expected_segment[0] != actual_segment[0]
            or len(expected_segment) != len(actual_segment)
        ):
            return {
                "status": "FAIL",
                "reason": "source-silhouette-command-signature",
            }
        for expected_value, actual_value in zip(
            expected_segment[1:], actual_segment[1:], strict=True
        ):
            maximum = max(maximum, abs(float(expected_value) - float(actual_value)))
    return {
        "status": "PASS" if maximum <= SOURCE_PATH_TOLERANCE_PX else "FAIL",
        "reason": (
            None
            if maximum <= SOURCE_PATH_TOLERANCE_PX
            else "source-silhouette-coordinate-drift"
        ),
        "max_coordinate_error_px": round(maximum, 6),
        "tolerance_px": SOURCE_PATH_TOLERANCE_PX,
    }


def _source_geometry_comparison(
    scene_element: dict[str, Any],
    spec: dict[str, Any],
    svg_inventory: dict[str, list[dict[str, Any]]],
) -> dict[str, Any]:
    element_id = scene_element.get("id")
    matches = svg_inventory.get(element_id, [])
    if len(matches) != 1:
        result = {
            "status": "INCONCLUSIVE",
            "reason": "source-svg-identity-missing" if not matches else "source-svg-identity-duplicate",
            "match_count": len(matches),
        }
        return {"status": "FAIL", "scene": result, "svg": result}

    match = matches[0]
    tag = match["tag"]
    matrix = match["matrix"]
    svg_element = match["element"]
    scene_geometry = scene_element.get("geometry")
    if not isinstance(scene_geometry, dict):
        scene_result = {"status": "INCONCLUSIVE", "reason": "scene-geometry-unavailable"}
    elif scene_element.get("svg_tag") != tag:
        scene_result = {"status": "FAIL", "reason": "scene-svg-tag-drift"}
    elif (scene_geometry.get("transform") or "") != (svg_element.get("transform") or ""):
        scene_result = {"status": "FAIL", "reason": "scene-transform-drift"}
    else:
        try:
            scene_segments = _source_segments(tag, scene_geometry, matrix)
            scene_path = (
                None
                if spec.get("representation") == "block_arrow"
                else path_from_segments(scene_segments)
            )
            scene_result = (
                _source_silhouette_comparison(spec.get("silhouette_path"), scene_segments)
                if spec.get("representation") == "block_arrow"
                else _source_path_comparison(spec.get("path"), scene_path)
            )
        except (IndexError, KeyError, TypeError, ValueError) as exc:
            scene_result = {
                "status": "INCONCLUSIVE",
                "reason": "scene-geometry-invalid",
                "detail": str(exc),
            }

    try:
        svg_segments = _source_segments(tag, svg_element, matrix)
        svg_path = (
            None
            if spec.get("representation") == "block_arrow"
            else path_from_segments(svg_segments)
        )
        svg_result = (
            _source_silhouette_comparison(spec.get("silhouette_path"), svg_segments)
            if spec.get("representation") == "block_arrow"
            else _source_path_comparison(spec.get("path"), svg_path)
        )
    except (IndexError, KeyError, TypeError, ValueError) as exc:
        svg_result = {
            "status": "INCONCLUSIVE",
            "reason": "source-svg-geometry-invalid",
            "detail": str(exc),
        }
    return {
        "status": (
            "PASS"
            if scene_result.get("status") == "PASS" and svg_result.get("status") == "PASS"
            else "FAIL"
        ),
        "scene": scene_result,
        "svg": svg_result,
    }


def _geometry_comparison(
    expected: dict[str, Any] | None,
    actual: dict[str, Any] | None,
    diagonal: float,
) -> dict[str, Any]:
    if actual is None:
        return {"status": "INCONCLUSIVE", "reason": "path-readback-unavailable"}
    try:
        expected_signature, expected_points = _path_control_points(expected)
        actual_signature, actual_points = _path_control_points(actual)
        expected_start, expected_end = path_endpoints(expected)
        actual_start, actual_end = path_endpoints(actual)
        expected_tangents = path_tangents(expected)
        actual_tangents = path_tangents(actual)
    except (KeyError, TypeError, ValueError) as exc:
        return {
            "status": "INCONCLUSIVE",
            "reason": "path-readback-invalid",
            "detail": str(exc),
        }
    if expected_signature != actual_signature or len(expected_points) != len(actual_points):
        return {
            "status": "FAIL",
            "reason": "path-command-signature",
            "expected_signature": list(expected_signature),
            "actual_signature": list(actual_signature),
        }
    start_error = math.dist(expected_start, actual_start)
    end_error = math.dist(expected_end, actual_end)
    point_errors = [
        math.dist(expected_point, actual_point)
        for expected_point, actual_point in zip(expected_points, actual_points, strict=True)
    ]
    full_path_error = max(point_errors, default=math.inf)
    expected_angles = (_angle(expected_tangents[0]), _angle(expected_tangents[1]))
    actual_angles = (_angle(actual_tangents[0]), _angle(actual_tangents[1]))
    start_angle_error = _angle_difference(expected_angles[0], actual_angles[0])
    end_angle_error = _angle_difference(expected_angles[1], actual_angles[1])
    tangents_available = all(
        (expected_angle is None) == (actual_angle is None)
        for expected_angle, actual_angle in zip(expected_angles, actual_angles, strict=True)
    )
    endpoint_tolerance = diagonal * 0.0025
    full_path_tolerance = diagonal * 0.0035
    angle_tolerance = 3.0
    passed = (
        max(start_error, end_error) <= endpoint_tolerance
        and full_path_error <= full_path_tolerance
        and tangents_available
        and start_angle_error is not None
        and start_angle_error <= angle_tolerance
        and end_angle_error is not None
        and end_angle_error <= angle_tolerance
    )
    return {
        "status": "PASS" if passed else "FAIL",
        "expected_start": [round(value, 6) for value in expected_start],
        "expected_end": [round(value, 6) for value in expected_end],
        "actual_start": [round(value, 6) for value in actual_start],
        "actual_end": [round(value, 6) for value in actual_end],
        "start_error_px": round(start_error, 4),
        "end_error_px": round(end_error, 4),
        "full_path_max_point_error_px": round(full_path_error, 4),
        "path_signature": list(expected_signature),
        "start_angle_error_deg": None if start_angle_error is None else round(start_angle_error, 4),
        "end_angle_error_deg": None if end_angle_error is None else round(end_angle_error, 4),
        "tangents_available": tangents_available,
        "tolerances": {
            "endpoint_px": round(endpoint_tolerance, 4),
            "full_path_px": round(full_path_tolerance, 4),
            "angle_deg": angle_tolerance,
        },
    }


def _heads_match(expected: dict[str, Any], actual: dict[str, Any]) -> bool:
    return all(expected.get(key) == actual.get(key) for key in ("type", "width", "length"))


def _silhouette_comparison(
    expected: dict[str, Any] | None,
    actual: list[tuple] | None,
) -> dict[str, Any]:
    expected_segments = [] if expected is None else expected.get("segments", [])
    if not isinstance(actual, list):
        return {"status": "INCONCLUSIVE", "reason": "silhouette-readback-unavailable"}
    if len(expected_segments) != len(actual):
        return {
            "status": "FAIL",
            "reason": "silhouette-command-count",
            "expected": len(expected_segments),
            "actual": len(actual),
        }
    maximum = 0.0
    for expected_segment, actual_segment in zip(expected_segments, actual, strict=True):
        if expected_segment[0] != actual_segment[0] or len(expected_segment) != len(actual_segment):
            return {"status": "FAIL", "reason": "silhouette-command-signature"}
        for expected_value, actual_value in zip(expected_segment[1:], actual_segment[1:], strict=True):
            maximum = max(maximum, abs(float(expected_value) - float(actual_value)))
    return {
        "status": "PASS" if maximum <= 0.25 else "FAIL",
        "max_coordinate_error_px": round(maximum, 4),
        "tolerance_px": 0.25,
    }


def _autoshape_comparison(
    expected: dict[str, Any] | None,
    actual: dict[str, Any] | None,
) -> dict[str, Any]:
    if not isinstance(expected, dict) or not isinstance(actual, dict):
        return {"status": "INCONCLUSIVE", "reason": "autoshape-readback-unavailable"}
    expected_bbox = expected.get("bbox")
    actual_bbox_payload = actual.get("bbox")
    actual_bbox = (
        None
        if not isinstance(actual_bbox_payload, dict)
        else [
            actual_bbox_payload.get("x"),
            actual_bbox_payload.get("y"),
            actual_bbox_payload.get("width"),
            actual_bbox_payload.get("height"),
        ]
    )
    if not (
        isinstance(expected_bbox, list)
        and len(expected_bbox) == 4
        and isinstance(actual_bbox, list)
        and all(isinstance(value, (int, float)) for value in actual_bbox)
    ):
        return {"status": "INCONCLUSIVE", "reason": "autoshape-bbox-unavailable"}
    expected_adjustments = expected.get("adjustments")
    actual_adjustments = actual.get("adjustments")
    if not (
        isinstance(expected_adjustments, list)
        and isinstance(actual_adjustments, list)
        and len(expected_adjustments) == len(actual_adjustments)
    ):
        return {
            "status": "FAIL",
            "reason": "autoshape-adjustment-count",
            "expected": expected_adjustments,
            "actual": actual_adjustments,
        }
    bbox_errors = [
        abs(float(expected_value) - float(actual_value))
        for expected_value, actual_value in zip(expected_bbox, actual_bbox, strict=True)
    ]
    adjustment_errors = [
        abs(float(expected_value) - float(actual_value))
        for expected_value, actual_value in zip(
            expected_adjustments, actual_adjustments, strict=True
        )
    ]
    subtype_pass = expected.get("subtype") == actual.get("subtype")
    maximum_bbox_error = max(bbox_errors, default=math.inf)
    maximum_adjustment_error = max(adjustment_errors, default=math.inf)
    passed = (
        subtype_pass
        and maximum_bbox_error <= 0.25
        and maximum_adjustment_error <= 1e-5
    )
    return {
        "status": "PASS" if passed else "FAIL",
        "expected_subtype": expected.get("subtype"),
        "actual_subtype": actual.get("subtype"),
        "subtype_pass": subtype_pass,
        "expected_bbox": expected_bbox,
        "actual_bbox": actual_bbox,
        "max_bbox_error_px": round(maximum_bbox_error, 6),
        "bbox_tolerance_px": 0.25,
        "expected_adjustments": expected_adjustments,
        "actual_adjustments": actual_adjustments,
        "max_adjustment_error": round(maximum_adjustment_error, 8),
        "adjustment_tolerance": 1e-5,
    }


def _autoshape_centerline(actual: dict[str, Any] | None) -> dict[str, Any] | None:
    """Recover the semantic axis represented by a vetted block AutoShape."""

    if not isinstance(actual, dict) or actual.get("subtype") != "leftRightArrow":
        return None
    bbox = actual.get("bbox")
    if not isinstance(bbox, dict):
        return None
    try:
        x = float(bbox["x"])
        y = float(bbox["y"])
        width = float(bbox["width"])
        height = float(bbox["height"])
    except (KeyError, TypeError, ValueError):
        return None
    if not all(math.isfinite(value) for value in (x, y, width, height)):
        return None
    if width <= 0 or height <= 0:
        return None
    center_y = y + height / 2.0
    return {
        "kind": "straight",
        "coordinate_space": "canvas",
        "points": [
            {"x": x, "y": center_y},
            {"x": x + width, "y": center_y},
        ],
    }


def write_arrow_reports(run: common.Run) -> tuple[dict[str, Any], dict[str, Any]]:
    scene = read_json(run.scene_path)
    bindings = read_json(run.bindings_path)
    meta = run.load_meta()
    svg_inventory = _svg_geometry_inventory(run.redraw_svg)
    artifact_hash = common.sha256_file(run.pptx_path)
    bindings_by_element: dict[str, list[dict[str, Any]]] = {}
    for binding in bindings.get("bindings", []):
        bindings_by_element.setdefault(binding.get("element_id", ""), []).append(binding)
    arrows = [item for item in scene.get("elements", []) if isinstance(item.get("arrow_spec"), dict)]
    missing_specs = [
        item.get("id", "[missing-id]")
        for item in scene.get("elements", [])
        if item.get("kind") == "edge" and not isinstance(item.get("arrow_spec"), dict)
    ]

    source_composition_findings = find_reciprocal_arrow_overlaps(arrows)
    compile_records: list[dict[str, Any]] = []
    compile_blockers: list[str] = [f"arrow:A11_SPEC_MISSING:{item}" for item in missing_specs]
    compile_blockers.extend(
        f"arrow:A12_SPEC_INVALID:scene:{error}"
        for error in validate_scene_arrow_specs(scene)
    )
    compile_blockers.extend(
        "arrow:A21_SOURCE_RECIPROCAL_OVERLAP:" + ":".join(finding["arrow_ids"])
        for finding in source_composition_findings
    )
    for element in arrows:
        element_id = element["id"]
        spec = element["arrow_spec"]
        errors = validate_arrow_spec(
            spec,
            expected_input_route=meta.get("input_route"),
            expected_reference_sha256=meta.get("source_sha256"),
        )
        source_geometry = _source_geometry_comparison(element, spec, svg_inventory)
        strategy = compiler_strategy(spec)
        rows = bindings_by_element.get(element_id, [])
        visible_rows = [
            row
            for row in rows
            if row.get("object_kind")
            in {
                "connector",
                "line",
                "freeform-arrow",
                "freeform",
                "block-arrow-autoshape",
                "arrowhead-fallback",
            }
        ]
        fallback = strategy in {"legacy-grouped-fallback", "unsupported"} or any(
            row.get("object_kind") in {"arrowhead-fallback", "arrow-group"} for row in rows
        )
        if errors:
            compile_blockers.extend(
                f"arrow:A12_SPEC_INVALID:{element_id}:{error}" for error in errors
            )
        for source_name in ("scene", "svg"):
            if source_geometry[source_name].get("status") != "PASS":
                compile_blockers.append(
                    f"arrow:A10_SOURCE_PATH_DRIFT:{element_id}:{source_name}"
                )
        if fallback:
            compile_blockers.append(f"arrow:A13_COMPILE_FIDELITY_LOSS:{element_id}:{strategy}")
        if len(visible_rows) != 1:
            compile_blockers.append(
                f"arrow:A14_MULTI_VISIBLE_OBJECT:{element_id}:{len(visible_rows)}"
            )
        compile_records.append(
            {
                "element_id": element_id,
                "arrow_spec_sha256": spec_sha256(spec),
                "strategy": strategy,
                "shape_ids": [row.get("shape_id") for row in rows],
                "shape_names": [row.get("shape_name") for row in rows],
                "object_kinds": [row.get("object_kind") for row in rows],
                "visible_object_count": len(visible_rows),
                "fallback": fallback,
                "fidelity_loss": "grouped-visible-arrow-parts" if fallback else None,
                "source_geometry": source_geometry,
                "validation_errors": errors,
                "status": (
                    "PASS"
                    if not errors
                    and source_geometry["status"] == "PASS"
                    and not fallback
                    and len(visible_rows) == 1
                    else "FAIL"
                ),
            }
        )
    compile_report = {
        "schema_version": "1.0.0",
        "kind": "arrow_compile_report",
        "created_at": bindings.get("updated_at") or utc_now(),
        "case": run.root.name,
        "scene_sha256": common.sha256_file(run.scene_path),
        "artifact_sha256": artifact_hash,
        "arrow_count": len(arrows),
        "records": compile_records,
        "blockers": list(dict.fromkeys(compile_blockers)),
        "pass": not compile_blockers,
    }
    write_json(run.qa_dir / "arrow-compile-report.json", compile_report)

    inventory = _parse_inventory(run.pptx_path)
    by_identity = {
        (item["shape_id"], item["shape_name"]): item
        for item in inventory
        if item["shape_id"] and item["shape_name"]
    }
    canvas = scene.get("canvas", {})
    diagonal = math.hypot(float(canvas.get("width", 0)), float(canvas.get("height", 0)))
    readback_records: list[dict[str, Any]] = []
    readback_blockers: list[str] = []
    readback_composition_specs: list[tuple[str, dict[str, Any]]] = []
    readback_composition_inconclusive: list[str] = []
    for element in arrows:
        element_id = element["id"]
        spec = element["arrow_spec"]
        rows = bindings_by_element.get(element_id, [])
        primary = next(
            (
                row
                for row in rows
                if row.get("object_kind")
                in {
                    "connector",
                    "line",
                    "freeform-arrow",
                    "freeform",
                    "block-arrow-autoshape",
                }
            ),
            None,
        )
        item = None
        if primary is not None:
            item = by_identity.get(
                (primary.get("shape_id"), primary.get("shape_name"))
            )
        if item is None:
            readback_blockers.append(f"arrow:A15_READBACK_MISSING:{element_id}")
            if spec.get("representation") == "line_arrow":
                readback_composition_inconclusive.append(element_id)
            readback_records.append(
                {
                    "element_id": element_id,
                    "arrow_spec_sha256": spec_sha256(spec),
                    "status": "INCONCLUSIVE",
                    "reason": "primary-object-readback-missing",
                }
            )
            continue

        strategy = compiler_strategy(spec)
        line = item.get("line") or {
            "start_head": {"type": "none", "width": None, "length": None},
            "end_head": {"type": "none", "width": None, "length": None},
        }
        if spec.get("representation") == "line_arrow":
            actual_path = item.get("path")
            if isinstance(actual_path, dict):
                readback_composition_specs.append(
                    (
                        element_id,
                        {
                            **spec,
                            "path": actual_path,
                            "start_head": line["start_head"],
                            "end_head": line["end_head"],
                        },
                    )
                )
            else:
                readback_composition_inconclusive.append(element_id)
        embedded_head_strategy = strategy in {
            "single-closed-freeform",
            "native-block-autoshape",
        }
        if embedded_head_strategy:
            start_pass = end_pass = True
        else:
            start_pass = _heads_match(spec["start_head"], line["start_head"])
            end_pass = _heads_match(spec["end_head"], line["end_head"])
        body = spec.get("body", {})
        if strategy == "native-block-autoshape":
            autoshape_actual = item.get("autoshape") or {}
            fill_actual = item.get("fill") or {}
            actual_adjustments = autoshape_actual.get("adjustments") or []
            actual_bbox = autoshape_actual.get("bbox") or {}
            actual_width = (
                float(actual_bbox.get("height")) * float(actual_adjustments[0])
                if actual_bbox.get("height") is not None and actual_adjustments
                else None
            )
            color_pass = body.get("color") in {None, fill_actual.get("color")}
            width_pass = (
                actual_width is not None
                and abs(float(body.get("width_px", 0)) - actual_width) <= 0.02
            )
            dash_pass = cap_pass = join_pass = True
        elif strategy == "single-closed-freeform" and spec.get("representation") == "block_arrow":
            fill_actual = item.get("fill") or {}
            actual_width = None
            color_pass = body.get("color") in {None, fill_actual.get("color")}
            # The single closed silhouette is the body-width evidence for an
            # arbitrary freeform; there is no independent centerline width.
            width_pass = dash_pass = cap_pass = join_pass = True
        else:
            color_pass = body.get("color") in {None, line.get("color")}
            actual_width = line.get("width_px")
            width_pass = actual_width is not None and abs(float(body.get("width_px", 0)) - actual_width) <= 0.02
            dash_pass = body.get("dash", "solid") == line.get("dash", "solid")
            cap_pass = body.get("line_cap", "butt") == line.get("line_cap", "butt")
            join_pass = body.get("line_join", "miter") == line.get("line_join", "miter")
        try:
            embedded = json.loads(item.get("description") or "{}")
        except json.JSONDecodeError:
            embedded = {}
        tag_pass = embedded.get("arrow_spec_sha256") == spec_sha256(spec)
        semantic_centerline = None
        autoshape_geometry = None
        if strategy == "single-closed-freeform":
            geometry = _silhouette_comparison(
                spec.get("silhouette_path"), item.get("segments")
            )
            if spec.get("representation") == "block_arrow":
                semantic_centerline = _geometry_comparison(
                    spec.get("path"),
                    embedded.get("arrow_semantic_centerline"),
                    diagonal,
                )
                semantic_centerline["evidence"] = (
                    "saved-pptx-cNvPr-description-bound-to-read-back-silhouette"
                )
        elif strategy == "native-block-autoshape":
            autoshape_geometry = _autoshape_comparison(
                spec.get("autoshape"), item.get("autoshape")
            )
            geometry = _geometry_comparison(
                spec.get("path"),
                _autoshape_centerline(item.get("autoshape")),
                diagonal,
            )
            semantic_centerline = {
                **geometry,
                "evidence": "saved-pptx-native-block-autoshape-geometry",
            }
        elif strategy in {
            "native-connector-line-end",
            "native-line-line-end",
            "native-freeform-line-end",
        }:
            geometry = _geometry_comparison(spec.get("path"), item.get("path"), diagonal)
        else:
            geometry = {
                "status": "INCONCLUSIVE",
                "reason": f"unsupported-compiler-strategy:{strategy}",
            }
        topology = spec.get("topology", {})
        topology_pass = True
        topology_actual = item.get("connections")
        if topology.get("mode") == "attached":
            source_shape_ids = {
                row.get("shape_id")
                for row in bindings_by_element.get(topology.get("source_id"), [])
            }
            target_shape_ids = {
                row.get("shape_id")
                for row in bindings_by_element.get(topology.get("target_id"), [])
            }
            topology_pass = bool(
                topology_actual
                and topology_actual.get("source_shape_id") in source_shape_ids
                and topology_actual.get("target_shape_id") in target_shape_ids
                and topology_actual.get("source_site") == topology.get("source_site")
                and topology_actual.get("target_site") == topology.get("target_site")
            )
        semantic_centerline_pass = (
            semantic_centerline is None
            or semantic_centerline.get("status") == "PASS"
        )
        geometry_pass = geometry.get("status") == "PASS" and semantic_centerline_pass and (
            autoshape_geometry is None
            or autoshape_geometry.get("status") == "PASS"
        )
        status = "PASS" if all(
            (
                start_pass,
                end_pass,
                color_pass,
                width_pass,
                dash_pass,
                cap_pass,
                join_pass,
                tag_pass,
                topology_pass,
                geometry_pass,
            )
        ) else "FAIL"
        if not (start_pass and end_pass):
            readback_blockers.append(f"arrow:A16_READBACK_HEAD:{element_id}")
        if not (color_pass and width_pass and dash_pass and cap_pass and join_pass and tag_pass):
            readback_blockers.append(f"arrow:A17_READBACK_BODY:{element_id}")
        if not geometry_pass:
            readback_blockers.append(f"arrow:A18_READBACK_PATH:{element_id}")
        if not topology_pass:
            readback_blockers.append(f"arrow:A19_READBACK_TOPOLOGY:{element_id}")
        def semantic_head_readback(side: str) -> dict[str, Any]:
            expected = spec[side]
            physical = line[side]
            if not embedded_head_strategy or expected.get("type") == "none":
                return physical
            return {
                **expected,
                "representation": "embedded-silhouette",
                "ooxml_line_end": physical,
            }

        readback_records.append(
            {
                "element_id": element_id,
                "arrow_spec_sha256": spec_sha256(spec),
                "shape_id": item["shape_id"],
                "shape_name": item["shape_name"],
                "ooxml_kind": item["ooxml_kind"],
                "start_head": {
                    "expected": spec["start_head"],
                    "actual": semantic_head_readback("start_head"),
                    "pass": start_pass,
                },
                "end_head": {
                    "expected": spec["end_head"],
                    "actual": semantic_head_readback("end_head"),
                    "pass": end_pass,
                },
                "body": {
                    "expected": body,
                    "actual": {
                        "color": (
                            (item.get("fill") or {}).get("color")
                            if spec.get("representation") == "block_arrow"
                            else line.get("color")
                        ),
                        "width_px": actual_width,
                    "dash": None if strategy == "native-block-autoshape" else line.get("dash"),
                    "line_cap": None if strategy == "native-block-autoshape" else line.get("line_cap"),
                    "line_join": None if strategy == "native-block-autoshape" else line.get("line_join"),
                    },
                    "color_pass": color_pass,
                    "width_pass": width_pass,
                    "dash_pass": dash_pass,
                    "line_cap_pass": cap_pass,
                    "line_join_pass": join_pass,
                },
                "path_geometry": geometry,
                "semantic_centerline": semantic_centerline,
                "autoshape": {
                    "expected": spec.get("autoshape"),
                    "actual": item.get("autoshape"),
                    "geometry": autoshape_geometry,
                    "pass": autoshape_geometry.get("status") == "PASS",
                } if strategy == "native-block-autoshape" else None,
                "embedded_arrow_spec_sha256": embedded.get("arrow_spec_sha256"),
                "embedded_tag_pass": tag_pass,
                "topology": {
                    "expected": topology,
                    "actual": topology_actual,
                    "pass": topology_pass,
                },
                "status": status,
            }
        )
    readback_composition_findings = find_reciprocal_arrow_overlaps(
        readback_composition_specs
    )
    readback_blockers.extend(
        "arrow:A22_READBACK_RECIPROCAL_OVERLAP:" + ":".join(finding["arrow_ids"])
        for finding in readback_composition_findings
    )
    readback_blockers.extend(
        f"arrow:A23_RECIPROCAL_READBACK_INCONCLUSIVE:{element_id}"
        for element_id in readback_composition_inconclusive
    )
    readback_report = {
        "schema_version": "1.0.0",
        "kind": "powerpoint_arrow_readback",
        "created_at": bindings.get("updated_at") or utc_now(),
        "case": run.root.name,
        "scene_sha256": common.sha256_file(run.scene_path),
        "artifact_sha256": artifact_hash,
        "bindings_artifact_sha256": bindings.get("artifact_sha256"),
        "saved_reopened": bindings.get("saved_reopened") is True,
        "inventory_count": len(inventory),
        "arrow_count": len(arrows),
        "records": readback_records,
        "blockers": list(dict.fromkeys(readback_blockers)),
        "pass": not readback_blockers,
    }
    write_json(run.qa_dir / "powerpoint-arrow-readback.json", readback_report)
    composition_blockers = [
        *(
            "arrow:A21_SOURCE_RECIPROCAL_OVERLAP:" + ":".join(finding["arrow_ids"])
            for finding in source_composition_findings
        ),
        *(
            "arrow:A22_READBACK_RECIPROCAL_OVERLAP:" + ":".join(finding["arrow_ids"])
            for finding in readback_composition_findings
        ),
        *(
            f"arrow:A23_RECIPROCAL_READBACK_INCONCLUSIVE:{element_id}"
            for element_id in readback_composition_inconclusive
        ),
    ]
    composition_report = {
        "schema_version": "1.0.0",
        "kind": "arrow_composition_audit",
        "created_at": bindings.get("updated_at") or utc_now(),
        "case": run.root.name,
        "scene_sha256": common.sha256_file(run.scene_path),
        "artifact_sha256": artifact_hash,
        "source_findings": source_composition_findings,
        "readback_findings": readback_composition_findings,
        "readback_inconclusive": readback_composition_inconclusive,
        "blockers": list(dict.fromkeys(composition_blockers)),
        "pass": not composition_blockers,
    }
    write_json(run.qa_dir / "arrow-composition-audit.json", composition_report)
    return compile_report, readback_report


def strict_blockers(run: common.Run) -> list[str]:
    expected_hash = common.sha256_file(run.pptx_path)
    blockers: list[str] = []
    for filename, prefix in (
        ("arrow-compile-report.json", "arrow-compile"),
        ("powerpoint-arrow-readback.json", "arrow-readback"),
        ("arrow-composition-audit.json", "arrow-composition"),
    ):
        path = run.qa_dir / filename
        if not path.is_file():
            blockers.append(f"{prefix}:missing")
            continue
        report = read_json(path)
        if report.get("artifact_sha256") != expected_hash:
            blockers.append("arrow:A20_ARTIFACT_IDENTITY")
        if report.get("scene_sha256") != common.sha256_file(run.scene_path):
            blockers.append("arrow:A20_ARTIFACT_IDENTITY")
        blockers.extend(report.get("blockers", []))
    bindings = read_json(run.bindings_path)
    if bindings.get("artifact_sha256") != expected_hash:
        blockers.append("arrow:A20_ARTIFACT_IDENTITY")
    readback_path = run.powerpoint_arrow_readback_path
    if readback_path.is_file():
        readback = read_json(readback_path)
        if readback.get("bindings_artifact_sha256") != expected_hash:
            blockers.append("arrow:A20_ARTIFACT_IDENTITY")
    return list(dict.fromkeys(blockers))
