"""Region-level visual gates; global metrics never override a critical failure."""

from __future__ import annotations

import hashlib
import json
import math
from typing import Any

import numpy as np
from PIL import Image

from tools import common
from tools.contracts import SCHEMA_VERSION, read_json, write_json

CRITICAL_SSIM_FLOOR = 0.85
CRITICAL_EDGE_IOU_FLOOR = 0.75
SUBJECT_MASK_IOU_FLOOR = 0.65
OBSTACLE_MASK_IOU_FLOOR = 0.65
THIN_STROKE_MASK_IOU_FLOOR = 0.55
OBSTACLE_COLOR_COSINE_FLOOR = 0.90


def _gate_sha256(value: Any) -> str:
    canonical = json.dumps(
        value, sort_keys=True, separators=(",", ":"), ensure_ascii=False
    )
    return hashlib.sha256(canonical.encode("utf-8")).hexdigest()


def _crop(array: np.ndarray, bbox: list[int]) -> np.ndarray:
    if (
        not isinstance(bbox, list)
        or len(bbox) != 4
        or any(not isinstance(value, int) or isinstance(value, bool) for value in bbox)
    ):
        raise ValueError(f"region bbox must contain four integers: {bbox}")
    x, y, width, height = bbox
    if x < 0 or y < 0 or width <= 0 or height <= 0:
        raise ValueError(f"invalid region bbox: {bbox}")
    if x + width > array.shape[1] or y + height > array.shape[0]:
        raise ValueError(f"region bbox outside canvas: {bbox}")
    return array[y : y + height, x : x + width]


def _ssim(reference: np.ndarray, render: np.ndarray) -> float:
    from skimage.metrics import structural_similarity

    smallest = min(reference.shape[0], reference.shape[1])
    win_size = min(7, smallest if smallest % 2 else smallest - 1)
    if win_size < 3:
        return 1.0 if np.array_equal(reference, render) else 0.0
    return float(
        structural_similarity(reference, render, channel_axis=2, data_range=255, win_size=win_size)
    )


def _edge_iou(reference: np.ndarray, render: np.ndarray) -> float:
    from skimage.feature import canny
    from scipy.ndimage import binary_dilation

    ref_gray = np.dot(reference[..., :3], np.array([0.299, 0.587, 0.114])) / 255.0
    ren_gray = np.dot(render[..., :3], np.array([0.299, 0.587, 0.114])) / 255.0
    ref_edges = canny(ref_gray, sigma=1.0)
    ren_edges = canny(ren_gray, sigma=1.0)
    ref_count = int(ref_edges.sum())
    ren_count = int(ren_edges.sum())
    if ref_count == 0 and ren_count == 0:
        return 1.0
    if ref_count == 0 or ren_count == 0:
        return 0.0
    # PowerPoint and PNG references can rasterize the same vector boundary on
    # adjacent pixels. A two-pixel symmetric tolerance is 0.13% of the case-01
    # diagonal: it measures topology rather than anti-aliasing phase while still
    # penalizing missing or extra edges.
    ref_match = int(
        np.logical_and(ref_edges, binary_dilation(ren_edges, iterations=2)).sum()
    )
    ren_match = int(
        np.logical_and(ren_edges, binary_dilation(ref_edges, iterations=2)).sum()
    )
    # A fuzzy match still needs a one-to-one count ceiling. Averaging the two
    # directional match counts makes a doubled nearby contour score 1.0 because
    # every extra render edge can reuse the same reference edge. Limiting the
    # intersection to the smaller directional count preserves subpixel-shift
    # tolerance while penalizing missing or duplicated boundaries.
    intersection = float(min(ref_match, ren_match))
    union = ref_count + ren_count - intersection
    return float(intersection / union) if union > 0 else 1.0


def _fuzzy_mask_iou(reference: np.ndarray, render: np.ndarray) -> float:
    from scipy.ndimage import binary_dilation

    reference_count = int(reference.sum())
    render_count = int(render.sum())
    if reference_count == 0 and render_count == 0:
        return 1.0
    if reference_count == 0 or render_count == 0:
        return 0.0
    reference_match = int(
        np.logical_and(reference, binary_dilation(render)).sum()
    )
    render_match = int(
        np.logical_and(render, binary_dilation(reference)).sum()
    )
    intersection = float(min(reference_match, render_match))
    union = reference_count + render_count - intersection
    return float(intersection / union) if union > 0 else 1.0


def _mean_rgb(array: np.ndarray, point: list[int], radius: int) -> np.ndarray:
    x, y = (int(v) for v in point)
    y0, y1 = max(0, y - radius), min(array.shape[0], y + radius + 1)
    x0, x1 = max(0, x - radius), min(array.shape[1], x + radius + 1)
    if x1 <= x0 or y1 <= y0:
        raise ValueError(f"color probe outside canvas: {point}")
    return array[y0:y1, x0:x1, :3].mean(axis=(0, 1))


def _delta_e(reference_rgb: np.ndarray, render_rgb: np.ndarray) -> float:
    from skimage.color import deltaE_ciede2000, rgb2lab

    pair = np.array([[reference_rgb, render_rgb]], dtype=float) / 255.0
    lab = rgb2lab(pair)
    return float(deltaE_ciede2000(lab[:, :1], lab[:, 1:])[0, 0])


def _evaluate_probes(
    reference: np.ndarray,
    render: np.ndarray,
    probes: list[dict[str, Any]],
) -> tuple[list[dict[str, Any]], bool]:
    results: list[dict[str, Any]] = []
    passed = True
    for probe in probes:
        radius = int(probe.get("radius", 2))
        ref_rgb = _mean_rgb(reference, probe["point"], radius)
        ren_rgb = _mean_rgb(render, probe["point"], radius)
        delta = _delta_e(ref_rgb, ren_rgb)
        maximum = float(probe.get("max_delta_e", 5.0))
        item_pass = delta <= maximum
        passed = passed and item_pass
        results.append(
            {
                "id": probe.get("id", f"probe-{len(results) + 1}"),
                "point": probe["point"],
                "delta_e00": round(delta, 4),
                "max_delta_e": maximum,
                "pass": item_pass,
            }
        )
    return results, passed


def _foreground_mask(array: np.ndarray, background_rgb: list[int], tolerance: float) -> np.ndarray:
    background = np.asarray(background_rgb, dtype=float)
    if background.shape != (3,) or not np.all(np.isfinite(background)):
        raise ValueError("ink_contract background_rgb must contain three finite values")
    delta = np.max(np.abs(array[..., :3].astype(float) - background), axis=2)
    return delta > tolerance


def _core_foreground_rgb(
    array: np.ndarray,
    mask: np.ndarray,
    background_rgb: list[int],
) -> np.ndarray:
    """Return a robust ink-core color, excluding antialiased edge pixels.

    Rasterizers blend vector ink toward the background at object boundaries. A
    mean over every foreground pixel therefore compares antialiasing coverage as
    if it were a semantic color change. The highest-contrast foreground
    quartile is a renderer-stable estimate of the actual vector ink color.
    """

    pixels = array[..., :3][mask].astype(float)
    if not len(pixels):
        raise ValueError("core foreground color requires nonempty ink")
    background = np.asarray(background_rgb, dtype=float)
    contrast = np.linalg.norm(background - pixels, axis=1)
    if not np.all(np.isfinite(contrast)):
        raise ValueError("core foreground contrast must be finite")
    # The darkest decile is the closest raster evidence to the authored vector
    # color.  A broader quartile is unstable for one- and two-pixel dots because
    # it can be dominated by antialias coverage instead of actual ink.
    cutoff = float(np.quantile(contrast, 0.90))
    core = pixels[contrast >= cutoff]
    if not len(core):
        maximum_index = int(np.argmax(contrast))
        core = pixels[maximum_index : maximum_index + 1]
    return np.median(core, axis=0)


def _target_color_mask(
    array: np.ndarray,
    candidate_mask: np.ndarray,
    background_rgb: list[int],
    reference_core_rgb: np.ndarray,
    minimum_cosine: float,
) -> np.ndarray:
    """Select ink whose background-to-pixel direction matches the target color.

    Antialiasing changes coverage but keeps pixels on the line between the ink
    and background colors. Comparing that direction rather than raw RGB keeps
    those pixels while rejecting a same-size object rendered in a different
    semantic color.
    """

    background = np.asarray(background_rgb, dtype=float)
    target = background - np.asarray(reference_core_rgb, dtype=float)
    target_norm = float(np.linalg.norm(target))
    if target_norm <= 0:
        raise ValueError("target color must differ from the background")
    vectors = background - array[..., :3].astype(float)
    magnitudes = np.linalg.norm(vectors, axis=2)
    cosine = np.zeros(magnitudes.shape, dtype=float)
    valid = candidate_mask & (magnitudes > 0)
    cosine[valid] = (
        np.sum(vectors[valid] * target, axis=1)
        / (magnitudes[valid] * target_norm)
    )
    return valid & (cosine >= minimum_cosine)


def _mask_geometry(mask: np.ndarray) -> dict[str, Any] | None:
    rows, columns = np.nonzero(mask)
    if not len(rows):
        return None
    left = int(columns.min())
    top = int(rows.min())
    right = int(columns.max())
    bottom = int(rows.max())
    return {
        "bbox": [left, top, right - left + 1, bottom - top + 1],
        "edges": [left, top, right, bottom],
        "center": [(left + right) / 2.0, (top + bottom) / 2.0],
        "area_px": int(mask.sum()),
    }


def _obstacle_mask_floor(
    reference_geometry: dict[str, Any], configured: Any = None
) -> tuple[float, float, bool]:
    """Choose a renderer-aware mask floor without weakening solid objects."""

    reference_bbox = reference_geometry["bbox"]
    reference_fill_fraction = int(reference_geometry["area_px"]) / max(
        int(reference_bbox[2]) * int(reference_bbox[3]), 1
    )
    thin_open_reference = (
        max(int(reference_bbox[2]), int(reference_bbox[3])) >= 8
        and reference_fill_fraction <= 0.20
    )
    renderer_floor = (
        THIN_STROKE_MASK_IOU_FLOOR
        if thin_open_reference
        else OBSTACLE_MASK_IOU_FLOOR
    )
    minimum = renderer_floor if configured is None else max(float(configured), renderer_floor)
    if not 0 <= minimum <= 1:
        raise ValueError("obstacle mask IoU tolerance must be in [0, 1]")
    return minimum, reference_fill_fraction, thin_open_reference


def _evaluate_ink_contract(
    reference: np.ndarray,
    render: np.ndarray,
    contract: dict[str, Any] | None,
) -> tuple[dict[str, Any] | None, bool]:
    if contract is None:
        return None, True
    if not isinstance(contract, dict):
        return {"pass": False, "error": "ink_contract must be an object"}, False
    try:
        background_rgb = contract.get("background_rgb", [255, 255, 255])
        background_tolerance = float(contract.get("background_tolerance", 24.0))
        bbox_tolerance = float(contract.get("bbox_tolerance_px", 1.0))
        center_tolerance = float(contract.get("center_tolerance_px", 1.0))
        area_relative_tolerance = float(contract.get("area_relative_tolerance", 0.10))
    except (TypeError, ValueError) as exc:
        return {"pass": False, "error": f"invalid ink_contract value: {exc}"}, False
    if background_tolerance < 0 or bbox_tolerance < 0 or center_tolerance < 0:
        return {"pass": False, "error": "ink_contract tolerances must be nonnegative"}, False
    if not 0 <= area_relative_tolerance <= 1:
        return {
            "pass": False,
            "error": "ink_contract area_relative_tolerance must be within [0, 1]",
        }, False
    try:
        reference_geometry = _mask_geometry(
            _foreground_mask(reference, background_rgb, background_tolerance)
        )
        render_geometry = _mask_geometry(
            _foreground_mask(render, background_rgb, background_tolerance)
        )
    except (TypeError, ValueError) as exc:
        return {"pass": False, "error": str(exc)}, False
    if reference_geometry is None or render_geometry is None:
        return {
            "pass": False,
            "error": "ink_contract requires foreground pixels in reference and render",
            "reference": reference_geometry,
            "render": render_geometry,
        }, False
    edge_errors = [
        abs(float(left) - float(right))
        for left, right in zip(
            reference_geometry["edges"], render_geometry["edges"], strict=True
        )
    ]
    center_error = float(
        np.linalg.norm(
            np.asarray(reference_geometry["center"], dtype=float)
            - np.asarray(render_geometry["center"], dtype=float)
        )
    )
    reference_area = int(reference_geometry["area_px"])
    render_area = int(render_geometry["area_px"])
    area_relative_error = abs(render_area - reference_area) / max(reference_area, 1)
    passed = (
        max(edge_errors) <= bbox_tolerance
        and center_error <= center_tolerance
        and area_relative_error <= area_relative_tolerance
    )
    return {
        "background_rgb": background_rgb,
        "background_tolerance": background_tolerance,
        "reference": reference_geometry,
        "render": render_geometry,
        "bbox_edge_errors_px": [round(value, 4) for value in edge_errors],
        "bbox_tolerance_px": bbox_tolerance,
        "center_error_px": round(center_error, 4),
        "center_tolerance_px": center_tolerance,
        "area_relative_error": round(area_relative_error, 4),
        "area_relative_tolerance": area_relative_tolerance,
        "pass": passed,
    }, passed


def _dominant_color_mask(
    array: np.ndarray, channel_name: str, minimum_dominance: float
) -> np.ndarray:
    channel_index = {"red": 0, "green": 1, "blue": 2}.get(channel_name)
    if channel_index is None:
        raise ValueError("color clearance subject_dominant_channel must be red, green, or blue")
    rgb = array[..., :3].astype(float)
    subject = rgb[..., channel_index]
    others = np.delete(rgb, channel_index, axis=2).max(axis=2)
    return subject - others >= minimum_dominance


def _components_touching_bbox(
    candidate: np.ndarray, bbox_mask: np.ndarray
) -> np.ndarray:
    from scipy.ndimage import label

    components, _ = label(candidate, structure=np.ones((3, 3), dtype=np.uint8))
    touching_labels = np.unique(components[bbox_mask])
    touching_labels = touching_labels[touching_labels != 0]
    return np.isin(components, touching_labels)


def _outside_component_delta(
    reference: np.ndarray, render: np.ndarray
) -> dict[str, Any]:
    """Match legal same-color objects by component, not antialiasing pixels."""

    from scipy.ndimage import binary_dilation, label

    structure = np.ones((3, 3), dtype=np.uint8)
    reference_labels, reference_count = label(reference, structure=structure)
    render_labels, render_count = label(render, structure=structure)

    def records(labels: np.ndarray, count: int) -> list[dict[str, Any]]:
        output: list[dict[str, Any]] = []
        for component_id in range(1, count + 1):
            mask = labels == component_id
            geometry = _mask_geometry(mask)
            if geometry is not None:
                output.append(
                    {
                        "component_id": component_id,
                        "mask": mask,
                        "geometry": geometry,
                    }
                )
        return output

    reference_records = records(reference_labels, reference_count)
    render_records = records(render_labels, render_count)
    used_reference: set[int] = set()
    matches: list[dict[str, Any]] = []
    unmatched: list[dict[str, Any]] = []
    for render_record in render_records:
        render_geometry = render_record["geometry"]
        candidates: list[tuple[float, int, dict[str, Any]]] = []
        for reference_index, reference_record in enumerate(reference_records):
            if reference_index in used_reference:
                continue
            reference_geometry = reference_record["geometry"]
            edge_errors = [
                abs(float(left) - float(right))
                for left, right in zip(
                    reference_geometry["edges"],
                    render_geometry["edges"],
                    strict=True,
                )
            ]
            center_error = math.dist(
                reference_geometry["center"], render_geometry["center"]
            )
            area_error = abs(
                int(reference_geometry["area_px"])
                - int(render_geometry["area_px"])
            ) / max(int(reference_geometry["area_px"]), 1)
            reference_mask = reference_record["mask"]
            render_mask = render_record["mask"]
            ref_left, ref_top, ref_width, ref_height = reference_geometry["bbox"]
            ren_left, ren_top, ren_width, ren_height = render_geometry["bbox"]
            reference_tight = reference_mask[
                ref_top : ref_top + ref_height,
                ref_left : ref_left + ref_width,
            ]
            render_tight = render_mask[
                ren_top : ren_top + ren_height,
                ren_left : ren_left + ren_width,
            ]
            canvas_height = max(reference_tight.shape[0], render_tight.shape[0])
            canvas_width = max(reference_tight.shape[1], render_tight.shape[1])
            reference_canvas = np.zeros((canvas_height, canvas_width), dtype=bool)
            render_canvas = np.zeros((canvas_height, canvas_width), dtype=bool)
            reference_canvas[
                : reference_tight.shape[0], : reference_tight.shape[1]
            ] = reference_tight
            render_canvas[: render_tight.shape[0], : render_tight.shape[1]] = (
                render_tight
            )
            reference_match = int(
                np.logical_and(
                    reference_canvas, binary_dilation(render_canvas)
                ).sum()
            )
            render_match = int(
                np.logical_and(
                    render_canvas, binary_dilation(reference_canvas)
                ).sum()
            )
            intersection = min(reference_match, render_match)
            union = (
                int(reference_canvas.sum())
                + int(render_canvas.sum())
                - intersection
            )
            fuzzy_iou = float(intersection / union) if union else 1.0
            if (
                max(edge_errors) <= 4.0
                and center_error <= 4.0
                and area_error <= 0.50
                and fuzzy_iou >= 0.50
            ):
                score = max(edge_errors) + center_error + area_error + (1 - fuzzy_iou)
                candidates.append(
                    (
                        score,
                        reference_index,
                        {
                            "reference_component_id": reference_record["component_id"],
                            "render_component_id": render_record["component_id"],
                            "reference": reference_geometry,
                            "render": render_geometry,
                            "bbox_edge_errors_px": [
                                round(value, 4) for value in edge_errors
                            ],
                            "center_error_px": round(center_error, 4),
                            "area_relative_error": round(area_error, 4),
                            "fuzzy_mask_iou": round(fuzzy_iou, 4),
                        },
                    )
                )
        if not candidates:
            unmatched.append(
                {
                    "render_component_id": render_record["component_id"],
                    "render": render_geometry,
                }
            )
            continue
        _, reference_index, match = min(candidates, key=lambda item: item[0])
        used_reference.add(reference_index)
        matches.append(match)
    return {
        "reference_component_count": len(reference_records),
        "render_component_count": len(render_records),
        "matches": matches,
        "unmatched_render_components": unmatched,
        "new_render_area_px": sum(
            int(item["render"]["area_px"]) for item in unmatched
        ),
        "pass": not unmatched,
    }


def _pixel_clearance(subject: np.ndarray, obstacle: np.ndarray) -> float | None:
    if not subject.any() or not obstacle.any():
        return None
    from scipy.ndimage import distance_transform_edt

    center_distance = float(distance_transform_edt(~obstacle)[subject].min())
    # Adjacent differently colored pixels have a center distance of one pixel
    # and therefore zero empty raster rows/columns between their ink masks.
    return max(0.0, center_distance - 1.0)


def _local_bbox_mask(
    shape: tuple[int, ...], bbox: Any, *, label: str
) -> tuple[np.ndarray, list[int]]:
    if not isinstance(bbox, (list, tuple)) or len(bbox) != 4:
        raise ValueError(f"{label} must be [x, y, width, height]")
    try:
        x, y, width, height = (int(value) for value in bbox)
    except (TypeError, ValueError) as exc:
        raise ValueError(f"{label} must contain integer-like values") from exc
    if width <= 0 or height <= 0:
        raise ValueError(f"{label} width and height must be positive")
    if x < 0 or y < 0 or x + width > shape[1] or y + height > shape[0]:
        raise ValueError(f"{label} lies outside its region crop")
    mask = np.zeros(shape[:2], dtype=bool)
    mask[y : y + height, x : x + width] = True
    return mask, [x, y, width, height]


def _evaluate_color_clearance_contracts(
    reference: np.ndarray,
    render: np.ndarray,
    contracts: Any,
    region_element_ids: Any = None,
) -> tuple[list[dict[str, Any]], bool]:
    if contracts is None:
        return [], True
    if not isinstance(contracts, list):
        return [
            {
                "id": "[invalid]",
                "pass": False,
                "error": "color_clearance_contracts must be an array",
            }
        ], False
    results: list[dict[str, Any]] = []
    all_pass = True
    for index, contract in enumerate(contracts):
        if not isinstance(contract, dict):
            results.append(
                {
                    "id": f"clearance-{index + 1}",
                    "pass": False,
                    "error": "color clearance contract must be an object",
                }
            )
            all_pass = False
            continue
        contract_id = contract.get("id", f"clearance-{index + 1}")
        try:
            scoped_element_ids = (
                set(region_element_ids) if isinstance(region_element_ids, list) else set()
            )
            subject_element_ids = contract.get("subject_element_ids")
            if (
                not isinstance(subject_element_ids, list)
                or not subject_element_ids
                or any(
                    not isinstance(value, str) or not value
                    for value in subject_element_ids
                )
                or len(subject_element_ids) != len(set(subject_element_ids))
                or not set(subject_element_ids) <= scoped_element_ids
            ):
                raise ValueError(
                    "color clearance subject_element_ids must be a nonempty "
                    "unique subset of the region element_ids"
                )
            channel_name = contract.get("subject_dominant_channel")
            minimum_dominance = float(contract.get("minimum_dominance", 20.0))
            background_rgb = contract.get("background_rgb", [255, 255, 255])
            background_tolerance = float(contract.get("background_tolerance", 24.0))
            absolute_minimum = float(contract.get("absolute_min_px", 1.0))
            reference_loss_tolerance = float(
                contract.get("reference_loss_tolerance_px", 1.0)
            )
            if minimum_dominance < 0 or background_tolerance < 0:
                raise ValueError("color clearance mask tolerances must be nonnegative")
            if absolute_minimum < 0 or reference_loss_tolerance < 0:
                raise ValueError("color clearance distance tolerances must be nonnegative")
            subject_bbox_mask, subject_bbox = _local_bbox_mask(
                reference.shape,
                contract.get("subject_bbox"),
                label=f"{contract_id}.subject_bbox",
            )
            obstacles = contract.get("obstacles")
            if not isinstance(obstacles, list) or not obstacles:
                raise ValueError(
                    "color clearance requires a nonempty obstacles array with explicit bboxes"
                )
            if any(not isinstance(obstacle, dict) for obstacle in obstacles):
                raise ValueError("every color-clearance obstacle must be an object")
            for obstacle in obstacles:
                obstacle_element_ids = obstacle.get("element_ids")
                if (
                    not isinstance(obstacle_element_ids, list)
                    or not obstacle_element_ids
                    or any(
                        not isinstance(value, str) or not value
                        for value in obstacle_element_ids
                    )
                    or len(obstacle_element_ids) != len(set(obstacle_element_ids))
                    or not set(obstacle_element_ids) <= scoped_element_ids
                    or set(obstacle_element_ids) & set(subject_element_ids)
                ):
                    raise ValueError(
                        "every obstacle element_ids must be a nonempty unique "
                        "region subset disjoint from subject_element_ids"
                    )
            reference_dominant = _dominant_color_mask(
                reference, channel_name, minimum_dominance
            )
            render_dominant = _dominant_color_mask(
                render, channel_name, minimum_dominance
            )
            reference_subject = _components_touching_bbox(
                reference_dominant, subject_bbox_mask
            )
            render_subject = _components_touching_bbox(
                render_dominant, subject_bbox_mask
            )
            reference_foreground = _foreground_mask(
                reference, background_rgb, background_tolerance
            )
            render_foreground = _foreground_mask(render, background_rgb, background_tolerance)
        except (TypeError, ValueError) as exc:
            results.append({"id": contract_id, "pass": False, "error": str(exc)})
            all_pass = False
            continue
        if not reference_subject.any() or not render_subject.any():
            results.append(
                {
                    "id": contract_id,
                    "pass": False,
                    "error": "color clearance requires subject ink in both images",
                    "subject_bbox": subject_bbox,
                }
            )
            all_pass = False
            continue
        try:
            subject_bbox_tolerance = float(
                contract.get("subject_bbox_tolerance_px", 1.5)
            )
            subject_center_tolerance = float(
                contract.get("subject_center_tolerance_px", 1.0)
            )
            subject_area_tolerance = float(
                contract.get("subject_area_relative_tolerance", 0.15)
            )
            subject_bbox_halo = min(
                int(contract.get("subject_bbox_halo_px", 1)), 2
            )
            subject_overflow_tolerance = min(
                int(contract.get("subject_overflow_tolerance_px", 0)), 0
            )
            if subject_bbox_tolerance < 0 or subject_center_tolerance < 0:
                raise ValueError("subject bbox/center tolerances must be nonnegative")
            if not 0 <= subject_area_tolerance <= 1:
                raise ValueError("subject area tolerance must lie within [0, 1]")
            if subject_bbox_halo < 0 or subject_overflow_tolerance < 0:
                raise ValueError("subject overflow tolerances must be nonnegative")
            from scipy.ndimage import binary_dilation

            allowed_subject_mask = (
                binary_dilation(subject_bbox_mask, iterations=subject_bbox_halo)
                if subject_bbox_halo
                else subject_bbox_mask
            )
            reference_outside_subject_color = (
                reference_dominant & ~allowed_subject_mask
            )
            render_outside_subject_color = render_dominant & ~allowed_subject_mask
            # Existing same-color objects outside the subject are legal. Match
            # them as whole components so renderer antialiasing cannot create
            # false additions, while a new disconnected glyph remains visible.
            outside_component_audit = _outside_component_delta(
                reference_outside_subject_color, render_outside_subject_color
            )
            new_outside_subject_color_px = int(
                outside_component_audit["new_render_area_px"]
            )
            outside_subject_color_pass = bool(outside_component_audit["pass"])
            reference_subject_overflow = int(
                np.logical_and(reference_subject, ~allowed_subject_mask).sum()
            )
            render_subject_overflow = int(
                np.logical_and(render_subject, ~allowed_subject_mask).sum()
            )
            subject_overflow_pass = (
                render_subject_overflow
                <= reference_subject_overflow + subject_overflow_tolerance
            )
            reference_subject_geometry = _mask_geometry(reference_subject)
            render_subject_geometry = _mask_geometry(render_subject)
            assert reference_subject_geometry is not None
            assert render_subject_geometry is not None
            subject_bbox_errors = [
                abs(float(left) - float(right))
                for left, right in zip(
                    reference_subject_geometry["edges"],
                    render_subject_geometry["edges"],
                    strict=True,
                )
            ]
            subject_center_error = math.dist(
                reference_subject_geometry["center"],
                render_subject_geometry["center"],
            )
            subject_area_error = abs(
                int(render_subject_geometry["area_px"])
                - int(reference_subject_geometry["area_px"])
            ) / max(int(reference_subject_geometry["area_px"]), 1)
            subject_x, subject_y, subject_width, subject_height = subject_bbox
            reference_subject_crop = reference[
                subject_y : subject_y + subject_height,
                subject_x : subject_x + subject_width,
            ]
            render_subject_crop = render[
                subject_y : subject_y + subject_height,
                subject_x : subject_x + subject_width,
            ]
            subject_ssim = _ssim(reference_subject_crop, render_subject_crop)
            subject_edge_iou = _edge_iou(reference_subject_crop, render_subject_crop)
            subject_mask_iou = _fuzzy_mask_iou(reference_subject, render_subject)
            subject_mask_iou_min = max(
                float(
                    contract.get(
                        "subject_mask_iou_min", SUBJECT_MASK_IOU_FLOOR
                    )
                ),
                SUBJECT_MASK_IOU_FLOOR,
            )
            subject_edge_iou_min = max(
                float(
                    contract.get(
                        "subject_edge_iou_min", CRITICAL_EDGE_IOU_FLOOR
                    )
                ),
                CRITICAL_EDGE_IOU_FLOOR,
            )
            subject_color_delta_max = min(
                float(contract.get("subject_foreground_delta_e_max", 12.0)),
                12.0,
            )
            if subject_color_delta_max < 0:
                raise ValueError("subject color tolerance must be nonnegative")
            subject_color_delta = _delta_e(
                reference[reference_subject].astype(float).mean(axis=0),
                render[render_subject].astype(float).mean(axis=0),
            )
            subject_pass = (
                max(subject_bbox_errors) <= subject_bbox_tolerance
                and subject_center_error <= subject_center_tolerance
                and subject_area_error <= subject_area_tolerance
                and subject_mask_iou >= subject_mask_iou_min
                and subject_edge_iou >= subject_edge_iou_min
                and subject_color_delta <= subject_color_delta_max
                and subject_overflow_pass
                and outside_subject_color_pass
            )
        except (AssertionError, TypeError, ValueError) as exc:
            results.append(
                {
                    "id": contract_id,
                    "pass": False,
                    "error": f"invalid subject geometry contract: {exc}",
                }
            )
            all_pass = False
            continue
        obstacle_results: list[dict[str, Any]] = []
        contract_pass = subject_pass
        for obstacle_index, obstacle in enumerate(obstacles, start=1):
            obstacle_id = str(obstacle.get("id") or f"obstacle-{obstacle_index}")
            try:
                obstacle_bbox_mask, obstacle_bbox = _local_bbox_mask(
                    reference.shape,
                    obstacle.get("bbox"),
                    label=f"{contract_id}.{obstacle_id}.bbox",
                )
                area_tolerance = float(
                    obstacle.get("area_relative_tolerance", 0.25)
                )
                bbox_tolerance = float(obstacle.get("bbox_tolerance_px", 2.0))
                obstacle_absolute_minimum = float(
                    obstacle.get("absolute_min_px", absolute_minimum)
                )
                obstacle_reference_loss = float(
                    obstacle.get(
                        "reference_loss_tolerance_px", reference_loss_tolerance
                    )
                )
                if not 0 <= area_tolerance <= 1 or bbox_tolerance < 0:
                    raise ValueError("obstacle geometry tolerances are invalid")
                if obstacle_absolute_minimum < 0 or obstacle_reference_loss < 0:
                    raise ValueError("obstacle clearance tolerances must be nonnegative")
                obstacle_edge_iou_min = max(
                    float(
                        obstacle.get("edge_iou_min", CRITICAL_EDGE_IOU_FLOOR)
                    ),
                    CRITICAL_EDGE_IOU_FLOOR,
                )
                obstacle_color_delta_max = min(
                    float(obstacle.get("foreground_delta_e_max", 12.0)), 12.0
                )
                obstacle_color_cosine_min = max(
                    float(
                        obstacle.get(
                            "target_color_cosine_min",
                            OBSTACLE_COLOR_COSINE_FLOOR,
                        )
                    ),
                    OBSTACLE_COLOR_COSINE_FLOOR,
                )
                if (
                    obstacle_color_delta_max < 0
                    or not 0 <= obstacle_color_cosine_min <= 1
                ):
                    raise ValueError("obstacle pixel tolerances must be nonnegative")
                reference_obstacle_ink = (
                    reference_foreground & obstacle_bbox_mask & ~reference_subject
                )
                render_obstacle_ink = (
                    render_foreground & obstacle_bbox_mask & ~render_subject
                )
                if not reference_obstacle_ink.any() or not render_obstacle_ink.any():
                    raise ValueError("obstacle ink is missing in reference or render")
                reference_core_rgb = _core_foreground_rgb(
                    reference, reference_obstacle_ink, background_rgb
                )
                render_core_rgb = _core_foreground_rgb(
                    render, render_obstacle_ink, background_rgb
                )
                reference_obstacle = _target_color_mask(
                    reference,
                    reference_obstacle_ink,
                    background_rgb,
                    reference_core_rgb,
                    obstacle_color_cosine_min,
                )
                render_obstacle = _target_color_mask(
                    render,
                    render_obstacle_ink,
                    background_rgb,
                    reference_core_rgb,
                    obstacle_color_cosine_min,
                )
                reference_geometry = _mask_geometry(reference_obstacle)
                render_geometry = _mask_geometry(render_obstacle)
                if reference_geometry is None:
                    raise ValueError("reference obstacle has no target-color ink")
                (
                    obstacle_mask_iou_min,
                    reference_fill_fraction,
                    thin_open_reference,
                ) = _obstacle_mask_floor(
                    reference_geometry, obstacle.get("mask_iou_min")
                )
                reference_clearance = _pixel_clearance(
                    reference_subject, reference_obstacle_ink
                )
                render_clearance = _pixel_clearance(
                    render_subject, render_obstacle_ink
                )
                if reference_clearance is None or render_clearance is None:
                    raise ValueError("subject-to-obstacle clearance is not measurable")
                required = max(
                    obstacle_absolute_minimum,
                    reference_clearance - obstacle_reference_loss,
                )
                if render_geometry is None:
                    area_relative_error = 1.0
                    bbox_edge_errors: list[float] | None = None
                else:
                    area_relative_error = abs(
                        int(render_geometry["area_px"])
                        - int(reference_geometry["area_px"])
                    ) / max(int(reference_geometry["area_px"]), 1)
                    bbox_edge_errors = [
                        abs(float(left) - float(right))
                        for left, right in zip(
                            reference_geometry["edges"],
                            render_geometry["edges"],
                            strict=True,
                        )
                    ]
                obstacle_x, obstacle_y, obstacle_width, obstacle_height = obstacle_bbox
                reference_obstacle_crop = reference[
                    obstacle_y : obstacle_y + obstacle_height,
                    obstacle_x : obstacle_x + obstacle_width,
                ]
                render_obstacle_crop = render[
                    obstacle_y : obstacle_y + obstacle_height,
                    obstacle_x : obstacle_x + obstacle_width,
                ]
                obstacle_ssim = _ssim(reference_obstacle_crop, render_obstacle_crop)
                obstacle_edge_iou = _edge_iou(
                    reference_obstacle_crop, render_obstacle_crop
                )
                obstacle_mask_iou = _fuzzy_mask_iou(
                    reference_obstacle, render_obstacle
                )
                obstacle_mean_delta = float(
                    np.abs(
                        reference_obstacle_crop.astype(float)
                        - render_obstacle_crop.astype(float)
                    ).mean()
                )
                obstacle_color_delta = _delta_e(
                    reference_core_rgb, render_core_rgb
                )
                obstacle_pass = (
                    render_clearance >= required
                    and area_relative_error <= area_tolerance
                    and bbox_edge_errors is not None
                    and max(bbox_edge_errors) <= bbox_tolerance
                    and obstacle_mask_iou >= obstacle_mask_iou_min
                    and obstacle_edge_iou >= obstacle_edge_iou_min
                    and obstacle_color_delta <= obstacle_color_delta_max
                )
                obstacle_results.append(
                    {
                        "id": obstacle_id,
                        "element_ids": obstacle.get("element_ids"),
                        "bbox": obstacle_bbox,
                        "reference": reference_geometry,
                        "render": render_geometry,
                        "reference_clearance_px": round(reference_clearance, 4),
                        "render_clearance_px": round(render_clearance, 4),
                        "required_clearance_px": round(required, 4),
                        "area_relative_error": round(area_relative_error, 4),
                        "area_relative_tolerance": area_tolerance,
                        "bbox_edge_errors_px": (
                            [round(value, 4) for value in bbox_edge_errors]
                            if bbox_edge_errors is not None
                            else None
                        ),
                        "bbox_tolerance_px": bbox_tolerance,
                        "mask_iou": round(obstacle_mask_iou, 4),
                        "mask_iou_min": obstacle_mask_iou_min,
                        "reference_fill_fraction": round(
                            reference_fill_fraction, 4
                        ),
                        "thin_open_reference": thin_open_reference,
                        "edge_iou": round(obstacle_edge_iou, 4),
                        "edge_iou_min": obstacle_edge_iou_min,
                        "rgb_ssim_diagnostic": round(obstacle_ssim, 4),
                        "mean_abs_rgb_delta_diagnostic": round(
                            obstacle_mean_delta, 4
                        ),
                        "reference_core_rgb": [
                            round(float(value), 4) for value in reference_core_rgb
                        ],
                        "render_core_rgb": [
                            round(float(value), 4) for value in render_core_rgb
                        ],
                        "target_color_cosine_min": obstacle_color_cosine_min,
                        "core_foreground_delta_e00": round(
                            obstacle_color_delta, 4
                        ),
                        "core_foreground_delta_e_max": obstacle_color_delta_max,
                        "foreground_delta_e00": round(obstacle_color_delta, 4),
                        "foreground_delta_e_max": obstacle_color_delta_max,
                        "pass": obstacle_pass,
                    }
                )
            except (TypeError, ValueError) as exc:
                obstacle_pass = False
                obstacle_results.append(
                    {"id": obstacle_id, "pass": False, "error": str(exc)}
                )
            contract_pass = contract_pass and obstacle_pass
        reference_clearances = [
            item["reference_clearance_px"]
            for item in obstacle_results
            if item.get("reference_clearance_px") is not None
        ]
        render_clearances = [
            item["render_clearance_px"]
            for item in obstacle_results
            if item.get("render_clearance_px") is not None
        ]
        required_clearances = [
            item["required_clearance_px"]
            for item in obstacle_results
            if item.get("required_clearance_px") is not None
        ]
        passed = contract_pass and len(obstacle_results) == len(obstacles)
        all_pass = all_pass and passed
        results.append(
            {
                "id": contract_id,
                "subject_dominant_channel": channel_name,
                "subject_element_ids": subject_element_ids,
                "subject_bbox": subject_bbox,
                "subject": {
                    "reference": reference_subject_geometry,
                    "render": render_subject_geometry,
                    "bbox_edge_errors_px": [
                        round(value, 4) for value in subject_bbox_errors
                    ],
                    "bbox_tolerance_px": subject_bbox_tolerance,
                    "center_error_px": round(subject_center_error, 4),
                    "center_tolerance_px": subject_center_tolerance,
                    "area_relative_error": round(subject_area_error, 4),
                    "area_relative_tolerance": subject_area_tolerance,
                    "rgb_ssim_diagnostic": round(subject_ssim, 4),
                    "mask_iou": round(subject_mask_iou, 4),
                    "mask_iou_min": subject_mask_iou_min,
                    "edge_iou": round(subject_edge_iou, 4),
                    "edge_iou_min": subject_edge_iou_min,
                    "foreground_delta_e00": round(subject_color_delta, 4),
                    "foreground_delta_e_max": subject_color_delta_max,
                    "reference_overflow_px": reference_subject_overflow,
                    "render_overflow_px": render_subject_overflow,
                    "subject_bbox_halo_px": subject_bbox_halo,
                    "overflow_tolerance_px": subject_overflow_tolerance,
                    "overflow_pass": subject_overflow_pass,
                    "reference_outside_subject_color_px": int(
                        reference_outside_subject_color.sum()
                    ),
                    "render_outside_subject_color_px": int(
                        render_outside_subject_color.sum()
                    ),
                    "new_outside_subject_color_px": new_outside_subject_color_px,
                    "outside_subject_color_pass": outside_subject_color_pass,
                    "outside_subject_color_components": outside_component_audit,
                    "pass": subject_pass,
                },
                "reference_clearance_px": min(reference_clearances, default=None),
                "render_clearance_px": min(render_clearances, default=None),
                "required_clearance_px": min(required_clearances, default=None),
                "absolute_min_px": absolute_minimum,
                "reference_loss_tolerance_px": reference_loss_tolerance,
                "obstacles": obstacle_results,
                "pass": passed,
            }
        )
    return results, all_pass


def _audit_critical_region_expectation(
    payload: dict[str, Any],
) -> tuple[dict[str, Any] | None, list[str]]:
    """Verify an exact inventory and gate shape for required critical regions."""

    expectation = payload.get("critical_region_expectation")
    critical_regions = [
        region
        for region in payload.get("regions", [])
        if isinstance(region, dict) and region.get("critical") is True
    ]
    if expectation is None:
        if critical_regions:
            return {
                "pass": False,
                "error": "critical_region_expectation is required for critical regions",
            }, ["regions:expectation:missing"]
        return None, []
    if not isinstance(expectation, dict):
        return {
            "pass": False,
            "error": "critical_region_expectation must be an object",
        }, ["regions:expectation:invalid"]
    expected_contracts = expectation.get("contracts")
    expected_count = expectation.get("count")
    if not critical_regions and expected_count == 0 and expected_contracts == []:
        return {
            "expected_count": 0,
            "actual_count": 0,
            "expected_region_ids": [],
            "actual_region_ids": [],
            "contracts": [],
            "pass": True,
        }, []
    if (
        not isinstance(expected_contracts, list)
        or not expected_contracts
        or any(not isinstance(value, dict) for value in expected_contracts)
        or not isinstance(expected_count, int)
        or isinstance(expected_count, bool)
        or expected_count <= 0
        or expected_count != len(expected_contracts)
    ):
        return {
            "pass": False,
            "error": (
                "critical_region_expectation requires a positive count and the "
                "same number of region contract objects"
            ),
        }, ["regions:expectation:invalid"]
    allowed_gates = {
        "pixel_metrics",
        "ink_contract",
        "color_clearance_contracts",
        "color_probes",
        "required_relations",
    }
    expected_ids: list[str] = []
    normalized_contracts: list[dict[str, Any]] = []
    for contract in expected_contracts:
        region_id = contract.get("id")
        bbox = contract.get("bbox")
        pixel_bbox = contract.get("pixel_bbox")
        element_ids = contract.get("element_ids")
        required_gates = contract.get("required_gates")
        gate_hashes = contract.get("gate_sha256")
        hashed_gates = set(required_gates or []) - {"pixel_metrics"}
        if (
            not isinstance(region_id, str)
            or not region_id
            or not isinstance(bbox, list)
            or len(bbox) != 4
            or any(not isinstance(value, (int, float)) for value in bbox)
            or (
                pixel_bbox is not None
                and (
                    not isinstance(pixel_bbox, list)
                    or len(pixel_bbox) != 4
                    or any(
                        not isinstance(value, (int, float))
                        for value in pixel_bbox
                    )
                )
            )
            or not isinstance(element_ids, list)
            or not element_ids
            or any(not isinstance(value, str) or not value for value in element_ids)
            or len(element_ids) != len(set(element_ids))
            or not isinstance(required_gates, list)
            or not required_gates
            or any(gate not in allowed_gates for gate in required_gates)
            or len(required_gates) != len(set(required_gates))
            or "pixel_metrics" not in required_gates
            or not isinstance(gate_hashes, dict)
            or set(gate_hashes) != hashed_gates
            or any(
                not isinstance(value, str)
                or len(value) != 64
                or any(character not in "0123456789abcdef" for character in value)
                for value in gate_hashes.values()
            )
        ):
            return {
                "pass": False,
                "error": (
                    "every critical-region expectation needs id, bbox, exact "
                    "nonempty element_ids, unique required_gates including "
                    "pixel_metrics, and exact hashes for every object-level gate"
                ),
            }, ["regions:expectation:invalid"]
        expected_ids.append(region_id)
        normalized_contract = {
                "id": region_id,
                "bbox": bbox,
                "element_ids": element_ids,
                "required_gates": required_gates,
                "gate_sha256": gate_hashes,
            }
        if pixel_bbox is not None:
            normalized_contract["pixel_bbox"] = pixel_bbox
        normalized_contracts.append(normalized_contract)
    if len(expected_ids) != len(set(expected_ids)):
        return {
            "pass": False,
            "error": "critical-region expectation ids must be unique",
        }, ["regions:expectation:invalid"]
    actual_ids = [
        str(region.get("id") or "")
        for region in payload.get("regions", [])
        if isinstance(region, dict) and region.get("critical") is True
    ]
    duplicates = sorted({value for value in actual_ids if actual_ids.count(value) > 1})
    missing = sorted(set(expected_ids) - set(actual_ids))
    unexpected = sorted(set(actual_ids) - set(expected_ids))
    blockers: list[str] = []
    if len(actual_ids) != expected_count:
        blockers.append("regions:expectation:count-mismatch")
    if missing:
        blockers.append("regions:expectation:missing-regions")
    if unexpected:
        blockers.append("regions:expectation:unexpected-regions")
    if duplicates:
        blockers.append("regions:expectation:duplicate-regions")
    actual_by_id = {
        str(region.get("id")): region
        for region in payload.get("regions", [])
        if isinstance(region, dict)
        and isinstance(region.get("id"), str)
        and region.get("critical") is True
    }
    contract_records: list[dict[str, Any]] = []
    for contract in normalized_contracts:
        region = actual_by_id.get(contract["id"])
        if region is None:
            contract_records.append({**contract, "pass": False, "error": "missing region"})
            continue
        actual_element_ids = region.get("element_ids", [])
        bbox_pass = region.get("bbox") == contract["bbox"]
        expected_pixel_bbox = contract.get("pixel_bbox")
        actual_pixel_bbox = region.get("pixel_bbox")
        pixel_bbox_pass = actual_pixel_bbox == expected_pixel_bbox
        element_ids_pass = actual_element_ids == contract["element_ids"]
        gate_passes = {
            gate: (
                True
                if gate == "pixel_metrics"
                else isinstance(region.get(gate), dict)
                if gate == "ink_contract"
                else isinstance(region.get(gate), list) and bool(region.get(gate))
            )
            for gate in contract["required_gates"]
        }
        actual_gate_hashes = {
            gate: _gate_sha256(region.get(gate))
            for gate in contract["required_gates"]
            if gate != "pixel_metrics" and gate_passes[gate]
        }
        gate_hash_passes = {
            gate: actual_gate_hashes.get(gate) == expected_hash
            for gate, expected_hash in contract["gate_sha256"].items()
        }
        contract_pass = (
            bbox_pass
            and pixel_bbox_pass
            and element_ids_pass
            and all(gate_passes.values())
            and all(gate_hash_passes.values())
        )
        if not bbox_pass:
            blockers.append(f"regions:expectation:{contract['id']}:bbox-mismatch")
        if not pixel_bbox_pass:
            blockers.append(
                f"regions:expectation:{contract['id']}:pixel-bbox-mismatch"
            )
        if not element_ids_pass:
            blockers.append(f"regions:expectation:{contract['id']}:element-ids-mismatch")
        for gate, gate_pass in gate_passes.items():
            if not gate_pass:
                blockers.append(
                    f"regions:expectation:{contract['id']}:missing-gate:{gate}"
                )
        for gate, hash_pass in gate_hash_passes.items():
            if not hash_pass:
                blockers.append(
                    f"regions:expectation:{contract['id']}:gate-hash-mismatch:{gate}"
                )
        contract_records.append(
            {
                **contract,
                "actual_bbox": region.get("bbox"),
                "actual_pixel_bbox": actual_pixel_bbox,
                "actual_element_ids": actual_element_ids,
                "gate_passes": gate_passes,
                "actual_gate_sha256": actual_gate_hashes,
                "gate_hash_passes": gate_hash_passes,
                "pass": contract_pass,
            }
        )
    return {
        "expected_count": expected_count,
        "actual_count": len(actual_ids),
        "expected_region_ids": expected_ids,
        "actual_region_ids": actual_ids,
        "missing_region_ids": missing,
        "unexpected_region_ids": unexpected,
        "duplicate_region_ids": duplicates,
        "contracts": contract_records,
        "pass": not blockers,
    }, blockers


def build_critical_region_expectation(payload: dict[str, Any]) -> dict[str, Any]:
    """Freeze the current critical-region inventory after author review."""

    contracts: list[dict[str, Any]] = []
    object_gates = (
        "ink_contract",
        "color_clearance_contracts",
        "color_probes",
        "required_relations",
    )
    for region in payload.get("regions", []):
        if not isinstance(region, dict) or region.get("critical") is not True:
            continue
        element_ids = region.get("element_ids")
        if not isinstance(element_ids, list) or not element_ids:
            raise ValueError(
                f"critical region {region.get('id')!r} needs nonempty element_ids"
            )
        required_gates = ["pixel_metrics"]
        for gate in object_gates:
            value = region.get(gate)
            if isinstance(value, dict) or isinstance(value, list) and value:
                required_gates.append(gate)
        contract = {
                "id": region["id"],
                "bbox": region["bbox"],
                "element_ids": element_ids,
                "required_gates": required_gates,
                "gate_sha256": {
                    gate: _gate_sha256(region[gate])
                    for gate in required_gates
                    if gate != "pixel_metrics"
                },
            }
        if "pixel_bbox" in region:
            contract["pixel_bbox"] = region["pixel_bbox"]
        contracts.append(contract)
    return {"count": len(contracts), "contracts": contracts}


def _bbox_union(bboxes: list[list[float]]) -> list[float] | None:
    if not bboxes:
        return None
    left = min(bbox[0] for bbox in bboxes)
    top = min(bbox[1] for bbox in bboxes)
    right = max(bbox[0] + bbox[2] for bbox in bboxes)
    bottom = max(bbox[1] + bbox[3] for bbox in bboxes)
    return [left, top, right - left, bottom - top]


def _binding_bboxes(run: common.Run) -> dict[str, list[list[float]]]:
    binding_document = read_json(run.bindings_path)
    bindings = binding_document.get("bindings", [])
    shape_bboxes: dict[int, tuple[str, list[float]]] = {}
    if run.pptx_path.exists():
        try:
            from pptx import Presentation

            presentation = Presentation(run.pptx_path)
            width_px, height_px = common.image_size(run.source_png)
            scale_x = width_px / float(presentation.slide_width)
            scale_y = height_px / float(presentation.slide_height)

            def visit(shapes: Any) -> None:
                for shape in shapes:
                    shape_bboxes[int(shape.shape_id)] = (
                        str(shape.name),
                        [
                            float(shape.left) * scale_x,
                            float(shape.top) * scale_y,
                            float(shape.width) * scale_x,
                            float(shape.height) * scale_y,
                        ],
                    )
                    if getattr(shape, "shapes", None) is not None:
                        visit(shape.shapes)

            for slide in presentation.slides:
                visit(slide.shapes)
        except (KeyError, TypeError, ValueError, OSError):
            shape_bboxes = {}
    result: dict[str, list[list[float]]] = {}
    for binding in bindings:
        if not isinstance(binding, dict) or not isinstance(
            binding.get("element_id"), str
        ):
            continue
        shape_id = binding.get("shape_id")
        if not isinstance(shape_id, int) or shape_id not in shape_bboxes:
            continue
        actual_name, normalized = shape_bboxes[shape_id]
        if binding.get("shape_name") != actual_name:
            continue
        result.setdefault(binding["element_id"], []).append(normalized)
    for binding in binding_document.get("logical_group_bindings", []):
        if not isinstance(binding, dict) or not isinstance(
            binding.get("element_id"), str
        ):
            continue
        backend_ids = binding.get("backend_object_ids")
        backend_names = binding.get("backend_object_names")
        if (
            not isinstance(backend_ids, list)
            or not isinstance(backend_names, list)
            or len(backend_ids) != len(backend_names)
        ):
            continue
        group_bboxes: list[list[float]] = []
        for shape_id, shape_name in zip(backend_ids, backend_names, strict=True):
            if not isinstance(shape_id, int) or shape_id not in shape_bboxes:
                continue
            actual_name, normalized = shape_bboxes[shape_id]
            if shape_name == actual_name:
                group_bboxes.append(normalized)
        if len(group_bboxes) == len(backend_ids) and group_bboxes:
            result[binding["element_id"]] = group_bboxes
    return result


def _bbox_identity_metrics(expected: list[float], actual: list[float] | None) -> dict[str, Any]:
    """Bind a logical target to a frozen native-object envelope.

    ``expected`` is deliberately not a visual-ink box.  It is the case-owned
    native PowerPoint binding envelope, frozen through the critical-region gate
    hash.  Every actual object edge (or both endpoints of an axis-aligned line)
    must remain inside that envelope with at most a three-pixel raster/host halo.
    """

    if actual is None:
        return {"expected_bbox": expected, "actual_bbox": None, "pass": False}
    halo = 3.0
    expected_left, expected_top, expected_width, expected_height = expected
    actual_left, actual_top, actual_width, actual_height = actual
    if (actual_width == 0) != (actual_height == 0):
        vertical = actual_width == 0
        line_length = actual_height if vertical else actual_width
        expected_length = expected_height if vertical else expected_width
        center_error = math.dist(
            (expected_left + expected_width / 2, expected_top + expected_height / 2),
            (actual_left + actual_width / 2, actual_top + actual_height / 2),
        )
        if vertical:
            endpoints = [
                (actual_left, actual_top),
                (actual_left, actual_top + actual_height),
            ]
        else:
            endpoints = [
                (actual_left, actual_top),
                (actual_left + actual_width, actual_top),
            ]
        endpoints_inside = all(
            expected_left - halo <= x <= expected_left + expected_width + halo
            and expected_top - halo <= y <= expected_top + expected_height + halo
            for x, y in endpoints
        )
        length_ratio = line_length / max(expected_length, 1e-9)
        passed = (
            endpoints_inside
            and center_error <= halo
            and 0.70 <= length_ratio <= 1.30
        )
        return {
            "expected_bbox": [round(value, 4) for value in expected],
            "actual_bbox": [round(value, 4) for value in actual],
            "geometry_kind": "vertical-line" if vertical else "horizontal-line",
            "center_error_px": round(center_error, 4),
            "center_tolerance_px": halo,
            "line_endpoints": [
                [round(coordinate, 4) for coordinate in point] for point in endpoints
            ],
            "line_endpoints_inside_native_bbox_with_3px_halo": endpoints_inside,
            "length_ratio": round(length_ratio, 4),
            "length_ratio_min": 0.70,
            "length_ratio_max": 1.30,
            "pass": passed,
        }
    intersection_width = max(
        0.0,
        min(expected_left + expected_width, actual_left + actual_width)
        - max(expected_left, actual_left),
    )
    intersection_height = max(
        0.0,
        min(expected_top + expected_height, actual_top + actual_height)
        - max(expected_top, actual_top),
    )
    intersection = intersection_width * intersection_height
    expected_area = expected_width * expected_height
    actual_area = actual_width * actual_height
    overlap = intersection / max(min(expected_area, actual_area), 1e-9)
    center_error = math.dist(
        (expected_left + expected_width / 2, expected_top + expected_height / 2),
        (actual_left + actual_width / 2, actual_top + actual_height / 2),
    )
    center_tolerance = min(
        12.0, max(3.0, math.hypot(expected_width, expected_height) * 0.12)
    )
    area_ratio = max(expected_area, actual_area) / max(
        min(expected_area, actual_area), 1e-9
    )
    contained = (
        actual_left >= expected_left - halo
        and actual_top >= expected_top - halo
        and actual_left + actual_width <= expected_left + expected_width + halo
        and actual_top + actual_height <= expected_top + expected_height + halo
    )
    passed = (
        overlap >= 0.35
        and center_error <= center_tolerance
        and area_ratio <= 4.0
        and contained
    )
    return {
        "expected_bbox": [round(value, 4) for value in expected],
        "actual_bbox": [round(value, 4) for value in actual],
        "intersection_over_smaller": round(overlap, 4),
        "intersection_over_smaller_min": 0.35,
        "center_error_px": round(center_error, 4),
        "center_tolerance_px": round(center_tolerance, 4),
        "area_ratio": round(area_ratio, 4),
        "area_ratio_max": 4.0,
        "actual_fully_inside_native_bbox_with_3px_halo": contained,
        "pass": passed,
    }


def _bbox_containment_metrics(
    expected: list[float], actual: list[float] | None
) -> dict[str, Any]:
    """Require one member object to stay inside its logical union envelope."""

    if actual is None:
        return {"expected_bbox": expected, "actual_bbox": None, "pass": False}
    halo = 3.0
    expected_left, expected_top, expected_width, expected_height = expected
    actual_left, actual_top, actual_width, actual_height = actual
    if (actual_width == 0) != (actual_height == 0):
        if actual_width == 0:
            points = [
                (actual_left, actual_top),
                (actual_left, actual_top + actual_height),
            ]
        else:
            points = [
                (actual_left, actual_top),
                (actual_left + actual_width, actual_top),
            ]
        contained = all(
            expected_left - halo <= x <= expected_left + expected_width + halo
            and expected_top - halo <= y <= expected_top + expected_height + halo
            for x, y in points
        )
        return {
            "expected_bbox": [round(value, 4) for value in expected],
            "actual_bbox": [round(value, 4) for value in actual],
            "line_endpoints": [
                [round(coordinate, 4) for coordinate in point] for point in points
            ],
            "line_endpoints_inside_native_bbox_with_3px_halo": contained,
            "pass": contained,
        }
    contained = (
        actual_left >= expected_left - halo
        and actual_top >= expected_top - halo
        and actual_left + actual_width <= expected_left + expected_width + halo
        and actual_top + actual_height <= expected_top + expected_height + halo
    )
    return {
        "expected_bbox": [round(value, 4) for value in expected],
        "actual_bbox": [round(value, 4) for value in actual],
        "actual_fully_inside_native_bbox_with_3px_halo": contained,
        "pass": contained,
    }


def _audit_clearance_binding_bboxes(
    region: dict[str, Any],
    binding_bboxes: dict[str, list[list[float]]],
) -> tuple[list[dict[str, Any]], bool]:
    audits: list[dict[str, Any]] = []
    all_pass = True
    region_bbox = region.get("bbox")
    if not isinstance(region_bbox, list) or len(region_bbox) != 4:
        return [], False
    region_x, region_y = float(region_bbox[0]), float(region_bbox[1])
    for contract in region.get("color_clearance_contracts", []):
        if not isinstance(contract, dict):
            audits.append({"id": "[invalid]", "pass": False})
            all_pass = False
            continue
        targets: list[dict[str, Any]] = []
        targets.append(
            {
                "role": "subject",
                "id": contract.get("id"),
                "element_ids": contract.get("subject_element_ids", []),
                "native_binding_bbox": contract.get("native_binding_bbox"),
            }
        )
        for obstacle in contract.get("obstacles", []):
            if isinstance(obstacle, dict):
                targets.append(
                    {
                        "role": "obstacle",
                        "id": obstacle.get("id"),
                        "element_ids": obstacle.get("element_ids", []),
                        "native_binding_bbox": obstacle.get("native_binding_bbox"),
                    }
                )
        target_audits: list[dict[str, Any]] = []
        for target in targets:
            local_bbox = target["native_binding_bbox"]
            element_ids = target["element_ids"]
            if (
                not isinstance(local_bbox, list)
                or len(local_bbox) != 4
                or any(
                    not isinstance(value, (int, float))
                    or isinstance(value, bool)
                    or not math.isfinite(float(value))
                    for value in local_bbox
                )
                or float(local_bbox[0]) < 0
                or float(local_bbox[1]) < 0
                or float(local_bbox[2]) < 0
                or float(local_bbox[3]) < 0
                or float(local_bbox[2]) == float(local_bbox[3]) == 0
            ):
                target_audits.append(
                    {
                        "role": target["role"],
                        "id": target["id"],
                        "element_ids": element_ids,
                        "native_binding_bbox": local_bbox,
                        "pass": False,
                        "error": (
                            "native_binding_bbox is required and must be a finite "
                            "[x, y, width, height] native-object envelope"
                        ),
                    }
                )
                continue
            expected = [
                region_x + float(local_bbox[0]),
                region_y + float(local_bbox[1]),
                float(local_bbox[2]),
                float(local_bbox[3]),
            ]
            missing_element_ids = [
                element_id
                for element_id in element_ids
                if not binding_bboxes.get(element_id)
            ]
            resolved_bboxes = {
                element_id: binding_bboxes.get(element_id, [])
                for element_id in element_ids
            }
            actual = _bbox_union(
                [bbox for bboxes in resolved_bboxes.values() for bbox in bboxes]
            )
            metrics = _bbox_identity_metrics(expected, actual)
            per_element = []
            for element_id, bboxes in resolved_bboxes.items():
                element_bbox = _bbox_union(bboxes)
                if element_bbox is None:
                    per_element.append(
                        {"element_id": element_id, "actual_bbox": None, "pass": False}
                    )
                    continue
                element_metrics = _bbox_containment_metrics(expected, element_bbox)
                per_element.append(
                    {
                        "element_id": element_id,
                        **element_metrics,
                    }
                )
            target_pass = (
                metrics["pass"]
                and not missing_element_ids
                and all(item["pass"] for item in per_element)
            )
            target_audits.append(
                {
                    "role": target["role"],
                    "id": target["id"],
                    "element_ids": element_ids,
                    "native_binding_bbox": local_bbox,
                    "missing_element_ids": missing_element_ids,
                    "per_element": per_element,
                    **metrics,
                    "pass": target_pass,
                }
            )
        passed = bool(target_audits) and all(item["pass"] for item in target_audits)
        audits.append(
            {
                "id": contract.get("id"),
                "targets": target_audits,
                "pass": passed,
            }
        )
        all_pass = all_pass and passed
    return audits, all_pass


def evaluate_regions(run: common.Run) -> dict[str, Any]:
    payload = read_json(run.regions_path)
    from tools.region_contract import audit_region_contract

    contract_audit = audit_region_contract(run, payload)
    expectation_audit, expectation_blockers = _audit_critical_region_expectation(payload)
    with Image.open(run.source_png) as ref_image, Image.open(run.render_png) as ren_image:
        reference = np.asarray(ref_image.convert("RGB"), dtype=np.uint8)
        render = np.asarray(ren_image.convert("RGB"), dtype=np.uint8)
    if reference.shape != render.shape:
        raise common.fail(f"reference/render size mismatch: {reference.shape} != {render.shape}")

    defaults = payload.get("defaults", {})
    results: list[dict[str, Any]] = []
    binding_bboxes = _binding_bboxes(run)
    critical_count = 0
    critical_pass = True
    for region in payload.get("regions", []):
        critical = bool(region.get("critical", False))
        try:
            ref_crop = _crop(reference, region["bbox"])
            ren_crop = _crop(render, region["bbox"])
            pixel_bbox = region.get("pixel_bbox", region["bbox"])
            ref_pixel_crop = _crop(reference, pixel_bbox)
            ren_pixel_crop = _crop(render, pixel_bbox)
        except (KeyError, TypeError, ValueError) as exc:
            if critical:
                critical_count += 1
                critical_pass = False
            results.append(
                {
                    "id": str(region.get("id") or "[invalid-region]"),
                    "label": region.get("label", region.get("id", "[invalid-region]")),
                    "bbox": region.get("bbox"),
                    "critical": critical,
                    "pass": False,
                    "error": str(exc),
                }
            )
            continue
        mean_delta = float(
            np.abs(
                ref_pixel_crop.astype(float) - ren_pixel_crop.astype(float)
            ).mean()
        )
        ssim = _ssim(ref_pixel_crop, ren_pixel_crop)
        edge_iou = _edge_iou(ref_pixel_crop, ren_pixel_crop)
        thresholds = {**defaults, **region.get("thresholds", {})}
        if critical:
            thresholds["ssim_min"] = max(
                float(thresholds.get("ssim_min", CRITICAL_SSIM_FLOOR)),
                CRITICAL_SSIM_FLOOR,
            )
            thresholds["edge_iou_min"] = max(
                float(thresholds.get("edge_iou_min", CRITICAL_EDGE_IOU_FLOOR)),
                CRITICAL_EDGE_IOU_FLOOR,
            )
        region_pass = (
            ssim >= float(thresholds.get("ssim_min", 0.85))
            and edge_iou >= float(thresholds.get("edge_iou_min", 0.75))
        )
        probes, probes_pass = _evaluate_probes(reference, render, region.get("color_probes", []))
        ink_contract, ink_pass = _evaluate_ink_contract(
            ref_crop, ren_crop, region.get("ink_contract")
        )
        color_clearances, clearances_pass = _evaluate_color_clearance_contracts(
            ref_crop,
            ren_crop,
            region.get("color_clearance_contracts"),
            region.get("element_ids"),
        )
        clearance_bindings, clearance_bindings_pass = _audit_clearance_binding_bboxes(
            region, binding_bboxes
        )
        for clearance, binding_audit in zip(
            color_clearances, clearance_bindings, strict=False
        ):
            clearance["binding_bbox_audit"] = binding_audit
            clearance["pass"] = bool(clearance.get("pass")) and binding_audit["pass"]
        region_pass = (
            region_pass
            and probes_pass
            and ink_pass
            and clearances_pass
            and clearance_bindings_pass
        )
        if critical:
            critical_count += 1
            critical_pass = critical_pass and region_pass
        results.append(
            {
                "id": region["id"],
                "label": region.get("label", region["id"]),
                "bbox": region["bbox"],
                "pixel_bbox": pixel_bbox,
                "critical": critical,
                "mean_abs_rgb_delta": round(mean_delta, 4),
                "ssim": round(ssim, 4),
                "edge_iou": round(edge_iou, 4),
                "thresholds": thresholds,
                "color_probes": probes,
                "ink_contract": ink_contract,
                "color_clearance_contracts": color_clearances,
                "pass": region_pass,
            }
        )
    report = {
        "schema_version": SCHEMA_VERSION,
        "kind": "region_evaluation",
        "reference_sha256": run.load_meta()["source_sha256"],
        "critical_regions": critical_count,
        "strict_pass": (
            critical_count > 0
            and critical_pass
            and contract_audit["pass"]
            and not expectation_blockers
        ),
        "blockers": ([] if critical_count else ["regions:no-critical-regions"])
        + [f"region:{item['id']}" for item in results if item["critical"] and not item["pass"]]
        + contract_audit["blockers"]
        + expectation_blockers,
        "contract_audit": contract_audit,
        "critical_region_expectation": expectation_audit,
        "regions": results,
    }
    write_json(run.qa_dir / "regions-report.json", report)
    return report


def normalized_distance_limit(width: int, height: int, fraction: float) -> float:
    """Convert a canvas-diagonal fraction to pixels for arrow acceptance tests."""
    return math.hypot(width, height) * fraction
