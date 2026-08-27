"""convert 合同测试:atomic-vector 分支(vtracer 描摹片段 → 原生 freeform group)。

全部 case-neutral:合成最小 SVG 片段与 assets.json 条目,走完整 shadow convert。
"""

from __future__ import annotations

import hashlib
import zipfile
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from tools import common
from tools.core.contracts import read_json, write_json
from tools.convert import convert

EMU_PER_PX = 9525


@pytest.fixture()
def run_factory(tmp_path: Path):
    def make(svg: str, size: tuple[int, int] = (200, 100)) -> common.Run:
        source = tmp_path / "ref.png"
        Image.new("RGB", size, (240, 240, 240)).save(source)
        run = common.create_run(
            source,
            case="case",
            cases_root=tmp_path / "examples",
            input_route="svg-seeded",
        )
        run.qa_dir.mkdir(exist_ok=True)
        run.redraw_svg.write_text(svg, encoding="utf-8")
        return run

    return make


def _raster_entry(asset_id: str) -> dict:
    return {
        "id": asset_id,
        "authorized": True,
        "authorization_basis": "User explicitly authorized tight crops from this case's own reference PNG.",
        "rights_status": "unknown; authorization records workflow permission, not third-party copyright clearance",
        "editable": False,
        "raster_reason": "Creative microasset whose native redraw materially reduces reference fidelity.",
        "decomposition_note": "Only the microasset is raster; labels and arrows remain native editable objects.",
        "source": "reference_crop",
        "source_sha256": "0" * 64,
        "bbox": [20, 20, 40, 30],
        "source_tightly_cropped": True,
        "atomic_raster_unit": True,
        "contains_reconstructable_content": False,
    }


def _vector_entry(asset_id: str, rel_path: str, sha256: str, fallback_id: str) -> dict:
    return {
        "id": asset_id,
        "editable": True,
        "source": "vtracer-trace",
        "vector_source_svg": {"path": rel_path, "sha256": sha256},
        "trace_method": "vtracer-color-stacked-spline",
        "trace_engine_version": "0.6.15",
        "authorization_basis": "User explicitly authorized tight crops from this case's own reference PNG.",
        "rights_status": "unknown; authorization records workflow permission, not third-party copyright clearance",
        "fallback_atomic_raster": fallback_id,
        "ink_contract_region_id": "region-icon-a",
        "trace_eligibility": "flat-illustration",
    }


def _write_fragment(run: common.Run, rel_path: str, body: str) -> str:
    target = run.root / rel_path
    target.parent.mkdir(parents=True, exist_ok=True)
    target.write_text(body, encoding="utf-8")
    return hashlib.sha256(target.read_bytes()).hexdigest()


def _install_vector_asset(
    run: common.Run,
    asset_id: str = "atomic:icon-a",
    rel_path: str = "vector-sources/icon-a.svg",
    fragment: str | None = None,
    with_fallback: bool = True,
) -> str:
    fragment = fragment if fragment is not None else (
        '<svg xmlns="http://www.w3.org/2000/svg" width="40" height="30" viewBox="0 0 40 30">'
        '<path d="M0 0 L40 0 L40 30 L0 30 Z" fill="#112233" transform="translate(0,0)"/>'
        '<g transform="translate(10,5)">'
        '<path d="M0 0 C6 0 12 4 12 10 C12 16 6 20 0 20 Z" fill="#445566"/>'
        "</g></svg>"
    )
    sha256 = _write_fragment(run, rel_path, fragment)
    fallback_id = f"{asset_id}-raster"
    assets = read_json(run.assets_path)
    entries = [_vector_entry(asset_id, rel_path, sha256, fallback_id)]
    if with_fallback:
        entries.append(_raster_entry(fallback_id))
    assets["assets"] = entries
    write_json(run.assets_path, assets)
    return sha256


def _all_shapes(run: common.Run):
    def walk(shapes):
        for shape in shapes:
            yield shape
            children = getattr(shape, "shapes", None)
            if children is not None:
                yield from walk(children)

    return list(walk(Presentation(run.pptx_path).slides[0].shapes))


ATOMIC_SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" '
    'viewBox="0 0 200 100"><rect id="atomic:icon-a" x="20" y="20" '
    'width="40" height="30"/></svg>'
)


def test_atomic_vector_compiles_to_native_freeform_group(run_factory):
    run = run_factory(ATOMIC_SVG)
    sha256 = _install_vector_asset(run)
    convert(run)

    (group,) = Presentation(run.pptx_path).slides[0].shapes
    assert group.shape_type == MSO_SHAPE_TYPE.GROUP
    members = list(group.shapes)
    assert len(members) == 2
    assert all(member.shape_type == MSO_SHAPE_TYPE.FREEFORM for member in members)
    assert not any(
        shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in _all_shapes(run)
    )

    # viewBox 0 0 40 30 与占位符 bbox (20,20,40,30) 1:1,group 即占位框;
    # 第二个成员带 g transform=translate(10,5),落到 (30,25,12,20)。
    assert (group.left, group.top) == (20 * EMU_PER_PX, 20 * EMU_PER_PX)
    assert (group.width, group.height) == (40 * EMU_PER_PX, 30 * EMU_PER_PX)
    assert (members[1].left, members[1].top) == (30 * EMU_PER_PX, 25 * EMU_PER_PX)
    assert (members[1].width, members[1].height) == (12 * EMU_PER_PX, 20 * EMU_PER_PX)
    first_fill = members[0]._element.spPr.find(
        f"{qn('a:solidFill')}/{qn('a:srgbClr')}"
    )
    assert first_fill is not None and first_fill.get("val") == "112233"

    bindings = read_json(run.bindings_path)["bindings"]
    (binding,) = [b for b in bindings if b["element_id"] == "atomic:icon-a"]
    assert binding["object_kind"] == "atomic-vector"
    assert binding["editable"] is True
    assert not any(b["object_kind"] == "atomic-raster" for b in bindings)

    entries = {
        item["id"]: item for item in read_json(run.assets_path)["assets"]
    }
    vector = entries["atomic:icon-a"]
    assert set(vector) == {
        "id",
        "editable",
        "source",
        "vector_source_svg",
        "trace_method",
        "trace_engine_version",
        "authorization_basis",
        "rights_status",
        "fallback_atomic_raster",
        "ink_contract_region_id",
        "trace_eligibility",
    }
    assert vector["editable"] is True
    assert vector["source"] == "vtracer-trace"
    assert vector["vector_source_svg"]["sha256"] == sha256
    # 回退位图条目原样保留,不被矢量编译改写。
    assert entries["atomic:icon-a-raster"]["source"] == "reference_crop"
    assert entries["atomic:icon-a-raster"]["editable"] is False

    with zipfile.ZipFile(run.pptx_path) as package:
        tag_parts = [name for name in package.namelist() if name.startswith("ppt/tags/tag")]
        assert len(tag_parts) == 1
        tags = package.read(tag_parts[0]).decode("utf-8")
        slide = package.read("ppt/slides/slide1.xml").decode("utf-8")
    assert 'val="atomic:icon-a"' in tags
    assert f'val="{sha256}"' in tags
    assert 'val="atomic-vector"' in tags
    assert 'val="vtracer-provider"' in tags
    assert "AISCIENTIFICILLUSTRATOREDITABLE" in tags
    assert 'val="vtracer-color-stacked-spline"' in tags
    assert "custDataLst" in slide and "p:tags" in slide


def test_atomic_vector_maps_viewbox_onto_element_bbox(run_factory):
    run = run_factory(ATOMIC_SVG)
    _install_vector_asset(
        run,
        fragment=(
            '<svg xmlns="http://www.w3.org/2000/svg" width="10" height="10" '
            'viewBox="0 0 10 10">'
            '<path d="M0 0 L10 0 L10 10 L0 10 Z" fill="#AABBCC"/></svg>'
        ),
    )
    convert(run)
    (group,) = Presentation(run.pptx_path).slides[0].shapes
    # viewBox 10x10 非均匀映射到 40x30 占位框,成员铺满整个 bbox。
    (member,) = list(group.shapes)
    assert (member.left, member.top) == (20 * EMU_PER_PX, 20 * EMU_PER_PX)
    assert (member.width, member.height) == (40 * EMU_PER_PX, 30 * EMU_PER_PX)


def test_atomic_vector_hash_drift_fails_closed(run_factory):
    run = run_factory(ATOMIC_SVG)
    _install_vector_asset(run)
    fragment = run.root / "vector-sources" / "icon-a.svg"
    fragment.write_text(fragment.read_text(encoding="utf-8") + "\n", encoding="utf-8")
    with pytest.raises(SystemExit, match="SHA-256"):
        convert(run)
    assert not run.pptx_path.exists()


def test_atomic_vector_missing_fragment_fails_closed(run_factory):
    run = run_factory(ATOMIC_SVG)
    _install_vector_asset(run)
    (run.root / "vector-sources" / "icon-a.svg").unlink()
    with pytest.raises(SystemExit, match="不存在"):
        convert(run)
    assert not run.pptx_path.exists()


def test_atomic_vector_contract_subset_violation_fails_closed(run_factory):
    run = run_factory(ATOMIC_SVG)
    _install_vector_asset(
        run,
        fragment=(
            '<svg xmlns="http://www.w3.org/2000/svg" width="40" height="30" '
            'viewBox="0 0 40 30">'
            '<image href="icon.png" x="0" y="0" width="40" height="30"/></svg>'
        ),
    )
    with pytest.raises(SystemExit, match="合同子集"):
        convert(run)
    assert not run.pptx_path.exists()


def test_atomic_placeholder_without_vector_entry_stays_raster(run_factory):
    run = run_factory(ATOMIC_SVG)
    assets = read_json(run.assets_path)
    assets["assets"] = [_raster_entry("atomic:icon-a")]
    write_json(run.assets_path, assets)
    convert(run)

    (shape,) = Presentation(run.pptx_path).slides[0].shapes
    assert shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    bindings = read_json(run.bindings_path)["bindings"]
    (binding,) = [b for b in bindings if b["element_id"] == "atomic:icon-a"]
    assert binding["object_kind"] == "atomic-raster"
    assert binding["editable"] is False
    entry = read_json(run.assets_path)["assets"][0]
    assert entry["source"] == "reference_crop"
    assert entry["editable"] is False


LINKED_VECTOR_REL_PATH = "vector-sources/icon-a-vector.svg"


def _install_linked_assets(run: common.Run) -> str:
    """id 模型链接形态:元素/位图条目 atomic:icon-a,矢量条目 atomic:icon-a-vector。"""
    fragment = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="40" height="30" '
        'viewBox="0 0 40 30">'
        '<path d="M0 0 L40 0 L40 30 L0 30 Z" fill="#112233"/>'
        '<path d="M10 5 C16 5 22 9 22 15 C22 21 16 25 10 25 Z" fill="#445566"/>'
        "</svg>"
    )
    sha256 = _write_fragment(run, LINKED_VECTOR_REL_PATH, fragment)
    assets = read_json(run.assets_path)
    assets["assets"] = [
        _raster_entry("atomic:icon-a"),
        _vector_entry("atomic:icon-a-vector", LINKED_VECTOR_REL_PATH, sha256, "atomic:icon-a"),
    ]
    write_json(run.assets_path, assets)
    return sha256


def test_atomic_vector_linked_via_fallback_atomic_raster(run_factory):
    run = run_factory(ATOMIC_SVG)
    sha256 = _install_linked_assets(run)
    convert(run)

    (group,) = Presentation(run.pptx_path).slides[0].shapes
    assert group.shape_type == MSO_SHAPE_TYPE.GROUP
    assert len(list(group.shapes)) == 2
    assert all(member.shape_type == MSO_SHAPE_TYPE.FREEFORM for member in group.shapes)
    assert not any(
        shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in _all_shapes(run)
    )

    # 绑定行以场景元素 id(位图条目 id)登记,kind 为 atomic-vector。
    bindings = read_json(run.bindings_path)["bindings"]
    (binding,) = [b for b in bindings if b["element_id"] == "atomic:icon-a"]
    assert binding["object_kind"] == "atomic-vector"
    assert binding["editable"] is True
    assert not any(b["object_kind"] == "atomic-raster" for b in bindings)

    entries = {item["id"]: item for item in read_json(run.assets_path)["assets"]}
    vector = entries["atomic:icon-a-vector"]
    assert vector["source"] == "vtracer-trace"
    assert vector["editable"] is True
    assert vector["fallback_atomic_raster"] == "atomic:icon-a"
    assert vector["vector_source_svg"]["sha256"] == sha256
    assert entries["atomic:icon-a"]["source"] == "reference_crop"
    assert entries["atomic:icon-a"]["editable"] is False

    with zipfile.ZipFile(run.pptx_path) as package:
        tag_parts = [name for name in package.namelist() if name.startswith("ppt/tags/tag")]
        assert len(tag_parts) == 1
        tags = package.read(tag_parts[0]).decode("utf-8")
    assert 'val="atomic:icon-a-vector"' in tags
    assert f'val="{sha256}"' in tags
    assert 'val="atomic-vector"' in tags
    assert 'val="vtracer-provider"' in tags


def test_atomic_vector_ambiguous_linkage_fails_closed(run_factory):
    run = run_factory(ATOMIC_SVG)
    fragment = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="40" height="30" '
        'viewBox="0 0 40 30"><path d="M0 0 L40 0 L40 30 L0 30 Z" fill="#112233"/></svg>'
    )
    exact_sha256 = _write_fragment(run, "vector-sources/icon-a.svg", fragment)
    linked_sha256 = _write_fragment(run, LINKED_VECTOR_REL_PATH, fragment)
    assets = read_json(run.assets_path)
    assets["assets"] = [
        _raster_entry("atomic:icon-a-raster"),
        # 精确命中:条目 id == 元素 id。
        _vector_entry("atomic:icon-a", "vector-sources/icon-a.svg", exact_sha256, "atomic:icon-a-raster"),
        # 链接命中:fallback_atomic_raster == 元素 id;与精确命中不同条目,歧义。
        _vector_entry("atomic:icon-a-vector", LINKED_VECTOR_REL_PATH, linked_sha256, "atomic:icon-a"),
    ]
    write_json(run.assets_path, assets)
    with pytest.raises(SystemExit, match="歧义"):
        convert(run)
    assert not run.pptx_path.exists()
