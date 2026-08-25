from __future__ import annotations

import hashlib
import io
import json
import xml.etree.ElementTree as ET
from pathlib import Path

from PIL import Image, ImageChops
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE

from tools import common
from tools.layout import audit_layout


CASE = (
    Path(__file__).resolve().parents[1]
    / "examples"
    / "svg-seeded"
    / "01-modular-agent"
)
REFERENCE_ONLY_CASE = (
    Path(__file__).resolve().parents[1]
    / "examples"
    / "reference-only"
    / "01-modular-agent-reference-only"
)


def test_case01_mapping_arrow_is_above_imagination_panel():
    root = ET.parse(CASE / "redraw.svg").getroot()
    direct_ids = [element.get("id") for element in root]
    assert direct_ids.index("mapping-to-imagination") > direct_ids.index(
        "task-conditioned-imagination"
    )

    shapes = list(Presentation(CASE / "redraw.pptx").slides[0].shapes)
    z_order = {shape.name: index for index, shape in enumerate(shapes)}
    assert z_order["af-mapping-to-imagination-connector-01"] > z_order[
        "af-task-conditioned-imagination-rect-01"
    ]


def test_case01_globe_is_exact_authorized_reference_crop():
    assets = json.loads((CASE / "assets.json").read_text(encoding="utf-8"))
    asset = next(
        item for item in assets["assets"] if item["id"] == "atomic:environment-globe"
    )
    assert asset["authorized"] is True
    assert asset["editable"] is False
    assert asset["source"] == "reference_crop"

    presentation = Presentation(CASE / "redraw.pptx")
    picture = next(
        shape
        for shape in presentation.slides[0].shapes
        if shape.name == "af-atomic-environment-globe-atomic-raster-01"
    )
    assert picture.shape_type == MSO_SHAPE_TYPE.PICTURE
    embedded = picture.image.blob
    assert hashlib.sha256(embedded).hexdigest() == asset["source_sha256"]

    x, y, width, height = asset["bbox"]
    with Image.open(CASE / "reference.png") as reference:
        expected_image = reference.crop((x, y, x + width, y + height)).convert("RGBA")
        expected_buffer = io.BytesIO()
        expected_image.save(expected_buffer, format="PNG")
    with Image.open(io.BytesIO(embedded)) as actual:
        difference = ImageChops.difference(actual.convert("RGBA"), expected_image)
    assert difference.getbbox() is None


def test_case01_mapping_and_repeated_vectors_pass_explicit_layout_contracts():
    report = audit_layout(common.open_run(CASE))
    assert report["findings"] == []
    assert report["pass"] is True

    containment = {item["element"]: item for item in report["containment"]}
    for element_id in ("task-mapping-label", "task-mapping-formula"):
        row = containment[element_id]
        assert max(row["backend"]["overflow_px"].values()) <= row["tolerance_px"]

    groups = {item["id"]: item for item in report["repeat_groups"]}
    assert groups["e-v-stack"]["source"]["steps_px"] == [38.0, 37.0]
    assert groups["e-v-stack"]["backend"]["steps_px"] == [38.0, 37.0]
    assert groups["s-t-stack"]["source"]["steps_px"] == [38.0, 37.0]
    assert groups["s-t-stack"]["backend"]["steps_px"] == [38.0, 37.0]


def test_both_case01_routes_have_reference_bound_horizontal_gap_arrows():
    for case in (CASE, REFERENCE_ONLY_CASE):
        report = audit_layout(common.open_run(case))
        rows = {item["element"]: item for item in report["gap_arrows"]}
        assert set(rows) == {
            "allocator-task1-route-1",
            "allocator-task1-route-2",
            "allocator-task1-skip-edge",
        }
        assert [rows[key]["source"]["actual_start"] for key in rows] == [
            [604.0, 347.0],
            [680.0, 347.0],
            [758.0, 347.0],
        ]
        assert [rows[key]["source"]["actual_end"] for key in rows] == [
            [628.0, 347.0],
            [706.0, 347.0],
            [801.0, 347.0],
        ]
        assert all(
            item["backend"]["topology_matches_gap_contract"]
            for item in rows.values()
        )
        assert all(
            item["backend"]["endpoint_metrics"]["start_delta"]["distance"] == 0.0
            and item["backend"]["endpoint_metrics"]["end_delta"]["distance"] == 0.0
            for item in rows.values()
        )
        assert all(item["backend"]["path_and_head_readback_pass"] for item in rows.values())


def test_both_case01_routes_preserve_rotated_word_semantics():
    for case in (CASE, REFERENCE_ONLY_CASE):
        report = audit_layout(common.open_run(case))
        flows = {item["element"]: item for item in report["text_flow_contracts"]}
        assert len(flows) == 8
        assert {item["mode"] for item in flows.values()} == {"rotated-word"}
        assert {item["backend"]["rotation_deg"] for item in flows.values()} == {90.0}
        assert all(len(item["backend"]["paragraphs"]) == 1 for item in flows.values())
        assert all(item["backend"]["contract_pass"] for item in flows.values())


def test_reference_only_case01_encoder_pair_and_feedback_arrow_are_native():
    report = audit_layout(common.open_run(REFERENCE_ONLY_CASE))
    peers = {item["id"]: item for item in report["peer_size_groups"]}
    encoder_pair = peers["encoder-pair"]
    assert encoder_pair["source"]["heights_px"] == [78.0, 78.0]
    assert encoder_pair["backend"]["heights_px"] == [78.0, 78.0]

    scene = json.loads((REFERENCE_ONLY_CASE / "scene.json").read_text(encoding="utf-8"))
    arrow = next(
        item["arrow_spec"]
        for item in scene["edges"]
        if item["id"] == "joint-loop-arrow"
    )
    assert arrow["representation"] == "line_arrow"
    assert arrow["body"]["width_px"] == 10.0
    assert arrow["body"]["line_cap"] == "round"
    assert arrow["body"]["line_join"] == "round"
    assert arrow["end_head"]["type"] == "triangle"
    assert arrow["topology"]["target_id"] == "joint-panel"
    assert arrow["single_visible_object"] is True


def test_reference_only_case01_has_three_native_upward_reward_arrows():
    expected = {
        "rollout-z0-to-r0": ("rollout-z0", "reward-r0", [1033.0, 296.0], [1033.0, 269.0]),
        "rollout-z1-to-r1": ("rollout-z1", "reward-r1", [1129.0, 296.0], [1129.0, 268.0]),
        "rollout-z2-to-r2": ("rollout-z2", "reward-r2", [1229.0, 296.0], [1229.0, 268.0]),
    }
    scene = json.loads((REFERENCE_ONLY_CASE / "scene.json").read_text(encoding="utf-8"))
    edges = {item["id"]: item for item in scene["edges"]}
    for element_id, (source, target, start, end) in expected.items():
        edge = edges[element_id]
        spec = edge["arrow_spec"]
        assert edge["source"] == source
        assert edge["target"] == target
        assert spec["path"]["points"] == [
            {"x": start[0], "y": start[1]},
            {"x": end[0], "y": end[1]},
        ]
        assert end[1] < start[1]
        assert spec["routing"] == "host"
        assert spec["topology"] == {
            "mode": "attached",
            "source_id": source,
            "target_id": target,
            "source_site": 0,
            "target_site": 2,
        }
        assert spec["body"]["color"] == "#7E6000"
        assert spec["body"]["width_px"] == 2.0
        assert spec["end_head"]["type"] == "triangle"
        assert spec["end_head"]["width"] == "med"
        assert spec["end_head"]["length"] == "med"
        assert spec["fallback_policy"] == "strict_fail"
        assert spec["single_visible_object"] is True

    readback = json.loads(
        (REFERENCE_ONLY_CASE / "qa" / "powerpoint-arrow-readback.json").read_text(
            encoding="utf-8"
        )
    )
    records = {item["element_id"]: item for item in readback["records"]}
    for element_id in expected:
        record = records[element_id]
        assert record["ooxml_kind"] == "connector"
        assert record["end_head"]["pass"] is True
        assert record["path_geometry"]["status"] == "PASS"
        assert record["topology"]["pass"] is True
        assert record["status"] == "PASS"


def test_reference_only_case01_interaction_is_one_native_bidirectional_block_arrow():
    svg_text = (REFERENCE_ONLY_CASE / "redraw.svg").read_text(encoding="utf-8")
    assert 'id="interaction-left"' not in svg_text
    assert 'id="interaction-right"' not in svg_text
    root = ET.fromstring(svg_text)
    glyphs = [
        element
        for element in root.iter()
        if element.get("id") == "interaction-exchange"
    ]
    assert len(glyphs) == 1
    glyph = glyphs[0]
    assert glyph.tag.endswith("polygon")
    assert glyph.get("data-arrow-representation") == "block_arrow"
    assert glyph.get("data-ppt-autoshape") == "leftRightArrow"
    assert glyph.get("data-start-head-type") == "triangle"
    assert glyph.get("data-end-head-type") == "triangle"

    scene = json.loads((REFERENCE_ONLY_CASE / "scene.json").read_text(encoding="utf-8"))
    edges = [item for item in scene["edges"] if item["id"] == "interaction-exchange"]
    assert len(edges) == 1
    spec = edges[0]["arrow_spec"]
    assert spec["representation"] == "block_arrow"
    assert spec["single_visible_object"] is True
    assert spec["autoshape"] == {
        "subtype": "leftRightArrow",
        "adjustments": [0.529411765, 0.470588235],
        "bbox": [1014.0, 563.5, 160.0, 17.0],
    }
    assert spec["start_head"]["type"] == "triangle"
    assert spec["end_head"]["type"] == "triangle"
    assert spec["path"] == {
        "kind": "straight",
        "coordinate_space": "canvas",
        "points": [
            {"x": 1014.0, "y": 572.0},
            {"x": 1174.0, "y": 572.0},
        ],
    }

    bindings = json.loads(
        (REFERENCE_ONLY_CASE / "bindings.json").read_text(encoding="utf-8")
    )
    bound = [
        item
        for item in bindings["bindings"]
        if item["element_id"] == "interaction-exchange"
    ]
    assert len(bound) == 1
    assert bound[0]["object_kind"] == "block-arrow-autoshape"
    assert bound[0]["single_visible_object"] is True

    presentation = Presentation(REFERENCE_ONLY_CASE / "redraw.pptx")
    shapes = [
        shape
        for shape in presentation.slides[0].shapes
        if shape.name == "af-interaction-exchange-block-arrow-autoshape-01"
    ]
    assert len(shapes) == 1
    assert shapes[0].shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    assert shapes[0].auto_shape_type == MSO_AUTO_SHAPE_TYPE.LEFT_RIGHT_ARROW

    readback = json.loads(
        (REFERENCE_ONLY_CASE / "qa" / "powerpoint-arrow-readback.json").read_text(
            encoding="utf-8"
        )
    )
    record = next(
        item for item in readback["records"] if item["element_id"] == "interaction-exchange"
    )
    assert record["autoshape"]["pass"] is True
    assert record["path_geometry"]["status"] == "PASS"
    assert record["status"] == "PASS"

    composition = json.loads(
        (REFERENCE_ONLY_CASE / "qa" / "arrow-composition-audit.json").read_text(
            encoding="utf-8"
        )
    )
    assert composition["artifact_sha256"] == common.sha256_file(
        REFERENCE_ONLY_CASE / "redraw.pptx"
    )
    assert composition["source_findings"] == []
    assert composition["readback_findings"] == []
    assert composition["pass"] is True

    regions = json.loads(
        (REFERENCE_ONLY_CASE / "qa" / "regions-report.json").read_text(
            encoding="utf-8"
        )
    )
    region = next(
        item
        for item in regions["regions"]
        if item["id"] == "interaction-exchange-block-arrow"
    )
    assert region["ssim"] >= 0.9
    assert region["edge_iou"] >= 0.85
    assert region["pass"] is True


def test_svg_seeded_case01_mixture_arrows_are_above_panel_and_below_target():
    report = audit_layout(common.open_run(CASE))
    rows = {item["element"]: item for item in report["z_order_contracts"]}
    for element_id in ("top-mixture-to-rollout", "bottom-mixture-to-rollout"):
        row = rows[element_id]
        assert row["above"] == ["rollout-panel"]
        assert row["below"] == ["rollout-z0"]
        assert row["source_pass"] is True
        assert row["backend_pass"] is True
