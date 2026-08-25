from __future__ import annotations

import copy
from pathlib import Path

import numpy as np
import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu

from tools import common
from tools.contracts import read_json, write_json
from tools.regions import (
    _bbox_identity_metrics,
    _edge_iou,
    _obstacle_mask_floor,
    build_critical_region_expectation,
    evaluate_regions,
)


def test_obstacle_mask_floor_is_renderer_aware_only_for_sparse_open_strokes():
    thin_floor, thin_fill, thin = _obstacle_mask_floor(
        {"bbox": [0, 0, 108, 16], "area_px": 241}
    )
    solid_floor, solid_fill, solid = _obstacle_mask_floor(
        {"bbox": [0, 0, 17, 3], "area_px": 19}
    )

    assert thin is True
    assert thin_fill <= 0.20
    assert thin_floor == 0.55
    assert solid is False
    assert solid_fill > 0.20
    assert solid_floor == 0.65
    assert _obstacle_mask_floor(
        {"bbox": [0, 0, 108, 16], "area_px": 241}, 0.75
    )[0] == 0.75


def _bind_synthetic_elements(run: common.Run, element_ids: list[str]) -> None:
    scene = read_json(run.scene_path)
    scene["elements"] = [
        {"id": element_id, "kind": "shape"} for element_id in element_ids
    ]
    scene["edges"] = []
    write_json(run.scene_path, scene)
    bindings = read_json(run.bindings_path)
    bindings["bindings"] = [
        {
            "element_id": element_id,
            "shape_id": index,
            "shape_name": f"af-{element_id}",
            "readback_found": True,
        }
        for index, element_id in enumerate(element_ids, start=1)
    ]
    write_json(run.bindings_path, bindings)


def _freeze(run: common.Run, regions: dict, element_ids: str | list[str]) -> None:
    frozen_ids = [element_ids] if isinstance(element_ids, str) else element_ids
    _bind_synthetic_elements(run, frozen_ids)
    for region in regions["regions"]:
        if region.get("critical") is True:
            region["element_ids"] = frozen_ids
    regions["critical_region_expectation"] = build_critical_region_expectation(regions)


def _run(tmp_path: Path) -> common.Run:
    reference = tmp_path / "reference-source.png"
    Image.new("RGB", (30, 20), (240, 120, 80)).save(reference)
    run = common.create_run(
        reference,
        case="case",
        cases_root=tmp_path / "examples",
        input_route="svg-seeded",
    )
    run.qa_dir.mkdir(exist_ok=True)
    Image.open(run.source_png).save(run.render_png)
    regions = read_json(run.regions_path)
    regions["regions"] = [
        {
            "id": "critical",
            "bbox": [5, 5, 10, 10],
            "critical": True,
            "color_probes": [
                {"id": "center", "point": [10, 10], "radius": 1, "max_delta_e": 5}
            ],
        },
        {"id": "whole", "bbox": [0, 0, 30, 20], "critical": False},
    ]
    _freeze(run, regions, "synthetic-critical")
    write_json(run.regions_path, regions)
    return run


def test_identical_critical_region_passes(tmp_path: Path):
    run = _run(tmp_path)
    report = evaluate_regions(run)
    assert report["strict_pass"] is True
    assert report["blockers"] == []
    assert report["regions"][0]["color_probes"][0]["delta_e00"] == 0.0


def test_local_failure_blocks_even_when_whole_canvas_is_diagnostic(tmp_path: Path):
    run = _run(tmp_path)
    render = Image.open(run.render_png)
    for y in range(5, 15):
        for x in range(5, 15):
            render.putpixel((x, y), (0, 0, 0))
    render.save(run.render_png)
    report = evaluate_regions(run)
    assert report["strict_pass"] is False
    assert report["blockers"] == ["region:critical"]


def test_frozen_pixel_bbox_separates_pixel_roi_from_object_analysis_bbox(
    tmp_path: Path,
):
    run = _run(tmp_path)
    regions = read_json(run.regions_path)
    region = regions["regions"][0]
    region.pop("color_probes")
    region["pixel_bbox"] = [5, 5, 5, 5]
    regions["critical_region_expectation"] = build_critical_region_expectation(
        regions
    )
    write_json(run.regions_path, regions)
    render = Image.open(run.render_png)
    for y in range(10, 15):
        for x in range(10, 15):
            render.putpixel((x, y), (0, 0, 0))
    render.save(run.render_png)

    report = evaluate_regions(run)

    assert report["strict_pass"] is True
    assert report["regions"][0]["bbox"] == [5, 5, 10, 10]
    assert report["regions"][0]["pixel_bbox"] == [5, 5, 5, 5]
    assert report["regions"][0]["ssim"] == 1.0


def test_frozen_pixel_bbox_tamper_is_an_expectation_blocker(tmp_path: Path):
    run = _run(tmp_path)
    regions = read_json(run.regions_path)
    region = regions["regions"][0]
    region["pixel_bbox"] = [5, 5, 5, 5]
    regions["critical_region_expectation"] = build_critical_region_expectation(
        regions
    )
    region["pixel_bbox"] = [5, 5, 6, 5]
    write_json(run.regions_path, regions)

    report = evaluate_regions(run)

    assert report["strict_pass"] is False
    assert (
        "regions:expectation:critical:pixel-bbox-mismatch"
        in report["blockers"]
    )


def test_no_critical_region_is_an_explicit_strict_blocker(tmp_path: Path):
    reference = tmp_path / "no-critical.png"
    Image.new("RGB", (30, 20), "white").save(reference)
    run = common.create_run(
        reference,
        case="no-critical",
        cases_root=tmp_path / "examples",
        input_route="svg-seeded",
    )
    Image.open(run.source_png).save(run.render_png)
    report = evaluate_regions(run)
    assert report["strict_pass"] is False
    assert report["blockers"] == ["regions:no-critical-regions"]


def _ink_run(tmp_path: Path) -> common.Run:
    reference = tmp_path / "ink-reference.png"
    image = Image.new("RGB", (40, 30), "white")
    for y in range(8, 18):
        for x in range(10, 20):
            image.putpixel((x, y), (40, 40, 40))
    image.save(reference)
    run = common.create_run(
        reference,
        case="ink",
        cases_root=tmp_path / "examples",
        input_route="svg-seeded",
    )
    Image.open(run.source_png).save(run.render_png)
    regions = read_json(run.regions_path)
    regions["regions"] = [
        {
            "id": "icon",
            "bbox": [5, 3, 25, 22],
            "critical": True,
            "thresholds": {"ssim_min": 0, "edge_iou_min": 0},
            "ink_contract": {
                "background_rgb": [255, 255, 255],
                "background_tolerance": 24,
                "bbox_tolerance_px": 1,
                "center_tolerance_px": 1,
                "area_relative_tolerance": 0.1,
            },
        }
    ]
    _freeze(run, regions, "synthetic-icon")
    write_json(run.regions_path, regions)
    return run


def test_identical_tight_foreground_contract_passes(tmp_path: Path):
    run = _ink_run(tmp_path)
    report = evaluate_regions(run)
    contract = report["regions"][0]["ink_contract"]
    assert contract["pass"] is True
    assert contract["reference"]["bbox"] == [5, 5, 10, 10]
    assert report["strict_pass"] is True


def test_foreground_contract_blocks_undersized_shifted_icon(tmp_path: Path):
    run = _ink_run(tmp_path)
    render = Image.new("RGB", (40, 30), "white")
    for y in range(9, 16):
        for x in range(12, 19):
            render.putpixel((x, y), (40, 40, 40))
    render.save(run.render_png)

    report = evaluate_regions(run)

    contract = report["regions"][0]["ink_contract"]
    assert contract["pass"] is False
    assert contract["render"]["bbox"] == [7, 6, 7, 7]
    assert contract["area_relative_error"] == 0.51
    assert report["blockers"] == ["region:icon"]


def _clearance_run(tmp_path: Path) -> common.Run:
    reference = tmp_path / "clearance-reference.png"
    image = Image.new("RGB", (60, 30), "white")
    for y in range(10, 20):
        for x in range(10, 20):
            image.putpixel((x, y), (240, 0, 0))
        for x in range(25, 30):
            image.putpixel((x, y), (40, 40, 40))
    image.save(reference)
    run = common.create_run(
        reference,
        case="clearance",
        cases_root=tmp_path / "examples",
        input_route="svg-seeded",
    )
    Image.open(run.source_png).save(run.render_png)
    regions = read_json(run.regions_path)
    regions["regions"] = [
        {
            "id": "caption-clearance",
            "bbox": [5, 5, 35, 20],
            "critical": True,
            "thresholds": {"ssim_min": 0, "edge_iou_min": 0},
            "color_clearance_contracts": [
                {
                    "id": "caption-to-brace",
                    "subject_element_ids": ["synthetic-caption"],
                    "subject_dominant_channel": "red",
                    "minimum_dominance": 20,
                    "background_rgb": [255, 255, 255],
                    "background_tolerance": 24,
                    "absolute_min_px": 1,
                    "reference_loss_tolerance_px": 1,
                    "subject_bbox": [4, 4, 16, 12],
                    "native_binding_bbox": [4, 4, 11, 12],
                    "obstacles": [
                        {
                            "id": "brace",
                            "element_ids": ["synthetic-brace"],
                            "bbox": [15, 4, 12, 12],
                            "native_binding_bbox": [20, 4, 5, 12],
                            "area_relative_tolerance": 0.1,
                            "bbox_tolerance_px": 2,
                        }
                    ],
                }
            ],
        }
    ]
    _freeze(run, regions, ["synthetic-caption", "synthetic-brace"])
    presentation = Presentation()
    presentation.slide_width = Emu(60 * 91440)
    presentation.slide_height = Emu(30 * 91440)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shapes = {}
    for element_id, bbox in (
        ("synthetic-caption", [9, 9, 11, 12]),
        ("synthetic-brace", [25, 9, 5, 12]),
    ):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(bbox[0] * 91440),
            Emu(bbox[1] * 91440),
            Emu(bbox[2] * 91440),
            Emu(bbox[3] * 91440),
        )
        shape.name = f"af-{element_id}"
        shapes[element_id] = shape
    presentation.save(run.pptx_path)
    bindings = read_json(run.bindings_path)
    for binding in bindings["bindings"]:
        shape = shapes[binding["element_id"]]
        binding["shape_id"] = shape.shape_id
        binding["shape_name"] = shape.name
    write_json(run.bindings_path, bindings)
    write_json(run.regions_path, regions)
    return run


def test_reference_derived_color_clearance_passes_identical_render(tmp_path: Path):
    run = _clearance_run(tmp_path)

    report = evaluate_regions(run)

    result = report["regions"][0]["color_clearance_contracts"][0]
    assert report["regions"][0]["thresholds"] == {
        "ssim_min": 0.85,
        "edge_iou_min": 0.75,
    }
    assert result["reference_clearance_px"] == 5.0
    assert result["render_clearance_px"] == 5.0
    assert result["pass"] is True


def test_color_clearance_blocks_caption_that_moves_too_close_to_obstacle(tmp_path: Path):
    run = _clearance_run(tmp_path)
    render = Image.open(run.render_png)
    for y in range(10, 20):
        for x in range(25, 30):
            render.putpixel((x, y), (255, 255, 255))
        for x in range(21, 26):
            render.putpixel((x, y), (40, 40, 40))
    render.save(run.render_png)

    report = evaluate_regions(run)

    result = report["regions"][0]["color_clearance_contracts"][0]
    assert result["reference_clearance_px"] == 5.0
    assert result["render_clearance_px"] == 1.0
    assert result["required_clearance_px"] == 4.0
    assert result["pass"] is False
    assert report["blockers"] == ["region:caption-clearance"]


def test_color_clearance_blocks_when_the_named_obstacle_disappears(tmp_path: Path):
    run = _clearance_run(tmp_path)
    render = Image.open(run.render_png)
    for y in range(10, 20):
        for x in range(25, 30):
            render.putpixel((x, y), (255, 255, 255))
        for x in range(35, 40):
            render.putpixel((x, y), (40, 40, 40))
    render.save(run.render_png)

    report = evaluate_regions(run)

    result = report["regions"][0]["color_clearance_contracts"][0]
    assert result["pass"] is False
    assert "obstacle ink is missing" in result["obstacles"][0]["error"]


def test_color_clearance_blocks_obstacle_size_growth(tmp_path: Path):
    run = _clearance_run(tmp_path)
    render = Image.open(run.render_png)
    for y in range(10, 20):
        for x in range(30, 32):
            render.putpixel((x, y), (40, 40, 40))
    render.save(run.render_png)

    report = evaluate_regions(run)

    obstacle = report["regions"][0]["color_clearance_contracts"][0]["obstacles"][0]
    assert obstacle["area_relative_error"] == 0.4
    assert obstacle["area_relative_error"] > obstacle["area_relative_tolerance"]
    assert obstacle["pass"] is False


def test_color_clearance_rejects_subject_colored_pixels_impersonating_obstacle(
    tmp_path: Path,
):
    run = _clearance_run(tmp_path)
    render = Image.open(run.render_png)
    for y in range(18, 20):
        for x in range(25, 30):
            render.putpixel((x, y), (240, 0, 0))
    render.save(run.render_png)

    report = evaluate_regions(run)

    obstacle = report["regions"][0]["color_clearance_contracts"][0]["obstacles"][0]
    assert obstacle["pass"] is False
    assert obstacle["area_relative_error"] == 0.2
    assert obstacle["mask_iou"] < 1.0


def test_color_clearance_rejects_black_obstacle_recolored_red(tmp_path: Path):
    run = _clearance_run(tmp_path)
    render = Image.open(run.render_png)
    for y in range(10, 20):
        for x in range(25, 30):
            render.putpixel((x, y), (240, 0, 0))
    render.save(run.render_png)

    report = evaluate_regions(run)

    obstacle = report["regions"][0]["color_clearance_contracts"][0]["obstacles"][0]
    assert obstacle["pass"] is False
    assert obstacle["render"] is None
    assert obstacle["mask_iou"] == 0.0
    assert (
        obstacle["core_foreground_delta_e00"]
        > obstacle["core_foreground_delta_e_max"]
    )


def test_obstacle_gate_tolerates_antialiased_edges_without_rgb_white_dilution(
    tmp_path: Path,
):
    run = _clearance_run(tmp_path)
    render = Image.open(run.render_png)
    for y in range(10, 20):
        for x in range(25, 30):
            if x in {25, 29} or y in {10, 19}:
                render.putpixel((x, y), (180, 180, 180))
    render.save(run.render_png)

    report = evaluate_regions(run)

    obstacle = report["regions"][0]["color_clearance_contracts"][0]["obstacles"][0]
    assert obstacle["mean_abs_rgb_delta_diagnostic"] > 24.0
    assert obstacle["mask_iou"] == 1.0
    assert obstacle["edge_iou"] >= obstacle["edge_iou_min"]
    assert obstacle["core_foreground_delta_e00"] == 0.0
    assert obstacle["pass"] is True


def test_color_clearance_rejects_subject_ink_that_overflows_its_frozen_bbox(
    tmp_path: Path,
):
    run = _clearance_run(tmp_path)
    render = Image.open(run.render_png)
    for y in range(10, 20):
        for x in range(20, 30):
            render.putpixel((x, y), (240, 0, 0))
    render.save(run.render_png)

    report = evaluate_regions(run)

    subject = report["regions"][0]["color_clearance_contracts"][0]["subject"]
    assert subject["render_overflow_px"] > subject["reference_overflow_px"]
    assert subject["overflow_pass"] is False


def test_clearance_target_ids_must_resolve_to_the_declared_obstacle_bbox(
    tmp_path: Path,
):
    run = _clearance_run(tmp_path)
    presentation = Presentation(run.pptx_path)
    shape = next(
        item
        for item in presentation.slides[0].shapes
        if item.name == "af-synthetic-brace"
    )
    shape.left = Emu(45 * 91440)
    presentation.save(run.pptx_path)

    report = evaluate_regions(run)

    audit = report["regions"][0]["color_clearance_contracts"][0][
        "binding_bbox_audit"
    ]
    obstacle = next(item for item in audit["targets"] if item["role"] == "obstacle")
    assert obstacle["pass"] is False


def test_native_binding_bbox_rejects_centered_two_times_larger_box(tmp_path: Path):
    run = _clearance_run(tmp_path)
    presentation = Presentation(run.pptx_path)
    shape = next(
        item
        for item in presentation.slides[0].shapes
        if item.name == "af-synthetic-caption"
    )
    shape.left = Emu(4 * 91440)
    shape.top = Emu(3 * 91440)
    shape.width = Emu(22 * 91440)
    shape.height = Emu(24 * 91440)
    presentation.save(run.pptx_path)

    report = evaluate_regions(run)

    audit = report["regions"][0]["color_clearance_contracts"][0][
        "binding_bbox_audit"
    ]
    subject = next(item for item in audit["targets"] if item["role"] == "subject")
    assert subject["center_error_px"] <= subject["center_tolerance_px"]
    assert subject["actual_fully_inside_native_bbox_with_3px_halo"] is False
    assert subject["pass"] is False


def test_native_binding_bbox_rejects_four_times_longer_line():
    metrics = _bbox_identity_metrics(
        [10.0, 10.0, 0.0, 10.0],
        [10.0, -5.0, 0.0, 40.0],
    )

    assert metrics["length_ratio"] == 4.0
    assert metrics["length_ratio_max"] == 1.30
    assert metrics["line_endpoints_inside_native_bbox_with_3px_halo"] is False
    assert metrics["pass"] is False


@pytest.mark.parametrize("target", ["subject", "obstacle"])
def test_native_binding_bbox_is_required_for_every_clearance_target(
    tmp_path: Path, target: str
):
    run = _clearance_run(tmp_path)
    regions = read_json(run.regions_path)
    contract = regions["regions"][0]["color_clearance_contracts"][0]
    if target == "subject":
        del contract["native_binding_bbox"]
    else:
        del contract["obstacles"][0]["native_binding_bbox"]
    write_json(run.regions_path, regions)

    report = evaluate_regions(run)

    audit = report["regions"][0]["color_clearance_contracts"][0][
        "binding_bbox_audit"
    ]
    failed = next(item for item in audit["targets"] if item["role"] == target)
    assert "native_binding_bbox is required" in failed["error"]
    assert failed["pass"] is False
    assert report["strict_pass"] is False


def test_native_binding_bbox_is_frozen_by_critical_gate_hash(tmp_path: Path):
    run = _clearance_run(tmp_path)
    regions = read_json(run.regions_path)
    regions["regions"][0]["color_clearance_contracts"][0][
        "native_binding_bbox"
    ] = [5, 4, 11, 12]
    write_json(run.regions_path, regions)

    report = evaluate_regions(run)

    assert (
        "regions:expectation:caption-clearance:gate-hash-mismatch:"
        "color_clearance_contracts"
    ) in report["blockers"]


def test_subject_bbox_keeps_gold_arrow_out_of_red_caption_mask(tmp_path: Path):
    run = _clearance_run(tmp_path)
    for path in (run.source_png, run.render_png):
        image = Image.open(path)
        for y in range(10, 20):
            for x in range(25, 30):
                image.putpixel((x, y), (141, 106, 0))
        image.save(path)

    report = evaluate_regions(run)

    result = report["regions"][0]["color_clearance_contracts"][0]
    assert result["pass"] is True
    assert result["obstacles"][0]["reference"]["area_px"] == 50


def test_disconnected_same_color_glyph_outside_subject_bbox_is_blocked(
    tmp_path: Path,
):
    run = _clearance_run(tmp_path)
    render = Image.open(run.render_png)
    for y in range(6, 9):
        for x in range(35, 38):
            render.putpixel((x, y), (240, 0, 0))
    render.save(run.render_png)

    report = evaluate_regions(run)

    subject = report["regions"][0]["color_clearance_contracts"][0]["subject"]
    assert subject["new_outside_subject_color_px"] == 9
    assert subject["outside_subject_color_pass"] is False
    assert report["strict_pass"] is False


def test_reference_legal_same_color_object_outside_subject_is_not_an_addition(
    tmp_path: Path,
):
    run = _clearance_run(tmp_path)
    for path in (run.source_png, run.render_png):
        image = Image.open(path)
        for y in range(6, 9):
            for x in range(35, 38):
                image.putpixel((x, y), (240, 0, 0))
        image.save(path)

    report = evaluate_regions(run)

    subject = report["regions"][0]["color_clearance_contracts"][0]["subject"]
    assert subject["reference_outside_subject_color_px"] == 9
    assert subject["render_outside_subject_color_px"] == 9
    assert subject["new_outside_subject_color_px"] == 0
    assert subject["outside_subject_color_pass"] is True
    assert report["strict_pass"] is True


def test_subject_only_geometry_catches_caption_growth_even_if_total_ink_is_balanced(
    tmp_path: Path,
):
    run = _clearance_run(tmp_path)
    render = Image.open(run.render_png)
    for y in range(10, 20):
        for x in range(20, 25):
            render.putpixel((x, y), (240, 0, 0))
        for x in range(25, 30):
            render.putpixel((x, y), (255, 255, 255))
    render.save(run.render_png)

    report = evaluate_regions(run)

    result = report["regions"][0]["color_clearance_contracts"][0]
    assert result["pass"] is False
    assert result["subject"]["area_relative_error"] == 0.5


def test_subject_pixel_gate_catches_different_caption_ink_with_same_bbox_and_area(
    tmp_path: Path,
):
    run = _clearance_run(tmp_path)
    for path, reverse in ((run.source_png, False), (run.render_png, True)):
        image = Image.open(path)
        for y in range(10, 20):
            for x in range(10, 20):
                image.putpixel((x, y), (255, 255, 255))
        for offset in range(10):
            x = 19 - offset if reverse else 10 + offset
            image.putpixel((x, 10 + offset), (240, 0, 0))
        extra_points = ((10, 10), (19, 19)) if reverse else ((10, 19), (19, 10))
        for point in extra_points:
            image.putpixel(point, (240, 0, 0))
        image.save(path)

    report = evaluate_regions(run)

    subject = report["regions"][0]["color_clearance_contracts"][0]["subject"]
    assert subject["reference"]["bbox"] == subject["render"]["bbox"]
    assert subject["reference"]["area_px"] == subject["render"]["area_px"]
    assert subject["pass"] is False


def test_malformed_object_contracts_fail_closed_instead_of_crashing(tmp_path: Path):
    run = _clearance_run(tmp_path)
    regions = read_json(run.regions_path)
    region = regions["regions"][0]
    region["ink_contract"] = {"bbox_tolerance_px": "not-a-number"}
    region["color_clearance_contracts"][0]["absolute_min_px"] = "not-a-number"
    write_json(run.regions_path, regions)

    report = evaluate_regions(run)

    audited = report["regions"][0]
    assert audited["ink_contract"]["pass"] is False
    assert audited["color_clearance_contracts"][0]["pass"] is False
    assert "region:caption-clearance" in report["blockers"]
    assert any("gate-hash-mismatch" in blocker for blocker in report["blockers"])


def test_critical_region_inventory_blocks_silent_gate_removal(tmp_path: Path):
    run = _clearance_run(tmp_path)
    regions = read_json(run.regions_path)
    original = regions["critical_region_expectation"]["contracts"][0]
    missing = copy.deepcopy(original)
    missing["id"] = "required-caption-ink"
    missing["required_gates"] = ["pixel_metrics", "ink_contract"]
    missing["gate_sha256"] = {"ink_contract": "0" * 64}
    regions["critical_region_expectation"] = {
        "count": 2,
        "contracts": [original, missing],
    }
    write_json(run.regions_path, regions)

    report = evaluate_regions(run)

    expectation = report["critical_region_expectation"]
    assert expectation["pass"] is False
    assert expectation["missing_region_ids"] == ["required-caption-ink"]
    assert "regions:expectation:count-mismatch" in report["blockers"]
    assert "regions:expectation:missing-regions" in report["blockers"]


def test_critical_region_inventory_blocks_removing_a_required_gate(tmp_path: Path):
    run = _clearance_run(tmp_path)
    regions = read_json(run.regions_path)
    del regions["regions"][0]["color_clearance_contracts"]
    write_json(run.regions_path, regions)

    report = evaluate_regions(run)

    assert (
        "regions:expectation:caption-clearance:missing-gate:color_clearance_contracts"
        in report["blockers"]
    )


def test_critical_region_inventory_blocks_truncating_the_frozen_bbox(tmp_path: Path):
    run = _clearance_run(tmp_path)
    regions = read_json(run.regions_path)
    regions["regions"][0]["bbox"] = [5, 5, 20, 20]
    write_json(run.regions_path, regions)

    report = evaluate_regions(run)

    assert "regions:expectation:caption-clearance:bbox-mismatch" in report["blockers"]


def test_critical_region_inventory_hash_blocks_inner_obstacle_bbox_drift(
    tmp_path: Path,
):
    run = _clearance_run(tmp_path)
    regions = read_json(run.regions_path)
    regions["regions"][0]["color_clearance_contracts"][0]["obstacles"][0][
        "bbox"
    ] = [16, 4, 11, 12]
    write_json(run.regions_path, regions)

    report = evaluate_regions(run)

    assert (
        "regions:expectation:caption-clearance:gate-hash-mismatch:"
        "color_clearance_contracts"
    ) in report["blockers"]


def test_fuzzy_edge_iou_penalizes_a_doubled_nearby_contour():
    reference = np.full((50, 50, 3), 255, dtype=np.uint8)
    render = reference.copy()
    reference[8:42, 25] = 0
    render[8:42, 23] = 0
    render[8:42, 27] = 0

    assert _edge_iou(reference, render) < 0.75


def test_fuzzy_edge_iou_tolerates_one_pixel_rasterization_phase_shift():
    reference = np.full((50, 50, 3), 255, dtype=np.uint8)
    render = reference.copy()
    reference[8:42, 25] = 0
    render[8:42, 26] = 0

    assert _edge_iou(reference, render) >= 0.95
