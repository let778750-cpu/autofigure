from __future__ import annotations

import json
import zipfile
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation

from tools.core import common
from tools.core.contracts import read_json, transition, write_json
from tools.pipeline.convert import convert
from tools.arrows.pptx_arrows import write_arrow_reports
from tools.providers.providers import write_case_capabilities
from tools.regions.regions import evaluate_regions
from tools.repair.repair import (
    _inventory_sha256,
    _pptx_roundtrip_signature,
    _publish_file_set_atomically,
    _verify_pptx_save_reopen_structure,
    build_live_request,
    ingest_live_evidence,
    live_evidence_passes,
    main as repair_main,
    publish_live_candidate,
    publish_live_save_reopen_candidate,
)

XSL_PATH = Path(r"C:\Program Files\Microsoft Office\root\Office16\MML2OMML.XSL")


def _native_math_engine_ready() -> bool:
    try:
        import latex2mathml  # noqa: F401
    except ImportError:
        return False
    return XSL_PATH.is_file()


requires_native_math_engine = pytest.mark.skipif(
    not _native_math_engine_ready(),
    reason="需要 latex2mathml 与本机 Office MML2OMML.XSL",
)


def test_cli_save_reopen_only_routes_to_intermediate_publisher(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    run = _run(tmp_path)
    captured: dict[str, object] = {}

    def fake_publisher(
        opened: common.Run,
        evidence: Path,
        candidate: Path,
        *,
        reopened_path: Path | None,
        render_path: Path | None,
    ) -> None:
        captured.update(
            {
                "root": opened.root,
                "evidence": evidence,
                "candidate": candidate,
                "reopened": reopened_path,
                "render": render_path,
            }
        )

    monkeypatch.setattr(
        "tools.repair.repair.publish_live_save_reopen_candidate", fake_publisher
    )
    result = repair_main(
        [
            str(run.root),
            "--save-reopen-only",
            "--evidence",
            "evidence.json",
            "--candidate",
            "candidate.pptx",
            "--reopened",
            "reopened.pptx",
        ]
    )

    assert result == 0
    assert captured == {
        "root": run.root,
        "evidence": Path("evidence.json"),
        "candidate": Path("candidate.pptx"),
        "reopened": Path("reopened.pptx"),
        "render": None,
    }


def _run(tmp_path: Path) -> common.Run:
    reference = tmp_path / "reference-source.png"
    Image.new("RGB", (120, 100), "white").save(reference)
    run = common.create_run(
        reference,
        case="case",
        cases_root=tmp_path / "examples",
        input_route="svg-seeded",
    )
    run.redraw_svg.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="120" height="100" viewBox="0 0 120 100">'
        '<rect id="box" x="10" y="10" width="100" height="80" fill="#ffffff" stroke="#111111"/>'
        '</svg>',
        encoding="utf-8",
    )
    convert(run)
    Image.new("RGB", (120, 100), "white").save(run.render_png)
    regions = read_json(run.regions_path)
    regions["regions"] = [
        {
            "id": "failed",
            "bbox": [0, 0, 10, 10],
            "critical": True,
            "element_ids": ["box"],
        }
    ]
    write_json(run.regions_path, regions)
    return run


def _machine_evidence(
    run: common.Run,
    candidate: Path,
    *,
    regions: dict[str, str] | None = None,
    operation: str = "finalize_target",
) -> tuple[Path, Path, Path]:
    session_id = "11111111-2222-3333-4444-555555555555"
    candidate_hash = common.sha256_file(candidate)
    render = run.live_case_dir / "build" / "candidates" / "live.png"
    render.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", common.image_size(run.source_png), "white").save(render)
    inventory = {
        "presentation": {
            "project_id": run.root.name,
            "revision": 0,
        },
        "slides": [{"slide_id": 256, "slide_index": 1, "shapes": []}],
    }
    inventory_hash = _inventory_sha256(inventory)
    operation_log = (
        run.live_case_dir
        / "build"
        / "sessions"
        / "powerpoint"
        / session_id
        / "operation-log.ndjson"
    )
    operation_log.parent.mkdir(parents=True, exist_ok=True)
    bridge = read_json(run.live_bridge_path)
    save_event = (
        {
            "event": "finalize_target_prepared",
            "candidate_path": candidate.resolve()
            .relative_to(run.live_case_dir.resolve())
            .as_posix(),
            "candidate_sha256": candidate_hash,
            "reopened_inventory_sha256": inventory_hash,
            "preview_path": render.resolve()
            .relative_to(run.live_case_dir.resolve())
            .as_posix(),
            "preview_sha256": common.sha256_file(render),
        }
        if operation == "finalize_target"
        else {
            "event": "save_candidate",
            "tool": "powerpoint_save_candidate",
            "revision": 0,
            "candidate_path": candidate.resolve()
            .relative_to(run.live_case_dir.resolve())
            .as_posix(),
            "candidate_sha256": candidate_hash,
            "reopened_inventory_sha256": inventory_hash,
        }
    )
    events = [
        {
            "event": "begin_session",
            "revision": 0,
            "target_id": "autofigure-pptx",
            "scene_graph_sha256": bridge["adapter_scene_sha256"],
            "contract_hashes": {
                "sourceManifestSha256": bridge["contract_files"][
                    "input/source_manifest.json"
                ]
            },
        },
        save_event,
    ]
    operation_log.write_text(
        "".join(f"{json.dumps(event)}\n" for event in events),
        encoding="utf-8",
    )
    evidence_path = render.parent / "evidence.json"
    write_json(
        evidence_path,
        {
            "provider": "powerpoint-live",
            "reference_sha256": run.load_meta()["source_sha256"],
            "target_id": "autofigure-pptx",
            "saved_reopened": True,
            "bindings_complete": True,
            "candidate_sha256": candidate_hash,
            "reopened_artifact_sha256": candidate_hash,
            "live_inventory_sha256": inventory_hash,
            "reopened_inventory_sha256": inventory_hash,
            "live_inventory": inventory,
            "reopened_inventory": inventory,
            "render_sha256": common.sha256_file(render),
            "session_id": session_id,
            "host_reopen_method": "powerpoint-live",
            "arrow_mutations": False,
            "regions": regions or {"failed": "REGION_PASS"},
        },
    )
    return evidence_path, render, operation_log


def test_live_request_and_evidence_are_hash_bound(tmp_path: Path):
    run = _run(tmp_path)
    transition(run, "repairing", "live-request")
    request = build_live_request(run)
    assert request["visible"] is True
    assert request["failed_regions"] == ["failed"]
    assert request["failed_region_tasks"][0]["allowed_element_ids"] == ["box"]
    assert request["failed_region_tasks"][0]["all_other_elements_protected"] is True
    assert request["failed_region_tasks"][0]["manual_scope_required"] is False
    assert "save-reopen" in request["required_capabilities"]
    assert "finalize-target" in request["required_capabilities"]
    assert request["completion_contract"]["render_must_be_finalizer_bound"] is True
    assert request["scene_compatibility"]["source_schema_version"] == "4.0.0"
    assert request["scene_compatibility"]["adapter_schema_version"] == "2.1.0"
    assert (run.root / request["case_root"] / "project_state.json").is_file()
    assert (run.root / request["template_path"]).is_file()
    evaluate_regions(run)
    candidate = run.root / request["template_path"]
    candidate_hash = common.sha256_file(candidate)

    evidence_path, render, _ = _machine_evidence(run, candidate)
    publish_live_candidate(
        run,
        evidence_path,
        candidate,
        reopened_path=candidate,
        render_path=render,
    )
    assert live_evidence_passes(run, ["failed"]) == (True, [])
    bindings = read_json(run.bindings_path)
    assert bindings["saved_reopened"] is True
    assert bindings["artifact_sha256"] == candidate_hash

    published_evidence = read_json(run.live_evidence_path)
    live_summary = read_json(run.qa_dir / "live-save-reopen-summary.json")
    assert live_summary["saved_reopened"] is True
    assert live_summary["bindings_complete"] is True
    assert live_summary["published_to_case_root"] is True
    assert live_summary["live_attempt_inventory_equal"] is True
    assert live_summary["live_evidence_sha256"] == common.sha256_file(
        run.live_evidence_path
    )
    assert live_summary["current_root_candidate_sha256"] == candidate_hash
    assert live_summary["render_finalizer_bound"] is True
    assert live_summary["strict_live_blockers"] == []
    assert published_evidence["render_finalizer_bound"] is True
    assert published_evidence["publication_mode"] == "finalizer"
    finalizer_receipt = run.qa_dir / "powerpoint-live-operation-receipt.json"
    assert published_evidence["operation_receipt_sha256"] == common.sha256_file(
        finalizer_receipt
    )
    assert live_summary["operation_receipt_sha256"] == common.sha256_file(
        finalizer_receipt
    )
    assert read_json(run.arrow_compile_report_path)["case"] == run.root.name
    assert read_json(run.provider_capabilities_path)["case"] == run.root.name

    write_arrow_reports(run)
    write_case_capabilities(run)

    assert common.sha256_file(run.powerpoint_arrow_readback_path) == (
        published_evidence["arrow_readback_sha256"]
    )
    assert common.sha256_file(run.provider_capabilities_path) == (
        published_evidence["provider_capabilities_sha256"]
    )
    assert common.sha256_file(run.bindings_path) == published_evidence["bindings_sha256"]
    source_scene = run.qa_dir / "powerpoint-live-source-scene.json"
    inventory_path = run.qa_dir / "powerpoint-live-inventory.json"
    assert common.sha256_file(source_scene) == published_evidence[
        "source_scene_snapshot_sha256"
    ]
    assert common.sha256_file(inventory_path) == published_evidence[
        "powerpoint_live_inventory_sha256"
    ]
    assert read_json(run.live_bridge_path)["source_scene_sha256"] == common.sha256_file(
        source_scene
    )
    assert live_evidence_passes(run, ["failed"]) == (True, [])

    inventory_report = read_json(inventory_path)
    inventory_report["live_inventory"]["presentation"]["revision"] = 1
    write_json(inventory_path, inventory_report)
    passed, blockers = live_evidence_passes(run, ["failed"])
    assert passed is False
    assert "live-evidence-inventory-file-mismatch" in blockers


def test_live_evidence_mismatch_is_rejected(tmp_path: Path):
    run = _run(tmp_path)
    evidence_path = tmp_path / "bad-evidence.json"
    write_json(
        evidence_path,
        {
            "provider": "powerpoint-live",
            "reference_sha256": "wrong",
            "target_id": "autofigure-pptx",
            "saved_reopened": True,
            "bindings_complete": True,
            "regions": {},
        },
    )
    with pytest.raises(SystemExit, match="contract mismatch"):
        ingest_live_evidence(run, evidence_path)


def test_live_evidence_from_another_artifact_is_rejected(tmp_path: Path):
    run = _run(tmp_path)
    evaluate_regions(run)
    root_hash = common.sha256_file(run.pptx_path)
    bindings = read_json(run.bindings_path)
    bindings["saved_reopened"] = True
    write_json(run.bindings_path, bindings)
    evidence_path = tmp_path / "other-artifact.json"
    write_json(
        evidence_path,
        {
            "provider": "powerpoint-live",
            "reference_sha256": run.load_meta()["source_sha256"],
            "target_id": "autofigure-pptx",
            "saved_reopened": True,
            "bindings_complete": True,
            "candidate_sha256": "0" * 64,
            "reopened_artifact_sha256": "0" * 64,
            "binding_artifact_sha256": root_hash,
            "arrow_mutations": False,
            "regions": {},
        },
    )
    with pytest.raises(SystemExit, match="artifact identity mismatch"):
        ingest_live_evidence(run, evidence_path)


def test_live_publish_is_atomic_when_powerpoint_drops_bound_objects(tmp_path: Path):
    run = _run(tmp_path)
    build_live_request(run)
    evaluate_regions(run)
    original_hash = common.sha256_file(run.pptx_path)
    candidate_dir = run.live_case_dir / "build" / "candidates"
    candidate_dir.mkdir(parents=True)
    candidate = candidate_dir / "lossy.pptx"
    blank = Presentation()
    blank.slides.add_slide(blank.slide_layouts[6])
    blank.save(candidate)
    evidence_path, render, _ = _machine_evidence(run, candidate)

    with pytest.raises(SystemExit, match="bindings are incomplete"):
        publish_live_candidate(
            run,
            evidence_path,
            candidate,
            reopened_path=candidate,
            render_path=render,
        )

    assert common.sha256_file(run.pptx_path) == original_hash


def test_live_publish_requires_render_reopened_and_machine_log(tmp_path: Path):
    run = _run(tmp_path)
    request = build_live_request(run)
    candidate = run.root / request["template_path"]
    evidence_path, render, operation_log = _machine_evidence(run, candidate)

    with pytest.raises(SystemExit, match="explicit reopened artifact"):
        publish_live_candidate(run, evidence_path, candidate, render_path=render)
    with pytest.raises(SystemExit, match="requires a PowerPoint live render"):
        publish_live_candidate(
            run, evidence_path, candidate, reopened_path=candidate
        )
    operation_log.unlink()
    with pytest.raises(SystemExit, match="operation log"):
        publish_live_candidate(
            run,
            evidence_path,
            candidate,
            reopened_path=candidate,
            render_path=render,
        )


def test_live_publish_rejects_self_declared_render_not_in_finalizer_log(tmp_path: Path):
    run = _run(tmp_path)
    request = build_live_request(run)
    candidate = run.root / request["template_path"]
    evidence_path, render, _ = _machine_evidence(run, candidate)

    # Simulate a forged evidence bundle that updates its own render digest after
    # replacing the image, while the host-owned finalizer event remains bound to
    # the actual PowerPoint export.
    Image.new("RGB", common.image_size(run.source_png), "black").save(render)
    evidence = read_json(evidence_path)
    evidence["render_sha256"] = common.sha256_file(render)
    write_json(evidence_path, evidence)

    with pytest.raises(SystemExit, match="finalize_target candidate/render"):
        publish_live_candidate(
            run,
            evidence_path,
            candidate,
            reopened_path=candidate,
            render_path=render,
        )


def test_save_reopen_only_publication_is_honest_and_strictly_blocked(
    tmp_path: Path,
):
    run = _run(tmp_path)
    request = build_live_request(run)
    evaluate_regions(run)
    candidate = run.root / request["template_path"]
    candidate_hash = common.sha256_file(candidate)
    evidence_path, render, operation_log = _machine_evidence(
        run, candidate, operation="save_candidate"
    )

    published = publish_live_save_reopen_candidate(
        run,
        evidence_path,
        candidate,
        reopened_path=candidate,
        render_path=render,
    )

    assert common.sha256_file(run.pptx_path) == candidate_hash
    assert common.sha256_file(run.render_png) == common.sha256_file(render)
    assert published["publication_mode"] == "save-reopen-only"
    assert published["render_finalizer_bound"] is False
    assert published["live_render_published"] is True
    assert published["layout_audit_sha256"] == common.sha256_file(
        run.layout_audit_path
    )
    assert published["primitive_audit_sha256"] == common.sha256_file(
        run.primitive_audit_path
    )
    assert published["arrow_compile_report_sha256"] == common.sha256_file(
        run.arrow_compile_report_path
    )
    assert published["arrow_composition_audit_sha256"] == common.sha256_file(
        run.qa_dir / "arrow-composition-audit.json"
    )
    receipt_path = run.qa_dir / "powerpoint-live-operation-receipt.json"
    assert published["operation_receipt_sha256"] == common.sha256_file(
        receipt_path
    )
    receipt = read_json(receipt_path)
    assert receipt["begin_event_index"] < receipt["operation_event_index"]
    assert receipt["operation_log_sha256"] == common.sha256_file(operation_log)
    assert receipt["operation_event"]["event"] == "save_candidate"
    assert len(receipt["matching_begin_event_sha256"]) == 64
    assert len(receipt["matching_operation_event_sha256"]) == 64
    passed, blockers = live_evidence_passes(run, [])
    assert passed is False
    assert blockers == ["live-render-finalizer-unverified"]

    summary = read_json(run.qa_dir / "live-save-reopen-summary.json")
    assert summary["saved_reopened"] is True
    assert summary["bindings_complete"] is True
    assert summary["render_finalizer_bound"] is False
    assert summary["strict_live_blockers"] == [
        "live-render-finalizer-unverified"
    ]
    assert summary["operation_receipt_sha256"] == common.sha256_file(
        receipt_path
    )

    # build/sessions is transient.  The immutable formal receipt is enough to
    # retain the save/reopen audit after cleanup.
    operation_log.unlink()
    assert live_evidence_passes(run, []) == (
        False,
        ["live-render-finalizer-unverified"],
    )


@requires_native_math_engine
def test_save_reopen_publication_rebinds_math_summary_to_reopened_root(
    tmp_path: Path,
):
    reference = tmp_path / "math-reference.png"
    Image.new("RGB", (200, 100), "white").save(reference)
    run = common.create_run(
        reference,
        case="math-case",
        cases_root=tmp_path / "examples",
        input_route="svg-seeded",
    )
    run.redraw_svg.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" '
        'viewBox="0 0 200 100"><text x="20" y="50" font-size="18" '
        'font-style="italic">τ</text></svg>',
        encoding="utf-8",
    )
    convert(run)
    Image.new("RGB", (200, 100), "white").save(run.render_png)
    from tools.pipeline.math import upgrade as upgrade_math

    assert upgrade_math(run)["injected"] == 1
    request = build_live_request(run)
    candidate = run.root / request["template_path"]
    evidence_path, render, _ = _machine_evidence(
        run, candidate, operation="save_candidate"
    )

    published = publish_live_save_reopen_candidate(
        run,
        evidence_path,
        candidate,
        reopened_path=candidate,
        render_path=render,
    )

    math_summary = read_json(run.qa_dir / "math-summary.json")
    assert math_summary["pptx_sha256"] == common.sha256_file(run.pptx_path)
    assert math_summary["logical_formula_count"] == 1
    assert math_summary["omml_count"] == 1
    assert math_summary["verified"] == 1
    assert len(math_summary["formulas"]) == 1
    assert math_summary["saved_reopened"] is True
    assert published["math_summary_sha256"] == common.sha256_file(
        run.qa_dir / "math-summary.json"
    )


def test_save_reopen_only_rejects_candidate_prepared_before_root_math_rewrite(
    tmp_path: Path,
):
    run = _run(tmp_path)
    request = build_live_request(run)
    candidate = run.root / request["template_path"]
    evidence_path, render, _ = _machine_evidence(
        run, candidate, operation="save_candidate"
    )

    # Model the ordering bug that prompted this gate: Live was prepared from
    # the converter PPTX, then native-math (or any deterministic post-process)
    # rewrote root truth before the old Live candidate was published.
    presentation = Presentation(run.pptx_path)
    presentation.core_properties.subject = "post-bridge-native-math-rewrite"
    presentation.save(run.pptx_path)
    rewritten_hash = common.sha256_file(run.pptx_path)
    bindings = read_json(run.bindings_path)
    bindings["artifact_sha256"] = rewritten_hash
    write_json(run.bindings_path, bindings)
    write_json(
        run.qa_dir / "math-summary.json",
        {"injected": 1, "failed": 0, "pptx_sha256": rewritten_hash},
    )

    with pytest.raises(SystemExit, match="stale PowerPoint-live input candidate"):
        publish_live_save_reopen_candidate(
            run,
            evidence_path,
            candidate,
            reopened_path=candidate,
            render_path=render,
        )
    assert common.sha256_file(run.pptx_path) == rewritten_hash


def _write_math_signature_package(
    path: Path,
    formulas: list[tuple[int, str, str]],
    *,
    choice_copies: int = 1,
    fallback_copies: int = 1,
    fallback_identity: bool = True,
) -> None:
    alternate_contents = []
    for shape_id, shape_name, formula in formulas:
        carriers = []
        for carrier, copies in (
            ("Choice", choice_copies),
            ("Fallback", fallback_copies),
        ):
            requires = ' Requires="m"' if carrier == "Choice" else ""
            identity = (
                f'<p:cNvPr id="{shape_id}" name="{shape_name}"/>'
                if carrier == "Choice" or fallback_identity
                else ""
            )
            math_nodes = "".join(
                '<m:oMath><m:r>'
                f'<m:t>{formula}</m:t>'
                '</m:r></m:oMath>'
                for _ in range(copies)
            )
            carriers.append(
                f'<mc:{carrier}{requires}><p:sp><p:nvSpPr>'
                f'{identity}'
                '<p:cNvSpPr/><p:nvPr/></p:nvSpPr><p:spPr/>'
                '<p:txBody><a:bodyPr/><a:lstStyle/><a:p>'
                f'{math_nodes}'
                '</a:p></p:txBody></p:sp></mc:'
                f'{carrier}>'
            )
        alternate_contents.append(
            '<mc:AlternateContent>' + "".join(carriers) + '</mc:AlternateContent>'
        )
    slide_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006" '
        'xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math">'
        '<p:cSld><p:spTree>'
        '<p:sp><p:nvSpPr><p:cNvPr id="2" name="box"/>'
        '<p:cNvSpPr/><p:nvPr/></p:nvSpPr><p:spPr/></p:sp>'
        + "".join(alternate_contents)
        + '</p:spTree></p:cSld></p:sld>'
    )
    path.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(path, "w") as package:
        package.writestr("ppt/slides/slide1.xml", slide_xml)


def test_roundtrip_signature_counts_choice_fallback_as_one_logical_math(
    tmp_path: Path,
):
    package = tmp_path / "logical-math.pptx"
    _write_math_signature_package(package, [(10, "math:001", "x")])

    identities, native_math = _pptx_roundtrip_signature(package)

    assert ("ppt/slides/slide1.xml", 10, "math:001") in identities
    assert native_math == (
        ("ppt/slides/slide1.xml", 10, "math:001", ("x",)),
    )


@pytest.mark.parametrize(
    ("choice_copies", "fallback_copies", "fallback_identity"),
    [(2, 0, True), (2, 1, True), (1, 1, False)],
)
def test_roundtrip_signature_rejects_duplicate_or_unbound_math_branches(
    tmp_path: Path,
    choice_copies: int,
    fallback_copies: int,
    fallback_identity: bool,
):
    package = tmp_path / "invalid-logical-math.pptx"
    _write_math_signature_package(
        package,
        [(10, "math:001", "x")],
        choice_copies=choice_copies,
        fallback_copies=fallback_copies,
        fallback_identity=fallback_identity,
    )

    with pytest.raises(SystemExit, match="native-math branch"):
        _pptx_roundtrip_signature(package)


@pytest.mark.parametrize(
    ("source_formulas", "candidate_formulas"),
    [
        (
            [(10, "math:001", "x")],
            [(10, "math:001", "x"), (11, "math:002", "y")],
        ),
        (
            [(10, "math:001", "x"), (11, "math:002", "y")],
            [(10, "math:001", "x")],
        ),
        (
            [(10, "math:001", "x"), (11, "math:002", "y")],
            [(10, "math:001", "y"), (11, "math:002", "x")],
        ),
    ],
)
def test_roundtrip_structure_rejects_logical_math_add_delete_or_rebinding(
    tmp_path: Path,
    source_formulas: list[tuple[int, str, str]],
    candidate_formulas: list[tuple[int, str, str]],
):
    run = _run(tmp_path)
    live_input = run.live_case_dir / "input" / "candidate.pptx"
    candidate = run.live_case_dir / "build" / "candidates" / "candidate.pptx"
    _write_math_signature_package(live_input, source_formulas)
    _write_math_signature_package(candidate, candidate_formulas)

    with pytest.raises(
        SystemExit, match="changed shape identities or native-math inventory"
    ):
        _verify_pptx_save_reopen_structure(run, candidate, candidate)


def test_save_reopen_only_rejects_non_save_operation_and_scene_drift(
    tmp_path: Path,
):
    run = _run(tmp_path)
    request = build_live_request(run)
    candidate = run.root / request["template_path"]
    evidence_path, render, operation_log = _machine_evidence(run, candidate)

    with pytest.raises(SystemExit, match="save_candidate operation"):
        publish_live_save_reopen_candidate(
            run,
            evidence_path,
            candidate,
            reopened_path=candidate,
            render_path=render,
        )

    evidence_path, render, operation_log = _machine_evidence(
        run, candidate, operation="save_candidate"
    )
    events = [
        json.loads(line)
        for line in operation_log.read_text(encoding="utf-8").splitlines()
    ]
    events[0]["scene_graph_sha256"] = "0" * 64
    operation_log.write_text(
        "".join(f"{json.dumps(event)}\n" for event in events),
        encoding="utf-8",
    )
    with pytest.raises(SystemExit, match="current-scene begin_session"):
        publish_live_save_reopen_candidate(
            run,
            evidence_path,
            candidate,
            reopened_path=candidate,
            render_path=render,
        )


def test_save_reopen_only_requires_begin_before_save(tmp_path: Path):
    run = _run(tmp_path)
    request = build_live_request(run)
    candidate = run.root / request["template_path"]
    evidence_path, render, operation_log = _machine_evidence(
        run, candidate, operation="save_candidate"
    )
    events = [
        json.loads(line)
        for line in operation_log.read_text(encoding="utf-8").splitlines()
    ]
    operation_log.write_text(
        "".join(f"{json.dumps(event)}\n" for event in reversed(events)),
        encoding="utf-8",
    )

    with pytest.raises(SystemExit, match="before the matching save operation"):
        publish_live_save_reopen_candidate(
            run,
            evidence_path,
            candidate,
            reopened_path=candidate,
            render_path=render,
        )


def test_operation_receipt_missing_or_drifted_fails_closed(tmp_path: Path):
    run = _run(tmp_path)
    request = build_live_request(run)
    candidate = run.root / request["template_path"]
    evidence_path, render, _ = _machine_evidence(
        run, candidate, operation="save_candidate"
    )
    publish_live_save_reopen_candidate(
        run,
        evidence_path,
        candidate,
        reopened_path=candidate,
        render_path=render,
    )
    receipt_path = run.qa_dir / "powerpoint-live-operation-receipt.json"
    receipt = read_json(receipt_path)
    receipt["begin_event_index"] = receipt["operation_event_index"]
    write_json(receipt_path, receipt)
    passed, blockers = live_evidence_passes(run, [])
    assert passed is False
    assert "live-evidence-operation-receipt-mismatch" in blockers
    assert "live-operation-receipt-order-invalid" in blockers

    receipt_path.unlink()
    passed, blockers = live_evidence_passes(run, [])
    assert passed is False
    assert "live-operation-receipt-missing" in blockers


def test_save_reopen_only_can_publish_candidate_without_new_render(tmp_path: Path):
    run = _run(tmp_path)
    request = build_live_request(run)
    candidate = run.root / request["template_path"]
    original_render_hash = common.sha256_file(run.render_png)
    evidence_path, _, _ = _machine_evidence(
        run, candidate, operation="save_candidate"
    )
    evidence = read_json(evidence_path)
    evidence.pop("render_sha256")
    evidence.pop("regions")
    write_json(evidence_path, evidence)

    published = publish_live_save_reopen_candidate(
        run,
        evidence_path,
        candidate,
        reopened_path=candidate,
    )

    assert common.sha256_file(run.render_png) == original_render_hash
    assert published["live_render_published"] is False
    assert published["render_finalizer_bound"] is False
    assert live_evidence_passes(run, ["failed"])[0] is False
    assert "live-render-finalizer-unverified" in live_evidence_passes(
        run, ["failed"]
    )[1]


def test_file_set_publish_rolls_back_every_replaced_destination(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    source_one = tmp_path / "source-one"
    source_two = tmp_path / "source-two"
    destination_one = tmp_path / "destination-one"
    destination_two = tmp_path / "destination-two"
    source_one.write_text("new-one", encoding="utf-8")
    source_two.write_text("new-two", encoding="utf-8")
    destination_one.write_text("old-one", encoding="utf-8")
    destination_two.write_text("old-two", encoding="utf-8")
    from tools.repair import repair

    real_replace = repair.os.replace
    destination_replacements = 0

    def fail_once(source: Path, destination: Path) -> None:
        nonlocal destination_replacements
        if Path(destination) in {destination_one, destination_two}:
            destination_replacements += 1
            if destination_replacements == 2:
                raise OSError("injected publication failure")
        real_replace(source, destination)

    monkeypatch.setattr(repair.os, "replace", fail_once)
    with pytest.raises(OSError, match="injected publication failure"):
        _publish_file_set_atomically(
            [(source_one, destination_one), (source_two, destination_two)],
            tmp_path / "transaction",
        )

    assert destination_one.read_text(encoding="utf-8") == "old-one"
    assert destination_two.read_text(encoding="utf-8") == "old-two"
