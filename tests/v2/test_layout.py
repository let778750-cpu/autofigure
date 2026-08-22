from __future__ import annotations

from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation

from tools.v2 import common
from tools.v2.contracts import read_json
from tools.v2.convert import convert
from tools.v2.layout import audit_layout, strict_blockers


@pytest.fixture()
def run_factory(tmp_path: Path):
    def make(body: str) -> common.Run:
        reference = tmp_path / "reference.png"
        Image.new("RGB", (200, 120), "white").save(reference)
        run = common.create_run(
            reference,
            case="layout",
            cases_root=tmp_path / "examples",
            input_route="svg-seeded",
        )
        run.redraw_svg.write_text(
            '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="120" '
            f'viewBox="0 0 200 120">{body}</svg>',
            encoding="utf-8",
        )
        return run

    return make


def test_container_annotation_clips_selection_box_and_survives_readback(run_factory):
    run = run_factory(
        '<rect id="box" x="10" y="10" width="100" height="40" fill="#eeeeee"/>'
        '<text id="label" data-layout-container="box" data-layout-padding="2" '
        'x="18" y="34" font-size="16">contained</text>'
    )
    summary = convert(run)
    assert summary["layout_pass"] is True
    report = audit_layout(run)
    assert report["pass"] is True

    presentation = Presentation(run.pptx_path)
    shapes = {shape.name: shape for shape in presentation.slides[0].shapes}
    box = shapes["af-box-rect-01"]
    label = shapes["af-label-text-01"]
    emu_per_px = 9525
    assert label.left / emu_per_px >= box.left / emu_per_px + 2 - 0.01
    assert (label.top + label.height) / emu_per_px <= (
        (box.top + box.height) / emu_per_px - 2 + 0.01
    )
    scene = read_json(run.scene_path)
    label_record = next(item for item in scene["elements"] if item["id"] == "label")
    assert label_record["layout"]["container_id"] == "box"


def test_repeat_contract_detects_irregular_source_and_backend_spacing(run_factory):
    repeated = "".join(
        f'<circle id="dot-{index}" data-repeat-group="dots" data-repeat-axis="vertical" '
        f'data-repeat-order="{index}" cx="50" cy="{cy}" r="8" fill="#cc6600"/>'
        for index, cy in enumerate((20, 50, 70), start=1)
    )
    run = run_factory(repeated)
    convert(run)
    report = audit_layout(run)
    assert report["pass"] is False
    spacing_findings = [item for item in report["findings"] if item["code"] == "L9"]
    assert {item["stage"] for item in spacing_findings} == {"source", "backend"}
    assert any(blocker == "layout:L9:source:dots" for blocker in strict_blockers(report))


def test_repeat_contract_allows_one_pixel_reference_quantization(run_factory):
    repeated = "".join(
        f'<circle id="dot-{index}" data-repeat-group="dots" data-repeat-axis="vertical" '
        f'data-repeat-order="{index}" cx="50" cy="{cy}" r="8" fill="#cc6600"/>'
        for index, cy in enumerate((20, 38, 55), start=1)
    )
    run = run_factory(repeated)
    convert(run)
    report = audit_layout(run)
    group = report["repeat_groups"][0]
    assert group["source"]["steps_px"] == [18.0, 17.0]
    assert group["backend"]["steps_px"] == [18.0, 17.0]
    assert report["pass"] is True


def test_layout_annotation_without_stable_id_is_a_hard_finding(run_factory):
    run = run_factory(
        '<rect id="box" x="10" y="10" width="100" height="40" fill="#eeeeee"/>'
        '<text data-layout-container="box" x="18" y="34" font-size="16">anonymous</text>'
    )
    convert(run)
    report = audit_layout(run)
    assert any(item["code"] == "L1" for item in report["findings"])


def test_saved_powerpoint_object_must_remain_inside_slide_canvas(run_factory):
    run = run_factory('<rect id="box" x="10" y="10" width="40" height="30" fill="#eeeeee"/>')
    convert(run)
    presentation = Presentation(run.pptx_path)
    presentation.slides[0].shapes[0].left = -10 * 9525
    presentation.save(run.pptx_path)

    report = audit_layout(run)
    finding = next(item for item in report["findings"] if item["code"] == "L10")
    assert finding["stage"] == "backend"
    assert finding["target"] == "box"
    assert finding["metrics"]["overflow_px"]["left"] == pytest.approx(10.0)
    assert "layout:L10:backend:box" in strict_blockers(report)
