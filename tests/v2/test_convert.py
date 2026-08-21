"""convert 合同测试：SVG 元素 → 原生 PPTX 对象读回（不需要 PowerPoint）。"""

from __future__ import annotations

import json
import zipfile
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from tools.v2 import common
from tools.v2.contracts import read_json, write_json
from tools.v2.convert import convert

VALID_DASHES = {
    "solid", "dot", "sysDash", "sysDot", "sysDashDot", "sysDashDotDot",
    "lgDash", "lgDashDot", "lgDashDotDot", "dash", "dashDot",
}


@pytest.fixture()
def run_factory(tmp_path: Path):
    def make(svg: str, size: tuple[int, int] = (200, 100)) -> common.Run:
        source = tmp_path / "ref.png"
        Image.new("RGB", size, (240, 240, 240)).save(source)
        run = common.create_run(source, case="case", cases_root=tmp_path / "examples")
        run.qa_dir.mkdir(exist_ok=True)
        run.redraw_svg.write_text(svg, encoding="utf-8")
        return run

    return make


def _shapes(run: common.Run):
    return list(Presentation(run.pptx_path).slides[0].shapes)


def test_rect_with_radius_fill_and_stroke(run_factory):
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" viewBox="0 0 200 100">'
        '<rect x="10" y="10" width="50" height="30" rx="6" fill="#AABBCC" stroke="#112233" stroke-width="2"/></svg>'
    )
    convert(run)
    (shape,) = _shapes(run)
    xml = shape._element.xml
    assert 'prst="roundRect"' in xml
    assert 'srgbClr val="AABBCC"' in xml
    assert 'srgbClr val="112233"' in xml
    # spPr 子元素顺序：fill 必须在 effectLst 之前（theme 蓝色回归防护）
    sp_pr = shape._element.spPr
    tags = [child.tag for child in sp_pr]
    assert tags.index(qn("a:solidFill")) < tags.index(qn("a:effectLst"))


def test_linear_gradient_becomes_gradfill(run_factory):
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" viewBox="0 0 200 100">'
        "<defs><linearGradient id=\"g\" x1=\"0\" y1=\"0\" x2=\"0\" y2=\"1\">"
        '<stop offset="0" stop-color="#FF0000" stop-opacity="0.5"/>'
        '<stop offset="1" stop-color="#0000FF"/></linearGradient></defs>'
        '<rect x="0" y="0" width="100" height="50" fill="url(#g)"/></svg>'
    )
    convert(run)
    (shape,) = _shapes(run)
    grad = shape._element.spPr.find(qn("a:gradFill"))
    assert grad is not None
    stops = grad.findall(f"{qn('a:gsLst')}/{qn('a:gs')}")
    assert len(stops) == 2
    first = stops[0].find(qn("a:srgbClr"))
    assert first.get("val") == "FF0000"
    assert first.find(qn("a:alpha")) is not None  # stop-opacity 保留
    lin = grad.find(qn("a:lin"))
    assert lin is not None and lin.get("ang") == str(90 * 60000)  # 垂直向下


def test_text_runs_italic_bold_and_baseline_shift(run_factory):
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" viewBox="0 0 200 100">'
        '<text x="20" y="50" font-family="Times New Roman, serif" font-size="20" fill="#111111">'
        '<tspan font-style="italic">z</tspan>'
        '<tspan baseline-shift="super" font-size="14" font-style="italic">τ</tspan>'
        '<tspan baseline-shift="sub" font-size="14" font-style="italic">t+1</tspan></text></svg>'
    )
    convert(run)
    (shape,) = _shapes(run)
    assert shape.has_text_frame
    runs = shape.text_frame.paragraphs[0].runs
    assert [r.text for r in runs] == ["z", "τ", "t+1"]
    assert all(r.font.italic for r in runs)
    assert runs[0].font.name == "Times New Roman"
    assert runs[1]._r.get_or_add_rPr().get("baseline") == "30000"
    assert runs[2]._r.get_or_add_rPr().get("baseline") == "-25000"


def test_text_overflow_padding_preserves_middle_anchor(run_factory):
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" '
        'viewBox="0 0 200 100"><text x="100" y="50" font-size="16" '
        'text-anchor="middle">task</text></svg>'
    )
    convert(run)
    (shape,) = _shapes(run)
    center_x = (shape.left + shape.width / 2) / 9525
    assert center_x == pytest.approx(100, abs=0.01)
    assert shape.height / 9525 == pytest.approx(40, abs=0.01)
    assert shape.width / 9525 > 32


def test_text_padding_is_clipped_to_canvas_without_moving_start_anchor(run_factory):
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" '
        'viewBox="0 0 200 100"><text x="180" y="50" font-size="16" '
        'text-anchor="start">Edge</text></svg>'
    )
    convert(run)
    (shape,) = _shapes(run)
    assert shape.left / 9525 == pytest.approx(180, abs=0.01)
    assert (shape.left + shape.width) / 9525 <= 200.01


def test_dasharray_maps_to_valid_ooxml_preset(run_factory):
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" viewBox="0 0 200 100">'
        '<path d="M 0 50 C 50 0 150 100 200 50" stroke="#2d5ea8" stroke-width="2.8" fill="none"'
        ' stroke-linecap="round" stroke-dasharray="2 4"/></svg>'
    )
    convert(run)
    (shape,) = _shapes(run)
    dash = shape._element.spPr.find(f"{qn('a:ln')}/{qn('a:prstDash')}")
    assert dash is not None
    assert dash.get("val") in VALID_DASHES  # roundDot/squareDot 曾致 PowerPoint 判损坏


def test_freeform_bbox_includes_bezier_control_points(run_factory):
    # 控制点 (300,-50) 越出端点 bbox——曾致 PowerPoint 判文件损坏
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" viewBox="0 0 200 100">'
        '<path d="M 0 50 C 100 -200 300 200 199 50" stroke="#000000" fill="none"/></svg>'
    )
    convert(run)
    (shape,) = _shapes(run)
    sp_pr = shape._element.spPr
    path = sp_pr.find(f"{qn('a:custGeom')}/{qn('a:pathLst')}/{qn('a:path')}")
    assert path is not None
    width = int(path.get("w"))
    height = int(path.get("h"))
    xs = [int(pt.get("x")) for pt in path.iter(qn("a:pt"))]
    ys = [int(pt.get("y")) for pt in path.iter(qn("a:pt"))]
    assert min(xs) >= 0 and max(xs) <= width
    assert min(ys) >= 0 and max(ys) <= height


def test_cubic_preserved_as_cubicbezto(run_factory):
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" viewBox="0 0 200 100">'
        '<path d="M 0 0 C 50 50 100 50 150 0" stroke="#000000" fill="none"/></svg>'
    )
    convert(run)
    (shape,) = _shapes(run)
    cubic = shape._element.spPr.find(f".//{qn('a:cubicBezTo')}")
    assert cubic is not None


def test_atomic_placeholder_crops_reference(run_factory):
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" viewBox="0 0 200 100">'
        '<rect id="atomic:photo" x="20" y="20" width="40" height="30" '
        'fill="#EEEEEE" stroke="#999999" stroke-dasharray="4 3"/></svg>'
    )
    convert(run)
    (shape,) = _shapes(run)
    assert shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    assert shape.width > 0 and shape.height > 0


def test_authorized_atomic_asset_embeds_powerpoint_live_raster_tags(run_factory):
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" '
        'viewBox="0 0 200 100"><rect id="atomic:photo" x="20" y="20" '
        'width="40" height="30"/></svg>'
    )
    assets = read_json(run.assets_path)
    assets["assets"] = [
        {
            "id": "atomic:photo",
            "authorized": True,
            "authorization_basis": "user supplied reference crop",
            "editable": False,
        }
    ]
    write_json(run.assets_path, assets)
    convert(run)

    with zipfile.ZipFile(run.pptx_path) as package:
        tag_parts = [name for name in package.namelist() if name.startswith("ppt/tags/tag")]
        assert len(tag_parts) == 1
        tags = package.read(tag_parts[0]).decode("utf-8")
        slide = package.read("ppt/slides/slide1.xml").decode("utf-8")
        relationships = package.read("ppt/slides/_rels/slide1.xml.rels").decode("utf-8")
    assert "AISCIENTIFICILLUSTRATORASSETID" in tags
    assert 'val="atomic:photo"' in tags
    assert "AISCIENTIFICILLUSTRATORSOURCESHA256" in tags
    assert 'val="True"' in tags and 'val="False"' in tags
    assert "custDataLst" in slide and "p:tags" in slide
    assert "relationships/tags" in relationships
    asset = read_json(run.assets_path)["assets"][0]
    assert len(asset["source_sha256"]) == 64
    assert asset["atomic_raster_unit"] is True
    assert asset["contains_reconstructable_content"] is False


def test_simple_marker_becomes_native_powerpoint_arrowhead(run_factory):
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" viewBox="0 0 200 100">'
        '<defs><marker id="arr" markerWidth="10" markerHeight="10" refX="8" refY="5" orient="auto"'
        ' markerUnits="userSpaceOnUse"><path d="M1,1 L8,5 L1,9" fill="none" stroke="#000000"/></marker></defs>'
        '<line x1="10" y1="50" x2="190" y2="50" stroke="#000000" marker-end="url(#arr)"/></svg>'
    )
    convert(run)
    shapes = _shapes(run)
    assert len(shapes) == 1
    tail = shapes[0]._element.spPr.find(f"{qn('a:ln')}/{qn('a:tailEnd')}")
    assert tail is not None
    assert tail.get("type") == "arrow"


def _triangle_abs_points(shapes):
    """返回第一个 3 顶点自由三角形在画布上的绝对坐标（px）。"""
    emu = 9525.0
    for shape in shapes:
        sp_pr = shape._element.find(qn("p:spPr"))
        if sp_pr is None:
            continue
        geom = sp_pr.find(qn("a:custGeom"))
        if geom is None:
            continue
        path = geom.find(f"{qn('a:pathLst')}/{qn('a:path')}")
        pts = [(int(pt.get("x")), int(pt.get("y"))) for pt in path.iter(qn("a:pt"))]
        if len(pts) != 3:
            continue
        off = sp_pr.find(f"{qn('a:xfrm')}/{qn('a:off')}")
        ox, oy = int(off.get("x")), int(off.get("y"))
        return [((ox + x) / emu, (oy + y) / emu) for x, y in pts]
    return None


def _tip_and_base_mid(pts):
    cx = sum(p[0] for p in pts) / 3
    cy = sum(p[1] for p in pts) / 3
    tip = max(pts, key=lambda p: (p[0] - cx) ** 2 + (p[1] - cy) ** 2)
    base = [p for p in pts if p is not tip]
    mid = ((base[0][0] + base[1][0]) / 2, (base[0][1] + base[1][1]) / 2)
    return tip, mid


def test_curved_path_uses_native_end_that_follows_path_tangent(run_factory):
    # 末端切线 = end − 第二控制点 = (0,-80) 竖直向上；旧实现按弦方向（atan2(-80,160)≈-26.6°）
    # 放置导致三角横甩脱开（01 案例 π/a 圆间曲线箭头 43-47° 偏差的真实根因）
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" viewBox="0 0 200 100">'
        '<defs><marker id="arr" markerWidth="12" markerHeight="12" refX="10" refY="6" orient="auto"'
        ' markerUnits="userSpaceOnUse"><path d="M0,0 L12,6 L0,12 Z" fill="#000000" stroke="none"/></marker></defs>'
        '<path d="M 20 100 C 100 100 180 100 180 20" stroke="#000000" fill="none" marker-end="url(#arr)"/></svg>'
    )
    convert(run)
    (shape,) = _shapes(run)
    assert shape._element.spPr.find(f".//{qn('a:cubicBezTo')}") is not None
    tail = shape._element.spPr.find(f"{qn('a:ln')}/{qn('a:tailEnd')}")
    assert tail is not None and tail.get("type") == "triangle"


def test_marker_start_becomes_native_head_end(run_factory):
    # SVG 语义：orient="auto" 在 marker-start 处沿行进方向放置；
    # 定义朝 -x 的起始箭头（尖端 local x=0）→ 尖端应落在起点外侧（左侧）
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" viewBox="0 0 200 100">'
        '<defs><marker id="sarr" markerWidth="12" markerHeight="12" refX="2" refY="6" orient="auto"'
        ' markerUnits="userSpaceOnUse"><path d="M12,0 L0,6 L12,12 Z" fill="#000000" stroke="none"/></marker></defs>'
        '<line x1="10" y1="50" x2="190" y2="50" stroke="#000000" marker-start="url(#sarr)"/></svg>'
    )
    convert(run)
    (shape,) = _shapes(run)
    head = shape._element.spPr.find(f"{qn('a:ln')}/{qn('a:headEnd')}")
    assert head is not None and head.get("type") == "triangle"



def test_viewbox_mismatch_rejected(run_factory, tmp_path: Path):
    source = tmp_path / "ref2.png"
    Image.new("RGB", (100, 100), (255, 255, 255)).save(source)
    run = common.create_run(source, case="case2", cases_root=tmp_path / "examples2")
    run.redraw_svg.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="300" height="300" viewBox="0 0 300 300">'
        '<rect x="0" y="0" width="10" height="10"/></svg>',
        encoding="utf-8",
    )
    with pytest.raises(SystemExit, match="viewBox"):
        convert(run)


def test_summary_counts_and_readback(run_factory):
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" viewBox="0 0 200 100">'
        '<rect x="0" y="0" width="200" height="100" fill="#FFFFFF"/>'
        '<text x="10" y="50" font-size="16">hello</text></svg>'
    )
    summary = convert(run)
    assert summary["textbox_with_text"] == 1
    assert summary["shape_count"] == 2
    on_disk = json.loads((run.qa_dir / "convert-summary.json").read_text(encoding="utf-8"))
    assert on_disk["shape_count"] == 2


def test_data_uri_image_tolerated_as_atomic_crop(run_factory):
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" xmlns:xlink="http://www.w3.org/1999/xlink"'
        ' width="200" height="100" viewBox="0 0 200 100">'
        '<image x="20" y="20" width="40" height="30" preserveAspectRatio="none"'
        ' xlink:href="data:image/png;base64,iVBORw0KGgo="/></svg>'
    )
    summary = convert(run)
    (shape,) = _shapes(run)
    assert shape.shape_type == MSO_SHAPE_TYPE.PICTURE  # 按 bbox 从参考图裁剪，不读内嵌数据
    assert summary["emitted"].get("atomic") == 1
    assert any("<image>" in w for w in summary["warnings"])


def test_oversized_image_rejected_as_canvas_cheat(run_factory):
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" xmlns:xlink="http://www.w3.org/1999/xlink"'
        ' width="200" height="100" viewBox="0 0 200 100">'
        '<image x="0" y="0" width="200" height="100" xlink:href="data:image/png;base64,iVBORw0KGgo="/></svg>'
    )
    with pytest.raises(SystemExit, match="整图截图"):
        convert(run)


def test_stroke_none_removes_powerpoint_default_outline(run_factory):
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" '
        'viewBox="0 0 200 100"><rect x="10" y="10" width="50" height="30" '
        'fill="#EEEEEE" stroke="none"/></svg>'
    )
    convert(run)
    (shape,) = _shapes(run)
    assert shape._element.spPr.find(f"{qn('a:ln')}/{qn('a:noFill')}") is not None


def test_unsupported_marker_fallback_is_physically_grouped(run_factory):
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" '
        'viewBox="0 0 200 100"><defs><marker id="complex" refX="8" refY="5" '
        'orient="auto" markerUnits="userSpaceOnUse">'
        '<path d="M0,0 L8,5 L0,10 Z" fill="#111111"/>'
        '<path d="M1,5 L4,5" fill="none" stroke="#FFFFFF"/>'
        '</marker></defs><line id="edge" x1="10" y1="50" x2="190" y2="50" '
        'stroke="#111111" marker-end="url(#complex)"/></svg>'
    )
    summary = convert(run)
    (group,) = _shapes(run)
    assert getattr(group, "shapes", None) is not None
    assert len(group.shapes) == 3
    assert summary["emitted"]["arrow-group"] == 1
    bindings = json.loads(run.bindings_path.read_text(encoding="utf-8"))["bindings"]
    children = [item for item in bindings if item.get("group_shape_id")]
    assert children and all(item["physically_grouped"] for item in children)


def test_connector_identity_binding_and_ooxml_attachments(run_factory):
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" '
        'viewBox="0 0 200 100">'
        '<rect id="source" x="10" y="30" width="30" height="40"/>'
        '<rect id="target" x="160" y="30" width="30" height="40"/>'
        '<line id="edge" data-source-id="source" data-target-id="target" '
        'data-source-site="3" data-target-site="1" x1="40" y1="50" x2="160" y2="50" '
        'stroke="#111111"/></svg>'
    )
    convert(run)
    shapes = _shapes(run)
    connector = next(shape for shape in shapes if shape.name.startswith("af-edge-connector"))
    c_nv = connector._element.find(qn("p:nvCxnSpPr")).find(qn("p:cNvCxnSpPr"))
    start = c_nv.find(qn("a:stCxn"))
    end = c_nv.find(qn("a:endCxn"))
    assert start is not None and end is not None
    assert start.get("id") != end.get("id")
    bindings = json.loads(run.bindings_path.read_text(encoding="utf-8"))
    assert bindings["bindings_complete"] is True


def test_missing_viewbox_is_rejected(run_factory):
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100">'
        '<rect x="0" y="0" width="10" height="10"/></svg>'
    )
    with pytest.raises(SystemExit, match="viewBox"):
        convert(run)


def test_conversion_drops_stale_scene_elements(run_factory):
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" '
        'viewBox="0 0 200 100"><rect id="current" x="10" y="10" width="20" height="20"/></svg>'
    )
    scene = read_json(run.scene_path)
    scene["elements"] = [{"id": "stale", "kind": "shape", "editable": True, "z_index": 0}]
    write_json(run.scene_path, scene)
    convert(run)
    ids = {item["id"] for item in read_json(run.scene_path)["elements"]}
    assert "current" in ids
    assert "stale" not in ids


def test_reconversion_refreshes_scene_and_powerpoint_z_order(run_factory):
    first = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" '
        'viewBox="0 0 200 100"><line id="edge" x1="10" y1="10" '
        'x2="90" y2="90" stroke="#777777"/><rect id="panel" x="50" '
        'y="40" width="100" height="50" fill="#EEEEEE"/></svg>'
    )
    run = run_factory(first)
    convert(run)
    initial = {item["id"]: item for item in read_json(run.scene_path)["elements"]}
    assert initial["edge"]["z_index"] < initial["panel"]["z_index"]

    # Move the edge after the panel.  Existing scene records must not override
    # the new DOM/PowerPoint stack order during reconversion.
    run.redraw_svg.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" '
        'viewBox="0 0 200 100"><rect id="panel" x="50" y="40" '
        'width="100" height="50" fill="#EEEEEE"/><line id="edge" '
        'x1="10" y1="10" x2="90" y2="90" stroke="#777777"/></svg>',
        encoding="utf-8",
    )
    convert(run)

    refreshed = {item["id"]: item for item in read_json(run.scene_path)["elements"]}
    assert refreshed["edge"]["z_index"] > refreshed["panel"]["z_index"]
    bindings = {
        item["element_id"]: item
        for item in read_json(run.bindings_path)["bindings"]
    }
    assert bindings["edge"]["shape_id"] > bindings["panel"]["shape_id"]
