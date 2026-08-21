from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation

from tools.v2 import common
from tools.v2.contracts import read_json
from tools.v2.convert import convert
from tools.v2.ingest import main as ingest_main
from tools.v2.prepare import main as prepare_main


def _reference(tmp_path: Path) -> Path:
    path = tmp_path / "input-only.png"
    Image.new("RGB", (160, 100), "white").save(path)
    return path


def test_prepare_png_only_creates_tasks_without_web_svg_prerequisite(tmp_path: Path):
    cases_root = tmp_path / "examples"
    assert prepare_main(
        [
            str(_reference(tmp_path)),
            "--case",
            "png-only",
            "--cases-root",
            str(cases_root),
            "--source-mode",
            "png_reconstruct",
        ]
    ) == 0

    run = common.open_run(cases_root / "png-only")
    meta = run.load_meta()
    assert meta["source_mode"] == "png_reconstruct"
    assert meta["fidelity_profile"] == "hybrid_fidelity"
    assert meta["workflow"]["state"] == "prepared"
    assert not run.redraw_svg.exists()
    tasks = read_json(run.region_tasks_path)
    assert tasks["reference_sha256"] == meta["source_sha256"]
    assert tasks["result_contract"]["offline_initial_render_carrier"] == "svg"
    assert tasks["tasks"]
    prompt = run.prompt_md.read_text(encoding="utf-8")
    assert "不要求、也不依赖 GPT Web" in prompt
    assert "不能把“入口已连通”误写成“PNG 已自动一比一重建”" in prompt
    assert "data-layout-container" in prompt
    assert "data-repeat-group/data-repeat-axis/data-repeat-order" in prompt


def test_png_only_channel_ingests_agent_candidate_and_builds_editable_pptx(tmp_path: Path):
    cases_root = tmp_path / "examples"
    reference = _reference(tmp_path)
    prepare_main(
        [
            str(reference),
            "--case",
            "png-to-pptx",
            "--cases-root",
            str(cases_root),
            "--source-mode",
            "png_reconstruct",
        ]
    )
    run = common.open_run(cases_root / "png-to-pptx")
    candidate = tmp_path / "agent-candidate.svg"
    candidate.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="160" height="100" '
        'viewBox="0 0 160 100"><rect id="panel" x="10" y="10" width="140" '
        'height="80" rx="8" fill="#EEF4FB" stroke="#23456A"/>'
        '<text id="label" x="80" y="50" text-anchor="middle" font-size="16">'
        'PNG-only</text></svg>',
        encoding="utf-8",
    )

    assert ingest_main([str(run.root), str(candidate), "--kind", "svg"]) == 0
    assert run.load_meta()["workflow"]["state"] == "candidate"
    summary = convert(run)
    assert summary["shape_count"] == 2
    assert read_json(run.bindings_path)["bindings_complete"] is True
    presentation = Presentation(run.pptx_path)
    assert len(presentation.slides) == 1
    assert any(shape.has_text_frame and shape.text == "PNG-only" for shape in presentation.slides[0].shapes)


def test_rejected_web_svg_switches_to_png_reconstruct_tasks(tmp_path: Path):
    run = common.create_run(
        _reference(tmp_path),
        case="rejected-svg",
        cases_root=tmp_path / "examples",
    )
    assert ingest_main(
        [str(run.root), "--rejected", "--fallback", "png_reconstruct"]
    ) == 0
    meta = run.load_meta()
    assert meta["source_mode"] == "png_reconstruct"
    assert meta["fidelity_profile"] == "hybrid_fidelity"
    assert meta["workflow"]["state"] == "repairing"
    assert read_json(run.region_tasks_path)["tasks"]
