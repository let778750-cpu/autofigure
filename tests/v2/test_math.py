"""math 命令测试：检测规则与 LaTeX 重建（纯函数）+ pptx 级注入集成（需公式引擎环境）。"""

from __future__ import annotations

import json
import zipfile
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation

from tools.v2 import common
from tools.v2.convert import convert
from tools.v2.math import classify_runs, rebuild_latex, upgrade

XSL_PATH = Path(r"C:\Program Files\Microsoft Office\root\Office16\MML2OMML.XSL")


def _engine_ready() -> bool:
    try:
        import latex2mathml  # noqa: F401
    except ImportError:
        return False
    return XSL_PATH.is_file()


requires_engine = pytest.mark.skipif(
    not _engine_ready(), reason="需要 latex2mathml 与本机 Office MML2OMML.XSL"
)

# run 元组 = (text, italic, baseline)，baseline ∈ {None, '30000'(super), '-25000'(sub)}


# ---------------------------------------------------------------- 检测规则（纯函数）


def test_classify_strong_signal_on_baseline_runs():
    runs = [("z", True, None), ("τ", True, "30000"), ("t", True, "-25000")]
    assert classify_runs(runs) == "strong"
    assert classify_runs([("θ", True, "-25000")]) == "strong"


def test_classify_weak_signal_single_math_letter():
    assert classify_runs([("τ", True, None)]) == "weak"
    assert classify_runs([("π", True, None)]) == "weak"
    assert classify_runs([(" τ ", True, None)]) == "weak"  # 去空格判定
    assert classify_runs([("𝔼", True, None), ("x", True, None)]) == "weak"


def test_classify_excludes_plain_ascii_labels():
    assert classify_runs([("GT Answers", True, None)]) is None
    assert classify_runs([("B", True, None)]) is None
    assert classify_runs([("I", True, None)]) is None
    assert classify_runs([("Questions", True, None)]) is None


def test_classify_excludes_non_italic_and_long_text():
    assert classify_runs([("τ", False, None)]) is None  # 非斜体
    assert classify_runs([("τ", True, None), ("x", False, None)]) is None  # 部分非斜体
    assert classify_runs([("向量机模型", True, None)]) is None  # 去空格 >4 字符
    assert classify_runs([]) is None
    assert classify_runs([("   ", True, None)]) is None


# ---------------------------------------------------------------- LaTeX 重建（纯函数）


def test_rebuild_super_then_sub():
    runs = [("z", True, None), ("τ", True, "30000"), ("t", True, "-25000")]
    assert rebuild_latex(runs) == r"z^{\tau}_{t}"


def test_rebuild_hat_and_mathbb():
    runs = [("Î = 𝔼", True, None), ("θ", True, "-25000"), ("(I)", True, None)]
    assert rebuild_latex(runs) == r"\hat{I} = \mathbb{E}_{\theta}(I)"


def test_rebuild_angle_bracket_list():
    runs = [
        ("⟨q", True, None), ("1", True, "-25000"),
        (", q", True, None), ("2", True, "-25000"),
        (" ... q", True, None), ("n", True, "-25000"),
        ("⟩", True, None),
    ]
    assert rebuild_latex(runs) == r"\langle q_{1}, q_{2} ... q_{n}\rangle"


def test_rebuild_nabla_mathcal_with_two_subs():
    runs = [("∇", True, None), ("θ", True, "-25000"), (" ℒ", True, None), ("ext", True, "-25000")]
    assert rebuild_latex(runs) == r"\nabla_{\theta} \mathcal{L}_{ext}"


def test_rebuild_merges_consecutive_same_direction():
    runs = [("x", True, None), ("a", True, "30000"), ("b", True, "30000")]
    assert rebuild_latex(runs) == r"x^{ab}"


def test_rebuild_escapes_latex_specials():
    assert rebuild_latex([("a_b", True, None)]) == r"a\_b"
    assert rebuild_latex([("100%", True, None)]) == r"100\%"
    assert rebuild_latex([("a^b", True, None)]) == r"a\^{}b"
    assert rebuild_latex([("{x}", True, None)]) == r"\{x\}"


def test_rebuild_unmapped_unicode_retained():
    assert rebuild_latex([("x→y", True, None)]) == "x→y"  # 映射表外字符原样保留


def test_rebuild_command_followed_by_letter_gets_space():
    assert rebuild_latex([("τt", True, None)]) == r"\tau t"  # 防命令名粘连成 \taut


# ---------------------------------------------------------------- pptx 级集成


@pytest.fixture()
def run_factory(tmp_path: Path):
    def make(svg: str, size: tuple[int, int] = (400, 200)) -> common.Run:
        source = tmp_path / "ref.png"
        Image.new("RGB", size, (240, 240, 240)).save(source)
        run = common.create_run(
            source,
            case="case",
            cases_root=tmp_path / "examples",
            input_route="svg-seeded",
        )
        run.qa_dir.mkdir(exist_ok=True)
        run.redraw_svg.write_text(svg, encoding="utf-8")
        return run

    return make


FORMULA_SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" width="400" height="200" viewBox="0 0 400 200">'
    '<rect x="0" y="0" width="400" height="200" fill="#FFFFFF"/>'
    '<text x="20" y="50" font-size="17" font-style="italic" fill="#1F3864">'
    "<tspan>z</tspan>"
    '<tspan baseline-shift="super" font-size="11">τ</tspan>'
    '<tspan baseline-shift="sub" font-size="11">t</tspan></text>'
    '<text x="20" y="100" font-size="18" font-style="italic" fill="#000000">τ</text>'
    '<text x="20" y="150" font-size="18" font-style="italic" fill="#000000">GT Answers</text>'
    "</svg>"
)


def _slide_xml(run: common.Run) -> str:
    with zipfile.ZipFile(run.pptx_path) as package:
        return package.read("ppt/slides/slide1.xml").decode("utf-8")


def test_dry_run_does_not_touch_pptx(run_factory):
    run = run_factory(FORMULA_SVG)
    convert(run)
    before = run.pptx_path.read_bytes()
    summary = upgrade(run, dry_run=True)
    assert summary["detected"] == 2  # 强信号 z^τ_t + 弱信号 τ；GT Answers 排除
    assert summary["injected"] == 0
    assert run.pptx_path.read_bytes() == before
    assert not (run.qa_dir / "math").exists()
    on_disk = json.loads((run.qa_dir / "math-summary.json").read_text(encoding="utf-8"))
    assert on_disk["dry_run"] is True
    assert on_disk["formulas"][0]["latex"] == r"z^{\tau}_{t}"
    assert on_disk["formulas"][0]["status"] == "detected"


def test_no_formula_leaves_pptx_untouched(run_factory):
    run = run_factory(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" viewBox="0 0 200 100">'
        '<text x="10" y="50" font-size="16" font-style="italic">GT Answers</text></svg>',
        size=(200, 100),
    )
    convert(run)
    before = run.pptx_path.read_bytes()
    summary = upgrade(run)
    assert summary["detected"] == 0
    assert summary["formulas"] == []
    assert run.pptx_path.read_bytes() == before


@requires_engine
def test_upgrade_injects_omml(run_factory):
    run = run_factory(FORMULA_SVG)
    convert(run)
    summary = upgrade(run)
    assert summary["detected"] == 2
    assert summary["injected"] == 2
    assert summary["failed"] == 0

    reopened = Presentation(run.pptx_path)  # 注入后必须仍能打开
    assert len(reopened.slides._sldIdLst) == 1
    xml = _slide_xml(run)
    assert 'name="math:001"' in xml
    assert 'name="math:002"' in xml
    assert xml.count("<m:oMath") == 2  # 每个公式一个 OMML 根（inline）
    assert "a14:m" in xml and "AlternateContent" in xml
    assert "GT Answers" in xml  # 非公式标签原文保留、未包裹

    bindings = json.loads(run.bindings_path.read_text(encoding="utf-8"))
    assert bindings["artifact_sha256"] == common.sha256_file(run.pptx_path)
    assert bindings["bindings_complete"] is True
    native_math = [item for item in bindings["bindings"] if item.get("native_math")]
    assert {item["shape_name"] for item in native_math} == {"math:001", "math:002"}
    scene = json.loads(run.scene_path.read_text(encoding="utf-8"))
    assert scene["artifact"]["sha256"] == bindings["artifact_sha256"]
    assert sum(item.get("native_math") is True for item in scene["elements"]) == 2
    assert run.load_meta()["workflow"]["state"] == "candidate"

    plan = json.loads((run.qa_dir / "math" / "plan.json").read_text(encoding="utf-8"))
    assert plan["schema_version"] == "1.0"
    assert len(plan["operations"]) == 2
    first = plan["operations"][0]
    assert first["placeholder_name"] == "math:001"
    assert first["target_font_size_pt"] == 12.75
    assert first["target_font_color"] == "#1F3864"
    receipt = json.loads((run.qa_dir / "math" / "EQ001.json").read_text(encoding="utf-8"))
    assert receipt["status"] == "PASS"
    assert receipt["canonical_latex"] == r"z^{\tau}_{t}"
    on_disk = json.loads((run.qa_dir / "math-summary.json").read_text(encoding="utf-8"))
    assert on_disk["injected"] == 2


@requires_engine
def test_upgrade_rerun_is_noop(run_factory):
    run = run_factory(FORMULA_SVG)
    convert(run)
    assert upgrade(run)["injected"] == 2
    second = upgrade(run)  # 注入后形状在 AlternateContent 内，不再被检出 → 幂等
    assert second["detected"] == 0
    assert second["injected"] == 0
    assert Presentation(run.pptx_path) is not None


@requires_engine
def test_compile_failure_keeps_original_textbox(run_factory, monkeypatch):
    run = run_factory(FORMULA_SVG)
    convert(run)
    from tools import powerpoint_native_math

    real_compile = powerpoint_native_math.compile_formula

    def flaky(formula_id, latex, mode, *args, **kwargs):
        if latex == r"\tau":
            raise powerpoint_native_math.NativeMathError("probe failure")
        return real_compile(formula_id, latex, mode, *args, **kwargs)

    monkeypatch.setattr(powerpoint_native_math, "compile_formula", flaky)
    summary = upgrade(run)
    assert summary["injected"] == 1
    assert summary["failed"] == 1
    failed = [row for row in summary["formulas"] if row["status"] == "failed"]
    assert len(failed) == 1
    assert "probe failure" in failed[0]["error"]
    xml = _slide_xml(run)
    assert 'name="math:001"' in xml  # 成功框已改名注入
    assert "τ" in xml  # 失败框原文保留（未改名未注入）
