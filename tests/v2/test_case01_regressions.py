from __future__ import annotations

import hashlib
import io
import json
import xml.etree.ElementTree as ET
from pathlib import Path

from PIL import Image, ImageChops
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from tools.v2 import common
from tools.v2.layout import audit_layout


CASE = Path(__file__).resolve().parents[2] / "examples" / "01-modular-agent"


def test_case01_mapping_arrow_is_above_imagination_panel():
    root = ET.parse(CASE / "redraw.svg").getroot()
    direct_ids = [element.get("id") for element in root]
    assert direct_ids.index("mapping-to-imagination") > direct_ids.index(
        "task-conditioned-imagination"
    )

    shapes = list(Presentation(CASE / "redraw.pptx").slides[0].shapes)
    z_order = {shape.name: index for index, shape in enumerate(shapes)}
    assert z_order["af-mapping-to-imagination-connector-01"] > z_order[
        "af-task-conditioned-imagination-rect-01"
    ]


def test_case01_globe_is_exact_authorized_reference_crop():
    assets = json.loads((CASE / "assets.json").read_text(encoding="utf-8"))
    asset = next(
        item for item in assets["assets"] if item["id"] == "atomic:environment-globe"
    )
    assert asset["authorized"] is True
    assert asset["editable"] is False
    assert asset["source"] == "reference_crop"

    presentation = Presentation(CASE / "redraw.pptx")
    picture = next(
        shape
        for shape in presentation.slides[0].shapes
        if shape.name == "af-atomic-environment-globe-atomic-raster-01"
    )
    assert picture.shape_type == MSO_SHAPE_TYPE.PICTURE
    embedded = picture.image.blob
    assert hashlib.sha256(embedded).hexdigest() == asset["source_sha256"]

    x, y, width, height = asset["bbox"]
    with Image.open(CASE / "reference.png") as reference:
        expected_image = reference.crop((x, y, x + width, y + height)).convert("RGBA")
        expected_buffer = io.BytesIO()
        expected_image.save(expected_buffer, format="PNG")
    with Image.open(io.BytesIO(embedded)) as actual:
        difference = ImageChops.difference(actual.convert("RGBA"), expected_image)
    assert difference.getbbox() is None


def test_case01_mapping_and_repeated_vectors_pass_explicit_layout_contracts():
    report = audit_layout(common.open_run(CASE))
    assert report["findings"] == []
    assert report["pass"] is True

    containment = {item["element"]: item for item in report["containment"]}
    for element_id in ("task-mapping-label", "task-mapping-formula"):
        row = containment[element_id]
        assert max(row["backend"]["overflow_px"].values()) <= row["tolerance_px"]

    groups = {item["id"]: item for item in report["repeat_groups"]}
    assert groups["e-v-stack"]["source"]["steps_px"] == [38.0, 37.0]
    assert groups["e-v-stack"]["backend"]["steps_px"] == [38.0, 37.0]
    assert groups["s-t-stack"]["source"]["steps_px"] == [38.0, 37.0]
    assert groups["s-t-stack"]["backend"]["steps_px"] == [38.0, 37.0]
