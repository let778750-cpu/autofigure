from __future__ import annotations

from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE

from tools import common
from tools.contracts import read_json, write_json
from tools.convert import convert
from tools.arrow_spec import path_from_segments
from tools.pptx_arrows import (
    _geometry_comparison,
    _source_path_comparison,
    refresh_bindings,
    strict_blockers,
    write_arrow_reports,
)


def _run(tmp_path: Path, svg: str, size: tuple[int, int]) -> common.Run:
    reference = tmp_path / "reference.png"
    Image.new("RGB", size, "white").save(reference)
    run = common.create_run(
        reference,
        case="arrows",
        cases_root=tmp_path / "examples",
        input_route="svg-seeded",
    )
    run.redraw_svg.write_text(svg, encoding="utf-8")
    convert(run)
    return run


def test_polyline_readback_compares_every_internal_point():
    expected = path_from_segments(
        [("M", 0, 0), ("L", 10, 0), ("L", 20, 0), ("L", 30, 0), ("L", 40, 0)]
    )
    detour = path_from_segments(
        [("M", 0, 0), ("L", 10, 0), ("L", 20, 50), ("L", 30, 0), ("L", 40, 0)]
    )
    comparison = _geometry_comparison(expected, detour, diagonal=100)
    assert comparison["status"] == "FAIL"
    assert comparison["full_path_max_point_error_px"] == 50.0


def test_cubic_readback_compares_every_control_point():
    expected = path_from_segments(
        [
            ("M", 0, 0),
            ("C", 10, 0, 20, 0, 30, 0),
            ("C", 40, 0, 50, 0, 60, 0),
        ]
    )
    detour = path_from_segments(
        [
            ("M", 0, 0),
            ("C", 10, 0, 20, 50, 30, 50),
            ("C", 40, 50, 50, 0, 60, 0),
        ]
    )
    comparison = _geometry_comparison(expected, detour, diagonal=100)
    assert comparison["status"] == "FAIL"
    assert comparison["full_path_max_point_error_px"] == 50.0


def test_source_path_comparison_rejects_straight_endpoint_drift():
    expected = path_from_segments([("M", 0, 0), ("L", 40, 0)])
    drifted = path_from_segments([("M", 0, 0), ("L", 41, 0)])
    comparison = _source_path_comparison(expected, drifted)
    assert comparison["status"] == "FAIL"
    assert comparison["max_point_error_px"] == 1.0


def test_source_path_comparison_rejects_polyline_signature_and_point_drift():
    expected = path_from_segments(
        [("M", 0, 0), ("L", 20, 0), ("L", 40, 0)]
    )
    extra_point = path_from_segments(
        [("M", 0, 0), ("L", 10, 0), ("L", 20, 0), ("L", 40, 0)]
    )
    drifted = path_from_segments(
        [("M", 0, 0), ("L", 20, 1), ("L", 40, 0)]
    )
    assert _source_path_comparison(expected, extra_point)["reason"] == (
        "source-path-command-signature"
    )
    assert _source_path_comparison(expected, drifted)["status"] == "FAIL"


def test_source_path_comparison_rejects_cubic_control_point_drift():
    expected = path_from_segments(
        [("M", 0, 0), ("C", 10, 0, 30, 0, 40, 0)]
    )
    drifted = path_from_segments(
        [("M", 0, 0), ("C", 10, 0, 30, 1, 40, 0)]
    )
    comparison = _source_path_comparison(expected, drifted)
    assert comparison["status"] == "FAIL"
    assert comparison["max_point_error_px"] == 1.0


def test_ninety_native_endpoint_combinations_roundtrip_exactly(tmp_path: Path):
    heads = ("open", "triangle", "stealth", "diamond", "oval")
    sizes = ("sm", "med", "lg")
    lines: list[str] = []
    row = 0
    for side in ("start", "end"):
        for head_type in heads:
            for width in sizes:
                for length in sizes:
                    y = 10 + row * 10
                    lines.append(
                        f'<line id="a-{row}" x1="30" y1="{y}" x2="270" y2="{y}" '
                        f'stroke="#000000" stroke-width="1" marker-{side}="url(#head)" '
                        f'data-{side}-arrow-type="{head_type}" '
                        f'data-{side}-arrow-width="{width}" '
                        f'data-{side}-arrow-length="{length}"/>'
                    )
                    row += 1
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="300" height="920" viewBox="0 0 300 920">'
        '<defs><marker id="head" markerWidth="4" markerHeight="4" refX="4" refY="2" '
        'orient="auto-start-reverse" markerUnits="userSpaceOnUse" viewBox="0 0 4 4">'
        '<path d="M0 0 L4 2 L0 4 Z" fill="#000000"/></marker></defs>'
        + "".join(lines)
        + "</svg>"
    )
    run = _run(tmp_path, svg, (300, 920))
    compile_report = read_json(run.arrow_compile_report_path)
    readback = read_json(run.powerpoint_arrow_readback_path)
    assert compile_report["arrow_count"] == 90
    assert compile_report["pass"] is True
    assert readback["arrow_count"] == 90
    assert readback["pass"] is True
    assert all(row["embedded_tag_pass"] for row in readback["records"])


def test_all_twelve_dash_styles_compile_to_exact_body_semantics(tmp_path: Path):
    dashes = (
        "solid",
        "square_dot",
        "round_dot",
        "dash",
        "dash_dot",
        "dash_dot_dot",
        "long_dash",
        "long_dash_dot",
        "long_dash_dot_dot",
        "sys_dash",
        "sys_dot",
        "sys_dash_dot",
    )
    lines = [
        f'<line id="d-{index}" x1="20" y1="{10 + index * 10}" x2="180" '
        f'y2="{10 + index * 10}" stroke="#000000" marker-end="url(#head)" '
        f'data-ppt-dash="{dash}"/>'
        for index, dash in enumerate(dashes)
    ]
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="140" viewBox="0 0 200 140">'
        '<defs><marker id="head" markerWidth="4" markerHeight="4" refX="4" refY="2" '
        'orient="auto" markerUnits="userSpaceOnUse" viewBox="0 0 4 4">'
        '<path d="M0 0 L4 2 L0 4 Z" fill="#000000"/></marker></defs>'
        + "".join(lines)
        + "</svg>"
    )
    run = _run(tmp_path, svg, (200, 140))
    readback = read_json(run.powerpoint_arrow_readback_path)
    assert readback["pass"] is True
    actual = {row["body"]["actual"]["dash"] for row in readback["records"]}
    assert actual == set(dashes)


def test_color_mismatched_custom_head_fails_before_artifact_creation(tmp_path: Path):
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" viewBox="0 0 200 100">'
        '<defs><marker id="head" markerWidth="4" markerHeight="4" refX="4" refY="2" '
        'orient="auto" markerUnits="userSpaceOnUse" viewBox="0 0 4 4">'
        '<path d="M0 0 L4 2 L0 4 Z" fill="#FF0000"/></marker></defs>'
        '<line id="loss" x1="20" y1="50" x2="180" y2="50" stroke="#000000" '
        'marker-end="url(#head)"/></svg>'
    )
    reference = tmp_path / "reference.png"
    Image.new("RGB", (200, 100), "white").save(reference)
    run = common.create_run(
        reference,
        case="arrows",
        cases_root=tmp_path / "examples",
        input_route="svg-seeded",
    )
    run.redraw_svg.write_text(svg, encoding="utf-8")

    with pytest.raises(SystemExit, match="exactly one visible PowerPoint object"):
        convert(run)
    assert not run.pptx_path.exists()
    assert not run.arrow_compile_report_path.exists()


def test_closed_block_arrow_is_one_editable_freeform(tmp_path: Path):
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" viewBox="0 0 200 100">'
        '<polygon id="block" data-arrow-representation="block_arrow" '
        'data-arrow-centerline="M20 50 L180 50" data-end-head-type="custom" '
        'points="20,40 130,40 130,25 180,50 130,75 130,60 20,60" '
        'fill="#3366CC" stroke="#224488" stroke-width="1"/></svg>'
    )
    run = _run(tmp_path, svg, (200, 100))
    compile_report = read_json(run.arrow_compile_report_path)
    readback = read_json(run.powerpoint_arrow_readback_path)
    assert compile_report["pass"] is True
    assert compile_report["records"][0]["strategy"] == "single-closed-freeform"
    assert compile_report["records"][0]["visible_object_count"] == 1
    assert readback["pass"] is True
    head = readback["records"][0]["end_head"]
    assert head["expected"]["type"] == "custom"
    assert head["actual"]["type"] == "custom"
    assert head["actual"]["representation"] == "embedded-silhouette"
    assert head["actual"]["ooxml_line_end"]["type"] == "none"
    centerline = readback["records"][0]["semantic_centerline"]
    assert centerline["status"] == "PASS"
    assert centerline["actual_start"] == [20.0, 50.0]
    assert centerline["actual_end"] == [180.0, 50.0]
    assert centerline["evidence"].startswith("saved-pptx-cNvPr-description")


def test_fixed_block_arrow_keeps_declared_topology_in_semantic_tag(tmp_path: Path):
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" viewBox="0 0 200 100">'
        '<rect id="source" x="5" y="35" width="20" height="30" fill="#eeeeee"/>'
        '<rect id="target" x="175" y="35" width="20" height="30" fill="#eeeeee"/>'
        '<polygon id="block" data-arrow-representation="block_arrow" '
        'data-arrow-centerline="M25 50 L175 50" data-arrow-routing="fixed" '
        'data-arrow-topology="declared" data-source-id="source" data-target-id="target" '
        'data-end-head-type="custom" points="25,42 140,42 140,30 175,50 140,70 140,58 25,58" '
        'fill="#3366CC" stroke="#224488" stroke-width="1"/></svg>'
    )
    run = _run(tmp_path, svg, (200, 100))
    scene = read_json(run.scene_path)
    spec = next(item for item in scene["edges"] if item["id"] == "block")["arrow_spec"]
    assert spec["routing"] == "fixed"
    assert spec["topology"]["mode"] == "declared"
    readback = read_json(run.powerpoint_arrow_readback_path)
    record = next(item for item in readback["records"] if item["element_id"] == "block")
    assert record["topology"]["expected"]["mode"] == "declared"
    assert record["topology"]["pass"] is True
    assert record["embedded_tag_pass"] is True


def test_closed_block_freeform_without_explicit_centerline_fails_closed(
    tmp_path: Path,
):
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" '
        'viewBox="0 0 200 100"><polygon id="block" '
        'data-arrow-representation="block_arrow" '
        'points="20,40 130,40 130,25 180,50 130,75 130,60 20,60" '
        'fill="#3366CC" stroke="none"/></svg>'
    )

    with pytest.raises(SystemExit, match="requires data-arrow-centerline"):
        _run(tmp_path, svg, (200, 100))


def test_arrow_reports_are_bound_to_current_artifact_hash(tmp_path: Path):
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" viewBox="0 0 200 100">'
        '<defs><marker id="head" markerWidth="4" markerHeight="4" refX="4" refY="2" '
        'orient="auto" markerUnits="userSpaceOnUse" viewBox="0 0 4 4">'
        '<path d="M0 0 L4 2 L0 4 Z" fill="#000000"/></marker></defs>'
        '<line id="a" x1="20" y1="50" x2="180" y2="50" stroke="#000000" '
        'marker-end="url(#head)"/></svg>'
    )
    run = _run(tmp_path, svg, (200, 100))
    with run.pptx_path.open("ab") as stream:
        stream.write(b"artifact-drift")
    assert "arrow:A20_ARTIFACT_IDENTITY" in strict_blockers(run)


def test_refresh_bindings_requires_shape_id_and_name_pair(tmp_path: Path):
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="120" height="100" '
        'viewBox="0 0 120 100"><rect id="box" x="10" y="10" width="100" '
        'height="80" fill="#ffffff" stroke="#111111"/></svg>'
    )
    run = _run(tmp_path, svg, (120, 100))
    bindings = read_json(run.bindings_path)
    bindings["bindings"][0]["shape_name"] = "wrong-name-with-valid-id"
    write_json(run.bindings_path, bindings)

    summary = refresh_bindings(run, host_saved_reopened=True)

    assert summary["bindings_complete"] is False
    assert read_json(run.bindings_path)["bindings"][0]["readback_found"] is False


def test_arrow_readback_requires_shape_id_and_name_pair(tmp_path: Path):
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" '
        'viewBox="0 0 200 100"><defs><marker id="head" markerWidth="4" '
        'markerHeight="4" refX="4" refY="2" orient="auto" '
        'markerUnits="userSpaceOnUse" viewBox="0 0 4 4"><path '
        'd="M0 0 L4 2 L0 4 Z" fill="#000000"/></marker></defs><line id="a" '
        'x1="20" y1="50" x2="180" y2="50" stroke="#000000" '
        'marker-end="url(#head)"/></svg>'
    )
    run = _run(tmp_path, svg, (200, 100))
    bindings = read_json(run.bindings_path)
    arrow_binding = next(
        row for row in bindings["bindings"] if row.get("element_id") == "a"
    )
    arrow_binding["shape_name"] = "different-name-with-valid-id"
    write_json(run.bindings_path, bindings)

    _, readback = write_arrow_reports(run)

    assert readback["pass"] is False
    assert "arrow:A15_READBACK_MISSING:a" in readback["blockers"]


def test_arrow_source_evidence_must_match_current_run(tmp_path: Path):
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" '
        'viewBox="0 0 200 100"><defs><marker id="head" markerWidth="4" '
        'markerHeight="4" refX="4" refY="2" orient="auto" '
        'markerUnits="userSpaceOnUse" viewBox="0 0 4 4"><path '
        'd="M0 0 L4 2 L0 4 Z" fill="#000000"/></marker></defs><line id="a" '
        'x1="20" y1="50" x2="180" y2="50" stroke="#000000" '
        'marker-end="url(#head)"/></svg>'
    )
    run = _run(tmp_path, svg, (200, 100))
    scene = read_json(run.scene_path)
    for item in [*scene["elements"], *scene["edges"]]:
        if item.get("id") == "a":
            item["arrow_spec"]["source_evidence"]["reference_sha256"] = "f" * 64
    write_json(run.scene_path, scene)

    compile_report, _ = write_arrow_reports(run)

    assert compile_report["pass"] is False
    assert any(
        blocker.endswith("source-reference-sha256-mismatch")
        for blocker in compile_report["blockers"]
    )


def test_compile_report_blocks_scene_geometry_drift_even_when_spec_and_pptx_agree(
    tmp_path: Path,
):
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" '
        'viewBox="0 0 200 100"><defs><marker id="head" markerWidth="4" '
        'markerHeight="4" refX="4" refY="2" orient="auto" '
        'markerUnits="userSpaceOnUse" viewBox="0 0 4 4"><path '
        'd="M0 0 L4 2 L0 4 Z" fill="#000000"/></marker></defs><line id="a" '
        'x1="20" y1="50" x2="180" y2="50" stroke="#000000" '
        'marker-end="url(#head)"/></svg>'
    )
    run = _run(tmp_path, svg, (200, 100))
    scene = read_json(run.scene_path)
    arrow = next(item for item in scene["elements"] if item.get("id") == "a")
    arrow["geometry"]["x2"] = "181"
    write_json(run.scene_path, scene)

    compile_report, readback = write_arrow_reports(run)

    blocker = "arrow:A10_SOURCE_PATH_DRIFT:a:scene"
    assert blocker in compile_report["blockers"]
    assert compile_report["pass"] is False
    record = next(
        item for item in compile_report["records"] if item["element_id"] == "a"
    )
    assert record["source_geometry"]["scene"]["max_point_error_px"] == 1.0
    assert record["source_geometry"]["svg"]["status"] == "PASS"
    assert readback["pass"] is True


def test_idless_svg_arrows_use_the_converter_stable_source_identity(tmp_path: Path):
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" '
        'viewBox="0 0 200 100"><defs><marker id="head" markerWidth="4" '
        'markerHeight="4" refX="4" refY="2" orient="auto" '
        'markerUnits="userSpaceOnUse" viewBox="0 0 4 4"><path '
        'd="M0 0 L4 2 L0 4 Z" fill="#000000"/></marker></defs>'
        '<line x1="20" y1="30" x2="180" y2="30" stroke="#000000" '
        'marker-end="url(#head)"/><g transform="translate(0 20)"><line '
        'x1="20" y1="30" x2="180" y2="30" stroke="#000000" '
        'marker-end="url(#head)"/></g></svg>'
    )
    run = _run(tmp_path, svg, (200, 100))

    compile_report = read_json(run.arrow_compile_report_path)

    assert compile_report["pass"] is True
    assert [row["element_id"] for row in compile_report["records"]] == [
        "svg-line-0001",
        "svg-line-0002",
    ]
    assert all(
        row["source_geometry"]["status"] == "PASS"
        for row in compile_report["records"]
    )


def test_native_left_right_block_autoshape_roundtrips_exactly(tmp_path: Path):
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" '
        'viewBox="0 0 200 100"><polygon id="block" data-role="edge" '
        'data-arrow-representation="block_arrow" data-ppt-autoshape="leftRightArrow" '
        'data-ppt-adjustments="0.529411765 0.470588235" '
        'data-arrow-body-width="9" data-start-head-type="triangle" '
        'data-end-head-type="triangle" data-start-arrow-width="lg" '
        'data-end-arrow-width="lg" data-start-arrow-length="sm" '
        'data-end-arrow-length="sm" '
        'points="20,50 28,41.5 28,45.5 172,45.5 172,41.5 180,50 '
        '172,58.5 172,54.5 28,54.5 28,58.5" '
        'fill="#767171" stroke="none"/></svg>'
    )
    run = _run(tmp_path, svg, (200, 100))
    compile_report = read_json(run.arrow_compile_report_path)
    scene = read_json(run.scene_path)
    readback = read_json(run.powerpoint_arrow_readback_path)
    record = next(item for item in readback["records"] if item["element_id"] == "block")
    spec = next(item for item in scene["edges"] if item["id"] == "block")[
        "arrow_spec"
    ]

    assert compile_report["pass"] is True
    assert compile_report["records"][0]["strategy"] == "native-block-autoshape"
    assert compile_report["records"][0]["visible_object_count"] == 1
    assert record["ooxml_kind"] == "shape"
    assert record["start_head"]["expected"]["type"] == "triangle"
    assert record["end_head"]["expected"]["type"] == "triangle"
    assert record["autoshape"]["actual"]["subtype"] == "leftRightArrow"
    assert record["autoshape"]["pass"] is True
    assert record["body"]["actual"]["color"] == "#767171"
    assert record["body"]["width_pass"] is True
    assert spec["path"] == {
        "kind": "straight",
        "coordinate_space": "canvas",
        "points": [{"x": 20.0, "y": 50.0}, {"x": 180.0, "y": 50.0}],
    }
    assert readback["pass"] is True


def test_native_block_autoshape_readback_rejects_centerline_tamper(tmp_path: Path):
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" '
        'viewBox="0 0 200 100"><polygon id="block" data-role="edge" '
        'data-arrow-representation="block_arrow" data-ppt-autoshape="leftRightArrow" '
        'data-ppt-adjustments="0.529411765 0.470588235" '
        'data-arrow-body-width="9" data-start-head-type="triangle" '
        'data-end-head-type="triangle" data-start-arrow-width="lg" '
        'data-end-arrow-width="lg" data-start-arrow-length="sm" '
        'data-end-arrow-length="sm" '
        'points="20,50 28,41.5 28,45.5 172,45.5 172,41.5 180,50 '
        '172,58.5 172,54.5 28,54.5 28,58.5" '
        'fill="#767171" stroke="none"/></svg>'
    )
    run = _run(tmp_path, svg, (200, 100))
    scene = read_json(run.scene_path)
    tampered_path = {
        "kind": "straight",
        "coordinate_space": "canvas",
        "points": [{"x": 20.0, "y": 60.0}, {"x": 180.0, "y": 60.0}],
    }
    for collection in ("elements", "edges"):
        edge = next(item for item in scene[collection] if item.get("id") == "block")
        edge["arrow_spec"]["path"] = tampered_path
    write_json(run.scene_path, scene)

    _, readback = write_arrow_reports(run)
    record = next(item for item in readback["records"] if item["element_id"] == "block")

    assert record["path_geometry"]["status"] == "FAIL"
    assert record["path_geometry"]["start_error_px"] == pytest.approx(10.0, abs=0.001)
    assert "arrow:A18_READBACK_PATH:block" in readback["blockers"]
    assert readback["pass"] is False


def test_left_right_block_autoshape_rejects_non_axis_aligned_silhouette(
    tmp_path: Path,
):
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" '
        'viewBox="0 0 200 100"><polygon id="block" data-role="edge" '
        'data-arrow-representation="block_arrow" data-ppt-autoshape="leftRightArrow" '
        'data-ppt-adjustments="0.5 0.5" data-arrow-body-width="9" '
        'data-start-head-type="triangle" data-end-head-type="triangle" '
        'points="20,50 28,40 30,46 172,44 172,40 180,50 '
        '172,60 172,54 28,56 28,60" fill="#767171" stroke="none"/></svg>'
    )
    with pytest.raises(SystemExit, match="axis-aligned"):
        _run(tmp_path, svg, (200, 100))


def test_overlaid_reciprocal_native_connectors_fail_composition_audit(tmp_path: Path):
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" '
        'viewBox="0 0 200 100"><defs><marker id="head" markerWidth="4" '
        'markerHeight="4" refX="4" refY="2" orient="auto" '
        'markerUnits="userSpaceOnUse" viewBox="0 0 4 4"><path '
        'd="M0 0 L4 2 L0 4 Z" fill="#7F7F7F"/></marker></defs>'
        '<line id="left" x1="180" y1="50" x2="20" y2="50" '
        'stroke="#7F7F7F" stroke-width="7" marker-end="url(#head)"/>'
        '<line id="right" x1="20" y1="50" x2="180" y2="50" '
        'stroke="#7F7F7F" stroke-width="7" marker-end="url(#head)"/></svg>'
    )
    run = _run(tmp_path, svg, (200, 100))
    compile_report = read_json(run.arrow_compile_report_path)
    readback = read_json(run.powerpoint_arrow_readback_path)
    composition = read_json(run.qa_dir / "arrow-composition-audit.json")

    assert "arrow:A21_SOURCE_RECIPROCAL_OVERLAP:left:right" in compile_report["blockers"]
    assert "arrow:A22_READBACK_RECIPROCAL_OVERLAP:left:right" in readback["blockers"]
    assert composition["pass"] is False
    assert composition["source_findings"][0]["arrow_ids"] == ["left", "right"]
    assert composition["readback_findings"][0]["arrow_ids"] == ["left", "right"]


def test_refresh_bindings_rejects_unbound_visible_object(tmp_path: Path):
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="120" height="100" '
        'viewBox="0 0 120 100"><rect id="box" x="10" y="10" width="100" '
        'height="80" fill="#ffffff" stroke="#111111"/></svg>'
    )
    run = _run(tmp_path, svg, (120, 100))
    presentation = Presentation(run.pptx_path)
    presentation.slides[0].shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 100, 100)
    presentation.save(run.pptx_path)

    summary = refresh_bindings(run, host_saved_reopened=True)

    assert summary["bindings_complete"] is False
    assert summary["unbound_object_count"] == 1
