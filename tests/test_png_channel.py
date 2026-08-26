from __future__ import annotations

from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation

from tools.core import common
from tools.core.contracts import read_json, write_json
from tools.pipeline.convert import convert
from tools.pipeline.ingest import main as ingest_main
from tools.pipeline.prepare import SVG_AUTHORING_CONTRACT
from tools.pipeline.prepare import main as prepare_main
from tools.assets.reference_inventory import OBJECT_KINDS, freeze_inventory


def _reference(tmp_path: Path) -> Path:
    path = tmp_path / "input-only.png"
    Image.new("RGB", (160, 100), "white").save(path)
    return path


def _seed(tmp_path: Path) -> Path:
    path = tmp_path / "external-seed.svg"
    path.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="160" height="100" '
        'viewBox="0 0 160 100"><rect id="panel" x="10" y="10" '
        'width="140" height="80"/></svg>',
        encoding="utf-8",
    )
    return path


def _freeze_panel_inventory(run: common.Run) -> None:
    regions = read_json(run.regions_path)
    regions["regions"] = [
        {
            "id": "panel-region",
            "label": "Panel and label",
            "bbox": [0, 0, 160, 100],
            "critical": True,
            "relations_exhaustive": True,
            "element_ids": ["panel", "label"],
        }
    ]
    inventory = regions["reference_inventory"]
    inventory["expected_counts"] = {kind: 0 for kind in OBJECT_KINDS}
    inventory["expected_counts"].update({"text": 1, "shape": 1})
    inventory["zero_count_authorizations"] = [
        {
            "kind": kind,
            "basis": "full-reference-review",
            "reviewer": "test-reviewer",
            "reference_sha256": run.load_meta()["source_sha256"],
        }
        for kind in ("arrow", "icon", "brace")
    ]
    inventory["objects"] = [
        {
            "id": "panel",
            "kind": "shape",
            "bbox": [10, 10, 140, 80],
            "element_ids": ["panel"],
            "critical_region_ids": ["panel-region"],
        },
        {
            "id": "label",
            "kind": "text",
            "bbox": [45, 35, 70, 22],
            "element_ids": ["label"],
            "critical_region_ids": ["panel-region"],
            "typography": {
                "exact_text": "PNG-only",
                "font_family": "Arial",
                "font_size_px": 16,
                "font_weight": "normal",
                "font_style": "normal",
                "line_count": 1,
                "alignment": "center",
                "bbox_tolerance_px": 2,
                "font_size_tolerance_px": 0.5,
            },
        },
    ]
    write_json(run.regions_path, regions)
    freeze_inventory(run)


def test_prepare_png_only_creates_tasks_without_web_svg_prerequisite(tmp_path: Path):
    cases_root = tmp_path / "examples"
    assert prepare_main(
        [
            str(_reference(tmp_path)),
            "--case",
            "png-only",
            "--cases-root",
            str(cases_root),
            "--input-route",
            "reference-only",
        ]
    ) == 0

    run = common.open_run(cases_root / "reference-only" / "png-only")
    meta = run.load_meta()
    assert meta["input_route"] == "reference-only"
    assert meta["processing_mode"] == "png_reconstruct"
    assert "source_mode" not in meta
    assert meta["fidelity_profile"] == "hybrid_fidelity"
    assert meta["workflow"]["state"] == "prepared"
    assert not run.redraw_svg.exists()
    tasks = read_json(run.region_tasks_path)
    assert tasks["reference_sha256"] == meta["source_sha256"]
    assert tasks["result_contract"]["offline_initial_render_carrier"] == "svg"
    assert tasks["result_contract"]["svg_authoring_contract"] == "prompt.md"
    assert tasks["tasks"]
    prompt = run.prompt_md.read_text(encoding="utf-8")
    assert "不要求、也不依赖 GPT Web" in prompt
    assert "不能把“入口已连通”误写成“PNG 已自动一比一重建”" in prompt
    assert 'viewBox="0 0 160 100"' in prompt
    assert "箭头的粗细、头部样式（实心/开放/块状）" in prompt
    assert "atomic:" in prompt
    assert "data-layout-container" in prompt
    assert "data-repeat-group" in prompt
    assert "data-repeat-axis" in prompt
    assert "data-repeat-order" in prompt
    assert "required_relations" in prompt
    assert (
        "source_id/target_id/direction/start_head_type/end_head_type/"
        "representation/visible_object_count"
        in prompt
    )
    assert "一个双端 ArrowSpec 和一个 PowerPoint 可见对象" in prompt
    assert "禁止拆成两条共线反向单头箭头" in prompt


def test_png_only_channel_ingests_agent_candidate_and_builds_editable_pptx(tmp_path: Path):
    cases_root = tmp_path / "examples"
    reference = _reference(tmp_path)
    prepare_main(
        [
            str(reference),
            "--case",
            "png-to-pptx",
            "--cases-root",
            str(cases_root),
            "--input-route",
            "reference-only",
        ]
    )
    run = common.open_run(cases_root / "reference-only" / "png-to-pptx")
    _freeze_panel_inventory(run)
    candidate = tmp_path / "agent-candidate.svg"
    candidate.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="160" height="100" '
        'viewBox="0 0 160 100"><rect id="panel" x="10" y="10" width="140" '
        'height="80" rx="8" fill="#EEF4FB" stroke="#23456A"/>'
        '<text id="label" x="80" y="50" text-anchor="middle" font-size="16">'
        'PNG-only</text></svg>',
        encoding="utf-8",
    )

    assert ingest_main(
        [
            str(run.root),
            str(candidate),
            "--kind",
            "svg",
            "--candidate-origin",
            "codex",
        ]
    ) == 0
    assert run.load_meta()["workflow"]["state"] == "candidate"
    provenance = read_json(run.provenance_path)
    assert provenance["external_svg_seed"] is None
    assert provenance["candidate_history"][0]["role"] == "reconstruction-candidate"
    summary = convert(run)
    assert summary["shape_count"] == 2
    assert read_json(run.bindings_path)["bindings_complete"] is True
    presentation = Presentation(run.pptx_path)
    assert len(presentation.slides) == 1
    assert any(shape.has_text_frame and shape.text == "PNG-only" for shape in presentation.slides[0].shapes)


def test_rejected_web_svg_switches_to_png_reconstruct_tasks(tmp_path: Path):
    run = common.create_run(
        _reference(tmp_path),
        case="rejected-svg",
        cases_root=tmp_path / "examples",
        input_route="svg-seeded",
    )
    assert ingest_main(
        [str(run.root), "--rejected", "--fallback", "png_reconstruct"]
    ) == 0
    meta = run.load_meta()
    assert meta["input_route"] == "svg-seeded"
    assert meta["processing_mode"] == "png_reconstruct"
    assert meta["fidelity_profile"] == "hybrid_fidelity"
    assert meta["workflow"]["state"] == "repairing"
    assert read_json(run.region_tasks_path)["tasks"]


def test_prepare_requires_explicit_input_route(tmp_path: Path):
    with pytest.raises(SystemExit) as exc:
        prepare_main(
            [
                str(_reference(tmp_path)),
                "--case",
                "missing-route",
                "--cases-root",
                str(tmp_path / "examples"),
            ]
        )
    assert exc.value.code == 2


def test_reference_only_rejects_external_seed_before_copy(tmp_path: Path):
    run = common.create_run(
        _reference(tmp_path),
        case="no-seed",
        cases_root=tmp_path / "examples",
        input_route="reference-only",
    )
    candidate = tmp_path / "external.svg"
    candidate.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="160" height="100" '
        'viewBox="0 0 160 100"/>',
        encoding="utf-8",
    )
    with pytest.raises(SystemExit, match="不能摄取 external-seed"):
        ingest_main(
            [
                str(run.root),
                str(candidate),
                "--candidate-role",
                "external-seed",
            ]
        )
    assert not run.redraw_svg.exists()


def test_both_input_routes_embed_the_same_svg_authoring_contract(tmp_path: Path):
    cases_root = tmp_path / "examples"
    reference = _reference(tmp_path)
    seed = _seed(tmp_path)
    for route, case in (
        ("svg-seeded", "seeded-contract-smoke"),
        ("reference-only", "png-only-contract-smoke"),
    ):
        arguments = [
            str(reference),
            "--case",
            case,
            "--cases-root",
            str(cases_root),
            "--input-route",
            route,
        ]
        if route == "svg-seeded":
            arguments.extend(["--seed", str(seed)])
        assert (
            prepare_main(arguments)
            == 0
        )

    contract = SVG_AUTHORING_CONTRACT.format(width=160, height=100)
    seeded_prompt = (
        cases_root / "svg-seeded" / "seeded-contract-smoke" / "prompt.md"
    ).read_text(encoding="utf-8")
    png_prompt = (
        cases_root / "reference-only" / "png-only-contract-smoke" / "prompt.md"
    ).read_text(encoding="utf-8")
    assert contract in seeded_prompt
    assert contract in png_prompt
