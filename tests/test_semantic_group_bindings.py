from __future__ import annotations

from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation

from tools import common
from tools.arrow_spec import validate_scene_arrow_specs
from tools.contracts import read_json, write_json
from tools.convert import convert
from tools.live_bridge import build_powerpoint_live_bridge
from tools.math import upgrade
from tools.pptx_arrows import refresh_bindings
from tools.region_contract import audit_region_contract


SVG = """<svg xmlns="http://www.w3.org/2000/svg" width="240" height="120"
  viewBox="0 0 240 120">
  <defs>
    <marker id="head" markerWidth="4" markerHeight="4" refX="4" refY="2"
      orient="auto" markerUnits="userSpaceOnUse" viewBox="0 0 4 4">
      <path d="M0 0 L4 2 L0 4 Z" fill="#223344"/>
    </marker>
  </defs>
  <g id="source-node" data-role="node">
    <rect id="source-body" x="15" y="30" width="55" height="55"
      rx="4" fill="#DDEEFF" stroke="#223344"/>
    <g id="source-icon" data-role="icon">
      <circle id="source-dot" cx="42.5" cy="57.5" r="8"
        fill="#88AACC" stroke="#223344"/>
    </g>
  </g>
  <g id="target-node" data-role="node">
    <rect id="target-body" x="170" y="30" width="55" height="55"
      rx="4" fill="#FFEEDD" stroke="#223344"/>
  </g>
  <line id="flow" x1="70" y1="57.5" x2="170" y2="57.5"
    stroke="#223344" stroke-width="2" marker-end="url(#head)"
    data-source-id="source-node" data-target-id="target-node"
    data-arrow-topology="declared"/>
</svg>"""

XSL_PATH = Path(r"C:\Program Files\Microsoft Office\root\Office16\MML2OMML.XSL")


def _math_engine_ready() -> bool:
    try:
        import latex2mathml  # noqa: F401
    except ImportError:
        return False
    return XSL_PATH.is_file()


def _run(tmp_path: Path) -> common.Run:
    reference = tmp_path / "reference.png"
    Image.new("RGB", (240, 120), "white").save(reference)
    run = common.create_run(
        reference,
        case="semantic-groups",
        cases_root=tmp_path / "examples",
        input_route="svg-seeded",
    )
    run.redraw_svg.write_text(SVG, encoding="utf-8")
    convert(run)
    return run


def _region_payload() -> dict:
    return {
        "regions": [
            {
                "id": "group-flow",
                "critical": True,
                "element_ids": ["source-node", "target-node", "flow"],
                "relations_exhaustive": True,
                "required_relations": [
                    {
                        "id": "flow",
                        "source_id": "source-node",
                        "target_id": "target-node",
                        "direction": "forward",
                        "start_head_type": "none",
                        "end_head_type": "triangle",
                        "representation": "line_arrow",
                        "visible_object_count": 1,
                    }
                ],
            }
        ]
    }


def test_svg_semantic_groups_close_scene_topology_and_composite_readback(
    tmp_path: Path,
) -> None:
    run = _run(tmp_path)
    scene = read_json(run.scene_path)
    by_id = {item["id"]: item for item in scene["elements"]}

    assert by_id["source-node"]["kind"] == "logical_group"
    assert by_id["source-node"]["member_ids"] == ["source-body", "source-dot"]
    assert by_id["source-node"]["logical_descendant_group_ids"] == [
        "source-icon"
    ]
    assert by_id["source-icon"]["member_ids"] == ["source-dot"]
    assert by_id["target-node"]["member_ids"] == ["target-body"]
    assert validate_scene_arrow_specs(scene) == []

    bindings = read_json(run.bindings_path)
    physical_ids = [row["element_id"] for row in bindings["bindings"]]
    assert physical_ids == ["source-body", "source-dot", "target-body", "flow"]
    assert "source-node" not in physical_ids
    assert "target-node" not in physical_ids
    source_dot = next(
        row for row in bindings["bindings"] if row["element_id"] == "source-dot"
    )
    assert source_dot["logical_group_ids"] == ["source-node", "source-icon"]

    groups = {
        row["element_id"]: row
        for row in bindings["logical_group_bindings"]
    }
    assert set(groups) == {"source-node", "source-icon", "target-node"}
    assert groups["source-node"]["binding_kind"] == "logical-group-composite"
    assert groups["source-node"]["member_element_ids"] == [
        "source-body",
        "source-dot",
    ]
    assert groups["source-node"]["visible_object_count"] == 2
    assert "shape_id" not in groups["source-node"]
    assert "shape_name" not in groups["source-node"]
    assert all(row["readback_found"] is True for row in groups.values())

    physical_shape_count = len(Presentation(run.pptx_path).slides[0].shapes)
    assert physical_shape_count == len(bindings["bindings"])
    assert sum(row["element_id"] == "flow" for row in bindings["bindings"]) == 1

    refreshed = refresh_bindings(run, host_saved_reopened=False)
    assert refreshed["bindings_complete"] is True
    assert refreshed["unbound_object_count"] == 0
    assert refreshed["binding_count"] == physical_shape_count
    assert refreshed["logical_group_binding_count"] == 3

    manifest = build_powerpoint_live_bridge(run)
    bridge_scene = read_json(
        run.root / manifest["case_root"] / "design" / "scene_graph.json"
    )
    source_node = next(
        item for item in bridge_scene["nodes"] if item["id"] == "source-node"
    )
    flow_edge = next(item for item in bridge_scene["edges"] if item["id"] == "flow")
    assert source_node["composite"] is True
    assert source_node["memberIds"] == ["source-body", "source-dot"]
    assert flow_edge["source"] == "source-node"
    assert flow_edge["target"] == "target-node"

    report = audit_region_contract(run, _region_payload())
    assert report["pass"] is True, report["blockers"]
    assert report["blockers"] == []


def test_logical_group_composite_identity_drift_fails_closed(tmp_path: Path) -> None:
    run = _run(tmp_path)
    bindings = read_json(run.bindings_path)
    source_group = next(
        row
        for row in bindings["logical_group_bindings"]
        if row["element_id"] == "source-node"
    )
    source_group["backend_object_names"][0] = "wrong-member-name"
    write_json(run.bindings_path, bindings)

    refreshed = refresh_bindings(run, host_saved_reopened=False)
    assert refreshed["bindings_complete"] is False
    assert refreshed["unbound_object_count"] == 0

    report = audit_region_contract(run, _region_payload())
    assert report["pass"] is False
    assert (
        "region-contract:group-flow:logical-group-identity-record-drift:source-node"
        in report["blockers"]
    )
    assert (
        "region-contract:group-flow:logical-group-backend-member-drift:source-node"
        in report["blockers"]
    )
    assert (
        "region-contract:group-flow:binding-readback-missing:source-node"
        in report["blockers"]
    )


def test_arrow_semantics_cannot_hide_a_multi_object_group_fallback(
    tmp_path: Path,
) -> None:
    reference = tmp_path / "arrow-group-reference.png"
    Image.new("RGB", (120, 100), "white").save(reference)
    run = common.create_run(
        reference,
        case="forbidden-arrow-group",
        cases_root=tmp_path / "examples",
        input_route="svg-seeded",
    )
    run.redraw_svg.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="120" height="100" '
        'viewBox="0 0 120 100"><g id="split-arrow" data-role="arrow">'
        '<line id="shaft" x1="10" y1="50" x2="90" y2="50" '
        'stroke="#222222"/><polygon id="head" points="90,44 110,50 90,56" '
        'fill="#222222"/></g></svg>',
        encoding="utf-8",
    )

    with pytest.raises(SystemExit, match="one logical arrow must compile"):
        convert(run)
    assert not run.pptx_path.exists()


@pytest.mark.skipif(
    not _math_engine_ready(),
    reason="requires latex2mathml and the local Office MML2OMML transform",
)
def test_native_math_rename_propagates_into_logical_group_composite(
    tmp_path: Path,
) -> None:
    reference = tmp_path / "math-reference.png"
    Image.new("RGB", (180, 100), "white").save(reference)
    run = common.create_run(
        reference,
        case="semantic-group-math",
        cases_root=tmp_path / "examples",
        input_route="svg-seeded",
    )
    run.redraw_svg.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="180" height="100" '
        'viewBox="0 0 180 100"><g id="formula-cluster" data-role="formula">'
        '<circle id="formula-badge" cx="40" cy="50" r="20" fill="#DDEEFF"/>'
        '<text id="formula" x="30" y="57" font-size="18" font-style="italic">'
        '<tspan>z</tspan><tspan baseline-shift="super" font-size="11">τ</tspan>'
        '</text></g></svg>',
        encoding="utf-8",
    )
    convert(run)

    summary = upgrade(run)
    assert summary["injected"] == 1
    bindings = read_json(run.bindings_path)
    formula = next(
        row for row in bindings["bindings"] if row["element_id"] == "formula"
    )
    group = next(
        row
        for row in bindings["logical_group_bindings"]
        if row["element_id"] == "formula-cluster"
    )
    assert formula["shape_name"] == "math:001"
    assert group["backend_object_names"] == [
        next(
            row["shape_name"]
            for row in bindings["bindings"]
            if row["element_id"] == "formula-badge"
        ),
        "math:001",
    ]
    assert group["backend_object_identities"] == [
        {"shape_id": shape_id, "shape_name": shape_name}
        for shape_id, shape_name in zip(
            group["backend_object_ids"],
            group["backend_object_names"],
            strict=True,
        )
    ]
    assert group["readback_found"] is True
    assert refresh_bindings(run, host_saved_reopened=False)["bindings_complete"] is True
