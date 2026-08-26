from __future__ import annotations

from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Pt

from tools.core import common
from tools.core.contracts import read_json, write_json
from tools.pipeline.convert import convert
from tools.assets.reference_inventory import topology_contracts_sha256
from tools.regions.visual_contracts import evaluate_visual_contracts, strict_blockers


@pytest.fixture()
def run_factory(tmp_path: Path):
    def make(
        *,
        icon_width: int = 20,
        font_size: int = 20,
        icon_transform: str | None = None,
    ) -> common.Run:
        reference = tmp_path / "reference.png"
        Image.new("RGB", (160, 120), "white").save(reference)
        run = common.create_run(
            reference,
            case="visual-contract",
            cases_root=tmp_path / "examples",
            input_route="svg-seeded",
        )
        transform_attr = (
            f'transform="{icon_transform}" ' if icon_transform else ""
        )
        run.redraw_svg.write_text(
            '<svg xmlns="http://www.w3.org/2000/svg" width="160" height="120" '
            'viewBox="0 0 160 120">'
            f'<text id="label" x="10" y="30" data-visual-bbox="10 14 45 20" '
            f'font-family="Arial" font-size="{font_size}">Label</text>'
            f'<rect id="icon" x="100" y="10" width="{icon_width}" height="20" '
            f'{transform_attr}'
            'fill="#4A86E8"/>'
            '</svg>',
            encoding="utf-8",
        )
        regions = read_json(run.regions_path)
        regions["reference_inventory"] = {
            "schema_version": "1.0.0",
            "required": True,
            "status": "frozen",
            "reference_sha256": run.load_meta()["source_sha256"],
            "expected_counts": {
                "text": 1,
                "formula": 0,
                "arrow": 0,
                "icon": 1,
                "brace": 0,
                "plot": 0,
                "shape": 0,
            },
            "zero_count_authorizations": [],
            "objects": [
                {
                    "id": "label-object",
                    "kind": "text",
                    "bbox": [10, 14, 45, 20],
                    "element_ids": ["label"],
                    "typography": {
                        "exact_text": "Label",
                        "font_family": "Arial",
                        "font_size_px": 20,
                        "font_weight": 400,
                        "font_style": "normal",
                        "line_count": 1,
                        "alignment": "left",
                        "bbox_tolerance_px": 1,
                        "font_size_tolerance_px": 0.1,
                    },
                },
                {
                    "id": "icon-object",
                    "kind": "icon",
                    "bbox": [100, 10, 20, 20],
                    "element_ids": ["icon"],
                    "visual": {
                        "bbox_tolerance_px": 1,
                        "aspect_ratio_tolerance": 0.01,
                    },
                },
            ],
        }
        regions["visual_contracts"] = {
            "collision_tolerance_px": 0.25,
            "clearances": [
                {
                    "id": "label-icon-gap",
                    "a": "label-object",
                    "b": "icon-object",
                    "axis": "x",
                    "min_px": 40,
                }
            ],
        }
        write_json(run.regions_path, regions)
        convert(run)
        return run

    return make


def test_frozen_visual_inventory_passes_source_and_powerpoint(run_factory):
    run = run_factory()
    report = evaluate_visual_contracts(run)
    assert report["pass"] is True
    assert report["object_count"] == 2
    assert report["typography"][0]["backend"]["pass"] is True
    assert report["clearances"][0]["pass"] is True
    assert strict_blockers(report) == []


def test_icon_scale_drift_is_a_blocker_before_white_canvas_can_hide_it(run_factory):
    run = run_factory(icon_width=30)
    report = evaluate_visual_contracts(run)
    assert report["pass"] is False
    assert "visual-contract:V5:icon-object" in report["blockers"]
    assert "visual-contract:V7:icon-object" in report["blockers"]


def test_source_anisotropy_blocks_even_when_svg_and_powerpoint_bboxes_close(
    run_factory,
):
    run = run_factory(
        icon_width=10,
        icon_transform="matrix(2 0 0 1 -100 0)",
    )

    report = evaluate_visual_contracts(run)
    record = next(item for item in report["objects"] if item["id"] == "icon-object")
    transform = next(
        item
        for item in report["source_transform_anisotropy"]
        if item["id"] == "icon-object"
    )

    assert record["source_geometry_pass"] is True
    assert record["backend_pass"] is True
    assert record["source_transform_pass"] is False
    assert transform["offenders"][0]["anisotropy"] == 1.0
    assert "visual-contract:V38:icon-object" in report["blockers"]
    assert "visual-contract:V5:icon-object" not in report["blockers"]
    assert "visual-contract:V7:icon-object" not in report["blockers"]


def test_source_anisotropy_requires_explicit_hash_bound_allowance(run_factory):
    run = run_factory(
        icon_width=10,
        icon_transform="matrix(2 0 0 1 -100 0)",
    )
    payload = read_json(run.regions_path)
    visual = payload["reference_inventory"]["objects"][1]["visual"]
    visual["allow_source_anisotropic_scale"] = True
    visual["source_anisotropy_basis"] = "reference-measured-affine-distortion"
    write_json(run.regions_path, payload)

    report = evaluate_visual_contracts(run)
    transform = next(
        item
        for item in report["source_transform_anisotropy"]
        if item["id"] == "icon-object"
    )

    assert report["pass"] is True
    assert transform["policy"] == "explicitly_allowed"
    assert transform["pass"] is True


def test_rotation_does_not_count_as_source_anisotropy(run_factory):
    run = run_factory(icon_transform="matrix(0 1 -1 0 130 -90)")

    report = evaluate_visual_contracts(run)
    transform = next(
        item
        for item in report["source_transform_anisotropy"]
        if item["id"] == "icon-object"
    )

    assert report["pass"] is True
    assert transform["elements"][0]["anisotropy"] == 0.0


def test_arrow_bbox_is_delegated_to_the_physical_arrow_gate(run_factory):
    run = run_factory(icon_width=30)
    payload = read_json(run.regions_path)
    arrow = payload["reference_inventory"]["objects"][1]
    arrow["kind"] = "arrow"
    arrow["contract_refs"] = {
        "arrow_visual": {"contract_id": "arrow-object"}
    }
    write_json(run.regions_path, payload)

    report = evaluate_visual_contracts(run)
    record = next(item for item in report["objects"] if item["id"] == "icon-object")
    assert record["pass"] is True
    assert record["geometry_authority"] == "arrow_visual_physical_gate"


def test_font_size_drift_is_an_object_level_blocker(run_factory):
    run = run_factory(font_size=28)
    report = evaluate_visual_contracts(run)
    assert "visual-contract:V12:label-object" in report["blockers"]
    assert "visual-contract:V13:label-object" in report["blockers"]


def test_saved_powerpoint_tampering_breaks_scale_and_font_readback(run_factory):
    run = run_factory()
    presentation = Presentation(run.pptx_path)
    shapes = {shape.name: shape for shape in presentation.slides[0].shapes}
    shapes["af-icon-rect-01"].width += 8 * 9525
    shapes["af-label-text-01"].text_frame.paragraphs[0].runs[0].font.size = Pt(30)
    presentation.save(run.pptx_path)

    report = evaluate_visual_contracts(run)
    assert "visual-contract:V7:icon-object" in report["blockers"]
    assert "visual-contract:V13:label-object" in report["blockers"]


@pytest.fixture()
def topology_run_factory(tmp_path: Path):
    def make(*, mutation: str | None = None) -> common.Run:
        reference = tmp_path / "topology-reference.png"
        Image.new("RGB", (180, 120), "white").save(reference)
        run = common.create_run(
            reference,
            case="topology-contract",
            cases_root=tmp_path / "examples",
            input_route="svg-seeded",
        )
        molecule_bond_02 = (
            ""
            if mutation == "molecule-missing-bond"
            else '<line id="molecule-bond-02" x1="20" y1="20" x2="30" y2="20" '
            'data-source-id="molecule-atom-02" data-target-id="molecule-atom-03" '
            'data-topology-relation="bond" stroke="#777"/>'
        )
        dna_strand_b = (
            ""
            if mutation == "dna-missing-strand"
            else '<line id="dna-strand-b" x1="90" y1="15" x2="90" y2="45" '
            'stroke="#446"/>'
        )
        dna_extra_rung = (
            '<line id="dna-rung-03" x1="70" y1="28" x2="90" y2="28" '
            'stroke="#88a"/>'
            if mutation == "dna-extra-rung"
            else ""
        )
        plot_body_02 = (
            ""
            if mutation == "plot-missing-body"
            else '<rect id="plot-body-02" x="146" y="20" width="8" height="10" '
            'data-pair-with="plot-wick-02" fill="#fff" stroke="#222"/>'
        )
        plot_extra_wick = (
            '<line id="plot-wick-03" x1="140" y1="18" x2="140" y2="32" '
            'stroke="#222"/>'
            if mutation == "plot-extra-wick"
            else ""
        )
        run.redraw_svg.write_text(
            '<svg xmlns="http://www.w3.org/2000/svg" width="180" height="120" '
            'viewBox="0 0 180 120">'
            '<g id="molecule">'
            '<line id="molecule-bond-01" x1="10" y1="20" x2="20" y2="20" '
            'data-source-id="molecule-atom-01" data-target-id="molecule-atom-02" '
            'data-topology-relation="bond" stroke="#777"/>'
            f'{molecule_bond_02}'
            '<circle id="molecule-atom-01" cx="10" cy="20" r="3" fill="#d66"/>'
            '<circle id="molecule-atom-02" cx="20" cy="20" r="3" fill="#6a6"/>'
            '<circle id="molecule-atom-03" cx="30" cy="20" r="3" fill="#66d"/>'
            '</g>'
            '<g id="dna">'
            '<line id="dna-strand-a" x1="70" y1="15" x2="70" y2="45" '
            'stroke="#446"/>'
            f'{dna_strand_b}'
            '<line id="dna-rung-01" x1="70" y1="20" x2="90" y2="20" '
            'data-source-id="dna-strand-a" data-target-id="dna-strand-b" '
            'data-topology-relation="bridge" stroke="#88a"/>'
            '<line id="dna-rung-02" x1="70" y1="35" x2="90" y2="35" '
            'data-source-id="dna-strand-a" data-target-id="dna-strand-b" '
            'data-topology-relation="bridge" stroke="#88a"/>'
            f'{dna_extra_rung}'
            '</g>'
            '<g id="plot">'
            '<line id="plot-wick-01" x1="130" y1="15" x2="130" y2="35" '
            'data-pair-with="plot-body-01" stroke="#222"/>'
            '<rect id="plot-body-01" x="126" y="20" width="8" height="10" '
            'data-pair-with="plot-wick-01" fill="#fff" stroke="#222"/>'
            '<line id="plot-wick-02" x1="150" y1="15" x2="150" y2="35" '
            'data-pair-with="plot-body-02" stroke="#222"/>'
            f'{plot_body_02}{plot_extra_wick}'
            '</g>'
            '</svg>',
            encoding="utf-8",
        )
        inventory = {
            "schema_version": "1.0.0",
            "required": True,
            "status": "frozen",
            "reference_sha256": run.load_meta()["source_sha256"],
            "expected_counts": {
                "text": 0,
                "formula": 0,
                "arrow": 0,
                "icon": 2,
                "brace": 0,
                "plot": 1,
                "shape": 0,
            },
            "zero_count_authorizations": [],
            "objects": [
                {
                    "id": "molecule",
                    "kind": "icon",
                    "bbox": [7, 17, 26, 6],
                    "element_ids": [
                        "molecule-bond-01",
                        "molecule-bond-02",
                        "molecule-atom-01",
                        "molecule-atom-02",
                        "molecule-atom-03",
                    ],
                    "visual": {
                        "bbox_tolerance_px": 1,
                        "aspect_ratio_tolerance": 0.1,
                    },
                    "topology_contract": {
                        "role_counts": {"atom": 3, "bond": 2},
                        "role_mapping": {
                            "molecule-atom-01": "atom",
                            "molecule-atom-02": "atom",
                            "molecule-atom-03": "atom",
                            "molecule-bond-01": "bond",
                            "molecule-bond-02": "bond",
                        },
                        "required_pairs": [],
                        "relations": [
                            {
                                "id": "molecule-bond-01",
                                "source_id": "molecule-atom-01",
                                "target_id": "molecule-atom-02",
                                "relation": "bond",
                            },
                            {
                                "id": "molecule-bond-02",
                                "source_id": "molecule-atom-02",
                                "target_id": "molecule-atom-03",
                                "relation": "bond",
                            },
                        ],
                        "component_count": 1,
                    },
                },
                {
                    "id": "dna",
                    "kind": "icon",
                    "bbox": [70, 15, 20, 30],
                    "element_ids": [
                        "dna-strand-a",
                        "dna-strand-b",
                        "dna-rung-01",
                        "dna-rung-02",
                    ],
                    "visual": {
                        "bbox_tolerance_px": 1,
                        "aspect_ratio_tolerance": 0.1,
                    },
                    "topology_contract": {
                        "role_counts": {
                            "strand": {
                                "count": 2,
                                "element_id_pattern": "dna-strand-[ab]",
                            },
                            "rung": {
                                "count": 2,
                                "element_id_pattern": "dna-rung-[0-9]{2}",
                            },
                        },
                        "required_pairs": [],
                        "required_relations": [
                            {
                                "id": "dna-rung-01",
                                "source": "dna-strand-a",
                                "target": "dna-strand-b",
                                "kind": "bridge",
                            },
                            {
                                "id": "dna-rung-02",
                                "source": "dna-strand-a",
                                "target": "dna-strand-b",
                                "kind": "bridge",
                            },
                        ],
                        "component_count": 1,
                    },
                },
                {
                    "id": "plot",
                    "kind": "plot",
                    "bbox": [126, 15, 28, 20],
                    "element_ids": [
                        "plot-wick-01",
                        "plot-body-01",
                        "plot-wick-02",
                        "plot-body-02",
                    ],
                    "visual": {
                        "bbox_tolerance_px": 1,
                        "aspect_ratio_tolerance": 0.1,
                    },
                    "topology_contract": {
                        "role_counts": {"wick": 2, "body": 2},
                        "role_patterns": {
                            "wick": "plot-wick-[0-9]{2}",
                            "body": "plot-body-[0-9]{2}",
                        },
                        "required_pairs": [
                            ["plot-wick-01", "plot-body-01"],
                            ["plot-wick-02", "plot-body-02"],
                        ],
                        "relations": [],
                        "component_count": 2,
                    },
                },
            ],
        }
        regions = read_json(run.regions_path)
        regions["reference_inventory"] = inventory
        write_json(run.regions_path, regions)
        convert(run)
        return run

    return make


def test_topology_contract_passes_for_svg_and_powerpoint(topology_run_factory):
    run = topology_run_factory()

    report = evaluate_visual_contracts(run)

    assert report["pass"] is True
    assert report["topology_contract_count"] == 3
    assert report["topology_contracts_sha256"] == topology_contracts_sha256(
        read_json(run.regions_path)["reference_inventory"]
    )
    assert all(row["source"]["pass"] for row in report["topology"])
    assert all(row["backend"]["pass"] for row in report["topology"])


def test_topology_contract_rejects_missing_source_relation_metadata(
    topology_run_factory,
):
    run = topology_run_factory()
    source = run.redraw_svg.read_text(encoding="utf-8")
    metadata = (
        'data-source-id="molecule-atom-01" '
        'data-target-id="molecule-atom-02" '
        'data-topology-relation="bond" '
    )
    assert metadata in source
    run.redraw_svg.write_text(source.replace(metadata, "", 1), encoding="utf-8")

    report = evaluate_visual_contracts(run)
    topology = next(row for row in report["topology"] if row["id"] == "molecule")

    assert topology["source"]["pass"] is False
    assert topology["source"]["missing_relation_metadata"] == ["molecule-bond-01"]
    assert "visual-contract:V34:molecule" in report["blockers"]


def test_topology_contract_rejects_missing_source_pair_metadata(
    topology_run_factory,
):
    run = topology_run_factory()
    source = run.redraw_svg.read_text(encoding="utf-8")
    metadata = 'data-pair-with="plot-body-01" '
    assert metadata in source
    run.redraw_svg.write_text(source.replace(metadata, "", 1), encoding="utf-8")

    report = evaluate_visual_contracts(run)
    topology = next(row for row in report["topology"] if row["id"] == "plot")

    assert topology["source"]["pass"] is False
    assert topology["source"]["missing_pair_metadata"] == ["pair-1"]
    assert "visual-contract:V34:plot" in report["blockers"]


def test_topology_contract_rejects_reversed_source_relation_mutation(
    topology_run_factory,
):
    run = topology_run_factory()
    source = run.redraw_svg.read_text(encoding="utf-8")
    expected = (
        'data-source-id="molecule-atom-01" '
        'data-target-id="molecule-atom-02" '
    )
    reversed_relation = (
        'data-source-id="molecule-atom-02" '
        'data-target-id="molecule-atom-01" '
    )
    assert expected in source
    run.redraw_svg.write_text(
        source.replace(expected, reversed_relation, 1),
        encoding="utf-8",
    )

    report = evaluate_visual_contracts(run)
    topology = next(row for row in report["topology"] if row["id"] == "molecule")

    assert topology["source"]["pass"] is False
    assert topology["source"]["malformed_relations"] == ["molecule-bond-01"]
    assert "visual-contract:V34:molecule" in report["blockers"]


@pytest.mark.parametrize(
    ("mutation", "object_id"),
    [
        ("dna-missing-strand", "dna"),
        ("dna-extra-rung", "dna"),
        ("molecule-missing-bond", "molecule"),
        ("plot-extra-wick", "plot"),
        ("plot-missing-body", "plot"),
    ],
)
def test_topology_contract_rejects_internal_count_drift(
    topology_run_factory,
    mutation: str,
    object_id: str,
):
    run = topology_run_factory(mutation=mutation)

    report = evaluate_visual_contracts(run)
    topology = next(row for row in report["topology"] if row["id"] == object_id)

    assert report["pass"] is False
    assert topology["source"]["pass"] is False
    assert topology["backend"]["pass"] is False
    assert f"visual-contract:V31:{object_id}" in report["blockers"]
    assert f"visual-contract:V32:{object_id}" in report["blockers"]


def test_topology_contract_rejects_missing_powerpoint_readback(topology_run_factory):
    run = topology_run_factory()
    bindings = read_json(run.bindings_path)
    binding = next(
        row
        for row in bindings["bindings"]
        if row["element_id"] == "dna-strand-b"
    )
    binding["readback_found"] = False
    write_json(run.bindings_path, bindings)

    report = evaluate_visual_contracts(run)
    topology = next(row for row in report["topology"] if row["id"] == "dna")

    assert topology["source"]["pass"] is True
    assert topology["backend"]["pass"] is False
    assert "visual-contract:V33:dna" in report["blockers"]


def test_topology_contract_rejects_extra_powerpoint_role_binding(topology_run_factory):
    run = topology_run_factory()
    bindings = read_json(run.bindings_path)
    binding = next(
        row
        for row in bindings["bindings"]
        if row["element_id"] == "dna-rung-02"
    )
    bindings["bindings"].append({**binding, "element_id": "dna-rung-03"})
    write_json(run.bindings_path, bindings)

    report = evaluate_visual_contracts(run)
    topology = next(row for row in report["topology"] if row["id"] == "dna")

    assert topology["source"]["pass"] is True
    assert topology["backend"]["pass"] is False
    assert topology["backend"]["extra_ids"]["rung"] == ["dna-rung-03"]
    assert "visual-contract:V32:dna" in report["blockers"]
