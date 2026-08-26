from __future__ import annotations

from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Pt

from tools.core import common
from tools.core.contracts import read_json, write_json
from tools.pipeline.convert import convert
from tools.pipeline.layout import audit_layout, strict_blockers
from tools.arrows.pptx_arrows import refresh_bindings, write_arrow_reports


@pytest.fixture()
def run_factory(tmp_path: Path):
    def make(body: str) -> common.Run:
        reference = tmp_path / "reference.png"
        Image.new("RGB", (200, 120), "white").save(reference)
        run = common.create_run(
            reference,
            case="layout",
            cases_root=tmp_path / "examples",
            input_route="svg-seeded",
        )
        run.redraw_svg.write_text(
            '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="120" '
            f'viewBox="0 0 200 120">{body}</svg>',
            encoding="utf-8",
        )
        return run

    return make


def test_container_annotation_clips_selection_box_and_survives_readback(run_factory):
    run = run_factory(
        '<rect id="box" x="10" y="10" width="100" height="40" fill="#eeeeee"/>'
        '<text id="label" data-layout-container="box" data-layout-padding="2" '
        'x="18" y="34" font-size="16">contained</text>'
    )
    summary = convert(run)
    assert summary["layout_pass"] is True
    report = audit_layout(run)
    assert report["pass"] is True

    presentation = Presentation(run.pptx_path)
    shapes = {shape.name: shape for shape in presentation.slides[0].shapes}
    box = shapes["af-box-rect-01"]
    label = shapes["af-label-text-01"]
    emu_per_px = 9525
    assert label.left / emu_per_px >= box.left / emu_per_px + 2 - 0.01
    assert (label.top + label.height) / emu_per_px <= (
        (box.top + box.height) / emu_per_px - 2 + 0.01
    )
    scene = read_json(run.scene_path)
    label_record = next(item for item in scene["elements"] if item["id"] == "label")
    assert label_record["layout"]["container_id"] == "box"


def test_repeat_contract_detects_irregular_source_and_backend_spacing(run_factory):
    repeated = "".join(
        f'<circle id="dot-{index}" data-repeat-group="dots" data-repeat-axis="vertical" '
        f'data-repeat-order="{index}" cx="50" cy="{cy}" r="8" fill="#cc6600"/>'
        for index, cy in enumerate((20, 50, 70), start=1)
    )
    run = run_factory(repeated)
    convert(run)
    report = audit_layout(run)
    assert report["pass"] is False
    spacing_findings = [item for item in report["findings"] if item["code"] == "L9"]
    assert {item["stage"] for item in spacing_findings} == {"source", "backend"}
    assert any(blocker == "layout:L9:source:dots" for blocker in strict_blockers(report))


def test_repeat_contract_allows_one_pixel_reference_quantization(run_factory):
    repeated = "".join(
        f'<circle id="dot-{index}" data-repeat-group="dots" data-repeat-axis="vertical" '
        f'data-repeat-order="{index}" cx="50" cy="{cy}" r="8" fill="#cc6600"/>'
        for index, cy in enumerate((20, 38, 55), start=1)
    )
    run = run_factory(repeated)
    convert(run)
    report = audit_layout(run)
    group = report["repeat_groups"][0]
    assert group["source"]["steps_px"] == [18.0, 17.0]
    assert group["backend"]["steps_px"] == [18.0, 17.0]
    assert report["pass"] is True


def test_layout_annotation_without_stable_id_is_a_hard_finding(run_factory):
    run = run_factory(
        '<rect id="box" x="10" y="10" width="100" height="40" fill="#eeeeee"/>'
        '<text data-layout-container="box" x="18" y="34" font-size="16">anonymous</text>'
    )
    convert(run)
    report = audit_layout(run)
    assert any(item["code"] == "L1" for item in report["findings"])


def test_saved_powerpoint_object_must_remain_inside_slide_canvas(run_factory):
    run = run_factory('<rect id="box" x="10" y="10" width="40" height="30" fill="#eeeeee"/>')
    convert(run)
    presentation = Presentation(run.pptx_path)
    presentation.slides[0].shapes[0].left = -10 * 9525
    presentation.save(run.pptx_path)

    report = audit_layout(run)
    finding = next(item for item in report["findings"] if item["code"] == "L10")
    assert finding["stage"] == "backend"
    assert finding["target"] == "box"
    assert finding["metrics"]["overflow_px"]["left"] == pytest.approx(10.0)
    assert "layout:L10:backend:box" in strict_blockers(report)


def test_peer_size_contract_fails_closed_for_unequal_semantic_peers(run_factory):
    run = run_factory(
        '<rect id="a" data-peer-size-group="peers" x="10" y="10" '
        'width="40" height="30" fill="#eeeeee"/>'
        '<rect id="b" data-peer-size-group="peers" x="70" y="10" '
        'width="44" height="30" fill="#eeeeee"/>'
    )
    convert(run)
    report = audit_layout(run)
    findings = [item for item in report["findings"] if item["code"] in {"L11", "L12"}]
    assert {item["stage"] for item in findings} == {"source", "backend"}
    assert any(blocker == "layout:L12:source:peers" for blocker in strict_blockers(report))


def test_peer_size_group_cannot_relax_tolerance_from_one_member(run_factory):
    run = run_factory(
        '<rect id="a" data-peer-size-group="peers" data-peer-size-tolerance="0.25" '
        'x="10" y="10" width="40" height="30" fill="#eeeeee"/>'
        '<rect id="b" data-peer-size-group="peers" data-peer-size-tolerance="100" '
        'x="70" y="10" width="40" height="30" fill="#eeeeee"/>'
    )
    convert(run)
    report = audit_layout(run)
    finding = next(item for item in report["findings"] if item["code"] == "L11")
    assert finding["stage"] == "source"
    assert finding["metrics"]["member_tolerances_px"] == [0.25, 100.0]


def test_group_level_xy_repeat_contract_uses_descendant_union_in_both_artifacts(
    run_factory,
):
    run = run_factory(
        '<g id="motif-a" data-repeat-group="motifs" data-repeat-axis="xy" '
        'data-repeat-order="1">'
        '<circle id="motif-a-dot" cx="20" cy="20" r="5" fill="#cc6600"/>'
        '<rect id="motif-a-bar" x="28" y="17" width="12" height="6" fill="#444444"/>'
        '</g>'
        '<g id="motif-b" data-repeat-group="motifs" data-repeat-axis="xy" '
        'data-repeat-order="2" transform="translate(70 50)">'
        '<circle id="motif-b-dot" cx="20" cy="20" r="5" fill="#cc6600"/>'
        '<rect id="motif-b-bar" x="28" y="17" width="12" height="6" fill="#444444"/>'
        '</g>'
    )
    convert(run)
    report = audit_layout(run)

    assert report["pass"] is True
    assert report["annotation_coverage"]["logical_svg_groups"] == 2
    group = report["repeat_groups"][0]
    assert group["axis"] == "xy"
    assert group["source"]["members"] == ["motif-a", "motif-b"]
    assert group["source"]["widths_px"] == [25.0, 25.0]
    assert group["backend"]["widths_px"] == [25.0, 25.0]
    assert group["source"]["steps_xy_px"] == [[70.0, 50.0]]


def test_group_level_container_contract_audits_union_and_padding(run_factory):
    run = run_factory(
        '<rect id="panel" x="10" y="10" width="80" height="60" fill="#eeeeee"/>'
        '<g id="overflowing" data-layout-container="panel" data-layout-padding="4">'
        '<rect id="overflowing-body" x="80" y="30" width="15" height="10" '
        'fill="#cc6600"/>'
        '</g>'
    )
    convert(run)
    report = audit_layout(run)

    row = next(item for item in report["containment"] if item["element"] == "overflowing")
    assert row["source"]["overflow_px"]["right"] == pytest.approx(9.0)
    assert row["backend"]["overflow_px"]["right"] == pytest.approx(9.0)
    assert {item["code"] for item in report["findings"]} >= {"L3", "L5"}


def test_group_level_peer_contract_checks_axis_order_and_spacing(run_factory):
    peers = "".join(
        '<g id="peer-{index}" data-peer-group="row" data-peer-axis="x" '
        'data-peer-order="{order}" transform="translate({offset} 0)">'
        '<rect id="peer-{index}-body" x="10" y="20" width="10" height="10" '
        'fill="#336699"/></g>'.format(index=index, order=index, offset=offset)
        for index, offset in enumerate((0, 30, 60), start=1)
    )
    run = run_factory(peers)
    convert(run)
    report = audit_layout(run)

    assert report["pass"] is True
    group = report["peer_groups"][0]
    assert group["axis"] == "horizontal"
    assert group["source"]["steps_px"] == [30.0, 30.0]
    assert group["backend"]["steps_px"] == [30.0, 30.0]


def test_incomplete_group_annotations_are_strict_blockers(run_factory):
    run = run_factory(
        '<g id="peer-a" data-peer-group="row" data-peer-axis="x">'
        '<rect id="peer-a-body" x="10" y="20" width="10" height="10"/>'
        '</g>'
        '<g id="peer-b" data-peer-group="row" data-peer-axis="x">'
        '<rect id="peer-b-body" x="40" y="20" width="10" height="10"/>'
        '</g>'
        '<rect id="orphan-repeat" data-repeat-axis="x" data-repeat-order="1" '
        'x="80" y="20" width="10" height="10"/>'
    )
    convert(run)
    report = audit_layout(run)

    assert any(item["code"] == "L17" and item["target"] == "row" for item in report["findings"])
    assert any(item["code"] == "L6" and item["target"] == "orphan-repeat" for item in report["findings"])
    blockers = strict_blockers(report)
    assert "layout:L17:source:row" in blockers
    assert "layout:L6:source:orphan-repeat" in blockers


def test_layout_rejects_split_shape_id_and_name_identity(run_factory):
    run = run_factory(
        '<rect id="a" data-peer-size-group="peers" x="10" y="10" '
        'width="40" height="30" fill="#eeeeee"/>'
        '<rect id="b" data-peer-size-group="peers" x="70" y="10" '
        'width="40" height="30" fill="#eeeeee"/>'
    )
    convert(run)
    bindings = read_json(run.bindings_path)
    rows = {item["element_id"]: item for item in bindings["bindings"]}
    rows["a"]["shape_name"] = rows["b"]["shape_name"]
    write_json(run.bindings_path, bindings)

    report = audit_layout(run)
    assert "a" in report["missing_annotated_backend_objects"]
    finding = next(
        item
        for item in report["findings"]
        if item["code"] == "L11" and item["stage"] == "backend"
    )
    assert "a" in finding["message"]


def test_gap_arrow_contract_derives_endpoints_from_peer_boundaries(run_factory):
    run = run_factory(
        '<defs><marker id="head" markerWidth="8" markerHeight="8" refX="8" '
        'refY="4" orient="auto" markerUnits="userSpaceOnUse">'
        '<path d="M0 0 L8 4 L0 8 Z" fill="#111111"/></marker></defs>'
        '<rect id="left" x="10" y="20" width="40" height="70" fill="#eeeeee"/>'
        '<rect id="right" x="80" y="20" width="40" height="70" fill="#eeeeee"/>'
        '<line id="gap" x1="52" y1="55" x2="78" y2="55" stroke="#111111" '
        'stroke-width="2" marker-end="url(#head)" data-source-id="left" '
        'data-target-id="right" data-source-gap="2" data-target-gap="2" '
        'data-attach="false" data-gap-source-id="left" data-gap-target-id="right" '
        'data-gap-axis="horizontal" data-gap-start-inset="2" data-gap-end-inset="2" '
        'data-gap-cross-position="55" data-gap-tolerance="0.25"/>'
    )
    summary = convert(run)
    assert summary["layout_pass"] is True
    report = audit_layout(run)
    row = report["gap_arrows"][0]
    assert row["source"]["start_delta"]["distance"] == 0
    assert row["source"]["end_delta"]["distance"] == 0
    assert row["backend"]["path_and_head_readback_pass"] is True


def test_gap_block_arrow_uses_declared_semantic_centerline(run_factory):
    run = run_factory(
        '<rect id="left" x="10" y="20" width="40" height="70" fill="#eeeeee"/>'
        '<rect id="right" x="80" y="20" width="40" height="70" fill="#eeeeee"/>'
        '<path id="gap" d="M52 50 H68 V46 L78 55 L68 64 V60 H52 Z" '
        'fill="#cc6600" stroke="#663300" data-role="arrow" '
        'data-arrow-representation="block_arrow" '
        'data-arrow-centerline="M52 55 L78 55" data-arrow-body-width="10" '
        'data-start-head-type="none" data-end-head-type="custom" '
        'data-arrow-routing="fixed" data-arrow-topology="declared" '
        'data-source-id="left" data-target-id="right" '
        'data-gap-source-id="left" data-gap-target-id="right" '
        'data-gap-axis="horizontal" data-gap-start-inset="2" '
        'data-gap-end-inset="2" data-gap-cross-position="55" '
        'data-gap-tolerance="0.25"/>'
    )

    summary = convert(run)
    assert summary["layout_pass"] is True
    row = audit_layout(run)["gap_arrows"][0]
    assert row["source"]["actual_start"] == [52.0, 55.0]
    assert row["source"]["actual_end"] == [78.0, 55.0]
    assert row["backend"]["path_and_head_readback_pass"] is True
    assert row["backend"]["endpoint_evidence"].startswith(
        "saved-pptx-cNvPr-description"
    )


def test_gap_arrow_contract_rejects_stale_fixed_geometry(run_factory):
    run = run_factory(
        '<defs><marker id="head" markerWidth="8" markerHeight="8" refX="8" '
        'refY="4" orient="auto" markerUnits="userSpaceOnUse">'
        '<path d="M0 0 L8 4 L0 8 Z" fill="#111111"/></marker></defs>'
        '<rect id="left" x="10" y="20" width="40" height="70" fill="#eeeeee"/>'
        '<rect id="right" x="80" y="20" width="40" height="70" fill="#eeeeee"/>'
        '<line id="gap" x1="56" y1="55" x2="78" y2="55" stroke="#111111" '
        'stroke-width="2" marker-end="url(#head)" data-source-id="left" '
        'data-target-id="right" data-source-gap="6" data-target-gap="2" '
        'data-attach="false" data-gap-source-id="left" data-gap-target-id="right" '
        'data-gap-axis="horizontal" data-gap-start-inset="2" data-gap-end-inset="2" '
        'data-gap-cross-position="55" data-gap-tolerance="0.25"/>'
    )
    convert(run)
    report = audit_layout(run)
    finding = next(item for item in report["findings"] if item["code"] == "L13")
    assert finding["stage"] == "source"
    assert finding["target"] == "gap"


def test_gap_arrow_cross_position_must_be_inside_peer_overlap(run_factory):
    run = run_factory(
        '<defs><marker id="head" markerWidth="8" markerHeight="8" refX="8" '
        'refY="4" orient="auto" markerUnits="userSpaceOnUse">'
        '<path d="M0 0 L8 4 L0 8 Z" fill="#111111"/></marker></defs>'
        '<rect id="left" x="10" y="20" width="40" height="30" fill="#eeeeee"/>'
        '<rect id="right" x="80" y="30" width="40" height="30" fill="#eeeeee"/>'
        '<line id="gap" x1="52" y1="70" x2="78" y2="70" stroke="#111111" '
        'stroke-width="2" marker-end="url(#head)" data-source-id="left" '
        'data-target-id="right" data-source-gap="2" data-target-gap="2" '
        'data-attach="false" data-gap-source-id="left" data-gap-target-id="right" '
        'data-gap-axis="horizontal" data-gap-start-inset="2" data-gap-end-inset="2" '
        'data-gap-cross-position="70" data-gap-tolerance="0.25"/>'
    )
    convert(run)
    report = audit_layout(run)

    finding = next(item for item in report["findings"] if item["code"] == "L13")
    assert finding["stage"] == "source"
    assert "cross-axis overlap" in finding["message"]
    assert report["pass"] is False


def test_gap_arrow_tolerance_cannot_exceed_safe_cap(run_factory):
    run = run_factory(
        '<defs><marker id="head" markerWidth="8" markerHeight="8" refX="8" '
        'refY="4" orient="auto" markerUnits="userSpaceOnUse">'
        '<path d="M0 0 L8 4 L0 8 Z" fill="#111111"/></marker></defs>'
        '<rect id="left" x="10" y="20" width="40" height="70" fill="#eeeeee"/>'
        '<rect id="right" x="80" y="20" width="40" height="70" fill="#eeeeee"/>'
        '<line id="gap" x1="60" y1="55" x2="70" y2="55" stroke="#111111" '
        'stroke-width="2" marker-end="url(#head)" data-source-id="left" '
        'data-target-id="right" data-source-gap="10" data-target-gap="10" '
        'data-attach="false" data-gap-source-id="left" data-gap-target-id="right" '
        'data-gap-axis="horizontal" data-gap-start-inset="2" data-gap-end-inset="2" '
        'data-gap-cross-position="55" data-gap-tolerance="100"/>'
    )
    convert(run)
    report = audit_layout(run)

    finding = next(item for item in report["findings"] if item["code"] == "L13")
    assert finding["stage"] == "source"
    assert finding["metrics"]["declared_tolerance_px"] == 100
    assert "<= 1 px" in finding["message"]
    assert report["pass"] is False


def test_gap_arrow_backend_recomputes_against_moved_peer(run_factory):
    run = run_factory(
        '<defs><marker id="head" markerWidth="8" markerHeight="8" refX="8" '
        'refY="4" orient="auto" markerUnits="userSpaceOnUse">'
        '<path d="M0 0 L8 4 L0 8 Z" fill="#111111"/></marker></defs>'
        '<rect id="left" x="10" y="20" width="40" height="70" fill="#eeeeee"/>'
        '<rect id="right" x="80" y="20" width="40" height="70" fill="#eeeeee"/>'
        '<line id="gap" x1="52" y1="55" x2="78" y2="55" stroke="#111111" '
        'stroke-width="2" marker-end="url(#head)" data-source-id="left" '
        'data-target-id="right" data-source-site="3" data-target-site="1" '
        'data-attach="false" data-gap-source-id="left" data-gap-target-id="right" '
        'data-gap-axis="horizontal" data-gap-start-inset="2" data-gap-end-inset="2" '
        'data-gap-cross-position="55" data-gap-tolerance="0.25"/>'
    )
    convert(run)
    presentation = Presentation(run.pptx_path)
    right = next(
        shape
        for shape in presentation.slides[0].shapes
        if shape.name == "af-right-rect-01"
    )
    right.left += 10 * 9525
    presentation.save(run.pptx_path)
    refresh_bindings(run, host_saved_reopened=True)
    write_arrow_reports(run)

    report = audit_layout(run)
    finding = next(item for item in report["findings"] if item["code"] == "L14")
    assert finding["stage"] == "backend"
    row = report["gap_arrows"][0]
    assert row["backend"]["endpoint_metrics"]["end_delta"]["distance"] == 10.0


def test_relative_z_order_contract_checks_svg_and_powerpoint(run_factory):
    run = run_factory(
        '<line id="edge" data-z-above="panel" x1="10" y1="50" x2="90" y2="50" '
        'stroke="#111111" stroke-width="2"/>'
        '<rect id="panel" x="40" y="20" width="80" height="70" fill="#eeeeee"/>'
    )
    convert(run)
    report = audit_layout(run)
    findings = [item for item in report["findings"] if item["code"] == "L15"]
    assert {item["stage"] for item in findings} == {"source", "backend"}


def test_explicit_rotated_and_stacked_text_flows_survive_readback(run_factory):
    run = run_factory(
        '<rect id="slot" x="10" y="10" width="37" height="100" fill="#eeeeee"/>'
        '<rect id="stack-slot" x="80" y="10" width="40" height="100" fill="#eeeeee"/>'
        '<text id="rotated" x="29" y="60" font-size="18" text-anchor="middle" '
        'transform="rotate(90 29 60)" data-text-flow="rotated-word" '
        'data-text-container="slot" data-text-box-height="36">Sem</text>'
        '<text id="stacked" x="100" y="40" font-size="12" text-anchor="middle" '
        'data-text-flow="stacked-characters" data-text-container="stack-slot" '
        'data-text-stack-step="14">Dyn</text>'
    )
    summary = convert(run)
    assert summary["layout_pass"] is True
    flows = {item["element"]: item for item in audit_layout(run)["text_flow_contracts"]}
    assert flows["rotated"]["backend"]["paragraphs"] == ["Sem"]
    assert flows["rotated"]["backend"]["rotation_deg"] == 90
    assert flows["stacked"]["backend"]["paragraphs"] == ["D", "y", "n"]
    assert flows["stacked"]["backend"]["rotation_deg"] == 0


def test_multiline_text_flow_preserves_lines_and_container_clearance(run_factory):
    run = run_factory(
        '<rect id="slot" x="20" y="10" width="160" height="100" fill="#eeeeee"/>'
        '<text id="multiline" x="100" y="35" font-size="18" text-anchor="middle" '
        'data-text-flow="multiline" data-layout-container="slot">'
        '<tspan x="100" y="35" font-weight="700">Pareto-Conditioned</tspan>'
        '<tspan x="100" y="60">Diffusion Model</tspan>'
        '</text>'
    )
    summary = convert(run)
    assert summary["layout_pass"] is True
    report = audit_layout(run)
    flow = next(item for item in report["text_flow_contracts"] if item["element"] == "multiline")
    assert flow["backend"]["paragraphs"] == [
        "Pareto-Conditioned",
        "Diffusion Model",
    ]
    assert flow["backend"]["contract_pass"] is True


def test_vertical_text_requires_container_and_stack_spacing_readback(run_factory):
    missing_container = run_factory(
        '<text id="stacked" x="100" y="40" font-size="12" text-anchor="middle" '
        'data-text-flow="stacked-characters" data-text-stack-step="14">Dyn</text>'
    )
    with pytest.raises(SystemExit, match="data-text-container"):
        convert(missing_container)

    run = missing_container
    run.redraw_svg.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="120" '
        'viewBox="0 0 200 120"><rect id="slot" x="80" y="10" width="40" '
        'height="100" fill="#eeeeee"/><text id="stacked" x="100" y="40" '
        'font-size="12" text-anchor="middle" data-text-flow="stacked-characters" '
        'data-text-container="slot" data-text-stack-step="14">Dyn</text></svg>',
        encoding="utf-8",
    )
    convert(run)
    presentation = Presentation(run.pptx_path)
    stacked = next(
        shape
        for shape in presentation.slides[0].shapes
        if shape.name == "af-stacked-text-01"
    )
    for paragraph in stacked.text_frame.paragraphs:
        paragraph.line_spacing = Pt(30)
    presentation.save(run.pptx_path)
    refresh_bindings(run, host_saved_reopened=True)
    report = audit_layout(run)
    finding = next(item for item in report["findings"] if item["code"] == "L16")
    assert finding["stage"] == "backend"
